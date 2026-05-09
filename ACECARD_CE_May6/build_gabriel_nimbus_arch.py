"""Infrastructure Architecture — Gabriel Nimbus Deployment
Single-page streaming pipeline. Same look-and-feel as v12 main deck."""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

OUT = os.path.join(os.path.dirname(__file__), "Gabriel_Nimbus_Architecture_v2.pptx")

NAVY   = RGBColor(0x0D, 0x1B, 0x2A)
BLUE   = RGBColor(0x1B, 0x4F, 0x72)
TEAL   = RGBColor(0x0E, 0x6B, 0x8A)
GOLD   = RGBColor(0xB7, 0x95, 0x0B)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
INK    = RGBColor(0x2C, 0x3E, 0x50)
LGRAY  = RGBColor(0xF7, 0xF8, 0xFA)
MGRAY  = RGBColor(0xE8, 0xEB, 0xEF)
DGRAY  = RGBColor(0x6C, 0x75, 0x7D)
LBLUE  = RGBColor(0xEB, 0xF5, 0xFB)
LTEAL  = RGBColor(0xE8, 0xF4, 0xF8)
LGOLD  = RGBColor(0xFE, 0xF9, 0xE7)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
slide_count = [0]


def blank():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    slide_count[0] += 1
    return s

def brand_bar(s):
    rect(s, Inches(0), Inches(7.15), W, Inches(0.35), fill=NAVY)
    txt(s, Inches(0.5), Inches(7.17), Inches(6), Inches(0.3),
        '22nd Century Technologies  |  PROPRIETARY  |  ravindra.shukla@tscti.com', 8, RGBColor(0x8A, 0x9B, 0xAE))
    txt(s, Inches(11.5), Inches(7.17), Inches(1.5), Inches(0.3),
        str(slide_count[0]), 8, RGBColor(0x8A, 0x9B, 0xAE), al=PP_ALIGN.RIGHT)

def rect(s, x, y, w, h, fill=None, border=None):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else: sh.fill.background()
    if border: sh.line.color.rgb = border; sh.line.width = Pt(1)
    else: sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def rrect(s, x, y, w, h, fill=None, border=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else: sh.fill.background()
    if border: sh.line.color.rgb = border; sh.line.width = Pt(1)
    else: sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def hline(s, x, y, w, color=MGRAY, thick=Pt(1)):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, thick)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background(); sh.shadow.inherit = False

def txt(s, x, y, w, h, text, sz=13, c=INK, b=False, al=PP_ALIGN.LEFT, fn='Segoe UI'):
    tf = s.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(sz); p.font.color.rgb = c
    p.font.bold = b; p.font.name = fn; p.alignment = al
    return tf

def add_p(tf, text, sz=13, c=INK, b=False, sp=Pt(8)):
    p = tf.add_paragraph()
    p.text = text; p.font.size = Pt(sz); p.font.color.rgb = c
    p.font.bold = b; p.font.name = 'Segoe UI'; p.space_before = sp
    return p

def slide_title(s, title, takeaway=None):
    txt(s, Inches(1.2), Inches(0.35), Inches(10.5), Inches(0.7),
        title, 24, NAVY, True, fn='Georgia')
    hline(s, Inches(1.2), Inches(1.05), Inches(2.0), BLUE, Pt(3))
    if takeaway:
        txt(s, Inches(1.2), Inches(1.15), Inches(10.5), Inches(0.35),
            takeaway, 12, BLUE)
    brand_bar(s)


# ════════════════════════════════════════════════════════════════════════
# STREAMING ARCHITECTURE ON GABRIEL NIMBUS
# ════════════════════════════════════════════════════════════════════════
s = blank()
slide_title(s, 'Streaming Architecture on Gabriel Nimbus',
    'Cloud-native behavioral detection pipeline. End-to-end LOTL detection in 1-5 seconds.')

# ── FOUR BLOCKS: TELEMETRY → BUS → FLINK → OUTPUT ──

# Block 1: DODIN Telemetry (left)
rrect(s, Inches(1.2), Inches(1.55), Inches(2.0), Inches(3.5), fill=LBLUE, border=MGRAY)
rect(s, Inches(1.2), Inches(1.55), Inches(2.0), Inches(0.4), fill=BLUE)
txt(s, Inches(1.3), Inches(1.6), Inches(1.8), Inches(0.3),
    'DODIN TELEMETRY', 10, WHITE, True, PP_ALIGN.CENTER)
tf = txt(s, Inches(1.4), Inches(2.1), Inches(1.6), Inches(2.8),
    'NGFW / IPS Logs', 11, INK, True)
add_p(tf, 'Identity (AD, Okta, Entra, TACACS+)', 11, INK)
add_p(tf, 'Endpoint (EDR / Sysmon)', 11, INK)
add_p(tf, 'Network (NetFlow, DNS, VPC)', 11, INK)
add_p(tf, 'Cloud (CloudTrail, Azure)', 11, INK)

# Arrow 1
txt(s, Inches(3.25), Inches(3.1), Inches(0.3), Inches(0.3),
    '→', 16, BLUE, True, PP_ALIGN.CENTER)

# Block 2: Message Bus
rrect(s, Inches(3.6), Inches(1.55), Inches(1.4), Inches(3.5), fill=BLUE)
txt(s, Inches(3.65), Inches(1.7), Inches(1.3), Inches(0.5),
    'Kafka\n/ Kinesis', 15, WHITE, True, PP_ALIGN.CENTER, fn='Georgia')
hline(s, Inches(3.75), Inches(2.45), Inches(1.1), RGBColor(0x2A, 0x6F, 0x97))
tf = txt(s, Inches(3.7), Inches(2.6), Inches(1.2), Inches(2.3),
    'Per-source topics', 10, RGBColor(0xA0, 0xC8, 0xE0))
add_p(tf, 'Independent consumer groups', 10, RGBColor(0xA0, 0xC8, 0xE0))
add_p(tf, '', 6)
add_p(tf, 'Event-time watermarks', 9, RGBColor(0x7A, 0x9B, 0xBE))

# Arrow 2
txt(s, Inches(5.05), Inches(3.1), Inches(0.3), Inches(0.3),
    '→', 16, BLUE, True, PP_ALIGN.CENTER)

# Block 3: Flink Cluster (hero — widest block)
rrect(s, Inches(5.4), Inches(1.55), Inches(3.6), Inches(3.5), fill=NAVY, border=MGRAY)
txt(s, Inches(5.5), Inches(1.62), Inches(3.4), Inches(0.3),
    'Flink Cluster', 13, WHITE, True, fn='Georgia')

stages = [
    ('1.  Kafka source', 'Event-time, watermarks'),
    ('2.  OCSF normalize', 'Per-source parsers'),
    ('3.  Enrich & classify', 'Identity, asset, tier, IOC'),
    ('4.  CEP + ML detect', 'NFA, Sigma, baselines'),
    ('5.  ACECARD behavioral', 'Digital Twin, drift, MITRE'),
    ('6.  Sink fanout', 'Alerts, events, archive'),
]
py = Inches(2.05)
for title, detail in stages:
    rrect(s, Inches(5.6), py, Inches(3.2), Inches(0.48), fill=BLUE, border=RGBColor(0x2A, 0x6F, 0x97))
    txt(s, Inches(5.7), py + Inches(0.04), Inches(1.6), Inches(0.2),
        title, 10, WHITE, True)
    txt(s, Inches(7.3), py + Inches(0.04), Inches(1.4), Inches(0.2),
        detail, 9, RGBColor(0xA0, 0xC8, 0xE0))
    py += Inches(0.52)

txt(s, Inches(5.5), Inches(4.7), Inches(3.4), Inches(0.2),
    '3-15 ms per event  ·  CEP NFA under 1 ms', 8, RGBColor(0x7A, 0x9B, 0xBE), al=PP_ALIGN.CENTER)

# Arrow 3
txt(s, Inches(9.05), Inches(3.1), Inches(0.3), Inches(0.3),
    '→', 16, BLUE, True, PP_ALIGN.CENTER)

# Block 4: Store & Analyst (right)
rrect(s, Inches(9.4), Inches(1.55), Inches(2.8), Inches(3.5), fill=LGRAY, border=MGRAY)

# Store sub-section
rect(s, Inches(9.4), Inches(1.55), Inches(2.8), Inches(0.4), fill=TEAL)
txt(s, Inches(9.5), Inches(1.6), Inches(2.6), Inches(0.3),
    'STORE', 10, WHITE, True, PP_ALIGN.CENTER)
tf = txt(s, Inches(9.6), Inches(2.1), Inches(2.4), Inches(1.2),
    'ClickHouse / Pinot (30-90d hot)', 10, INK)
add_p(tf, 'Vector DB / pgvector (behavioral)', 10, INK)
add_p(tf, 'Iceberg / S3 (cold archive)', 10, INK)

# Divider
hline(s, Inches(9.55), Inches(3.15), Inches(2.5), MGRAY)

# Analyst sub-section
txt(s, Inches(9.6), Inches(3.25), Inches(2.4), Inches(0.2),
    'ANALYST', 9, TEAL, True)
tf = txt(s, Inches(9.6), Inches(3.5), Inches(2.4), Inches(1.4),
    'SIEM / SOAR (alerts)', 10, INK)
add_p(tf, 'SOC Dashboard (real-time)', 10, INK)
add_p(tf, 'Hunt UI (SQL hunts)', 10, INK)
add_p(tf, 'ServiceNow / JIRA (ticketing)', 10, INK)

# ── BOTTOM ROW: SUPPORTING INFRASTRUCTURE ──
# Preemptive Engine
rrect(s, Inches(1.2), Inches(5.25), Inches(3.8), Inches(0.55), fill=LBLUE, border=MGRAY)
rect(s, Inches(1.2), Inches(5.25), Inches(0.06), Inches(0.55), fill=BLUE)
txt(s, Inches(1.4), Inches(5.28), Inches(3.4), Inches(0.2),
    'Preemptive Verification Engine', 10, NAVY, True)
txt(s, Inches(1.4), Inches(5.48), Inches(3.4), Inches(0.2),
    'Formal math model  ·  hourly + on config change', 8, DGRAY)

# RocksDB State
rrect(s, Inches(5.2), Inches(5.25), Inches(3.8), Inches(0.55), fill=LTEAL, border=MGRAY)
rect(s, Inches(5.2), Inches(5.25), Inches(0.06), Inches(0.55), fill=TEAL)
txt(s, Inches(5.4), Inches(5.28), Inches(3.4), Inches(0.2),
    'RocksDB Keyed State', 10, NAVY, True)
txt(s, Inches(5.4), Inches(5.48), Inches(3.4), Inches(0.2),
    'NVMe local + S3 incremental checkpoints', 8, DGRAY)

# Remediation
rrect(s, Inches(9.2), Inches(5.25), Inches(3.0), Inches(0.55), fill=LGOLD, border=MGRAY)
rect(s, Inches(9.2), Inches(5.25), Inches(0.06), Inches(0.55), fill=GOLD)
txt(s, Inches(9.4), Inches(5.28), Inches(2.6), Inches(0.2),
    'Remediation Intelligence', 10, NAVY, True)
txt(s, Inches(9.4), Inches(5.48), Inches(2.6), Inches(0.2),
    'Zero-FP playbooks  ·  auto-ticketing', 8, DGRAY)

# ── PERFORMANCE CALLOUT (gold bar — matches v12 bottom bars) ──
rrect(s, Inches(1.2), Inches(5.95), Inches(11.0), Inches(0.35), fill=LGOLD)
txt(s, Inches(1.5), Inches(5.98), Inches(10.5), Inches(0.3),
    'End-to-end LOTL detection: 1-5 seconds.  Source ~100ms + Kafka ~10ms + Flink 3-15ms + alert ~50ms.',
    10, NAVY, True, PP_ALIGN.CENTER)

# ── GABRIEL NIMBUS PLATFORM SERVICES BAR ──
rrect(s, Inches(1.2), Inches(6.4), Inches(11.0), Inches(0.55), fill=LGRAY, border=MGRAY)
rect(s, Inches(1.2), Inches(6.4), Inches(11.0), Inches(0.06), fill=GOLD)
txt(s, Inches(1.5), Inches(6.5), Inches(10.5), Inches(0.2),
    'Gabriel Nimbus Platform Services', 10, NAVY, True, PP_ALIGN.CENTER)
txt(s, Inches(1.5), Inches(6.7), Inches(10.5), Inches(0.2),
    'Kubernetes Orchestration  |  Object Store (S3)  |  Platform IAM  |  Big Data Analytics  |  Monitoring  |  Service Mesh',
    9, DGRAY, al=PP_ALIGN.CENTER)


prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
