"""CFIC Final Presentation - McKinsey Style (v4 - visual overhaul)
Professional icons (Wingdings/Segoe UI Symbol), generous whitespace,
varied layouts, rounded icon badges, flow arrows, branding bar.
22nd Century Technologies Inc. - May 6, 2026"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

OUT = os.path.join(os.path.dirname(__file__), "CFIC_Final_Presentation_v7.pptx")

# McKinsey palette - restrained, professional
INK     = RGBColor(0x2C, 0x3E, 0x50)
NAVY    = RGBColor(0x0D, 0x1B, 0x2A)
BLUE    = RGBColor(0x1B, 0x4F, 0x72)
LGRAY   = RGBColor(0xF7, 0xF8, 0xFA)
MGRAY   = RGBColor(0xE8, 0xEB, 0xEF)
DGRAY   = RGBColor(0x6C, 0x75, 0x7D)
RED     = RGBColor(0xC0, 0x39, 0x2B)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BLACK   = RGBColor(0x00, 0x00, 0x00)
GOLD    = RGBColor(0xB7, 0x95, 0x0B)
TEAL    = RGBColor(0x0E, 0x6B, 0x8A)
LTEAL   = RGBColor(0xE8, 0xF4, 0xF8)
LBLUE   = RGBColor(0xEB, 0xF5, 0xFB)
LRED    = RGBColor(0xFD, 0xED, 0xEC)
LGOLD   = RGBColor(0xFE, 0xF9, 0xE7)
# Architecture layer colors
L5_BLUE = RGBColor(0x2E, 0x86, 0xAB)
L4_BLUE = RGBColor(0x1B, 0x6F, 0x92)
L3_BLUE = RGBColor(0x14, 0x5D, 0x82)
L2_BLUE = RGBColor(0x0F, 0x3D, 0x5C)
L1_BLUE = RGBColor(0x0D, 0x1B, 0x2A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

slide_count = [0]

# ────────────────────────────────────────────────────────────────
# HELPERS
# ────────────────────────────────────────────────────────────────

def blank():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    slide_count[0] += 1
    return s

def brand_bar(s):
    """Bottom branding strip on every content slide"""
    rect(s, Inches(0), Inches(7.15), W, Inches(0.35), fill=NAVY)
    txt(s, Inches(0.5), Inches(7.17), Inches(5), Inches(0.3),
        '22nd Century Technologies  |  PROPRIETARY', 8, RGBColor(0x8A, 0x9B, 0xAE), False)
    txt(s, Inches(11.5), Inches(7.17), Inches(1.5), Inches(0.3),
        str(slide_count[0]), 8, RGBColor(0x8A, 0x9B, 0xAE), False, PP_ALIGN.RIGHT)

def hline(s, x, y, w, color=MGRAY, thick=Pt(1)):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, thick)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def rect(s, x, y, w, h, fill=None, border=None):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if border:
        sh.line.color.rgb = border
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def rrect(s, x, y, w, h, fill=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def oval(s, x, y, size, fill):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def chevron_right(s, x, y, w, h, fill):
    sh = s.shapes.add_shape(MSO_SHAPE.CHEVRON, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def txt(s, x, y, w, h, text, sz=11, c=INK, b=False, al=PP_ALIGN.LEFT, fn='Calibri'):
    tf = s.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = c
    p.font.bold = b
    p.font.name = fn
    p.alignment = al
    return tf

def add_p(tf, text, sz=11, c=INK, b=False, sp=Pt(6), fn='Calibri'):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = c
    p.font.bold = b
    p.font.name = fn
    p.space_before = sp
    return p

def icon_circle(s, x, y, icon_char, color, size=Inches(0.5)):
    """Colored circle with a Wingdings/Segoe icon inside"""
    oval(s, x, y, size, color)
    tf = s.shapes.add_textbox(x, y, size, size).text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = icon_char
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = 'Segoe UI Symbol'
    p.alignment = PP_ALIGN.CENTER

def section_header(title, subtitle=None):
    s = blank()
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = NAVY
    # Accent bar at top
    rect(s, Inches(0), Inches(0), W, Inches(0.06), fill=GOLD)
    # Large section number + title
    txt(s, Inches(1.8), Inches(2.6), Inches(10), Inches(1.2),
        title, 38, WHITE, True, PP_ALIGN.LEFT, 'Georgia')
    if subtitle:
        txt(s, Inches(1.8), Inches(3.7), Inches(9.5), Inches(0.6),
            subtitle, 15, RGBColor(0xA0, 0xB0, 0xC0), False)
    # Thin gold line under title
    hline(s, Inches(1.8), Inches(3.5), Inches(2.5), GOLD, Pt(2))
    brand_bar(s)
    return s

def slide_title(s, title, takeaway=None):
    """Standard slide title with blue accent line"""
    txt(s, Inches(1.2), Inches(0.4), Inches(10.5), Inches(0.6),
        title, 22, NAVY, True, fn='Georgia')
    hline(s, Inches(1.2), Inches(1.0), Inches(2.0), BLUE, Pt(2.5))
    if takeaway:
        txt(s, Inches(1.2), Inches(1.1), Inches(10.5), Inches(0.45),
            takeaway, 11.5, BLUE, False)
    brand_bar(s)
    return s


# Wingdings / Segoe UI Symbol icons that reliably render in PowerPoint
ICON_SHIELD   = '⛨'   # shield
ICON_LOCK     = '\U0001F512' # lock
ICON_WARNING  = '⚠'   # warning triangle
ICON_CHECK    = '✓'   # checkmark
ICON_CROSS    = '✗'   # cross mark
ICON_ARROW_R  = '▶'   # right-pointing triangle
ICON_GRAPH    = '≡'   # graph/bars
ICON_EYE      = '◉'   # eye/target
ICON_GEAR     = '⚙'   # gear
ICON_NET      = '◈'   # diamond / network
ICON_BRAIN    = '⦿'   # circled bullet / AI
ICON_BOLT     = '⚡'   # lightning bolt
ICON_CHAIN    = '⛓'   # chain links
ICON_CLOCK    = '⏱'   # stopwatch
ICON_STAR     = '★'   # star


# ════════════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ════════════════════════════════════════════════════════════════════════
s = blank()
s.background.fill.solid()
s.background.fill.fore_color.rgb = NAVY
# Top accent
rect(s, Inches(0), Inches(0), W, Inches(0.06), fill=GOLD)
# Company
txt(s, Inches(1.8), Inches(1.8), Inches(10), Inches(0.4),
    '22nd Century Technologies Inc.', 13, GOLD, True, PP_ALIGN.LEFT, 'Georgia')
# Main title
txt(s, Inches(1.8), Inches(2.6), Inches(10), Inches(1.5),
    'Containerized Behavioral\nCyber Defense for Gabriel Nimbus', 36, WHITE, True, fn='Georgia')
# Subtitle
txt(s, Inches(1.8), Inches(4.5), Inches(10), Inches(0.5),
    'Preemptive Posture Verification  +  ACECARD UEBA Behavioral Intelligence', 14, RGBColor(0xA0, 0xB8, 0xD0))
# Gold divider
hline(s, Inches(1.8), Inches(5.3), Inches(3.5), GOLD, Pt(2))
# Event details
txt(s, Inches(1.8), Inches(5.6), Inches(10), Inches(0.4),
    'CFIC Collaboration Event  |  Georgia Cyber Innovation & Training Center', 11, RGBColor(0x7A, 0x8B, 0x9E))
txt(s, Inches(1.8), Inches(6.0), Inches(10), Inches(0.3),
    'Augusta, GA  |  May 6, 2026', 11, RGBColor(0x7A, 0x8B, 0x9E))
# Classification
txt(s, Inches(1.8), Inches(6.7), Inches(10), Inches(0.3),
    'PROPRIETARY  //  NOT FOR PUBLIC RELEASE', 9, RGBColor(0x5A, 0x6B, 0x7E))
brand_bar(s)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 2: EXECUTIVE SUMMARY
# ════════════════════════════════════════════════════════════════════════
s = blank()
slide_title(s, 'Executive Summary')

# Three columns with colored top bars
col_w = Inches(3.5)
col_h = Inches(5.0)
gap = Inches(0.3)
start_x = Inches(1.2)
start_y = Inches(1.6)

# SITUATION
x = start_x
rect(s, x, start_y, col_w, col_h, fill=LRED)
rect(s, x, start_y, col_w, Inches(0.5), fill=RED)
txt(s, x+Inches(0.1), start_y+Inches(0.08), col_w-Inches(0.2), Inches(0.4),
    '  SITUATION', 13, WHITE, True)
tf = txt(s, x+Inches(0.3), start_y+Inches(0.7), col_w-Inches(0.6), Inches(4.0),
    'Nation-state adversaries (Volt Typhoon, Salt Typhoon) have dwelled inside U.S. critical infrastructure for years.', 11, INK)
add_p(tf, '', 8)
add_p(tf, 'AI-enabled reconnaissance and evasion make these campaigns harder to detect with each cycle.', 11, INK)
add_p(tf, '', 8)
add_p(tf, 'ARCYBER needs behavioral intelligence inside the Gabriel Nimbus classified big data platform.', 11, INK)

# COMPLICATION
x = start_x + col_w + gap
rect(s, x, start_y, col_w, col_h, fill=LBLUE)
rect(s, x, start_y, col_w, Inches(0.5), fill=NAVY)
txt(s, x+Inches(0.1), start_y+Inches(0.08), col_w-Inches(0.2), Inches(0.4),
    '  COMPLICATION', 13, WHITE, True)
tf = txt(s, x+Inches(0.3), start_y+Inches(0.7), col_w-Inches(0.6), Inches(4.0),
    'Config gaps are never PROVEN closed. Pentests sample; nobody exhaustively verifies.', 11, INK)
add_p(tf, '', 8)
add_p(tf, 'Current UEBA uses fixed thresholds. LOTL tradecraft operates below every threshold.', 11, INK)
add_p(tf, '', 8)
add_p(tf, 'Identity fragmented across 10+ systems. No cross-domain behavioral trajectory.', 11, INK)
add_p(tf, '', 8)
add_p(tf, 'No containerized, K8s-native UEBA for classified enclaves.', 11, INK)

# RESOLUTION
x = start_x + 2*(col_w + gap)
rect(s, x, start_y, col_w, col_h, fill=LTEAL)
rect(s, x, start_y, col_w, Inches(0.5), fill=TEAL)
txt(s, x+Inches(0.1), start_y+Inches(0.08), col_w-Inches(0.2), Inches(0.4),
    '  RESOLUTION', 13, WHITE, True)
tf = txt(s, x+Inches(0.3), start_y+Inches(0.7), col_w-Inches(0.6), Inches(4.0),
    '22CT Preemptive: Formal verification of every defensive config. Zero FP/FN within model scope.', 11, INK)
add_p(tf, '', 8)
add_p(tf, 'ACECARD UEBA: 1536-d embeddings + CUSUM change-point + MITRE-mapped drift direction.', 11, INK)
add_p(tf, '', 8)
add_p(tf, 'Together: attack surface never exists OR behavioral drift detected in hours, not years.', 11, NAVY, True)
add_p(tf, '', 8)
add_p(tf, 'K8s-native, Iron Bank, IL5/IL6/JWICS. 90 days to first detection.', 11, INK)

# Compliance footer
txt(s, Inches(1.2), Inches(6.8), Inches(11.0), Inches(0.3),
    'DoD Zero Trust RA (Pillars 4 & 7)  |  NIST 800-53 Rev 5  |  cATO path  |  RMF/eMASS ready',
    9.5, DGRAY, False, PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════
# SECTION 1: PROBLEM STATEMENT
# ════════════════════════════════════════════════════════════════════════
section_header('1. The Problem',
    'AI-enabled adversaries dwell for years. Current defenses detect neither the entry nor the evolution.')

# SLIDE 4: AI-Enabled Cyber Attacks
s = blank()
slide_title(s, 'AI-Enabled Adversaries Have Changed the Game',
    'Nation-state actors use AI for reconnaissance, evasion, and adaptive persistence.')

cols = [
    (ICON_BOLT, 'AI-Powered\nReconnaissance', RED,
     'LLMs generate exploit code from CVE descriptions in minutes\n'
     'AI scans 100K+ endpoints for misconfiguration patterns\n'
     'Automated supply-chain mapping identifies weakest links'),
    (ICON_EYE, 'AI-Assisted\nEvasion', RED,
     'Adversary AI mutates C2 traffic to evade signature detection\n'
     'ML models test payloads against defender rules pre-campaign\n'
     'Generative AI produces convincing phishing at scale'),
    (ICON_CLOCK, 'AI-Extended\nDwell Time', RED,
     'Adaptive beaconing avoids pattern-based detection\n'
     'AI decides when to move laterally based on defender activity\n'
     'Automated credential rotation maintains persistence'),
]
x_positions = [Inches(1.2), Inches(5.1), Inches(9.0)]
for i, (icon, title, color, body) in enumerate(cols):
    x = x_positions[i]
    rrect(s, x, Inches(1.6), Inches(3.6), Inches(3.3), fill=LGRAY)
    icon_circle(s, x+Inches(0.2), Inches(1.8), icon, color)
    txt(s, x+Inches(0.85), Inches(1.75), Inches(2.5), Inches(0.7), title, 12, NAVY, True)
    tf = txt(s, x+Inches(0.3), Inches(2.6), Inches(3.0), Inches(2.2), '', 10.5, INK)
    for line_text in body.split('\n'):
        add_p(tf, f'•  {line_text}', 10.5, INK, sp=Pt(4))

# Impact callout
rrect(s, Inches(1.2), Inches(5.2), Inches(11.0), Inches(0.9), fill=LRED)
rect(s, Inches(1.2), Inches(5.2), Inches(0.06), Inches(0.9), fill=RED)
tf = txt(s, Inches(1.5), Inches(5.3), Inches(10.5), Inches(0.7),
    'Two active campaigns operated inside U.S. critical infrastructure for years:', 11.5, RED, True)
add_p(tf, 'Volt Typhoon (G1017): active since mid-2021, pre-positioned in power, water, telecom   |   '
      'Salt Typhoon: 9+ U.S. telecoms, multi-year dwell', 10.5, INK, sp=Pt(4))

txt(s, Inches(1.2), Inches(6.3), Inches(11.0), Inches(0.3),
    'Sources: CISA AA24-038A  |  FBI/NSA/CISA Joint Advisory (Dec 2024)  |  Microsoft Threat Intel  |  Cisco Talos',
    9, DGRAY, False, PP_ALIGN.LEFT)

# SLIDE 5: Volt Typhoon
s = blank()
slide_title(s, 'Case Study: Volt Typhoon (MITRE G1017)',
    'Pre-positioned in U.S. critical infrastructure for future sabotage. Active since at least mid-2021.')

phases = [
    (ICON_BOLT, 'Initial Access', 'Exploits unpatched Fortinet, Ivanti, Cisco IOS XE, PAN edge devices. CVE-2023-20198, CVE-2024-3400.'),
    (ICON_LOCK, 'Credential Access', 'ntdsutil.exe dumps NTDS.dit. LSASS memory extraction. All tools are legitimate Windows admin binaries.'),
    (ICON_NET, 'Lateral Movement', 'RDP, WinRM, SMB to domain controllers. Routes C2 through KV botnet (compromised SOHO routers).'),
    (ICON_CHAIN, 'Persistence', 'Web shells (.aspx/.jspx). Scheduled tasks. Registry modification. Multi-year pre-positioning.'),
    (ICON_EYE, 'Defense Evasion', 'Clears event logs. Deletes MRU keys. Living-off-the-Land only — no malware to detect.'),
]
y = Inches(1.6)
for i, (icon, name, desc) in enumerate(phases):
    bg = LGRAY if i % 2 == 0 else WHITE
    rrect(s, Inches(1.2), y, Inches(11.0), Inches(0.85), fill=bg)
    icon_circle(s, Inches(1.4), y+Inches(0.17), icon, NAVY, Inches(0.5))
    txt(s, Inches(2.0), y+Inches(0.08), Inches(1.6), Inches(0.4), f'Phase {i+1}', 9, DGRAY)
    txt(s, Inches(2.0), y+Inches(0.35), Inches(1.6), Inches(0.4), name, 11.5, NAVY, True)
    txt(s, Inches(3.8), y+Inches(0.15), Inches(8.2), Inches(0.6), desc, 10.5, INK)
    y += Inches(0.92)

# Key insight
rrect(s, Inches(1.2), Inches(6.3), Inches(11.0), Inches(0.55), fill=LRED)
icon_circle(s, Inches(1.4), Inches(6.38), ICON_WARNING, RED, Inches(0.4))
txt(s, Inches(1.95), Inches(6.35), Inches(10.0), Inches(0.45),
    'Every tool is a legitimate admin binary. Every login uses valid credentials. No malware signature exists. Detection is structurally blind.', 10.5, RED, True)

# SLIDE 6: Salt Typhoon
s = blank()
slide_title(s, 'Case Study: Salt Typhoon',
    'Carrier-scale telecom intrusions. 9+ U.S. telecoms compromised. Multi-year dwell. CALEA access.')

phases = [
    (ICON_BOLT, 'Edge Exploitation', 'CVE-2023-20198 (Cisco IOS XE), CVE-2018-0171, CVE-2024-21887 (Ivanti). 1,000+ Cisco devices.'),
    (ICON_CHAIN, 'Persistence', 'Local privileged accounts on routers. ACL whitelisting. Guest Shell enabled. Config-resident.'),
    (ICON_LOCK, 'TACACS+ Harvest', 'On-box packet capture of TCP/49. Harvests network engineer credentials. JumbledPath chains.'),
    (ICON_NET, 'Collection', 'GRE tunnels siphon CDR/metadata. Targets CALEA lawful-intercept. GhostSpider + Demodex (Trend Micro).'),
    (ICON_CLOCK, 'Long-Term Espionage', 'Slow, persistent CDR exfiltration over months/years. Encrypted C2. LOTL on network gear.'),
]
y = Inches(1.6)
for i, (icon, name, desc) in enumerate(phases):
    bg = LGRAY if i % 2 == 0 else WHITE
    rrect(s, Inches(1.2), y, Inches(11.0), Inches(0.85), fill=bg)
    icon_circle(s, Inches(1.4), y+Inches(0.17), icon, NAVY, Inches(0.5))
    txt(s, Inches(2.0), y+Inches(0.08), Inches(1.6), Inches(0.4), f'Phase {i+1}', 9, DGRAY)
    txt(s, Inches(2.0), y+Inches(0.35), Inches(1.6), Inches(0.4), name, 11.5, NAVY, True)
    txt(s, Inches(3.8), y+Inches(0.15), Inches(8.2), Inches(0.6), desc, 10.5, INK)
    y += Inches(0.92)

rrect(s, Inches(1.2), Inches(6.3), Inches(11.0), Inches(0.55), fill=LRED)
icon_circle(s, Inches(1.4), Inches(6.38), ICON_WARNING, RED, Inches(0.4))
txt(s, Inches(1.95), Inches(6.35), Inches(10.0), Inches(0.45),
    'Salt Typhoon lives in router configs and ACLs — not in files. Their persistence IS the network configuration.', 10.5, RED, True)


# ════════════════════════════════════════════════════════════════════════
# SECTION 2: WHY CURRENT TOOLS FAIL
# ════════════════════════════════════════════════════════════════════════
section_header('2. Why Current Tools Fail',
    'Neither configuration-only nor behavior-only defenses can detect these campaigns.')

# SLIDE 8: Four Structural Gaps
s = blank()
slide_title(s, 'Four Structural Gaps Exploited by Both Campaigns',
    'Each gap is a design limitation, not a tuning problem. No amount of rules or thresholds closes them.')

gaps = [
    (ICON_CROSS, 'Configuration Gaps Are Never PROVEN Closed',
     '99% of firewall breaches are misconfiguration (Gartner). Pentests sample. ASM enumerates. Nobody PROVES coverage across the combinatorial config state space.'),
    (ICON_CROSS, 'Behavioral Detection Has No Direction',
     '"User is 87% anomalous" — anomalous HOW? LOTL tools (ntdsutil, wmic, netsh) are legitimate admin binaries. Volt Typhoon operated BELOW all thresholds for years.'),
    (ICON_CROSS, 'Identity Is Fragmented Across 10+ Systems',
     'AD, AAD, Okta, AWS IAM, K8s, CrowdStrike, VPN, Splunk, PKI, TACACS+. Cross-domain lateral reads as 4 unrelated weak signals. No composite trajectory.'),
    (ICON_CROSS, 'No Cloud-Native UEBA for Gabriel Nimbus',
     'Current UEBA is SaaS-only or legacy on-prem. No containerized, K8s-native behavioral platform for IL5/IL6/JWICS with cATO alignment.'),
]
y = Inches(1.6)
for icon, title, body in gaps:
    rrect(s, Inches(1.2), y, Inches(11.0), Inches(1.15), fill=LGRAY)
    rect(s, Inches(1.2), y, Inches(0.06), Inches(1.15), fill=RED)
    icon_circle(s, Inches(1.5), y+Inches(0.32), icon, RED, Inches(0.45))
    txt(s, Inches(2.1), y+Inches(0.1), Inches(9.9), Inches(0.35), title, 12, NAVY, True)
    txt(s, Inches(2.1), y+Inches(0.5), Inches(9.9), Inches(0.6), body, 10.5, INK)
    y += Inches(1.25)

# SLIDE 9: Why SIEM/EDR/UEBA fail
s = blank()
slide_title(s, 'Existing Security Stack vs. These Adversaries',
    'Tool-by-tool assessment against Volt Typhoon and Salt Typhoon tradecraft.')

tools = [
    ('SIEM', 'Correlates known patterns. Cannot detect novel combinations of legitimate tool usage. Alert fatigue from 10K+ daily events.', ICON_CROSS),
    ('EDR / XDR', 'Catches known malware. LOTL binaries are allowlisted. Cannot reason about cross-host behavioral evolution. Blind on network gear.', ICON_CROSS),
    ('Commercial UEBA', 'Fixed thresholds. Per-source silos. No cross-system identity fusion. No drift direction. Cannot explain WHY.', ICON_CROSS),
    ('Pentest / Red Team', 'Point-in-time sample. Cannot cover combinatorial state space. Exercises known TTPs only. Annual cadence vs. continuous adversary.', ICON_CROSS),
    ('ASM / BAS', 'Enumerates attack surface. Cannot PROVE controls block every path. Statistical confidence, not mathematical certainty.', ICON_CROSS),
]
y = Inches(1.6)
for tool, desc, icon in tools:
    rrect(s, Inches(1.2), y, Inches(11.0), Inches(0.88), fill=LGRAY)
    rect(s, Inches(1.2), y, Inches(2.0), Inches(0.88), fill=NAVY)
    txt(s, Inches(1.3), y+Inches(0.2), Inches(1.8), Inches(0.5), tool, 11, WHITE, True, PP_ALIGN.CENTER)
    txt(s, Inches(3.4), y+Inches(0.15), Inches(8.0), Inches(0.6), desc, 10.5, INK)
    icon_circle(s, Inches(11.6), y+Inches(0.2), icon, RED, Inches(0.4))
    y += Inches(0.95)

rrect(s, Inches(1.2), Inches(6.4), Inches(11.0), Inches(0.5), fill=LBLUE)
txt(s, Inches(1.5), Inches(6.45), Inches(10.5), Inches(0.4),
    'The gap is structural, not operational. A fundamentally different approach is required — one that PROVES and one that detects DIRECTION.',
    10.5, NAVY, True, PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════
# SECTION 3: THE SOLUTION
# ════════════════════════════════════════════════════════════════════════
section_header('3. Our Solution: Two Layers',
    'Preemptive verification + behavioral trajectory intelligence. Together, they close all four gaps.')

# SLIDE 11: Two-Layer Overview
s = blank()
slide_title(s, 'Two-Layer Defense: Preemptive + Behavioral',
    'Layer 1 eliminates attack surface mathematically. Layer 2 detects what remains behaviorally.')

# Layer 1 card
lw = Inches(5.4)
lh = Inches(4.6)
rect(s, Inches(1.2), Inches(1.6), lw, lh, fill=LBLUE)
rect(s, Inches(1.2), Inches(1.6), lw, Inches(0.55), fill=BLUE)
icon_circle(s, Inches(1.4), Inches(1.68), ICON_SHIELD, BLUE, Inches(0.4))
txt(s, Inches(1.9), Inches(1.68), Inches(4.5), Inches(0.4),
    'LAYER 1 — 22CT PREEMPTIVE', 14, WHITE, True)
txt(s, Inches(1.5), Inches(2.3), Inches(4.8), Inches(0.3),
    'Closes Gap #1: Configuration never proven', 10, BLUE, False)

tf = txt(s, Inches(1.5), Inches(2.7), Inches(4.8), Inches(3.2),
    '•  Formal mathematical model of every NGFW, IPS, IdP, SASE, WAF', 10.5, INK)
add_p(tf, '•  Exhaustively reasons over combinatorial config state space', 10.5, INK)
add_p(tf, '•  Three Intelligences: Attack, Defense, Remediation', 10.5, INK)
add_p(tf, '•  Re-verified hourly + on every config change', 10.5, INK)
add_p(tf, '•  Supports PAN, Fortinet, Cisco, Check Point', 10.5, INK)
add_p(tf, '', 10)
add_p(tf, 'GUARANTEE: Zero FP/FN within formal model scope.', 11, BLUE, True)
add_p(tf, 'Complete reasoning, not sampling.', 11, NAVY, True)
add_p(tf, '', 6)
add_p(tf, 'Powered by Rigor AI (innovation partner).', 9.5, DGRAY)

# Layer 2 card
rect(s, Inches(6.9), Inches(1.6), lw, lh, fill=LTEAL)
rect(s, Inches(6.9), Inches(1.6), lw, Inches(0.55), fill=TEAL)
icon_circle(s, Inches(7.1), Inches(1.68), ICON_BRAIN, TEAL, Inches(0.4))
txt(s, Inches(7.6), Inches(1.68), Inches(4.5), Inches(0.4),
    'LAYER 2 — ACECARD UEBA', 14, WHITE, True)
txt(s, Inches(7.2), Inches(2.3), Inches(4.8), Inches(0.3),
    'Closes Gaps #2, #3, #4', 10, TEAL, False)

tf = txt(s, Inches(7.2), Inches(2.7), Inches(4.8), Inches(3.2),
    '•  5 signals embedded into 1536-d vectors per entity/hour', 10.5, INK)
add_p(tf, '•  CUSUM change-point catches sub-threshold drift', 10.5, INK)
add_p(tf, '•  Drift DIRECTION with MITRE ATT&CK mapping', 10.5, INK)
add_p(tf, '•  Entity fusion across 10+ identity systems', 10.5, INK)
add_p(tf, '•  Containerized K8s-native for Gabriel Nimbus', 10.5, INK)
add_p(tf, '', 10)
add_p(tf, 'INNOVATION: Not what is anomalous —', 11, TEAL, True)
add_p(tf, 'what they are BECOMING. Direction, not magnitude.', 11, NAVY, True)

# Bottom integration line
rrect(s, Inches(1.2), Inches(6.45), Inches(11.0), Inches(0.45), fill=LGOLD)
txt(s, Inches(1.5), Inches(6.5), Inches(10.5), Inches(0.35),
    'Together: attack surface never exists (Preemptive) OR behavioral drift detected and contained within hours (ACECARD + ABAC)',
    10.5, NAVY, True, PP_ALIGN.CENTER)

# SLIDE 12: Preemptive - Three Intelligences
s = blank()
slide_title(s, 'Layer 1: Three Intelligences',
    'From threat intel ingestion to verified remediation. Continuous, not periodic.')

cols = [
    (ICON_BOLT, 'Attack\nIntelligence', BLUE,
     ['MITRE-enriched attack graphs',
      'Built by ML/NLP/LLM pipeline',
      'Ingests: CVE, CISA, CTI, IOCs',
      'Full TTP graph for G1017 + Salt',
      'Updates on every new advisory']),
    (ICON_SHIELD, 'Defense\nIntelligence', BLUE,
     ['Symbolic Model of Computation',
      'Every config as formal logic',
      'Exhaustive state-space reasoning',
      'PROVES every path blocked',
      'Or surfaces the exact gap']),
    (ICON_GEAR, 'Remediation\nIntelligence', BLUE,
     ['Guardrailed agentic-AI reasoners',
      'Risk-prioritized, vendor-specific',
      'Validates fix introduces no gaps',
      'SOC playbooks auto-generated',
      'ServiceNow/JIRA integration']),
]
x_positions = [Inches(1.2), Inches(5.1), Inches(9.0)]
for i, (icon, title, color, items) in enumerate(cols):
    x = x_positions[i]
    rrect(s, x, Inches(1.6), Inches(3.6), Inches(4.6), fill=LGRAY)
    rect(s, x, Inches(1.6), Inches(3.6), Inches(0.06), fill=color)
    icon_circle(s, x+Inches(1.5), Inches(1.8), icon, color, Inches(0.55))
    txt(s, x+Inches(0.3), Inches(2.45), Inches(3.0), Inches(0.7), title, 13, NAVY, True, PP_ALIGN.CENTER)
    tf = txt(s, x+Inches(0.3), Inches(3.3), Inches(3.0), Inches(2.5), '', 10.5, INK)
    for item in items:
        add_p(tf, f'•  {item}', 10.5, INK, sp=Pt(5))

rrect(s, Inches(1.2), Inches(6.4), Inches(11.0), Inches(0.5), fill=LGRAY)
txt(s, Inches(1.5), Inches(6.45), Inches(10.5), Inches(0.4),
    'OUTPUT:  Mathematical proof every TTP path is blocked. Re-verified hourly. Config drift caught <1hr. Vendor-specific fix generated.',
    10, NAVY, True, PP_ALIGN.CENTER)

# SLIDE 13: ACECARD - Core Innovations
s = blank()
slide_title(s, 'Layer 2: ACECARD UEBA — Core Innovations',
    'Behavioral trajectory intelligence. Not anomaly scores — directional, MITRE-mapped drift detection.')

innovations = [
    (ICON_GRAPH, '1536-d Behavioral Embeddings', TEAL,
     '5 signals (Auth, Process, Network, File, Identity) summarized as text, embedded into unified vector space. '
     'Captures the SHAPE of behavior. 6 vectors stored per entity per window.'),
    (ICON_STAR, 'CUSUM Change-Point Detection', TEAL,
     "Page's CUSUM (1954): accumulates small drifts that NEVER cross fixed thresholds. "
     'Minimax optimal for i.i.d. shift detection (Moustakides 1986). Alarm in hours, not years.'),
    (ICON_ARROW_R, 'Drift Direction + MITRE Mapping', TEAL,
     'Drift vector projected onto 8 threat concepts via cosine similarity. Output: "drifting toward lateral_movement (0.78, T1021)" '
     '— not "87% anomalous."'),
    (ICON_CHAIN, 'Entity Fusion (10+ Systems)', TEAL,
     'One entity_uuid resolves AD, AAD, Okta, AWS, K8s, CrowdStrike, VPN, Splunk, PKI, TACACS+. '
     'Without fusion: 4 weak signals. With fusion: 1 strong trajectory, HIGH alert.'),
]
y = Inches(1.6)
for icon, title, color, body in innovations:
    rrect(s, Inches(1.2), y, Inches(11.0), Inches(1.1), fill=LGRAY)
    rect(s, Inches(1.2), y, Inches(0.06), Inches(1.1), fill=color)
    icon_circle(s, Inches(1.5), y+Inches(0.3), icon, color, Inches(0.45))
    txt(s, Inches(2.1), y+Inches(0.08), Inches(9.9), Inches(0.3), title, 12, NAVY, True)
    txt(s, Inches(2.1), y+Inches(0.42), Inches(9.9), Inches(0.6), body, 10, INK)
    y += Inches(1.18)

rrect(s, Inches(1.2), Inches(6.4), Inches(11.0), Inches(0.45), fill=LTEAL)
txt(s, Inches(1.5), Inches(6.45), Inches(10.5), Inches(0.35),
    'No commercial UEBA combines: behavioral embeddings + CUSUM on vectors + drift direction onto MITRE + entity fusion across 10+ sources.',
    10, TEAL, True, PP_ALIGN.CENTER)

# SLIDE 14: Five Signals
s = blank()
slide_title(s, 'Five Behavioral Signals',
    'Each signal captures a distinct attack dimension. Together they form a complete behavioral fingerprint.')

signals = [
    (ICON_LOCK, 'AUTH', 'Win 4624/4625/4768, Okta, AAD, CloudTrail',
     'logon_count, failure_rate, unique_hosts, off_hours_ratio, impossible_travel, mfa_skip_ratio'),
    (ICON_GEAR, 'PROCESS', 'Sysmon EID 1/3, CrowdStrike, Defender, auditd',
     'unique_processes, cmdline_entropy, lolbin_count, parent_child_depth, unsigned_ratio, encoded_cmd'),
    (ICON_NET, 'NETWORK', 'NetFlow/IPFIX, Zeek, PAN-OS, FortiGate, VPC Flow',
     'unique_dest_ips, bytes_out_ratio, beacon_score, dns_query_rate, geo_anomaly, admin_share'),
    (ICON_GRAPH, 'FILE', 'Sysmon EID 11, Win 4663, DLP, cloud storage',
     'files_created, files_deleted, sensitive_access, archive_creates, ext_changes, shadow_copy_ops'),
    (ICON_STAR, 'IDENTITY', 'Win 4672/4728, CloudTrail IAM, K8s audit, AD change',
     'priv_escalations, group_adds, mfa_bypass, role_changes, service_acct_use, admin_actions'),
]
y = Inches(1.6)
for icon, name, sources, features in signals:
    rrect(s, Inches(1.2), y, Inches(11.0), Inches(0.9), fill=LGRAY)
    icon_circle(s, Inches(1.4), y+Inches(0.2), icon, TEAL, Inches(0.45))
    rect(s, Inches(2.0), y+Inches(0.15), Inches(1.4), Inches(0.55), fill=NAVY)
    txt(s, Inches(2.05), y+Inches(0.22), Inches(1.3), Inches(0.4), name, 11, WHITE, True, PP_ALIGN.CENTER)
    txt(s, Inches(3.6), y+Inches(0.08), Inches(8.4), Inches(0.3), sources, 9.5, DGRAY)
    txt(s, Inches(3.6), y+Inches(0.42), Inches(8.4), Inches(0.4), features, 10, INK, False, fn='Consolas')
    y += Inches(0.97)

rrect(s, Inches(1.2), Inches(6.5), Inches(11.0), Inches(0.35), fill=LTEAL)
txt(s, Inches(1.5), Inches(6.53), Inches(10.5), Inches(0.3),
    'Air-gapped default: local embedding model (E5-Mistral / nomic-embed-text).  IL5 option: Azure OpenAI GovCloud.',
    9.5, TEAL, False, PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════
# SECTION 4: HOW IT STOPS THE ATTACKS
# ════════════════════════════════════════════════════════════════════════
section_header('4. How It Stops These Attacks',
    'Phase-by-phase: what Preemptive eliminates and what ACECARD catches if preemption is bypassed.')

# SLIDE 16: Volt Typhoon - Both Layers
s = blank()
slide_title(s, 'Stopping Volt Typhoon: Both Layers',
    'Each kill-chain phase addressed by the appropriate layer. Defense-in-depth.')

# Table header
rect(s, Inches(1.2), Inches(1.5), Inches(1.3), Inches(0.4), fill=NAVY)
txt(s, Inches(1.25), Inches(1.53), Inches(1.2), Inches(0.35), 'PHASE', 9, WHITE, True, PP_ALIGN.CENTER)
rect(s, Inches(2.5), Inches(1.5), Inches(3.0), Inches(0.4), fill=NAVY)
txt(s, Inches(2.55), Inches(1.53), Inches(2.9), Inches(0.35), 'ADVERSARY ACTION', 9, WHITE, True, PP_ALIGN.CENTER)
rect(s, Inches(5.5), Inches(1.5), Inches(3.4), Inches(0.4), fill=BLUE)
txt(s, Inches(5.55), Inches(1.53), Inches(3.3), Inches(0.35), '22CT PREEMPTIVE', 9, WHITE, True, PP_ALIGN.CENTER)
rect(s, Inches(8.9), Inches(1.5), Inches(3.3), Inches(0.4), fill=TEAL)
txt(s, Inches(8.95), Inches(1.53), Inches(3.2), Inches(0.35), 'ACECARD DETECTS', 9, WHITE, True, PP_ALIGN.CENTER)

rows = [
    ('Initial Access', 'Edge CVE exploit\n(Fortinet, Cisco, PAN)',
     'IPS + WAF verified;\ncompensating controls proven',
     'Auth: anomalous mgmt-plane\nsource IP, drift 0.08'),
    ('Credential\nAccess', 'ntdsutil NTDS.dit dump;\nLSASS extraction',
     'EDR coverage verified;\nscript-control blocks ntdsutil',
     'Process: lolbin_count 0→7,\nCUSUM alarm @ T+47min'),
    ('Lateral\nMovement', 'SMB/RDP/WinRM to DCs;\ncross-tenant pivots',
     'SMB shadow detected;\nmicro-segmentation proven',
     'Network: unique_dest 3→18,\ncohort z=6.17, HIGH'),
    ('Persistence', 'Web shells; local accounts;\nACL modification',
     'Config Drift catches\nbaseline deviation in <1hr',
     'Identity: svc_acct spikes,\ndrift toward T1078'),
    ('Pre-position', 'Credential harvest;\nC2 to KV botnet',
     'Egress filtering proven;\nC2 IPs blocked all zones',
     'File: sensitive_access rises,\nCRITICAL → ABAC BLOCKED'),
]
y = Inches(1.95)
for i, (phase, attack, preempt, acecard) in enumerate(rows):
    bg = LGRAY if i % 2 == 0 else WHITE
    rect(s, Inches(1.2), y, Inches(11.0), Inches(0.88), fill=bg)
    txt(s, Inches(1.3), y+Inches(0.1), Inches(1.1), Inches(0.7), phase, 9.5, NAVY, True, PP_ALIGN.CENTER)
    txt(s, Inches(2.6), y+Inches(0.1), Inches(2.8), Inches(0.7), attack, 9.5, INK)
    txt(s, Inches(5.6), y+Inches(0.1), Inches(3.2), Inches(0.7), preempt, 9.5, BLUE)
    txt(s, Inches(9.0), y+Inches(0.1), Inches(3.1), Inches(0.7), acecard, 9.5, TEAL)
    y += Inches(0.92)

# SLIDE 17: Salt Typhoon
s = blank()
slide_title(s, 'Stopping Salt Typhoon: Both Layers',
    'Config-resident adversary. Preemptive catches config drift. ACECARD catches behavioral drift.')

rect(s, Inches(1.2), Inches(1.5), Inches(1.3), Inches(0.4), fill=NAVY)
txt(s, Inches(1.25), Inches(1.53), Inches(1.2), Inches(0.35), 'PHASE', 9, WHITE, True, PP_ALIGN.CENTER)
rect(s, Inches(2.5), Inches(1.5), Inches(3.0), Inches(0.4), fill=NAVY)
txt(s, Inches(2.55), Inches(1.53), Inches(2.9), Inches(0.35), 'ADVERSARY ACTION', 9, WHITE, True, PP_ALIGN.CENTER)
rect(s, Inches(5.5), Inches(1.5), Inches(3.4), Inches(0.4), fill=BLUE)
txt(s, Inches(5.55), Inches(1.53), Inches(3.3), Inches(0.35), '22CT PREEMPTIVE', 9, WHITE, True, PP_ALIGN.CENTER)
rect(s, Inches(8.9), Inches(1.5), Inches(3.3), Inches(0.4), fill=TEAL)
txt(s, Inches(8.95), Inches(1.53), Inches(3.2), Inches(0.35), 'ACECARD DETECTS', 9, WHITE, True, PP_ALIGN.CENTER)

rows = [
    ('Edge CVE\nExploit', 'CVE-2023-20198 IOS XE;\nlevel-15 admin obtained',
     'WAF blocks /webui/ public;\nmgmt-plane ACL verified',
     'Network: anomalous mgmt\nsource IP, drift 0.09'),
    ('Account\nCreation', 'Local priv account;\nAAA config modified',
     'Config Drift detects\nbaseline diff in <1 hour',
     'Identity: graph hash change;\nnew SID linked to router'),
    ('TACACS+\nCapture', 'On-box PCAP TCP/49;\ncreds harvested',
     'Capture restricted to\nauthorized roles; mgmt VRF',
     'Process: tcpdump outside\nchange window, CUSUM alarm'),
    ('JumbledPath\nPivot', 'ELF binary chains across\ncompromised routers',
     'Outbound to C2 blocked\nall internal zones',
     'Composite: same uuid\nacross 4 router hops'),
    ('GRE Tunnel\n+ Exfil', 'GRE tunnel; slow CDR\nexfil over months',
     'GRE peer auth verified;\nCALEA-zone seg proven',
     'data_exfiltration 0.44;\nslow drift: CUSUM alarm'),
]
y = Inches(1.95)
for i, (phase, attack, preempt, acecard) in enumerate(rows):
    bg = LGRAY if i % 2 == 0 else WHITE
    rect(s, Inches(1.2), y, Inches(11.0), Inches(0.88), fill=bg)
    txt(s, Inches(1.3), y+Inches(0.1), Inches(1.1), Inches(0.7), phase, 9.5, NAVY, True, PP_ALIGN.CENTER)
    txt(s, Inches(2.6), y+Inches(0.1), Inches(2.8), Inches(0.7), attack, 9.5, INK)
    txt(s, Inches(5.6), y+Inches(0.1), Inches(3.2), Inches(0.7), preempt, 9.5, BLUE)
    txt(s, Inches(9.0), y+Inches(0.1), Inches(3.1), Inches(0.7), acecard, 9.5, TEAL)
    y += Inches(0.92)

# SLIDE 18: CUSUM Deep Dive
s = blank()
slide_title(s, 'Why CUSUM Catches What Thresholds Cannot',
    'Volt Typhoon never crosses any fixed threshold. CUSUM accumulates small drifts until alarm.')

# Left panel: Algorithm
rrect(s, Inches(1.2), Inches(1.6), Inches(5.5), Inches(2.5), fill=LTEAL)
rect(s, Inches(1.2), Inches(1.6), Inches(5.5), Inches(0.06), fill=TEAL)
icon_circle(s, Inches(1.5), Inches(1.8), ICON_STAR, TEAL, Inches(0.4))
txt(s, Inches(2.0), Inches(1.75), Inches(4.5), Inches(0.35), 'THE ALGORITHM', 12, TEAL, True)
tf = txt(s, Inches(1.5), Inches(2.3), Inches(4.9), Inches(1.5),
    'S(t) = max(0, S(t-1) + drift(t) - μ - k)', 13, NAVY, True, fn='Consolas')
add_p(tf, '', 6)
add_p(tf, 'S(t) = cumulative sum statistic', 10, INK)
add_p(tf, 'drift(t) = cosine distance from baseline', 10, INK)
add_p(tf, 'μ = expected baseline drift (calibrated)', 10, INK)
add_p(tf, 'k = 0.5σ (allowable noise)', 10, INK)
add_p(tf, 'ALARM when S(t) > h = 4σ', 11, RED, True)

# Right panel: Example
rrect(s, Inches(6.9), Inches(1.6), Inches(5.5), Inches(2.5), fill=LBLUE)
rect(s, Inches(6.9), Inches(1.6), Inches(5.5), Inches(0.06), fill=BLUE)
icon_circle(s, Inches(7.2), Inches(1.8), ICON_BOLT, BLUE, Inches(0.4))
txt(s, Inches(7.7), Inches(1.75), Inches(4.5), Inches(0.35), 'VOLT TYPHOON EXAMPLE', 12, BLUE, True)
tf = txt(s, Inches(7.2), Inches(2.3), Inches(4.9), Inches(1.5),
    'Baseline: μ=0.03, σ=0.03', 10.5, INK)
add_p(tf, 'k = 0.015  |  h = 0.12', 10.5, NAVY, True)
add_p(tf, '', 4)
add_p(tf, 'APT drift/window = 0.06 (below ALL thresholds)', 10.5, INK)
add_p(tf, 'Fixed threshold = 0.15 → NEVER fires', 10.5, RED, True)
add_p(tf, '', 4)
add_p(tf, 'CUSUM: excess = 0.06 - 0.03 - 0.015 = 0.015', 10.5, INK)
add_p(tf, 'Alarm in 0.12/0.015 = 8 windows = 8 hours', 10.5, TEAL, True)

# Bottom comparison
rrect(s, Inches(1.2), Inches(4.4), Inches(5.5), Inches(2.0), fill=LRED)
rect(s, Inches(1.2), Inches(4.4), Inches(5.5), Inches(0.06), fill=RED)
icon_circle(s, Inches(1.5), Inches(4.6), ICON_CROSS, RED, Inches(0.4))
txt(s, Inches(2.0), Inches(4.55), Inches(4.5), Inches(0.35), 'FIXED THRESHOLD (Current UEBA)', 11, RED, True)
tf = txt(s, Inches(1.5), Inches(5.1), Inches(4.9), Inches(1.0),
    'Volt Typhoon per-window drift: 0.01 – 0.04', 10.5, INK)
add_p(tf, 'NEVER crosses 0.15. Attacker dwells for YEARS.', 10.5, RED, True)
add_p(tf, 'Every commercial UEBA uses fixed thresholds.', 10.5, INK)

rrect(s, Inches(6.9), Inches(4.4), Inches(5.5), Inches(2.0), fill=LTEAL)
rect(s, Inches(6.9), Inches(4.4), Inches(5.5), Inches(0.06), fill=TEAL)
icon_circle(s, Inches(7.2), Inches(4.6), ICON_CHECK, TEAL, Inches(0.4))
txt(s, Inches(7.7), Inches(4.55), Inches(4.5), Inches(0.35), 'CUSUM (Our Approach)', 11, TEAL, True)
tf = txt(s, Inches(7.2), Inches(5.1), Inches(4.9), Inches(1.0),
    'Same drift ACCUMULATES over windows.', 10.5, INK)
add_p(tf, 'CUSUM alarm in 4–8 hours. Dwell: YEARS → HOURS.', 10.5, TEAL, True)
add_p(tf, 'Minimax optimal (Lorden 1971, Moustakides 1986).', 10.5, INK)


# ════════════════════════════════════════════════════════════════════════
# SECTION 5: HOW THEY WORK TOGETHER
# ══════════════���═════════════════════════════════════════════════════════
section_header('5. Architecture & Integration',
    'How the two layers feed each other. High-level architecture and data flow.')

# SLIDE 20: Integration Points
s = blank()
slide_title(s, 'Four Integration Points: Better Together',
    'Each layer makes the other stronger. Shared vocabulary: MITRE ATT&CK technique IDs.')

points = [
    (ICON_ARROW_R, 'Preemptive Coverage Map → ACECARD Sensitivity', BLUE,
     'Preemptive publishes which MITRE techniques are config-blocked. ACECARD raises thresholds on covered concepts (less noise) and lowers on uncovered (more sensitivity).'),
    (ICON_SHIELD, 'Preemptive Gaps → ACECARD Trust Floor', BLUE,
     'Unmitigated control gap found → ACECARD lowers trust floor for entities in that segment. ELEVATED_WATCH until gap closed. Step-up MFA applied automatically.'),
    (ICON_BRAIN, 'ACECARD Confirmed Alerts → Preemptive Threat Feed', TEAL,
     'Analyst-confirmed TPs (with MITRE IDs + entity context) feed into Preemptive Attack Intelligence. Field-observed TTPs drive preemptive policy changes.'),
    (ICON_CHAIN, 'Shared Vocabulary: MITRE IDs + Entity UUIDs', GOLD,
     'Both layers use MITRE ATT&CK. ACECARD outputs T1003, T1021, T1041 — same IDs Preemptive uses. Entity UUIDs become the SOAR join key.'),
]
y = Inches(1.6)
for icon, title, color, body in points:
    rrect(s, Inches(1.2), y, Inches(11.0), Inches(1.05), fill=LGRAY)
    icon_circle(s, Inches(1.5), y+Inches(0.28), icon, color, Inches(0.5))
    txt(s, Inches(2.2), y+Inches(0.08), Inches(9.8), Inches(0.3), title, 12, NAVY, True)
    txt(s, Inches(2.2), y+Inches(0.42), Inches(9.8), Inches(0.55), body, 10.5, INK)
    y += Inches(1.15)

rrect(s, Inches(1.2), Inches(6.3), Inches(11.0), Inches(0.5), fill=LGOLD)
txt(s, Inches(1.5), Inches(6.35), Inches(10.5), Inches(0.4),
    '22CT Preemptive + ACECARD is a system, not a bundle of two products.',
    11, NAVY, True, PP_ALIGN.CENTER)

# SLIDE 21: Architecture
s = blank()
slide_title(s, 'High-Level Architecture: Gabriel Nimbus Deployment',
    'Five-layer stack. K8s-native. Iron Bank hardened. IL5 / IL6 / JWICS.')

layers = [
    ('L5', 'PRESENTATION', 'Plotly Dash / React  |  7-chart dashboard  |  Alert feed  |  Kill-chain  |  ABAC admin  |  CAC/PIV', L5_BLUE),
    ('L4', 'API + INTEGRATION', 'FastAPI (async)  |  REST + WS  |  Istio mTLS  |  SOAR webhooks  |  ServiceNow/JIRA  |  RBAC', L4_BLUE),
    ('L3', 'BEHAVIORAL ENGINE', '5 Signals  |  Embed  |  Trajectory  |  CUSUM  |  Drift Direction  |  Entity Fusion  |  ABAC', L3_BLUE),
    ('L2', 'DATA PLANE', 'PostgreSQL + pgvector (HNSW)  |  TimescaleDB  |  Redis  |  S3  |  Encryption at rest (AES-256)', L2_BLUE),
    ('L1', 'INFRASTRUCTURE', 'K8s (Big Bang / RKE2)  |  Iron Bank  |  Istio  |  Fluentd  |  Prometheus + Grafana  |  FIPS 140-2/3', L1_BLUE),
]
y = Inches(1.6)
for num, name, components, color in layers:
    rect(s, Inches(1.2), y, Inches(11.0), Inches(0.88), fill=LGRAY)
    rrect(s, Inches(1.2), y, Inches(1.6), Inches(0.88), fill=color)
    txt(s, Inches(1.25), y+Inches(0.1), Inches(1.5), Inches(0.3), num, 10, WHITE, True, PP_ALIGN.CENTER)
    txt(s, Inches(1.25), y+Inches(0.4), Inches(1.5), Inches(0.35), name, 8.5, WHITE, True, PP_ALIGN.CENTER)
    txt(s, Inches(3.0), y+Inches(0.2), Inches(9.0), Inches(0.5), components, 10, INK)
    y += Inches(0.96)

rrect(s, Inches(1.2), Inches(6.4), Inches(11.0), Inches(0.45), fill=LGRAY)
txt(s, Inches(1.5), Inches(6.45), Inches(10.5), Inches(0.35),
    'Single Helm chart.  Air-gapped: same containers, local model, zero external calls.  cATO: OSCAL SSP, eMASS, STIG scans, POA&M auto.',
    9.5, NAVY, True, PP_ALIGN.CENTER)

# SLIDE 22: Data Flow
s = blank()
slide_title(s, 'Data Flow: Raw Telemetry to Actionable Alert',
    'End-to-end inside Gabriel Nimbus. No data leaves the enclave.')

steps = [
    (ICON_ARROW_R, 'INGEST', 'ECS / OCSF / STIX 2.1 from SIEM, EDR, NetFlow, IdP, K8s audit', 'Kafka / Fluentd'),
    (ICON_CLOCK, 'WINDOW', '1-hour windows per entity_uuid. Entity fusion resolves 10+ identity systems.', 'TimescaleDB'),
    (ICON_GEAR, 'SUMMARIZE', '5 signal summarizers produce structured text (metrics only). No PII.', 'Python workers'),
    (ICON_BRAIN, 'EMBED', 'Concatenated text → 1536-d vector. Default: local model. IL5: Azure OpenAI GovCloud.', 'pgvector HNSW'),
    (ICON_GRAPH, 'ANALYZE', 'Trajectory: cosine drift, velocity, accel. CUSUM 4σ. Drift direction → 8 concepts.', 'Per entity/window'),
    (ICON_SHIELD, 'ALERT', 'Alert tiers (CRITICAL→INFO) with MITRE IDs, cohort z-score. ABAC trust update.', 'WS + SOAR'),
]
y = Inches(1.6)
for i, (icon, name, desc, infra) in enumerate(steps):
    bg = LGRAY if i % 2 == 0 else WHITE
    rrect(s, Inches(1.2), y, Inches(11.0), Inches(0.75), fill=bg)
    icon_circle(s, Inches(1.4), y+Inches(0.13), icon, BLUE, Inches(0.45))
    txt(s, Inches(2.0), y+Inches(0.05), Inches(1.2), Inches(0.35), name, 11, NAVY, True)
    txt(s, Inches(3.3), y+Inches(0.1), Inches(6.5), Inches(0.5), desc, 10, INK)
    txt(s, Inches(10.0), y+Inches(0.15), Inches(2.0), Inches(0.4), infra, 9, DGRAY, False, PP_ALIGN.RIGHT)
    # Flow arrow between steps (except last)
    if i < len(steps) - 1:
        txt(s, Inches(1.55), y+Inches(0.7), Inches(0.3), Inches(0.3), '▼', 10, BLUE, False, PP_ALIGN.CENTER)
    y += Inches(0.83)

rrect(s, Inches(1.2), Inches(6.6), Inches(11.0), Inches(0.3), fill=LGRAY)
txt(s, Inches(1.5), Inches(6.62), Inches(10.5), Inches(0.25),
    '<3 sec/entity (API)  |  10K entities/hr per 4-vCPU pod  |  2K/hr (local, CPU)  |  Horizontal autoscale',
    9, NAVY, True, PP_ALIGN.CENTER)

# SLIDE 23: Coverage Matrix
s = blank()
slide_title(s, 'Coverage Matrix: Which Layer Owns What',
    'Preemptive handles config-resident threats. ACECARD handles behavior-resident threats.')

# Header row
rect(s, Inches(1.2), Inches(1.5), Inches(4.3), Inches(0.4), fill=NAVY)
txt(s, Inches(1.4), Inches(1.53), Inches(4.0), Inches(0.35), 'Adversary Capability', 10, WHITE, True)
rect(s, Inches(5.5), Inches(1.5), Inches(1.5), Inches(0.4), fill=NAVY)
txt(s, Inches(5.55), Inches(1.53), Inches(1.4), Inches(0.35), 'Owner', 10, WHITE, True, PP_ALIGN.CENTER)
rect(s, Inches(7.0), Inches(1.5), Inches(5.2), Inches(0.4), fill=NAVY)
txt(s, Inches(7.1), Inches(1.53), Inches(5.0), Inches(0.35), 'Why This Layer', 10, WHITE, True)

mx = [
    ('Edge-device CVE exploitation', 'PREEMPT', 'Formal proof: patch enforced or compensating control', BLUE),
    ('Misconfigured / shadowed FW rules', 'PREEMPT', 'Math model finds shadows across multi-vendor config', BLUE),
    ('Router ACL tampering / config drift', 'PREEMPT', 'Config drift re-verifies hourly vs approved baseline', BLUE),
    ('Compromised valid accounts', 'ACECARD', 'Policy-allowed traffic; only drift reveals abuse', TEAL),
    ('Living-off-the-Land binaries', 'ACECARD', 'lotl_execution concept + process signal structural shift', TEAL),
    ('Cross-domain lateral movement', 'ACECARD', 'Entity fusion joins identities into one drift trajectory', TEAL),
    ('Slow APT dwell (months/years)', 'ACECARD', 'CUSUM accumulates sub-threshold drift; alarm in hours', TEAL),
    ('C2 to known malicious infra', 'BOTH', 'Preemptive blocks known; ACECARD catches unknown via beacon', NAVY),
    ('Zero-day / unmapped TTPs', 'ACECARD', 'Behavioral anomaly detected even with no MITRE label', TEAL),
    ('Insider threat / data hoarding', 'ACECARD', 'insider_hoarding (T1074) + cohort outlier detection', TEAL),
]
y = Inches(1.95)
for i, (cap, layer, why, color) in enumerate(mx):
    bg = LGRAY if i % 2 == 0 else WHITE
    rect(s, Inches(1.2), y, Inches(11.0), Inches(0.42), fill=bg)
    txt(s, Inches(1.4), y+Pt(3), Inches(4.0), Inches(0.36), cap, 10, INK)
    rrect(s, Inches(5.55), y+Inches(0.05), Inches(1.4), Inches(0.28), fill=color)
    txt(s, Inches(5.6), y+Inches(0.06), Inches(1.3), Inches(0.26), layer, 9, WHITE, True, PP_ALIGN.CENTER)
    txt(s, Inches(7.1), y+Pt(3), Inches(5.0), Inches(0.36), why, 9.5, INK)
    y += Inches(0.44)


# ════════════════════════════════════════════════════════════════════════
# SECTION 6: TIMELINE
# ════════════════════════════════════════════════════════════════════════
section_header('6. MVP Timeline & Qualifications',
    '90 days from authorization to first detection. Phased deployment with incremental capability.')

# SLIDE 25: Timeline
s = blank()
slide_title(s, '90 Days: Authorization to First Detection',
    'Assumes AO-approved pilot authorization or existing cATO reciprocity.')

phases = [
    ('PHASE 1', 'Weeks 1–4', 'FOUNDATION', TEAL,
     ['Deploy K8s containers (Iron Bank) into Gabriel Nimbus',
      'Connect SIEM/EDR/NetFlow via ECS/OCSF connectors',
      'Configure entity fusion (AD/AAD/Okta/AWS)',
      'Data flowing, entities resolved, baselines computing']),
    ('PHASE 2', 'Weeks 5–8', 'DETECTION', BLUE,
     ['Enable 5-signal embedding pipeline + trajectory',
      'Activate CUSUM 4σ + drift direction (8 concepts)',
      'Configure cohort definitions from AD/LDAP',
      'Alerts firing, drift direction operational']),
    ('PHASE 3', 'Weeks 9–12', 'OPERATIONS', NAVY,
     ['SOC integration (SOAR, ServiceNow)',
      'Threshold tuning via analyst TP/FP feedback',
      'ABAC enabled + Volt/Salt scenario replay',
      'Full capability, ABAC enforcing, validated']),
    ('PHASE 4', 'Week 13+', 'PREEMPTIVE', GOLD,
     ['Integrate 22CT Preemptive config verification',
      'Connect NGFW/IPS/SASE for formal analysis',
      'Trust-floor coupling + continuous re-verification',
      'Both layers operational, closed-loop defense']),
]
y = Inches(1.5)
for phase, weeks, name, color, items in phases:
    rrect(s, Inches(1.2), y, Inches(11.0), Inches(1.15), fill=LGRAY)
    rect(s, Inches(1.2), y, Inches(0.08), Inches(1.15), fill=color)
    # Phase badge
    rrect(s, Inches(1.5), y+Inches(0.08), Inches(2.0), Inches(0.35), fill=color)
    txt(s, Inches(1.55), y+Inches(0.1), Inches(1.9), Inches(0.3), f'{phase}: {weeks}', 9.5, WHITE, True, PP_ALIGN.CENTER)
    txt(s, Inches(3.7), y+Inches(0.1), Inches(2.0), Inches(0.3), name, 12, color, True)
    # Bullet items
    tf = txt(s, Inches(1.5), y+Inches(0.48), Inches(10.5), Inches(0.65), '', 10, INK)
    for j, item in enumerate(items):
        prefix = '✓  ' if j == len(items)-1 else '•  '
        c = NAVY if j == len(items)-1 else INK
        b = j == len(items)-1
        add_p(tf, f'{prefix}{item}', 10, c, b, Pt(3))
    y += Inches(1.22)

rrect(s, Inches(1.2), Inches(6.4), Inches(11.0), Inches(0.45), fill=LRED)
txt(s, Inches(1.5), Inches(6.45), Inches(10.5), Inches(0.35),
    'RISKS:  ATO timeline (mitigated: cATO + reciprocity)  |  Data quality (mitigated: 2-week connector sprint)  |  Red team within 60 days',
    9.5, RED, False, PP_ALIGN.CENTER)

# SLIDE 26: Key Numbers + Qualifications
s = blank()
slide_title(s, 'Key Numbers & 22CT Qualifications')

# Left: Key numbers with icon badges
left = [
    (ICON_WARNING, '99%', 'FW breaches = misconfiguration (Gartner)'),
    (ICON_CHECK, '0 / 0', 'FP / FN within formal model scope'),
    (ICON_CLOCK, '<1 hour', 'Config drift detection time'),
    (ICON_BRAIN, '1,536', 'Embedding dimensions per signal'),
    (ICON_GRAPH, '5', 'Behavioral signals per entity'),
    (ICON_NET, '8', 'MITRE threat concepts'),
    (ICON_STAR, '4σ', 'CUSUM alarm threshold (configurable)'),
    (ICON_BOLT, '<3 sec', 'Event-to-analysis per entity'),
    (ICON_SHIELD, '90 days', 'Authorization to first detection'),
]
y = Inches(1.4)
for icon, num, meaning in left:
    icon_circle(s, Inches(1.3), y+Inches(0.05), icon, BLUE, Inches(0.38))
    rect(s, Inches(1.8), y, Inches(1.5), Inches(0.48), fill=NAVY)
    txt(s, Inches(1.85), y+Inches(0.07), Inches(1.4), Inches(0.35),
        num, 11.5, WHITE, True, PP_ALIGN.CENTER, fn='Consolas')
    txt(s, Inches(3.5), y+Inches(0.1), Inches(3.0), Inches(0.35), meaning, 10, INK)
    y += Inches(0.55)

# Right: Qualifications
right_quals = [
    (ICON_SHIELD, '$90M Army SOC/MDR Contract',
     '800+ cleared analysts. 24x7x365. Proven IR at scale.'),
    (ICON_STAR, 'USAF $108M | FBI TSC $56M | NAVAIR $145M',
     'Active federal cyber contracts across DoD and IC.'),
    (ICON_LOCK, 'Cleared Workforce (Secret – TS/SCI)',
     'MITRE ATT&CK, NIST 800-53, DoD RMF trained.'),
    (ICON_GEAR, 'IL5 / IL6 / JWICS Experience',
     'ATO artifacts, air-gapped ops. cATO alignment.'),
    (ICON_CHECK, 'ACECARD: Built and Tested',
     'Working prototype. Synthetic Volt/Salt replay validated.'),
    (ICON_BRAIN, 'Innovation Partner: Rigor AI',
     'Preemptive layer powered by Rigor AI technology.'),
]
y = Inches(1.4)
for icon, title, desc in right_quals:
    rrect(s, Inches(7.0), y, Inches(5.3), Inches(0.75), fill=LGRAY)
    icon_circle(s, Inches(7.2), y+Inches(0.15), icon, GOLD, Inches(0.4))
    txt(s, Inches(7.75), y+Inches(0.05), Inches(4.3), Inches(0.3), title, 10.5, NAVY, True)
    txt(s, Inches(7.75), y+Inches(0.38), Inches(4.3), Inches(0.3), desc, 9.5, INK)
    y += Inches(0.82)

# Compliance bar
rrect(s, Inches(1.2), Inches(6.5), Inches(11.0), Inches(0.35), fill=LGRAY)
txt(s, Inches(1.5), Inches(6.53), Inches(10.5), Inches(0.3),
    'NIST 800-53 Rev 5 (SI-4, AU-6, IR-4, SC-7)  |  DoD ZT RA Pillars 4 & 7  |  DISA STIGs  |  RMF/eMASS  |  OSCAL SSP  |  DCO-IDM',
    9, NAVY, True, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 27: NEITHER LAYER ALONE IS SUFFICIENT
# ════════════════════════════════════════════════════════════════════════
s = blank()
slide_title(s, 'Why Two Layers — Neither Alone Is Sufficient')

# Left: Preemptive alone
rrect(s, Inches(1.2), Inches(1.5), Inches(5.3), Inches(2.2), fill=LBLUE)
rect(s, Inches(1.2), Inches(1.5), Inches(5.3), Inches(0.5), fill=BLUE)
icon_circle(s, Inches(1.4), Inches(1.58), ICON_SHIELD, BLUE, Inches(0.35))
txt(s, Inches(1.85), Inches(1.55), Inches(4.5), Inches(0.4),
    'PREEMPTIVE ALONE', 12, WHITE, True)
tf = txt(s, Inches(1.5), Inches(2.1), Inches(4.8), Inches(1.5),
    '', 11, INK)
add_p(tf, ICON_CHECK + '  Eliminates config gaps (PROVEN)', 11, TEAL, sp=Pt(4))
add_p(tf, ICON_CHECK + '  Blocks known attack paths exhaustively', 11, TEAL, sp=Pt(4))
add_p(tf, ICON_CHECK + '  Zero FP/FN within model scope', 11, TEAL, sp=Pt(4))
add_p(tf, '', 6)
add_p(tf, ICON_CROSS + '  Cannot detect valid-credential abuse', 11, RED, sp=Pt(4))
add_p(tf, ICON_CROSS + '  Blind to LOTL tradecraft (tools are legit)', 11, RED, sp=Pt(4))
add_p(tf, ICON_CROSS + '  Cannot see behavioral evolution', 11, RED, sp=Pt(4))

# Right: ACECARD alone
rrect(s, Inches(6.8), Inches(1.5), Inches(5.3), Inches(2.2), fill=LTEAL)
rect(s, Inches(6.8), Inches(1.5), Inches(5.3), Inches(0.5), fill=TEAL)
icon_circle(s, Inches(7.0), Inches(1.58), ICON_BRAIN, TEAL, Inches(0.35))
txt(s, Inches(7.45), Inches(1.55), Inches(4.5), Inches(0.4),
    'ACECARD ALONE', 12, WHITE, True)
tf = txt(s, Inches(7.1), Inches(2.1), Inches(4.8), Inches(1.5),
    '', 11, INK)
add_p(tf, ICON_CHECK + '  Detects behavioral evolution over time', 11, TEAL, sp=Pt(4))
add_p(tf, ICON_CHECK + '  Catches LOTL via process/network drift', 11, TEAL, sp=Pt(4))
add_p(tf, ICON_CHECK + '  Works even without signatures', 11, TEAL, sp=Pt(4))
add_p(tf, '', 6)
add_p(tf, ICON_CROSS + '  Only detects AFTER adversary is inside', 11, RED, sp=Pt(4))
add_p(tf, ICON_CROSS + '  Cannot prevent config gaps from existing', 11, RED, sp=Pt(4))
add_p(tf, ICON_CROSS + '  Reactive, not preemptive', 11, RED, sp=Pt(4))

# Center: TOGETHER
rrect(s, Inches(1.2), Inches(4.0), Inches(11.0), Inches(2.5), fill=LGOLD)
rect(s, Inches(1.2), Inches(4.0), Inches(11.0), Inches(0.06), fill=GOLD)
txt(s, Inches(1.5), Inches(4.2), Inches(10.5), Inches(0.4),
    'TOGETHER: DEFENSE IN DEPTH', 14, NAVY, True, PP_ALIGN.CENTER)

tf = txt(s, Inches(1.5), Inches(4.7), Inches(10.5), Inches(1.5), '', 12, INK)
add_p(tf, 'Preemptive eliminates the attack surface in configuration — adversary cannot enter through config gaps.', 12, INK, sp=Pt(6))
add_p(tf, 'ACECARD detects behavioral evolution — catches adversaries who bypass via valid credentials & legit tools.', 12, INK, sp=Pt(6))
add_p(tf, '', 8)
add_p(tf, 'Attack surface NEVER EXISTS  ...or...  attack detected in HOURS, not years.', 13, NAVY, True, sp=Pt(4))

txt(s, Inches(1.5), Inches(6.5), Inches(10.5), Inches(0.4),
    'Preemptive closes the door.  ACECARD watches the room.  Nothing gets through unnoticed.',
    11, GOLD, True, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 28: HOW ACECARD SEES BEHAVIOR (Simple explanation)
# ════════════════════════════════════════════════════════════════════════
s = blank()
slide_title(s, 'How ACECARD Reads Behavior: Five Questions',
    'All tied to the same person via their login session. Together = complete behavioral fingerprint.')

# Simple 5-signal table
signals_simple = [
    (ICON_LOCK, 'AUTH', 'WHO logged in?', 'Where from, when, how often, did they fail?', BLUE),
    (ICON_GEAR, 'PROCESS', 'WHAT did they run?', 'Which programs, scripts, commands on their machines', BLUE),
    (ICON_NET, 'NETWORK', 'WHERE did they connect?', 'Which IPs, how much data sent, any periodic beaconing?', TEAL),
    (ICON_GRAPH, 'FILE', 'WHAT did they touch?', 'Files opened, copied, deleted, archived, encrypted', TEAL),
    (ICON_STAR, 'IDENTITY', 'WHAT POWER did they gain?', 'New privileges, group memberships, role changes', NAVY),
]
y = Inches(1.6)
for icon, sig, question, detail, color in signals_simple:
    rrect(s, Inches(1.2), y, Inches(11.0), Inches(0.85), fill=LGRAY)
    icon_circle(s, Inches(1.5), y+Inches(0.18), icon, color, Inches(0.48))
    rect(s, Inches(2.15), y+Inches(0.12), Inches(1.6), Inches(0.55), fill=color)
    txt(s, Inches(2.2), y+Inches(0.2), Inches(1.5), Inches(0.4), sig, 12, WHITE, True, PP_ALIGN.CENTER)
    txt(s, Inches(3.95), y+Inches(0.08), Inches(3.5), Inches(0.35), question, 13, NAVY, True)
    txt(s, Inches(3.95), y+Inches(0.45), Inches(7.0), Inches(0.35), detail, 10.5, INK)
    y += Inches(0.92)

# Normal vs compromised
rrect(s, Inches(1.2), Inches(6.3), Inches(5.2), Inches(0.55), fill=LTEAL)
txt(s, Inches(1.5), Inches(6.35), Inches(4.8), Inches(0.45),
    'Normal: Office, SharePoint, reports, 9-5', 10.5, TEAL, True)
rrect(s, Inches(6.7), Inches(6.3), Inches(5.5), Inches(0.55), fill=LRED)
txt(s, Inches(7.0), Inches(6.35), Inches(5.1), Inches(0.45),
    'Compromised: ntdsutil, 18 servers, NTDS.dit, 2am', 10.5, RED, True)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 29: LIVING-OFF-THE-LAND EXPLAINED
# ════════════════════════════════════════════════════════════════════════
s = blank()
slide_title(s, 'Living-off-the-Land: The Invisible Attack',
    'Volt Typhoon brings NO malware. They use tools Microsoft already installed. Antivirus cannot see them.')

# Left: Traditional vs LOTL
rrect(s, Inches(1.2), Inches(1.5), Inches(5.3), Inches(1.6), fill=LTEAL)
rect(s, Inches(1.2), Inches(1.5), Inches(5.3), Inches(0.06), fill=TEAL)
icon_circle(s, Inches(1.5), Inches(1.7), ICON_CHECK, TEAL, Inches(0.4))
txt(s, Inches(2.0), Inches(1.65), Inches(4.3), Inches(0.35),
    'TRADITIONAL HACKER', 12, TEAL, True)
tf = txt(s, Inches(1.5), Inches(2.1), Inches(4.8), Inches(0.8), '', 11, INK)
add_p(tf, 'Brings their own weapon (malware, virus, backdoor)', 11, INK, sp=Pt(4))
add_p(tf, 'Antivirus sees unknown file ' + ICON_ARROW_R + ' BLOCKED', 11, TEAL, True, sp=Pt(4))

rrect(s, Inches(6.8), Inches(1.5), Inches(5.3), Inches(1.6), fill=LRED)
rect(s, Inches(6.8), Inches(1.5), Inches(5.3), Inches(0.06), fill=RED)
icon_circle(s, Inches(7.1), Inches(1.7), ICON_WARNING, RED, Inches(0.4))
txt(s, Inches(7.6), Inches(1.65), Inches(4.3), Inches(0.35),
    'LIVING-OFF-THE-LAND (VOLT TYPHOON)', 11, RED, True)
tf = txt(s, Inches(7.1), Inches(2.1), Inches(4.8), Inches(0.8), '', 11, INK)
add_p(tf, 'Brings NOTHING. Uses tools Windows already has.', 11, INK, sp=Pt(4))
add_p(tf, 'Antivirus sees Microsoft binary ' + ICON_ARROW_R + ' "It\'s fine"', 11, RED, True, sp=Pt(4))

# Tool table
tools_data = [
    ('ntdsutil.exe', 'Backup Active Directory', 'Steal ALL domain passwords'),
    ('wmic.exe', 'Query system info remotely', 'Run commands on other machines'),
    ('netsh.exe', 'Configure network settings', 'Open firewall holes'),
    ('certutil.exe', 'Manage certificates', 'Download files from internet'),
    ('powershell.exe', 'Automation / scripting', 'Run attack scripts'),
    ('rundll32.exe', 'Run DLL files', 'Execute malicious code'),
]
# Header
rect(s, Inches(1.2), Inches(3.3), Inches(3.0), Inches(0.35), fill=NAVY)
txt(s, Inches(1.3), Inches(3.33), Inches(2.8), Inches(0.3), 'Tool (signed by Microsoft)', 9, WHITE, True)
rect(s, Inches(4.2), Inches(3.3), Inches(3.8), Inches(0.35), fill=TEAL)
txt(s, Inches(4.3), Inches(3.33), Inches(3.6), Inches(0.3), 'Legitimate Purpose', 9, WHITE, True)
rect(s, Inches(8.0), Inches(3.3), Inches(4.2), Inches(0.35), fill=RED)
txt(s, Inches(8.1), Inches(3.33), Inches(4.0), Inches(0.3), 'What Attacker Uses It For', 9, WHITE, True)

y = Inches(3.65)
for i, (tool, legit, attack) in enumerate(tools_data):
    bg = LGRAY if i % 2 == 0 else WHITE
    rrect(s, Inches(1.2), y, Inches(11.0), Inches(0.35), fill=bg)
    txt(s, Inches(1.3), y+Inches(0.03), Inches(2.8), Inches(0.3), tool, 10, NAVY, True, fn='Consolas')
    txt(s, Inches(4.3), y+Inches(0.03), Inches(3.6), Inches(0.3), legit, 10, INK)
    txt(s, Inches(8.1), y+Inches(0.03), Inches(4.0), Inches(0.3), attack, 10, RED)
    y += Inches(0.37)

# Bottom insight
rrect(s, Inches(1.2), Inches(5.95), Inches(11.0), Inches(0.9), fill=LGOLD)
rect(s, Inches(1.2), Inches(5.95), Inches(0.06), Inches(0.9), fill=GOLD)
tf = txt(s, Inches(1.5), Inches(6.0), Inches(10.5), Inches(0.8),
    'The TOOLS are legitimate.  The PERSON running them is not an admin.', 12, NAVY, True)
add_p(tf, 'ACECARD tracks lolbin_count per user. John the accountant running ntdsutil = alarm.  John the sysadmin = normal.', 11, INK, sp=Pt(4))
add_p(tf, 'It\'s not about WHAT ran.  It\'s about WHO ran it.', 11, GOLD, True, sp=Pt(4))

# ════════════════════════════════════════════════════════════════════════
# SLIDE 30: WHAT THE METRICS MEAN (simplified walkthrough)
# ════════════════════════════════════════════════════════════════════════
s = blank()
slide_title(s, 'What the Alert Numbers Actually Mean',
    'Plain-English walkthrough of three real detection metrics.')

# Metric 1: unique_dest_ips
rrect(s, Inches(1.2), Inches(1.5), Inches(11.0), Inches(1.5), fill=LBLUE)
rect(s, Inches(1.2), Inches(1.5), Inches(0.06), Inches(1.5), fill=BLUE)
icon_circle(s, Inches(1.5), Inches(1.65), ICON_NET, BLUE, Inches(0.45))
txt(s, Inches(2.1), Inches(1.55), Inches(4.0), Inches(0.35),
    'unique_dest_ips:  3 ' + ICON_ARROW_R + ' 18', 12, NAVY, True)
tf = txt(s, Inches(2.1), Inches(1.95), Inches(9.8), Inches(0.9), '', 10.5, INK)
add_p(tf, 'How many different machines did this person REACH OUT TO.  Not where they logged in from — where they connected TO.', 10.5, INK, sp=Pt(3))
add_p(tf, 'Normal: John talks to 3 servers daily (file, email, SharePoint). Always the same 3.', 10.5, TEAL, sp=Pt(3))
add_p(tf, 'Attack: John suddenly connects to 18 servers — domain controllers, HR, finance, database servers. Lateral movement.', 10.5, RED, sp=Pt(3))

# Metric 2: lolbin_count
rrect(s, Inches(1.2), Inches(3.2), Inches(11.0), Inches(1.5), fill=LTEAL)
rect(s, Inches(1.2), Inches(3.2), Inches(0.06), Inches(1.5), fill=TEAL)
icon_circle(s, Inches(1.5), Inches(3.35), ICON_GEAR, TEAL, Inches(0.45))
txt(s, Inches(2.1), Inches(3.25), Inches(4.0), Inches(0.35),
    'lolbin_count:  0 ' + ICON_ARROW_R + ' 7', 12, NAVY, True)
tf = txt(s, Inches(2.1), Inches(3.65), Inches(9.8), Inches(0.9), '', 10.5, INK)
add_p(tf, 'How many Living-off-the-Land binaries did this person run (legit tools used for attack — ntdsutil, wmic, netsh, etc.)', 10.5, INK, sp=Pt(3))
add_p(tf, 'Normal: John the accountant runs 0 admin tools. Ever. That\'s his baseline.', 10.5, TEAL, sp=Pt(3))
add_p(tf, 'Attack: 7 admin tools in one hour. Same tools. But John never ran them before. WHO ran it matters, not WHAT ran.', 10.5, RED, sp=Pt(3))

# Metric 3: service_account_usage
rrect(s, Inches(1.2), Inches(4.9), Inches(11.0), Inches(1.5), fill=LRED)
rect(s, Inches(1.2), Inches(4.9), Inches(0.06), Inches(1.5), fill=RED)
icon_circle(s, Inches(1.5), Inches(5.05), ICON_LOCK, RED, Inches(0.45))
txt(s, Inches(2.1), Inches(4.95), Inches(5.0), Inches(0.35),
    'service_account_usage:  0 ' + ICON_ARROW_R + ' 4', 12, NAVY, True)
tf = txt(s, Inches(2.1), Inches(5.35), Inches(9.8), Inches(0.9), '', 10.5, INK)
add_p(tf, 'Service accounts (svc_backup, svc_sql) are for applications, not people. High privilege + weak security (often no MFA).', 10.5, INK, sp=Pt(3))
add_p(tf, 'Normal: John never uses service accounts — he uses his own account.', 10.5, TEAL, sp=Pt(3))
add_p(tf, 'Attack: John logs in as svc_backup 4 times. Getting high privilege through the back door. Persistence.', 10.5, RED, sp=Pt(3))

# Bottom
rrect(s, Inches(1.2), Inches(6.55), Inches(11.0), Inches(0.35), fill=LGRAY)
txt(s, Inches(1.5), Inches(6.58), Inches(10.5), Inches(0.3),
    'Each number alone = slightly weird.  All three moving in the same direction = attack trajectory.  CUSUM catches the consistency.',
    10, NAVY, True, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 31: CUSUM — THE BUCKET ANALOGY
# ════════════════════════════════════════════════════════════════════════
s = blank()
slide_title(s, 'CUSUM: The Bucket Under a Dripping Faucet',
    'Why current tools miss slow attackers and CUSUM catches them.')

# Left: the analogy
rrect(s, Inches(1.2), Inches(1.5), Inches(5.3), Inches(3.0), fill=LTEAL)
rect(s, Inches(1.2), Inches(1.5), Inches(5.3), Inches(0.06), fill=TEAL)
txt(s, Inches(1.5), Inches(1.7), Inches(4.8), Inches(0.35),
    'THE BUCKET ANALOGY', 12, TEAL, True)
tf = txt(s, Inches(1.5), Inches(2.1), Inches(4.8), Inches(2.2), '', 11, INK)
add_p(tf, 'Bucket = CUSUM statistic (starts empty)', 11, INK, sp=Pt(6))
add_p(tf, 'Drip = behavioral drift each hour', 11, INK, sp=Pt(6))
add_p(tf, 'Drain hole = allowance for normal noise', 11, INK, sp=Pt(6))
add_p(tf, 'Overflow line = alarm threshold', 11, INK, sp=Pt(6))
add_p(tf, '', 10)
add_p(tf, 'Normal: drip <= drain. Bucket never fills.', 11, TEAL, True, sp=Pt(6))
add_p(tf, 'Attack: drip slightly > drain. Fills up.', 11, RED, True, sp=Pt(6))
add_p(tf, 'After 8 hours: OVERFLOW = ALARM.', 11, RED, True, sp=Pt(6))

# Right: current tools
rrect(s, Inches(6.8), Inches(1.5), Inches(5.3), Inches(3.0), fill=LRED)
rect(s, Inches(6.8), Inches(1.5), Inches(5.3), Inches(0.06), fill=RED)
txt(s, Inches(7.1), Inches(1.7), Inches(4.8), Inches(0.35),
    'CURRENT TOOLS: NO BUCKET', 12, RED, True)
tf = txt(s, Inches(7.1), Inches(2.1), Inches(4.8), Inches(2.2), '', 11, INK)
add_p(tf, 'They look at each drip individually:', 11, INK, sp=Pt(6))
add_p(tf, '', 6)
add_p(tf, '"Is THIS drip a flood? No."', 11, DGRAY, False, sp=Pt(4))
add_p(tf, '"Is THIS drip a flood? No."', 11, DGRAY, False, sp=Pt(4))
add_p(tf, '"Is THIS drip a flood? No."', 11, DGRAY, False, sp=Pt(4))
add_p(tf, '...forever "No."', 11, DGRAY, False, sp=Pt(4))
add_p(tf, '', 8)
add_p(tf, 'They never accumulate evidence.', 11, RED, True, sp=Pt(4))
add_p(tf, 'Slow attacker = forever invisible.', 11, RED, True, sp=Pt(4))

# Bottom: the key difference
rrect(s, Inches(1.2), Inches(4.8), Inches(11.0), Inches(1.8), fill=LGRAY)
rect(s, Inches(1.2), Inches(4.8), Inches(11.0), Inches(0.06), fill=NAVY)
txt(s, Inches(1.5), Inches(5.0), Inches(10.5), Inches(0.35),
    'THE FUNDAMENTAL DIFFERENCE', 12, NAVY, True, PP_ALIGN.CENTER)

# Two columns inside bottom box
txt(s, Inches(1.5), Inches(5.4), Inches(5.0), Inches(0.3),
    'Fixed Threshold = alarm on POSITION', 11, RED, True)
txt(s, Inches(1.5), Inches(5.7), Inches(5.0), Inches(0.4),
    '"How far are you from normal right now?"', 10.5, INK)
txt(s, Inches(1.5), Inches(6.1), Inches(5.0), Inches(0.3),
    'Attacker can control position (stay close)', 10, DGRAY)

txt(s, Inches(6.8), Inches(5.4), Inches(5.0), Inches(0.3),
    'CUSUM = alarm on INTENT', 11, TEAL, True)
txt(s, Inches(6.8), Inches(5.7), Inches(5.0), Inches(0.4),
    '"Are you consistently moving away, every hour?"', 10.5, INK)
txt(s, Inches(6.8), Inches(6.1), Inches(5.0), Inches(0.3),
    'Attacker CANNOT hide intent (must keep moving)', 10, TEAL)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 30: DRIFT DIRECTION — NORMAL VS ATTACKER
# ════════════════════════════════════════════════════════════════════════
s = blank()
slide_title(s, 'Continuous Drift in One Direction = Problem',
    'Normal users drift randomly. Attackers MUST keep moving toward their objective.')

# Normal user visual
rrect(s, Inches(1.2), Inches(1.5), Inches(5.3), Inches(2.3), fill=LTEAL)
rect(s, Inches(1.2), Inches(1.5), Inches(5.3), Inches(0.06), fill=TEAL)
icon_circle(s, Inches(1.5), Inches(1.7), ICON_CHECK, TEAL, Inches(0.4))
txt(s, Inches(2.0), Inches(1.7), Inches(4.3), Inches(0.35),
    'NORMAL USER', 12, TEAL, True)
tf = txt(s, Inches(1.5), Inches(2.2), Inches(4.8), Inches(1.4), '', 11, INK)
add_p(tf, 'Drift pattern:  ' + ICON_ARROW_R + '  ' + '←' + '  ' + ICON_ARROW_R + '  ' + ICON_ARROW_R + '  ' + '←' + '  ' + '←' + '  ' + ICON_ARROW_R + '  ' + '←', 13, TEAL, True, sp=Pt(8))
add_p(tf, '', 6)
add_p(tf, 'Random direction each day', 11, INK, sp=Pt(4))
add_p(tf, 'Cancels out over time', 11, INK, sp=Pt(4))
add_p(tf, 'CUSUM stays near zero (resets on negative)', 11, INK, sp=Pt(4))

# Attacker visual
rrect(s, Inches(6.8), Inches(1.5), Inches(5.3), Inches(2.3), fill=LRED)
rect(s, Inches(6.8), Inches(1.5), Inches(5.3), Inches(0.06), fill=RED)
icon_circle(s, Inches(7.1), Inches(1.7), ICON_WARNING, RED, Inches(0.4))
txt(s, Inches(7.6), Inches(1.7), Inches(4.3), Inches(0.35),
    'ATTACKER (VOLT TYPHOON)', 12, RED, True)
tf = txt(s, Inches(7.1), Inches(2.2), Inches(4.8), Inches(1.4), '', 11, INK)
add_p(tf, 'Drift pattern:  ' + ICON_ARROW_R + '  ' + ICON_ARROW_R + '  ' + ICON_ARROW_R + '  ' + ICON_ARROW_R + '  ' + ICON_ARROW_R + '  ' + ICON_ARROW_R + '  ' + ICON_ARROW_R + '  ' + ICON_ARROW_R, 13, RED, True, sp=Pt(8))
add_p(tf, '', 6)
add_p(tf, 'Same direction every hour', 11, INK, sp=Pt(4))
add_p(tf, 'MUST keep moving toward objective', 11, INK, sp=Pt(4))
add_p(tf, 'CUSUM accumulates ' + ICON_ARROW_R + ' ALARM in 8 hours', 11, RED, True, sp=Pt(4))

# Why attackers can't avoid this
rrect(s, Inches(1.2), Inches(4.1), Inches(11.0), Inches(2.5), fill=LGRAY)
rect(s, Inches(1.2), Inches(4.1), Inches(11.0), Inches(0.06), fill=NAVY)
txt(s, Inches(1.5), Inches(4.3), Inches(10.5), Inches(0.35),
    'WHY ATTACKERS CANNOT AVOID THIS', 12, NAVY, True, PP_ALIGN.CENTER)

tf = txt(s, Inches(1.5), Inches(4.7), Inches(10.5), Inches(1.6), '', 11, INK)
add_p(tf, 'To steal credentials: MUST run credential tools  ' + ICON_ARROW_R + '  Process signal drifts toward credential_theft', 11, INK, sp=Pt(6))
add_p(tf, 'To move laterally: MUST connect to new hosts  ' + ICON_ARROW_R + '  Network signal drifts toward lateral_movement', 11, INK, sp=Pt(6))
add_p(tf, 'To exfiltrate data: MUST send data outbound  ' + ICON_ARROW_R + '  File + Network signals drift toward exfiltration', 11, INK, sp=Pt(6))
add_p(tf, '', 8)
add_p(tf, 'They can go SLOW (and avoid thresholds). They cannot avoid MOVING.', 12, NAVY, True, sp=Pt(4))
add_p(tf, 'CUSUM catches the movement. Drift direction names the objective.', 12, TEAL, True, sp=Pt(4))


# ════════════════════════════════════════════════════════════════════════
# SLIDE 33: DUAL EMBEDDING ARCHITECTURE (from Rigor_ACECARD_Volt_Salt_Typhoon)
# ════════════════════════════════════════════════════════════════════════
s = blank()
rect(s, Inches(0), Inches(0), W, Inches(0.3), fill=TEAL)
txt(s, Inches(0.5), Inches(0.0), Inches(12.3), Inches(0.3),
    'PART 2  --  ACECARD UEBA  --  THE MATH', 11, WHITE, True)
txt(s, Inches(0.5), Inches(0.5), Inches(12.3), Inches(0.6),
    'Dual Embedding Architecture: 6 Vectors = 9,216 Dimensions', 24, NAVY, True, fn='Georgia')
txt(s, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.3),
    'Each signal individually embedded (1536-d) + one fused composite. Per-signal drift tracking.', 13, DGRAY)

# Key insight box
ACE_BG = RGBColor(0xE0, 0xF2, 0xF7)
CARD   = RGBColor(0xF7, 0xF8, 0xFA)
ORG    = RGBColor(0xF2, 0x6B, 0x1F)
BODY   = RGBColor(0x2D, 0x37, 0x48)
SUB    = RGBColor(0x5A, 0x68, 0x78)

rrect(s, Inches(0.5), Inches(1.5), Inches(12.1), Inches(0.8), fill=ACE_BG)
txt(s, Inches(0.7), Inches(1.55), Inches(11.7), Inches(0.25), 'THE KEY INSIGHT', 12, TEAL, True)
txt(s, Inches(0.7), Inches(1.8), Inches(11.7), Inches(0.4),
    'Each of the 5 signals is INDIVIDUALLY embedded at 1536-d (per-signal drift tracking). '
    'All 5 are also fused into one COMPOSITE 1536-d vector (holistic entity state). '
    'Total: 6 vectors = 9,216 dimensions per entity per window.', 11, BODY)

# 5 signal color strips
sig_colors = [RGBColor(0x1E,0x4D,0x8C), RGBColor(0x6C,0x2E,0xD9), TEAL, ORG, RED]
sig_names = ['AUTH', 'PROCESS', 'NETWORK', 'FILE', 'IDENTITY']
xp = [Inches(0.5), Inches(3.1), Inches(5.7), Inches(8.3), Inches(10.9)]
for i in range(5):
    rect(s, xp[i], Inches(2.5), Inches(2.4), Inches(0.35), fill=sig_colors[i])
    txt(s, xp[i], Inches(2.52), Inches(2.4), Inches(0.3), sig_names[i], 11, WHITE, True, PP_ALIGN.CENTER)
    rrect(s, xp[i], Inches(2.9), Inches(2.4), Inches(0.45), fill=CARD)
    txt(s, xp[i]+Inches(0.1), Inches(2.95), Inches(2.2), Inches(0.35),
        'Individual 1536-d\nPer-signal drift', 9, BODY, False, PP_ALIGN.CENTER)

# Composite bar
rect(s, Inches(0.5), Inches(3.5), Inches(12.1), Inches(0.35), fill=NAVY)
txt(s, Inches(0.7), Inches(3.52), Inches(11.7), Inches(0.3),
    'All 5 fused -> 1 COMPOSITE 1536-d vector (holistic entity behavioral state)', 11, GOLD, True, PP_ALIGN.CENTER)

# Three info boxes
rrect(s, Inches(0.5), Inches(4.1), Inches(3.8), Inches(1.5), fill=CARD)
rect(s, Inches(0.5), Inches(4.1), Inches(0.1), Inches(1.5), fill=TEAL)
txt(s, Inches(0.8), Inches(4.15), Inches(3.3), Inches(0.25), 'DIMENSIONS', 11, TEAL, True)
tf = txt(s, Inches(0.8), Inches(4.4), Inches(3.3), Inches(1.1),
    '5 individual: 5 x 1,536 = 7,680', 10, BODY)
add_p(tf, '1 composite: 1 x 1,536 = 1,536', 10, BODY)
add_p(tf, 'TOTAL: 6 x 1,536 = 9,216 dims', 11, NAVY, True)

rrect(s, Inches(4.6), Inches(4.1), Inches(4.0), Inches(1.5), fill=CARD)
rect(s, Inches(4.6), Inches(4.1), Inches(0.1), Inches(1.5), fill=ORG)
txt(s, Inches(4.9), Inches(4.15), Inches(3.5), Inches(0.25), 'WHY INDIVIDUAL + COMPOSITE?', 11, ORG, True)
tf = txt(s, Inches(4.9), Inches(4.4), Inches(3.5), Inches(1.1),
    'Composite: "entity behavior changed"', 10, BODY)
add_p(tf, 'Individual: "AUTH drove 72% of change"', 10, BODY)
add_p(tf, 'Direction: "PROCESS signal drifting', 10, BODY)
add_p(tf, '  toward lotl_execution (T1059)"', 10, NAVY, True)

rrect(s, Inches(8.8), Inches(4.1), Inches(3.8), Inches(1.5), fill=CARD)
rect(s, Inches(8.8), Inches(4.1), Inches(0.1), Inches(1.5), fill=NAVY)
txt(s, Inches(9.1), Inches(4.15), Inches(3.3), Inches(0.25), 'DEPLOYMENT OPTIONS', 11, NAVY, True)
tf = txt(s, Inches(9.1), Inches(4.4), Inches(3.3), Inches(1.1),
    'IL5: OpenAI text-embedding-3-small', 10, BODY)
add_p(tf, 'IL6: Local Phi-4 (14B) or Mistral 7B', 10, BODY)
add_p(tf, 'JWICS: Same SLM, air-gapped', 10, BODY)
add_p(tf, 'Switch: env var only, no code change', 10, NAVY, True)

# Privacy note
rrect(s, Inches(0.5), Inches(5.8), Inches(12.1), Inches(0.5), fill=CARD)
rect(s, Inches(0.5), Inches(5.8), Inches(0.1), Inches(0.5), fill=GOLD)
txt(s, Inches(0.8), Inches(5.85), Inches(11.6), Inches(0.4),
    'PRIVACY: No raw event content sent to embedding provider. Only summarized metrics. No usernames, IPs, paths. JWICS-safe.',
    10, GOLD, True)

# Bottom bar
rect(s, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.5), fill=NAVY)
txt(s, Inches(0.7), Inches(6.8), Inches(12.0), Inches(0.5),
    'NOBODY IN THE INDUSTRY embeds behavioral signals into a unified vector space with per-signal drift decomposition.',
    11.5, GOLD, True, PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 34: THE MATH OF BEHAVIORAL DRIFT (from Rigor_ACECARD_Volt_Salt_Typhoon)
# ════════════════════════════════════════════════════════════════════════
s = blank()
rect(s, Inches(0), Inches(0), W, Inches(0.3), fill=TEAL)
txt(s, Inches(0.5), Inches(0.0), Inches(12.3), Inches(0.3),
    'PART 2  --  ACECARD UEBA  --  THE MATH', 11, WHITE, True)
txt(s, Inches(0.5), Inches(0.5), Inches(12.3), Inches(0.6),
    'The Math of Behavioral Drift: E(t) - E(t-1)', 24, NAVY, True, fn='Georgia')
txt(s, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.3),
    'Each piece is well-known math. The innovation is the integration across 1536 dimensions.', 13, DGRAY)

# Four math boxes A-D (exact replica of reference deck layout)
math_boxes = [
    ('A', 'Behavioral Embedding',
     'V_t = embed("auth: logon_count=47, failure_rate=0.032 | process: lolbin_count=0, cmdline_entropy=1.4 | '
     'network: unique_dest=3, beacon_score=0.05 | file: sensitive_access=0 | identity: priv_esc=0")',
     sig_colors[0]),
    ('B', 'Cosine Drift',
     'drift(t) = 1 - (V_t . V_baseline) / (||V_t|| * ||V_baseline||)\n'
     'Range: [0, 2]    Typical healthy: [0, 0.05]    Compromised: [0.06, 0.15+]',
     sig_colors[1]),
    ('C', 'CUSUM (Page 1954)',
     'S_t = max(0, S_{t-1} + (drift_t - mu - k))\n'
     'k = 0.5 * sigma  (slack)    h = 4 * sigma  (alarm threshold)\n'
     'ALARM when S_t > h. Catches sustained sub-threshold drift accumulated over windows.',
     sig_colors[2]),
    ('D', 'Drift Direction = E(t) - E(t-1) projected onto threat concepts',
     'drift_vec = V_current - V_baseline\n'
     'proj(c) = cosine(drift_vec, embed(concept_c))    for c in {credential_dumping, lateral_movement, ...}\n'
     'Returns: argmax top-3 concepts with MITRE technique IDs',
     sig_colors[3]),
]

y = Inches(1.5)
for label, title, formula, color in math_boxes:
    rrect(s, Inches(0.5), y, Inches(12.1), Inches(1.2), fill=CARD)
    rect(s, Inches(0.5), y, Inches(0.5), Inches(1.2), fill=color)
    txt(s, Inches(0.52), y+Inches(0.35), Inches(0.46), Inches(0.5), label, 18, WHITE, True, PP_ALIGN.CENTER, fn='Georgia')
    txt(s, Inches(1.15), y+Inches(0.05), Inches(10.9), Inches(0.25), title, 12, NAVY, True)
    txt(s, Inches(1.15), y+Inches(0.35), Inches(10.9), Inches(0.75), formula, 9, BODY, fn='Cascadia Code')
    y += Inches(1.3)

# Bottom bar
rect(s, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.5), fill=NAVY)
txt(s, Inches(0.7), Inches(6.8), Inches(12.0), Inches(0.5),
    'Cosine drift + CUSUM + drift direction: three formulas that together catch Volt + Salt Typhoon where thresholds fail.',
    11.5, GOLD, True, PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 35: CLOSING
# ════════════════════════════════════════════════════════════════════════
s = blank()
s.background.fill.solid()
s.background.fill.fore_color.rgb = NAVY
rect(s, Inches(0), Inches(0), W, Inches(0.06), fill=GOLD)

txt(s, Inches(1.8), Inches(1.8), Inches(10), Inches(0.4),
    '22nd Century Technologies Inc.', 13, GOLD, True, fn='Georgia')
txt(s, Inches(1.8), Inches(2.6), Inches(10), Inches(1.0),
    'From Detection to Preemption\nin 90 Days', 34, WHITE, True, fn='Georgia')
txt(s, Inches(1.8), Inches(4.0), Inches(10), Inches(0.5),
    'Preemptive closes the door.  ACECARD watches the room.', 15, RGBColor(0xA0, 0xB8, 0xD0))

hline(s, Inches(1.8), Inches(4.7), Inches(3.5), GOLD, Pt(2))

tf = txt(s, Inches(1.8), Inches(5.0), Inches(10), Inches(2.0),
    '', 12, WHITE)
add_p(tf, '1.  Pilot 22CT Preemptive — Formal analysis vs. Volt + Salt TTPs', 12, WHITE, sp=Pt(8))
add_p(tf, '2.  Pilot ACECARD UEBA — Replay scenarios on your telemetry', 12, WHITE, sp=Pt(8))
add_p(tf, '3.  Integrate — Preemptive findings drive ACECARD trust floor', 12, WHITE, sp=Pt(8))
add_p(tf, '4.  Deploy to Gabriel Nimbus — K8s / Iron Bank / IL5-IL6-JWICS / cATO', 12, WHITE, sp=Pt(8))

txt(s, Inches(1.8), Inches(6.3), Inches(10), Inches(0.3),
    'Ravindra Shukla, Head of AI  |  22nd Century Technologies Inc.', 11, RGBColor(0xA0, 0xB0, 0xC0))
txt(s, Inches(1.8), Inches(6.7), Inches(10), Inches(0.3),
    'CFIC Collaboration Event  |  Augusta, GA  |  May 6, 2026', 10, RGBColor(0x7A, 0x8B, 0x9E))
brand_bar(s)


# ════════════════════════════════════════════════════════════════════════
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
