"""CFIC Collaboration Event - Final Presentation (20 slides)
Problem Statement Briefs session: 1420-1535, May 6 2026
Gabriel Nimbus / containerized cloud-native prototype
22nd Century Technologies Inc. - CONFIDENTIAL"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

OUT = os.path.join(os.path.dirname(__file__), "CFIC_Final_Presentation.pptx")

DKNAVY  = RGBColor(0x07, 0x14, 0x2A)
NAVY    = RGBColor(0x0B, 0x1F, 0x3A)
BODY    = RGBColor(0x2D, 0x37, 0x48)
SUB     = RGBColor(0x5A, 0x68, 0x78)
ORG     = RGBColor(0xF2, 0x6B, 0x1F)
TEAL    = RGBColor(0x08, 0x91, 0xB2)
CRED    = RGBColor(0xC0, 0x39, 0x2B)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
CARD    = RGBColor(0xF7, 0xF8, 0xFA)
RIGOR_BG = RGBColor(0xFD, 0xE5, 0xD4)
ACE_BG  = RGBColor(0xE0, 0xF2, 0xF7)
GOLD    = RGBColor(0xB8, 0x86, 0x2A)
GREEN   = RGBColor(0x16, 0x6B, 0x3A)
LGREEN  = RGBColor(0xE8, 0xF5, 0xE9)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

def _fill(s, c):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = c

def _rect(s, l, t, w, h, fill=None):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.background()
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def _box(s, l, t, w, h, fill=None, border=None, accent=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.background()
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if border:
        sh.line.color.rgb = border
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    if accent:
        _rect(s, l, t, Inches(0.1), h, fill=accent)
    return sh

def _txt(s, l, t, w, h, text, sz=11, c=BODY, b=False, al=PP_ALIGN.LEFT, fn='Calibri'):
    tf = s.shapes.add_textbox(l, t, w, h).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = c
    p.font.bold = b
    p.font.name = fn
    p.alignment = al
    return tf

def _p(tf, text, sz=11, c=BODY, b=False, sp=Pt(4)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = c
    p.font.bold = b
    p.font.name = 'Calibri'
    p.space_before = sp
    return p

def _cs(part_text, title, subtitle=None, bar_color=ORG):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _fill(s, WHITE)
    _rect(s, Inches(0), Inches(0), W, Inches(0.3), fill=bar_color)
    _txt(s, Inches(0.5), Inches(0.0), Inches(12.3), Inches(0.3),
         part_text, 11, WHITE, True)
    _txt(s, Inches(0.5), Inches(0.5), Inches(12.3), Inches(0.6),
         title, 24, NAVY, True, fn='Georgia')
    if subtitle:
        _txt(s, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.3),
             subtitle, 13, SUB)
    return s

def _bb(s, text, y=7.0):
    _rect(s, Inches(0.5), Inches(y), Inches(12.3), Inches(0.4), fill=NAVY)
    _txt(s, Inches(0.7), Inches(y), Inches(12.0), Inches(0.4),
         text, 10.5, GOLD, True, PP_ALIGN.CENTER)


# ================================================================
# SLIDE 1: TITLE
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_fill(s, DKNAVY)
_txt(s, Inches(1), Inches(1.5), Inches(11.3), Inches(0.5),
     '22nd Century Technologies Inc.', 20, GOLD, True, PP_ALIGN.CENTER, 'Georgia')
_txt(s, Inches(1), Inches(2.2), Inches(11.3), Inches(1.2),
     'CONTAINERIZED BEHAVIORAL CYBER DEFENSE\nFOR GABRIEL NIMBUS', 36, WHITE, True, PP_ALIGN.CENTER, 'Georgia')
_txt(s, Inches(1), Inches(3.8), Inches(11.3), Inches(0.5),
     'Preemptive Posture Verification  +  ACECARD UEBA Behavioral Intelligence', 16, RGBColor(0xA0, 0xB0, 0xC0), False, PP_ALIGN.CENTER)

_box(s, Inches(1.5), Inches(4.8), Inches(4.8), Inches(1.0), fill=None, border=ORG)
_txt(s, Inches(1.7), Inches(4.85), Inches(4.4), Inches(0.25), 'PREEMPTIVE LAYER', 10, ORG, True, PP_ALIGN.CENTER)
_txt(s, Inches(1.7), Inches(5.15), Inches(4.4), Inches(0.5),
     'Mathematically verifies every defensive control\ncloses every known TTP', 11, RGBColor(0xCC, 0xDD, 0xEE), False, PP_ALIGN.CENTER)

_box(s, Inches(7.0), Inches(4.8), Inches(4.8), Inches(1.0), fill=None, border=TEAL)
_txt(s, Inches(7.2), Inches(4.85), Inches(4.4), Inches(0.25), 'BEHAVIORAL LAYER -- ACECARD UEBA', 10, TEAL, True, PP_ALIGN.CENTER)
_txt(s, Inches(7.2), Inches(5.15), Inches(4.4), Inches(0.5),
     'Continuous 1536-d behavioral trajectories\nwith MITRE-mapped drift direction', 11, RGBColor(0xCC, 0xDD, 0xEE), False, PP_ALIGN.CENTER)

_txt(s, Inches(1), Inches(6.2), Inches(11.3), Inches(0.3),
     'Cloud-Native  |  Containerized (K8s/Iron Bank)  |  IL5/IL6/JWICS  |  cATO-Aligned', 12, GOLD, True, PP_ALIGN.CENTER)
_txt(s, Inches(1), Inches(6.8), Inches(11.3), Inches(0.4),
     'CFIC Collaboration Event  |  Augusta, GA  |  May 6, 2026  |  CONFIDENTIAL', 11, RGBColor(0x70, 0x80, 0x90), False, PP_ALIGN.CENTER)


# ================================================================
# SLIDE 2: THE PROBLEM STATEMENT
# ================================================================
s = _cs('PROBLEM STATEMENT', 'Four Limiting Factors in Current Cyber Defense',
        'What Volt Typhoon (5+ yr dwell) and Salt Typhoon (200+ telcos) exploit -- and why current tools miss them.', CRED)

problems = [
    ('01', 'Configuration Gaps Are Never PROVEN Closed', CRED,
     '99% of firewall breaches are misconfiguration (Gartner). Pentests sample. ASM enumerates. Nobody PROVES. '
     '2^8000 config state space cannot be sampled adequately.',
     'Both Typhoons enter via unpatched/misconfigured edge devices (Fortinet, Cisco IOS XE, Ivanti, PAN).'),
    ('02', 'Behavioral Detection Has No Direction', CRED,
     '"User is 87% anomalous" -- anomalous HOW? Threshold-based scoring misses Living-off-the-Land because '
     'every tool used (wmic, ntdsutil, netsh) is a legitimate admin binary.',
     'Volt Typhoon operated BELOW every threshold for 5+ years. No single metric ever crossed a line.'),
    ('03', 'Identity Is Fragmented Across 10+ Systems', CRED,
     'AD, AAD, Okta, AWS IAM, K8s, CrowdStrike, VPN, Splunk, PKI, TACACS+. Each sees one identifier. '
     'Cross-domain lateral movement reads as 4 unrelated weak signals.',
     'One attacker = 10 unrelated identities. No correlation. No composite trajectory.'),
    ('04', 'No Cloud-Native Behavioral Platform for Gabriel Nimbus', CRED,
     'Current UEBA tools are SaaS-only or legacy on-prem. No containerized, K8s-native solution exists that '
     'can deploy into Gabriel Nimbus at IL5/IL6 with cATO alignment.',
     'ARCYBER needs behavioral intelligence INSIDE the classified big data platform, not outside it.'),
]
y = Inches(1.4)
for num, title, color, desc, consequence in problems:
    _box(s, Inches(0.5), y, Inches(12.1), Inches(1.25), fill=CARD, accent=color)
    _rect(s, Inches(0.5), y, Inches(0.7), Inches(1.25), fill=CRED)
    _txt(s, Inches(0.55), y+Inches(0.15), Inches(0.6), Inches(0.4), num, 20, WHITE, True, PP_ALIGN.CENTER, 'Georgia')
    _txt(s, Inches(1.4), y+Inches(0.05), Inches(10.9), Inches(0.3), title, 13, NAVY, True)
    _txt(s, Inches(1.4), y+Inches(0.35), Inches(10.9), Inches(0.4), desc, 9.5, BODY)
    _txt(s, Inches(1.4), y+Inches(0.85), Inches(10.9), Inches(0.35),
         'CONSEQUENCE:  ' + consequence, 9, CRED, True)
    y += Inches(1.32)

_bb(s, 'These four gaps define the attack surface Volt Typhoon and Salt Typhoon exploit. Our solution closes all four.')


# ================================================================
# SLIDE 3: OUR ANSWER -- TWO LAYERS
# ================================================================
s = _cs('OUR SOLUTION', 'The Two-Layer Answer: Preemptive + Behavioral',
        'Preemptive verification shrinks the attack surface mathematically. ACECARD detects what remains behaviorally.', NAVY)

_box(s, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.3), fill=RIGOR_BG)
_rect(s, Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.5), fill=ORG)
_txt(s, Inches(0.7), Inches(1.52), Inches(5.4), Inches(0.45),
     'LAYER 1 -- PREEMPTIVE POSTURE VERIFICATION', 13, WHITE, True)
_txt(s, Inches(0.8), Inches(2.1), Inches(5.3), Inches(0.25), 'CLOSES PROBLEM 1', 10, ORG, True)
tf = _txt(s, Inches(0.8), Inches(2.4), Inches(5.3), Inches(3.5),
     'Formal mathematical model of every NGFW, IPS, IdP, SASE, WAF.', 11, BODY)
_p(tf, 'Reasons exhaustively over 2^8000 state space.', 11, BODY)
_p(tf, '', 6, SUB)
_p(tf, 'THREE INTELLIGENCES:', 10, ORG, True)
_p(tf, 'Attack: MITRE-enriched attack graphs from CVE/CTI/IOC feeds', 10, BODY)
_p(tf, 'Defense: Symbolic Model of Computation -- exhaustive reasoning', 10, BODY)
_p(tf, 'Remediation: Vendor-specific fixes, no new errors', 10, BODY)
_p(tf, '', 6, SUB)
_p(tf, 'MATHEMATICAL GUARANTEE:', 10, NAVY, True)
_p(tf, 'Zero false positives. Zero false negatives.', 11, NAVY, True)
_p(tf, 'Complete reasoning, not sampling.', 11, NAVY, True)
_p(tf, 'Re-verified hourly + on every config change.', 10, BODY)

_box(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.3), fill=ACE_BG)
_rect(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(0.5), fill=TEAL)
_txt(s, Inches(7.0), Inches(1.52), Inches(5.4), Inches(0.45),
     'LAYER 2 -- ACECARD UEBA (BEHAVIORAL)', 13, WHITE, True)
_txt(s, Inches(7.1), Inches(2.1), Inches(5.3), Inches(0.25), 'CLOSES PROBLEMS 2, 3, 4', 10, TEAL, True)
tf = _txt(s, Inches(7.1), Inches(2.4), Inches(5.3), Inches(3.5),
     'Embeds 5 behavioral signals into 1536-d vectors per entity per hour.', 11, BODY)
_p(tf, 'Tracks behavioral TRAJECTORY, not just anomaly score.', 11, BODY)
_p(tf, '', 6, SUB)
_p(tf, 'FOUR INNOVATIONS:', 10, TEAL, True)
_p(tf, '1. CUSUM change-point catches sub-threshold drift (Problem 2)', 10, BODY)
_p(tf, '2. Drift DIRECTION with MITRE mapping -- "drifting toward T1021" (Problem 2)', 10, BODY)
_p(tf, '3. Entity fusion across 10+ identity systems (Problem 3)', 10, BODY)
_p(tf, '4. Containerized K8s-native for Gabriel Nimbus (Problem 4)', 10, BODY)
_p(tf, '', 6, SUB)
_p(tf, 'STRUCTURAL INNOVATION:', 10, NAVY, True)
_p(tf, 'Not what is anomalous -- what they are BECOMING.', 11, NAVY, True)
_p(tf, 'Direction, not just magnitude.', 11, NAVY, True)

_bb(s, 'Together: the attack surface never existed (Preemptive) OR the attack is detected within hours (ACECARD + ABAC)')


# ================================================================
# SLIDE 4: PREEMPTIVE -- HOW IT WORKS
# ================================================================
s = _cs('LAYER 1 -- PREEMPTIVE', 'Close the Door Before They Knock',
        'Three intelligences working together -- from threat intel to verified remediation.', ORG)

cols = [(Inches(0.5), 'Analyze All Known\nThreat Intelligence',
         'Terabytes of CVE, advisories, CTI, IOCs, IOAs, PSIRTs',
         'Attack Intelligence',
         'MITRE-enriched attack graphs via ML/NLP/LLM pipeline. Full TTP coverage of Volt Typhoon (G1017) and Salt Typhoon. Auto-updates on new advisory.'),
        (Inches(4.8), 'Identify Every\nDefensive Control Gap',
         'Exponential combination of NGFW, IdP, IPS, SASE, WAF configs',
         'Defense Intelligence',
         'Symbolic Model of Computation. Formal mathematical model -- exhaustively reasons over 2^8000 state space. No sampling. Every path analyzed.'),
        (Inches(9.1), 'Prescribe Correct\nConfigurations',
         'Without creating new errors or business disruption',
         'Remediation Intelligence',
         'Guardrailed agentic-AI reasoners. Risk-prioritized, vendor-specific fixes. No false positives. No false negatives.')]

for x, prob, prob_desc, sol, sol_desc in cols:
    _box(s, x, Inches(1.5), Inches(4.0), Inches(1.8), fill=CARD, accent=CRED)
    _txt(s, x+Inches(0.3), Inches(1.55), Inches(3.5), Inches(0.2), 'PROBLEM', 9, CRED, True)
    _txt(s, x+Inches(0.3), Inches(1.75), Inches(3.5), Inches(0.5), prob, 12, NAVY, True)
    _txt(s, x+Inches(0.3), Inches(2.35), Inches(3.5), Inches(0.6), prob_desc, 9.5, BODY)

    _box(s, x, Inches(3.5), Inches(4.0), Inches(3.0), fill=CARD, accent=ORG)
    _txt(s, x+Inches(0.3), Inches(3.55), Inches(3.5), Inches(0.2), 'SOLUTION', 9, ORG, True)
    _txt(s, x+Inches(0.3), Inches(3.75), Inches(3.5), Inches(0.3), sol, 12, NAVY, True)
    _txt(s, x+Inches(0.3), Inches(4.15), Inches(3.5), Inches(2.1), sol_desc, 10, BODY)

_bb(s, 'Supports PAN, Fortinet, Cisco, Check Point. SaaS, on-prem, and air-gapped deployments.')


# ================================================================
# SLIDE 5: PREEMPTIVE -- STOPPING BOTH TYPHOONS
# ================================================================
s = _cs('LAYER 1 -- PREEMPTIVE', 'How Preemptive Stops Both Typhoons -- Phase by Phase',
        'Configuration weakness is the entry point for both groups. We mathematically prove those weaknesses do not exist.', ORG)

_rect(s, Inches(0.5), Inches(1.4), Inches(0.7), Inches(0.35), fill=NAVY)
_txt(s, Inches(0.55), Inches(1.42), Inches(0.6), Inches(0.3), 'PHASE', 9, WHITE, True, PP_ALIGN.CENTER)
_rect(s, Inches(1.2), Inches(1.4), Inches(3.5), Inches(0.35), fill=CRED)
_txt(s, Inches(1.3), Inches(1.42), Inches(3.4), Inches(0.3), 'ATTACKER BEHAVIOR', 9, WHITE, True)
_rect(s, Inches(4.7), Inches(1.4), Inches(4.0), Inches(0.35), fill=ORG)
_txt(s, Inches(4.8), Inches(1.42), Inches(3.9), Inches(0.3), 'HOW WE PREEMPT IT', 9, WHITE, True)
_rect(s, Inches(8.7), Inches(1.4), Inches(3.9), Inches(0.35), fill=TEAL)
_txt(s, Inches(8.8), Inches(1.42), Inches(3.8), Inches(0.3), 'IF BYPASSED: ACECARD DETECTS', 9, WHITE, True)

rows = [
    ('Initial\nAccess', 'Edge CVE exploit (Fortinet, Cisco IOS XE, Ivanti, PAN)',
     'IPS signature + WAF verified; compensating controls proven',
     'Auth: anomalous mgmt-plane source IP, drift 0.08'),
    ('Credential\nAccess', 'ntdsutil NTDS.dit dump; LSASS extraction; TACACS+ capture',
     'EDR coverage verified; script-control blocks ntdsutil for non-admin',
     'Process: lolbin_count 0->7, CUSUM alarm @ T+47min'),
    ('Lateral\nMovement', 'SMB/RDP/WinRM to DCs; cross-tenant pivots',
     'SMB inspection shadow detected + remediated; micro-segmentation proven',
     'Network: unique_dest 3->18, cohort z=6.17, HIGH alert'),
    ('Persistence', 'Web shells; local accounts; ACL modification; GRE tunnels',
     'Config Drift catches baseline deviation within 1 hour',
     'Identity: svc_acct spikes, group_adds, drift toward T1078'),
    ('Exfiltration', 'CDR/metadata exfil; C2 to KV botnet / JumbledPath',
     'Egress filtering proven; C2 IPs blocked all zones',
     'File: sensitive_access rises, CRITICAL -> ABAC BLOCKED'),
]
y = Inches(1.85)
for i, (phase, attack, preempt, acecard) in enumerate(rows):
    bg = CARD if i % 2 == 0 else WHITE
    _box(s, Inches(0.5), y, Inches(12.1), Inches(0.95), fill=bg)
    _txt(s, Inches(0.55), y+Pt(2), Inches(0.65), Inches(0.9), phase, 8.5, NAVY, True, PP_ALIGN.CENTER)
    _txt(s, Inches(1.3), y+Pt(2), Inches(3.35), Inches(0.9), attack, 9, BODY)
    _txt(s, Inches(4.8), y+Pt(2), Inches(3.8), Inches(0.9), preempt, 9, RGBColor(0x1A, 0x4D, 0x8C))
    _txt(s, Inches(8.8), y+Pt(2), Inches(3.8), Inches(0.9), acecard, 9, TEAL)
    y += Inches(0.98)

_bb(s, 'Both Typhoons defeated: config gaps never exist (Preemptive) + behavioral drift caught in hours (ACECARD)')


# ================================================================
# SLIDE 6: ACECARD -- 5 SIGNALS + EMBEDDING
# ================================================================
s = _cs('LAYER 2 -- ACECARD UEBA', 'Five Behavioral Signals Into One 1536-d Vector',
        'Each signal individually embedded (1536-d) + one fused composite. 6 vectors = 9,216 dimensions per entity.', TEAL)

sig_colors = [RGBColor(0x1E,0x4D,0x8C), RGBColor(0x6C,0x2E,0xD9), TEAL, ORG, CRED]
sig_names = ['AUTH', 'PROCESS', 'NETWORK', 'FILE', 'IDENTITY']
sig_sources = [
    'Win 4624/4625/4768, Okta, AAD, AWS CloudTrail',
    'Sysmon EID 1/3, CrowdStrike, Defender, auditd',
    'NetFlow/IPFIX, Zeek, PAN-OS, FortiGate, VPC Flow',
    'Sysmon EID 11, Win 4663, DLP, cloud storage',
    'Win 4672/4728, CloudTrail IAM, K8s audit, AD change',
]
sig_features = [
    'logon_count, failure_rate, unique_hosts, off_hours_ratio, impossible_travel, mfa_skip_ratio',
    'unique_processes, cmdline_entropy, lolbin_count, parent_child_depth, unsigned_ratio, encoded_cmd',
    'unique_dest_ips, bytes_out_ratio, beacon_score, dns_query_rate, geo_anomaly, admin_share_access',
    'files_created, files_deleted, sensitive_access, archive_creates, extension_changes, shadow_copy_ops',
    'priv_escalations, group_adds, mfa_bypass, role_changes, service_acct_use, admin_actions',
]

xp = [Inches(0.5), Inches(3.1), Inches(5.7), Inches(8.3), Inches(10.9)]
for i in range(5):
    _rect(s, xp[i], Inches(1.5), Inches(2.4), Inches(0.35), fill=sig_colors[i])
    _txt(s, xp[i], Inches(1.52), Inches(2.4), Inches(0.3), sig_names[i], 11, WHITE, True, PP_ALIGN.CENTER)
    _box(s, xp[i], Inches(1.9), Inches(2.4), Inches(1.4), fill=CARD, border=sig_colors[i])
    _txt(s, xp[i]+Inches(0.1), Inches(1.95), Inches(2.2), Inches(0.35), sig_sources[i], 8, SUB)
    _txt(s, xp[i]+Inches(0.1), Inches(2.35), Inches(2.2), Inches(0.9), sig_features[i], 8, NAVY, True, fn='Cascadia Code')

_rect(s, Inches(0.5), Inches(3.5), Inches(12.1), Inches(0.4), fill=NAVY)
_txt(s, Inches(0.7), Inches(3.52), Inches(11.7), Inches(0.35),
     'All 5 concatenated -> 1 COMPOSITE 1536-d vector (holistic entity behavioral state)  |  Total: 6 x 1,536 = 9,216 dims',
     11, GOLD, True, PP_ALIGN.CENTER)

# Three info boxes
_box(s, Inches(0.5), Inches(4.1), Inches(3.8), Inches(2.5), fill=CARD, accent=TEAL)
_txt(s, Inches(0.8), Inches(4.15), Inches(3.3), Inches(0.25), 'WHY EMBEDDING?', 11, TEAL, True)
tf = _txt(s, Inches(0.8), Inches(4.4), Inches(3.3), Inches(2.0),
     'Captures the SHAPE of behavior, not just metrics.', 10, BODY)
_p(tf, 'Drift in shape = attack even if no metric spikes.', 10, BODY)
_p(tf, 'Direction: "drifting toward T1021" not "87% anomalous"', 10, NAVY, True)
_p(tf, '', 4, SUB)
_p(tf, 'Traditional UEBA: rules + thresholds', 10, CRED)
_p(tf, 'ACECARD: vector geometry + MITRE mapping', 10, TEAL, True)

_box(s, Inches(4.6), Inches(4.1), Inches(4.0), Inches(2.5), fill=CARD, accent=ORG)
_txt(s, Inches(4.9), Inches(4.15), Inches(3.5), Inches(0.25), 'ENTITY FUSION (PROBLEM 3)', 11, ORG, True)
tf = _txt(s, Inches(4.9), Inches(4.4), Inches(3.5), Inches(2.0),
     'One entity_uuid resolves 10+ identity systems:', 10, BODY)
_p(tf, 'AD, AAD, Okta, AWS, K8s, CrowdStrike,', 10, BODY)
_p(tf, 'VPN, Splunk, PKI, TACACS+', 10, BODY)
_p(tf, '', 4, SUB)
_p(tf, 'Without fusion: 4 weak signals, no alert', 10, CRED)
_p(tf, 'With fusion: 1 strong trajectory, HIGH alert', 10, TEAL, True)

_box(s, Inches(8.8), Inches(4.1), Inches(3.8), Inches(2.5), fill=CARD, accent=NAVY)
_txt(s, Inches(9.1), Inches(4.15), Inches(3.3), Inches(0.25), 'DEPLOYMENT OPTIONS', 11, NAVY, True)
tf = _txt(s, Inches(9.1), Inches(4.4), Inches(3.3), Inches(2.0),
     'IL5: OpenAI text-embedding-3-small', 10, BODY)
_p(tf, 'IL6: Local Phi-4 (14B) or Mistral 7B', 10, BODY)
_p(tf, 'JWICS: Same SLM, air-gapped', 10, BODY)
_p(tf, 'Switch: env var only, no code change', 10, NAVY, True)
_p(tf, '', 4, SUB)
_p(tf, 'PRIVACY: No raw event content sent.', 10, GOLD, True)
_p(tf, 'Only summarized metrics. No PII.', 10, GOLD)

_bb(s, 'NOBODY embeds behavioral signals into a unified vector space with per-signal drift decomposition and entity fusion.')


# ================================================================
# SLIDE 7: CUSUM -- THE KEY INNOVATION
# ================================================================
s = _cs('LAYER 2 -- ACECARD UEBA -- CUSUM', 'CUSUM: Catching What Thresholds Miss',
        "Page's Cumulative Sum detects small, sustained behavioral shifts -- the Volt Typhoon signature.", TEAL)

_box(s, Inches(0.5), Inches(1.5), Inches(5.8), Inches(2.5), fill=CARD, accent=TEAL)
_txt(s, Inches(0.8), Inches(1.55), Inches(5.3), Inches(0.25), 'THE ALGORITHM', 11, TEAL, True)
tf = _txt(s, Inches(0.8), Inches(1.85), Inches(5.3), Inches(2.0),
     'S(t) = max(0, S(t-1) + drift(t) - (mu + k))', 12, NAVY, True, fn='Cascadia Code')
_p(tf, '', 4, SUB)
_p(tf, 'S(t) = CUSUM statistic at time t', 10, BODY)
_p(tf, 'drift(t) = cosine distance from baseline', 10, BODY)
_p(tf, 'mu = expected mean drift (30-day rolling)', 10, BODY)
_p(tf, 'k = slack = 0.5 * sigma (allowable noise)', 10, BODY)
_p(tf, 'ALARM when S(t) > h = 4 * sigma', 11, CRED, True)
_p(tf, '', 4, SUB)
_p(tf, 'Resets to 0 when drift returns to normal.', 10, BODY)

_box(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(2.5), fill=CARD, accent=ORG)
_txt(s, Inches(7.1), Inches(1.55), Inches(5.3), Inches(0.25), 'WORKED EXAMPLE: VOLT TYPHOON', 11, ORG, True)
tf = _txt(s, Inches(7.1), Inches(1.85), Inches(5.3), Inches(2.0),
     'Baseline: mu=0.03, sigma=0.03', 10, BODY)
_p(tf, 'k = 0.5 * 0.03 = 0.015', 10, NAVY, True)
_p(tf, 'h = 4 * 0.03 = 0.12', 10, NAVY, True)
_p(tf, '', 4, SUB)
_p(tf, 'APT drift per window = 0.06 (below threshold!)', 10, BODY)
_p(tf, 'Fixed threshold = 0.15 --> NEVER fires', 10, CRED, True)
_p(tf, '', 4, SUB)
_p(tf, 'CUSUM accumulation:', 10, BODY)
_p(tf, '  Excess per window: 0.06 - 0.03 - 0.015 = 0.015', 9, BODY)
_p(tf, '  ALARM in 0.12/0.015 = 8 windows (8 hours)', 10, CRED, True)

_box(s, Inches(0.5), Inches(4.2), Inches(12.1), Inches(2.5), fill=CARD, accent=NAVY)
_txt(s, Inches(0.8), Inches(4.25), Inches(11.6), Inches(0.3),
     'WHY CUSUM CATCHES VOLT TYPHOON BUT THRESHOLDS DO NOT', 13, NAVY, True)

_rect(s, Inches(0.8), Inches(4.7), Inches(5.5), Inches(0.3), fill=CRED)
_txt(s, Inches(0.85), Inches(4.72), Inches(5.4), Inches(0.26), 'FIXED THRESHOLD = 0.15', 10, WHITE, True)
tf = _txt(s, Inches(0.8), Inches(5.05), Inches(5.5), Inches(1.3),
     'Volt Typhoon drift per window = 0.01-0.04', 10, BODY)
_p(tf, 'NEVER crosses 0.15. Ever.', 10, CRED, True)
_p(tf, 'Attacker operates for 5+ YEARS undetected.', 10, CRED, True)
_p(tf, 'Every commercial UEBA uses fixed thresholds.', 10, CRED)

_rect(s, Inches(6.8), Inches(4.7), Inches(5.5), Inches(0.3), fill=TEAL)
_txt(s, Inches(6.85), Inches(4.72), Inches(5.4), Inches(0.26), 'CUSUM (h = 4*sigma = 0.12)', 10, WHITE, True)
tf = _txt(s, Inches(6.8), Inches(5.05), Inches(5.5), Inches(1.3),
     'Same drift 0.03-0.04 ACCUMULATES over windows.', 10, BODY)
_p(tf, 'CUSUM alarm in 4-8 days.', 10, TEAL, True)
_p(tf, 'Dwell time: YEARS --> DAYS.', 10, TEAL, True)
_p(tf, 'Minimax optimal (Lorden 1971, Moustakides 1986).', 10, TEAL)

_bb(s, 'NO COMMERCIAL UEBA uses CUSUM on behavioral embeddings. This is original research applied to cybersecurity.')


# ================================================================
# SLIDE 8: DRIFT DIRECTION -- 8 CONCEPTS
# ================================================================
s = _cs('LAYER 2 -- ACECARD UEBA -- DRIFT DIRECTION', 'Drift Direction: What Is the Entity Becoming?',
        'Cosine projection onto 8 MITRE-mapped threat concepts. Not "87% anomalous" -- "drifting toward T1021 at 0.82."', TEAL)

concepts = [
    ('credential_dumping', 'T1003', RGBColor(0x1E,0x4D,0x8C), 'LSASS access, SAM reads, ntdsutil, DCSync, Kerberoasting'),
    ('lateral_movement', 'T1021', RGBColor(0x6C,0x2E,0xD9), 'RDP/WinRM/SSH to new hosts, admin shares, PtH/PtT'),
    ('data_exfiltration', 'T1041', TEAL, 'Bytes-out inversion, archive staging, DNS tunneling, GRE tunneling'),
    ('c2_beaconing', 'T1071', ORG, 'Regular-interval callbacks, DGA domains, KV botnet, JumbledPath'),
    ('lotl_execution', 'T1059', CRED, 'wmic/ntdsutil/netsh/rundll32 chains, encoded PowerShell'),
    ('privilege_escalation', 'T1078', RGBColor(0x8B,0x5C,0xF6), 'Group adds (Domain Admins), sudo abuse, token manipulation'),
    ('defense_evasion', 'T1562', RGBColor(0xDB,0x27,0x77), 'AV exclusions, log clearing, agent tampering, ACL modification'),
    ('insider_hoarding', 'T1074', RGBColor(0x06,0xB6,0xD4), 'Steady sensitive access increase, bulk downloads, data staging'),
]

y = Inches(1.5)
for name, tid, color, desc in concepts:
    _box(s, Inches(0.5), y, Inches(12.1), Inches(0.55), fill=CARD)
    _rect(s, Inches(0.5), y+Inches(0.05), Inches(2.4), Inches(0.25), fill=color)
    _txt(s, Inches(0.55), y+Inches(0.06), Inches(2.3), Inches(0.22), name, 9.5, WHITE, True, PP_ALIGN.CENTER)
    _txt(s, Inches(3.0), y+Inches(0.06), Inches(0.8), Inches(0.22), tid, 10, NAVY, True, fn='Cascadia Code')
    _txt(s, Inches(3.9), y+Inches(0.02), Inches(8.5), Inches(0.5), desc, 9.5, BODY)
    y += Inches(0.58)

_box(s, Inches(0.5), Inches(6.2), Inches(12.1), Inches(0.6), fill=CARD, accent=GOLD)
_txt(s, Inches(0.8), Inches(6.22), Inches(11.6), Inches(0.55),
     'projection = cosine(drift_vector, concept_embedding). Top 3 returned with MITRE IDs. '
     'Result: "Entity drifting toward lateral_movement (0.82, T1021) + lotl_execution (0.71, T1059)" -- actionable, mappable, auditable.',
     10, NAVY, True)

_bb(s, 'NO OTHER UEBA projects drift direction onto MITRE techniques. This is the core innovation.')


# ================================================================
# SLIDE 9: PEER COHORT + ALERT TIERS
# ================================================================
s = _cs('LAYER 2 -- ACECARD UEBA -- COHORT + ALERTING',
        'Peer Cohort Comparison + Calibrated Alert Tiers',
        'A sysadmin logging into 30 hosts is normal. A finance analyst doing the same is anomalous.', TEAL)

# Left: Peer Cohort
_box(s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(5.5), fill=CARD, accent=TEAL)
_rect(s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(0.35), fill=NAVY)
_txt(s, Inches(0.7), Inches(1.42), Inches(5.6), Inches(0.3), 'PEER COHORT', 11, WHITE, True)

_txt(s, Inches(0.8), Inches(1.9), Inches(5.5), Inches(0.35),
     'z = (entity_drift - cohort_mean) / cohort_std', 12, TEAL, True, fn='Cascadia Code')
_txt(s, Inches(0.8), Inches(2.3), Inches(5.5), Inches(0.25),
     'Flag outlier if z > 2.0  |  Cohort = role / OU / security group from AD/LDAP/IAM', 9, SUB)

_txt(s, Inches(0.8), Inches(2.7), Inches(5.5), Inches(0.3),
     'EXAMPLE -- Lateral Movement Detection', 11, ORG, True)
_txt(s, Inches(0.8), Inches(3.0), Inches(5.5), Inches(0.25),
     'Entity: jsmith (financial_analyst, cohort = 42 peers)', 9.5, BODY)

metrics = [('jsmith drift:', '0.142', '(38 failed logons to 12 new hosts)'),
           ('Cohort mean:', '0.031', '(peers avg 2.1 failed, 0.3 new hosts)'),
           ('Cohort std:', '0.018', ''),
           ('z-score:', '6.17', '>>> OUTLIER (z > 2.0)')]
y = Inches(3.4)
for label, val, note in metrics:
    _txt(s, Inches(0.8), y, Inches(1.8), Inches(0.25), label, 10, BODY)
    zc = CRED if val == '6.17' else ORG
    _txt(s, Inches(2.8), y, Inches(1.0), Inches(0.25), val, 13, zc, True, fn='Cascadia Code')
    if note:
        _txt(s, Inches(3.9), y, Inches(2.5), Inches(0.25), note, 8.5, SUB)
    y += Inches(0.3)

_box(s, Inches(0.8), Inches(4.7), Inches(5.5), Inches(1.8), fill=ACE_BG)
tf = _txt(s, Inches(1.0), Inches(4.75), Inches(5.1), Inches(1.6),
     'Drift direction:  lateral_movement (0.78), credential_dumping (0.45)', 9.5, TEAL, True)
_p(tf, 'MITRE:  T1021 (Remote Services), T1003 (Credential Dumping)', 9.5, NAVY, True)
_p(tf, 'Action:  HIGH severity -> SOAR playbook triggered', 9.5, ORG, True)

# Right: Alert Tiers
_box(s, Inches(6.8), Inches(1.4), Inches(5.8), Inches(5.5), fill=CARD, accent=NAVY)
_rect(s, Inches(6.8), Inches(1.4), Inches(5.8), Inches(0.35), fill=NAVY)
_txt(s, Inches(7.0), Inches(1.42), Inches(5.4), Inches(0.3), 'ALERT TIERS', 11, WHITE, True)

tiers = [
    ('CRITICAL', CRED, 'Health < 40', 'Immediate SOC investigation',
     'Multiple signals anomalous, CUSUM triggered, drift > 0.15'),
    ('HIGH', ORG, 'Health < 40 OR Velocity z > 3', 'Tier-2 triage in 15 min',
     'Rapid behavioral divergence or composite health collapse'),
    ('MEDIUM', GOLD, 'Health < 70 + change-point', 'Tier-1 review in 1 hour',
     'Recent structural shift detected but not yet severe'),
    ('LOW', TEAL, 'Health 70-85 AND cohort z > 2', 'Watchlist; auto-enrich',
     'Peer outlier but within normal drift range'),
    ('INFO', SUB, 'Any CUSUM trigger (4-sigma)', 'Log for trend analysis',
     'Structural change; may be benign (role change, new project)'),
]
y = Inches(1.9)
for label, color, trigger, action, desc in tiers:
    _box(s, Inches(7.0), y, Inches(5.4), Inches(0.95), fill=WHITE)
    _rect(s, Inches(7.0), y+Inches(0.05), Inches(1.1), Inches(0.85), fill=color)
    _txt(s, Inches(7.05), y+Inches(0.15), Inches(1.0), Inches(0.4), label, 11, WHITE, True, PP_ALIGN.CENTER)
    _txt(s, Inches(8.2), y+Inches(0.05), Inches(4.1), Inches(0.25), trigger, 9, NAVY, True, fn='Cascadia Code')
    _txt(s, Inches(8.2), y+Inches(0.3), Inches(4.1), Inches(0.25), action, 9, color, True)
    _txt(s, Inches(8.2), y+Inches(0.55), Inches(4.1), Inches(0.35), desc, 8.5, SUB)
    y += Inches(0.98)

_bb(s, 'Analyst TP/FP feedback per alert feeds back as labeled data -> threshold auto-tuning + ABAC trust tags')


# ================================================================
# SLIDE 10: ABAC TRUST LOOP
# ================================================================
s = _cs('LAYER 2 -- ACECARD UEBA -- ABAC', 'ABAC Trust State Machine: Detection to Enforcement in <5 Minutes',
        'Not binary lockout. Graduated, reversible, proportional response with analyst feedback loop.', TEAL)

# The Loop
_rect(s, Inches(0.5), Inches(1.4), Inches(12.1), Inches(0.35), fill=NAVY)
_txt(s, Inches(0.7), Inches(1.42), Inches(11.7), Inches(0.3),
     'THE LOOP -- closes in <5 minutes from TP/FP click to ABAC tag update', 11, WHITE, True, PP_ALIGN.CENTER)

loop_steps = [
    ('1', 'ALERT', 'CUSUM/health\ntriggers alert'),
    ('2', 'TRIAGE', 'Analyst reviews\n7-chart dashboard'),
    ('3', 'VERDICT', 'Analyst marks\nTP or FP'),
    ('4', 'LEARN', 'Labeled data feeds\nthreshold tuning'),
    ('5', 'TAG', 'Entity trust state\nupdated in ABAC'),
    ('6', 'ENFORCE', 'ABAC adjusts:\nMFA, restrict, block'),
]
x = Inches(0.5)
for num, name, desc in loop_steps:
    _box(s, x, Inches(1.9), Inches(1.9), Inches(1.2), fill=CARD, border=TEAL)
    _rect(s, x+Inches(0.05), Inches(1.95), Inches(0.35), Inches(0.35), fill=TEAL)
    _txt(s, x+Inches(0.08), Inches(1.97), Inches(0.3), Inches(0.3), num, 14, WHITE, True, PP_ALIGN.CENTER)
    _txt(s, x+Inches(0.45), Inches(1.97), Inches(1.3), Inches(0.3), name, 10, NAVY, True)
    _txt(s, x+Inches(0.45), Inches(2.35), Inches(1.3), Inches(0.6), desc, 8.5, BODY)
    x += Inches(2.0)

# ABAC Trust States
_txt(s, Inches(0.5), Inches(3.3), Inches(12.1), Inches(0.3),
     'ABAC TRUST STATES -- proportional response, not binary', 12, NAVY, True)

states = [
    ('TRUSTED', TEAL, 'Health >= 70, no recent alerts', 'Normal access, standard MFA'),
    ('ELEVATED_WATCH', GOLD, 'Health 40-70 OR cohort z > 2', 'Step-up MFA on sensitive resources'),
    ('RESTRICTED', ORG, 'Health < 40 OR TP alert confirmed', 'Read-only access, no lateral movement'),
    ('BLOCKED', CRED, 'Active incident confirmed', 'Account disabled, session terminated'),
]
y = Inches(3.7)
for label, color, condition, enforcement in states:
    _box(s, Inches(0.5), y, Inches(12.1), Inches(0.7), fill=CARD)
    _rect(s, Inches(0.5), y+Inches(0.1), Inches(2.3), Inches(0.3), fill=color)
    _txt(s, Inches(0.55), y+Inches(0.12), Inches(2.2), Inches(0.26), label, 11, WHITE, True, PP_ALIGN.CENTER)
    _txt(s, Inches(3.0), y+Inches(0.05), Inches(1.2), Inches(0.2), 'CONDITION:', 8, color, True)
    _txt(s, Inches(3.0), y+Inches(0.25), Inches(4.5), Inches(0.4), condition, 10, BODY)
    _txt(s, Inches(7.5), y+Inches(0.05), Inches(1.5), Inches(0.2), 'ENFORCEMENT:', 8, color, True)
    _txt(s, Inches(7.5), y+Inches(0.25), Inches(4.8), Inches(0.4), enforcement, 10, BODY)
    y += Inches(0.73)

_bb(s, 'Self-reinforcing: TP label strengthens signals. FP raises thresholds for that role. ABAC enforcement is automatic.')


# ================================================================
# SLIDE 11: GABRIEL NIMBUS DEPLOYMENT ARCHITECTURE
# ================================================================
s = _cs('DEPLOYMENT -- GABRIEL NIMBUS', 'Containerized Cloud-Native Architecture for Gabriel Nimbus',
        'Kubernetes-native. Iron Bank hardened images. IL5/IL6/JWICS. cATO-aligned.', NAVY)

# 5 layers
layers = [
    ('LAYER 5: PRESENTATION', RGBColor(0x6C,0x2E,0xD9),
     'Plotly Dash / React UI  |  7-chart entity dashboard  |  Alert feed  |  Kill-chain visualization  |  ABAC admin',
     'CAC/PIV auth  |  Role-based access'),
    ('LAYER 4: API GATEWAY', TEAL,
     'FastAPI (async)  |  REST + WebSocket  |  RBAC middleware  |  Rate limiting  |  OpenAPI docs',
     'Istio service mesh  |  mTLS between services'),
    ('LAYER 3: BEHAVIORAL ENGINE', ORG,
     'Signal Summarizers (5)  |  Embedding Pipeline  |  Trajectory Analysis  |  CUSUM  |  Drift Direction  |  Entity Fusion  |  ABAC Engine',
     'Horizontal pod autoscaling  |  10K entities/hr per 4-vCPU pod'),
    ('LAYER 2: DATA PLANE', NAVY,
     'PostgreSQL + pgvector (HNSW)  |  TimescaleDB (time-series)  |  Redis (cache + pub/sub)  |  S3-compatible object store',
     'Persistent volumes  |  Automated backups  |  Encryption at rest'),
    ('LAYER 1: INFRASTRUCTURE', GREEN,
     'Kubernetes (Big Bang / RKE2)  |  Iron Bank hardened containers  |  Istio  |  Fluentd  |  Prometheus + Grafana',
     'cATO artifacts auto-generated  |  STIG-compliant  |  FIPS 140-2'),
]

y = Inches(1.4)
for name, color, components, infra in layers:
    _box(s, Inches(0.5), y, Inches(12.1), Inches(1.0), fill=CARD, accent=color)
    _rect(s, Inches(0.5), y, Inches(2.8), Inches(0.35), fill=color)
    _txt(s, Inches(0.55), y+Inches(0.02), Inches(2.7), Inches(0.3), name, 9, WHITE, True, PP_ALIGN.CENTER)
    _txt(s, Inches(3.5), y+Inches(0.05), Inches(8.9), Inches(0.35), components, 9, BODY)
    _txt(s, Inches(3.5), y+Inches(0.45), Inches(8.9), Inches(0.45), infra, 8.5, SUB)
    y += Inches(1.05)

_bb(s, 'Entire stack deploys via helm chart. Air-gapped: same containers, local SLM (Phi-4/Mistral), no external calls.')


# ================================================================
# SLIDE 12: DATA FLOW IN GABRIEL NIMBUS
# ================================================================
s = _cs('DEPLOYMENT -- DATA FLOW', 'Data Flow: From Raw Telemetry to Actionable Alert',
        'End-to-end pipeline inside Gabriel Nimbus. No data leaves the enclave.', NAVY)

steps = [
    ('1', 'INGEST', GREEN,
     'ECS / OCSF / STIX 2.1 normalized events arrive from SIEM, EDR, NetFlow, IdP, K8s audit logs.',
     'Kafka / Fluentd input'),
    ('2', 'WINDOW', TEAL,
     'Aggregate into 1-hour time windows per entity_uuid. Entity fusion resolves 10+ identity systems.',
     'TimescaleDB hypertable'),
    ('3', 'SUMMARIZE', RGBColor(0x6C,0x2E,0xD9),
     '5 signal summarizers produce structured text (metrics only). No raw content. No PII.',
     'Python workers (horizontal autoscale)'),
    ('4', 'EMBED', ORG,
     'Concatenated signal text embedded into 1536-d vector. IL5: OpenAI API. IL6/JWICS: local Phi-4.',
     'pgvector HNSW index'),
    ('5', 'ANALYZE', NAVY,
     'Trajectory analysis: cosine drift, velocity, acceleration, CUSUM 4-sigma, drift direction onto 8 concepts.',
     'Per-entity, per-window computation'),
    ('6', 'ALERT', CRED,
     'Alert tiers (CRITICAL through INFO) with MITRE technique IDs, cohort z-score, ABAC trust state update.',
     'WebSocket push to dashboard + SOAR webhook'),
]

y = Inches(1.4)
for num, name, color, desc, infra in steps:
    _box(s, Inches(0.5), y, Inches(12.1), Inches(0.85), fill=CARD)
    _rect(s, Inches(0.5), y+Inches(0.05), Inches(0.6), Inches(0.75), fill=color)
    _txt(s, Inches(0.55), y+Inches(0.1), Inches(0.5), Inches(0.3), num, 18, WHITE, True, PP_ALIGN.CENTER, 'Georgia')
    _txt(s, Inches(1.2), y+Inches(0.05), Inches(1.2), Inches(0.3), name, 12, color, True)
    _txt(s, Inches(2.5), y+Inches(0.05), Inches(7.5), Inches(0.4), desc, 9.5, BODY)
    _txt(s, Inches(2.5), y+Inches(0.5), Inches(7.5), Inches(0.3), infra, 8.5, SUB)
    _txt(s, Inches(10.2), y+Inches(0.15), Inches(2.2), Inches(0.5),
         '<3 sec per entity' if name != 'EMBED' else 'env var switches model', 9, color, True, PP_ALIGN.CENTER)
    y += Inches(0.88)

_bb(s, 'Raw event to actionable alert: <3 seconds per entity. 10K entities/hour per 4-vCPU pod. Horizontal autoscale.')


# ================================================================
# SLIDE 13: INTEGRATION ARCHITECTURE
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_fill(s, WHITE)
_rect(s, Inches(0), Inches(0), Inches(6.67), Inches(0.3), fill=ORG)
_rect(s, Inches(6.67), Inches(0), Inches(6.67), Inches(0.3), fill=TEAL)
_txt(s, Inches(0.5), Inches(0.0), Inches(12.3), Inches(0.3),
     'INTEGRATION -- COMBINED ARCHITECTURE', 11, WHITE, True)
_txt(s, Inches(0.5), Inches(0.5), Inches(12.3), Inches(0.5),
     'End-to-End Architecture: Preemptive + ACECARD UEBA', 24, NAVY, True, fn='Georgia')
_txt(s, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.3),
     'Preemptive proves what cannot happen. ACECARD detects what is happening. Shared MITRE vocabulary.', 13, SUB)

# Inputs column
_box(s, Inches(0.5), Inches(1.5), Inches(2.5), Inches(5.3), fill=CARD, accent=NAVY)
_rect(s, Inches(0.5), Inches(1.5), Inches(2.5), Inches(0.35), fill=NAVY)
_txt(s, Inches(0.6), Inches(1.52), Inches(2.3), Inches(0.3), 'INPUTS', 11, WHITE, True, PP_ALIGN.CENTER)
inputs = [('Threat Intel', 'MITRE ATT&CK, CISA, FBI, NSA, CTI feeds'),
          ('Defense Config', 'NGFW, IPS, IdP, IAM, SASE, WAF rules'),
          ('Telemetry', 'Win EventLog, Sysmon, EDR, NetFlow, Zeek, Okta, K8s'),
          ('Schema', 'ECS, OCSF, STIX 2.1 normalized')]
y = Inches(2.0)
for name, desc in inputs:
    _txt(s, Inches(0.7), y, Inches(2.1), Inches(0.2), name, 10, ORG, True)
    _txt(s, Inches(0.7), y+Inches(0.2), Inches(2.1), Inches(0.35), desc, 8, SUB)
    y += Inches(0.6)

# Layer 1
_box(s, Inches(3.3), Inches(1.5), Inches(4.5), Inches(5.3), fill=RIGOR_BG, accent=ORG)
_rect(s, Inches(3.3), Inches(1.5), Inches(4.5), Inches(0.35), fill=ORG)
_txt(s, Inches(3.4), Inches(1.52), Inches(4.3), Inches(0.3), 'LAYER 1 -- PREEMPTIVE', 11, WHITE, True, PP_ALIGN.CENTER)
l1_items = ['Attack Intelligence\nMITRE-enriched attack graphs',
            'Defense Intelligence\nSymbolic Model of Computation',
            'Remediation Intelligence\nRisk-prioritized config fixes']
y = Inches(2.0)
for item in l1_items:
    _box(s, Inches(3.5), y, Inches(4.1), Inches(0.8), fill=WHITE)
    _txt(s, Inches(3.7), y+Inches(0.05), Inches(3.7), Inches(0.7), item, 9.5, BODY)
    y += Inches(0.85)

_box(s, Inches(3.5), Inches(4.6), Inches(4.1), Inches(1.0), fill=WHITE, border=ORG)
_txt(s, Inches(3.7), Inches(4.62), Inches(3.7), Inches(0.2), 'OUTPUT', 9, ORG, True)
tf = _txt(s, Inches(3.7), Inches(4.85), Inches(3.7), Inches(0.7),
     'Hardened defense posture', 9, BODY)
_p(tf, 'Verified SOC playbooks', 9, BODY)
_p(tf, 'MITRE TTP coverage proof', 9, BODY)

# Layer 2
_box(s, Inches(8.1), Inches(1.5), Inches(4.5), Inches(5.3), fill=ACE_BG, accent=TEAL)
_rect(s, Inches(8.1), Inches(1.5), Inches(4.5), Inches(0.35), fill=TEAL)
_txt(s, Inches(8.2), Inches(1.52), Inches(4.3), Inches(0.3), 'LAYER 2 -- ACECARD UEBA', 11, WHITE, True, PP_ALIGN.CENTER)
l2_items = ['5 Signal Summarizers -> 1536-d\nAuth, Process, Network, File, Identity',
            'Trajectory Analysis\nCosine drift + velocity + CUSUM 4-sigma',
            'Drift Direction -> 8 MITRE Concepts\nTop-3 projections -> ATT&CK technique IDs']
y = Inches(2.0)
for item in l2_items:
    _box(s, Inches(8.3), y, Inches(4.1), Inches(0.8), fill=WHITE)
    _txt(s, Inches(8.5), y+Inches(0.05), Inches(3.7), Inches(0.7), item, 9.5, BODY)
    y += Inches(0.85)

_box(s, Inches(8.3), Inches(4.6), Inches(4.1), Inches(1.0), fill=WHITE, border=TEAL)
_txt(s, Inches(8.5), Inches(4.62), Inches(3.7), Inches(0.2), 'OUTPUT', 9, TEAL, True)
tf = _txt(s, Inches(8.5), Inches(4.85), Inches(3.7), Inches(0.7),
     '5-tier alerts (CRITICAL -> INFO)', 9, BODY)
_p(tf, 'ABAC trust state per entity', 9, BODY)
_p(tf, '7-chart dashboard + SOAR hooks', 9, BODY)

# Shared vocab
_rect(s, Inches(3.3), Inches(5.8), Inches(9.3), Inches(0.35), fill=NAVY)
_txt(s, Inches(3.5), Inches(5.82), Inches(8.9), Inches(0.3),
     'Shared vocabulary: MITRE ATT&CK technique IDs  (T1003, T1021, T1041, T1059, T1071, T1078, T1190, T1556 ...)',
     9, GOLD, True, PP_ALIGN.CENTER)

_txt(s, Inches(0.5), Inches(6.3), Inches(12.1), Inches(0.2),
     'Iron Bank  |  Big Bang  |  IL5/IL6/JWICS  |  CAC/PIV  |  cATO-aligned', 9, SUB, False, PP_ALIGN.CENTER)

_bb(s, 'Both layers share MITRE ATT&CK as their analytical vocabulary. One alert = what, which TTP, which entity, is config covered.')


# ================================================================
# SLIDE 14: COVERAGE MATRIX
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_fill(s, WHITE)
_rect(s, Inches(0), Inches(0), Inches(6.67), Inches(0.3), fill=ORG)
_rect(s, Inches(6.67), Inches(0), Inches(6.67), Inches(0.3), fill=TEAL)
_txt(s, Inches(0.5), Inches(0.0), Inches(12.3), Inches(0.3),
     'COVERAGE MATRIX -- WHICH LAYER OWNS WHAT', 11, WHITE, True)
_txt(s, Inches(0.5), Inches(0.5), Inches(12.3), Inches(0.5),
     'Combined Coverage: Preemptive First, Behavioral as Backstop', 24, NAVY, True, fn='Georgia')

_rect(s, Inches(0.5), Inches(1.2), Inches(12.1), Inches(0.35), fill=NAVY)
_txt(s, Inches(0.6), Inches(1.22), Inches(3.8), Inches(0.3), 'Adversary Capability', 9, WHITE, True)
_txt(s, Inches(4.5), Inches(1.22), Inches(0.5), Inches(0.3), 'VT', 9, WHITE, True, PP_ALIGN.CENTER)
_txt(s, Inches(5.1), Inches(1.22), Inches(0.5), Inches(0.3), 'ST', 9, WHITE, True, PP_ALIGN.CENTER)
_txt(s, Inches(5.7), Inches(1.22), Inches(1.3), Inches(0.3), 'Owner', 9, WHITE, True, PP_ALIGN.CENTER)
_txt(s, Inches(7.1), Inches(1.22), Inches(5.5), Inches(0.3), 'Why This Layer', 9, WHITE, True)

mx = [
    ('Edge-device CVE exploitation', 'Y', 'Y', 'PREEMPT', 'Formal proof patch enforced or compensating control', ORG),
    ('Misconfigured / shadowed FW rules', 'Y', 'Y', 'PREEMPT', 'Math model finds shadows across multi-vendor', ORG),
    ('Router ACL tampering / drift', '', 'Y', 'PREEMPT', 'Config drift re-verifies hourly vs baseline', ORG),
    ('C2 to known malicious infra', 'Y', 'Y', 'BOTH', 'Preemptive blocks known; ACECARD catches unknown', RGBColor(0x6C,0x2E,0xD9)),
    ('Compromised valid accounts', 'Y', 'Y', 'ACECARD', 'Policy-allowed traffic; only drift reveals abuse', TEAL),
    ('Living-off-the-Land (T1059)', 'Y', 'Y', 'ACECARD', 'lotl_execution concept + process signal shift', TEAL),
    ('Cross-domain lateral movement', 'Y', 'Y', 'ACECARD', 'Entity fusion joins all identities -> one drift', TEAL),
    ('Slow APT (months/years)', 'Y', 'Y', 'ACECARD', 'CUSUM accumulates sub-threshold drift', TEAL),
    ('Zero-day / unknown TTPs', 'Y', 'Y', 'ACECARD', 'Behavioral anomaly even with no MITRE label', TEAL),
    ('Insider threat / data hoarding', '', '', 'ACECARD', 'insider_hoarding (T1074) + cohort z-score', TEAL),
]

y = Inches(1.6)
for i, (cap, vt, st, layer, why, color) in enumerate(mx):
    bg = CARD if i % 2 == 0 else WHITE
    _box(s, Inches(0.5), y, Inches(12.1), Inches(0.45), fill=bg)
    _txt(s, Inches(0.6), y+Pt(2), Inches(3.8), Inches(0.4), cap, 9, NAVY, True)
    _txt(s, Inches(4.5), y+Pt(2), Inches(0.5), Inches(0.4), vt, 9, CRED, al=PP_ALIGN.CENTER)
    _txt(s, Inches(5.1), y+Pt(2), Inches(0.5), Inches(0.4), st, 9, CRED, al=PP_ALIGN.CENTER)
    _rect(s, Inches(5.8), y+Inches(0.07), Inches(1.1), Inches(0.24), fill=color)
    _txt(s, Inches(5.8), y+Inches(0.08), Inches(1.1), Inches(0.22), layer, 8, WHITE, True, PP_ALIGN.CENTER)
    _txt(s, Inches(7.1), y+Pt(2), Inches(5.4), Inches(0.4), why, 9, BODY)
    y += Inches(0.47)

_bb(s, 'Together: attack surface never existed (Preemptive) OR attack detected and contained within hours (ACECARD + ABAC)')


# ================================================================
# SLIDE 15: INNOVATION PROOF
# ================================================================
s = _cs('INNOVATION -- NOBODY ELSE HAS THIS', 'What Exists vs. What We Built',
        'Systematic comparison of our innovations against every commercial UEBA and preemptive platform.', NAVY)

innovations = [
    ('Formal verification of 2^8000 config state', 'Pentests sample. ASM enumerates. Nobody PROVES.', '22CT PREEMPT', ORG),
    ('Zero false positives / zero false negatives', 'Mathematical guarantee, not probabilistic.', '22CT PREEMPT', ORG),
    ('Behavioral embedding in unified 1536-d space', 'Industry uses features/rules. We use vectors.', 'ACECARD', TEAL),
    ('Per-signal drift decomposition (5 individual)', 'Others embed once. We embed 6x for attribution.', 'ACECARD', TEAL),
    ('CUSUM on embedding trajectories', 'Nobody applies sequential change-point to vectors.', 'ACECARD', TEAL),
    ('Drift direction with MITRE technique mapping', '"Drifting toward T1021" -- no other UEBA does this.', 'ACECARD', TEAL),
    ('Entity fusion across 10+ identity systems', 'Commercial UEBA: 1-3 sources max. We fuse 10+.', 'ACECARD', TEAL),
    ('HMAC-chained forensic evidence', 'Tamper-evident behavioral chains. Novel.', 'ACECARD', TEAL),
    ('ABAC trust loop with behavioral triggers', 'Detection + enforcement in one system. <5 min.', 'ACECARD', TEAL),
    ('Containerized K8s-native for Gabriel Nimbus', 'No existing UEBA ships as Iron Bank container.', 'COMBINED', NAVY),
]

y = Inches(1.4)
for innovation, why_unique, who, color in innovations:
    _box(s, Inches(0.5), y, Inches(12.1), Inches(0.45), fill=CARD)
    _rect(s, Inches(0.5), y, Inches(0.08), Inches(0.45), fill=color)
    _txt(s, Inches(0.7), y+Pt(2), Inches(5.5), Inches(0.4), innovation, 10, NAVY, True)
    _txt(s, Inches(6.3), y+Pt(2), Inches(4.0), Inches(0.4), why_unique, 9, BODY)
    _rect(s, Inches(10.5), y+Inches(0.08), Inches(2.0), Inches(0.25), fill=color)
    _txt(s, Inches(10.5), y+Inches(0.09), Inches(2.0), Inches(0.23), who, 8.5, WHITE, True, PP_ALIGN.CENTER)
    y += Inches(0.48)

_bb(s, 'Combined: 10 innovations that do not exist anywhere else in the industry. This is not incremental improvement.')


# ================================================================
# SLIDE 16: KEY NUMBERS
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_fill(s, WHITE)
_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.3), fill=NAVY)
_txt(s, Inches(0.5), Inches(0.0), Inches(12.3), Inches(0.3),
     'KEY NUMBERS TO MEMORIZE', 11, WHITE, True)
_txt(s, Inches(0.5), Inches(0.5), Inches(12.3), Inches(0.5),
     'Key Numbers', 24, NAVY, True, fn='Georgia')

left = [
    ('5+ years', 'Volt Typhoon dwell time'),
    ('4+ years', 'Salt Typhoon dwell time'),
    ('200+', 'Telcos compromised by Salt Typhoon'),
    ('99%', 'FW breaches = misconfig (Gartner)'),
    ('2^8000', 'Config state space (cannot be sampled)'),
    ('50M+', 'Rule combinations reasoned over'),
    ('0', 'False positives (mathematical proof)'),
    ('0', 'False negatives (exhaustive reasoning)'),
    ('<1 hour', 'Config drift detection time'),
]
right = [
    ('1,536', 'Embedding dims per signal'),
    ('9,216', 'Total dims (6 x 1536) per entity'),
    ('5', 'Behavioral signals'),
    ('8', 'MITRE threat concepts'),
    ('4-sigma', 'CUSUM alarm threshold'),
    ('<3 sec', 'Event to analysis per entity'),
    ('10K/hr', 'Entities per 4-vCPU worker'),
    ('90 days', 'Contract to first detection'),
    ('IL5/IL6/JWICS', 'Three enclave support'),
]

for nums, xs in [(left, Inches(0.5)), (right, Inches(6.8))]:
    y = Inches(1.2)
    for num, meaning in nums:
        _box(s, xs, y, Inches(6.0), Inches(0.45), fill=CARD)
        _rect(s, xs+Inches(0.03), y+Inches(0.05), Inches(1.3), Inches(0.3), fill=NAVY)
        _txt(s, xs+Inches(0.08), y+Inches(0.07), Inches(1.2), Inches(0.26),
             num, 10.5, WHITE, True, PP_ALIGN.CENTER, fn='Cascadia Code')
        _txt(s, xs+Inches(1.5), y+Inches(0.07), Inches(4.3), Inches(0.26), meaning, 10, BODY)
        y += Inches(0.5)


# ================================================================
# SLIDE 17: 22CT QUALIFICATIONS
# ================================================================
s = _cs('22ND CENTURY TECHNOLOGIES', 'Why 22CT: Proven at Scale in Federal Cyber Operations',
        'Operating the largest Army SOC/MDR contract with 800+ cleared analysts.', NAVY)

quals = [
    ('$90M Army SOC/MDR Contract', ORG,
     'Active Army SOC and Managed Detection & Response. 800+ cleared analysts across multiple installations. '
     '24x7x365 operations. Proven incident response at enterprise scale.'),
    ('Cleared Workforce at Scale', TEAL,
     '800+ cleared analysts (Secret through TS/SCI). Trained on MITRE ATT&CK, NIST 800-53, DoD RMF. '
     'Ready to operate ACECARD from Day 1 -- no hiring ramp.'),
    ('Gabriel Nimbus Experience', NAVY,
     'Team members with direct experience operating within classified big data platforms. '
     'Understand the deployment constraints, accreditation requirements, and operational workflows.'),
    ('IL5/IL6/JWICS Deployment', GREEN,
     'Existing infrastructure and ATO artifacts for IL5 and IL6 deployments. '
     'Air-gapped operation validated. cATO artifacts auto-generated by platform.'),
    ('ACECARD: Built and Tested', RGBColor(0x6C,0x2E,0xD9),
     'ACECARD UEBA is not a concept -- it is a working prototype with synthetic Volt/Salt Typhoon replay scenarios, '
     'full trajectory analytics, drift direction, entity fusion, and ABAC enforcement.'),
]

y = Inches(1.4)
for title, color, desc in quals:
    _box(s, Inches(0.5), y, Inches(12.1), Inches(0.95), fill=CARD, accent=color)
    _txt(s, Inches(0.8), y+Inches(0.05), Inches(11.6), Inches(0.3), title, 13, color, True)
    _txt(s, Inches(0.8), y+Inches(0.35), Inches(11.6), Inches(0.55), desc, 10, BODY)
    y += Inches(1.05)

_bb(s, '22CT brings the analysts, the clearances, the infrastructure, and the technology. Ready to deploy.')


# ================================================================
# SLIDE 18: TIMELINE
# ================================================================
s = _cs('DEPLOYMENT TIMELINE', '90 Days: Contract to First Detection',
        'Phased deployment into Gabriel Nimbus with incremental capability delivery.', NAVY)

phases = [
    ('PHASE 1', 'Weeks 1-4', 'FOUNDATION', GREEN,
     ['Deploy K8s containers into Gabriel Nimbus (Iron Bank images)',
      'Connect SIEM/EDR/NetFlow telemetry via ECS/OCSF connectors',
      'Configure entity fusion rules (AD/AAD/Okta mapping)',
      'Baseline 30-day rolling centroid per entity',
      'DELIVERABLE: Data flowing, entities resolved, baselines computing']),
    ('PHASE 2', 'Weeks 5-8', 'DETECTION', TEAL,
     ['Enable 5-signal embedding pipeline',
      'Activate trajectory analysis (drift, velocity, CUSUM)',
      'Deploy 8 threat concept vectors + drift direction',
      'Configure cohort definitions from AD/LDAP',
      'DELIVERABLE: Alerts firing, drift direction operational']),
    ('PHASE 3', 'Weeks 9-12', 'OPERATIONS', ORG,
     ['Integrate with SOC workflows (SOAR webhooks, ServiceNow)',
      'Tune alert thresholds with analyst TP/FP feedback',
      'Enable ABAC trust state machine',
      'Replay Volt + Salt Typhoon scenarios for validation',
      'DELIVERABLE: Full operational capability, ABAC enforcing']),
    ('PHASE 4', 'Week 13+', 'PREEMPTIVE', NAVY,
     ['Integrate preemptive config verification layer',
      'Connect NGFW/IPS/SASE configs for formal analysis',
      'Enable trust-floor coupling (preemptive gaps -> ABAC)',
      'Continuous re-verification on config/threat-feed changes',
      'DELIVERABLE: Both layers operational, closed-loop defense']),
]

y = Inches(1.4)
for phase, weeks, name, color, items in phases:
    _box(s, Inches(0.5), y, Inches(12.1), Inches(1.3), fill=CARD, accent=color)
    _rect(s, Inches(0.5), y, Inches(1.5), Inches(0.35), fill=color)
    _txt(s, Inches(0.55), y+Inches(0.02), Inches(1.4), Inches(0.3), phase, 10, WHITE, True, PP_ALIGN.CENTER)
    _txt(s, Inches(2.1), y+Inches(0.02), Inches(1.3), Inches(0.3), weeks, 10, color, True)
    _txt(s, Inches(3.5), y+Inches(0.02), Inches(2.5), Inches(0.3), name, 12, NAVY, True)
    tf = _txt(s, Inches(0.8), y+Inches(0.4), Inches(11.5), Inches(0.8), items[0], 9, BODY)
    for item in items[1:]:
        c = NAVY if item.startswith('DELIVERABLE') else BODY
        b = item.startswith('DELIVERABLE')
        _p(tf, item, 9, c, b)
    y += Inches(1.35)


# ================================================================
# SLIDE 19: VERDICT
# ================================================================
s = _cs('VERDICT', 'What the Combined Solution Delivers',
        'Mapped against the full Volt + Salt Typhoon kill chain and the four limiting factors.', NAVY)

# Left: Covers
_box(s, Inches(0.5), Inches(1.4), Inches(5.8), Inches(5.3), fill=LGREEN, accent=GREEN)
_txt(s, Inches(0.8), Inches(1.45), Inches(5.3), Inches(0.3), 'WHAT IT COVERS', 14, GREEN, True)

covers = [
    'Every known MITRE TTP for Volt + Salt Typhoon -- preempted at config layer (Problem 1)',
    'Living-off-the-Land tradecraft -- caught as behavioral drift, not tool signatures (Problem 2)',
    'Cross-domain identity fusion -- one trajectory across AD/AAD/AWS/Okta/K8s (Problem 3)',
    'Containerized K8s-native deployment for Gabriel Nimbus at IL5/IL6/JWICS (Problem 4)',
    'Zero-day & unmapped TTPs -- caught as deviation from baseline even with no MITRE label',
    'Insider threat & valid-account abuse -- cohort z-score + cosine drift surfaces misuse',
    'Long-dwell stealth (Salt Typhoon signature) -- CUSUM 4-sigma on rolling 30-day baseline',
    'Mathematical proof of NIST 800-53 / SP800-207 / D3FEND control coverage',
]
y = Inches(1.8)
for cover in covers:
    _txt(s, Inches(0.8), y, Inches(5.3), Inches(0.5), cover, 9.5, BODY)
    y += Inches(0.55)

# Right: Remaining gaps
_box(s, Inches(6.8), Inches(1.4), Inches(5.8), Inches(3.8), fill=RGBColor(0xFD,0xEB,0xEB), accent=CRED)
_txt(s, Inches(7.1), Inches(1.45), Inches(5.3), Inches(0.3), 'REMAINING GAPS (BY DESIGN)', 14, CRED, True)

gaps = [
    ('Endpoint runtime enforcement', 'Pair with CrowdStrike / MS Defender / SentinelOne'),
    ('Email & phishing initial access', 'Pair with Proofpoint / Defender for O365'),
    ('OT/ICS-specific protocols', 'Pair with Dragos / Claroty / Nozomi / TXOne'),
    ('Physical & supply-chain', 'Pair with TPRM platforms'),
    ('Human factors', 'Awareness training + tabletops'),
]
y = Inches(1.85)
for gap, pair in gaps:
    _txt(s, Inches(7.1), y, Inches(5.3), Inches(0.2), gap, 10, CRED, True)
    _txt(s, Inches(7.1), y+Inches(0.22), Inches(5.3), Inches(0.2), pair, 9, SUB)
    y += Inches(0.55)

# Recommended stack
_box(s, Inches(6.8), Inches(5.4), Inches(5.8), Inches(1.3), fill=CARD, accent=NAVY)
_txt(s, Inches(7.1), Inches(5.45), Inches(5.3), Inches(0.25), 'RECOMMENDED FULL STACK', 11, NAVY, True)
tf = _txt(s, Inches(7.1), Inches(5.75), Inches(5.3), Inches(0.8),
     'Preemptive + ACECARD UEBA (core)', 10, NAVY, True)
_p(tf, '+ EDR/XDR (endpoint runtime)', 9, BODY)
_p(tf, '+ Email Security (initial access vector)', 9, BODY)
_p(tf, '+ NDR + SIEM (east-west visibility + log fabric)', 9, BODY)

_bb(s, 'Preemptive + ACECARD UEBA = comprehensive defense against Volt + Salt Typhoon. Containerized for Gabriel Nimbus.')


# ================================================================
# SLIDE 20: CLOSING -- NEXT STEPS
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_fill(s, DKNAVY)
_txt(s, Inches(1), Inches(1.0), Inches(11.3), Inches(0.5),
     '22nd Century Technologies Inc.', 18, GOLD, True, PP_ALIGN.CENTER, 'Georgia')
_txt(s, Inches(1), Inches(1.8), Inches(11.3), Inches(1.0),
     'BUILD YOUR DOME', 40, WHITE, True, PP_ALIGN.CENTER, 'Georgia')
_txt(s, Inches(1), Inches(3.0), Inches(11.3), Inches(0.5),
     'Preemptive closes the door. ACECARD watches the room.\nTogether: comprehensive defense for Gabriel Nimbus.', 16,
     RGBColor(0xA0, 0xB0, 0xC0), False, PP_ALIGN.CENTER)

_txt(s, Inches(1), Inches(4.0), Inches(11.3), Inches(0.4),
     'NEXT STEPS', 16, GOLD, True, PP_ALIGN.CENTER, 'Georgia')

steps = [
    ('1.', 'Pilot Preemptive', 'Formal Policy Analysis on edge defense estate against Volt + Salt Typhoon TTPs'),
    ('2.', 'Pilot ACECARD UEBA', 'Replay Volt + Salt Typhoon scenarios on your telemetry; tune thresholds + cohorts'),
    ('3.', 'Integrate', 'Preemptive verified findings raise ACECARD trust floor on entities traversing closed gaps'),
    ('4.', 'Deploy to Gabriel Nimbus', 'Containerized K8s / Iron Bank / IL5-IL6-JWICS / cATO / continuous re-verification'),
]
y = Inches(4.5)
for num, title, desc in steps:
    _txt(s, Inches(2.5), y, Inches(0.5), Inches(0.35), num, 14, GOLD, True, fn='Georgia')
    _txt(s, Inches(3.0), y, Inches(2.5), Inches(0.35), title, 13, WHITE, True)
    _txt(s, Inches(5.5), y, Inches(6.0), Inches(0.35), desc, 11, RGBColor(0xA0, 0xB0, 0xC0))
    y += Inches(0.45)

_txt(s, Inches(1), Inches(6.5), Inches(11.3), Inches(0.3),
     'Ravindra Shukla, Head of AI  |  ravindra.shukla@gmail.com  |  22nd Century Technologies Inc. (TSCTI)', 12, RGBColor(0xA0, 0xB0, 0xC0), False, PP_ALIGN.CENTER)
_txt(s, Inches(1), Inches(6.9), Inches(11.3), Inches(0.3),
     'CFIC Collaboration Event  |  Augusta, GA  |  May 6, 2026  |  CONFIDENTIAL', 10, RGBColor(0x70, 0x80, 0x90), False, PP_ALIGN.CENTER)


# ================================================================
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
