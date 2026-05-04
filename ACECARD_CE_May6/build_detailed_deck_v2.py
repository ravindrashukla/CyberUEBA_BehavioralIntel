"""Build detailed deep-dive deck v2: LIGHT BACKGROUND.
Part 1: Preemptive Cyber (Rigor AI) — from reference deck pages 4-13
Part 2: Behavioral (ACECARD UEBA) — signals, metrics, thresholds
Part 3: Integration — entity fusion, dimensions, signature, CUSUM, drift direction, ABAC"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

OUT = os.path.join(os.path.dirname(__file__), "06_Detailed_Deep_Dive_Deck.pptx")

NAVY   = RGBColor(0x00, 0x2B, 0x5C)
GOLD   = RGBColor(0xB8, 0x86, 0x2A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BG     = RGBColor(0xF7, 0xF8, 0xFA)
CARD   = RGBColor(0xFF, 0xFF, 0xFF)
BLUE   = RGBColor(0x1E, 0x4D, 0x8C)
PURPLE = RGBColor(0x6C, 0x2E, 0xD9)
TEAL   = RGBColor(0x0D, 0x8B, 0x7C)
ORANGE = RGBColor(0xD9, 0x7A, 0x06)
RED    = RGBColor(0xCC, 0x33, 0x33)
CYBER  = RGBColor(0x00, 0x99, 0x7A)
TXT    = RGBColor(0x22, 0x22, 0x33)  # main text
SUB    = RGBColor(0x55, 0x5E, 0x70)  # subtitle
BDR    = RGBColor(0xDD, 0xDD, 0xE5)  # borders
ALT    = RGBColor(0xF0, 0xF2, 0xF5)  # alternating row

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width


def _fill(s, c):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = c

def _box(s, l, t, w, h, fill=None, border=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.background()
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    ln = sh.line
    if border:
        ln.color.rgb = border
        ln.width = Pt(1)
    else:
        ln.fill.background()
    sh.shadow.inherit = False
    return sh

def _txt(s, l, t, w, h, text, sz=14, c=TXT, b=False, al=PP_ALIGN.LEFT, fn='Segoe UI'):
    tf = s.shapes.add_textbox(l, t, w, h).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = c
    p.font.bold = b
    p.font.name = fn
    p.alignment = al
    return tf

def _p(tf, text, sz=14, c=TXT, b=False, sp=Pt(4)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = c
    p.font.bold = b
    p.font.name = 'Segoe UI'
    p.space_before = sp
    return p

def _title(s, title, sub=None, part=None):
    _box(s, Inches(0), Inches(0), W, Inches(1.1), fill=NAVY)
    _txt(s, Inches(0.6), Inches(0.15), Inches(10), Inches(0.5), title, 26, WHITE, True)
    if sub:
        _txt(s, Inches(0.6), Inches(0.6), Inches(10), Inches(0.4), sub, 13, RGBColor(0xCC, 0xDD, 0xEE))
    if part:
        _txt(s, Inches(10.2), Inches(0.25), Inches(2.8), Inches(0.5), part, 12, GOLD, True, PP_ALIGN.RIGHT)

def ns(title, sub=None, part='DEEP DIVE'):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _fill(s, BG)
    _title(s, title, sub, part)
    return s

def mtable(s, x, y, w, rows, hc=NAVY):
    rh = Inches(0.32)
    hh = Inches(0.28)
    cw = [w*0.17, w*0.33, w*0.17, w*0.16, w*0.17]
    hdrs = ['METRIC', 'WHAT IT MEASURES', 'NORMAL', 'WARNING', 'CRITICAL']
    _box(s, x, y, w, hh, fill=hc)
    cx = x
    for i, h in enumerate(hdrs):
        _txt(s, cx+Pt(4), y+Pt(1), cw[i]-Pt(8), hh, h, 8, WHITE, True)
        cx += cw[i]
    for ri, rd in enumerate(rows):
        ry = y + hh + ri * rh
        bg = CARD if ri % 2 == 0 else ALT
        _box(s, x, ry, w, rh, fill=bg, border=BDR)
        cx = x
        for ci, v in enumerate(rd):
            clr = NAVY if ci == 0 else (RED if ci == 4 else (ORANGE if ci == 3 else TXT))
            _txt(s, cx+Pt(4), ry+Pt(1), cw[ci]-Pt(8), rh, v, 9, clr, ci==0, fn='Cascadia Code' if ci==0 else 'Segoe UI')
            cx += cw[ci]
    return y + hh + len(rows) * rh


# ================================================================
# SLIDE 1: TITLE
# ================================================================
s = ns('Detailed Technical Reference', 'Preemptive Cyber (Rigor AI) + Behavioral Intelligence (ACECARD UEBA)')
_txt(s, Inches(1), Inches(2.0), Inches(11), Inches(0.5),
     'Supporting Material for ACECARD CE Solution Brief', 22, NAVY, True, PP_ALIGN.CENTER)
_txt(s, Inches(1), Inches(3.0), Inches(5.5), Inches(3.5),
     'PART 1: PREEMPTIVE CYBER (RIGOR AI)\n\n'
     '  1. The Rigor AI Approach\n'
     '  2. RigorOS Operational Flow\n'
     '  3. How Rigor Stops Volt Typhoon (Phase by Phase)\n'
     '  4. How Rigor Stops Salt Typhoon\n'
     '  5. Volt Typhoon TTP Coverage (MITRE)\n'
     '  6. Salt Typhoon TTP Coverage (MITRE)\n'
     '  7. Nine Rigor Use Cases\n'
     '  8. Where Rigor Needs Help (Why ACECARD)\n'
     '  9. Rigor Production Proof Points',
     13, TXT)
_txt(s, Inches(6.8), Inches(3.0), Inches(5.5), Inches(3.5),
     'PART 2: BEHAVIORAL (ACECARD UEBA)\n\n'
     '  10. Dual Embedding Architecture\n'
     '  11-15. All 5 Signals (metrics + thresholds)\n'
     '  16. Why Thresholds Cannot Catch LoTL\n\n'
     'PART 3: INTEGRATION\n\n'
     '  17. Entity Fusion\n'
     '  18. 7 Entity Dimensions\n'
     '  19. 4-Component Signature\n'
     '  20. CUSUM Algorithm\n'
     '  21. Drift Direction (8 Concepts)\n'
     '  22. ABAC Trust States\n'
     '  23. How We Solve Rule Decay\n'
     '  24. Combined Coverage Matrix\n'
     '  25. Key Numbers',
     13, TXT)

# ================================================================
# PART 1: RIGOR AI — PREEMPTIVE CYBER
# ================================================================

# SLIDE 2: RIGOR APPROACH
s = ns('The Rigor AI Approach: Close the Door Before They Knock',
       'Mathematically rigorous, preemptive cyber defense -- Complete, Continuous, Verifiable',
       'PART 1: RIGOR AI')

problems = [
    ('Analyze All Known Threat Intelligence', 'Terabytes of CVE, advisories, CTI, IOCs, IOAs, PSIRTs',
     'Attack Intelligence', 'MITRE-enriched attack graphs built by ML/NLP/LLM pipeline. Full TTP coverage of Volt Typhoon (G1017) and Salt Typhoon.', BLUE),
    ('Identify Every Defensive Control Gap', 'Exponential combination of NGFW, IdP, IPS, SASE, WAF configs',
     'Defense Intelligence', 'Symbolic Model of Computation. Formal mathematical model -- exhaustively reasons over the 2^8000 state space. No sampling.', PURPLE),
    ('Prescribe Correct Configurations', 'Without creating new errors or business disruption',
     'Remediation Intelligence', 'Grounded agentic-AI reasoners. Risk-prioritized, vendor-specific fixes. No false positives. No false negatives.', TEAL),
]
y = Inches(1.3)
for prob, probdesc, sol, soldesc, color in problems:
    _box(s, Inches(0.3), y, Inches(6.2), Inches(1.75), fill=CARD, border=BDR)
    _box(s, Inches(0.35), y+Inches(0.05), Inches(1.2), Inches(0.22), fill=RED)
    _txt(s, Inches(0.4), y+Inches(0.06), Inches(1.1), Inches(0.2), 'PROBLEM', 9, WHITE, True, PP_ALIGN.CENTER)
    _txt(s, Inches(1.7), y+Inches(0.05), Inches(4.6), Inches(0.25), prob, 12, TXT, True)
    _txt(s, Inches(1.7), y+Inches(0.35), Inches(4.6), Inches(0.6), probdesc, 10, SUB)

    _box(s, Inches(6.8), y, Inches(6.2), Inches(1.75), fill=CARD, border=color)
    _box(s, Inches(6.85), y+Inches(0.05), Inches(1.2), Inches(0.22), fill=color)
    _txt(s, Inches(6.9), y+Inches(0.06), Inches(1.1), Inches(0.2), 'SOLUTION', 9, WHITE, True, PP_ALIGN.CENTER)
    _txt(s, Inches(8.2), y+Inches(0.05), Inches(4.6), Inches(0.25), sol, 12, color, True)
    _txt(s, Inches(8.2), y+Inches(0.35), Inches(4.6), Inches(1.2), soldesc, 10, TXT)
    y += Inches(1.85)

_box(s, Inches(0.3), Inches(6.9), Inches(12.7), Inches(0.4), fill=NAVY)
_txt(s, Inches(0.5), Inches(6.92), Inches(12.3), Inches(0.35),
     "The world's first Mathematically Rigorous -- Complete, Preemptive, Continuous, Verifiable -- Cyber Defense Platform.",
     12, GOLD, True, PP_ALIGN.CENTER)


# SLIDE 3: RIGOR OPERATIONAL FLOW
s = ns('How RigorOS Operates: Inputs -> Reasoning -> Continuous Remediation',
       'SaaS, on-prem, and air-gapped deployments. Integrates with ServiceNow, JIRA, and major SOC ticketing.',
       'PART 1: RIGOR AI')

# Inputs
_box(s, Inches(0.3), Inches(1.3), Inches(3.5), Inches(2.5), fill=CARD, border=BLUE)
_box(s, Inches(0.3), Inches(1.3), Inches(3.5), Inches(0.35), fill=BLUE)
_txt(s, Inches(0.5), Inches(1.32), Inches(3.1), Inches(0.3), 'INPUTS', 14, WHITE, True)
tf = _txt(s, Inches(0.5), Inches(1.75), Inches(3.1), Inches(1.8),
     'Threat Intelligence', 11, TXT, True)
_p(tf, 'CVE, advisories, CTI, IOCs, IOAs, PSIRTs', 10, SUB)
_p(tf, 'Tech Controls Config', 11, TXT, True)
_p(tf, 'NGFW, IdP, IAM, IPS, EPP, SASE, WAF rules', 10, SUB)
_p(tf, 'Grounded World Models', 11, TXT, True)
_p(tf, 'Vendor capabilities, MITRE ATT&CK / D3FEND', 10, SUB)

# Engine
_box(s, Inches(4.1), Inches(1.3), Inches(4.8), Inches(2.5), fill=CARD, border=NAVY)
_box(s, Inches(4.1), Inches(1.3), Inches(4.8), Inches(0.35), fill=NAVY)
_txt(s, Inches(4.3), Inches(1.32), Inches(4.4), Inches(0.3), 'RIGOR ENGINE', 14, WHITE, True)
steps = [
    ('1. Identify Risk-Prioritized Control Gaps', 'Model attacks from threat intel; model customer cyber defenses; reason about control gaps'),
    ('2. Generate Remediations', 'Risk-prioritized fixes; vendor-specific playbooks; alerts to SOC + Control Admin'),
    ('3. Fix Control Gaps', 'Apply remediations (optional); validate; re-verify continuously'),
]
y = Inches(1.75)
for title, desc in steps:
    _txt(s, Inches(4.3), y, Inches(4.4), Inches(0.2), title, 11, NAVY, True)
    _txt(s, Inches(4.3), y+Inches(0.22), Inches(4.4), Inches(0.35), desc, 9, SUB)
    y += Inches(0.55)

# Outputs
_box(s, Inches(9.2), Inches(1.3), Inches(3.8), Inches(2.5), fill=CARD, border=TEAL)
_box(s, Inches(9.2), Inches(1.3), Inches(3.8), Inches(0.35), fill=TEAL)
_txt(s, Inches(9.4), Inches(1.32), Inches(3.4), Inches(0.3), 'CONTINUOUS ASSURANCE', 14, WHITE, True)
tf = _txt(s, Inches(9.4), Inches(1.75), Inches(3.4), Inches(1.8),
     'Risk-prioritized control gap findings', 10, TXT)
_p(tf, 'Vendor-specific remediation playbooks', 10, TXT)
_p(tf, 'Mathematical proof -- no false positives', 10, NAVY, True)
_p(tf, 'Mathematical proof -- no false negatives', 10, NAVY, True)
_p(tf, 'SOC and control admin workflows', 10, TXT)
_p(tf, 'ServiceNow / JIRA integration', 10, TXT)
_p(tf, 'Auto re-verify on every config change', 10, TXT)
_p(tf, 'Auto re-verify on every threat-feed update', 10, TXT)

# Supported vendors
_box(s, Inches(0.3), Inches(4.1), Inches(12.7), Inches(0.5), fill=CARD, border=BDR)
_txt(s, Inches(0.5), Inches(4.15), Inches(12.3), Inches(0.4),
     'Currently supports: Palo Alto Networks, Fortinet, Cisco, Check Point (N-S and E-W NGFWs). Identity, SASE, and additional vendors on roadmap.',
     11, SUB)

# What it PROVES
_box(s, Inches(0.3), Inches(4.8), Inches(12.7), Inches(2.4), fill=CARD, border=NAVY)
_box(s, Inches(0.3), Inches(4.8), Inches(12.7), Inches(0.35), fill=NAVY)
_txt(s, Inches(0.5), Inches(4.82), Inches(12.3), Inches(0.3), 'WHAT RIGOR PROVES (MATHEMATICAL CERTAINTY)', 14, WHITE, True)

proofs = [
    ('Edge-device exploit paths are blocked', 'CVE-2023-20198, CVE-2024-21887, CVE-2024-3400 etc. verified closed'),
    ('Threat-feed C2 IPs/domains denied', 'Across 50M+ rule combinations. Known KV botnet, JumbledPath C2'),
    ('Config drift from baseline is caught', 'Every hour or per config change. Salt Typhoon ACL tampering = caught'),
    ('Critical segmentation holds end-to-end', 'User VLAN cannot reach DC mgmt. CALEA zone isolated. Zero-trust verified.'),
    ('Shadowed/zombie rules surfaced', '3 higher-priority rules shadowing app-control (real customer finding)'),
    ('Compensating controls for unpatchable CVEs', 'When no patch exists, identifies alternate control that closes the path'),
]
y = Inches(5.25)
for title, desc in proofs:
    _txt(s, Inches(0.5), y, Inches(4.5), Inches(0.22), title, 11, NAVY, True)
    _txt(s, Inches(5.1), y, Inches(7.7), Inches(0.22), desc, 10, SUB)
    y += Inches(0.3)


# SLIDE 4: RIGOR vs VOLT TYPHOON — PHASE BY PHASE
s = ns('How Rigor AI Stops Volt Typhoon -- Phase by Phase',
       'Volt Typhoon needs configuration weakness to enter and pivot. Rigor proves those weaknesses do not exist.',
       'PART 1: RIGOR AI')

cols_x = [Inches(0.3), Inches(2.2), Inches(5.8), Inches(9.6)]
cols_w = [Inches(1.8), Inches(3.5), Inches(3.7), Inches(3.4)]
hdrs = ['PHASE', 'VOLT TYPHOON BEHAVIOR', 'HOW RIGOR AI STOPS IT', 'CVEs / TTPs']

_box(s, Inches(0.25), Inches(1.25), Inches(12.8), Inches(0.38), fill=NAVY)
for i, h in enumerate(hdrs):
    _txt(s, cols_x[i], Inches(1.27), cols_w[i], Inches(0.34), h, 10, WHITE, True)

rows_vt = [
    ('Phase 1\nInitial Access\n0-12h',
     'Scans for unpatched Fortinet, Ivanti, Versa, PAN edge devices. Exploits CVE to land on perimeter.',
     'Attack Intel ingests CISA/NSA/vendor PSIRT advisories. Defense Intel proves every NGFW/VPN/SASE has IPS signature or threat-feed deny rule enforced. If no patch, identifies compensating control.',
     'CVE-2023-20198\nCVE-2024-21887\nCVE-2024-3400\nT1190'),
    ('Phase 2\nPersistence\n12-48h',
     'Drops .aspx/.jspx web shells. Routes C2 through KV botnet of compromised SOHO routers.',
     'Models every threat-feed IP, domain, ASN, SOHO fingerprint as denied destination. Proves 50M+ rule combinations block these flows. WAF verified for web shell upload inspection.',
     'KV Botnet\nT1505 (Web Shell)\nT1071 (App Layer C2)'),
    ('Phase 3\nCredential Access\n48-96h',
     'ntdsutil.exe to dump NTDS.dit. LSASS memory extraction. Exfiltrates AD credentials to C2.',
     'Verifies EDR/EPP coverage enabled on every DC. Script-control profile blocks ntdsutil by non-admin. Egress filtering blocks credential exfil to known C2 ranges. Missing = exact vendor fix.',
     'T1003 (Cred Dump)\nntdsutil, LSASS\nT1059 (Scripting)'),
    ('Phase 4\nLateral Movement\n96h+',
     'SMB/RDP/WinRM to domain controllers and file servers. Admin share access.',
     'Real finding: 3 higher-priority rules shadowed SMB IPS profile -- Rigor closed lateral path. Micro-segmentation verified user VLAN cannot reach DC management ports.',
     'T1021.002 (SMB)\nT1021.001 (RDP)\nT1078 (Valid Accts)'),
    ('Phase 5\nPre-positioning\nWeeks-Years',
     'Credential harvest + network mapping for future sabotage during crisis.',
     'Continuous re-verification: every config change triggers re-analysis. Zero-trust segmentation re-proven hourly. IdP/IAM coverage modeled; trust paths constrained.',
     'T1133 (Ext Remote)\nT1556 (Modify Auth)\n5+ year dwell'),
]

y = Inches(1.7)
for i, (phase, attack, rigor, ttps) in enumerate(rows_vt):
    bg = CARD if i % 2 == 0 else ALT
    _box(s, Inches(0.25), y, Inches(12.8), Inches(1.05), fill=bg, border=BDR)
    _txt(s, cols_x[0], y+Pt(2), cols_w[0], Inches(1.0), phase, 9, NAVY, True)
    _txt(s, cols_x[1], y+Pt(2), cols_w[1], Inches(1.0), attack, 9, TXT)
    _txt(s, cols_x[2], y+Pt(2), cols_w[2], Inches(1.0), rigor, 9, BLUE)
    _txt(s, cols_x[3], y+Pt(2), cols_w[3], Inches(1.0), ttps, 8, RED, fn='Cascadia Code')
    y += Inches(1.08)


# SLIDE 5: RIGOR vs SALT TYPHOON
s = ns('How Rigor AI Stops Salt Typhoon -- Carrier-Scale Edge Defense',
       'Salt Typhoon lives on misconfigured Cisco/Ivanti/PAN edge gear. Rigor verifies every config every hour.',
       'PART 1: RIGOR AI')

_box(s, Inches(0.25), Inches(1.25), Inches(12.8), Inches(0.38), fill=NAVY)
for i, h in enumerate(hdrs):
    _txt(s, cols_x[i], Inches(1.27), cols_w[i], Inches(0.34), h, 10, WHITE, True)

rows_st = [
    ('Phase 1\nEdge CVE Exploit\n0-24h',
     'Exploits Cisco IOS XE web UI (CVE-2023-20198) for level-15 admin. Also Ivanti, PAN, Citrix. 1000+ Cisco devices targeted Dec 2024-Jan 2025.',
     'CVE-to-control mapping is automatic. Every device exposing IOS XE web UI to internet = flagged. Proves WAF/mgmt-plane ACL blocks /webui/. For unpatched CVE, identifies compensating control.',
     'CVE-2023-20198\nCVE-2018-0171\nCVE-2024-3400\nT1190'),
    ('Phase 2\nPersistence\n1-7d',
     'Creates local privileged accounts. Modifies router ACLs to whitelist attacker IPs. Enables Guest Shell, sshd_operns. 3+ years access in some victims.',
     'Config Drift Management: takes approved baseline of every router ACL, AAA config, mgmt-plane policy. Compares live state every hour. Any unauthorized change = immediate gap finding with exact diff.',
     'T1556 (Modify Auth)\nLocal account creation\nACL tampering\n3+ year dwell'),
    ('Phase 3\nTACACS+ Capture\n7-30d',
     'Captures TACACS+ traffic on TCP/49 via on-box packet capture. Harvests credentials. Uses JumbledPath to chain remote connections.',
     'Mgmt-plane segmentation verified. On-box capture restricted to authorized roles. Outbound rules to JumbledPath C2 IPs blocked across all internal zones.',
     'JumbledPath malware\nTACACS+ TCP/49\nT1040 (Sniffing)'),
    ('Phase 4\nGRE Tunnel + C2\n30-90d',
     'GRE tunnel configured. Encrypted C2. Multi-hop connections across compromised routers.',
     'GRE peer authorization verified. Unauthorized tunnel endpoints detected in routing config. Egress policy verified: only approved protocols/ports permitted outbound.',
     'GRE tunneling\nT1572 (Protocol Tunnel)\nT1071 (App Layer C2)'),
    ('Phase 5\nCALEA + CDR Exfil\n90d+',
     'Accesses CALEA lawful intercept systems. Steady CDR/metadata exfiltration over months. Targets: government officials.',
     'CALEA-zone segmentation intent proven end-to-end. Privilege verification on interception systems. Outbound restriction policies verified for CDR stores. DLP coverage modeled.',
     'CALEA access\nCDR exfiltration\nT1041 (Exfil over C2)\n200+ telcos'),
]

y = Inches(1.7)
for i, (phase, attack, rigor, ttps) in enumerate(rows_st):
    bg = CARD if i % 2 == 0 else ALT
    _box(s, Inches(0.25), y, Inches(12.8), Inches(1.05), fill=bg, border=BDR)
    _txt(s, cols_x[0], y+Pt(2), cols_w[0], Inches(1.0), phase, 9, NAVY, True)
    _txt(s, cols_x[1], y+Pt(2), cols_w[1], Inches(1.0), attack, 9, TXT)
    _txt(s, cols_x[2], y+Pt(2), cols_w[2], Inches(1.0), rigor, 9, BLUE)
    _txt(s, cols_x[3], y+Pt(2), cols_w[3], Inches(1.0), ttps, 8, RED, fn='Cascadia Code')
    y += Inches(1.08)


# SLIDE 6: VOLT TYPHOON TTP COVERAGE
s = ns('Rigor AI Preempts Volt Typhoon -- TTP by TTP',
       'Each MITRE technique -> formal control gap query -> mathematically verified blocking configuration',
       'PART 1: RIGOR AI')

ttp_hdrs = ['MITRE Technique', 'Volt Typhoon Behavior', 'Rigor AI Preemptive Action']
ttp_x = [Inches(0.3), Inches(2.8), Inches(7.0)]
ttp_w = [Inches(2.4), Inches(4.1), Inches(6.0)]

_box(s, Inches(0.25), Inches(1.25), Inches(12.8), Inches(0.38), fill=NAVY)
for i, h in enumerate(ttp_hdrs):
    _txt(s, ttp_x[i], Inches(1.27), ttp_w[i], Inches(0.34), h, 10, WHITE, True)

ttp_vt = [
    ('T1190\nExploit Public App', 'Exploits unpatched Fortinet, Versa Director, Ivanti, PAN edge devices',
     'Formally verifies every NGFW/SASE/WAF rule blocks known CVE traffic; flags weak/missing inspection profiles'),
    ('T1133\nExternal Remote Svc', 'Abuses VPN, RDP gateways with stolen valid accounts',
     'Proves zero-trust segmentation intent; surfaces shadowed allow-rules that bypass MFA enforcement'),
    ('T1021.002\nSMB / Admin Shares', 'Lateral movement via SMBv1, admin shares between branches',
     'Real customer case: 3 higher-priority rules shadowed SMB IPS profile -- Rigor closed the lateral path'),
    ('T1003\nOS Credential Dump', 'ntdsutil, LSASS dumps; credentials exfiltrated to C2',
     'Verifies egress filtering blocks C2 destinations; ensures EDR/IPS profiles cover credential-theft signatures'),
    ('T1071\nC2 over SOHO Routers', 'KV botnet relays C2 through compromised home/SMB routers',
     'Models threat-feed IPs/domains; proves outbound rules block them across all internal zones'),
    ('T1059\nPowerShell / LOLBins', 'In-memory exec via wmic, ntdsutil, netsh; log evasion',
     'Verifies EDR/EPP coverage in policy; identifies hosts where script-control profile is not applied'),
    ('T1078\nValid Accounts', 'Pivots laterally with stolen credentials across trust zones',
     'IdP/IAM coverage modeled; inter-zone trust paths surfaced and constrained -> handed to ACECARD for behavioral check'),
]

y = Inches(1.7)
for i, (tech, behavior, action) in enumerate(ttp_vt):
    bg = CARD if i % 2 == 0 else ALT
    _box(s, Inches(0.25), y, Inches(12.8), Inches(0.72), fill=bg, border=BDR)
    _txt(s, ttp_x[0], y+Pt(2), ttp_w[0], Inches(0.68), tech, 10, RED, True, fn='Cascadia Code')
    _txt(s, ttp_x[1], y+Pt(2), ttp_w[1], Inches(0.68), behavior, 9, TXT)
    _txt(s, ttp_x[2], y+Pt(2), ttp_w[2], Inches(0.68), action, 9, BLUE)
    y += Inches(0.75)

_txt(s, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.3),
     'Source: MITRE ATT&CK Group G1017 (Volt Typhoon); CISA AA24-038A; Microsoft Threat Intelligence', 9, SUB)


# SLIDE 7: NINE RIGOR USE CASES
s = ns('Nine Rigor AI Use Cases',
       'Pre-built reasoning templates. Five directly preempt Volt + Salt Typhoon TTPs.',
       'PART 1: RIGOR AI')

ucs = [
    ('01', 'Preemptive Cyber Defense', 'Continuous threat-intel + config -> proof that every TTP path is blocked',
     'Volt T1190, T1133, T1071 / Salt T1190, T1556', BLUE),
    ('02', 'Configuration Drift Mgmt', 'Live config vs approved baseline every hour. Surfaces unauthorized changes.',
     'Salt Typhoon ACL tampering, Guest Shell, local accounts', BLUE),
    ('03', 'Formal Policy Analysis', 'Reasons over 2^8000 firewall rule state space. Shadow rules, conflicts, weak inspection.',
     'SMB lateral paths, shadowed inspection profiles', BLUE),
    ('04', 'Preemptive ASM (External)', 'Scans external attack surface. Maps exposed assets to compensating controls.',
     'Edge-device CVE exploitation (CVE-2023-20198, etc.)', BLUE),
    ('05', 'Internal Attack Surface Mgmt', 'Models internal paths from any compromised host to crown-jewel assets.',
     'Cross-tenant pivots (Salt Typhoon telco -> enterprise)', BLUE),
    ('06', 'NIST 800-53 / Zero-Trust', 'Proves zero-trust segmentation holds end-to-end. Maps to NIST CSF / D3FEND.',
     'Continuous SP800-207 verification', TEAL),
    ('07', 'Vendor Migration Safety', 'Models before/after config during NGFW vendor migration. Proves no regression.',
     'PAN->Fortinet migrations, Cisco->Palo refresh cycles', TEAL),
    ('08', 'M&A Due Diligence', 'Models combined defense posture pre-merger. Surfaces gaps from acquired entity.',
     'Risk assessment before integration', TEAL),
    ('09', 'Third-Party Risk', 'Models third-party access paths and verifies isolation controls.',
     'Supply chain and vendor risk management', TEAL),
]

y = Inches(1.3)
for num, title, desc, stops, color in ucs:
    _box(s, Inches(0.3), y, Inches(12.7), Inches(0.62), fill=CARD, border=BDR)
    _box(s, Inches(0.35), y+Inches(0.05), Inches(0.5), Inches(0.28), fill=color)
    _txt(s, Inches(0.4), y+Inches(0.07), Inches(0.4), Inches(0.24), num, 11, WHITE, True, PP_ALIGN.CENTER)
    _txt(s, Inches(1.0), y+Inches(0.05), Inches(2.5), Inches(0.24), title, 12, NAVY, True)
    _txt(s, Inches(3.6), y+Inches(0.05), Inches(5.0), Inches(0.5), desc, 10, TXT)
    _txt(s, Inches(8.8), y+Inches(0.05), Inches(4.0), Inches(0.5), stops, 9, RED)
    y += Inches(0.65)

_box(s, Inches(0.3), Inches(7.1), Inches(12.7), Inches(0.25), fill=NAVY)
_txt(s, Inches(0.5), Inches(7.12), Inches(12.3), Inches(0.2),
     'Use cases 01-05 (blue) directly preempt Volt + Salt Typhoon TTPs', 10, GOLD, True, PP_ALIGN.CENTER)


# SLIDE 8: WHERE RIGOR NEEDS HELP
s = ns('Where Rigor Needs Help -- Why ACECARD Belongs Next to It',
       'Rigor closes the configuration attack surface. Three classes of attack remain.',
       'PART 1: RIGOR AI')

gaps = [
    ('1. Zero-Day & Unknown TTPs', RED,
     'Rigor preempts every TTP mapped to MITRE ATT&CK. A genuine zero-day with no signature, no IOC, no MITRE mapping yet has nothing for Rigor to prove against.',
     'Example: A Volt Typhoon exploit chain published the day after Rigor\'s threat-intel sync. In the gap, behavioral detection is the only line of defense.',
     'ACECARD detects the resulting behavioral drift even with no MITRE label. Post-exploitation behavior (lateral movement, credential access) always creates drift.'),
    ('2. Compromised Valid Accounts', ORANGE,
     'Both Typhoons abuse stolen, valid credentials. A correctly-configured firewall, IdP, and EDR will ALLOW that traffic -- by design. Configuration is not the attack surface; behavior is.',
     'Example: Salt Typhoon harvests TACACS+ credentials from router PCAP, then logs in as a legitimate network admin. Every Rigor-verified control says "approved."',
     'ACECARD\'s cohort z-score and cosine drift catch the abuse -- same admin credentials, structurally different behavior from the real admin.'),
    ('3. Living-off-the-Land Tradecraft', PURPLE,
     'ntdsutil, wmic, netsh, rundll32, PowerShell -- every binary Volt Typhoon uses is a legitimate Windows admin tool. Restricting them entirely is operationally infeasible.',
     'Example: Volt Typhoon runs ntdsutil.exe ifm "create full c:\\temp" -- a real administrative command. Rigor verifies EDR coverage is on. The command itself is the threat.',
     'ACECARD\'s 5-signal embedding captures the COMBINATION of LoTL tools, timing, and targets. Drift direction: lotl_execution (T1059). Only behavioral context separates admin from adversary.'),
]

y = Inches(1.3)
for title, color, problem, example, acecard in gaps:
    _box(s, Inches(0.3), y, Inches(12.7), Inches(1.8), fill=CARD, border=color)
    _box(s, Inches(0.35), y+Inches(0.05), Inches(12.55), Inches(0.3), fill=color)
    _txt(s, Inches(0.5), y+Inches(0.07), Inches(12.3), Inches(0.26), title, 14, WHITE, True)
    _txt(s, Inches(0.5), y+Inches(0.4), Inches(12.3), Inches(0.4), problem, 10, TXT)
    _txt(s, Inches(0.5), y+Inches(0.85), Inches(12.3), Inches(0.35), example, 9, SUB, False)
    _box(s, Inches(0.5), y+Inches(1.25), Inches(0.08), Inches(0.4), fill=CYBER)
    _txt(s, Inches(0.7), y+Inches(1.25), Inches(12.0), Inches(0.5),
         'HANDED TO ACECARD UEBA: ' + acecard, 10, CYBER, True)
    y += Inches(1.88)


# SLIDE 9: RIGOR PROOF POINTS
s = ns('Rigor in Production: Real Customer Findings',
       'Selected findings from financial services, telecom, government, and critical infrastructure',
       'PART 1: RIGOR AI')

findings = [
    ('FinTech | Asia', 'Closed SMB Lateral Movement Path', BLUE,
     'SMB app-control profiles completely shadowed by 3 higher-priority rules -- allowed unprotected SMBv1 between branches. Rigor formal model exposed the multi-rule shadow.',
     'Shadow rule detection across multi-vendor NGFW stack'),
    ('Communications | NA', 'Eliminated Known CVE Exposure', PURPLE,
     'IPS filters were missing on rsync traffic to backup servers. Exposed enterprise to lateral malware movement and remote code execution against critical infrastructure.',
     'Missing IPS inspection on specific protocol'),
    ('Regional Bank | E. Asia', 'Closed RCE Attack Paths', TEAL,
     'Both PANW and Fortinet NGFWs allowed Windows Update service traffic from trusted zones with no IPS -- exposed bank to RCE on critical multi-vendor infrastructure.',
     'Cross-vendor policy inconsistency'),
    ('Critical Infra | SE Asia', 'Closed Application Obfuscation Paths', ORANGE,
     'Incorrectly configured vendor capabilities on DMZ rules -- not all applications on the DMZ were blocked as intended. Susceptible to port-hopping evasion.',
     'DMZ rule misconfiguration'),
    ('Finance | E. Asia', 'NIST 800-53 SMTP Compliance', RED,
     'SMTP traffic from trusted to untrusted zones had no threat inspection enabled, despite NGFW built-in capability -- exposed customer to phishing and credential harvest.',
     'Compliance gap on email traffic'),
    ('EU Sovereign Govt', 'Preemptive ASM for National Agency', NAVY,
     'EU country cyber defense umbrella -- Rigor formally extended EASM scan coverage and proved compliance of external attack surface across 12 government agencies.',
     'National-scale external attack surface management'),
]

y = Inches(1.3)
for customer, title, color, desc, finding_type in findings:
    _box(s, Inches(0.3), y, Inches(12.7), Inches(0.9), fill=CARD, border=BDR)
    _box(s, Inches(0.35), y+Inches(0.05), Inches(2.0), Inches(0.28), fill=color)
    _txt(s, Inches(0.4), y+Inches(0.07), Inches(1.9), Inches(0.24), customer, 9, WHITE, True, PP_ALIGN.CENTER)
    _txt(s, Inches(2.5), y+Inches(0.05), Inches(5.0), Inches(0.25), title, 12, NAVY, True)
    _txt(s, Inches(7.8), y+Inches(0.05), Inches(5.0), Inches(0.25), finding_type, 10, color, True)
    _txt(s, Inches(2.5), y+Inches(0.35), Inches(10.0), Inches(0.5), desc, 10, TXT)
    y += Inches(0.95)

_box(s, Inches(0.3), Inches(7.0), Inches(12.7), Inches(0.35), fill=NAVY)
_txt(s, Inches(0.5), Inches(7.02), Inches(12.3), Inches(0.3),
     '20+ design partners across F500, federal, critical infrastructure, MSSPs, nation states. Engagements in NA, EU, ME, Asia, Pacific.',
     11, GOLD, True, PP_ALIGN.CENTER)


# ================================================================
# PART 2: ACECARD UEBA — BEHAVIORAL
# ================================================================

# SLIDE 10: DUAL EMBEDDING ARCHITECTURE
s = ns('Dual Embedding Architecture: 6 Vectors = 9,216 Dimensions',
       'Each signal individually embedded (1536-d) + one fused composite. Per-signal drift tracking.',
       'PART 2: ACECARD UEBA')

_box(s, Inches(0.3), Inches(1.3), Inches(12.7), Inches(0.8), fill=CARD, border=NAVY)
_txt(s, Inches(0.5), Inches(1.35), Inches(12.3), Inches(0.25), 'THE KEY INSIGHT', 13, NAVY, True)
_txt(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(0.45),
     'Each signal is INDIVIDUALLY embedded at 1536-d (per-signal drift tracking). All 5 are also concatenated into one COMPOSITE 1536-d vector (holistic entity state). '
     'Total: 6 vectors = 9,216 dimensions per entity per window. Analyst sees WHICH signal drives drift + WHAT direction.',
     11, TXT)

sig_colors = [BLUE, PURPLE, TEAL, ORANGE, RED]
sig_names = ['AUTH', 'PROCESS', 'NETWORK', 'FILE', 'IDENTITY']
xp = [Inches(0.3), Inches(2.9), Inches(5.5), Inches(8.1), Inches(10.7)]
for i in range(5):
    bw = Inches(2.4)
    _box(s, xp[i], Inches(2.4), bw, Inches(0.35), fill=sig_colors[i])
    _txt(s, xp[i], Inches(2.42), bw, Inches(0.3), f'Signal {i+1}: {sig_names[i]}', 11, WHITE, True, PP_ALIGN.CENTER)
    _box(s, xp[i], Inches(2.8), bw, Inches(0.4), fill=CARD, border=sig_colors[i])
    _txt(s, xp[i]+Inches(0.1), Inches(2.82), bw-Inches(0.2), Inches(0.35),
         'Individual 1536-d embedding\n-> per-signal drift', 9, TXT, al=PP_ALIGN.CENTER)

_box(s, Inches(0.3), Inches(3.5), Inches(12.7), Inches(0.4), fill=CARD, border=GOLD)
_txt(s, Inches(0.5), Inches(3.52), Inches(12.3), Inches(0.35),
     'All 5 signal texts concatenated -> 1 COMPOSITE 1536-d vector (holistic entity behavioral state)', 11, GOLD, True, PP_ALIGN.CENTER)

_box(s, Inches(0.3), Inches(4.2), Inches(4.0), Inches(1.5), fill=CARD, border=NAVY)
_txt(s, Inches(0.5), Inches(4.25), Inches(3.6), Inches(0.25), 'DIMENSIONS', 12, NAVY, True)
tf = _txt(s, Inches(0.5), Inches(4.55), Inches(3.6), Inches(1.0), '5 individual: 5 x 1,536 = 7,680', 11, TXT)
_p(tf, '1 composite: 1 x 1,536 = 1,536', 11, TXT)
_p(tf, 'TOTAL: 6 vectors = 9,216 dimensions', 12, NAVY, True)

_box(s, Inches(4.5), Inches(4.2), Inches(4.2), Inches(1.5), fill=CARD, border=TEAL)
_txt(s, Inches(4.7), Inches(4.25), Inches(3.8), Inches(0.25), 'WHY INDIVIDUAL + COMPOSITE?', 12, TEAL, True)
tf = _txt(s, Inches(4.7), Inches(4.55), Inches(3.8), Inches(1.0),
     'Composite: "entity behavior changed"', 10, TXT)
_p(tf, 'Individual: "AUTH drove 72% of change"', 10, TXT)
_p(tf, 'Per-signal direction: "PROCESS signal', 10, TXT)
_p(tf, '  drifting toward lotl_execution (T1059)"', 10, NAVY, True)

_box(s, Inches(8.9), Inches(4.2), Inches(4.1), Inches(1.5), fill=CARD, border=BLUE)
_txt(s, Inches(9.1), Inches(4.25), Inches(3.7), Inches(0.25), 'EMBEDDING OPTIONS', 12, BLUE, True)
tf = _txt(s, Inches(9.1), Inches(4.55), Inches(3.7), Inches(1.0),
     'IL5 (NIPR): OpenAI text-embedding-3-small', 10, TXT)
_p(tf, 'IL6 (SIPR): Local Phi-4 (14B) or Mistral 7B', 10, TXT)
_p(tf, 'JWICS: Same local SLM, air-gapped', 10, TXT)
_p(tf, 'Switch: env var only, no code change', 10, NAVY, True)

# Privacy note
_box(s, Inches(0.3), Inches(6.0), Inches(12.7), Inches(0.5), fill=CARD, border=GOLD)
_txt(s, Inches(0.5), Inches(6.02), Inches(12.3), Inches(0.45),
     'No raw event content sent to embedding provider -- only summarized metrics. No usernames, no IP addresses, no file paths. Privacy-preserving by design. JWICS-safe with local SLM.',
     10, GOLD, True)


# SLIDES 11-15: EACH SIGNAL DETAIL (compact — one slide per signal)
signal_data = [
    ('Signal 1: AUTH -- Authentication Behavior', BLUE,
     'WHO logs in, WHERE, WHEN, HOW -- across all identity providers',
     ['Win 4624 (logon), 4625 (failed), 4768/4769 (Kerberos)', 'Okta syslog, Azure AD sign-in, Linux auth.log/PAM', 'AWS CloudTrail (ConsoleLogin, AssumeRole)'],
     ['Brute force (failure_rate spike)', 'Pass-the-Hash (TGS without TGT)', 'Lateral movement (unique_hosts jump)', 'Impossible travel (geo violation)', 'MFA bypass (stolen session tokens)'],
     [('logon_count', 'Auth events per window', '5-50/hr', '>200/hr', '>500/hr'),
      ('failure_rate', 'Failed/total logons', '0.01-0.05', '>0.15', '>0.40'),
      ('unique_hosts', 'Distinct machines authed to', '1-3/hr', '>8/hr', '>15/hr'),
      ('kerberos_ticket_types', 'Unusual TGS without TGT', '0', '>2/hr', '>5/hr'),
      ('off_hours_ratio', 'Off-hours logins / total', '0.0-0.10', '>0.35', '>0.60'),
      ('impossible_travel', 'Geo-impossible logins', '0', '1 event', '>1 event'),
     ]),
    ('Signal 2: PROCESS -- Execution Behavior', PURPLE,
     'WHAT programs run, HOW invoked, process trees, signing, LoTL binaries',
     ['Sysmon EID 1 (process create), EID 3 (net conn), EID 8 (thread inject)', 'CrowdStrike Falcon, MS Defender for Endpoint', 'Linux auditd EXECVE'],
     ['LoTL chain (wmic+ntdsutil+netsh = Volt Typhoon)', 'Encoded PowerShell (obfuscated payloads)', 'Code injection (CreateRemoteThread)', 'Unsigned binary execution', 'Deep process trees (staged execution)'],
     [('unique_processes', 'Distinct process names', '10-40/hr', '>80/hr', '>150/hr'),
      ('cmdline_entropy', 'Shannon entropy of cmdline', '2.0-3.5', '>4.5', '>5.5'),
      ('lolbin_count', 'LoTL binaries used', '0-1/hr', '>3/hr', '>6/hr'),
      ('parent_child_depth', 'Max process tree depth', '2-4', '>6', '>8'),
      ('unsigned_ratio', 'Unsigned binaries / total', '0.0-0.05', '>0.15', '>0.30'),
      ('new_process_rate', 'New procs (vs 30d baseline)', '0.0-0.02', '>0.08', '>0.15'),
     ]),
    ('Signal 3: NETWORK -- Traffic Behavior', TEAL,
     'WHERE data flows, HOW MUCH, HOW REGULARLY (beacon), unusual destinations',
     ['NetFlow/IPFIX (router-level flows)', 'Zeek (full protocol analysis)', 'Palo Alto PAN-OS, Fortinet, Cisco, AWS VPC Flow Logs'],
     ['C2 beaconing (regular-interval callbacks)', 'Data exfiltration (bytes_out inversion)', 'DNS tunneling (high-entropy domains)', 'Network scanning (unique_dest_ips explosion)', 'GRE tunneling (Salt Typhoon)'],
     [('unique_dest_ips', 'Distinct destination IPs', '5-20/hr', '>50/hr', '>100/hr'),
      ('bytes_out', 'Outbound/inbound byte ratio', '0.05-0.30', '>0.60', '>0.85'),
      ('beacon_score', 'Connection interval regularity', '0.0-0.15', '>0.45', '>0.70'),
      ('dns_query_rate', 'DNS queries per minute', '2-15/min', '>40/min', '>80/min'),
      ('geo_anomaly_count', 'Connections to new countries', '0', '>2/hr', '>5/hr'),
      ('admin_share_access', 'C$/ADMIN$/IPC$ connections', '0', '>2/hr', '>5/hr'),
     ]),
    ('Signal 4: FILE -- File System Behavior', ORANGE,
     'WHAT files accessed, created, modified, deleted -- staging, encryption, exfiltration',
     ['Sysmon EID 11 (file created)', 'Win 4663 (object access audit)', 'DLP (Symantec/Forcepoint/Purview), cloud storage audit'],
     ['Ransomware (extension_changes spike)', 'Pre-ransomware (vssadmin delete shadows)', 'Exfil staging (archive creation)', 'Insider theft (sensitive_access increase)', 'Credential theft (SAM, NTDS.dit access)'],
     [('files_created', 'New files per window', '0-20/hr', '>100/hr', '>500/hr'),
      ('files_deleted', 'Files removed per window', '0-5/hr', '>50/hr', '>200/hr'),
      ('sensitive_access_count', 'Restricted path access', '0-2/hr', '>5/hr', '>10/hr'),
      ('archive_creates', 'zip/tar/rar creations', '0-1/day', '>3/hr', '>5/hr'),
      ('extension_change_rate', 'Renamed file extensions', '0', '>10/hr', '>50/hr'),
      ('shadow_copy_ops', 'vssadmin/wbadmin deletes', '0', '1 event', '>1 event'),
     ]),
    ('Signal 5: IDENTITY -- Privilege & Role Behavior', RED,
     'Privilege escalation, group membership, MFA bypass, admin actions',
     ['Win 4672 (special privs), 4728/4732/4756 (group adds)', 'AWS CloudTrail IAM events', 'K8s audit log (RBAC bindings), AD change log'],
     ['Privilege escalation (Domain Admins adds)', 'MFA bypass (stolen session tokens)', 'Service account abuse (interactive logon)', 'Golden ticket (Kerberos delegation)', 'Insider (admin actions from non-admin role)'],
     [('priv_escalations', 'Privilege elevation events', '0-2/hr', '>5/hr', '>10/hr'),
      ('group_adds', 'Security group changes', '0/day', '>1/hr', '>3/hr'),
      ('mfa_bypass_attempts', 'Auth without MFA challenge', '0', '>1/hr', '>3/hr'),
      ('role_changes', 'IAM role assumptions', '0-1/day', '>3/hr', '>5/hr'),
      ('service_acct_use', 'Svc acct interactive logon', '0', '>1/day', '>1/hr'),
      ('admin_actions', 'Account/policy/GPO changes', '0/day', '>2/hr', '>5/hr'),
     ]),
]

for sig_title, sig_color, sig_sub, sources, attacks, metrics in signal_data:
    s = ns(sig_title, sig_sub, 'PART 2: ACECARD UEBA')

    _box(s, Inches(0.3), Inches(1.3), Inches(6.3), Inches(1.6), fill=CARD, border=sig_color)
    _txt(s, Inches(0.5), Inches(1.35), Inches(5.9), Inches(0.25), 'DATA SOURCES', 12, sig_color, True)
    y_src = Inches(1.65)
    for src in sources:
        _txt(s, Inches(0.5), y_src, Inches(5.9), Inches(0.25), src, 10, TXT)
        y_src += Inches(0.28)

    _box(s, Inches(6.8), Inches(1.3), Inches(6.2), Inches(1.6), fill=CARD, border=sig_color)
    _txt(s, Inches(7.0), Inches(1.35), Inches(5.8), Inches(0.25), 'ATTACKS THIS CATCHES', 12, sig_color, True)
    y_atk = Inches(1.65)
    for atk in attacks:
        _txt(s, Inches(7.0), y_atk, Inches(5.8), Inches(0.25), atk, 10, NAVY, True)
        y_atk += Inches(0.28)

    mtable(s, Inches(0.3), Inches(3.1), Inches(12.7), metrics, sig_color)


# ================================================================
# PART 3: INTEGRATION
# ================================================================

# SLIDE 16: COMBINED COVERAGE MATRIX (from ref deck page 28)
s = ns('Combined Coverage Matrix',
       'Each adversary capability owned by the right layer -- preemption first, behavioral as backstop',
       'PART 3: INTEGRATION')

mx_hdrs = ['Adversary Capability', 'Volt', 'Salt', 'Layer', 'Why This Layer Owns It']
mx_x = [Inches(0.3), Inches(4.3), Inches(5.0), Inches(5.7), Inches(7.0)]
mx_w = [Inches(3.9), Inches(0.6), Inches(0.6), Inches(1.2), Inches(6.0)]

_box(s, Inches(0.25), Inches(1.25), Inches(12.8), Inches(0.35), fill=NAVY)
for i, h in enumerate(mx_hdrs):
    _txt(s, mx_x[i], Inches(1.27), mx_w[i], Inches(0.3), h, 9, WHITE, True)

mx_rows = [
    ('Edge-device CVE exploitation', 'Y', 'Y', 'RIGOR', 'Configuration is the gap -- formal verification proves patch/inspection enforced', BLUE),
    ('Misconfigured / shadowed FW rules', 'Y', 'Y', 'RIGOR', 'Mathematical model finds shadows, conflicts, weak rules across multi-vendor', BLUE),
    ('Router ACL tampering / config drift', '', 'Y', 'RIGOR', 'Config Drift use case re-verifies live state vs approved baseline continuously', BLUE),
    ('C2 to known malicious infra', 'Y', 'Y', 'BOTH', 'Rigor blocks known IPs; ACECARD c2_beaconing catches new infra', PURPLE),
    ('Compromised valid accounts (T1078)', 'Y', 'Y', 'ACECARD', 'Traffic is policy-allowed; only behavioral drift reveals abuse', TEAL),
    ('Living-off-the-Land (T1059)', 'Y', 'Y', 'ACECARD', 'lotl_execution concept catches LoTL chain -- process signal shifts', TEAL),
    ('Cross-domain lateral movement', 'Y', 'Y', 'ACECARD', 'Composite Digital Signature joins identity across AD/AAD/AWS/Okta/K8s', TEAL),
    ('Slow lateral movement (months)', 'Y', 'Y', 'ACECARD', 'CUSUM 4-sigma surfaces gradual structural shifts over time', TEAL),
    ('Zero-day / unknown TTPs', 'Y', 'Y', 'ACECARD', 'Not in MITRE yet -- only behavioral anomaly flags deviation from normal', TEAL),
    ('Insider threat / data hoarding', '', '', 'ACECARD', 'insider_data_hoarding concept (T1074) + peer cohort z-score', TEAL),
]

y = Inches(1.65)
for i, (cap, volt, salt, layer, why, color) in enumerate(mx_rows):
    bg = CARD if i % 2 == 0 else ALT
    _box(s, Inches(0.25), y, Inches(12.8), Inches(0.5), fill=bg, border=BDR)
    _txt(s, mx_x[0], y+Pt(2), mx_w[0], Inches(0.45), cap, 10, TXT, True)
    _txt(s, mx_x[1], y+Pt(2), mx_w[1], Inches(0.45), volt, 10, RED if volt else SUB, al=PP_ALIGN.CENTER)
    _txt(s, mx_x[2], y+Pt(2), mx_w[2], Inches(0.45), salt, 10, RED if salt else SUB, al=PP_ALIGN.CENTER)
    _box(s, mx_x[3]+Inches(0.05), y+Inches(0.08), mx_w[3]-Inches(0.1), Inches(0.22), fill=color)
    _txt(s, mx_x[3], y+Inches(0.09), mx_w[3], Inches(0.2), layer, 9, WHITE, True, PP_ALIGN.CENTER)
    _txt(s, mx_x[4], y+Pt(2), mx_w[4], Inches(0.45), why, 9, TXT)
    y += Inches(0.53)

_box(s, Inches(0.3), Inches(6.95), Inches(12.7), Inches(0.35), fill=NAVY)
_txt(s, Inches(0.5), Inches(6.97), Inches(12.3), Inches(0.3),
     'Together: either the path never existed (Rigor) OR the attack is detected and contained within hours (ACECARD + ABAC)',
     11, GOLD, True, PP_ALIGN.CENTER)


# SLIDE 17: KEY NUMBERS
s = ns('Key Numbers to Memorize', 'Have these ready for Q&A', 'REFERENCE')

left = [
    ('5+ years', 'Volt Typhoon dwell time'),
    ('4+ years', 'Salt Typhoon dwell time'),
    ('200+', 'Telcos compromised by Salt Typhoon'),
    ('99%', 'FW breaches = misconfig (Gartner)'),
    ('~20%', 'Attacks stopped by KEV patching'),
    ('2^8000', 'Config state space (ASM samples)'),
    ('50M+', 'Rule combinations Rigor reasons over'),
    ('1,536', 'Embedding dims per signal'),
    ('9,216', 'Total dims (6 x 1536)'),
    ('5', 'Behavioral signals'),
    ('8', 'MITRE threat concepts'),
    ('4-sigma', 'CUSUM alarm threshold'),
]
right = [
    ('<3 sec', 'Event to analysis per entity'),
    ('10K/hr', 'Entities per 4-vCPU worker'),
    ('270 days', 'Synthetic data in prototype'),
    ('6', 'Attack scenarios detectable'),
    ('7', 'SOC dashboard charts'),
    ('30-45 min', 'Current analyst triage time'),
    ('<3 sec', 'Our triage time'),
    ('$90M', '22CT Army SOC/MDR contract'),
    ('800+', 'Cleared analysts'),
    ('20+', 'Rigor AI design partners'),
    ('90 days', 'Contract to first detection'),
    ('IL5/IL6/JWICS', 'Three enclave support'),
]

for nums, xs in [(left, Inches(0.3)), (right, Inches(6.8))]:
    y = Inches(1.3)
    for num, meaning in nums:
        _box(s, xs, y, Inches(6.2), Inches(0.42), fill=CARD, border=BDR)
        _box(s, xs+Inches(0.05), y+Inches(0.04), Inches(1.4), Inches(0.3), fill=NAVY)
        _txt(s, xs+Inches(0.1), y+Inches(0.06), Inches(1.3), Inches(0.26),
             num, 11, WHITE, True, PP_ALIGN.CENTER, fn='Cascadia Code')
        _txt(s, xs+Inches(1.6), y+Inches(0.06), Inches(4.4), Inches(0.28),
             meaning, 10, TXT)
        y += Inches(0.44)


# ================================================================
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
