"""Build SEPARATE Rigor AI Preemptive Cyber Defense deck.
Expands the Rigor_ACECARD_Volt_Salt_Typhoon reference material.
LIGHT background. 17 slides."""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

OUT = os.path.join(os.path.dirname(__file__), "07_Rigor_Preemptive_Cyber_Deck.pptx")

NAVY   = RGBColor(0x00, 0x2B, 0x5C)
GOLD   = RGBColor(0xB8, 0x86, 0x2A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BG     = RGBColor(0xF7, 0xF8, 0xFA)
CARD   = RGBColor(0xFF, 0xFF, 0xFF)
BLUE   = RGBColor(0x1E, 0x4D, 0x8C)
PURPLE = RGBColor(0x6C, 0x2E, 0xD9)
TEAL   = RGBColor(0x0D, 0x8B, 0x7C)
ORANGE = RGBColor(0xD9, 0x7A, 0x06)
RED    = RGBColor(0xCC, 0x33, 0x33)
CYBER  = RGBColor(0x00, 0x99, 0x7A)
TXT    = RGBColor(0x22, 0x22, 0x33)
SUB    = RGBColor(0x55, 0x5E, 0x70)
BDR    = RGBColor(0xDD, 0xDD, 0xE5)
ALT    = RGBColor(0xF0, 0xF2, 0xF5)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width


def _fill(s, c):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = c

def _box(s, l, t, w, h, fill=None, border=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.background()
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    ln = sh.line
    if border:
        ln.color.rgb = border
        ln.width = Pt(1)
    else:
        ln.fill.background()
    sh.shadow.inherit = False
    return sh

def _txt(s, l, t, w, h, text, sz=14, c=TXT, b=False, al=PP_ALIGN.LEFT, fn='Segoe UI'):
    tf = s.shapes.add_textbox(l, t, w, h).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = c
    p.font.bold = b
    p.font.name = fn
    p.alignment = al
    return tf

def _p(tf, text, sz=14, c=TXT, b=False, sp=Pt(4)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = c
    p.font.bold = b
    p.font.name = 'Segoe UI'
    p.space_before = sp
    return p

def _title(s, title, sub=None, part=None):
    _box(s, Inches(0), Inches(0), W, Inches(1.1), fill=NAVY)
    _txt(s, Inches(0.6), Inches(0.15), Inches(10), Inches(0.5), title, 26, WHITE, True)
    if sub:
        _txt(s, Inches(0.6), Inches(0.6), Inches(10), Inches(0.4), sub, 13, RGBColor(0xCC, 0xDD, 0xEE))
    if part:
        _txt(s, Inches(10.2), Inches(0.25), Inches(2.8), Inches(0.5), part, 12, GOLD, True, PP_ALIGN.RIGHT)

def ns(title, sub=None, part='RIGOR AI'):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _fill(s, BG)
    _title(s, title, sub, part)
    return s


# ================================================================
# SLIDE 1: TITLE
# ================================================================
s = ns('Rigor AI: Preemptive Cyber Defense',
       'Mathematically Rigorous, Complete, Continuous, Verifiable')
_txt(s, Inches(1), Inches(2.0), Inches(11), Inches(0.5),
     'Detailed Technical Reference -- Expanding the ACECARD Preemptive Layer', 20, NAVY, True, PP_ALIGN.CENTER)

_box(s, Inches(0.5), Inches(3.0), Inches(5.8), Inches(3.8), fill=CARD, border=BDR)
_txt(s, Inches(0.7), Inches(3.05), Inches(5.4), Inches(0.3), 'THIS DECK COVERS:', 13, NAVY, True)
tf = _txt(s, Inches(0.7), Inches(3.4), Inches(5.4), Inches(3.3),
     '1. The Rigor AI Approach (3 Intelligences)', 12, TXT)
_p(tf, '2. How RigorOS Operates (Inputs -> Engine -> Assurance)', 12, TXT)
_p(tf, '3. What Rigor Mathematically Proves', 12, TXT)
_p(tf, '4. Volt Typhoon: Phase-by-Phase Defense', 12, TXT)
_p(tf, '5. Salt Typhoon: Phase-by-Phase Defense', 12, TXT)
_p(tf, '6. Volt Typhoon TTP Coverage Matrix', 12, TXT)
_p(tf, '7. Salt Typhoon TTP Coverage Matrix', 12, TXT)
_p(tf, '8. The Formal Model: 2^8000 State Space', 12, TXT)
_p(tf, '9. Nine Use Cases', 12, TXT)

_box(s, Inches(6.8), Inches(3.0), Inches(5.8), Inches(3.8), fill=CARD, border=BDR)
tf = _txt(s, Inches(7.0), Inches(3.05), Inches(5.4), Inches(3.3),
     '10. Real Customer Findings (6 Proof Points)', 12, TXT)
_p(tf, '11. Where Rigor Needs Help (-> ACECARD)', 12, TXT)
_p(tf, '12. Four Integration Points (Rigor + ACECARD)', 12, TXT)
_p(tf, '13. Combined Detection Timeline', 12, TXT)
_p(tf, '14. Deployment Models (SaaS/On-Prem/Air-Gap)', 12, TXT)
_p(tf, '15. Combined Coverage Matrix', 12, TXT)
_p(tf, '16. Key Numbers to Memorize', 12, TXT)
_p(tf, '', 12, TXT)
_p(tf, 'Based on: Rigor_ACECARD_Volt_Salt_Typhoon.pdf', 11, SUB)
_p(tf, 'Expanded with detailed technical content', 11, SUB)

_box(s, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.3), fill=NAVY)
_txt(s, Inches(0.7), Inches(7.02), Inches(11.9), Inches(0.26),
     'CFIC Augusta GA | May 6, 2026 | ACECARD CE Event | Solution-Focused Technical Reference',
     11, GOLD, True, PP_ALIGN.CENTER)


# ================================================================
# SLIDE 2: THE RIGOR AI APPROACH
# ================================================================
s = ns('The Rigor AI Approach: Close the Door Before They Knock',
       'Three intelligences. Mathematical certainty. No sampling. No false positives.')

problems = [
    ('Analyze All Known Threat Intelligence', 'Terabytes of CVE, advisories, CTI, IOCs, IOAs, PSIRTs -- ingested continuously',
     'Attack Intelligence', 'MITRE-enriched attack graphs built by ML/NLP/LLM pipeline. Full TTP coverage of Volt Typhoon (G1017) and Salt Typhoon. Auto-updates on new advisories.', BLUE),
    ('Identify Every Defensive Control Gap', 'Exponential combination of NGFW, IdP, IPS, SASE, WAF configs -- 2^8000 state space',
     'Defense Intelligence', 'Symbolic Model of Computation. Formal mathematical model exhaustively reasons over the 2^8000 state space. No sampling. Complete reasoning. Every path analyzed.', PURPLE),
    ('Prescribe Correct Configurations', 'Without creating new errors, breaking connectivity, or business disruption',
     'Remediation Intelligence', 'Grounded agentic-AI reasoners. Risk-prioritized, vendor-specific fixes. No false positives (mathematically proven). No false negatives (exhaustive).', TEAL),
]
y = Inches(1.3)
for prob, probdesc, sol, soldesc, color in problems:
    _box(s, Inches(0.3), y, Inches(6.2), Inches(1.75), fill=CARD, border=BDR)
    _box(s, Inches(0.35), y+Inches(0.05), Inches(1.2), Inches(0.22), fill=RED)
    _txt(s, Inches(0.4), y+Inches(0.06), Inches(1.1), Inches(0.2), 'PROBLEM', 9, WHITE, True, PP_ALIGN.CENTER)
    _txt(s, Inches(1.7), y+Inches(0.05), Inches(4.6), Inches(0.25), prob, 12, TXT, True)
    _txt(s, Inches(1.7), y+Inches(0.35), Inches(4.6), Inches(1.2), probdesc, 10, SUB)

    _box(s, Inches(6.8), y, Inches(6.2), Inches(1.75), fill=CARD, border=color)
    _box(s, Inches(6.85), y+Inches(0.05), Inches(1.2), Inches(0.22), fill=color)
    _txt(s, Inches(6.9), y+Inches(0.06), Inches(1.1), Inches(0.2), 'SOLUTION', 9, WHITE, True, PP_ALIGN.CENTER)
    _txt(s, Inches(8.2), y+Inches(0.05), Inches(4.6), Inches(0.25), sol, 12, color, True)
    _txt(s, Inches(8.2), y+Inches(0.35), Inches(4.6), Inches(1.2), soldesc, 10, TXT)
    y += Inches(1.85)

_box(s, Inches(0.3), Inches(6.9), Inches(12.7), Inches(0.4), fill=NAVY)
_txt(s, Inches(0.5), Inches(6.92), Inches(12.3), Inches(0.35),
     "The world's first Mathematically Rigorous -- Complete, Preemptive, Continuous, Verifiable -- Cyber Defense Platform.",
     12, GOLD, True, PP_ALIGN.CENTER)


# ================================================================
# SLIDE 3: HOW RIGOROS OPERATES
# ================================================================
s = ns('How RigorOS Operates: Inputs -> Reasoning -> Continuous Remediation',
       'SaaS, on-prem, and air-gapped deployments. Integrates with ServiceNow, JIRA, and SOC ticketing.')

_box(s, Inches(0.3), Inches(1.3), Inches(3.5), Inches(2.8), fill=CARD, border=BLUE)
_box(s, Inches(0.3), Inches(1.3), Inches(3.5), Inches(0.35), fill=BLUE)
_txt(s, Inches(0.5), Inches(1.32), Inches(3.1), Inches(0.3), 'INPUTS', 14, WHITE, True)
tf = _txt(s, Inches(0.5), Inches(1.75), Inches(3.1), Inches(2.2),
     'Threat Intelligence', 11, TXT, True)
_p(tf, 'CVE, advisories, CTI, IOCs, IOAs, PSIRTs', 10, SUB)
_p(tf, 'CISA KEV, NSA advisories, vendor PSIRTs', 10, SUB)
_p(tf, '', 8, SUB)
_p(tf, 'Tech Controls Configuration', 11, TXT, True)
_p(tf, 'NGFW, IdP, IAM, IPS, EPP, SASE, WAF', 10, SUB)
_p(tf, 'Full running-config of every device', 10, SUB)
_p(tf, '', 8, SUB)
_p(tf, 'Grounded World Models', 11, TXT, True)
_p(tf, 'Vendor capability models (PAN, FTN, Cisco, CP)', 10, SUB)
_p(tf, 'MITRE ATT&CK / D3FEND frameworks', 10, SUB)

_box(s, Inches(4.1), Inches(1.3), Inches(4.8), Inches(2.8), fill=CARD, border=NAVY)
_box(s, Inches(4.1), Inches(1.3), Inches(4.8), Inches(0.35), fill=NAVY)
_txt(s, Inches(4.3), Inches(1.32), Inches(4.4), Inches(0.3), 'RIGOR ENGINE', 14, WHITE, True)
steps = [
    ('1. Model Attacks from Threat Intel', 'Build attack graphs from CVE/advisory data; map to MITRE techniques; identify which paths target THIS customer'),
    ('2. Model Customer Cyber Defenses', 'Symbolic model of every NGFW rule, ACL, IPS profile, identity policy; 2^8000 state combinations'),
    ('3. Reason About Control Gaps', 'Formal proof: does any attack path succeed against current config? If yes -> exact gap + vendor fix'),
    ('4. Generate Risk-Prioritized Remediations', 'Vendor-specific playbooks; risk-ranked; no false positives; alert to SOC + Control Admin'),
    ('5. Apply + Re-verify Continuously', 'Optional auto-apply; validate fix; re-verify every hour and on every config change'),
]
y2 = Inches(1.75)
for title, desc in steps:
    _txt(s, Inches(4.3), y2, Inches(4.4), Inches(0.2), title, 10, NAVY, True)
    _txt(s, Inches(4.3), y2+Inches(0.2), Inches(4.4), Inches(0.3), desc, 8, SUB)
    y2 += Inches(0.48)

_box(s, Inches(9.2), Inches(1.3), Inches(3.8), Inches(2.8), fill=CARD, border=TEAL)
_box(s, Inches(9.2), Inches(1.3), Inches(3.8), Inches(0.35), fill=TEAL)
_txt(s, Inches(9.4), Inches(1.32), Inches(3.4), Inches(0.3), 'CONTINUOUS ASSURANCE', 14, WHITE, True)
tf = _txt(s, Inches(9.4), Inches(1.75), Inches(3.4), Inches(2.2),
     'Risk-prioritized control gap findings', 10, TXT)
_p(tf, 'Vendor-specific remediation playbooks', 10, TXT)
_p(tf, 'Mathematical proof: no false positives', 10, NAVY, True)
_p(tf, 'Mathematical proof: no false negatives', 10, NAVY, True)
_p(tf, 'SOC and control admin workflows', 10, TXT)
_p(tf, 'ServiceNow / JIRA integration', 10, TXT)
_p(tf, 'Auto re-verify on every config change', 10, TXT)
_p(tf, 'Auto re-verify on every threat-feed update', 10, TXT)
_p(tf, 'Hourly continuous re-verification', 10, TXT)

_box(s, Inches(0.3), Inches(4.3), Inches(12.7), Inches(0.45), fill=CARD, border=BDR)
_txt(s, Inches(0.5), Inches(4.33), Inches(12.3), Inches(0.38),
     'Currently supports: Palo Alto Networks, Fortinet, Cisco, Check Point (N-S and E-W NGFWs). Identity, SASE, and additional vendors on roadmap.',
     11, SUB)

_box(s, Inches(0.3), Inches(4.95), Inches(12.7), Inches(2.3), fill=CARD, border=NAVY)
_box(s, Inches(0.3), Inches(4.95), Inches(12.7), Inches(0.35), fill=NAVY)
_txt(s, Inches(0.5), Inches(4.97), Inches(12.3), Inches(0.3), 'WHAT RIGOR PROVES (MATHEMATICAL CERTAINTY)', 14, WHITE, True)

proofs = [
    ('Edge-device exploit paths are blocked', 'CVE-2023-20198, CVE-2024-21887, CVE-2024-3400 etc. verified closed'),
    ('Threat-feed C2 IPs/domains denied', 'Across 50M+ rule combinations. Known KV botnet, JumbledPath C2 infrastructure'),
    ('Config drift from baseline is caught', 'Every hour or per config change. Salt Typhoon ACL tampering = caught within 1 hour'),
    ('Critical segmentation holds end-to-end', 'User VLAN cannot reach DC mgmt. CALEA zone isolated. Zero-trust verified mathematically.'),
    ('Shadowed/zombie rules surfaced', 'Real finding: 3 higher-priority rules shadowing SMB app-control profile'),
    ('Compensating controls for unpatchable CVEs', 'When no patch exists, identifies alternate control that closes the attack path'),
]
y2 = Inches(5.4)
for title, desc in proofs:
    _txt(s, Inches(0.5), y2, Inches(4.5), Inches(0.22), title, 11, NAVY, True)
    _txt(s, Inches(5.1), y2, Inches(7.7), Inches(0.22), desc, 10, SUB)
    y2 += Inches(0.32)


# ================================================================
# SLIDE 4: VOLT TYPHOON PHASE BY PHASE
# ================================================================
s = ns('How Rigor AI Stops Volt Typhoon -- Phase by Phase',
       'Volt Typhoon needs configuration weakness to enter and pivot. Rigor proves those weaknesses do not exist.',
       'VOLT TYPHOON')

cols_x = [Inches(0.3), Inches(2.2), Inches(5.8), Inches(9.6)]
cols_w = [Inches(1.8), Inches(3.5), Inches(3.7), Inches(3.4)]
hdrs = ['PHASE', 'VOLT TYPHOON BEHAVIOR', 'HOW RIGOR AI STOPS IT', 'CVEs / TTPs']

_box(s, Inches(0.25), Inches(1.25), Inches(12.8), Inches(0.38), fill=NAVY)
for i, h in enumerate(hdrs):
    _txt(s, cols_x[i], Inches(1.27), cols_w[i], Inches(0.34), h, 10, WHITE, True)

rows_vt = [
    ('Phase 1\nInitial Access\n0-12h',
     'Scans for unpatched Fortinet, Ivanti, Versa, PAN edge devices. Exploits CVE to land on perimeter.',
     'Attack Intel ingests CISA/NSA/vendor PSIRT advisories. Defense Intel proves every NGFW/VPN/SASE has IPS signature or threat-feed deny rule enforced. If no patch, identifies compensating control.',
     'CVE-2023-20198\nCVE-2024-21887\nCVE-2024-3400\nT1190'),
    ('Phase 2\nPersistence\n12-48h',
     'Drops .aspx/.jspx web shells. Routes C2 through KV botnet of compromised SOHO routers.',
     'Models every threat-feed IP, domain, ASN, SOHO fingerprint as denied destination. Proves 50M+ rule combinations block these flows. WAF verified for web shell upload inspection.',
     'KV Botnet\nT1505 (Web Shell)\nT1071 (App Layer C2)'),
    ('Phase 3\nCredential Access\n48-96h',
     'ntdsutil.exe to dump NTDS.dit. LSASS memory extraction. Exfiltrates AD credentials to C2.',
     'Verifies EDR/EPP coverage enabled on every DC. Script-control profile blocks ntdsutil by non-admin. Egress filtering blocks credential exfil to known C2 ranges. Missing = exact vendor fix.',
     'T1003 (Cred Dump)\nntdsutil, LSASS\nT1059 (Scripting)'),
    ('Phase 4\nLateral Movement\n96h+',
     'SMB/RDP/WinRM to domain controllers and file servers. Admin share access.',
     'Real finding: 3 higher-priority rules shadowed SMB IPS profile -- Rigor closed lateral path. Micro-segmentation verified: user VLAN cannot reach DC management ports.',
     'T1021.002 (SMB)\nT1021.001 (RDP)\nT1078 (Valid Accts)'),
    ('Phase 5\nPre-positioning\nWeeks-Years',
     'Credential harvest + network mapping for future sabotage during geopolitical crisis.',
     'Continuous re-verification: every config change triggers re-analysis. Zero-trust segmentation re-proven hourly. IdP/IAM coverage modeled; trust paths constrained.',
     'T1133 (Ext Remote)\nT1556 (Modify Auth)\n5+ year dwell'),
]

y = Inches(1.7)
for i, (phase, attack, rigor, ttps) in enumerate(rows_vt):
    bg = CARD if i % 2 == 0 else ALT
    _box(s, Inches(0.25), y, Inches(12.8), Inches(1.05), fill=bg, border=BDR)
    _txt(s, cols_x[0], y+Pt(2), cols_w[0], Inches(1.0), phase, 9, NAVY, True)
    _txt(s, cols_x[1], y+Pt(2), cols_w[1], Inches(1.0), attack, 9, TXT)
    _txt(s, cols_x[2], y+Pt(2), cols_w[2], Inches(1.0), rigor, 9, BLUE)
    _txt(s, cols_x[3], y+Pt(2), cols_w[3], Inches(1.0), ttps, 8, RED, fn='Cascadia Code')
    y += Inches(1.08)

_box(s, Inches(0.3), Inches(7.1), Inches(12.7), Inches(0.25), fill=NAVY)
_txt(s, Inches(0.5), Inches(7.12), Inches(12.3), Inches(0.2),
     'Volt Typhoon (G1017): PRC state-sponsored. 5+ year dwell. LoTL tradecraft. Critical infrastructure targeting.',
     10, GOLD, True, PP_ALIGN.CENTER)


# ================================================================
# SLIDE 5: SALT TYPHOON PHASE BY PHASE
# ================================================================
s = ns('How Rigor AI Stops Salt Typhoon -- Carrier-Scale Edge Defense',
       'Salt Typhoon lives on misconfigured Cisco/Ivanti/PAN edge gear. Rigor verifies every config every hour.',
       'SALT TYPHOON')

_box(s, Inches(0.25), Inches(1.25), Inches(12.8), Inches(0.38), fill=NAVY)
for i, h in enumerate(hdrs):
    _txt(s, cols_x[i], Inches(1.27), cols_w[i], Inches(0.34), h, 10, WHITE, True)

rows_st = [
    ('Phase 1\nEdge CVE Exploit\n0-24h',
     'Exploits Cisco IOS XE web UI (CVE-2023-20198) for level-15 admin. Also Ivanti, PAN, Citrix. 1000+ Cisco devices targeted Dec 2024-Jan 2025.',
     'CVE-to-control mapping is automatic. Every device exposing IOS XE web UI to internet = flagged. Proves WAF/mgmt-plane ACL blocks /webui/. For unpatched CVE, identifies compensating control.',
     'CVE-2023-20198\nCVE-2018-0171\nCVE-2024-3400\nT1190'),
    ('Phase 2\nPersistence\n1-7d',
     'Creates local privileged accounts. Modifies router ACLs to whitelist attacker IPs. Enables Guest Shell, sshd_operns. 3+ years access in some victims.',
     'Config Drift Management: takes approved baseline of every router ACL, AAA config, mgmt-plane policy. Compares live state every hour. Any unauthorized change = immediate gap finding with exact diff.',
     'T1556 (Modify Auth)\nLocal account creation\nACL tampering\n3+ year dwell'),
    ('Phase 3\nTACACS+ Capture\n7-30d',
     'Captures TACACS+ traffic on TCP/49 via on-box packet capture. Harvests credentials. Uses JumbledPath to chain remote connections.',
     'Mgmt-plane segmentation verified. On-box capture restricted to authorized roles. Outbound rules to JumbledPath C2 IPs blocked across all internal zones.',
     'JumbledPath malware\nTACACS+ TCP/49\nT1040 (Sniffing)'),
    ('Phase 4\nGRE Tunnel + C2\n30-90d',
     'GRE tunnel configured. Encrypted C2. Multi-hop connections across compromised routers.',
     'GRE peer authorization verified. Unauthorized tunnel endpoints detected in routing config. Egress policy verified: only approved protocols/ports permitted outbound.',
     'GRE tunneling\nT1572 (Protocol Tunnel)\nT1071 (App Layer C2)'),
    ('Phase 5\nCALEA + CDR Exfil\n90d+',
     'Accesses CALEA lawful intercept systems. Steady CDR/metadata exfiltration over months. Targets: government officials, political figures.',
     'CALEA-zone segmentation intent proven end-to-end. Privilege verification on interception systems. Outbound restriction policies verified for CDR stores. DLP coverage modeled.',
     'CALEA access\nCDR exfiltration\nT1041 (Exfil over C2)\n200+ telcos'),
]

y = Inches(1.7)
for i, (phase, attack, rigor, ttps) in enumerate(rows_st):
    bg = CARD if i % 2 == 0 else ALT
    _box(s, Inches(0.25), y, Inches(12.8), Inches(1.05), fill=bg, border=BDR)
    _txt(s, cols_x[0], y+Pt(2), cols_w[0], Inches(1.0), phase, 9, NAVY, True)
    _txt(s, cols_x[1], y+Pt(2), cols_w[1], Inches(1.0), attack, 9, TXT)
    _txt(s, cols_x[2], y+Pt(2), cols_w[2], Inches(1.0), rigor, 9, BLUE)
    _txt(s, cols_x[3], y+Pt(2), cols_w[3], Inches(1.0), ttps, 8, RED, fn='Cascadia Code')
    y += Inches(1.08)

_box(s, Inches(0.3), Inches(7.1), Inches(12.7), Inches(0.25), fill=NAVY)
_txt(s, Inches(0.5), Inches(7.12), Inches(12.3), Inches(0.2),
     'Salt Typhoon: PRC state-sponsored. 4+ years. 200+ telcos in 80+ countries. CVE-2023-20198. CALEA wiretap access.',
     10, GOLD, True, PP_ALIGN.CENTER)


# ================================================================
# SLIDE 6: VOLT TYPHOON TTP COVERAGE
# ================================================================
s = ns('Rigor AI Preempts Volt Typhoon -- TTP by TTP',
       'Each MITRE technique -> formal control gap query -> mathematically verified blocking configuration',
       'VOLT TYPHOON')

ttp_hdrs = ['MITRE Technique', 'Volt Typhoon Behavior', 'Rigor AI Preemptive Action']
ttp_x = [Inches(0.3), Inches(2.8), Inches(7.0)]
ttp_w = [Inches(2.4), Inches(4.1), Inches(6.0)]

_box(s, Inches(0.25), Inches(1.25), Inches(12.8), Inches(0.38), fill=NAVY)
for i, h in enumerate(ttp_hdrs):
    _txt(s, ttp_x[i], Inches(1.27), ttp_w[i], Inches(0.34), h, 10, WHITE, True)

ttp_vt = [
    ('T1190\nExploit Public App', 'Exploits unpatched Fortinet, Versa Director, Ivanti, PAN edge devices',
     'Formally verifies every NGFW/SASE/WAF rule blocks known CVE traffic; flags weak/missing inspection profiles'),
    ('T1133\nExternal Remote Svc', 'Abuses VPN, RDP gateways with stolen valid accounts',
     'Proves zero-trust segmentation intent; surfaces shadowed allow-rules that bypass MFA enforcement'),
    ('T1021.002\nSMB / Admin Shares', 'Lateral movement via SMBv1, admin shares between branches',
     'Real customer case: 3 higher-priority rules shadowed SMB IPS profile -- Rigor closed the lateral path'),
    ('T1003\nOS Credential Dump', 'ntdsutil, LSASS dumps; credentials exfiltrated to C2',
     'Verifies egress filtering blocks C2 destinations; ensures EDR/IPS profiles cover credential-theft signatures'),
    ('T1071\nC2 over SOHO Routers', 'KV botnet relays C2 through compromised home/SMB routers',
     'Models threat-feed IPs/domains; proves outbound rules block them across all internal zones'),
    ('T1059\nPowerShell / LOLBins', 'In-memory exec via wmic, ntdsutil, netsh; log evasion',
     'Verifies EDR/EPP coverage in policy; identifies hosts where script-control profile is not applied'),
    ('T1078\nValid Accounts', 'Pivots laterally with stolen credentials across trust zones',
     'IdP/IAM coverage modeled; inter-zone trust paths surfaced and constrained -> handed to ACECARD'),
]

y = Inches(1.7)
for i, (tech, behavior, action) in enumerate(ttp_vt):
    bg = CARD if i % 2 == 0 else ALT
    _box(s, Inches(0.25), y, Inches(12.8), Inches(0.72), fill=bg, border=BDR)
    _txt(s, ttp_x[0], y+Pt(2), ttp_w[0], Inches(0.68), tech, 10, RED, True, fn='Cascadia Code')
    _txt(s, ttp_x[1], y+Pt(2), ttp_w[1], Inches(0.68), behavior, 9, TXT)
    _txt(s, ttp_x[2], y+Pt(2), ttp_w[2], Inches(0.68), action, 9, BLUE)
    y += Inches(0.75)

_txt(s, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.3),
     'Source: MITRE ATT&CK Group G1017 (Volt Typhoon); CISA AA24-038A; Microsoft Threat Intelligence', 9, SUB)


# ================================================================
# SLIDE 7: SALT TYPHOON TTP COVERAGE
# ================================================================
s = ns('Rigor AI Preempts Salt Typhoon -- TTP by TTP',
       'Carrier/telco-focused techniques -> verified edge device and segmentation controls',
       'SALT TYPHOON')

_box(s, Inches(0.25), Inches(1.25), Inches(12.8), Inches(0.38), fill=NAVY)
for i, h in enumerate(ttp_hdrs):
    _txt(s, ttp_x[i], Inches(1.27), ttp_w[i], Inches(0.34), h, 10, WHITE, True)

ttp_st = [
    ('T1190\nExploit Cisco IOS XE', 'CVE-2023-20198 web UI exploit for level-15 admin; 1000+ devices Dec 2024',
     'Proves mgmt-plane ACL blocks /webui/ from external; IPS signatures for Cisco IOS exploitation enforced'),
    ('T1556\nModify Authentication', 'Creates local admin accounts; modifies AAA config; enables Guest Shell',
     'Config drift detection: hourly comparison to approved AAA baseline; any new local account = immediate alert'),
    ('T1040\nNetwork Sniffing', 'TACACS+ capture on TCP/49; harvests network admin credentials',
     'Mgmt-plane segmentation verified; on-box capture permissions modeled; TACACS+ encryption enforcement'),
    ('T1572\nProtocol Tunneling', 'GRE tunnel configured between compromised routers for C2 and data exfil',
     'GRE peer authorization verified in routing config; unauthorized tunnel endpoints flagged; egress policy proven'),
    ('T1041\nExfil Over C2', 'CDR/metadata exfiltration over encrypted channels to PRC infrastructure',
     'CALEA-zone segmentation proven end-to-end; DLP policies verified on CDR stores; outbound restrictions'),
    ('T1078\nValid Accounts', 'TACACS+ harvested credentials used as legitimate network admin sessions',
     'Trust path constrained; IdP coverage modeled -> passed to ACECARD for behavioral detection of credential abuse'),
    ('T1021\nRemote Services', 'SSH/Telnet lateral movement across carrier backbone using stolen credentials',
     'Inter-zone access policies verified; jump-box enforcement proven; direct access paths blocked'),
]

y = Inches(1.7)
for i, (tech, behavior, action) in enumerate(ttp_st):
    bg = CARD if i % 2 == 0 else ALT
    _box(s, Inches(0.25), y, Inches(12.8), Inches(0.72), fill=bg, border=BDR)
    _txt(s, ttp_x[0], y+Pt(2), ttp_w[0], Inches(0.68), tech, 10, RED, True, fn='Cascadia Code')
    _txt(s, ttp_x[1], y+Pt(2), ttp_w[1], Inches(0.68), behavior, 9, TXT)
    _txt(s, ttp_x[2], y+Pt(2), ttp_w[2], Inches(0.68), action, 9, BLUE)
    y += Inches(0.75)

_txt(s, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.3),
     'Source: Trend Micro, CISA, Cisco PSIRT; GhostSpider/Demodex/SnappyBee implant families; JumbledPath C2', 9, SUB)


# ================================================================
# SLIDE 8: THE FORMAL MODEL -- 2^8000 STATE SPACE
# ================================================================
s = ns('The Formal Model: Why 2^8000 Cannot Be Sampled',
       'Traditional scanners sample configurations. Rigor reasons EXHAUSTIVELY over the entire state space.',
       'FORMAL METHOD')

_box(s, Inches(0.3), Inches(1.3), Inches(6.2), Inches(2.8), fill=CARD, border=NAVY)
_txt(s, Inches(0.5), Inches(1.35), Inches(5.8), Inches(0.3), 'THE PROBLEM WITH SAMPLING', 13, NAVY, True)
tf = _txt(s, Inches(0.5), Inches(1.7), Inches(5.8), Inches(2.3),
     'A typical enterprise NGFW has:', 11, TXT)
_p(tf, '  5,000-50,000 firewall rules', 11, TXT)
_p(tf, '  Each rule: src zone, dst zone, src IP, dst IP, port, app,', 10, SUB)
_p(tf, '    user/group, schedule, action, inspection profile', 10, SUB)
_p(tf, '  Interactions between rules: 50,000^2 = 2.5 billion pairs', 11, TXT)
_p(tf, '  Multi-vendor: 3-10 firewalls in series/parallel', 11, TXT)
_p(tf, '', 8, SUB)
_p(tf, 'Total state space: ~2^8000 possible traffic paths', 12, NAVY, True)
_p(tf, '', 8, SUB)
_p(tf, 'Sampling 1 billion paths/second for 1 year =', 10, SUB)
_p(tf, '  3.15 x 10^16 samples (vs 2^8000 = 10^2408)', 10, RED, True)
_p(tf, '  Coverage: effectively 0%', 10, RED, True)

_box(s, Inches(6.8), Inches(1.3), Inches(6.2), Inches(2.8), fill=CARD, border=TEAL)
_txt(s, Inches(7.0), Inches(1.35), Inches(5.8), Inches(0.3), 'RIGOR: SYMBOLIC MODEL OF COMPUTATION', 13, TEAL, True)
tf = _txt(s, Inches(7.0), Inches(1.7), Inches(5.8), Inches(2.3),
     'Instead of sampling paths, Rigor builds a', 11, TXT)
_p(tf, 'mathematical MODEL of the entire config:', 11, TXT)
_p(tf, '', 8, SUB)
_p(tf, '1. Parse all rules into symbolic representation', 10, TXT)
_p(tf, '2. Build constraint-satisfaction model', 10, TXT)
_p(tf, '3. For each threat (CVE, C2 dest, lateral path):', 10, TXT)
_p(tf, '   "Does ANY rule combination ALLOW this traffic?"', 11, NAVY, True)
_p(tf, '4. Solver PROVES answer is NO (safe) or YES (gap)', 10, TXT)
_p(tf, '', 8, SUB)
_p(tf, 'Result: COMPLETE coverage of 2^8000 in seconds', 12, TEAL, True)
_p(tf, 'No sampling. No probabilistic reasoning.', 12, TEAL, True)
_p(tf, 'Mathematical certainty.', 12, TEAL, True)

_box(s, Inches(0.3), Inches(4.3), Inches(12.7), Inches(2.8), fill=CARD, border=GOLD)
_txt(s, Inches(0.5), Inches(4.35), Inches(12.3), Inches(0.3), 'WHAT THIS MEANS IN PRACTICE', 13, GOLD, True)

examples = [
    ('Traditional Pentest', '1-2 weeks, finds 5-15 paths, human effort', 'Point-in-time, incomplete, expensive', RED),
    ('Vulnerability Scanner', 'Samples known CVEs against config', 'Misses rule interactions, shadows, compensating controls', RED),
    ('Attack Surface Mgmt (ASM)', 'Enumerates external-facing assets', 'Sees what exists, not what is REACHABLE through config', RED),
    ('Rigor Formal Verification', 'Models ALL paths through ALL rules', 'Complete, continuous, no gaps, vendor-specific fixes', TEAL),
]
y2 = Inches(4.75)
for tool, method, limitation, color in examples:
    _box(s, Inches(0.5), y2, Inches(3.2), Inches(0.35), fill=color)
    _txt(s, Inches(0.55), y2+Inches(0.03), Inches(3.1), Inches(0.3), tool, 10, WHITE, True, PP_ALIGN.CENTER)
    _txt(s, Inches(3.85), y2+Inches(0.03), Inches(4.2), Inches(0.3), method, 10, TXT)
    _txt(s, Inches(8.2), y2+Inches(0.03), Inches(4.6), Inches(0.3), limitation, 10, color)
    y2 += Inches(0.45)

_box(s, Inches(0.3), Inches(7.05), Inches(12.7), Inches(0.3), fill=NAVY)
_txt(s, Inches(0.5), Inches(7.07), Inches(12.3), Inches(0.26),
     'Gartner 2024: "99% of firewall breaches are caused by misconfiguration, not firewall flaws." Rigor eliminates that 99%.',
     11, GOLD, True, PP_ALIGN.CENTER)


# ================================================================
# SLIDE 9: NINE USE CASES
# ================================================================
s = ns('Nine Rigor AI Use Cases',
       'Pre-built reasoning templates. Five directly preempt Volt + Salt Typhoon TTPs.',
       'USE CASES')

ucs = [
    ('01', 'Preemptive Cyber Defense', 'Continuous threat-intel + config -> proof that every TTP path is blocked',
     'Volt T1190, T1133, T1071 / Salt T1190, T1556', BLUE),
    ('02', 'Configuration Drift Mgmt', 'Live config vs approved baseline every hour. Surfaces unauthorized changes.',
     'Salt Typhoon ACL tampering, Guest Shell, local accounts', BLUE),
    ('03', 'Formal Policy Analysis', 'Reasons over 2^8000 firewall rule state space. Shadow rules, conflicts, weak inspection.',
     'SMB lateral paths, shadowed inspection profiles', BLUE),
    ('04', 'Preemptive ASM (External)', 'Scans external attack surface. Maps exposed assets to compensating controls.',
     'Edge-device CVE exploitation (CVE-2023-20198, etc.)', BLUE),
    ('05', 'Internal Attack Surface Mgmt', 'Models internal paths from any compromised host to crown-jewel assets.',
     'Cross-tenant pivots (Salt Typhoon telco -> enterprise)', BLUE),
    ('06', 'NIST 800-53 / Zero-Trust', 'Proves zero-trust segmentation holds end-to-end. Maps to NIST CSF / D3FEND.',
     'Continuous SP800-207 verification', TEAL),
    ('07', 'Vendor Migration Safety', 'Models before/after config during NGFW vendor migration. Proves no regression.',
     'PAN->Fortinet migrations, Cisco->Palo refresh cycles', TEAL),
    ('08', 'M&A Due Diligence', 'Models combined defense posture pre-merger. Surfaces gaps from acquired entity.',
     'Risk assessment before network integration', TEAL),
    ('09', 'Third-Party Risk', 'Models third-party access paths and verifies isolation controls.',
     'Supply chain and vendor risk management', TEAL),
]

y = Inches(1.3)
for num, title, desc, stops, color in ucs:
    _box(s, Inches(0.3), y, Inches(12.7), Inches(0.62), fill=CARD, border=BDR)
    _box(s, Inches(0.35), y+Inches(0.05), Inches(0.5), Inches(0.28), fill=color)
    _txt(s, Inches(0.4), y+Inches(0.07), Inches(0.4), Inches(0.24), num, 11, WHITE, True, PP_ALIGN.CENTER)
    _txt(s, Inches(1.0), y+Inches(0.05), Inches(2.5), Inches(0.24), title, 12, NAVY, True)
    _txt(s, Inches(3.6), y+Inches(0.05), Inches(5.0), Inches(0.5), desc, 10, TXT)
    _txt(s, Inches(8.8), y+Inches(0.05), Inches(4.0), Inches(0.5), stops, 9, RED)
    y += Inches(0.65)

_box(s, Inches(0.3), Inches(7.1), Inches(12.7), Inches(0.25), fill=NAVY)
_txt(s, Inches(0.5), Inches(7.12), Inches(12.3), Inches(0.2),
     'Use cases 01-05 (blue) directly preempt Volt + Salt Typhoon TTPs. 06-09 extend to compliance and risk.', 10, GOLD, True, PP_ALIGN.CENTER)


# ================================================================
# SLIDE 10: REAL CUSTOMER FINDINGS
# ================================================================
s = ns('Rigor in Production: Real Customer Findings',
       'Selected findings from financial services, telecom, government, and critical infrastructure',
       'PROOF POINTS')

findings = [
    ('FinTech | Asia', 'Closed SMB Lateral Movement Path', BLUE,
     'SMB app-control profiles completely shadowed by 3 higher-priority rules -- allowed unprotected SMBv1 between branches. Rigor formal model exposed the multi-rule shadow that no scanner found.',
     'Shadow rule detection across multi-vendor NGFW stack'),
    ('Communications | NA', 'Eliminated Known CVE Exposure', PURPLE,
     'IPS filters were missing on rsync traffic to backup servers. Exposed enterprise to lateral malware movement and remote code execution against critical infrastructure.',
     'Missing IPS inspection on specific protocol'),
    ('Regional Bank | E. Asia', 'Closed RCE Attack Paths', TEAL,
     'Both PANW and Fortinet NGFWs allowed Windows Update service traffic from trusted zones with no IPS -- exposed bank to RCE on critical multi-vendor infrastructure.',
     'Cross-vendor policy inconsistency'),
    ('Critical Infra | SE Asia', 'Closed Application Obfuscation Paths', ORANGE,
     'Incorrectly configured vendor capabilities on DMZ rules -- not all applications on the DMZ were blocked as intended. Susceptible to port-hopping evasion.',
     'DMZ rule misconfiguration'),
    ('Finance | E. Asia', 'NIST 800-53 SMTP Compliance', RED,
     'SMTP traffic from trusted to untrusted zones had no threat inspection enabled, despite NGFW built-in capability -- exposed customer to phishing and credential harvest.',
     'Compliance gap on email traffic'),
    ('EU Sovereign Govt', 'Preemptive ASM for National Agency', NAVY,
     'EU country cyber defense umbrella -- Rigor formally extended EASM scan coverage and proved compliance of external attack surface across 12 government agencies.',
     'National-scale external attack surface management'),
]

y = Inches(1.3)
for customer, title, color, desc, finding_type in findings:
    _box(s, Inches(0.3), y, Inches(12.7), Inches(0.9), fill=CARD, border=BDR)
    _box(s, Inches(0.35), y+Inches(0.05), Inches(2.0), Inches(0.28), fill=color)
    _txt(s, Inches(0.4), y+Inches(0.07), Inches(1.9), Inches(0.24), customer, 9, WHITE, True, PP_ALIGN.CENTER)
    _txt(s, Inches(2.5), y+Inches(0.05), Inches(5.0), Inches(0.25), title, 12, NAVY, True)
    _txt(s, Inches(7.8), y+Inches(0.05), Inches(5.0), Inches(0.25), finding_type, 10, color, True)
    _txt(s, Inches(2.5), y+Inches(0.35), Inches(10.0), Inches(0.5), desc, 10, TXT)
    y += Inches(0.95)

_box(s, Inches(0.3), Inches(7.0), Inches(12.7), Inches(0.35), fill=NAVY)
_txt(s, Inches(0.5), Inches(7.02), Inches(12.3), Inches(0.3),
     '20+ design partners across F500, federal, critical infrastructure, MSSPs, nation states. Engagements in NA, EU, ME, Asia, Pacific.',
     11, GOLD, True, PP_ALIGN.CENTER)


# ================================================================
# SLIDE 11: WHERE RIGOR NEEDS HELP
# ================================================================
s = ns('Where Rigor Needs Help -- Why ACECARD Belongs Next to It',
       'Rigor closes the configuration attack surface. Three classes of attack remain.',
       'INTEGRATION')

gaps = [
    ('1. Zero-Day & Unknown TTPs', RED,
     'Rigor preempts every TTP mapped to MITRE ATT&CK. A genuine zero-day with no signature, no IOC, no MITRE mapping yet has nothing for Rigor to prove against.',
     'A Volt Typhoon exploit published the day after Rigor\'s threat-intel sync. In the gap, behavioral detection is the only defense.',
     'ACECARD detects the resulting behavioral drift even with no MITRE label. Post-exploitation behavior always creates drift.'),
    ('2. Compromised Valid Accounts', ORANGE,
     'Both Typhoons abuse stolen, valid credentials. A correctly-configured firewall, IdP, and EDR will ALLOW that traffic -- by design. Configuration is not the attack surface; behavior is.',
     'Salt Typhoon harvests TACACS+ credentials, then logs in as a legitimate network admin. Every Rigor-verified control says "approved."',
     'ACECARD\'s cohort z-score and cosine drift catch the abuse -- same credentials, structurally different behavior from the real admin.'),
    ('3. Living-off-the-Land Tradecraft', PURPLE,
     'ntdsutil, wmic, netsh, rundll32, PowerShell -- every binary Volt Typhoon uses is a legitimate Windows admin tool. Restricting them entirely is operationally infeasible.',
     'Volt Typhoon runs ntdsutil.exe ifm "create full c:\\temp" -- a real admin command. Rigor verifies EDR coverage is on. The command itself is the threat.',
     'ACECARD\'s 5-signal embedding captures the COMBINATION of LoTL tools, timing, targets. Only behavioral context separates admin from adversary.'),
]

y = Inches(1.3)
for title, color, problem, example, acecard in gaps:
    _box(s, Inches(0.3), y, Inches(12.7), Inches(1.8), fill=CARD, border=color)
    _box(s, Inches(0.35), y+Inches(0.05), Inches(12.55), Inches(0.3), fill=color)
    _txt(s, Inches(0.5), y+Inches(0.07), Inches(12.3), Inches(0.26), title, 14, WHITE, True)
    _txt(s, Inches(0.5), y+Inches(0.4), Inches(12.3), Inches(0.4), problem, 10, TXT)
    _txt(s, Inches(0.5), y+Inches(0.85), Inches(12.3), Inches(0.35), example, 9, SUB)
    _box(s, Inches(0.5), y+Inches(1.25), Inches(0.08), Inches(0.4), fill=CYBER)
    _txt(s, Inches(0.7), y+Inches(1.25), Inches(12.0), Inches(0.5),
         'HANDED TO ACECARD UEBA: ' + acecard, 10, CYBER, True)
    y += Inches(1.88)


# ================================================================
# SLIDE 12: FOUR INTEGRATION POINTS
# ================================================================
s = ns('Four Integration Points: Rigor + ACECARD',
       'How the two layers communicate, share intelligence, and close the complete attack surface',
       'INTEGRATION')

integrations = [
    ('1. MITRE Coverage -> Concept Library', BLUE,
     'Rigor\'s Attack Intelligence builds complete MITRE technique coverage for Volt/Salt Typhoon. These verified TTP patterns feed directly into ACECARD\'s 8 threat concept embeddings.',
     'Rigor: "T1021.002 (SMB lateral) verified blocked by config"\nACECARD: "T1021 lateral_movement concept = behavioral detection if config bypassed"',
     'Shared MITRE technique IDs ensure both layers speak same language'),
    ('2. Rigor Findings -> ACECARD Trust Floor', PURPLE,
     'When Rigor discovers a verified control gap, ACECARD automatically lowers the trust floor for entities in that zone. A known-vulnerable segment starts at ELEVATED_WATCH, not TRUSTED.',
     'Rigor: "DMZ Zone 3 has shadowed SMB IPS profile -- gap open"\nACECARD: "All entities in Zone 3 start at ELEVATED_WATCH until gap closed"',
     'Risk-proportional: worse gaps = lower starting trust = faster detection'),
    ('3. ACECARD True Positives -> Rigor Threat Intel', TEAL,
     'When ACECARD confirms behavioral detection (analyst validates), the entity\'s drift pattern becomes threat intelligence that Rigor ingests for config verification.',
     'ACECARD: "Entity e7a2 confirmed lateral movement via SMB T1021"\nRigor: "Re-verify all SMB inspection profiles across all zones"',
     'Closed loop: behavioral detection improves config verification scope'),
    ('4. Shared Vocabulary: Entity UUIDs + MITRE IDs', GOLD,
     'Both systems reference the same entity_uuid (from entity fusion) and same MITRE technique IDs. A Rigor gap and an ACECARD alert about the same entity are automatically correlated.',
     'Kill-chain reconstruction: Rigor gap at T+0 -> ACECARD drift at T+4h -> same entity_uuid -> one incident',
     'Single pane of glass: SOC sees config gap + behavioral alert as ONE story'),
]

y = Inches(1.3)
for title, color, desc, example, impact in integrations:
    _box(s, Inches(0.3), y, Inches(12.7), Inches(1.35), fill=CARD, border=color)
    _box(s, Inches(0.35), y+Inches(0.05), Inches(12.55), Inches(0.28), fill=color)
    _txt(s, Inches(0.5), y+Inches(0.07), Inches(12.3), Inches(0.24), title, 12, WHITE, True)
    _txt(s, Inches(0.5), y+Inches(0.38), Inches(6.2), Inches(0.55), desc, 9, TXT)
    _txt(s, Inches(6.8), y+Inches(0.38), Inches(5.9), Inches(0.55), example, 8, SUB, fn='Cascadia Code')
    _txt(s, Inches(0.5), y+Inches(1.0), Inches(12.3), Inches(0.3), impact, 9, NAVY, True)
    y += Inches(1.4)


# ================================================================
# SLIDE 13: COMBINED DETECTION TIMELINE
# ================================================================
s = ns('Combined Detection Timeline: Rigor + ACECARD End-to-End',
       'Real scenario: Volt Typhoon campaign against DoD network. Both layers working together.',
       'INTEGRATION')

_box(s, Inches(0.3), Inches(1.25), Inches(12.7), Inches(0.35), fill=NAVY)
_txt(s, Inches(0.5), Inches(1.27), Inches(12.3), Inches(0.3),
     'SCENARIO: Volt Typhoon targets Air Force base network via Fortinet VPN CVE + LoTL credential access', 11, WHITE, True)

timeline = [
    ('T-30d', 'RIGOR', BLUE,
     'Continuous verification identifies Fortinet FortiOS CVE-2024-21887 IPS signature missing on VPN concentrator',
     'Gap finding pushed to ServiceNow. Patch scheduled but delayed by change window.'),
    ('T-30d', 'RIGOR->ACECARD', PURPLE,
     'Integration Point 2: Gap finding triggers trust floor reduction for all entities in VPN zone',
     'Entities entering via FortiGate start at ELEVATED_WATCH (not TRUSTED). Lower detection threshold active.'),
    ('T+0', 'ATTACKER', RED,
     'Volt Typhoon exploits CVE-2024-21887. Lands on FortiGate. Creates local admin. Begins LoTL recon.',
     'The gap Rigor found 30 days ago is now exploited. Config was not patched in time.'),
    ('T+2h', 'ACECARD', TEAL,
     'Entity enters via VPN at ELEVATED_WATCH. AUTH signal: new source, off-hours. Drift = 0.08 (elevated)',
     'Because trust floor is already lowered, ACECARD sensitivity is heightened. Watching closely.'),
    ('T+6h', 'ACECARD', TEAL,
     'PROCESS signal: wmic + ntdsutil chain. Drift = 0.14. CUSUM accumulating. Direction: lotl_execution (0.68)',
     'CUSUM at 0.07 (approaching 2-sigma). ELEVATED_WATCH -> RESTRICTED. Lateral movement blocked.'),
    ('T+8h', 'ACECARD', CYBER,
     'CUSUM crosses 4-sigma (0.125). ALARM. Entity -> BLOCKED. All access revoked. IR ticket auto-created.',
     'Total dwell: 8 hours (vs 5+ years without detection). Kill-chain: T1190 -> T1059 -> T1003.'),
    ('T+9h', 'ACECARD->RIGOR', GOLD,
     'Integration Point 3: Confirmed TP feeds back. Rigor re-verifies all FortiGate inspection profiles globally.',
     'Closed loop: attack detected -> config verification scope expanded -> similar gaps found and closed.'),
]

y = Inches(1.7)
for time, layer, color, event, impact in timeline:
    bg = CARD if timeline.index((time, layer, color, event, impact)) % 2 == 0 else ALT
    _box(s, Inches(0.25), y, Inches(12.85), Inches(0.72), fill=bg, border=BDR)
    _box(s, Inches(0.3), y+Inches(0.05), Inches(0.7), Inches(0.22), fill=NAVY)
    _txt(s, Inches(0.32), y+Inches(0.06), Inches(0.66), Inches(0.2), time, 8, WHITE, True, PP_ALIGN.CENTER, fn='Cascadia Code')
    _box(s, Inches(1.1), y+Inches(0.05), Inches(1.5), Inches(0.22), fill=color)
    _txt(s, Inches(1.12), y+Inches(0.06), Inches(1.46), Inches(0.2), layer, 8, WHITE, True, PP_ALIGN.CENTER)
    _txt(s, Inches(2.7), y+Inches(0.03), Inches(5.2), Inches(0.35), event, 8, TXT)
    _txt(s, Inches(8.0), y+Inches(0.03), Inches(5.0), Inches(0.35), impact, 8, SUB)
    y += Inches(0.75)

_box(s, Inches(0.3), Inches(6.95), Inches(12.7), Inches(0.4), fill=NAVY)
_txt(s, Inches(0.5), Inches(6.97), Inches(12.3), Inches(0.35),
     'Result: 8-hour dwell time (vs 5+ years). Rigor found the gap. ACECARD caught the exploitation. Together: complete defense.',
     11, GOLD, True, PP_ALIGN.CENTER)


# ================================================================
# SLIDE 14: DEPLOYMENT MODELS
# ================================================================
s = ns('Deployment Models: SaaS, On-Prem, Air-Gapped',
       'Rigor supports all three. ACECARD mirrors with IL5/IL6/JWICS enclave support.',
       'DEPLOYMENT')

models = [
    ('SaaS (IL5 / NIPR)', BLUE,
     ['Cloud-hosted Rigor Engine', 'Config ingestion via API / secure connector',
      'Threat intel auto-updated (NVD, CISA KEV, vendor PSIRTs)',
      'Results to ServiceNow / JIRA / SOC SOAR', 'Fastest time-to-value (days)',
      'ACECARD: OpenAI embeddings via API'],
     'Best for: NIPR networks, commercial enterprises, MSSPs'),
    ('On-Premises (IL6 / SIPR)', PURPLE,
     ['Rigor Engine deployed on-site (VM or container)',
      'Config read from local device management', 'Threat intel via classified feeds',
      'Air-gapped from internet', 'Requires local compute (8 vCPU, 32GB min)',
      'ACECARD: Local Phi-4 (14B) or Mistral 7B for embeddings'],
     'Best for: SIPR networks, classified environments, sovereign deployments'),
    ('Air-Gapped (JWICS / TS-SCI)', TEAL,
     ['Fully disconnected deployment', 'Threat intel via manual import (STIX/TAXII)',
      'Config via local file ingest or SNMP read',
      'No external network dependency whatsoever',
      'ACECARD: Same local SLM, physically isolated',
      'Switch between modes: environment variable only, no code change'],
     'Best for: TS/SCI, JWICS, intelligence community, nuclear facilities'),
]

y = Inches(1.3)
for title, color, bullets, best_for in models:
    _box(s, Inches(0.3), y, Inches(12.7), Inches(1.75), fill=CARD, border=color)
    _box(s, Inches(0.35), y+Inches(0.05), Inches(3.0), Inches(0.3), fill=color)
    _txt(s, Inches(0.4), y+Inches(0.07), Inches(2.9), Inches(0.26), title, 12, WHITE, True, PP_ALIGN.CENTER)
    by = y + Inches(0.4)
    for bullet in bullets:
        _txt(s, Inches(0.5), by, Inches(8.5), Inches(0.22), bullet, 10, TXT)
        by += Inches(0.2)
    _txt(s, Inches(9.0), y+Inches(0.7), Inches(3.8), Inches(0.8), best_for, 10, color, True)
    y += Inches(1.85)

_box(s, Inches(0.3), Inches(6.85), Inches(12.7), Inches(0.5), fill=CARD, border=GOLD)
_txt(s, Inches(0.5), Inches(6.87), Inches(12.3), Inches(0.45),
     'KEY: Same codebase across all three deployments. No forking. Environment variables control embedding provider (OpenAI vs local SLM) and connectivity mode. '
     'Security clearance requirement: deployment-model dependent, not code-dependent.',
     10, GOLD, True)


# ================================================================
# SLIDE 15: COMBINED COVERAGE MATRIX
# ================================================================
s = ns('Combined Coverage Matrix: Which Layer Owns What',
       'Each adversary capability owned by the right layer. Preemption first, behavioral as backstop.',
       'INTEGRATION')

mx_hdrs = ['Adversary Capability', 'VT', 'ST', 'Layer', 'Why This Layer Owns It']
mx_x = [Inches(0.3), Inches(4.3), Inches(5.0), Inches(5.7), Inches(7.0)]
mx_w = [Inches(3.9), Inches(0.6), Inches(0.6), Inches(1.2), Inches(6.0)]

_box(s, Inches(0.25), Inches(1.25), Inches(12.8), Inches(0.35), fill=NAVY)
for i, h in enumerate(mx_hdrs):
    _txt(s, mx_x[i], Inches(1.27), mx_w[i], Inches(0.3), h, 9, WHITE, True)

mx_rows = [
    ('Edge-device CVE exploitation', 'Y', 'Y', 'RIGOR', 'Configuration is the gap -- formal verification proves patch/inspection enforced', BLUE),
    ('Misconfigured / shadowed FW rules', 'Y', 'Y', 'RIGOR', 'Mathematical model finds shadows, conflicts, weak rules across multi-vendor', BLUE),
    ('Router ACL tampering / config drift', '', 'Y', 'RIGOR', 'Config Drift use case re-verifies live state vs approved baseline continuously', BLUE),
    ('C2 to known malicious infra', 'Y', 'Y', 'BOTH', 'Rigor blocks known IPs; ACECARD c2_beaconing catches new/unknown infra', PURPLE),
    ('Compromised valid accounts (T1078)', 'Y', 'Y', 'ACECARD', 'Traffic is policy-allowed; only behavioral drift reveals credential abuse', TEAL),
    ('Living-off-the-Land (T1059)', 'Y', 'Y', 'ACECARD', 'lotl_execution concept catches LoTL chain -- process signal embedding shifts', TEAL),
    ('Cross-domain lateral movement', 'Y', 'Y', 'ACECARD', 'Entity fusion joins identity across AD/AAD/AWS/Okta/K8s -- drift = one signal', TEAL),
    ('Slow lateral movement (months)', 'Y', 'Y', 'ACECARD', 'CUSUM 4-sigma surfaces gradual structural shifts that never cross fixed thresholds', TEAL),
    ('Zero-day / unknown TTPs', 'Y', 'Y', 'ACECARD', 'Not in MITRE yet -- only behavioral anomaly flags deviation from learned normal', TEAL),
    ('Insider threat / data hoarding', '', '', 'ACECARD', 'insider_data_hoarding concept (T1074) + peer cohort z-score over months', TEAL),
]

y = Inches(1.65)
for i, (cap, volt, salt, layer, why, color) in enumerate(mx_rows):
    bg = CARD if i % 2 == 0 else ALT
    _box(s, Inches(0.25), y, Inches(12.8), Inches(0.5), fill=bg, border=BDR)
    _txt(s, mx_x[0], y+Pt(2), mx_w[0], Inches(0.45), cap, 10, TXT, True)
    _txt(s, mx_x[1], y+Pt(2), mx_w[1], Inches(0.45), volt, 10, RED if volt else SUB, al=PP_ALIGN.CENTER)
    _txt(s, mx_x[2], y+Pt(2), mx_w[2], Inches(0.45), salt, 10, RED if salt else SUB, al=PP_ALIGN.CENTER)
    _box(s, mx_x[3]+Inches(0.05), y+Inches(0.08), mx_w[3]-Inches(0.1), Inches(0.22), fill=color)
    _txt(s, mx_x[3], y+Inches(0.09), mx_w[3], Inches(0.2), layer, 9, WHITE, True, PP_ALIGN.CENTER)
    _txt(s, mx_x[4], y+Pt(2), mx_w[4], Inches(0.45), why, 9, TXT)
    y += Inches(0.53)

_box(s, Inches(0.3), Inches(6.95), Inches(12.7), Inches(0.35), fill=NAVY)
_txt(s, Inches(0.5), Inches(6.97), Inches(12.3), Inches(0.3),
     'Together: either the path never existed (Rigor) OR the attack is detected and contained within hours (ACECARD + ABAC)',
     11, GOLD, True, PP_ALIGN.CENTER)


# ================================================================
# SLIDE 16: KEY NUMBERS
# ================================================================
s = ns('Key Numbers to Memorize', 'Have these ready for Q&A', 'REFERENCE')

left = [
    ('5+ years', 'Volt Typhoon dwell time'),
    ('4+ years', 'Salt Typhoon dwell time'),
    ('200+', 'Telcos compromised (Salt Typhoon)'),
    ('9', 'US telecoms compromised'),
    ('99%', 'FW breaches = misconfig (Gartner)'),
    ('~20%', 'Attacks stopped by KEV patching'),
    ('2^8000', 'Config state space Rigor reasons over'),
    ('50M+', 'Rule combinations analyzed'),
    ('0', 'False positives (mathematical proof)'),
    ('0', 'False negatives (exhaustive reasoning)'),
    ('<1 hour', 'Config drift detection time'),
    ('20+', 'Design partners (F500, federal, nation)'),
]
right = [
    ('1,536', 'Embedding dimensions per signal'),
    ('9,216', 'Total dims (6 x 1536) per entity'),
    ('5', 'Behavioral signals (AUTH/PROC/NET/FILE/ID)'),
    ('8', 'MITRE threat concepts'),
    ('4-sigma', 'CUSUM alarm threshold'),
    ('8 hours', 'Dwell time in timeline scenario'),
    ('10K/hr', 'Entities per 4-vCPU worker'),
    ('<3 sec', 'Event to analysis per entity'),
    ('$90M', '22CT Army SOC/MDR contract'),
    ('800+', 'Cleared analysts on contract'),
    ('90 days', 'Contract to first detection'),
    ('IL5/IL6/JWICS', 'Three enclave support'),
]

for nums, xs in [(left, Inches(0.3)), (right, Inches(6.8))]:
    y = Inches(1.3)
    for num, meaning in nums:
        _box(s, xs, y, Inches(6.2), Inches(0.42), fill=CARD, border=BDR)
        _box(s, xs+Inches(0.05), y+Inches(0.04), Inches(1.4), Inches(0.3), fill=NAVY)
        _txt(s, xs+Inches(0.1), y+Inches(0.06), Inches(1.3), Inches(0.26),
             num, 11, WHITE, True, PP_ALIGN.CENTER, fn='Cascadia Code')
        _txt(s, xs+Inches(1.6), y+Inches(0.06), Inches(4.4), Inches(0.28),
             meaning, 10, TXT)
        y += Inches(0.44)


# ================================================================
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
