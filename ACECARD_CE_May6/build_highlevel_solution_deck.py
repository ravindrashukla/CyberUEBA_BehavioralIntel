"""High-Level Solution Deck - 22CT Preemptive Cyber Defense
Consolidated 19-slide deck. Content-rich, no repetition.
Covers Volt/Salt Typhoon, Management Plane, ACECARD, Architecture."""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

OUT = os.path.join(os.path.dirname(__file__), "22CT_Preemptive_ACECARD_Volt_Salt_Typhoon_v12.pptx")

NAVY   = RGBColor(0x0D, 0x1B, 0x2A)
BLUE   = RGBColor(0x1B, 0x4F, 0x72)
TEAL   = RGBColor(0x0E, 0x6B, 0x8A)
GOLD   = RGBColor(0xB7, 0x95, 0x0B)
RED    = RGBColor(0xC0, 0x39, 0x2B)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
INK    = RGBColor(0x2C, 0x3E, 0x50)
LGRAY  = RGBColor(0xF7, 0xF8, 0xFA)
MGRAY  = RGBColor(0xE8, 0xEB, 0xEF)
DGRAY  = RGBColor(0x6C, 0x75, 0x7D)
LBLUE  = RGBColor(0xEB, 0xF5, 0xFB)
LTEAL  = RGBColor(0xE8, 0xF4, 0xF8)
LRED   = RGBColor(0xFD, 0xED, 0xEC)
LGOLD  = RGBColor(0xFE, 0xF9, 0xE7)
PURPLE = RGBColor(0x4A, 0x2C, 0x8A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
slide_count = [0]


def blank():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    slide_count[0] += 1
    return s

def brand_bar(s):
    rect(s, Inches(0), Inches(7.15), W, Inches(0.35), fill=NAVY)
    txt(s, Inches(0.5), Inches(7.17), Inches(6), Inches(0.3),
        '22nd Century Technologies  |  PROPRIETARY  |  ravindra.shukla@tscti.com', 8, RGBColor(0x8A, 0x9B, 0xAE))
    txt(s, Inches(11.5), Inches(7.17), Inches(1.5), Inches(0.3),
        str(slide_count[0]), 8, RGBColor(0x8A, 0x9B, 0xAE), al=PP_ALIGN.RIGHT)

def rect(s, x, y, w, h, fill=None, border=None):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if border:
        sh.line.color.rgb = border
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def rrect(s, x, y, w, h, fill=None, border=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if border:
        sh.line.color.rgb = border
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def hline(s, x, y, w, color=MGRAY, thick=Pt(1)):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, thick)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False

def txt(s, x, y, w, h, text, sz=13, c=INK, b=False, al=PP_ALIGN.LEFT, fn='Segoe UI'):
    tf = s.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = c
    p.font.bold = b
    p.font.name = fn
    p.alignment = al
    return tf

def add_p(tf, text, sz=13, c=INK, b=False, sp=Pt(8)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = c
    p.font.bold = b
    p.font.name = 'Segoe UI'
    p.space_before = sp
    return p

def slide_title(s, title, takeaway=None):
    txt(s, Inches(1.2), Inches(0.35), Inches(10.5), Inches(0.7),
        title, 24, NAVY, True, fn='Georgia')
    hline(s, Inches(1.2), Inches(1.05), Inches(2.0), BLUE, Pt(3))
    if takeaway:
        txt(s, Inches(1.2), Inches(1.15), Inches(10.5), Inches(0.35),
            takeaway, 12, BLUE)
    brand_bar(s)


IMG_DIR = os.path.join(os.path.dirname(__file__), "images")

# ═══════════════════════��══════════════════════��═══════════════════════��═
# SLIDE 1: TITLE
# ════════════════════════════��═══════════════════════════════════════════
s = blank()
s.background.fill.solid()
s.background.fill.fore_color.rgb = NAVY
rect(s, Inches(0), Inches(0), W, Inches(0.06), fill=GOLD)
banner_path = os.path.join(IMG_DIR, "slide1_Picture_4.png")
if os.path.exists(banner_path):
    s.shapes.add_picture(banner_path, Inches(0), Inches(0.1), Inches(6.5), Inches(1.5))
iso_path = os.path.join(IMG_DIR, "slide1_Picture_6.png")
if os.path.exists(iso_path):
    s.shapes.add_picture(iso_path, Inches(11.5), Inches(0.2), Inches(0.8), Inches(0.8))
txt(s, Inches(10.0), Inches(1.1), Inches(3.0), Inches(0.3),
    'ISO 9001 | 27001 | 28000 | 20000 | 17025 | 14001',
    7, RGBColor(0x8A, 0x9B, 0xAE), al=PP_ALIGN.RIGHT)
txt(s, Inches(7.5), Inches(0.2), Inches(3.7), Inches(0.35),
    'U.S. ARMY CYBER COMMAND  |  ArCTIC / CFIC', 9, GOLD, True, PP_ALIGN.CENTER, fn='Georgia')
txt(s, Inches(1.8), Inches(1.8), Inches(10), Inches(0.4),
    '22nd Century Technologies Inc.', 13, GOLD, True, fn='Georgia')
txt(s, Inches(1.8), Inches(2.5), Inches(10), Inches(1.5),
    'Preemptive Cyber Defense\n& Behavioral Intelligence Platform', 36, WHITE, True, fn='Georgia')
txt(s, Inches(1.8), Inches(4.3), Inches(10), Inches(0.5),
    'Provable Computing  |  DODIN Defense  |  Behavioral Intelligence', 14, RGBColor(0xA0, 0xB8, 0xD0))
hline(s, Inches(1.8), Inches(5.1), Inches(3.5), GOLD, Pt(2))
txt(s, Inches(1.8), Inches(5.3), Inches(10), Inches(0.4),
    'Formal Proof, Not Sampling — Defending the Warfighter Against Nation-State Adversaries', 13, RGBColor(0x7A, 0x8B, 0x9E))
txt(s, Inches(1.8), Inches(5.9), Inches(10), Inches(0.3),
    'Connecting & Serving Government & Citizens since 1997', 10, RGBColor(0x7A, 0x8B, 0x9E))
txt(s, Inches(1.8), Inches(6.2), Inches(10), Inches(0.3),
    'ArCTIC / CFIC Collaboration Event  |  Behavioral Anomaly Detection (ACECARD)  |  06 MAY 2026  |  Augusta, GA',
    11, RGBColor(0x7A, 0x8B, 0x9E))
txt(s, Inches(1.8), Inches(6.5), Inches(10), Inches(0.3),
    'PROPRIETARY  //  NOT FOR PUBLIC RELEASE', 9, RGBColor(0x5A, 0x6B, 0x7E))
brand_bar(s)


# ═════════════════════════════════════════���══════════════════════════════
# SLIDE 2: PROBLEM STATEMENT
# ═════════════════════════════════════════��══════════════════════════════
s = blank()
slide_title(s, 'Problem Statement',
    'Nation-state adversaries dwelling inside U.S. critical infrastructure — undetected 4–5+ years.')

rrect(s, Inches(1.2), Inches(1.6), Inches(11.0), Inches(1.6), fill=LRED, border=MGRAY)
rect(s, Inches(1.2), Inches(1.6), Inches(11.0), Inches(0.06), fill=RED)
tf = txt(s, Inches(1.5), Inches(1.75), Inches(10.5), Inches(1.4),
    'Two confirmed nation-state campaigns — Volt Typhoon and Salt Typhoon — have compromised', 13, INK, True)
add_p(tf, 'hundreds of critical infrastructure targets using techniques no existing tool can detect:', 13, INK, True)
add_p(tf, '•  Living-off-the-Land (LOTL): only legitimate admin tools — no malware signatures to find', 12, INK)
add_p(tf, '•  Configuration-resident persistence: the attack IS the network configuration itself', 12, INK)
add_p(tf, '•  Stolen valid credentials: every access control says "approved"', 12, INK)

rrect(s, Inches(1.2), Inches(3.5), Inches(11.0), Inches(1.6), fill=LGRAY, border=MGRAY)
rect(s, Inches(1.2), Inches(3.5), Inches(0.06), Inches(1.6), fill=NAVY)
tf = txt(s, Inches(1.5), Inches(3.6), Inches(10.5), Inches(1.4),
    'ARCYBER DEFEND-THE-DODIN CHALLENGE — LIMITING FACTORS:', 12, NAVY, True)
add_p(tf, '•  Tools detect known-bad signatures → adversaries use only known-good tools', 12, INK)
add_p(tf, '•  Fixed thresholds → LOTL operates below every one', 12, INK)
add_p(tf, '•  Point-in-time testing (annual pentests) → adversary is continuous', 12, INK)
add_p(tf, '•  No tool provides formal proof that all attack paths on the DODIN are blocked', 12, RED, True)

rrect(s, Inches(1.2), Inches(5.4), Inches(11.0), Inches(1.4), fill=LBLUE, border=MGRAY)
rect(s, Inches(1.2), Inches(5.4), Inches(0.06), Inches(1.4), fill=BLUE)
tf = txt(s, Inches(1.5), Inches(5.5), Inches(10.5), Inches(1.2),
    'OUR APPROACH — PROVABLE COMPUTING + BEHAVIORAL INTELLIGENCE:', 12, BLUE, True)
add_p(tf, '1.  Preemptive: Formal proof, not sampling — mathematical verification that configs block all known attack paths', 12, INK)
add_p(tf, '2.  Behavioral (ACECARD): Digital Twin of each entity — detects sub-threshold drift that no threshold system catches', 12, INK)
add_p(tf, 'Together: 95–98% DODIN coverage. Dwell time: 5 years → hours. Directly enhances ARCYBER DEFEND mission.', 12, NAVY, True)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 3: THREAT LANDSCAPE — VOLT & SALT TYPHOON
# ═══════════════════��═════════════════════════════��══════════════════════
s = blank()
slide_title(s, 'The Threat: Volt Typhoon & Salt Typhoon',
    'Two active PRC campaigns using techniques that bypass every existing security tool.')

cw = Inches(5.3)
rrect(s, Inches(1.2), Inches(1.5), cw, Inches(2.8), fill=LRED, border=MGRAY)
rect(s, Inches(1.2), Inches(1.5), cw, Inches(0.5), fill=RED)
txt(s, Inches(1.4), Inches(1.58), Inches(5.0), Inches(0.4), '  VOLT TYPHOON (G1017)', 13, WHITE, True)
tf = txt(s, Inches(1.5), Inches(2.1), Inches(4.8), Inches(2.0),
    'PRC state-sponsored. Active since mid-2021.', 12, INK)
add_p(tf, '•  Pre-positioned in power, water, telecom, transportation', 12, INK)
add_p(tf, '•  Living-off-the-Land: only legitimate admin tools', 12, INK)
add_p(tf, '•  No malware, no signatures, no IOCs', 12, INK)
add_p(tf, '•  5+ year undetected dwell time', 12, RED, True)
add_p(tf, '•  Purpose: future disruption during Taiwan crisis', 12, INK)

rrect(s, Inches(6.9), Inches(1.5), cw, Inches(2.8), fill=LBLUE, border=MGRAY)
rect(s, Inches(6.9), Inches(1.5), cw, Inches(0.5), fill=NAVY)
txt(s, Inches(7.1), Inches(1.58), Inches(5.0), Inches(0.4), '  SALT TYPHOON', 13, WHITE, True)
tf = txt(s, Inches(7.2), Inches(2.1), Inches(4.8), Inches(2.0),
    '200+ telecoms compromised across 80+ countries.', 12, INK)
add_p(tf, '•  9 major U.S. carriers. CALEA wiretap access.', 12, INK)
add_p(tf, '•  Persistence IS the router configuration itself', 12, INK)
add_p(tf, '•  Exploits Cisco IOS XE (CVE-2023-20198)', 12, INK)
add_p(tf, '•  4+ year undetected dwell time', 12, RED, True)
add_p(tf, '•  Steady CDR/metadata exfiltration over months', 12, INK)

rrect(s, Inches(1.2), Inches(4.6), Inches(11.0), Inches(2.2), fill=LGRAY, border=MGRAY)
rect(s, Inches(1.2), Inches(4.6), Inches(0.06), Inches(2.2), fill=RED)
tf = txt(s, Inches(1.5), Inches(4.7), Inches(10.5), Inches(2.0),
    'WHY CURRENT DEFENSES CANNOT STOP THEM:', 12, NAVY, True)
add_p(tf, '•  SIEM/EDR: detect known-bad signatures — these adversaries use only known-good tools', 12, INK)
add_p(tf, '•  Commercial UEBA: fixed thresholds — LOTL operates below every one', 12, INK)
add_p(tf, '•  Pentest/Red Team: annual point-in-time sampling vs. continuous adversary', 12, INK)
add_p(tf, '•  Config gaps never mathematically proven closed — pentests only sample attack paths', 12, INK)
add_p(tf, '•  Identity fragmented across 10+ systems — cross-domain movement reads as unrelated signals', 12, INK)

txt(s, Inches(1.2), Inches(6.9), Inches(11.0), Inches(0.2),
    'Sources: CISA AA24-038A  |  FBI/NSA/CISA Joint Advisory (Dec 2024)  |  Microsoft Threat Intel',
    9, DGRAY)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 4: SECURITY DEVICES — CAPABILITIES & GAPS
# ══════════════════���════════════════��════════════════════════════════════
s = blank()
slide_title(s, 'Security Devices: What They Do vs. What They Miss',
    'Each device solves one piece — no device verifies that ALL pieces work together.')

devices = [
    ('Next-Gen Firewall', 'Filters traffic by rules, app-aware policies, zones',
     'Rules only as good as config. No tool proves the ruleset covers all attack paths.'),
    ('IDS / IPS', 'Detects/blocks known attack signatures in traffic',
     'Signature-dependent. Volt Typhoon uses no malicious payloads — only legit admin commands.'),
    ('NAC (802.1X)', 'Controls which devices connect to the network',
     'Grants access, not intent. Lateral movement with valid credentials is invisible.'),
    ('WAF', 'Protects web apps from injection, XSS, known exploits',
     'HTTP layer only. Cannot see infrastructure-layer LOTL or router config manipulation.'),
    ('DLP', 'Detects sensitive data exfiltration patterns',
     'Pattern-based. Slow exfil over months using normal-volume channels bypasses thresholds.'),
    ('PAM / IAM', 'Manages privileged access, credential vaults, MFA',
     'Enforces policy — cannot detect stolen sessions used within authorized bounds.'),
    ('Email Gateway', 'Blocks phishing, malware attachments, spoofing',
     "Volt/Salt don't rely on phishing — they exploit infrastructure directly."),
    ('EDR / XDR', 'Monitors host processes, detects malicious behavior',
     "Can't deploy on routers/switches. LOTL uses only allowlisted binaries — nothing triggers."),
    ('Vuln Scanner', 'Enumerates assets and known CVEs',
     'Sees what EXISTS, not what is REACHABLE through config. Sampling, not proof.'),
]

# Header row
rrect(s, Inches(1.2), Inches(1.45), Inches(2.4), Inches(0.38), fill=NAVY)
txt(s, Inches(1.4), Inches(1.48), Inches(2.2), Inches(0.32), 'DEVICE', 10, WHITE, True, PP_ALIGN.CENTER)
rrect(s, Inches(3.6), Inches(1.45), Inches(3.6), Inches(0.38), fill=BLUE)
txt(s, Inches(3.8), Inches(1.48), Inches(3.4), Inches(0.32), 'WHAT IT DOES', 10, WHITE, True, PP_ALIGN.CENTER)
rrect(s, Inches(7.2), Inches(1.45), Inches(5.0), Inches(0.38), fill=RED)
txt(s, Inches(7.4), Inches(1.48), Inches(4.8), Inches(0.32), 'CAPABILITY GAP', 10, WHITE, True, PP_ALIGN.CENTER)

y = Inches(1.88)
row_h = Inches(0.53)
for i, (device, does, gap) in enumerate(devices):
    bg = LGRAY if i % 2 == 0 else WHITE
    rrect(s, Inches(1.2), y, Inches(11.0), row_h, fill=bg, border=MGRAY)
    txt(s, Inches(1.35), y + Inches(0.04), Inches(2.2), row_h, device, 10, NAVY, True)
    txt(s, Inches(3.7), y + Inches(0.04), Inches(3.4), row_h, does, 9.5, INK)
    txt(s, Inches(7.3), y + Inches(0.04), Inches(4.8), row_h, gap, 9.5, RED)
    y += row_h + Inches(0.02)

rrect(s, Inches(1.2), Inches(6.85), Inches(11.0), Inches(0.25), fill=LGOLD)
txt(s, Inches(1.5), Inches(6.86), Inches(10.5), Inches(0.22),
    'The gap is structural: no device validates total configuration across ALL devices blocks every attack path.',
    10, NAVY, True, PP_ALIGN.CENTER)


# ═════════════════���══════════════════════���═══════════════════════════════
# SLIDE 5: OUR APPROACH — TWO LAYERS + 80% FUNNEL
# ══════════════��══════════════════════════════════��══════════════════════
s = blank()
slide_title(s, 'Our Solution: Two Layers of DODIN Defense',
    'Formal proof, not sampling. 80% eliminated at network layer. Behavioral catches 15-18% that passes through.')

# Layer 1 block
rect(s, Inches(1.2), Inches(1.5), Inches(11.0), Inches(2.3), fill=LBLUE, border=MGRAY)
rect(s, Inches(1.2), Inches(1.5), Inches(11.0), Inches(0.45), fill=BLUE)
txt(s, Inches(1.5), Inches(1.55), Inches(5.0), Inches(0.35),
    'LAYER 1 — PREEMPTIVE VERIFICATION', 12, WHITE, True)
rect(s, Inches(9.5), Inches(1.55), Inches(2.5), Inches(0.35), fill=NAVY)
txt(s, Inches(9.55), Inches(1.57), Inches(2.4), Inches(0.3), '80% STOPPED HERE', 10, WHITE, True, PP_ALIGN.CENTER)
tf = txt(s, Inches(1.5), Inches(2.1), Inches(5.0), Inches(1.5),
    'Provable Computing: Formal Mathematical Model that', 12, INK, True)
add_p(tf, 'reasons exhaustively over the entire DODIN config space.', 12, INK)
add_p(tf, '', 6)
add_p(tf, '•  Ingests all threat intel continuously (MITRE attack graphs)', 11, INK)
add_p(tf, '•  Models every defensive control formally', 11, INK)
add_p(tf, '•  Proves whether ALL attack paths are blocked', 11, INK)
add_p(tf, '•  Prescribes vendor-specific remediations', 11, INK)
add_p(tf, '•  Re-verifies hourly + on every config change', 11, INK)

tf = txt(s, Inches(7.0), Inches(2.1), Inches(5.0), Inches(1.5),
    'What it stops:', 12, BLUE, True)
add_p(tf, '•  CVE exploits: proves IPS/threat-feed blocks every exploit path', 11, INK)
add_p(tf, '•  C2 traffic: 50M+ rule combinations verified', 11, INK)
add_p(tf, '•  Config drift: unauthorized changes caught within 1 hour', 11, INK)
add_p(tf, '•  Shadow rules, conflicts, weak inspection profiles', 11, INK)
add_p(tf, '', 6)
add_p(tf, 'Soundness within model scope. Zero FP. Zero FN.', 10, BLUE, True)

# Arrow
txt(s, Inches(5.5), Inches(3.85), Inches(2.5), Inches(0.3), '▼  20% passes through  ▼', 10, DGRAY, al=PP_ALIGN.CENTER)

# Layer 2 block
rect(s, Inches(2.0), Inches(4.15), Inches(9.4), Inches(2.3), fill=LTEAL, border=MGRAY)
rect(s, Inches(2.0), Inches(4.15), Inches(9.4), Inches(0.45), fill=TEAL)
txt(s, Inches(2.3), Inches(4.2), Inches(5.0), Inches(0.35),
    'LAYER 2 — BEHAVIORAL INTELLIGENCE (ACECARD)', 12, WHITE, True)
rect(s, Inches(9.0), Inches(4.2), Inches(2.2), Inches(0.35), fill=NAVY)
txt(s, Inches(9.05), Inches(4.22), Inches(2.1), Inches(0.3), '15-18% DETECTED', 10, WHITE, True, PP_ALIGN.CENTER)
tf = txt(s, Inches(2.3), Inches(4.75), Inches(4.5), Inches(1.5),
    'Digital Twin of each entity instance with', 12, INK, True)
add_p(tf, 'time dimension — behavioral evolution captured.', 12, INK)
add_p(tf, '', 6)
add_p(tf, '•  Multi-signal behavioral profiling per entity', 11, INK)
add_p(tf, '•  Advanced change-point detection (sub-threshold)', 11, INK)
add_p(tf, '•  Drift direction mapped to MITRE ATT&CK', 11, INK)
add_p(tf, '•  Cross-system entity fusion (10+ identity sources)', 11, INK)

tf = txt(s, Inches(7.0), Inches(4.75), Inches(4.2), Inches(1.5),
    'What it catches:', 12, TEAL, True)
add_p(tf, '•  Compromised valid credentials (behaviorally abnormal)', 11, INK)
add_p(tf, '•  LOTL: legit binaries in malicious combinations', 11, INK)
add_p(tf, '•  Zero-day post-exploitation (drift with no signature)', 11, INK)
add_p(tf, '•  Slow lateral movement (sub-threshold accumulation)', 11, INK)
add_p(tf, '', 6)
add_p(tf, 'DETECTS: what no threshold system can see.', 10, TEAL, True)

# Bottom result
rrect(s, Inches(1.2), Inches(6.6), Inches(11.0), Inches(0.35), fill=LGOLD)
txt(s, Inches(1.5), Inches(6.63), Inches(10.5), Inches(0.3),
    'Combined: 95-98% coverage.  2-5% residual risk.  Dwell time: years → hours.',
    11, NAVY, True, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════��════════════════════
# SLIDE 6: THREE INTELLIGENCES
# ══════════════════���══════════════════��══════════════════════════════════
s = blank()
slide_title(s, 'Preemptive Layer: Three Intelligences (Provable Computing)',
    'Attack Intelligence → Defense Intelligence → Remediation Intelligence. Continuous closed loop.')

col_w = Inches(3.4)
col_h = Inches(4.5)
gap_x = Inches(0.25)
x_start = Inches(1.2)

x = x_start
rrect(s, x, Inches(1.6), col_w, col_h, fill=LRED, border=MGRAY)
rect(s, x, Inches(1.6), col_w, Inches(0.5), fill=RED)
txt(s, x+Inches(0.2), Inches(1.68), Inches(3.0), Inches(0.4),
    '  ATTACK INTELLIGENCE', 12, WHITE, True)
tf = txt(s, x+Inches(0.25), Inches(2.25), Inches(2.9), Inches(3.7),
    'Analyze all known threat intelligence', 12, INK, True)
add_p(tf, '', 6)
add_p(tf, '•  Ingests CVE, advisories, CTI, IOCs, IOAs, PSIRTs continuously', 11, INK)
add_p(tf, '•  MITRE-enriched attack graphs by ML/NLP', 11, INK)
add_p(tf, '•  Full Volt/Salt Typhoon TTP coverage', 11, INK)
add_p(tf, '•  Auto-updates on every new advisory', 11, INK)
add_p(tf, '', 6)
add_p(tf, 'Sources:', 9, DGRAY, True)
add_p(tf, 'CISA KEV, NVD, NSA, vendor PSIRTs, MITRE ATT&CK, classified feeds', 9, DGRAY)

x = x_start + col_w + gap_x
rrect(s, x, Inches(1.6), col_w, col_h, fill=LBLUE, border=MGRAY)
rect(s, x, Inches(1.6), col_w, Inches(0.5), fill=BLUE)
txt(s, x+Inches(0.2), Inches(1.68), Inches(3.0), Inches(0.4),
    '  DEFENSE INTELLIGENCE', 12, WHITE, True)
tf = txt(s, x+Inches(0.25), Inches(2.25), Inches(2.9), Inches(3.7),
    'Formal model of every defensive control', 12, INK, True)
add_p(tf, '', 6)
add_p(tf, '•  Every NGFW rule, ACL, IPS, identity policy modeled', 11, INK)
add_p(tf, '•  Exhaustive reasoning over complete config state space', 11, INK)
add_p(tf, '•  No sampling — mathematical certainty', 11, INK)
add_p(tf, '•  Every path analyzed. Every combination verified.', 11, INK)
add_p(tf, '', 6)
add_p(tf, 'Supports:', 9, DGRAY, True)
add_p(tf, 'Palo Alto, Fortinet, Cisco, Check Point (N-S + E-W)', 9, DGRAY)

x = x_start + 2 * (col_w + gap_x)
rrect(s, x, Inches(1.6), col_w, col_h, fill=LTEAL, border=MGRAY)
rect(s, x, Inches(1.6), col_w, Inches(0.5), fill=TEAL)
txt(s, x+Inches(0.2), Inches(1.68), Inches(3.0), Inches(0.4),
    '  REMEDIATION INTELLIGENCE', 12, WHITE, True)
tf = txt(s, x+Inches(0.25), Inches(2.25), Inches(2.9), Inches(3.7),
    'Prescribe correct configurations', 12, INK, True)
add_p(tf, '', 6)
add_p(tf, '•  Risk-prioritized, vendor-specific playbooks', 11, INK)
add_p(tf, '•  Zero false positives (mathematically proven)', 11, INK)
add_p(tf, '•  Zero false negatives (exhaustive reasoning)', 11, INK)
add_p(tf, '•  Auto-push to ServiceNow, JIRA, SOC ticketing', 11, INK)
add_p(tf, '•  Optional auto-apply + re-verify', 11, INK)
add_p(tf, '', 6)
add_p(tf, 'Guarantees:', 9, DGRAY, True)
add_p(tf, 'No new errors, no broken connectivity, no disruption', 9, DGRAY)

txt(s, x_start + col_w, Inches(3.5), Inches(0.4), Inches(0.4), '→', 20, NAVY, True, PP_ALIGN.CENTER)
txt(s, x_start + col_w + col_w + gap_x, Inches(3.5), Inches(0.4), Inches(0.4), '→', 20, NAVY, True, PP_ALIGN.CENTER)

rrect(s, Inches(1.2), Inches(6.3), Inches(11.0), Inches(0.5), fill=LGOLD)
txt(s, Inches(1.5), Inches(6.33), Inches(10.5), Inches(0.4),
    'Continuous loop: new threat → model attack → verify defense → prescribe fix → re-verify. Every hour + on every config change.',
    10, NAVY, True, PP_ALIGN.CENTER)


# ════════════��══════════════════════════════════��════════════════════════
# SLIDE 7: MANAGEMENT PLANE
# ════════════════════════════��═══════════════════════════════════════════
s = blank()
slide_title(s, 'The Management Plane: Making All Devices Work Together',
    'A new layer ABOVE all security devices. Compositional reasoning across multi-vendor configurations at DODIN scale.')

# Problem statement
rrect(s, Inches(1.2), Inches(1.45), Inches(11.0), Inches(0.8), fill=LRED, border=MGRAY)
rect(s, Inches(1.2), Inches(1.45), Inches(0.08), Inches(0.8), fill=RED)
tf = txt(s, Inches(1.5), Inches(1.5), Inches(10.5), Inches(0.7),
    'THE PROBLEM: Devices are powerful individually but do not coordinate. Gaps emerge BETWEEN devices —', 12, INK, True)
add_p(tf, 'multi-rule, multi-device errors that no single-device tool can see.', 12, INK, True)

# Management Plane diagram
rect(s, Inches(1.2), Inches(2.5), Inches(11.0), Inches(1.5), fill=BLUE)
txt(s, Inches(1.5), Inches(2.55), Inches(3.5), Inches(0.35),
    'MANAGEMENT PLANE', 15, WHITE, True, fn='Georgia')
# Four steps inline
steps = [('INGEST', 'Pull configs\nfrom all devices'), ('MODEL', 'Build unified\nmath model'),
         ('VERIFY', 'Find multi-device\ngaps'), ('REMEDIATE', 'Push vendor-\nspecific fixes')]
sx = Inches(1.5)
for label, desc in steps:
    rrect(s, sx, Inches(3.0), Inches(2.4), Inches(0.85), fill=RGBColor(0x14, 0x3D, 0x5C), border=WHITE)
    txt(s, sx+Inches(0.1), Inches(3.05), Inches(2.2), Inches(0.3), label, 11, GOLD, True, PP_ALIGN.CENTER)
    txt(s, sx+Inches(0.1), Inches(3.35), Inches(2.2), Inches(0.45), desc, 9, WHITE, al=PP_ALIGN.CENTER)
    sx += Inches(2.6)
# Arrows between steps
for i in range(3):
    ax = Inches(1.5) + Inches(2.4) + Inches(2.6) * i
    txt(s, ax, Inches(3.3), Inches(0.25), Inches(0.3), '→', 14, WHITE, True)

# Device layer
txt(s, Inches(5.8), Inches(4.05), Inches(2.0), Inches(0.2), '▼  Verified configs  ▼', 8, NAVY, True, PP_ALIGN.CENTER)
rect(s, Inches(1.2), Inches(4.25), Inches(11.0), Inches(0.75), fill=LGRAY, border=MGRAY)
devices_row = ['NGFW\n(PAN, Fortinet\nCisco, CP)', 'IPS / WAF', 'Identity\n(Okta, Entra\nCyberArk)', 'SASE /\nSD-WAN', 'EDR / XDR', 'DLP / Email\nSecurity']
dx = Inches(1.4)
for dev in devices_row:
    rrect(s, dx, Inches(4.3), Inches(1.65), Inches(0.65), fill=WHITE, border=BLUE)
    txt(s, dx+Inches(0.05), Inches(4.33), Inches(1.55), Inches(0.6), dev, 8, NAVY, True, PP_ALIGN.CENTER)
    dx += Inches(1.78)

# Real findings
txt(s, Inches(1.2), Inches(5.2), Inches(11.0), Inches(0.25),
    'Production findings from Management Plane analysis:', 12, NAVY, True)
findings_mp = [
    ('Financial Services', 'Conflicting rules caused access to malicious IPs — found by cross-device analysis', BLUE),
    ('F100 Technology', '100+ config errors across multi-vendor NGFW estate — invisible to single-device tools', TEAL),
    ('F500 Data Center', 'Absence of data security profile on critical segments — no single tool could see it', GOLD),
    ('Supply Chain', 'Shadowed SMB app-control rules allowed lateral malware migration (SMBv1.0)', RED),
]
fy = Inches(5.5)
for cust, finding, color in findings_mp:
    rrect(s, Inches(1.2), fy, Inches(11.0), Inches(0.38), fill=LGRAY, border=MGRAY)
    rect(s, Inches(1.2), fy, Inches(0.06), Inches(0.38), fill=color)
    txt(s, Inches(1.5), fy+Inches(0.05), Inches(2.0), Inches(0.28), cust, 9, color, True)
    txt(s, Inches(3.6), fy+Inches(0.05), Inches(8.4), Inches(0.28), finding, 10, INK)
    fy += Inches(0.41)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 8: STOPPING VOLT TYPHOON (BOTH LAYERS)
# ═════════════════════════════════════════════════════════���══════════════
s = blank()
slide_title(s, 'Stopping Volt Typhoon: Both Layers',
    'Preemptive eliminates config weakness. ACECARD catches behavioral exploitation.')

# Preemptive section
rect(s, Inches(1.2), Inches(1.55), Inches(11.0), Inches(0.3), fill=BLUE)
txt(s, Inches(1.4), Inches(1.57), Inches(5.0), Inches(0.26), 'PREEMPTIVE LAYER — config verification', 10, WHITE, True)

phases_p = [
    ('Initial Access', 'Exploits unpatched Fortinet/Ivanti/PAN edge devices',
     'Maps every CVE to control. Proves IPS blocks exploit. Identifies compensating control.'),
    ('Persistence', 'Web shells. C2 via compromised SOHO routers',
     'Threat-feed IPs denied. 50M+ rule combos verified. WAF upload inspection proven.'),
    ('Lateral Movement', 'SMB/RDP/WinRM to DCs and file servers',
     'Real finding: 3 shadowed rules exposed SMB. Model closed path no scanner found.'),
]
y = Inches(1.9)
for phase, attack, defense in phases_p:
    rrect(s, Inches(1.2), y, Inches(11.0), Inches(0.6), fill=LGRAY, border=MGRAY)
    txt(s, Inches(1.4), y+Inches(0.05), Inches(1.8), Inches(0.5), phase, 10, NAVY, True)
    txt(s, Inches(3.2), y+Inches(0.05), Inches(3.3), Inches(0.5), attack, 10, RED)
    txt(s, Inches(6.6), y+Inches(0.05), Inches(5.4), Inches(0.5), defense, 10, INK)
    y += Inches(0.65)

# ACECARD section
rect(s, Inches(1.2), y+Inches(0.1), Inches(11.0), Inches(0.3), fill=TEAL)
txt(s, Inches(1.4), y+Inches(0.12), Inches(5.0), Inches(0.26), 'ACECARD LAYER — behavioral detection', 10, WHITE, True)

phases_a = [
    ('Credential abuse\n(T1078)', 'Valid creds — every control says "approved"',
     'Multi-signal: new source + off-hours + unusual targets. Structural shift triggers alarm.'),
    ('LOTL execution\n(T1059)', 'ntdsutil + wmic + netsh — all legit binaries',
     'Tool combination + sequence captured. Drift → "LOTL execution." Peer cohort: no match.'),
    ('Slow lateral\n(T1021)', 'Moves 1-2 hosts/week — below every threshold',
     'Change-point accumulates sub-threshold drift over time. Eventually crosses alarm boundary.'),
    ('Unknown C2', 'Routes through newly compromised SOHO routers',
     'Network signal: new persistent outbound to unusual destination. Direction: "C2 beaconing."'),
]
y += Inches(0.45)
for technique, misses, catches in phases_a:
    rrect(s, Inches(1.2), y, Inches(11.0), Inches(0.6), fill=LTEAL, border=MGRAY)
    txt(s, Inches(1.4), y+Inches(0.05), Inches(1.8), Inches(0.5), technique, 9.5, NAVY, True)
    txt(s, Inches(3.2), y+Inches(0.05), Inches(3.3), Inches(0.5), misses, 10, RED)
    txt(s, Inches(6.6), y+Inches(0.05), Inches(5.4), Inches(0.5), catches, 10, INK)
    y += Inches(0.65)

rrect(s, Inches(1.2), Inches(6.8), Inches(11.0), Inches(0.3), fill=LGOLD)
txt(s, Inches(1.5), Inches(6.82), Inches(10.5), Inches(0.26),
    'Volt Typhoon (G1017): PRC state-sponsored. 5+ year dwell. LOTL. Critical infrastructure. Both layers needed.',
    9.5, NAVY, True, PP_ALIGN.CENTER)


# ════════════════════════════════════════════════���═══════════════════════
# SLIDE 9: STOPPING SALT TYPHOON (BOTH LAYERS)
# ═══════════════════════════════════��═════════════════════════════════���══
s = blank()
slide_title(s, 'Stopping Salt Typhoon: Both Layers',
    'Salt Typhoon lives on misconfigured edge gear. Preemptive verifies. ACECARD detects abuse.')

# Preemptive section
rect(s, Inches(1.2), Inches(1.55), Inches(11.0), Inches(0.3), fill=BLUE)
txt(s, Inches(1.4), Inches(1.57), Inches(5.0), Inches(0.26), 'PREEMPTIVE LAYER — config verification', 10, WHITE, True)

phases_p = [
    ('Edge CVE Exploit', 'Cisco IOS XE web UI (CVE-2023-20198). 1000+ devices.',
     'Every exposed mgmt interface flagged. Proves ACL blocks /webui/. Compensating control for unpatched.'),
    ('Config Persistence', 'Creates local admin. Modifies router ACLs. Enables Guest Shell.',
     'Hourly drift comparison to approved AAA baseline. New account or ACL change = immediate finding.'),
    ('GRE Tunneling', 'Configures GRE tunnels between routers for C2.',
     'GRE peer authorization verified. Unauthorized tunnel endpoints flagged. Egress policy proven.'),
]
y = Inches(1.9)
for phase, attack, defense in phases_p:
    rrect(s, Inches(1.2), y, Inches(11.0), Inches(0.6), fill=LGRAY, border=MGRAY)
    txt(s, Inches(1.4), y+Inches(0.05), Inches(1.8), Inches(0.5), phase, 10, NAVY, True)
    txt(s, Inches(3.2), y+Inches(0.05), Inches(3.3), Inches(0.5), attack, 10, RED)
    txt(s, Inches(6.6), y+Inches(0.05), Inches(5.4), Inches(0.5), defense, 10, INK)
    y += Inches(0.65)

# ACECARD section
rect(s, Inches(1.2), y+Inches(0.1), Inches(11.0), Inches(0.3), fill=TEAL)
txt(s, Inches(1.4), y+Inches(0.12), Inches(5.0), Inches(0.26), 'ACECARD LAYER — behavioral detection', 10, WHITE, True)

phases_a = [
    ('TACACS+ reuse\n(T1078)', 'Harvested admin creds used as legit sessions',
     'Same creds, different trajectory: unusual device sequence, off-pattern timing. Multi-signal drift alarm.'),
    ('ACL tampering\n(T1556)', 'Modifies router ACLs — looks like routine maintenance',
     'Identity + network signals: config change correlates with new outbound. Diverges from peer cohort.'),
    ('GRE tunnels\n(T1572)', 'Configures tunnels between routers — normal for some operators',
     'This admin never creates tunnels; this device never had tunnel activity. Direction: "protocol tunneling."'),
    ('CDR exfil\n(T1041)', 'Slow steady data movement. Under volume thresholds.',
     'Change-point accumulates gradual shift in data access. Peer comparison: no other admin has this trajectory.'),
]
y += Inches(0.45)
for technique, misses, catches in phases_a:
    rrect(s, Inches(1.2), y, Inches(11.0), Inches(0.6), fill=LTEAL, border=MGRAY)
    txt(s, Inches(1.4), y+Inches(0.05), Inches(1.8), Inches(0.5), technique, 9.5, NAVY, True)
    txt(s, Inches(3.2), y+Inches(0.05), Inches(3.3), Inches(0.5), misses, 10, RED)
    txt(s, Inches(6.6), y+Inches(0.05), Inches(5.4), Inches(0.5), catches, 10, INK)
    y += Inches(0.65)

rrect(s, Inches(1.2), Inches(6.8), Inches(11.0), Inches(0.3), fill=LGOLD)
txt(s, Inches(1.5), Inches(6.82), Inches(10.5), Inches(0.26),
    'Salt Typhoon: 200+ telcos. 4+ year dwell. CALEA access. Config-resident. Both layers stop it.',
    9.5, NAVY, True, PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════��═══════════
# SLIDE 10: COMBINED DETECTION TIMELINE
# ═══════════════════════════════════════════════════════════════════════��
s = blank()
slide_title(s, 'Combined Detection Timeline',
    'Real scenario: Volt Typhoon campaign. Preemptive finds the gap. ACECARD catches the exploitation.')

rrect(s, Inches(1.2), Inches(1.45), Inches(11.0), Inches(0.35), fill=LRED)
txt(s, Inches(1.5), Inches(1.47), Inches(10.5), Inches(0.3),
    'Scenario: Volt Typhoon targets Air Force network via Fortinet VPN CVE + LOTL credential access', 11, RED, True)

timeline = [
    ('T-30d', 'PREEMPTIVE', BLUE, 'Identifies missing IPS signature on VPN concentrator. Gap finding pushed to ticketing.'),
    ('T-30d', 'INTEGRATION', GOLD, 'Gap triggers trust floor reduction. Behavioral sensitivity auto-heightened for VPN zone.'),
    ('T+0', 'ATTACKER', RED, 'Exploits CVE. Lands on device. Creates local admin. Begins LOTL reconnaissance.'),
    ('T+2h', 'ACECARD', TEAL, 'Entity enters with reduced trust. Auth signal: new source, off-hours. Drift accumulating.'),
    ('T+6h', 'ACECARD', TEAL, 'Process signal: wmic + ntdsutil chain. Change-point accumulating. Direction: LOTL execution.'),
    ('T+8h', 'ACECARD', TEAL, 'Change-point crosses alarm. Entity BLOCKED. Access revoked. IR ticket auto-created.'),
    ('T+9h', 'INTEGRATION', GOLD, 'Confirmed detection feeds back. Preemptive re-verifies all related configs. Similar gaps closed.'),
]

y = Inches(1.9)
for i, (time, layer, color, event) in enumerate(timeline):
    bg = LGRAY if i % 2 == 0 else WHITE
    rrect(s, Inches(1.2), y, Inches(11.0), Inches(0.6), fill=bg, border=MGRAY)
    rect(s, Inches(1.25), y+Inches(0.14), Inches(0.7), Inches(0.28), fill=NAVY)
    txt(s, Inches(1.27), y+Inches(0.16), Inches(0.66), Inches(0.24), time, 8.5, WHITE, True, PP_ALIGN.CENTER)
    rect(s, Inches(2.05), y+Inches(0.14), Inches(1.5), Inches(0.28), fill=color)
    txt(s, Inches(2.07), y+Inches(0.16), Inches(1.46), Inches(0.24), layer, 8, WHITE, True, PP_ALIGN.CENTER)
    txt(s, Inches(3.7), y+Inches(0.12), Inches(8.3), Inches(0.45), event, 11, INK)
    y += Inches(0.63)

rrect(s, Inches(1.2), Inches(6.4), Inches(11.0), Inches(0.5), fill=LGOLD)
txt(s, Inches(1.5), Inches(6.43), Inches(10.5), Inches(0.4),
    'Result: 8-hour dwell (vs. 5+ years).  Preemptive found the gap.  ACECARD caught the exploitation.  Complete defense.',
    11, NAVY, True, PP_ALIGN.CENTER)


# ════════════���═══════════════════════════════════════════════════════════
# SLIDE 11: HYPERSCALER PLATFORMS + WHAT WE DELIVER
# ════════════════��═══════════════════��═══════════════════════════════════
s = blank()
slide_title(s, '2026 Hyperscaler Platforms vs. What We Deliver',
    'They provide runtime plumbing. We provide the detection intelligence agents need to act on.')

# Top: Hyperscalers (compact)
txt(s, Inches(1.2), Inches(1.55), Inches(11.0), Inches(0.3),
    'What hyperscalers shipped in 2026 — runtime & orchestration:', 12, NAVY, True)
col_w = Inches(3.4)
x_start = Inches(1.2)
# AWS
x = x_start
rrect(s, x, Inches(1.75), col_w, Inches(1.3), fill=LGRAY, border=MGRAY)
rect(s, x, Inches(1.75), col_w, Inches(0.32), fill=RGBColor(0xFF, 0x99, 0x00))
txt(s, x+Inches(0.15), Inches(1.78), Inches(3.1), Inches(0.28), '  AWS Bedrock AgentCore', 10, WHITE, True)
tf = txt(s, x+Inches(0.2), Inches(2.15), Inches(3.0), Inches(0.85),
    '•  Multi-agent + MicroVM isolation', 10, INK)
add_p(tf, '•  Security Lake (OCSF native)', 10, INK)
add_p(tf, '•  Guardrails + Policy (GA 2026)', 10, INK)
# Microsoft
x = x_start + col_w + Inches(0.25)
rrect(s, x, Inches(1.75), col_w, Inches(1.3), fill=LGRAY, border=MGRAY)
rect(s, x, Inches(1.75), col_w, Inches(0.32), fill=RGBColor(0x00, 0x78, 0xD4))
txt(s, x+Inches(0.15), Inches(1.78), Inches(3.1), Inches(0.28), '  Microsoft Foundry + Agent 365', 10, WHITE, True)
tf = txt(s, x+Inches(0.2), Inches(2.15), Inches(3.0), Inches(0.85),
    '•  Security Copilot + Entra Agent IDs', 10, INK)
add_p(tf, '•  70+ pre-built agents in store', 10, INK)
add_p(tf, '•  Purview + Defender for AI', 10, INK)
# Google
x = x_start + 2 * (col_w + Inches(0.25))
rrect(s, x, Inches(1.75), col_w, Inches(1.3), fill=LGRAY, border=MGRAY)
rect(s, x, Inches(1.75), col_w, Inches(0.32), fill=RGBColor(0x34, 0xA8, 0x53))
txt(s, x+Inches(0.15), Inches(1.78), Inches(3.1), Inches(0.28), '  Google Gemini Enterprise Agent', 10, WHITE, True)
tf = txt(s, x+Inches(0.2), Inches(2.15), Inches(3.0), Inches(0.85),
    '•  Cryptographic agent identity', 10, INK)
add_p(tf, '•  Model Armor + zero-trust per step', 10, INK)
add_p(tf, '•  ADK multi-agent + SecOps', 10, INK)

# Divider
txt(s, Inches(1.2), Inches(3.25), Inches(11.0), Inches(0.3),
    'These platforms provide PLUMBING — not intelligence. What WE provide that NO platform ships:', 12, RED, True)

# Comparison table
rrect(s, Inches(1.2), Inches(3.6), Inches(11.0), Inches(0.35), fill=NAVY)
txt(s, Inches(1.4), Inches(3.63), Inches(4.5), Inches(0.3), 'CAPABILITY', 10, WHITE, True)
txt(s, Inches(5.9), Inches(3.63), Inches(2.0), Inches(0.3), 'HYPERSCALER', 10, WHITE, True)
txt(s, Inches(8.2), Inches(3.63), Inches(3.8), Inches(0.3), '22CT DELIVERS', 10, WHITE, True)

rows = [
    ('Prove all config paths block an attack', 'Cannot do', 'Formal model — mathematical proof'),
    ('Detect sub-threshold behavioral drift', 'Cannot do', 'Change-point on Digital Twin'),
    ('Map drift direction to MITRE technique', 'Cannot do', 'Drift direction classifier'),
    ('Fuse identity across 10+ systems → one entity', 'Partial (per-cloud)', 'Cross-cloud entity fusion'),
    ('Continuous re-verification (hourly + on change)', 'Cannot do', 'Automated closed-loop'),
    ('Vendor-specific remediation with zero-FP', 'Cannot do', 'Prescriptive playbooks'),
]
y = Inches(4.0)
for cap, hyper, ours in rows:
    bg = LGRAY if rows.index((cap, hyper, ours)) % 2 == 0 else WHITE
    rrect(s, Inches(1.2), y, Inches(11.0), Inches(0.38), fill=bg, border=MGRAY)
    txt(s, Inches(1.4), y+Inches(0.05), Inches(4.4), Inches(0.28), cap, 11, INK)
    txt(s, Inches(5.9), y+Inches(0.05), Inches(2.0), Inches(0.28), hyper, 10, RED, True)
    txt(s, Inches(8.2), y+Inches(0.05), Inches(3.8), Inches(0.28), ours, 10, BLUE, True)
    y += Inches(0.4)

rrect(s, Inches(1.2), Inches(6.5), Inches(11.0), Inches(0.4), fill=LGOLD)
txt(s, Inches(1.5), Inches(6.53), Inches(10.5), Inches(0.35),
    'Hyperscalers provide runtime. 22CT provides the domain intelligence that makes agents effective.',
    11, NAVY, True, PP_ALIGN.CENTER)


# ════════════��════════════════════���══════════════════════════════════��═══
# SLIDE 12: PLATFORM ARCHITECTURE
# ════════════════════════════════════════════���═══════════════════════════
s = blank()
slide_title(s, 'Platform Architecture: Five Layers',
    'Kubernetes-native. Iron Bank hardened. Same architecture across all classification levels.')

layers = [
    ('5', 'Presentation & Integration', 'SOC Dashboard  •  ServiceNow  •  JIRA  •  SOAR  •  API Gateway  •  ABAC Policy Engine',
     RGBColor(0x2E, 0x86, 0xAB)),
    ('4', 'Analytics & Reasoning', 'Preemptive Engine  •  Behavioral Engine  •  Kill-Chain Reconstruction  •  Alert Correlation',
     RGBColor(0x1B, 0x6F, 0x92)),
    ('3', 'Intelligence & Modeling', 'Threat Intel Pipeline  •  Math Model Builder  •  Behavioral Profiling  •  Entity Fusion  •  MITRE Mapping',
     RGBColor(0x14, 0x5D, 0x82)),
    ('2', 'Data & Storage', 'Vector Database (pgvector)  •  Time-Series Store  •  Config Repository  •  Baseline Library  •  Entity Graph',
     RGBColor(0x0F, 0x3D, 0x5C)),
    ('1', 'Infrastructure', 'Kubernetes (EKS/AKS/On-Prem)  •  Iron Bank Containers  •  Istio Service Mesh  •  Vault Secrets  •  Monitoring',
     RGBColor(0x0D, 0x1B, 0x2A)),
]

y = Inches(1.5)
for num, title, components, color in layers:
    rect(s, Inches(1.2), y, Inches(11.0), Inches(0.95), fill=color)
    txt(s, Inches(1.5), y+Inches(0.1), Inches(0.4), Inches(0.35), num, 14, WHITE, False)
    txt(s, Inches(2.0), y+Inches(0.1), Inches(4.0), Inches(0.35), title, 13, WHITE, True)
    txt(s, Inches(2.0), y+Inches(0.5), Inches(10.0), Inches(0.35), components, 11, RGBColor(0xCC, 0xDD, 0xEE))
    y += Inches(1.0)

rrect(s, Inches(1.2), Inches(6.7), Inches(11.0), Inches(0.3), fill=LGOLD)
txt(s, Inches(1.5), Inches(6.72), Inches(10.5), Inches(0.26),
    'All containers Iron Bank certified  •  Helm charts for repeatable deployment  •  Horizontal pod autoscaling',
    9.5, NAVY, True, PP_ALIGN.CENTER)


# ══════════════��════════════════════════════��════════════════════════════
# SLIDE 13: DEPLOYMENT MODELS + GABRIEL NIMBUS
# ═════════════════════════════════���══════════════════════════════��═══════
s = blank()
slide_title(s, 'Deployment Models & Gabriel Nimbus',
    'One codebase, multiple targets. Environment configuration only — no code forking.')

models = [
    ('COMMERCIAL\nIL5 / NIPR', BLUE, [
        'AWS GovCloud / Azure Gov',
        'API-based config ingestion',
        'Auto-updated threat intel',
        'Fastest: days to value',
    ]),
    ('ON-PREMISES\nIL6 / SIPR', TEAL, [
        'On-site K8s cluster',
        'Classified threat feeds',
        'Local AI inference only',
        '8 vCPU / 32GB per node',
    ]),
    ('AIR-GAPPED\nJWICS / TS-SCI', NAVY, [
        'Fully disconnected',
        'Manual STIX/TAXII import',
        'Local language model',
        'SCIF-deployable',
    ]),
    ('GABRIEL NIMBUS\nCloud-Native', GOLD, [
        'Containerized microservices',
        'Platform-native stores',
        'Big data analytics layer',
        'Scales with data volume',
    ]),
]

x = Inches(0.6)
for title, color, bullets in models:
    rrect(s, x, Inches(1.5), Inches(2.95), Inches(3.0), fill=LGRAY, border=MGRAY)
    rect(s, x, Inches(1.5), Inches(2.95), Inches(0.6), fill=color)
    txt(s, x+Inches(0.1), Inches(1.55), Inches(2.75), Inches(0.55), title, 10, WHITE, True, PP_ALIGN.CENTER)
    by = Inches(2.25)
    for bullet in bullets:
        txt(s, x+Inches(0.2), by, Inches(2.6), Inches(0.28), f'•  {bullet}', 11, INK)
        by += Inches(0.35)
    x += Inches(3.15)

# Bottom: shared characteristics
rrect(s, Inches(0.6), Inches(4.75), Inches(12.1), Inches(1.8), fill=LBLUE, border=MGRAY)
tf = txt(s, Inches(0.9), Inches(4.85), Inches(11.5), Inches(1.6),
    'Shared across ALL deployment models:', 12, NAVY, True)
add_p(tf, '•  Iron Bank certified containers  •  Helm charts  •  Istio zero-trust service mesh  •  Horizontal pod autoscaling', 11, INK)
add_p(tf, '•  Same security posture regardless of environment  •  Switch via environment variable — no code forking', 11, INK)
add_p(tf, '•  cATO pathway: security artifacts auto-generated  •  NIST 800-53 / SP 800-207 compliance built in', 11, INK)
add_p(tf, '', 8)
add_p(tf, 'Gabriel Nimbus: containerized application designed for cloud-native big data platform. Integration: ingestion connectors,', 11, TEAL)
add_p(tf, 'shared storage, platform IAM, analytics layer. Same Helm charts — deploys with environment configuration only.', 11, TEAL)


# ═════════════════════��══════════════════════════════════════════════════
# SLIDE 14: IMPLEMENTATION & ROADMAP
# ═════════���════════════════════════���═════════════════════════════════════
s = blank()
slide_title(s, 'Implementation Phases & Roadmap',
    'OTA-ready. 90-day prototype timeline aligns with Army Continuous Transformation (AD 24-02) cadence.')

# Implementation phases (top half)
phases = [
    ('Phase 1: Foundation\n0–90 days', BLUE,
     'K8s provisioning • Data connectors (NGFW, IPS, IdP) • Baseline config + model build • First preemptive findings • Entity fusion pipeline',
     'First verified gap findings'),
    ('Phase 2: Operational\n90–180 days', TEAL,
     'Behavioral baselines • Change-point detection active • SOC integration (SNOW/JIRA/SOAR) • Continuous re-verification • ABAC trust scoring',
     'Full two-layer detection'),
    ('Phase 3: Optimization\n180–270 days', GOLD,
     'Kill-chain reconstruction • Cohort tuning • Additional vendors • Compliance automation • SOC handoff',
     'Full capability + cATO'),
]

x = Inches(0.6)
for title, color, items, milestone in phases:
    rrect(s, x, Inches(1.4), Inches(3.9), Inches(2.6), fill=LGRAY, border=MGRAY)
    rect(s, x, Inches(1.4), Inches(3.9), Inches(0.5), fill=color)
    txt(s, x+Inches(0.15), Inches(1.45), Inches(3.6), Inches(0.45), title, 10, WHITE, True)
    tf = txt(s, x+Inches(0.15), Inches(2.0), Inches(3.6), Inches(1.5), '', 10, INK)
    for item in items.split(' • '):
        add_p(tf, f'•  {item}', 10, INK, sp=Pt(4))
    rrect(s, x+Inches(0.1), Inches(3.55), Inches(3.7), Inches(0.35), fill=LBLUE, border=MGRAY)
    txt(s, x+Inches(0.15), Inches(3.57), Inches(3.5), Inches(0.3), f'→ {milestone}', 10, BLUE, True)
    x += Inches(4.1)

# Roadmap (bottom half)
txt(s, Inches(0.6), Inches(4.15), Inches(12.0), Inches(0.3),
    'Product Roadmap:', 12, NAVY, True)

quarters = [
    ('Q2 2026 — Current', BLUE, 'PAN/Fortinet/Cisco/CP NGFW • N-S verification • UEBA core (5 signals) • Entity fusion • IL5'),
    ('Q3 2026 — Next', TEAL, 'E-W micro-segmentation • SASE/SD-WAN • More identity sources • 140+ TTPs • IL6 certified'),
    ('Q4 2026 — Planned', GOLD, 'Cloud posture (AWS/Azure/GCP) • WAF/API gateway • Auto-remediation • Kill-chain • JWICS'),
    ('2027 — Vision', NAVY, 'OT/ICS protocols • Autonomous threat hunting • Cross-org sharing • AI-verified remediation'),
]
y = Inches(4.5)
for title, color, items in quarters:
    rrect(s, Inches(0.6), y, Inches(12.1), Inches(0.5), fill=LGRAY, border=MGRAY)
    rect(s, Inches(0.6), y, Inches(0.06), Inches(0.5), fill=color)
    txt(s, Inches(0.9), y+Inches(0.08), Inches(2.5), Inches(0.35), title, 10, color, True)
    txt(s, Inches(3.5), y+Inches(0.1), Inches(9.0), Inches(0.35), items, 10, INK)
    y += Inches(0.55)


# ══════════════════��═══════════════════════��═════════════════════════════
# SLIDE 15: USE CASES
# ══════════���═════════════════════���════════════════════════════════��══════
s = blank()
slide_title(s, 'Mission Threads',
    'Nine operational mission threads spanning DCO, compliance, and risk management for DODIN defense.')

ucs = [
    ('01', 'Preemptive Cyber Defense', 'Continuous verification that every known TTP path is blocked', BLUE),
    ('02', 'Configuration Drift Mgmt', 'Hourly comparison to approved baseline; surfaces unauthorized changes', BLUE),
    ('03', 'Formal Policy Analysis', 'Shadow rules, conflicts, weak inspection across multi-vendor stack', BLUE),
    ('04', 'External Attack Surface', 'Maps exposed assets to compensating controls; proves coverage', BLUE),
    ('05', 'Internal Attack Surface', 'Models paths from compromised host to crown-jewel assets', BLUE),
    ('06', 'Behavioral Threat Detection', 'LOTL, credential abuse, insider threat, slow lateral movement', TEAL),
    ('07', 'Zero-Trust Verification', 'Proves segmentation holds end-to-end (NIST 800-53 / SP 800-207)', TEAL),
    ('08', 'Vendor Migration Safety', 'Proves no regression during NGFW vendor migration', GOLD),
    ('09', 'M&A / Third-Party Risk', 'Models combined posture pre-merger; verifies isolation controls', GOLD),
]

y = Inches(1.5)
for num, title, desc, color in ucs:
    bg = LGRAY if int(num) % 2 != 0 else WHITE
    rrect(s, Inches(1.2), y, Inches(11.0), Inches(0.55), fill=bg, border=MGRAY)
    rect(s, Inches(1.25), y+Inches(0.1), Inches(0.5), Inches(0.3), fill=color)
    txt(s, Inches(1.3), y+Inches(0.12), Inches(0.4), Inches(0.26), num, 9.5, WHITE, True, PP_ALIGN.CENTER)
    txt(s, Inches(1.9), y+Inches(0.08), Inches(3.0), Inches(0.4), title, 12, NAVY, True)
    txt(s, Inches(5.0), y+Inches(0.1), Inches(7.0), Inches(0.4), desc, 11, INK)
    y += Inches(0.58)


# ═══════════════════════════════════════���════════════════════════════════
# SLIDE 16: PROOF POINTS
# ══════════════════════════════════���═════════════════════════════════════
s = blank()
slide_title(s, 'Production Proof Points',
    'Selected findings from financial services, telecom, government, and critical infrastructure.')

findings = [
    ('FinTech | Asia', 'Closed lateral movement path hidden by shadowed rules across multi-vendor NGFW stack', BLUE),
    ('Communications | NA', 'Eliminated RCE exposure from missing IPS inspection on critical protocol', TEAL),
    ('Regional Bank | E. Asia', 'Closed cross-vendor policy inconsistency exposing PAN and Fortinet NGFWs', GOLD),
    ('Critical Infra | SE Asia', 'Detected application obfuscation paths from DMZ rule misconfiguration', RED),
    ('Finance | E. Asia', 'Identified NIST 800-53 compliance gap on email traffic inspection', NAVY),
    ('EU Sovereign Govt', 'National-scale attack surface management across 12 government agencies', BLUE),
]

y = Inches(1.5)
for customer, finding, color in findings:
    rrect(s, Inches(1.2), y, Inches(11.0), Inches(0.72), fill=LGRAY, border=MGRAY)
    rect(s, Inches(1.2), y, Inches(0.06), Inches(0.72), fill=color)
    rect(s, Inches(1.4), y+Inches(0.1), Inches(2.2), Inches(0.28), fill=color)
    txt(s, Inches(1.45), y+Inches(0.12), Inches(2.1), Inches(0.24), customer, 9, WHITE, True, PP_ALIGN.CENTER)
    txt(s, Inches(3.8), y+Inches(0.15), Inches(8.2), Inches(0.5), finding, 12, INK)
    y += Inches(0.78)

rrect(s, Inches(1.2), Inches(6.2), Inches(11.0), Inches(0.5), fill=LBLUE)
txt(s, Inches(1.5), Inches(6.25), Inches(10.5), Inches(0.4),
    '20+ design partners across F500, federal, critical infrastructure, MSSPs, and nation states (NA, EU, ME, Asia, Pacific).',
    10, NAVY, True, PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 17: 22CT COMPANY OVERVIEW (merged About Us + Credentials)
# ═══════════���═════════════════════��═════════════════════════════════════���
s = blank()
slide_title(s, '22nd Century Technologies: Company Overview',
    'Connecting & Serving Government & Citizens since 1997.')

# About Us image (top-left)
about_path = os.path.join(IMG_DIR, "slide3_Picture_2.png")
if os.path.exists(about_path):
    s.shapes.add_picture(about_path, Inches(0.5), Inches(1.3), Inches(6.5), Inches(2.8))

# What We Do image (top-right)
whatwedo_path = os.path.join(IMG_DIR, "slide4_Picture_3.png")
if os.path.exists(whatwedo_path):
    s.shapes.add_picture(whatwedo_path, Inches(7.0), Inches(1.3), Inches(5.8), Inches(2.8))

# Metrics bar
metrics = [
    ('$750M+', 'Revenue', BLUE),
    ('$1B+', 'Award Value', TEAL),
    ('7,500+', 'Total FTEs', GOLD),
    ('800+', 'Cleared Cyber Analysts', NAVY),
    ('IL5/IL6/JWICS', 'Three Enclaves', RED),
]
x = Inches(0.5)
for num, label, color in metrics:
    rrect(s, x, Inches(4.3), Inches(2.3), Inches(0.85), fill=LGRAY, border=MGRAY)
    rect(s, x, Inches(4.3), Inches(2.3), Inches(0.04), fill=color)
    txt(s, x+Inches(0.05), Inches(4.4), Inches(2.2), Inches(0.35), num, 14, color, True, PP_ALIGN.CENTER)
    txt(s, x+Inches(0.05), Inches(4.8), Inches(2.2), Inches(0.25), label, 8, INK, al=PP_ALIGN.CENTER)
    x += Inches(2.5)

# Key contracts + compliance (bottom)
rrect(s, Inches(0.5), Inches(5.35), Inches(6.0), Inches(1.4), fill=LGRAY, border=MGRAY)
tf = txt(s, Inches(0.7), Inches(5.4), Inches(5.6), Inches(1.3),
    'Key Contracts:', 10, NAVY, True)
add_p(tf, '•  U.S. Army SOC/MDR — $90M (800+ cleared analysts)', 9, INK, sp=Pt(3))
add_p(tf, '•  USAF Cybersecurity Operations — $108M', 9, INK, sp=Pt(3))
add_p(tf, '•  FBI TSC — $56M  •  NAVAIR/USMC — $145M', 9, INK, sp=Pt(3))
add_p(tf, '•  14/15 Cabinet Agencies  •  300+ Federal  •  50+ Cyber', 9, INK, sp=Pt(3))

rrect(s, Inches(6.7), Inches(5.35), Inches(6.0), Inches(1.4), fill=LGRAY, border=MGRAY)
tf = txt(s, Inches(6.9), Inches(5.4), Inches(5.6), Inches(1.3),
    'Compliance & Certifications:', 10, NAVY, True)
add_p(tf, '•  DoD Zero Trust RA: Pillars 4 & 7  •  NIST 800-53/207', 9, INK, sp=Pt(3))
add_p(tf, '•  cATO pathway  •  Iron Bank certified  •  eMASS/RMF', 9, INK, sp=Pt(3))
add_p(tf, '•  ISO 9001 | 27001 | 28000 | 20000 | 17025 | 14001', 9, INK, sp=Pt(3))
add_p(tf, '•  GSA MAS | CIO-SP3 | SEWP V  •  CMMI Level 3', 9, INK, sp=Pt(3))

# Recognition bar
rrect(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.25), fill=LGOLD)
txt(s, Inches(0.8), Inches(6.87), Inches(11.8), Inches(0.2),
    "America's Best Midsize Employer  •  Top 15 Country for Tech People  •  Inc. 5000  •  CRN Top 10 Cyber",
    8, NAVY, True, PP_ALIGN.CENTER)


# ════════════��═══════════════════════════════════════════════════════════
# SLIDE 18: WHAT NO EXISTING TOOL DELIVERS (merged Differentiators + Uniqueness)
# ════════��═══════════════════════���═════════════════════════════���═════════
s = blank()
slide_title(s, 'What No Existing Tool, Device, or Platform Delivers',
    'These capabilities do not exist in any product on the market today — at any price point.')

capabilities = [
    ('Exhaustive Mathematical Proof of Defense',
     'No SIEM, EDR, scanner, or pentest can prove ALL config paths block an attack. Zero FP. Zero FN.',
     'Not in: Splunk, CrowdStrike, Palo Alto, Tenable, Rapid7, Qualys, any SIEM/EDR/XDR', BLUE),
    ('Digital Twin with Behavioral Time Dimension',
     'Living entity model that captures sub-threshold drift and maps direction to MITRE ATT&CK.',
     'Not in: Exabeam, Securonix, Microsoft Sentinel, Gurucul, any commercial UEBA', TEAL),
    ('Cross-Domain Entity Fusion (10+ Systems)',
     'Identity stitched across AD, AAD, Okta, VPN, TACACS+, badge, endpoint into one behavioral trajectory.',
     'Not in: Any single-vendor identity solution or SIEM correlation engine', GOLD),
    ('Management Plane: Multi-Device Verification',
     'Reasons across ALL devices simultaneously. Finds multi-rule, multi-device gaps no single tool sees.',
     'Not in: Any firewall management platform (Panorama, FortiManager, SecureX)', RED),
    ('Continuous Re-Verification (Hourly + On Change)',
     'Re-proves entire defense posture every hour AND on every config change. Pentests run annually.',
     'Not in: Any pentest, red team, ASM, or vulnerability management platform', NAVY),
    ('Hyperscaler Platforms Cannot Replace This',
     'AWS/Microsoft/Google ship runtime. They do NOT ship detection intelligence or formal models.',
     'Not in: Bedrock AgentCore, Microsoft Foundry, Google Gemini Enterprise Agent Platform', PURPLE),
]

y = Inches(1.45)
for title, desc, not_in, color in capabilities:
    rrect(s, Inches(1.2), y, Inches(11.0), Inches(0.85), fill=LGRAY, border=MGRAY)
    rect(s, Inches(1.2), y, Inches(0.08), Inches(0.85), fill=color)
    txt(s, Inches(1.5), y+Inches(0.03), Inches(10.5), Inches(0.3), title, 12, color, True)
    txt(s, Inches(1.5), y+Inches(0.3), Inches(10.5), Inches(0.28), desc, 11, INK)
    txt(s, Inches(1.5), y+Inches(0.58), Inches(10.5), Inches(0.22), not_in, 9, DGRAY)
    y += Inches(0.88)

rrect(s, Inches(1.2), Inches(6.75), Inches(11.0), Inches(0.3), fill=LRED)
txt(s, Inches(1.5), Inches(6.77), Inches(10.5), Inches(0.26),
    'Not incremental improvement. A fundamentally new class of defense that does not exist elsewhere.',
    10, RED, True, PP_ALIGN.CENTER)


# ════════════════════════════════��═══════════════════════════════��═══════
# SLIDE 19: CLOSING — PROVE IT IN 4 WEEKS
# ═══════════════��═════════════════════════���══════════════════════════════
s = blank()
s.background.fill.solid()
s.background.fill.fore_color.rgb = NAVY
rect(s, Inches(0), Inches(0), W, Inches(0.06), fill=GOLD)

txt(s, Inches(1.8), Inches(1.3), Inches(10), Inches(0.5),
    '22nd Century Technologies Inc.', 13, GOLD, True, fn='Georgia')
txt(s, Inches(1.8), Inches(2.0), Inches(10), Inches(1.0),
    'Concrete Next Steps', 36, WHITE, True, fn='Georgia')
hline(s, Inches(1.8), Inches(2.9), Inches(3.0), GOLD, Pt(2))

txt(s, Inches(1.8), Inches(3.2), Inches(10), Inches(0.4),
    'Three paths forward — all OTA/PIA compatible:', 15, RGBColor(0xA0, 0xB8, 0xD0))

# Three options as cards
options = [
    ('A', 'Technical Deep-Dive', BLUE,
     '60-90 min follow-up with engineering\nteam and ARCYBER science advisor.\nArchitecture, integration points,\nformal methods detail.'),
    ('B', 'PM DCO Capability Brief', TEAL,
     '30-min brief tailored to specific\nDCO capability gap. CFIC introduces.\nAligned to BA-08 funding and\nContinuous Transformation cadence.'),
    ('C', '30-Day OTA Prototype', GOLD,
     'Scoped prototype on ARCYBER data\n(synthetic or DREN). Preemptive gap\nfindings + behavioral detections.\nMilestone-based, TRL 6 → 7.'),
]
ox = Inches(1.5)
for letter, title, color, desc in options:
    rrect(s, ox, Inches(3.7), Inches(3.3), Inches(2.2), fill=RGBColor(0x14, 0x24, 0x34), border=color)
    rect(s, ox, Inches(3.7), Inches(3.3), Inches(0.4), fill=color)
    txt(s, ox+Inches(0.1), Inches(3.73), Inches(0.4), Inches(0.35), letter, 16, WHITE, True, PP_ALIGN.CENTER)
    txt(s, ox+Inches(0.5), Inches(3.73), Inches(2.6), Inches(0.35), title, 12, WHITE, True)
    txt(s, ox+Inches(0.2), Inches(4.2), Inches(2.9), Inches(1.6), desc, 11, RGBColor(0xA0, 0xB8, 0xD0))
    ox += Inches(3.6)

txt(s, Inches(1.8), Inches(6.1), Inches(10), Inches(0.4),
    'Provable Computing is a stated CFIC focus area. We are the industry capability in that space.',
    13, GOLD, True)

hline(s, Inches(7.5), Inches(6.5), Inches(4.5), GOLD, Pt(1))
txt(s, Inches(7.5), Inches(6.6), Inches(5.0), Inches(0.4),
    'Ravindra Shukla  |  Head of AI', 12, WHITE, True)
txt(s, Inches(7.5), Inches(6.9), Inches(5.0), Inches(0.3),
    'ravindra.shukla@tscti.com  |  703-679-0117', 10, RGBColor(0xA0, 0xB8, 0xD0))
brand_bar(s)


# ══════════════��═════════════════════════════════════════════════════════
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
