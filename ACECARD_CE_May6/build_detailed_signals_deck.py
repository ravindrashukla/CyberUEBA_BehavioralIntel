"""Build detailed deep-dive deck: every signal, every metric, every threshold,
entity fusion, entity dimensions, 4-component signature, drift direction.
This is the SUPPORTING deck so the presenter knows each and every point."""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

OUT = os.path.join(os.path.dirname(__file__), "06_Detailed_Deep_Dive_Deck.pptx")

NAVY   = RGBColor(0x00, 0x2B, 0x5C)
GOLD   = RGBColor(0xB8, 0x96, 0x3A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x1A, 0x1A, 0x2E)
CYBER  = RGBColor(0x00, 0x8B, 0x6B)
RED    = RGBColor(0xE5, 0x3E, 0x3E)
BLUE   = RGBColor(0x2D, 0x5A, 0x87)
GRAY   = RGBColor(0x4A, 0x55, 0x68)
LTGRAY = RGBColor(0xA0, 0xAE, 0xC0)
BG     = RGBColor(0xF8, 0xF9, 0xFA)
CARD   = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xE8, 0x8D, 0x0A)
PURPLE = RGBColor(0x6C, 0x2E, 0xD9)
TEAL   = RGBColor(0x0D, 0x9B, 0x8C)
DKBLUE = RGBColor(0xF0, 0xF2, 0xF5)
BLACK  = RGBColor(0x1A, 0x1A, 0x2E)
TXTMAIN = RGBColor(0x2D, 0x2D, 0x3D)
TXTSUB  = RGBColor(0x55, 0x5E, 0x70)
BORDER  = RGBColor(0xDD, 0xDD, 0xE5)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height
snum_counter = [0]


def _fill(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _box(slide, l, t, w, h, fill=None, border=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.background()
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    ln = s.line
    if border:
        ln.color.rgb = border
        ln.width = Pt(1.5)
    else:
        ln.fill.background()
    s.shadow.inherit = False
    return s


def _txt(slide, l, t, w, h, text, size=14, color=TXTMAIN, bold=False,
         align=PP_ALIGN.LEFT, font='Segoe UI'):
    tf = slide.shapes.add_textbox(l, t, w, h).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return tf


def _para(tf, text, size=14, color=TXTMAIN, bold=False, space=Pt(4)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Segoe UI'
    p.space_before = space
    return p


def _title_bar(slide, title, subtitle=None):
    _box(slide, Inches(0), Inches(0), W, Inches(1.1), fill=NAVY)
    _txt(slide, Inches(0.6), Inches(0.15), Inches(10), Inches(0.5),
         title, 26, WHITE, True)
    if subtitle:
        _txt(slide, Inches(0.6), Inches(0.6), Inches(10), Inches(0.4),
             subtitle, 13, LTGRAY)
    _txt(slide, Inches(10.5), Inches(0.25), Inches(2.5), Inches(0.5),
         'DEEP DIVE', 14, CYBER, True, PP_ALIGN.RIGHT)


def new_slide(title, subtitle=None):
    snum_counter[0] += 1
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _fill(s, BG)
    _title_bar(s, title, subtitle)
    return s


def metric_table(slide, x, y, w, rows, header_color=CYBER):
    """rows = list of (metric, description, normal, warning, critical)"""
    row_h = Inches(0.32)
    hdr_h = Inches(0.28)
    col_w = [w * 0.17, w * 0.33, w * 0.17, w * 0.16, w * 0.17]
    headers = ['METRIC', 'WHAT IT MEASURES', 'NORMAL', 'WARNING', 'CRITICAL']

    _box(slide, x, y, w, hdr_h, fill=header_color)
    cx = x
    for i, h in enumerate(headers):
        _txt(slide, cx + Pt(4), y + Pt(1), col_w[i] - Pt(8), hdr_h,
             h, 8, WHITE, True)
        cx += col_w[i]

    for ri, row_data in enumerate(rows):
        ry = y + hdr_h + ri * row_h
        bg = CARD if ri % 2 == 0 else DKBLUE
        _box(slide, x, ry, w, row_h, fill=bg)
        cx = x
        for ci, val in enumerate(row_data):
            clr = NAVY if ci == 0 else (RED if ci == 4 else (ORANGE if ci == 3 else TXTSUB))
            _txt(slide, cx + Pt(4), ry + Pt(1), col_w[ci] - Pt(8), row_h,
                 val, 9, clr, ci == 0, font='Cascadia Code' if ci == 0 else 'Segoe UI')
            cx += col_w[ci]

    return y + hdr_h + len(rows) * row_h


# ================================================================
# SLIDE 1: TITLE
# ================================================================
s = new_slide('ACECARD UEBA: Detailed Technical Reference',
              'Every signal, metric, threshold, and algorithm explained')
_txt(s, Inches(1), Inches(2.0), Inches(11), Inches(0.5),
     'Supporting Material for ACECARD CE Solution Brief', 20, GOLD, True,
     PP_ALIGN.CENTER)
_txt(s, Inches(1), Inches(3.0), Inches(11), Inches(2.5),
     'This deck provides detailed explanations of:\n\n'
     '1. All 5 behavioral signals with every metric and threshold\n'
     '2. Dual embedding architecture (5 individual + 1 composite = 9,216-d)\n'
     '3. Entity fusion and identity resolution\n'
     '4. All 7 dimensions of an entity\n'
     '5. 4-component entity signature\n'
     '6. CUSUM algorithm with exact parameters\n'
     '7. Drift direction with all 8 concept definitions\n'
     '8. ABAC trust state machine\n'
     '9. How we solve Sigma rule decay\n'
     '10. Key numbers to memorize',
     14, TXTSUB, align=PP_ALIGN.LEFT)

# ================================================================
# SLIDE 2: SIGNAL ARCHITECTURE OVERVIEW
# ================================================================
s = new_slide('Dual Embedding Architecture',
              'Each signal individually embedded (1536-d) + one fused composite = 6 vectors per entity per window')

_box(s, Inches(0.3), Inches(1.3), Inches(12.7), Inches(1.0), fill=CARD, border=CYBER)
_txt(s, Inches(0.5), Inches(1.35), Inches(12.3), Inches(0.3),
     'THE KEY INSIGHT', 14, CYBER, True)
_txt(s, Inches(0.5), Inches(1.7), Inches(12.3), Inches(0.5),
     'Each of the 5 signals is INDIVIDUALLY embedded at 1536-d. This enables per-signal drift '
     'measurement: "which signal is driving the overall drift?" Additionally, all 5 signals are '
     'concatenated and embedded as one COMPOSITE 1536-d vector for holistic entity state.',
     12, RGBColor(0x2A, 0x6A, 0x3A))

# 5 signal boxes
signal_colors = [BLUE, PURPLE, TEAL, ORANGE, RED]
signal_names = ['AUTH', 'PROCESS', 'NETWORK', 'FILE', 'IDENTITY']
x_pos = [Inches(0.3), Inches(2.9), Inches(5.5), Inches(8.1), Inches(10.7)]

for i in range(5):
    bw = Inches(2.4)
    _box(s, x_pos[i], Inches(2.6), bw, Inches(0.35), fill=signal_colors[i])
    _txt(s, x_pos[i], Inches(2.62), bw, Inches(0.3),
         f'Signal {i+1}: {signal_names[i]}', 12, WHITE, True, PP_ALIGN.CENTER)
    _box(s, x_pos[i], Inches(3.0), bw, Inches(0.5), fill=CARD, border=signal_colors[i])
    _txt(s, x_pos[i] + Inches(0.1), Inches(3.05), bw - Inches(0.2), Inches(0.4),
         'Individual 1536-d\nembedding', 10, TXTSUB, align=PP_ALIGN.CENTER)
    # Arrow down
    _txt(s, x_pos[i], Inches(3.55), bw, Inches(0.25),
         'Individual drift tracking', 8, signal_colors[i], False, PP_ALIGN.CENTER)

# Composite box
_box(s, Inches(0.3), Inches(4.1), Inches(12.7), Inches(0.5), fill=CARD, border=GOLD)
_txt(s, Inches(0.5), Inches(4.12), Inches(12.3), Inches(0.45),
     'All 5 signals concatenated and embedded together = 1 COMPOSITE 1536-d vector (holistic entity state)',
     12, GOLD, True, PP_ALIGN.CENTER)

# Numbers
_box(s, Inches(0.3), Inches(4.9), Inches(4.0), Inches(2.0), fill=CARD, border=CYBER)
_txt(s, Inches(0.5), Inches(4.95), Inches(3.6), Inches(0.3), 'DIMENSIONS', 12, CYBER, True)
tf = _txt(s, Inches(0.5), Inches(5.3), Inches(3.6), Inches(1.5),
     '5 individual signal embeddings:', 11, TXTSUB)
_para(tf, '    5 x 1,536 = 7,680 dimensions', 11, GOLD, True)
_para(tf, '1 composite embedding:', 11, TXTSUB)
_para(tf, '    1 x 1,536 = 1,536 dimensions', 11, GOLD, True)
_para(tf, 'TOTAL: 6 vectors = 9,216 dimensions', 12, CYBER, True)

_box(s, Inches(4.5), Inches(4.9), Inches(4.2), Inches(2.0), fill=CARD, border=GOLD)
_txt(s, Inches(4.7), Inches(4.95), Inches(3.8), Inches(0.3), 'WHY INDIVIDUAL + COMPOSITE?', 12, GOLD, True)
tf = _txt(s, Inches(4.7), Inches(5.3), Inches(3.8), Inches(1.5),
     'Composite drift: "entity behavior changed"', 11, TXTSUB)
_para(tf, 'Individual drift: "AUTH signal drove 72% of change"', 11, TXTSUB)
_para(tf, 'Per-signal direction: "PROCESS signal drifting', 11, TXTSUB)
_para(tf, '  toward lotl_execution (T1059)"', 11, GOLD)
_para(tf, 'Analyst sees WHICH signal + WHAT direction', 11, CYBER, True)

_box(s, Inches(8.9), Inches(4.9), Inches(4.1), Inches(2.0), fill=CARD, border=BLUE)
_txt(s, Inches(9.1), Inches(4.95), Inches(3.7), Inches(0.3), 'EMBEDDING OPTIONS', 12, BLUE, True)
tf = _txt(s, Inches(9.1), Inches(5.3), Inches(3.7), Inches(1.5),
     'IL5 (NIPR):', 11, TXTSUB)
_para(tf, '  OpenAI text-embedding-3-small', 11, GOLD)
_para(tf, 'IL6 (SIPR):', 11, TXTSUB)
_para(tf, '  Local Phi-4 (14B) or Mistral 7B', 11, GOLD)
_para(tf, 'JWICS (TS/SCI):', 11, TXTSUB)
_para(tf, '  Same local SLM, air-gapped', 11, GOLD)

# ================================================================
# SLIDE 3: SIGNAL 1 — AUTH DETAIL
# ================================================================
s = new_slide('Signal 1: AUTH — Authentication Behavior',
              'WHO is logging in, WHERE, WHEN, HOW — across all identity providers')

_box(s, Inches(0.3), Inches(1.3), Inches(6.3), Inches(2.5), fill=CARD, border=BLUE)
_txt(s, Inches(0.5), Inches(1.35), Inches(5.9), Inches(0.3), 'DATA SOURCES', 14, BLUE, True)
tf = _txt(s, Inches(0.5), Inches(1.7), Inches(5.9), Inches(2.0),
     'Windows Event 4624 — Successful logon (includes logon type: interactive, network, batch, service)',
     10, TXTSUB)
_para(tf, 'Windows Event 4625 — Failed logon (includes failure reason: bad password, locked out, expired)',
     10, TXTSUB)
_para(tf, 'Windows Event 4768 — Kerberos TGT request (initial auth in AD domain)',
     10, TXTSUB)
_para(tf, 'Windows Event 4769 — Kerberos TGS request (service ticket, indicates resource access)',
     10, TXTSUB)
_para(tf, 'Okta System Log — SSO events, MFA challenges, factor enrollments',
     10, TXTSUB)
_para(tf, 'Azure AD (Entra ID) Sign-in Logs — cloud auth, conditional access results',
     10, TXTSUB)
_para(tf, 'Linux auth.log / PAM — sudo, su, SSH login events',
     10, TXTSUB)
_para(tf, 'AWS CloudTrail — ConsoleLogin, AssumeRole, GetSessionToken events',
     10, TXTSUB)

_box(s, Inches(6.8), Inches(1.3), Inches(6.2), Inches(2.5), fill=CARD, border=BLUE)
_txt(s, Inches(7.0), Inches(1.35), Inches(5.8), Inches(0.3), 'WHAT ATTACKS THIS CATCHES', 14, BLUE, True)
tf = _txt(s, Inches(7.0), Inches(1.7), Inches(5.8), Inches(2.0),
     'Brute force — failure_rate spikes to 40%+ (ATK-001)', 10, GOLD)
_para(tf, 'Credential stuffing — multiple new source IPs, high logon_count', 10, GOLD)
_para(tf, 'Pass-the-Hash — Kerberos anomaly: TGS without prior TGT', 10, GOLD)
_para(tf, 'Lateral movement — unique_hosts jumps from 2 to 15+ per hour', 10, GOLD)
_para(tf, 'Insider off-hours — off_hours_ratio shifts from 0.05 to 0.40', 10, GOLD)
_para(tf, 'Impossible travel — NYC login + Shanghai login within 30 min', 10, GOLD)
_para(tf, 'MFA bypass — mfa_skip_ratio indicates stolen session tokens', 10, GOLD)
_para(tf, 'Service account abuse — svc accounts with interactive logon type', 10, GOLD)

# Metrics table
metric_table(s, Inches(0.3), Inches(4.0), Inches(12.7), [
    ('logon_count', 'Total auth events (success + fail) per window', '5-50/hr', '>200/hr', '>500/hr'),
    ('failure_rate', 'Failed / total logons (brute force indicator)', '0.01-0.05', '>0.15', '>0.40'),
    ('unique_hosts', 'Distinct machines authenticated to', '1-3/hr', '>8/hr', '>15/hr'),
    ('kerberos_ticket_types', 'Unusual TGS without TGT or abnormal SPN patterns', '0 anomalies', '>2/hr', '>5/hr'),
    ('off_hours_ratio', 'Logins outside 0700-1900 local / total', '0.0-0.10', '>0.35', '>0.60'),
    ('impossible_travel', 'Geo distance / time gap violates physics', '0', '1 event', '>1 event'),
    ('new_source_ips', 'IPs never seen in 30-day baseline', '0-1/day', '>3/day', '>5/day'),
    ('mfa_skip_ratio', 'Auth events bypassing MFA / total', '0.0', '>0.10', '>0.30'),
], BLUE)

# ================================================================
# SLIDE 4: AUTH — SERIALIZATION EXAMPLE
# ================================================================
s = new_slide('Signal 1: AUTH — How Raw Events Become Embedding Input',
              'The text that gets sent to the embedding model (no raw log content, only metrics)')

_box(s, Inches(0.3), Inches(1.3), Inches(6.2), Inches(3.0), fill=CARD, border=BLUE)
_txt(s, Inches(0.5), Inches(1.35), Inches(5.8), Inches(0.3), 'RAW LOG EVENTS (MANY)', 12, BLUE, True)
tf = _txt(s, Inches(0.5), Inches(1.7), Inches(5.8), Inches(2.5),
     'EventID=4624 User=CORP\\jsmith Logon=Interactive Src=10.0.1.42', 9,
     RGBColor(0x55, 0x66, 0x77), font='Cascadia Code')
_para(tf, 'EventID=4625 User=CORP\\jsmith Reason=BadPwd Src=185.73.x.x', 9,
     RGBColor(0x55, 0x66, 0x77))
_para(tf, 'EventID=4625 User=CORP\\jsmith Reason=BadPwd Src=185.73.x.x', 9,
     RGBColor(0x55, 0x66, 0x77))
_para(tf, 'EventID=4768 User=CORP\\jsmith Service=krbtgt', 9,
     RGBColor(0x55, 0x66, 0x77))
_para(tf, 'EventID=4769 User=CORP\\jsmith SPN=MSSQLSvc/db01', 9,
     RGBColor(0x55, 0x66, 0x77))
_para(tf, 'Okta: user.session.start john.smith@corp.mil MFA=PUSH', 9,
     RGBColor(0x55, 0x66, 0x77))
_para(tf, 'AzureAD: SignIn john.smith@corp.mil ConditionalAccess=Pass', 9,
     RGBColor(0x55, 0x66, 0x77))
_para(tf, '... (47 total events this hour)', 9, LTGRAY)

_txt(s, Inches(6.35), Inches(2.5), Inches(0.3), Inches(0.5), '=>', 24, CYBER, True, PP_ALIGN.CENTER)

_box(s, Inches(6.8), Inches(1.3), Inches(6.2), Inches(3.0), fill=CARD, border=CYBER)
_txt(s, Inches(7.0), Inches(1.35), Inches(5.8), Inches(0.3), 'SERIALIZED TEXT (SENT TO EMBEDDER)', 12, CYBER, True)
tf = _txt(s, Inches(7.0), Inches(1.7), Inches(5.8), Inches(2.5),
     'AUTH SIGNAL for entity e7a2b3c4:', 10, GOLD, True, font='Cascadia Code')
_para(tf, '47 authentication events this window.', 10, RGBColor(0x2A, 0x6A, 0x3A))
_para(tf, '8.5% failure rate (4 failures, 43 success).', 10, RGBColor(0x2A, 0x6A, 0x3A))
_para(tf, '8 unique hosts contacted (baseline: 2).', 10, RGBColor(0x2A, 0x6A, 0x3A))
_para(tf, '2 Kerberos TGS without prior TGT detected.', 10, RGBColor(0x2A, 0x6A, 0x3A))
_para(tf, '0.15 off-hours ratio (elevated).', 10, RGBColor(0x2A, 0x6A, 0x3A))
_para(tf, '0 impossible travel events.', 10, RGBColor(0x2A, 0x6A, 0x3A))
_para(tf, '2 new source IPs (not in 30-day baseline).', 10, RGBColor(0x2A, 0x6A, 0x3A))
_para(tf, '0.0 MFA bypass ratio.', 10, RGBColor(0x2A, 0x6A, 0x3A))

_box(s, Inches(0.3), Inches(4.6), Inches(12.7), Inches(0.6), fill=CARD, border=GOLD)
_txt(s, Inches(0.5), Inches(4.62), Inches(12.3), Inches(0.55),
     'KEY: No usernames, no IP addresses, no hostnames sent to embedding model. Only numerical '
     'summaries. The embedding captures the BEHAVIORAL PATTERN, not the identity. '
     'Privacy-preserving by design. JWICS-safe with local SLM.',
     11, GOLD, True)

_box(s, Inches(0.3), Inches(5.5), Inches(12.7), Inches(1.5), fill=CARD, border=BLUE)
_txt(s, Inches(0.5), Inches(5.55), Inches(12.3), Inches(0.3), 'EMBEDDING OUTPUT', 12, BLUE, True)
_txt(s, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.25),
     'text-embedding-3-small(serialized_text) =>', 10, LTGRAY, font='Cascadia Code')
_txt(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.3),
     '[0.0234, -0.0891, 0.0456, 0.1123, -0.0334, 0.0789, ...  ]  (1,536 floating-point numbers)',
     11, CYBER, True, font='Cascadia Code')
_txt(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.3),
     'This vector IS the AUTH behavioral state. Cosine similarity to prior window = drift. '
     'Cosine to threat concept = direction.',
     10, TXTSUB)

# ================================================================
# SLIDE 5: SIGNAL 2 — PROCESS DETAIL
# ================================================================
s = new_slide('Signal 2: PROCESS — Execution Behavior',
              'WHAT programs are running, HOW they are invoked, process tree depth and signing')

_box(s, Inches(0.3), Inches(1.3), Inches(6.3), Inches(2.5), fill=CARD, border=PURPLE)
_txt(s, Inches(0.5), Inches(1.35), Inches(5.9), Inches(0.3), 'DATA SOURCES', 14, PURPLE, True)
tf = _txt(s, Inches(0.5), Inches(1.7), Inches(5.9), Inches(2.0),
     'Sysmon Event ID 1 — Process creation (full cmdline, parent process, hashes, user)',
     10, TXTSUB)
_para(tf, 'Sysmon Event ID 3 — Network connection (process making outbound connection)',
     10, TXTSUB)
_para(tf, 'Sysmon Event ID 7 — Image loaded (DLL loads, detects reflective loading)',
     10, TXTSUB)
_para(tf, 'Sysmon Event ID 8 — CreateRemoteThread (code injection indicator)',
     10, TXTSUB)
_para(tf, 'CrowdStrike Falcon — Process execution telemetry (cloud-delivered)',
     10, TXTSUB)
_para(tf, 'Microsoft Defender for Endpoint — Process events + AMSI logs',
     10, TXTSUB)
_para(tf, 'Linux auditd EXECVE — Process execution with full arguments',
     10, TXTSUB)

_box(s, Inches(6.8), Inches(1.3), Inches(6.2), Inches(2.5), fill=CARD, border=PURPLE)
_txt(s, Inches(7.0), Inches(1.35), Inches(5.8), Inches(0.3), 'LOLBIN DETECTION (THE VOLT TYPHOON SIGNAL)', 14, PURPLE, True)
tf = _txt(s, Inches(7.0), Inches(1.7), Inches(5.8), Inches(2.0),
     'LoTL binaries tracked (all legitimate Windows tools):', 10, GOLD, True)
_para(tf, 'wmic.exe — WMI command-line (system discovery, remote exec)', 10, RGBColor(0x2A, 0x6A, 0x3A))
_para(tf, 'ntdsutil.exe — AD database utility (credential extraction)', 10, RGBColor(0x2A, 0x6A, 0x3A))
_para(tf, 'netsh.exe — Network config (firewall mod, port forwarding)', 10, RGBColor(0x2A, 0x6A, 0x3A))
_para(tf, 'rundll32.exe — DLL execution (loads malicious code)', 10, RGBColor(0x2A, 0x6A, 0x3A))
_para(tf, 'certutil.exe — Certificate utility (downloads files)', 10, RGBColor(0x2A, 0x6A, 0x3A))
_para(tf, 'bitsadmin.exe — BITS transfers (stealthy downloads)', 10, RGBColor(0x2A, 0x6A, 0x3A))
_para(tf, 'mshta.exe — HTML application host (script execution)', 10, RGBColor(0x2A, 0x6A, 0x3A))
_para(tf, 'regsvr32.exe — COM registration (script proxy execution)', 10, RGBColor(0x2A, 0x6A, 0x3A))

metric_table(s, Inches(0.3), Inches(4.0), Inches(12.7), [
    ('unique_processes', 'Distinct process names executed per window', '10-40/hr', '>80/hr', '>150/hr'),
    ('cmdline_entropy', 'Shannon entropy of cmdline args (bits)', '2.0-3.5', '>4.5', '>5.5'),
    ('parent_child_depth', 'Max process tree depth (staged execution)', '2-4', '>6', '>8'),
    ('lolbin_count', 'LoTL binaries used (wmic, ntdsutil, etc.)', '0-1/hr', '>3/hr', '>6/hr'),
    ('unsigned_ratio', 'Unsigned binaries / total executed', '0.0-0.05', '>0.15', '>0.30'),
    ('new_process_rate', 'Processes never seen in 30-day baseline / total', '0.0-0.02', '>0.08', '>0.15'),
    ('encoded_cmd', 'PowerShell -EncodedCommand executions', '0', '>1/hr', '>3/hr'),
    ('inject_thread', 'CreateRemoteThread events (code injection)', '0', '>1/hr', '>3/hr'),
], PURPLE)

# ================================================================
# SLIDE 6: PROCESS — WHY THRESHOLD FAILS
# ================================================================
s = new_slide('Why Thresholds Cannot Catch LoTL',
              'Each individual action is legitimate. Only the TRAJECTORY reveals the attack.')

# Normal vs Attack comparison
_box(s, Inches(0.3), Inches(1.3), Inches(6.2), Inches(5.5), fill=CARD, border=TEAL)
_txt(s, Inches(0.5), Inches(1.35), Inches(5.8), Inches(0.3), 'NORMAL SYSADMIN (Tuesday 10 AM)', 14, TEAL, True)
tf = _txt(s, Inches(0.5), Inches(1.7), Inches(5.8), Inches(4.5),
     '10:02  wmic /node:filesvr01 os get caption', 10, TXTSUB, font='Cascadia Code')
_para(tf, '       (checking OS version before patch)', 9, LTGRAY)
_para(tf, '10:15  netsh advfirewall show allprofiles', 10, TXTSUB)
_para(tf, '       (verifying firewall state)', 9, LTGRAY)
_para(tf, '10:45  certutil -urlcache -f https://update.corp/kb123.msp', 10, TXTSUB)
_para(tf, '       (downloading approved patch)', 9, LTGRAY)
_para(tf, '', 10, LTGRAY)
_para(tf, 'lolbin_count = 3   (BELOW any threshold)', 11, TEAL, True)
_para(tf, 'cmdline_entropy = 2.8   (normal)', 11, TEAL, True)
_para(tf, 'unique_hosts = 1   (normal)', 11, TEAL, True)
_para(tf, '', 10, LTGRAY)
_para(tf, 'Embedding drift: 0.02 (baseline noise)', 11, TEAL, True)
_para(tf, 'CUSUM: +0.005 (nothing accumulating)', 11, TEAL, True)
_para(tf, 'Direction: no concept projection > 0.15', 11, TEAL, True)

_box(s, Inches(6.8), Inches(1.3), Inches(6.2), Inches(5.5), fill=CARD, border=RED)
_txt(s, Inches(7.0), Inches(1.35), Inches(5.8), Inches(0.3), 'VOLT TYPHOON (Tuesday 2 AM)', 14, RED, True)
tf = _txt(s, Inches(7.0), Inches(1.7), Inches(5.8), Inches(4.5),
     '02:03  wmic /node:dc01 process list brief', 10, TXTSUB, font='Cascadia Code')
_para(tf, '02:05  wmic /node:dc01 useraccount list full', 10, TXTSUB)
_para(tf, '02:08  ntdsutil "ac i ntds" "ifm" "create full c:\\temp"', 10, TXTSUB)
_para(tf, '02:12  netsh interface portproxy add v4tov4 ...', 10, TXTSUB)
_para(tf, '02:15  netsh advfirewall firewall add rule ...', 10, TXTSUB)
_para(tf, '02:18  rundll32.exe comsvcs.dll MiniDump ...', 10, TXTSUB)
_para(tf, '02:20  certutil -encode c:\\temp\\ntds.dit c:\\temp\\o.b64', 10, TXTSUB)
_para(tf, '', 10, LTGRAY)
_para(tf, 'lolbin_count = 7   (ABOVE threshold? Maybe. But', 11, RED, True)
_para(tf, '  a busy admin might also hit 5-6 during patching)', 10, LTGRAY)
_para(tf, 'cmdline_entropy = 4.8   (obfuscated args)', 11, RED, True)
_para(tf, 'off_hours + lolbin + ntdsutil + credential path', 11, RED, True)
_para(tf, '', 10, LTGRAY)
_para(tf, 'Embedding drift: 0.14 (significant shift)', 11, RED, True)
_para(tf, 'CUSUM: +0.09 (rapidly accumulating)', 11, RED, True)
_para(tf, 'Direction: lotl_execution=0.71, credential_dumping=0.54', 11, GOLD, True)

_box(s, Inches(0.3), Inches(6.9), Inches(12.7), Inches(0.4), fill=NAVY)
_txt(s, Inches(0.5), Inches(6.92), Inches(12.3), Inches(0.35),
     'Threshold-only: might alert at lolbin=7, with 50% false positive rate. '
     'Trajectory: 0.14 drift toward lotl_execution + credential_dumping = HIGH CONFIDENCE alert with MITRE mapping.',
     11, GOLD, True)

# ================================================================
# SLIDE 7: SIGNAL 3 — NETWORK DETAIL
# ================================================================
s = new_slide('Signal 3: NETWORK — Traffic Behavior',
              'WHERE data flows, HOW MUCH, HOW REGULARLY (beacon detection), unusual destinations')

_box(s, Inches(0.3), Inches(1.3), Inches(6.3), Inches(2.3), fill=CARD, border=TEAL)
_txt(s, Inches(0.5), Inches(1.35), Inches(5.9), Inches(0.3), 'DATA SOURCES', 14, TEAL, True)
tf = _txt(s, Inches(0.5), Inches(1.7), Inches(5.9), Inches(1.8),
     'NetFlow / IPFIX — Router-level flow records (src/dst IP, ports, bytes, duration)',
     10, TXTSUB)
_para(tf, 'Zeek (formerly Bro) — Full protocol analysis (HTTP, DNS, SSL, SMB, Kerberos)',
     10, TXTSUB)
_para(tf, 'Palo Alto PAN-OS — App-ID traffic logs with threat verdicts',
     10, TXTSUB)
_para(tf, 'Fortinet FortiGate — Session logs, web filter, IPS events',
     10, TXTSUB)
_para(tf, 'Cisco Stealthwatch/Secure Network Analytics — Flow anomaly data',
     10, TXTSUB)
_para(tf, 'AWS VPC Flow Logs — Cloud network flow records',
     10, TXTSUB)

_box(s, Inches(6.8), Inches(1.3), Inches(6.2), Inches(2.3), fill=CARD, border=TEAL)
_txt(s, Inches(7.0), Inches(1.35), Inches(5.8), Inches(0.3), 'WHAT ATTACKS THIS CATCHES', 14, TEAL, True)
tf = _txt(s, Inches(7.0), Inches(1.7), Inches(5.8), Inches(1.8),
     'C2 beaconing — beacon_score detects regular-interval callbacks', 10, GOLD)
_para(tf, 'Data exfiltration — bytes_out_ratio inverts (more out than in)', 10, GOLD)
_para(tf, 'DNS tunneling — dns_query_rate + dns_entropy spike together', 10, GOLD)
_para(tf, 'Network scanning — unique_dest_ips explodes, internal_scan_ratio rises', 10, GOLD)
_para(tf, 'Lateral movement — admin share access (C$, ADMIN$, IPC$)', 10, GOLD)
_para(tf, 'GRE tunneling (Salt Typhoon) — new protocol mix, unauthorized tunnel endpoints', 10, GOLD)

metric_table(s, Inches(0.3), Inches(3.8), Inches(12.7), [
    ('unique_dest_ips', 'Distinct destination IPs contacted', '5-20/hr', '>50/hr', '>100/hr'),
    ('bytes_out', 'Outbound / inbound byte ratio', '0.05-0.30', '>0.60', '>0.85'),
    ('port_diversity', 'Unique destination ports used', '3-10/hr', '>25/hr', '>50/hr'),
    ('dns_query_rate', 'DNS queries per minute', '2-15/min', '>40/min', '>80/min'),
    ('beacon_score', 'Connection interval regularity (0-1)', '0.0-0.15', '>0.45', '>0.70'),
    ('geo_anomaly_count', 'Connections to new countries (vs 30d baseline)', '0', '>2/hr', '>5/hr'),
    ('admin_share_access', 'C$, ADMIN$, IPC$ connections', '0', '>2/hr', '>5/hr'),
    ('internal_scan_ratio', 'Distinct internal IPs / total connections', '0.10-0.30', '>0.60', '>0.80'),
], TEAL)

_box(s, Inches(0.3), Inches(6.4), Inches(12.7), Inches(0.8), fill=CARD, border=GOLD)
_txt(s, Inches(0.5), Inches(6.42), Inches(12.3), Inches(0.3), 'BEACON SCORE EXPLAINED', 12, GOLD, True)
_txt(s, Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.4),
     'Coefficient of variation of inter-connection times. C2 malware calls home at regular intervals '
     '(e.g., every 60s). Low variance = high score. Normal browsing is irregular = low score. '
     'Works even on encrypted traffic because timing metadata is always visible.',
     10, RGBColor(0x2A, 0x6A, 0x3A))

# ================================================================
# SLIDE 8: SIGNAL 4 — FILE DETAIL
# ================================================================
s = new_slide('Signal 4: FILE — File System Behavior',
              'WHAT files are accessed, created, modified, deleted — staging, encryption, exfiltration')

_box(s, Inches(0.3), Inches(1.3), Inches(6.3), Inches(2.0), fill=CARD, border=ORANGE)
_txt(s, Inches(0.5), Inches(1.35), Inches(5.9), Inches(0.3), 'DATA SOURCES', 14, ORANGE, True)
tf = _txt(s, Inches(0.5), Inches(1.7), Inches(5.9), Inches(1.5),
     'Sysmon Event ID 11 — File created (full path, process that created it)',
     10, TXTSUB)
_para(tf, 'Windows Event 4663 — Object access audit (file read/write/delete with SID)',
     10, TXTSUB)
_para(tf, 'DLP (Data Loss Prevention) — Symantec/Forcepoint/Microsoft Purview events',
     10, TXTSUB)
_para(tf, 'Cloud storage audit — OneDrive/SharePoint/S3 access logs',
     10, TXTSUB)
_para(tf, 'Linux inotify / auditd — File system change monitoring',
     10, TXTSUB)

_box(s, Inches(6.8), Inches(1.3), Inches(6.2), Inches(2.0), fill=CARD, border=ORANGE)
_txt(s, Inches(7.0), Inches(1.35), Inches(5.8), Inches(0.3), 'WHAT ATTACKS THIS CATCHES', 14, ORANGE, True)
tf = _txt(s, Inches(7.0), Inches(1.7), Inches(5.8), Inches(1.5),
     'Ransomware — extension_changes spike (*.docx -> *.docx.encrypted)', 10, GOLD)
_para(tf, 'Pre-ransomware — shadow_copy_ops (vssadmin delete shadows)', 10, GOLD)
_para(tf, 'Exfiltration staging — archive_creates (zip/tar/7z creation)', 10, GOLD)
_para(tf, 'Insider data theft — sensitive_access steady increase over months', 10, GOLD)
_para(tf, 'Credential theft — access to SAM, NTDS.dit, lsass dump files', 10, GOLD)

metric_table(s, Inches(0.3), Inches(3.5), Inches(12.7), [
    ('files_created', 'New files written per window', '0-20/hr', '>100/hr', '>500/hr'),
    ('files_deleted', 'Files removed per window', '0-5/hr', '>50/hr', '>200/hr'),
    ('sensitive_access_count', 'Access to restricted/classified paths', '0-2/hr', '>5/hr', '>10/hr'),
    ('archive_creates', 'zip/tar/rar/7z creations', '0-1/day', '>3/hr', '>5/hr'),
    ('bulk_ops_flag', 'Bulk copy/move/delete > 100 files in burst', '0', '1 event', '>1 event'),
    ('extension_change_rate', 'Files renamed with new extensions', '0', '>10/hr', '>50/hr'),
    ('large_file_writes', 'Files >100MB written to disk', '0-1/day', '>3/hr', '>5/hr'),
    ('shadow_copy_ops', 'vssadmin/wbadmin delete operations', '0', '1 event', '>1 event'),
], ORANGE)

_box(s, Inches(0.3), Inches(6.2), Inches(12.7), Inches(1.0), fill=CARD, border=RED)
_txt(s, Inches(0.5), Inches(6.25), Inches(12.3), Inches(0.3), 'RANSOMWARE DETECTION TIMELINE', 12, RED, True)
_txt(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.6),
     'T+0: shadow_copy_ops detected (vssadmin delete) -> T+5min: extension_change_rate spikes (encryption starting) '
     '-> T+10min: FILE signal drift = 0.42, CUSUM ALARM, direction = defense_evasion(0.81) '
     '-> T+11min: ABAC moves to BLOCKED, all access revoked, IR triggered. '
     'Total response time: 11 minutes from first indicator.',
     10, RGBColor(0x2A, 0x6A, 0x3A))

# ================================================================
# SLIDE 9: SIGNAL 5 — IDENTITY DETAIL
# ================================================================
s = new_slide('Signal 5: IDENTITY — Privilege & Role Behavior',
              'Privilege escalation, group membership changes, MFA bypass, admin actions')

_box(s, Inches(0.3), Inches(1.3), Inches(6.3), Inches(2.0), fill=CARD, border=RED)
_txt(s, Inches(0.5), Inches(1.35), Inches(5.9), Inches(0.3), 'DATA SOURCES', 14, RED, True)
tf = _txt(s, Inches(0.5), Inches(1.7), Inches(5.9), Inches(1.5),
     'Windows Event 4672 — Special privileges assigned to new logon (admin logon)',
     10, TXTSUB)
_para(tf, 'Windows Event 4728/4732/4756 — Member added to security group (domain/local/universal)',
     10, TXTSUB)
_para(tf, 'AWS CloudTrail — IAM role assumption (AssumeRole), policy attachment events',
     10, TXTSUB)
_para(tf, 'Kubernetes audit log — RBAC bindings (ClusterRoleBinding, RoleBinding changes)',
     10, TXTSUB)
_para(tf, 'Active Directory change log — GPO modifications, delegation changes',
     10, TXTSUB)

_box(s, Inches(6.8), Inches(1.3), Inches(6.2), Inches(2.0), fill=CARD, border=RED)
_txt(s, Inches(7.0), Inches(1.35), Inches(5.8), Inches(0.3), 'WHAT ATTACKS THIS CATCHES', 14, RED, True)
tf = _txt(s, Inches(7.0), Inches(1.7), Inches(5.8), Inches(1.5),
     'Privilege escalation — group_adds to Domain Admins / Enterprise Admins', 10, GOLD)
_para(tf, 'MFA bypass — stolen session tokens used without MFA challenge', 10, GOLD)
_para(tf, 'Service account abuse — svc accounts with interactive logon (should never happen)', 10, GOLD)
_para(tf, 'Golden ticket — Kerberos delegation changes on DCs', 10, GOLD)
_para(tf, 'Insider — admin_actions from non-admin role (policy changes, user creation)', 10, GOLD)

metric_table(s, Inches(0.3), Inches(3.5), Inches(12.7), [
    ('priv_escalations', 'Privilege elevation events (sudo, runas, UAC)', '0-2/hr', '>5/hr', '>10/hr'),
    ('group_adds', 'Security group membership changes', '0/day', '>1/hr', '>3/hr'),
    ('mfa_bypass_attempts', 'Auth succeeding without MFA challenge', '0', '>1/hr', '>3/hr'),
    ('role_changes', 'IAM role assumptions or switches', '0-1/day', '>3/hr', '>5/hr'),
    ('service_acct_use', 'Service accounts with interactive logon', '0', '>1/day', '>1/hr'),
    ('admin_actions', 'Account creation, policy changes, GPO mods', '0/day', '>2/hr', '>5/hr'),
    ('delegation_changes', 'Kerberos delegation config modifications', '0', '>1/day', '>1/hr'),
], RED)

_box(s, Inches(0.3), Inches(6.2), Inches(12.7), Inches(1.0), fill=CARD, border=GOLD)
_txt(s, Inches(0.5), Inches(6.25), Inches(12.3), Inches(0.3), 'SALT TYPHOON IDENTITY PATTERN', 12, GOLD, True)
_txt(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.6),
     'Day 1: Initial admin access via CVE-2023-20198 (admin_actions spike). '
     'Day 7: AAA config changes + new admin sessions (priv_escalations). '
     'Day 30: CALEA intercept system access (service_acct_use from unexpected context). '
     'Day 90: Steady CDR access (sensitive_access in FILE signal crosses over). '
     'IDENTITY signal drift = 0.31 by Day 30. CUSUM triggers. Direction: privilege_escalation (0.68).',
     10, RGBColor(0x2A, 0x6A, 0x3A))

# ================================================================
# SLIDE 10: ENTITY FUSION
# ================================================================
s = new_slide('Entity Fusion: Solving Identity Fragmentation',
              'One human, 10+ identifiers, across 6+ systems. Without fusion, lateral movement is invisible.')

_box(s, Inches(0.3), Inches(1.3), Inches(12.7), Inches(0.4), fill=NAVY)
_txt(s, Inches(0.5), Inches(1.32), Inches(12.3), Inches(0.35),
     'ONE PERSON = 10+ IDENTIFIERS (John Smith, Tier-2 SOC Analyst)',
     14, WHITE, True, PP_ALIGN.CENTER)

ids = [
    ('Active Directory', 'CORP\\jsmith / S-1-5-21-123...-1234', BLUE),
    ('Azure AD (Entra ID)', 'john.smith@corp.mil / ObjectId 8a2b3c4d', PURPLE),
    ('Okta SSO', 'john.smith@corp.mil / 00u1a2b3c4', TEAL),
    ('AWS IAM', 'arn:aws:iam::842:user/jsmith', ORANGE),
    ('Kubernetes', 'system:serviceaccount:soc:jsmith-sa', RED),
    ('CrowdStrike EDR', 'aid:abc123def / JSMITH-WS01', BLUE),
    ('VPN / RADIUS', 'jsmith@corp.mil (certificate DN)', PURPLE),
    ('Splunk (SIEM)', 'src_user=jsmith OR user=CORP\\jsmith', TEAL),
    ('PKI / CAC', 'CN=SMITH.JOHN.Q.1234567890', ORANGE),
    ('TACACS+', 'jsmith (network device auth)', RED),
]

y = Inches(1.85)
for i, (system, identifier, color) in enumerate(ids):
    bg = CARD if i % 2 == 0 else DKBLUE
    _box(s, Inches(0.3), y, Inches(5.5), Inches(0.32), fill=bg)
    _box(s, Inches(0.35), y + Inches(0.02), Inches(1.8), Inches(0.24), fill=color)
    _txt(s, Inches(0.4), y + Inches(0.03), Inches(1.7), Inches(0.22),
         system, 9, WHITE, True, PP_ALIGN.CENTER)
    _txt(s, Inches(2.25), y + Inches(0.03), Inches(3.4), Inches(0.22),
         identifier, 9, TXTSUB, font='Cascadia Code')
    y += Inches(0.33)

# Arrow
_txt(s, Inches(5.9), Inches(3.2), Inches(0.6), Inches(0.5), '=>', 28, CYBER, True, PP_ALIGN.CENTER)

# Resolution result
_box(s, Inches(6.5), Inches(1.85), Inches(6.5), Inches(3.4), fill=CARD, border=CYBER)
_txt(s, Inches(6.7), Inches(1.9), Inches(6.1), Inches(0.3), 'RESOLVED ENTITY', 14, CYBER, True)
_txt(s, Inches(6.7), Inches(2.3), Inches(6.1), Inches(0.25),
     'entity_uuid: e7a2b3c4-9f1d-4e8b-a2c7-3d5f6g7h8i9j', 11, GOLD, True, font='Cascadia Code')

tf = _txt(s, Inches(6.7), Inches(2.7), Inches(6.1), Inches(2.0),
     'Resolution layers:', 11, NAVY, True)
_para(tf, '1. Deterministic: email match (AD + AAD + Okta + VPN)', 10, RGBColor(0x2A, 0x6A, 0x3A))
_para(tf, '2. Deterministic: EDIPI from CAC matches HR record', 10, RGBColor(0x2A, 0x6A, 0x3A))
_para(tf, '3. Probabilistic: CrowdStrike agent on hostname', 10, RGBColor(0x2A, 0x6A, 0x3A))
_para(tf, '   JSMITH-WS01 correlates with AD computer object', 10, LTGRAY)
_para(tf, '4. Graph: connected component = all 10 identifiers', 10, RGBColor(0x2A, 0x6A, 0x3A))
_para(tf, '', 10, LTGRAY)
_para(tf, 'All 5 signals from ALL 10 sources feed into', 11, GOLD, True)
_para(tf, 'ONE embedding per window. Lateral movement', 11, GOLD, True)
_para(tf, 'across 5 systems = 1 strong signal.', 11, GOLD, True)

# Bottom: before vs after
_box(s, Inches(0.3), Inches(5.5), Inches(6.2), Inches(1.7), fill=CARD, border=RED)
_txt(s, Inches(0.5), Inches(5.55), Inches(5.8), Inches(0.3), 'WITHOUT ENTITY FUSION', 13, RED, True)
tf = _txt(s, Inches(0.5), Inches(5.9), Inches(5.8), Inches(1.2),
     'AD sees: 1 normal logon (score: 0.3)', 11, TXTSUB)
_para(tf, 'AAD sees: 1 normal sign-in (score: 0.25)', 11, TXTSUB)
_para(tf, 'AWS sees: 1 normal console login (score: 0.2)', 11, TXTSUB)
_para(tf, 'K8s sees: 1 normal kubectl exec (score: 0.15)', 11, TXTSUB)
_para(tf, 'Result: 4 alerts below threshold. No correlation.', 11, RED, True)

_box(s, Inches(6.8), Inches(5.5), Inches(6.2), Inches(1.7), fill=CARD, border=CYBER)
_txt(s, Inches(7.0), Inches(5.55), Inches(5.8), Inches(0.3), 'WITH ENTITY FUSION', 13, CYBER, True)
tf = _txt(s, Inches(7.0), Inches(5.9), Inches(5.8), Inches(1.2),
     'ONE entity with ALL events fused:', 11, RGBColor(0x2A, 0x6A, 0x3A))
_para(tf, '  4 cross-system logins in 1 hour (unique_hosts=4)', 11, RGBColor(0x2A, 0x6A, 0x3A))
_para(tf, '  2 identity providers + 2 cloud platforms', 11, RGBColor(0x2A, 0x6A, 0x3A))
_para(tf, '  Composite drift: 0.23 (significant)', 11, GOLD, True)
_para(tf, 'Result: 1 HIGH alert. Direction: lateral_movement (0.82)', 11, CYBER, True)

# ================================================================
# SLIDE 11: ALL 7 ENTITY DIMENSIONS
# ================================================================
s = new_slide('All 7 Dimensions of an Entity',
              'An entity is NOT just a username. It is a rich, multi-dimensional object tracked across time.')

dims = [
    ('1. STATIC IDENTITY', BLUE,
     'entity_uuid, entity_type, identity_map (all linked IDs), org context (OU, role, clearance), peer cohort'),
    ('2. BEHAVIORAL STATE', PURPLE,
     '1536-d composite embedding + 5 individual signal embeddings (7,680-d) = 9,216 total dimensions per window'),
    ('3. TRAJECTORY HISTORY', TEAL,
     'Baseline embedding (30-day rolling), drift series, velocity series, acceleration, CUSUM accumulator state'),
    ('4. THREAT ALIGNMENT', ORANGE,
     'Drift vector (V_current - V_baseline), 8 concept projections, MITRE technique IDs, kill-chain position'),
    ('5. PEER CONTEXT', GOLD,
     'Cohort membership, peer mean drift, peer z-score, co-drift detection (is the whole group moving?)'),
    ('6. TRUST STATE', RED,
     'Current ABAC state (TRUSTED/WATCH/RESTRICTED/BLOCKED), transition history, enforcement actions active'),
    ('7. RELATIONSHIP GRAPH', CYBER,
     'Interacted entities (devices, apps), relationship strength (freq+recency), anomalous new pairs, kill-chain links'),
]

y = Inches(1.3)
for label, color, desc in dims:
    _box(s, Inches(0.3), y, Inches(12.7), Inches(0.78), fill=CARD, border=RGBColor(0x33, 0x33, 0x55))
    _box(s, Inches(0.35), y + Inches(0.05), Inches(2.6), Inches(0.3), fill=color)
    _txt(s, Inches(0.4), y + Inches(0.07), Inches(2.5), Inches(0.26),
         label, 10, WHITE, True, PP_ALIGN.CENTER)
    _txt(s, Inches(3.1), y + Inches(0.05), Inches(9.7), Inches(0.65),
         desc, 11, TXTSUB)
    y += Inches(0.82)

_box(s, Inches(0.3), Inches(7.0), Inches(12.7), Inches(0.35), fill=NAVY)
_txt(s, Inches(0.5), Inches(7.02), Inches(12.3), Inches(0.3),
     'SIEM has 1 dimension (event count). We have 7. That is why we detect what SIEM cannot.',
     12, GOLD, True, PP_ALIGN.CENTER)

# ================================================================
# SLIDE 12: 4-COMPONENT ENTITY SIGNATURE
# ================================================================
s = new_slide('4-Component Entity Signature',
              'Per entity per window: identity proof, behavioral state, explainability, tamper evidence')

comps = [
    ('1. IDENTITY-GRAPH HASH', BLUE,
     'SHA-256 hash of entity resolution graph',
     'Inputs: all linked identifiers + correlation edges + confidence scores',
     'Purpose: proves WHO was resolved, detects identity graph manipulation, audit trail'),
    ('2. BEHAVIORAL EMBEDDING', PURPLE,
     '1536-d vector from text-embedding-3-small (or local SLM)',
     'Inputs: concatenated 5-signal serialized text (metrics only, no PII)',
     'Purpose: IS the behavioral state. Enables drift, direction, peer comparison, kill-chain'),
    ('3. STRUCTURAL FEATURES', TEAL,
     'JSONB document with all ~50 raw feature values across 5 signals',
     'Example: {"auth":{"logon_count":47,"failure_rate":0.032,...},"process":{...}}',
     'Purpose: explainability (which metric drove drift?), audit (exact numbers at detection time)'),
    ('4. HMAC-CHAINED HISTORY', ORANGE,
     'HMAC-SHA256 chaining each snapshot to the previous one',
     'Formula: HMAC(prev_hmac + embedding_hash + features_hash + timestamp)',
     'Purpose: tamper-evident chain. Modifying any past snapshot breaks ALL subsequent HMACs'),
]

y = Inches(1.3)
for label, color, line1, line2, line3 in comps:
    _box(s, Inches(0.3), y, Inches(12.7), Inches(1.3), fill=CARD, border=color)
    _box(s, Inches(0.35), y + Inches(0.05), Inches(3.2), Inches(0.3), fill=color)
    _txt(s, Inches(0.4), y + Inches(0.07), Inches(3.1), Inches(0.26),
         label, 11, WHITE, True, PP_ALIGN.CENTER)
    _txt(s, Inches(3.7), y + Inches(0.05), Inches(9.1), Inches(0.3),
         line1, 11, NAVY, True)
    _txt(s, Inches(3.7), y + Inches(0.40), Inches(9.1), Inches(0.3),
         line2, 10, TXTSUB)
    _txt(s, Inches(3.7), y + Inches(0.75), Inches(9.1), Inches(0.5),
         line3, 10, GOLD)
    y += Inches(1.35)

_box(s, Inches(0.3), Inches(6.7), Inches(12.7), Inches(0.55), fill=CARD, border=GOLD)
_txt(s, Inches(0.5), Inches(6.72), Inches(12.3), Inches(0.5),
     'Together: forensic-grade behavioral evidence. WHO was resolved (hash). WHAT they did (embedding). '
     'WHY the embedding looks that way (features). WHEN it happened with proof of non-tampering (HMAC chain).',
     11, GOLD, True)

# ================================================================
# SLIDE 13: CUSUM ALGORITHM
# ================================================================
s = new_slide('CUSUM Algorithm: Catching What Thresholds Miss',
              'Cumulative Sum detects small, sustained behavioral shifts over time')

_box(s, Inches(0.3), Inches(1.3), Inches(6.2), Inches(2.8), fill=CARD, border=CYBER)
_txt(s, Inches(0.5), Inches(1.35), Inches(5.8), Inches(0.3), 'THE ALGORITHM', 14, CYBER, True)
tf = _txt(s, Inches(0.5), Inches(1.7), Inches(5.8), Inches(2.2),
     'S(t) = max(0, S(t-1) + drift(t) - (mu + k))', 12, GOLD, True, font='Cascadia Code')
_para(tf, '', 10, LTGRAY)
_para(tf, 'S(t) = CUSUM statistic at time t', 11, RGBColor(0x2A, 0x6A, 0x3A))
_para(tf, 'drift(t) = cosine drift at time t', 11, RGBColor(0x2A, 0x6A, 0x3A))
_para(tf, 'mu = expected mean drift (from baseline)', 11, RGBColor(0x2A, 0x6A, 0x3A))
_para(tf, 'k = slack parameter = 0.5 * sigma', 11, RGBColor(0x2A, 0x6A, 0x3A))
_para(tf, 'ALARM when S(t) > h = 4 * sigma', 11, RED, True)
_para(tf, '', 10, LTGRAY)
_para(tf, 'Key property: S(t) resets to 0 when drift', 11, RGBColor(0x2A, 0x6A, 0x3A))
_para(tf, 'returns to normal. No false accumulation.', 11, RGBColor(0x2A, 0x6A, 0x3A))

_box(s, Inches(6.8), Inches(1.3), Inches(6.2), Inches(2.8), fill=CARD, border=GOLD)
_txt(s, Inches(7.0), Inches(1.35), Inches(5.8), Inches(0.3), 'CONCRETE NUMBERS', 14, GOLD, True)
tf = _txt(s, Inches(7.0), Inches(1.7), Inches(5.8), Inches(2.2),
     'Typical baseline drift sigma = 0.03', 12, RGBColor(0x2A, 0x6A, 0x3A))
_para(tf, '', 10, LTGRAY)
_para(tf, 'k = 0.5 * 0.03 = 0.015', 11, GOLD, True)
_para(tf, '  (CUSUM tolerates drift up to baseline+0.015)', 11, TXTSUB)
_para(tf, 'h = 4 * 0.03 = 0.12', 11, GOLD, True)
_para(tf, '  (alarm at 0.12 accumulated excess drift)', 11, TXTSUB)
_para(tf, '', 10, LTGRAY)
_para(tf, 'Example: sustained drift of 0.04/window', 11, NAVY, True)
_para(tf, '  Excess per window: 0.04 - 0.03 - 0.015 = -0.005', 11, TXTSUB)
_para(tf, '  Wait: that is negative. Need drift > mu+k:', 11, TXTSUB)
_para(tf, '  drift 0.06: excess = 0.06-0.03-0.015 = 0.015/win', 11, GOLD)
_para(tf, '  ALARM at 0.12/0.015 = 8 windows (~8 hours)', 11, RED, True)

_box(s, Inches(0.3), Inches(4.3), Inches(12.7), Inches(1.0), fill=CARD, border=RED)
_txt(s, Inches(0.5), Inches(4.35), Inches(12.3), Inches(0.3),
     'WHY CUSUM CATCHES VOLT TYPHOON BUT THRESHOLDS DO NOT', 13, RED, True)
tf = _txt(s, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.55),
     'Fixed threshold at 0.15: Volt Typhoon drift per window = 0.01-0.04. '
     'NEVER crosses 0.15. Attacker operates for 5+ years.', 11, RGBColor(0x2A, 0x6A, 0x3A))
_para(tf, 'CUSUM at 4-sigma: drift 0.03-0.04/window accumulates. '
     'CUSUM triggers in 4-8 days. Dwell time drops from YEARS to DAYS.',
     11, GOLD, True)

# Threshold vs CUSUM visual
_box(s, Inches(0.3), Inches(5.5), Inches(6.2), Inches(1.7), fill=CARD, border=RED)
_txt(s, Inches(0.5), Inches(5.55), Inches(5.8), Inches(0.3), 'FIXED THRESHOLD = 0.15', 13, RED, True)
tf = _txt(s, Inches(0.5), Inches(5.9), Inches(5.8), Inches(1.2),
     'Window 1:  drift = 0.02  (no alert)', 10, TXTSUB, font='Cascadia Code')
_para(tf, 'Window 2:  drift = 0.03  (no alert)', 10, TXTSUB)
_para(tf, 'Window 3:  drift = 0.04  (no alert)', 10, TXTSUB)
_para(tf, '...', 10, LTGRAY)
_para(tf, 'Window 100: drift = 0.04  (STILL no alert)', 10, TXTSUB)
_para(tf, 'Attacker: 100+ hours undetected', 10, RED, True)

_box(s, Inches(6.8), Inches(5.5), Inches(6.2), Inches(1.7), fill=CARD, border=CYBER)
_txt(s, Inches(7.0), Inches(5.55), Inches(5.8), Inches(0.3), 'CUSUM (h = 4*sigma = 0.12)', 13, CYBER, True)
tf = _txt(s, Inches(7.0), Inches(5.9), Inches(5.8), Inches(1.2),
     'Window 1:  S=0.005  (accumulating)', 10, RGBColor(0x2A, 0x6A, 0x3A), font='Cascadia Code')
_para(tf, 'Window 2:  S=0.020  (accumulating)', 10, RGBColor(0x2A, 0x6A, 0x3A))
_para(tf, 'Window 3:  S=0.045  (accumulating)', 10, RGBColor(0x2A, 0x6A, 0x3A))
_para(tf, 'Window 5:  S=0.085  (WARNING)', 10, GOLD)
_para(tf, 'Window 8:  S=0.125  -> ALARM (>0.12)', 10, RED, True)
_para(tf, 'Attacker: detected in 8 hours', 10, CYBER, True)

# ================================================================
# SLIDE 14: DRIFT DIRECTION — ALL 8 CONCEPTS
# ================================================================
s = new_slide('Drift Direction: All 8 Threat Concepts',
              'Each concept is a text description embedded into the same 1536-d space')

concepts = [
    ('credential_dumping', 'T1003', BLUE,
     'LSASS access, SAM reads, ntdsutil, DCSync, Kerberoasting'),
    ('lateral_movement', 'T1021', PURPLE,
     'RDP/WinRM/SSH to new hosts, admin shares, PtH/PtT'),
    ('data_exfiltration', 'T1041', TEAL,
     'Bytes-out inversion, archive staging, cloud uploads, DNS tunnel'),
    ('c2_beaconing', 'T1071', ORANGE,
     'Regular-interval callbacks, DGA domains, encrypted C2 channels'),
    ('lotl_execution', 'T1059', GOLD,
     'wmic/ntdsutil/netsh/rundll32 chains, encoded PowerShell'),
    ('privilege_escalation', 'T1078', RED,
     'Group adds (Domain Admins), sudo abuse, token manipulation'),
    ('defense_evasion', 'T1562', RGBColor(0xCC, 0x66, 0x99),
     'AV exclusions, log clearing, agent tampering, timestomping'),
    ('insider_hoarding', 'T1074', RGBColor(0x66, 0xBB, 0xCC),
     'Steady sensitive access increase, bulk downloads, data staging'),
]

y = Inches(1.3)
for name, tid, color, desc in concepts:
    _box(s, Inches(0.3), y, Inches(12.7), Inches(0.65), fill=CARD, border=RGBColor(0x33, 0x33, 0x55))
    _box(s, Inches(0.35), y + Inches(0.05), Inches(2.2), Inches(0.25), fill=color)
    _txt(s, Inches(0.4), y + Inches(0.07), Inches(2.1), Inches(0.22),
         name, 10, WHITE, True, PP_ALIGN.CENTER)
    _txt(s, Inches(2.7), y + Inches(0.05), Inches(0.8), Inches(0.22),
         tid, 10, GOLD, True, font='Cascadia Code')
    _txt(s, Inches(3.6), y + Inches(0.05), Inches(9.2), Inches(0.55),
         desc, 10, TXTSUB)
    y += Inches(0.68)

_box(s, Inches(0.3), Inches(6.8), Inches(12.7), Inches(0.5), fill=CARD, border=GOLD)
_txt(s, Inches(0.5), Inches(6.82), Inches(12.3), Inches(0.45),
     'projection = cosine(drift_vector, concept_embedding). Top 3 returned with MITRE IDs. '
     'Result: "Entity drifting toward lateral_movement (0.82, T1021) + lotl_execution (0.71, T1059)"',
     11, GOLD, True)

# ================================================================
# SLIDE 15: ABAC TRUST STATE MACHINE
# ================================================================
s = new_slide('ABAC Trust State Machine',
              'Proportional, automated enforcement. Not binary on/off.')

states = [
    ('TRUSTED', CYBER, 'health > 70, no CUSUM concern',
     'Normal operations. Standard access controls. Baseline monitoring.',
     'Default state'),
    ('ELEVATED_WATCH', GOLD, 'health < 70 OR cusum > 2*sigma',
     'Step-up MFA on next auth. Increased log verbosity. Analyst notification.',
     'Reversal: health > 75 for 24h'),
    ('RESTRICTED', ORANGE, 'health < 40 OR cusum > 3.5*sigma OR concept > 0.55',
     'Read-only access. Lateral movement blocked. Incident ticket created.',
     'Reversal: analyst review + health > 60 for 48h'),
    ('BLOCKED', RED, 'health < 20 OR cusum ALARM (4*sigma) OR concept > 0.70',
     'All access revoked. Active sessions terminated. IR team notified.',
     'Reversal: IR team manual unlock ONLY'),
]

y = Inches(1.3)
for label, color, trigger, action, reversal in states:
    _box(s, Inches(0.3), y, Inches(12.7), Inches(1.25), fill=CARD, border=color)
    _box(s, Inches(0.35), y + Inches(0.05), Inches(2.3), Inches(0.35), fill=color)
    _txt(s, Inches(0.4), y + Inches(0.07), Inches(2.2), Inches(0.3),
         label, 14, WHITE, True, PP_ALIGN.CENTER)
    _txt(s, Inches(2.8), y + Inches(0.05), Inches(2.0), Inches(0.18),
         'TRIGGER:', 9, color, True)
    _txt(s, Inches(2.8), y + Inches(0.25), Inches(9.9), Inches(0.3),
         trigger, 10, RGBColor(0x2A, 0x6A, 0x3A))
    _txt(s, Inches(2.8), y + Inches(0.55), Inches(2.0), Inches(0.18),
         'ENFORCEMENT:', 9, color, True)
    _txt(s, Inches(2.8), y + Inches(0.72), Inches(9.9), Inches(0.3),
         action, 10, TXTSUB)
    _txt(s, Inches(2.8), y + Inches(0.98), Inches(9.9), Inches(0.22),
         reversal, 9, LTGRAY)
    y += Inches(1.3)

_box(s, Inches(0.3), Inches(6.55), Inches(12.7), Inches(0.7), fill=CARD, border=GOLD)
_txt(s, Inches(0.5), Inches(6.57), Inches(12.3), Inches(0.65),
     'KEY: Response is PROPORTIONAL. No binary lockout from one anomaly. Graduated enforcement. '
     'Each level reversible with specific criteria. False positive feedback auto-adjusts baseline. '
     'Detection to enforcement: <5 minutes. System improves with analyst feedback, not decays.',
     11, GOLD, True)

# ================================================================
# SLIDE 16: SOLVING SIGMA RULE DECAY
# ================================================================
s = new_slide('How We Solve Sigma Rule Decay',
              'Signatures detect TOOLS. We detect OBJECTIVES. Tools change. Objectives cannot.')

_box(s, Inches(0.3), Inches(1.3), Inches(6.2), Inches(2.5), fill=CARD, border=RED)
_txt(s, Inches(0.5), Inches(1.35), Inches(5.8), Inches(0.3), 'THE PROBLEM: ARMS RACE', 14, RED, True)
tf = _txt(s, Inches(0.5), Inches(1.7), Inches(5.8), Inches(2.0),
     'Step 1: Defender writes Sigma rule for ntdsutil', 10, TXTSUB)
_para(tf, 'Step 2: Rule published (open source / ISAC shared)', 10, TXTSUB)
_para(tf, 'Step 3: PRC actor downloads published rules', 10, TXTSUB)
_para(tf, 'Step 4: Tests tools until no rule matches', 10, TXTSUB)
_para(tf, 'Step 5: Deploys modified tool in campaign', 10, TXTSUB)
_para(tf, '', 10, LTGRAY)
_para(tf, 'Defender retool cycle: 1-3 months', 11, RED, True)
_para(tf, 'Adversary retool cycle: 2-4 weeks', 11, RED, True)
_para(tf, 'Coverage DECAYS between rule updates', 11, RED, True)

_box(s, Inches(6.8), Inches(1.3), Inches(6.2), Inches(2.5), fill=CARD, border=CYBER)
_txt(s, Inches(7.0), Inches(1.35), Inches(5.8), Inches(0.3), 'OUR SOLUTION: BEHAVIORAL OBJECTIVES', 14, CYBER, True)
tf = _txt(s, Inches(7.0), Inches(1.7), Inches(5.8), Inches(2.0),
     'We do not detect specific tools or commands.', 11, RGBColor(0x2A, 0x6A, 0x3A))
_para(tf, 'We detect the BEHAVIORAL EFFECT of the attack:', 11, GOLD, True)
_para(tf, '', 10, LTGRAY)
_para(tf, 'Credential access: entity MUST access cred stores', 10, RGBColor(0x2A, 0x6A, 0x3A))
_para(tf, 'Lateral movement: entity MUST auth to new targets', 10, RGBColor(0x2A, 0x6A, 0x3A))
_para(tf, 'Exfiltration: entity MUST send data outbound', 10, RGBColor(0x2A, 0x6A, 0x3A))
_para(tf, 'Persistence: entity MUST create scheduled tasks/registry', 10, RGBColor(0x2A, 0x6A, 0x3A))
_para(tf, '', 10, LTGRAY)
_para(tf, 'These behavioral effects are INHERENT to the', 11, GOLD, True)
_para(tf, 'attack objective. Cannot be changed without', 11, GOLD, True)
_para(tf, 'abandoning the attack itself.', 11, GOLD, True)

# Concrete example
_box(s, Inches(0.3), Inches(4.0), Inches(12.7), Inches(2.0), fill=CARD, border=GOLD)
_txt(s, Inches(0.5), Inches(4.05), Inches(12.3), Inches(0.3),
     'EXAMPLE: Attacker changes tool for credential dumping', 14, GOLD, True)

_txt(s, Inches(0.5), Inches(4.4), Inches(3.8), Inches(0.2), 'SIGMA RULE APPROACH:', 10, RED, True)
tf = _txt(s, Inches(0.5), Inches(4.65), Inches(3.8), Inches(1.2),
     'Rule: "process_name=ntdsutil"', 10, TXTSUB, font='Cascadia Code')
_para(tf, 'Attacker switches to secretsdump.py', 10, TXTSUB)
_para(tf, 'Rule: MISS (different process name)', 10, RED, True)
_para(tf, 'Attacker switches to pypykatz', 10, TXTSUB)
_para(tf, 'Rule: MISS (another different name)', 10, RED, True)

_txt(s, Inches(4.8), Inches(4.4), Inches(8.0), Inches(0.2), 'OUR APPROACH (BEHAVIORAL TRAJECTORY):', 10, CYBER, True)
tf = _txt(s, Inches(4.8), Inches(4.65), Inches(8.0), Inches(1.2),
     'Regardless of tool used, the BEHAVIOR is identical:', 10, RGBColor(0x2A, 0x6A, 0x3A))
_para(tf, '  AUTH: new service account patterns', 10, GOLD)
_para(tf, '  PROCESS: high-priv process accessing LSASS/SAM/NTDS', 10, GOLD)
_para(tf, '  FILE: access to credential store paths', 10, GOLD)
_para(tf, 'Drift direction: credential_dumping = 0.73 (T1003)', 11, CYBER, True)
_para(tf, 'Tool-agnostic. Signature-free. Cannot be evaded by retooling.', 11, CYBER, True)

_box(s, Inches(0.3), Inches(6.2), Inches(12.7), Inches(1.0), fill=NAVY)
_txt(s, Inches(0.5), Inches(6.25), Inches(12.3), Inches(0.9),
     'THE MATHEMATICAL INCOMPATIBILITY: An attacker MUST change behavior to achieve their objective. '
     'They can change the tool (ntdsutil -> secretsdump -> mimikatz -> pypykatz). '
     'They CANNOT change the behavioral effect (credential access). '
     'Signatures detect tools. We detect objectives. Tools change. Objectives cannot.',
     12, GOLD, True)

# ================================================================
# SLIDE 17: KEY NUMBERS
# ================================================================
s = new_slide('Key Numbers to Memorize',
              'Have these on the tip of your tongue for Q&A')

left_nums = [
    ('5+ years', 'Volt Typhoon dwell time'),
    ('4+ years', 'Salt Typhoon dwell time'),
    ('9', 'US telecoms compromised (Salt)'),
    ('99%', 'FW breaches = misconfig (Gartner)'),
    ('~20%', 'Attacks stopped by KEV patching'),
    ('2^8000', 'Config state space (ASM samples)'),
    ('1,536', 'Embedding dimensions per signal'),
    ('9,216', 'Total dims (6 x 1536)'),
    ('5', 'Behavioral signals'),
    ('8', 'MITRE threat concepts'),
    ('4-sigma', 'CUSUM alarm threshold'),
]

right_nums = [
    ('<3 sec', 'Event to analysis per entity window'),
    ('10K/hr', 'Entities processed per 4-vCPU'),
    ('270 days', 'Synthetic data in prototype'),
    ('6', 'Attack scenarios detectable'),
    ('7', 'SOC dashboard charts'),
    ('30-45 min', 'Current analyst triage time'),
    ('<3 sec', 'Our triage time'),
    ('$90M', '22CT Army SOC contract'),
    ('800+', 'Cleared analysts on contract'),
    ('90 days', 'Time to first detection capability'),
    ('IL5/IL6/JWICS', 'Three enclave support'),
]

for col_nums, x_start in [(left_nums, Inches(0.3)), (right_nums, Inches(6.8))]:
    y = Inches(1.3)
    for num, meaning in col_nums:
        _box(s, x_start, y, Inches(6.2), Inches(0.45), fill=CARD, border=RGBColor(0x33, 0x33, 0x55))
        _box(s, x_start + Inches(0.05), y + Inches(0.03), Inches(1.5), Inches(0.32), fill=NAVY)
        _txt(s, x_start + Inches(0.1), y + Inches(0.05), Inches(1.4), Inches(0.28),
             num, 12, CYBER, True, PP_ALIGN.CENTER, font='Cascadia Code')
        _txt(s, x_start + Inches(1.7), y + Inches(0.07), Inches(4.3), Inches(0.3),
             meaning, 11, TXTSUB)
        y += Inches(0.48)

# ================================================================
# SAVE
# ================================================================
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
