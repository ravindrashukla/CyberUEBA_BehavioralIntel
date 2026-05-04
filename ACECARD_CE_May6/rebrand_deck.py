"""Replace 'Rigor AI' with '22CT Preemptive' throughout the final presentation.
Rigor AI is acknowledged as innovation partner, not removed entirely."""
import sys
sys.stdout.reconfigure(encoding='utf-8')
import re
from pptx import Presentation

prs = Presentation('CFIC_Final_Presentation_new.pptx')

# Replacements - order matters (longer/more specific first)
replacements = [
    # Remove Rigor AI contact info
    ('Rigor AI:  sales-ext@rigor.ai', '22CT Preemptive (powered by Rigor AI, innovation partner)'),
    ('Sundar Iyer, CEO   ·   Navneet Yadav, CPO', 'Ravindra Shukla, Head of AI'),
    ('sales-ext@rigor.ai', ''),

    # Specific phrases
    ('Rigor AI is engaged with 20+ design partners across F500, federal, critical infrastructure, MSSPs, and nation states.', '22CT Preemptive is validated across F500, federal, critical infrastructure, and nation-state engagements (powered by Rigor AI, innovation partner).'),
    ('End-to-End Architecture: Rigor AI + ACECARD UEBA', 'End-to-End Architecture: 22CT Preemptive + ACECARD UEBA'),
    ('Rigor AI + ACECARD: The Data Flow', '22CT Preemptive + ACECARD: The Data Flow'),
    ('Rigor AI + ACECARD UEBA = comprehensive defense', '22CT Preemptive + ACECARD UEBA = comprehensive defense'),
    ('How RigorOS Operates: Inputs', 'How 22CT Preemptive Operates: Inputs'),
    ('How Rigor AI Preempts Volt Typhoon', 'How 22CT Preemptive Stops Volt Typhoon'),
    ('How Rigor AI Preempts Salt Typhoon', 'How 22CT Preemptive Stops Salt Typhoon'),
    ('HOW RIGOR AI STOPS IT', 'HOW 22CT PREEMPTIVE STOPS IT'),
    ('Pilot Rigor AI', 'Pilot 22CT Preemptive'),
    ('Rigor in Production: Real Customer Findings', '22CT Preemptive in Production: Real Findings'),

    # Part headers
    ('RIGOR AI  —  VOLT TYPHOON COVERAGE', '22CT PREEMPTIVE  —  VOLT TYPHOON COVERAGE'),
    ('RIGOR AI  —  SALT TYPHOON COVERAGE', '22CT PREEMPTIVE  —  SALT TYPHOON COVERAGE'),
    ('RIGOR AI  —  USE CASES', '22CT PREEMPTIVE  —  USE CASES'),
    ('RIGOR AI  —  HONEST ASSESSMENT', '22CT PREEMPTIVE  —  HONEST ASSESSMENT'),
    ('RIGOR AI  —  PROOF POINTS', '22CT PREEMPTIVE  —  PROOF POINTS'),

    # Layer labels
    ('LAYER 1 — RIGOR AI (PREEMPTIVE)', 'LAYER 1 — 22CT PREEMPTIVE'),
    ('PREEMPTIVE LAYER (RIGOR AI)', 'PREEMPTIVE LAYER (22CT)'),
    ('2 · PREEMPTIVE LAYER (RIGOR AI)', '2 · PREEMPTIVE LAYER (22CT PREEMPTIVE)'),

    # Possessive forms
    ("Rigor’s Attack Intel", "22CT Preemptive’s Attack Intel"),
    ("Rigor's Attack Intel", "22CT Preemptive's Attack Intel"),
    ("Rigor's CVE-to-control mapping", "22CT Preemptive's CVE-to-control mapping"),
    ("Rigor's Configuration Drift Management", "22CT Preemptive's Configuration Drift Management"),
    ("Rigor's threat-feed", "22CT Preemptive's threat-feed"),
    ("Rigor's MITRE Coverage Map", "22CT Preemptive's MITRE Coverage Map"),
    ("Rigor's Verified Findings", "22CT Preemptive's Verified Findings"),
    ("Rigor's Confirmed", "22CT Preemptive's Confirmed"),
    ("Rigor's preemptive verification", "22CT Preemptive's verification"),
    ("Rigor's hourly re-verification", "22CT Preemptive's hourly re-verification"),
    ("Rigor's NGFW pillar", "22CT Preemptive's NGFW pillar"),
    ("Rigor's third-party", "22CT Preemptive's third-party"),
    ("Rigor's continuous baseline", "22CT Preemptive's continuous baseline"),

    # Operational references
    ('Rigor closes the configuration', '22CT Preemptive closes the configuration'),
    ('Rigor + ACECARD is a system', '22CT Preemptive + ACECARD is a system'),
    ('Rigor + ACECARD', '22CT Preemptive + ACECARD'),
    ('Rigor governs', '22CT Preemptive governs'),
    ('Rigor pushes', '22CT Preemptive pushes'),
    ('Rigor publishes', '22CT Preemptive publishes'),

    # Action verbs
    ('Rigor mathematically proves', '22CT Preemptive mathematically proves'),
    ('Rigor mathematically', '22CT Preemptive mathematically'),
    ('Rigor formally proves', '22CT Preemptive formally proves'),
    ('Rigor formally extended', '22CT Preemptive formally extended'),
    ('Rigor formal model', '22CT Preemptive formal model'),
    ('Rigor formally', '22CT Preemptive formally'),
    ('Rigor verifies', '22CT Preemptive verifies'),
    ('Rigor proves', '22CT Preemptive proves'),
    ('Rigor models', '22CT Preemptive models'),
    ('Rigor exposed', '22CT Preemptive exposed'),
    ('Rigor identifies', '22CT Preemptive identifies'),
    ('Rigor takes', '22CT Preemptive takes'),
    ('Rigor generates', '22CT Preemptive generates'),
    ('Rigor makes', '22CT Preemptive makes'),
    ('Rigor catches', '22CT Preemptive catches'),
    ('Rigor surfaces', '22CT Preemptive surfaces'),

    # Product references
    ('RigorOS ships', '22CT Preemptive ships'),
    ('RigorOS', '22CT Preemptive Engine'),
    ('RIGOR ENGINE', '22CT PREEMPTIVE ENGINE'),
    ("Rigor doesn't detect", "22CT Preemptive doesn't detect"),
    ("Rigor doesn’t detect", "22CT Preemptive doesn’t detect"),

    # Generic catch-all for "Rigor AI"
    ('Rigor AI', '22CT Preemptive'),
]

count = 0
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    original = run.text
                    text = run.text
                    for old, new in replacements:
                        if old in text:
                            text = text.replace(old, new)
                    if text != original:
                        run.text = text
                        count += 1

# Second pass: catch remaining standalone 'Rigor' (not 'Rigorous', not 'rigor.ai')
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    text = run.text
                    if 'Rigor' in text and 'Rigorous' not in text and 'rigor.' not in text:
                        new_text = re.sub(r'\bRigor\b(?!ous|\.)', '22CT Preemptive', text)
                        if new_text != text:
                            run.text = new_text
                            count += 1

out = 'CFIC_Final_Presentation_v2.pptx'
prs.save(out)
print(f'Saved: {out}')
print(f'Text runs modified: {count}')
print(f'Slides: {len(prs.slides)}')
