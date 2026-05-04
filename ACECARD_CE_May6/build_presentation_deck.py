"""Build ACECARD CE Solution Brief presentation deck for May 6, 2026.
Focused on SOLUTION, not marketing. 10-12 minute presentation for the 1420-1535 slot."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

OUT = os.path.join(os.path.dirname(__file__), "ACECARD_CE_Solution_Brief_Deck.pptx")

NAVY   = RGBColor(0x00, 0x2B, 0x5C)
GOLD   = RGBColor(0xB8, 0x96, 0x3A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x1A, 0x1A, 0x2E)
CYBER  = RGBColor(0x00, 0xD4, 0xAA)
RED    = RGBColor(0xE5, 0x3E, 0x3E)
BLUE   = RGBColor(0x2D, 0x5A, 0x87)
GRAY   = RGBColor(0x4A, 0x55, 0x68)
LTGRAY = RGBColor(0xA0, 0xAE, 0xC0)
BG     = RGBColor(0x0F, 0x0F, 0x23)
CARD   = RGBColor(0x1A, 0x1A, 0x2E)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
TEAL   = RGBColor(0x14, 0xB8, 0xA6)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height
TOTAL = 12


def _fill(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _box(slide, l, t, w, h, fill=None, border=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.background()
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    ln = s.line
    if border:
        ln.color.rgb = border
        ln.width = Pt(1.5)
    else:
        ln.fill.background()
    s.shadow.inherit = False
    return s


def _txt(slide, l, t, w, h, text, size=14, color=WHITE, bold=False,
         align=PP_ALIGN.LEFT, font='Segoe UI'):
    tf = slide.shapes.add_textbox(l, t, w, h).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return tf


def _para(tf, text, size=14, color=WHITE, bold=False, space=Pt(4)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Segoe UI'
    p.space_before = space
    return p


def _title_bar(slide, title, subtitle=None):
    _box(slide, Inches(0), Inches(0), W, Inches(1.1), fill=NAVY)
    _txt(slide, Inches(0.6), Inches(0.15), Inches(10), Inches(0.5),
         title, 26, WHITE, True)
    if subtitle:
        _txt(slide, Inches(0.6), Inches(0.6), Inches(10), Inches(0.4),
             subtitle, 13, LTGRAY)
    _txt(slide, Inches(10.5), Inches(0.25), Inches(2.5), Inches(0.5),
         'ACECARD CE', 14, CYBER, True, PP_ALIGN.RIGHT)


def _snum(slide, n):
    _txt(slide, Inches(12.0), Inches(7.0), Inches(1.0), Inches(0.3),
         f'{n}/{TOTAL}', 10, LTGRAY, align=PP_ALIGN.RIGHT)


def _bullet_card(slide, l, t, w, h, title, bullets, title_color=CYBER,
                 border=None):
    _box(slide, l, t, w, h, fill=CARD, border=border or RGBColor(0x33, 0x33, 0x55))
    tf = _txt(slide, l + Inches(0.2), t + Inches(0.1), w - Inches(0.4),
              Inches(0.35), title, 14, title_color, True)
    y = t + Inches(0.5)
    for b in bullets:
        _txt(slide, l + Inches(0.3), y, w - Inches(0.5), Inches(0.26),
             f"•  {b}", 11, RGBColor(0xCC, 0xCC, 0xDD))
        y += Inches(0.26)


# ═══════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_fill(s, BG)
_box(s, Inches(0), Inches(0), W, Inches(0.06), fill=CYBER)
_box(s, Inches(0), Inches(7.44), W, Inches(0.06), fill=CYBER)

_txt(s, Inches(1), Inches(1.5), Inches(11), Inches(0.6),
     'ACECARD CE — Solution Brief', 48, CYBER, True, PP_ALIGN.CENTER)
_txt(s, Inches(1), Inches(2.3), Inches(11), Inches(0.5),
     'Behavioral Anomaly Detection Mission Application', 22, WHITE, False,
     PP_ALIGN.CENTER)
_txt(s, Inches(1), Inches(3.0), Inches(11), Inches(0.4),
     'Preemptive Cyber Defense + Behavioral Trajectory Intelligence', 16,
     GOLD, True, PP_ALIGN.CENTER)

_box(s, Inches(3), Inches(3.8), Inches(7.333), Inches(0.02),
     fill=RGBColor(0x33, 0x33, 0x55))

_txt(s, Inches(1), Inches(4.2), Inches(11), Inches(0.4),
     'A Containerized, Cloud-Native Prototype for Gabriel Nimbus', 15,
     LTGRAY, False, PP_ALIGN.CENTER)

_txt(s, Inches(1), Inches(5.2), Inches(11), Inches(0.4),
     'CFIC  |  U.S. Army Cyber Command (ARCYBER)  |  ArCTIC', 13,
     RGBColor(0x88, 0x88, 0xAA), False, PP_ALIGN.CENTER)
_txt(s, Inches(1), Inches(5.6), Inches(11), Inches(0.4),
     '22nd Century Technologies  |  Rigor AI Inc.', 13,
     RGBColor(0x77, 0x77, 0x99), False, PP_ALIGN.CENTER)
_txt(s, Inches(1), Inches(6.5), Inches(11), Inches(0.3),
     '06 MAY 2026  |  Solution-Focused  |  Working Prototype', 11,
     RGBColor(0x66, 0x66, 0x88), False, PP_ALIGN.CENTER)
_snum(s, 1)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2: THE PROBLEM
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_fill(s, BG)
_title_bar(s, 'The Problem', 'Four challenges that current SOC tooling cannot solve')

challenges = [
    ('Manual Context Reassembly', BLUE,
     'Analyst touches 6–12 data tables per alert',
     'Each table: different schema, different timestamps, different entity IDs',
     '30–45 min triage time per alert × 500 alerts/day = most never investigated'),
    ('Inconsistent Siloed Analytics', PURPLE,
     'AD, AAD, Okta, AWS, K8s each see 1 login',
     'Lateral movement across 5 systems = 5 weak signals, not 1 strong signal',
     'No current tool fuses these into a single trajectory per entity'),
    ('Scalability Constraints', TEAL,
     'Manual correlation caps at ~50 entities',
     'N log sources = N separate analytics pipelines',
     'Gabriel Nimbus scale: 100K+ entities. Current tools cannot keep up'),
    ('Operational Inefficiency', ORANGE,
     'Anomaly score "87%" — no direction, no MITRE mapping',
     'Static Sigma rules decay monthly (PRC actors test against them)',
     'Alert fatigue: analysts ignore most alerts. Attacker operates freely'),
]

x = Inches(0.3)
for title, color, l1, l2, l3 in challenges:
    bw = Inches(3.1)
    _box(s, x, Inches(1.3), bw, Inches(0.4), fill=color)
    _txt(s, x, Inches(1.32), bw, Inches(0.35), title, 13, WHITE, True, PP_ALIGN.CENTER)
    _box(s, x, Inches(1.75), bw, Inches(3.5), fill=CARD, border=color)
    y = Inches(1.85)
    for line in [l1, l2, l3]:
        _txt(s, x + Inches(0.15), y, bw - Inches(0.3), Inches(0.8),
             f"•  {line}", 11, RGBColor(0xBB, 0xCC, 0xDD))
        y += Inches(1.0)
    x += bw + Inches(0.15)

_txt(s, Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.8),
     'Result: Volt Typhoon dwelled 5+ years. Salt Typhoon dwelled 4+ years. '
     'Neither was caught by config-only or behavior-only defenses.',
     13, GOLD, True, PP_ALIGN.CENTER)
_txt(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.5),
     'Source: CISA AA24-038A (Volt Typhoon) | CISA AA25-239A (Salt Typhoon) | '
     'Gartner: 99% of firewall breaches = misconfiguration',
     10, LTGRAY, False, PP_ALIGN.CENTER)
_snum(s, 2)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3: WHY CURRENT TOOLS FAIL (8 GAPS)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_fill(s, BG)
_title_bar(s, 'Why Current Tools Structurally Cannot Solve This',
           'Two gap categories — one in preemption, one in detection')

# Left column: Preemptive
_box(s, Inches(0.3), Inches(1.3), Inches(6.3), Inches(0.45), fill=NAVY)
_txt(s, Inches(0.5), Inches(1.33), Inches(6), Inches(0.4),
     'PREEMPTIVE LAYER FAILS', 16, WHITE, True)

pre_gaps = [
    ('99% of firewall breaches = misconfiguration', 'Gartner. Both Typhoons exploit mis-set NGFWs and VPNs.'),
    ('KEV patching covers only ~20% of attacks', 'The other 80% — most Salt Typhoon paths — have no CVE patch.'),
    ('ASM/BAS only sample the 2^8000 state space', 'Cannot PROVE protection. Coverage is statistical, not deterministic.'),
    ('Config drift goes unverified', 'Salt Typhoon modifies router ACLs. Drift-detection is reactive, not preemptive.'),
]

y = Inches(1.85)
for title, desc in pre_gaps:
    _box(s, Inches(0.4), y, Inches(0.08), Inches(0.7), fill=RED)
    _txt(s, Inches(0.6), y, Inches(5.8), Inches(0.3), title, 12, WHITE, True)
    _txt(s, Inches(0.6), y + Inches(0.3), Inches(5.8), Inches(0.35), desc, 10, LTGRAY)
    y += Inches(0.78)

# Right column: Behavioral
_box(s, Inches(6.8), Inches(1.3), Inches(6.3), Inches(0.45), fill=NAVY)
_txt(s, Inches(7.0), Inches(1.33), Inches(6), Inches(0.4),
     'BEHAVIORAL LAYER FAILS', 16, WHITE, True)

beh_gaps = [
    ('Threshold scoring misses LoTL', 'Volt Typhoon\'s process tree was \'normal\' at every step. ntdsutil never crosses a threshold.'),
    ('Per-source siloed analytics fragment identity', 'AD, AAD, AWS, Okta, K8s each see one identifier — lateral movement reads as 4 weak signals.'),
    ('Static rules decay monthly', 'PRC actors test against published Sigma rules before campaigns. Coverage decays by month.'),
    ('\'Anomaly score\' without direction', 'Analyst gets \'87% anomalous\' — no hypothesis, no MITRE mapping. Triage time blows out.'),
]

y = Inches(1.85)
for title, desc in beh_gaps:
    _box(s, Inches(6.9), y, Inches(0.08), Inches(0.7), fill=RED)
    _txt(s, Inches(7.1), y, Inches(5.8), Inches(0.3), title, 12, WHITE, True)
    _txt(s, Inches(7.1), y + Inches(0.3), Inches(5.8), Inches(0.35), desc, 10, LTGRAY)
    y += Inches(0.78)

_snum(s, 3)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4: THE TWO-LAYER SOLUTION
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_fill(s, BG)
_title_bar(s, 'The Solution: Two Layers',
           'Rigor preempts what CAN be prevented. ACECARD detects what CANNOT.')

# Layer 1
_box(s, Inches(0.3), Inches(1.3), Inches(6.3), Inches(4.2), fill=CARD, border=BLUE)
_box(s, Inches(0.3), Inches(1.3), Inches(6.3), Inches(0.45), fill=BLUE)
_txt(s, Inches(0.5), Inches(1.33), Inches(6), Inches(0.4),
     'LAYER 1: Rigor AI — Preemptive Defense', 16, WHITE, True)

rigor_points = [
    'Formal mathematical model of ALL security controls (NGFW, IPS, IdP, SASE, WAF)',
    'Exhaustive analysis of 2^8000 state space — symbolic reasoning, NOT sampling',
    'PROVES attack paths are closed — mathematical certainty, not probability',
    'Continuous re-verification on every config change (hourly or per-change)',
    'Risk-prioritized remediation with vendor-specific fix instructions',
    'Solves the 99% misconfiguration problem: finds multi-rule shadow, zombie rules, ACL gaps',
]
y = Inches(1.9)
for p in rigor_points:
    _txt(s, Inches(0.5), y, Inches(5.9), Inches(0.3),
         f"•  {p}", 11, RGBColor(0xBB, 0xCC, 0xDD))
    y += Inches(0.35)

# Layer 2
_box(s, Inches(6.8), Inches(1.3), Inches(6.3), Inches(4.2), fill=CARD, border=CYBER)
_box(s, Inches(6.8), Inches(1.3), Inches(6.3), Inches(0.45), fill=RGBColor(0x00, 0x6B, 0x55))
_txt(s, Inches(7.0), Inches(1.33), Inches(6), Inches(0.4),
     'LAYER 2: ACECARD UEBA — Behavioral Detection', 16, WHITE, True)

acecard_points = [
    '1536-d behavioral embedding per entity per hour (OpenAI or local SLM)',
    '5 behavioral signals: auth, process, network, file, identity',
    'CUSUM change-point detection catches slow drift over weeks/months',
    'Drift direction: cosine projection onto 8 MITRE-mapped threat concepts',
    'Peer cohort z-score comparison (role/OU/group context)',
    'ABAC trust loop: auto-enforces step-up MFA, restrict, or block',
]
y = Inches(1.9)
for p in acecard_points:
    _txt(s, Inches(7.0), y, Inches(5.9), Inches(0.3),
         f"•  {p}", 11, RGBColor(0xBB, 0xDD, 0xCC))
    y += Inches(0.35)

# Bottom: why both
_box(s, Inches(0.3), Inches(5.7), Inches(12.7), Inches(1.2), fill=CARD, border=GOLD)
_txt(s, Inches(0.5), Inches(5.75), Inches(12.3), Inches(0.35),
     'Why Both Layers Are Required:', 14, GOLD, True)
_txt(s, Inches(0.5), Inches(6.15), Inches(12.3), Inches(0.7),
     'Rigor cannot detect attacks in progress (valid credentials, zero-days, insider threats). '
     'ACECARD cannot prevent misconfigured controls. '
     'Together: either the path never existed (Rigor closed it) OR the attack is detected and '
     'contained within hours (ACECARD + ABAC). Defense in depth.',
     12, RGBColor(0xCC, 0xDD, 0xCC))
_snum(s, 4)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5: HOW IT SOLVES THE 4 CFIC CHALLENGES
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_fill(s, BG)
_title_bar(s, 'How It Solves the Four CFIC Challenges',
           'Direct mapping: each challenge → specific capability')

challenges_solve = [
    ('Manual Context Reassembly', BLUE,
     'Unified 1536-d embedding fuses ALL log sources into one vector per entity. '
     'No manual joins across tables.',
     'One API call: GET /api/trajectory/{entity}/dashboard → drift, velocity, '
     'CUSUM, health, top-3 threat concepts, MITRE technique IDs, peer z-score.',
     'Analyst goes from alert to full context in <3 seconds, not 30–45 minutes.'),
    ('Inconsistent Siloed Analytics', PURPLE,
     '5 signal serializers normalize events from ANY source (ECS/OCSF/STIX 2.1) '
     'into one coherent pipeline.',
     'Entity resolution maps AD/AAD/Okta/AWS/K8s identities to single entity_uuid. '
     'All signals contribute to ONE embedding.',
     'Lateral movement across 5 identity systems = 1 strong signal, not 5 weak ones.'),
    ('Scalability Constraints', TEAL,
     'Containerized: Iron Bank base images, OCI-compliant, K8s Helm chart. '
     'Gabriel Nimbus native.',
     '10,000 entities/hour on 4-vCPU. Horizontal scaling via K8s replicas. '
     '<3 sec per entity window.',
     '1 unified pipeline handles ALL log types. New sources plug in via schema config, not code.'),
    ('Operational Inefficiency', ORANGE,
     'Fully automated: ingest → normalize → serialize → embed → '
     'analyze → alert → enforce. Zero manual wrangling.',
     'ABAC trust loop auto-restricts access proportionally: step-up MFA → '
     'read-only → block. <5 min from detection to enforcement.',
     'Analyst TP/FP feedback auto-tunes thresholds. System improves with use, not decays.'),
]

y_start = Inches(1.3)
for i, (title, color, l1, l2, l3) in enumerate(challenges_solve):
    y = y_start + Inches(i * 1.45)
    _box(s, Inches(0.3), y, Inches(12.7), Inches(1.35), fill=CARD, border=color)
    _box(s, Inches(0.35), y + Inches(0.05), Inches(2.8), Inches(0.3), fill=color)
    _txt(s, Inches(0.4), y + Inches(0.07), Inches(2.7), Inches(0.25),
         title, 11, WHITE, True, PP_ALIGN.CENTER)
    _txt(s, Inches(3.3), y + Inches(0.05), Inches(9.5), Inches(0.35),
         l1, 10, RGBColor(0xCC, 0xDD, 0xCC))
    _txt(s, Inches(3.3), y + Inches(0.42), Inches(9.5), Inches(0.35),
         l2, 10, RGBColor(0xBB, 0xCC, 0xDD))
    _txt(s, Inches(3.3), y + Inches(0.82), Inches(9.5), Inches(0.35),
         l3, 10, GOLD, True)
_snum(s, 5)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6: THE FIVE BEHAVIORAL SIGNALS
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_fill(s, BG)
_title_bar(s, 'Five Behavioral Signals → 6 Embeddings (9,216-d Total)',
           'Each signal individually embedded (1536-d) + one fused composite (1536-d) = per-signal drift analysis')

signals = [
    ('AUTH', BLUE,
     'Win 4624/4625/4768\nOkta, Azure AD',
     'logon_count, failure_rate, unique_hosts,\noff_hours_ratio, impossible_travel'),
    ('PROCESS', PURPLE,
     'Sysmon EID 1/3\nEDR, auditd',
     'unique_procs, lolbin_count, cmdline_entropy,\nparent_child_depth, unsigned_ratio'),
    ('NETWORK', TEAL,
     'NetFlow, Zeek\nPAN, Fortinet',
     'unique_dest_ips, bytes_out_ratio,\nbeacon_score, dns_rate, geo_anomaly'),
    ('FILE', ORANGE,
     'Sysmon EID 11\nFile audit, DLP',
     'files_created/deleted, sensitive_access,\narchive_creates, extension_changes'),
    ('IDENTITY', RED,
     'Win 4672, CloudTrail\nK8s audit, AD',
     'priv_escalations, group_adds,\nmfa_bypass, role_changes, admin_actions'),
]

x_positions = [Inches(0.3), Inches(2.9), Inches(5.5), Inches(8.1), Inches(10.7)]
for i, (label, color, sources, features) in enumerate(signals):
    x = x_positions[i]
    bw = Inches(2.4)
    _box(s, x, Inches(1.3), bw, Inches(0.38), fill=color)
    _txt(s, x, Inches(1.32), bw, Inches(0.35),
         f'Signal {i+1}: {label}', 12, WHITE, True, PP_ALIGN.CENTER)
    _box(s, x, Inches(1.75), bw, Inches(1.2), fill=CARD,
         border=RGBColor(0x33, 0x33, 0x55))
    _txt(s, x + Inches(0.08), Inches(1.78), bw - Inches(0.16), Inches(0.15),
         'SOURCES', 8, color, True)
    _txt(s, x + Inches(0.08), Inches(1.95), bw - Inches(0.16), Inches(0.4),
         sources, 9, RGBColor(0xAA, 0xBB, 0xCC))
    _box(s, x, Inches(3.05), bw, Inches(1.3), fill=CARD,
         border=RGBColor(0x33, 0x33, 0x55))
    _txt(s, x + Inches(0.08), Inches(3.08), bw - Inches(0.16), Inches(0.15),
         'FEATURES', 8, color, True)
    _txt(s, x + Inches(0.08), Inches(3.28), bw - Inches(0.16), Inches(1.0),
         features, 9, RGBColor(0xAA, 0xBB, 0xCC))

_box(s, Inches(0.3), Inches(4.6), Inches(12.7), Inches(0.65), fill=CARD, border=CYBER)
_txt(s, Inches(0.5), Inches(4.62), Inches(12.3), Inches(0.2),
     'DUAL EMBEDDING ARCHITECTURE', 10, CYBER, True)
_txt(s, Inches(0.5), Inches(4.8), Inches(12.3), Inches(0.45),
     'Each signal individually embedded → 5 x 1536-d vectors (per-signal drift tracking). '
     'All 5 signal texts also concatenated → 1 x 1536-d fused composite. '
     'Total: 6 vectors = 9,216 dimensions per (entity, window). '
     'No raw event content sent to provider — only summarized metrics. '
     'Local SLM (Phi-4/Mistral) for JWICS.',
     10, RGBColor(0xCC, 0xDD, 0xCC))

# CUSUM + Drift Direction summary
_box(s, Inches(0.3), Inches(5.3), Inches(6.2), Inches(1.8), fill=CARD, border=GOLD)
_txt(s, Inches(0.5), Inches(5.35), Inches(5.8), Inches(0.3),
     'Trajectory Analysis', 14, GOLD, True)
tf = _txt(s, Inches(0.5), Inches(5.7), Inches(5.8), Inches(1.3),
     '•  Cosine drift = 1 - cos(V_current, V_baseline)', 10,
     RGBColor(0xCC, 0xDD, 0xCC), font='Cascadia Code')
_para(tf, '•  Velocity = drift(t) - drift(t-1)', 10, RGBColor(0xCC, 0xDD, 0xCC))
_para(tf, '•  CUSUM: cumulative sum, alarm at 4σ', 10, RGBColor(0xCC, 0xDD, 0xCC))
_para(tf, '•  Health score: 0–100 (Critical <40)', 10, RGBColor(0xCC, 0xDD, 0xCC))

_box(s, Inches(6.7), Inches(5.3), Inches(6.3), Inches(1.8), fill=CARD, border=CYBER)
_txt(s, Inches(6.9), Inches(5.35), Inches(5.8), Inches(0.3),
     'Drift Direction (8 MITRE Concepts)', 14, CYBER, True)
tf = _txt(s, Inches(6.9), Inches(5.7), Inches(5.8), Inches(1.3),
     'projection = cosine(drift_vector, concept_embedding)', 10,
     RGBColor(0xCC, 0xDD, 0xCC), font='Cascadia Code')
_para(tf, 'credential_dumping (T1003) | lateral_movement (T1021)', 10, GOLD)
_para(tf, 'data_exfiltration (T1041) | c2_beaconing (T1071)', 10, GOLD)
_para(tf, 'lotl_execution (T1059) | privilege_escalation (T1078)', 10, GOLD)
_para(tf, 'defense_evasion (T1562) | insider_hoarding (T1074)', 10, GOLD)

_snum(s, 6)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7: VOLT TYPHOON WALKTHROUGH
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_fill(s, BG)
_title_bar(s, 'Volt Typhoon: Both Layers in Action',
           'Each phase: what Rigor preempts AND what ACECARD detects if preemption is bypassed')

cols = ['Phase', 'Attacker Behavior', 'Rigor Preempts', 'ACECARD Detects']
col_x = [Inches(0.3), Inches(2.2), Inches(5.5), Inches(9.3)]
col_w = [Inches(1.8), Inches(3.2), Inches(3.7), Inches(3.7)]

_box(s, Inches(0.25), Inches(1.25), Inches(12.8), Inches(0.38), fill=NAVY)
for i, h in enumerate(cols):
    _txt(s, col_x[i], Inches(1.27), col_w[i], Inches(0.34), h, 11, CYBER, True)

rows_vt = [
    ('Initial Access\n0–12h',
     'Compromised SOHO router → VPN credential reuse',
     'Zero-trust segmentation verified;\nMFA bypass paths surfaced (T1133)',
     'Auth: off_hours_ratio spikes,\nnew source IP — drift 0.08'),
    ('LOLBin Execution\n12–36h',
     'wmic, ntdsutil, netsh chain',
     'App-control policies verified;\nunsigned exec paths blocked (T1059)',
     'Process: lolbin_count 0→7,\ncmdline_entropy rises — drift 0.14'),
    ('Lateral Movement\n36–72h',
     'RDP + WinRM to domain controllers',
     'Micro-seg blocks DC access\nfrom user VLAN (T1021)',
     'Network: unique_dest 3→18,\nadmin shares — drift 0.19'),
    ('Persistence\n72–96h',
     'Scheduled tasks + registry mods',
     'Config drift detected;\nbaseline deviation flagged (T1053)',
     'Identity: svc_acct_use spikes\nCUSUM triggers at 4σ'),
    ('Pre-positioning\n96h+',
     'Credential harvest + network map',
     'Credential-guard policies\nverified on DCs (T1003)',
     'File: sensitive_access rises\nHealth: 28 (CRITICAL)'),
]

y = Inches(1.7)
for i, (phase, attack, rigor, acecard) in enumerate(rows_vt):
    bg = CARD if i % 2 == 0 else RGBColor(0x15, 0x15, 0x28)
    _box(s, Inches(0.25), y, Inches(12.8), Inches(0.95), fill=bg)
    _txt(s, col_x[0], y + Pt(2), col_w[0], Inches(0.9), phase, 10, GOLD, True)
    _txt(s, col_x[1], y + Pt(2), col_w[1], Inches(0.9), attack, 10, RGBColor(0xCC, 0xCC, 0xDD))
    _txt(s, col_x[2], y + Pt(2), col_w[2], Inches(0.9), rigor, 10, RGBColor(0xBB, 0xCC, 0xDD))
    _txt(s, col_x[3], y + Pt(2), col_w[3], Inches(0.9), acecard, 10, RGBColor(0xBB, 0xDD, 0xCC))
    y += Inches(0.98)

_snum(s, 7)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8: SALT TYPHOON WALKTHROUGH
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_fill(s, BG)
_title_bar(s, 'Salt Typhoon: Both Layers in Action',
           'Telecom-specific tradecraft. Rigor preempts + ACECARD detects.')

_box(s, Inches(0.25), Inches(1.25), Inches(12.8), Inches(0.38), fill=NAVY)
for i, h in enumerate(cols):
    _txt(s, col_x[i], Inches(1.27), col_w[i], Inches(0.34), h, 11, CYBER, True)

rows_st = [
    ('Edge CVE Exploit\n0–24h',
     'CVE-2023-20198 vs Cisco IOS XE;\nlevel-15 admin obtained',
     'WAF blocks /webui/ public access;\nmgmt-plane ACL verified (T1190)',
     'Network: anomalous mgmt-plane\nsource IP — drift 0.09'),
    ('ACL Tampering\n1–7d',
     'Modifies router ACLs + AAA;\ncreates persistent hidden access',
     'Config drift detected in <1 hour;\nbaseline deviation flagged (T1556)',
     'Auth: SNMP access pattern change;\nIdentity: new admin sessions'),
    ('GRE Tunneling\n7–30d',
     'GRE tunnel + encrypted C2\nchannels established',
     'Unauthorized tunnel endpoints\ndetected; egress verified (T1572)',
     'Network: beacon_score rises;\nnew protocol mix — CUSUM accumulating'),
    ('Wiretap Access\n30–90d',
     'CALEA lawful intercept\nsystem access',
     'Privilege verification on\ninterception systems (T1078)',
     'Identity: priv_esc on CALEA;\nadmin_actions spike — CUSUM triggers'),
    ('CDR Exfiltration\n90d+',
     'Call detail records stolen\nover months, slowly',
     'Outbound restriction policies\nverified for CDR stores',
     'File: steady bytes_out increase;\nHealth: 24 (CRITICAL)'),
]

y = Inches(1.7)
for i, (phase, attack, rigor, acecard) in enumerate(rows_st):
    bg = CARD if i % 2 == 0 else RGBColor(0x15, 0x15, 0x28)
    _box(s, Inches(0.25), y, Inches(12.8), Inches(0.95), fill=bg)
    _txt(s, col_x[0], y + Pt(2), col_w[0], Inches(0.9), phase, 10, GOLD, True)
    _txt(s, col_x[1], y + Pt(2), col_w[1], Inches(0.9), attack, 10, RGBColor(0xCC, 0xCC, 0xDD))
    _txt(s, col_x[2], y + Pt(2), col_w[2], Inches(0.9), rigor, 10, RGBColor(0xBB, 0xCC, 0xDD))
    _txt(s, col_x[3], y + Pt(2), col_w[3], Inches(0.9), acecard, 10, RGBColor(0xBB, 0xDD, 0xCC))
    y += Inches(0.98)

_snum(s, 8)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9: PROOF -- WORKING PROTOTYPE
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_fill(s, BG)
_title_bar(s, 'Proof: Working Prototype',
           'Not a concept. Not a proposal. A running system.')

_bullet_card(s, Inches(0.3), Inches(1.3), Inches(4.1), Inches(3.0),
    'Running Now', [
        'Docker-compose: PostgreSQL+pgvector, FastAPI, worker',
        '270 days synthetic event data',
        '6 injected attack scenarios (Volt/Salt Typhoon TTPs)',
        'Mock embedding pipeline: 5 signals → 1536-d vectors',
        'CUSUM detection + drift direction + MITRE mapping',
        'Kill-chain reconstruction (multi-phase narratives)',
    ], CYBER, CYBER)

_bullet_card(s, Inches(4.6), Inches(1.3), Inches(4.1), Inches(3.0),
    '7-Chart SOC Dashboard', [
        'Drift timeline (cosine drift per window)',
        'Velocity / acceleration (rate of change)',
        'Change-point timeline (CUSUM triggers)',
        'Health gauge (0–100)',
        'Drift-direction radar (8 MITRE concepts)',
        'Per-signal contribution (which signal drives drift)',
        'Peer-cohort scatter (outlier detection)',
    ], GOLD, GOLD)

_bullet_card(s, Inches(8.9), Inches(1.3), Inches(4.1), Inches(3.0),
    'Proven Foundation', [
        'E-11 Temporal Trajectory Intelligence at DLA',
        '500+ entities tracked over 23 monthly snapshots',
        '5 behavioral signals → CUSUM change-point detection',
        'Drift direction with concept alignment',
        'Same architecture, same algorithms',
        'Different signals: SCM → cyber defense',
    ], TEAL, TEAL)

# Rigor proof points
_box(s, Inches(0.3), Inches(4.5), Inches(12.7), Inches(2.4), fill=CARD, border=BLUE)
_txt(s, Inches(0.5), Inches(4.55), Inches(5), Inches(0.3),
     'Rigor AI — Production Findings', 14, BLUE, True)

findings = [
    ('FinTech | Asia', 'Closed SMB lateral movement path — 3 higher-priority rules shadowed app-control profiles'),
    ('Comms | NA', 'Eliminated CVE exposure — IPS filters missing on rsync traffic to backup servers'),
    ('Government', 'Found zombie rules — legacy permits for decommissioned systems still allowed traffic'),
    ('Critical Infra', 'Verified Zero Trust claims were incomplete — implicit trust in edge cases'),
]
y = Inches(4.95)
for customer, finding in findings:
    _txt(s, Inches(0.6), y, Inches(2.0), Inches(0.25), customer, 10, GOLD, True)
    _txt(s, Inches(2.7), y, Inches(10.0), Inches(0.25), finding, 10, RGBColor(0xBB, 0xCC, 0xDD))
    y += Inches(0.3)

_snum(s, 9)

# ═══════════════════════════════════════════════════════════════
# SLIDE 10: DEPLOYMENT READINESS
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_fill(s, BG)
_title_bar(s, 'Gabriel Nimbus Deployment Readiness',
           'Containerized, multi-enclave, cATO-aligned')

deploy_items = [
    ('Container', BLUE, 'Iron Bank-compatible base image, OCI-compliant, signed with SBOM, Big Bang Helm chart'),
    ('Enclaves', PURPLE, 'IL5 (NIPR) / IL6 (SIPR) / JWICS — replicated containers, no cross-domain data'),
    ('Embedding', TEAL, 'OpenAI text-embedding-3-small (IL5) | Local SLM: Phi-4 / Mistral for JWICS (env var swap, no code change)'),
    ('Performance', ORANGE, '<3s embed+analyze+interpret per entity window; 10K entities/hour on 4-vCPU; horizontal K8s scaling'),
    ('Auth', RED, 'Bearer token (dev) / CAC/PIV mTLS (production cluster); RBAC on API endpoints'),
    ('Schema', GOLD, 'ECS, OCSF, STIX 2.1 normalized input — new log sources plug in via schema config, not code'),
    ('Observability', CYBER, 'Structured JSON logs, Prometheus metrics, OpenTelemetry traces end-to-end'),
    ('CI/CD', LTGRAY, 'FluxCD GitOps, SBOM generation, vulnerability scanning (Anchore/Grype), cATO DoD CIO Feb 2022'),
]

y = Inches(1.4)
for label, color, desc in deploy_items:
    _box(s, Inches(0.4), y, Inches(12.5), Inches(0.62), fill=CARD, border=RGBColor(0x33, 0x33, 0x55))
    _box(s, Inches(0.45), y + Inches(0.05), Inches(1.6), Inches(0.28), fill=color)
    _txt(s, Inches(0.5), y + Inches(0.07), Inches(1.5), Inches(0.24),
         label, 11, WHITE, True, PP_ALIGN.CENTER)
    _txt(s, Inches(2.2), y + Inches(0.07), Inches(10.5), Inches(0.5),
         desc, 11, RGBColor(0xBB, 0xCC, 0xDD))
    y += Inches(0.66)

_snum(s, 10)

# ═══════════════════════════════════════════════════════════════
# SLIDE 11: TEAM CREDENTIALS
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_fill(s, BG)
_title_bar(s, 'Team Credentials',
           'Operational experience at scale in DoD cyber operations')

_box(s, Inches(0.3), Inches(1.3), Inches(6.2), Inches(4.5), fill=CARD, border=GOLD)
_txt(s, Inches(0.5), Inches(1.35), Inches(5.8), Inches(0.35),
     '22nd Century Technologies (TSCTI)', 18, GOLD, True)

contracts = [
    ('U.S. Army SOC/MDR', '$90M', '800+ cleared analysts'),
    ('U.S. Air Force', '$108M', 'Cybersecurity operations'),
    ('FBI TSC', '$56M', 'Cybersecurity support'),
    ('NAVAIR / USMC', '$145M', 'Cybersecurity operations'),
]
y = Inches(1.9)
for name, value, desc in contracts:
    _txt(s, Inches(0.6), y, Inches(2.5), Inches(0.3), name, 13, WHITE, True)
    _txt(s, Inches(3.2), y, Inches(1.2), Inches(0.3), value, 13, CYBER, True)
    _txt(s, Inches(4.5), y, Inches(1.8), Inches(0.3), desc, 11, LTGRAY)
    y += Inches(0.4)

_txt(s, Inches(0.6), Inches(3.6), Inches(5.7), Inches(1.5),
     'ACECARD UEBA is built by the team that operates SOC/MDR at scale '
     'for the U.S. military. This is not academic research — it is built '
     'by practitioners who triage alerts, investigate incidents, and '
     'respond to threats every day.',
     11, RGBColor(0xBB, 0xCC, 0xDD))

_box(s, Inches(6.8), Inches(1.3), Inches(6.2), Inches(4.5), fill=CARD, border=BLUE)
_txt(s, Inches(7.0), Inches(1.35), Inches(5.8), Inches(0.35),
     'Rigor AI Inc.', 18, BLUE, True)

_txt(s, Inches(7.0), Inches(1.9), Inches(5.8), Inches(3.5),
     '20+ design partners across:\n\n'
     '•  Fortune 500 enterprises\n'
     '•  Federal agencies\n'
     '•  Critical infrastructure operators\n'
     '•  MSSPs (Managed Security Service Providers)\n'
     '•  Nation states\n\n'
     'Production findings in financial services, telecom, '
     'and government. Formal verification engine with real '
     'customer results: closed SMB lateral movement paths, '
     'eliminated CVE exposure gaps, found zombie firewall '
     'rules, verified Zero Trust claims.',
     12, RGBColor(0xBB, 0xCC, 0xDD))

_snum(s, 11)

# ═══════════════════════════════════════════════════════════════
# SLIDE 12: CLOSING
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_fill(s, BG)
_box(s, Inches(0), Inches(0), W, Inches(0.06), fill=CYBER)

_txt(s, Inches(1), Inches(1.5), Inches(11), Inches(0.5),
     'We Have a Working Solution.', 36, CYBER, True, PP_ALIGN.CENTER)
_txt(s, Inches(1), Inches(2.2), Inches(11), Inches(0.5),
     'We Can Demo It Today.', 28, GOLD, True, PP_ALIGN.CENTER)

_box(s, Inches(2), Inches(3.2), Inches(9.333), Inches(0.02),
     fill=RGBColor(0x33, 0x33, 0x55))

_txt(s, Inches(1), Inches(3.5), Inches(11), Inches(1.5),
     'Rigor AI mathematically proves your configuration closes every '
     'known TTP path — before exploitation.\n\n'
     'ACECARD UEBA detects behavioral drift when Living-off-the-Land '
     'tradecraft slips through — with MITRE-mapped direction, not just '
     'anomaly scores.\n\n'
     'Together: preempt what can be prevented, detect what cannot, '
     'enforce proportionally via ABAC.',
     16, RGBColor(0xCC, 0xDD, 0xCC), align=PP_ALIGN.CENTER)

_box(s, Inches(2), Inches(5.5), Inches(9.333), Inches(0.02),
     fill=RGBColor(0x33, 0x33, 0x55))

_txt(s, Inches(1), Inches(5.8), Inches(11), Inches(0.4),
     'Containerized  |  Gabriel Nimbus Ready  |  IL5/IL6/JWICS  |  cATO-Aligned',
     14, LTGRAY, False, PP_ALIGN.CENTER)

_txt(s, Inches(1), Inches(6.4), Inches(11), Inches(0.4),
     '22nd Century Technologies  |  Rigor AI Inc.', 13,
     RGBColor(0x88, 0x88, 0xAA), False, PP_ALIGN.CENTER)

_box(s, Inches(0), Inches(7.44), W, Inches(0.06), fill=CYBER)
_snum(s, 12)

# ═══════════════════════════════════════════════════════════════
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
