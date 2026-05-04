"""Expanded 22CT Preemptive + ACECARD UEBA deck -- 22nd Century Technologies Inc.
is the solution owner. Matches the look/feel of the reference deck.
Adds deep detail on metrics, thresholds, algorithms, entity fusion,
4-component signature, CUSUM, drift direction, and innovation proof."""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

OUT = os.path.join(os.path.dirname(__file__),
                   "22CT_Preemptive_ACECARD_Volt_Salt_Typhoon_v2.pptx")

# === EXACT COLOR PALETTE from reference deck ===
DKNAVY  = RGBColor(0x07, 0x14, 0x2A)  # title/closing bg
NAVY    = RGBColor(0x0B, 0x1F, 0x3A)  # titles, bottom bars, dark accents
BODY    = RGBColor(0x2D, 0x37, 0x48)  # main body text
SUB     = RGBColor(0x5A, 0x68, 0x78)  # subtitle text
ORG     = RGBColor(0xF2, 0x6B, 0x1F)  # Part 1 bar, solution labels, phase labels
TEAL    = RGBColor(0x08, 0x91, 0xB2)  # Part 2 bar
CRED    = RGBColor(0xC0, 0x39, 0x2B)  # problem labels, consequence, red accents
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
CARD    = RGBColor(0xF7, 0xF8, 0xFA)  # card backgrounds
RIGOR_BG = RGBColor(0xFD, 0xE5, 0xD4) # Rigor card background (light peach)
ACE_BG  = RGBColor(0xE0, 0xF2, 0xF7)  # ACECARD card background (light cyan)
GOLD    = RGBColor(0xB8, 0x86, 0x2A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def _fill(s, c):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = c

def _rect(s, l, t, w, h, fill=None):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.background()
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def _box(s, l, t, w, h, fill=None, border=None, accent=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.background()
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if border:
        sh.line.color.rgb = border
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    if accent:
        _rect(s, l, t, Inches(0.1), h, fill=accent)
    return sh

def _txt(s, l, t, w, h, text, sz=11, c=BODY, b=False, al=PP_ALIGN.LEFT, fn='Calibri'):
    tf = s.shapes.add_textbox(l, t, w, h).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = c
    p.font.bold = b
    p.font.name = fn
    p.alignment = al
    return tf

def _p(tf, text, sz=11, c=BODY, b=False, sp=Pt(4)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = c
    p.font.bold = b
    p.font.name = 'Calibri'
    p.space_before = sp
    return p

def _title_slide(title, subtitle, bottom=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _fill(s, DKNAVY)
    _txt(s, Inches(1), Inches(2.2), Inches(11.3), Inches(1.0),
         title, 36, WHITE, True, PP_ALIGN.CENTER, 'Georgia')
    _txt(s, Inches(1), Inches(3.5), Inches(11.3), Inches(0.5),
         subtitle, 16, RGBColor(0xA0, 0xB0, 0xC0), False, PP_ALIGN.CENTER)
    if bottom:
        _txt(s, Inches(1), Inches(6.8), Inches(11.3), Inches(0.4),
             bottom, 11, RGBColor(0x70, 0x80, 0x90), False, PP_ALIGN.CENTER)
    return s

def _content_slide(part_text, title, subtitle=None, bar_color=ORG):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _fill(s, WHITE)
    _rect(s, Inches(0), Inches(0), W, Inches(0.3), fill=bar_color)
    _txt(s, Inches(0.5), Inches(0.0), Inches(12.3), Inches(0.3),
         part_text, 11, WHITE, True)
    _txt(s, Inches(0.5), Inches(0.5), Inches(12.3), Inches(0.6),
         title, 24, NAVY, True, fn='Georgia')
    if subtitle:
        _txt(s, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.3),
             subtitle, 13, SUB)
    return s

def _bottom_bar(s, text, y=6.8):
    _rect(s, Inches(0.5), Inches(y), Inches(12.3), Inches(0.5), fill=NAVY)
    _txt(s, Inches(0.7), Inches(y), Inches(12.0), Inches(0.5),
         text, 11.5, GOLD, True, PP_ALIGN.CENTER)


# ================================================================
# SLIDE 1: TITLE
# ================================================================
s = _title_slide(
    '22nd Century Technologies Inc.\nCOMPREHENSIVE DEFENSE AGAINST\nVOLT & SALT TYPHOON',
    'Preemptive Cyber Defense  +  ACECARD UEBA',
    '22nd Century Technologies Inc. (TSCTI)  |  CONFIDENTIAL')
_txt(s, Inches(1), Inches(4.8), Inches(5.5), Inches(1.5),
     'PREEMPTIVE LAYER\nMathematically verifies every defensive control\ncloses every known TTP',
     12, RGBColor(0xCC, 0xDD, 0xEE))
_txt(s, Inches(7), Inches(4.8), Inches(5.5), Inches(1.5),
     'BEHAVIORAL LAYER\nACECARD UEBA\nContinuous 1536-d behavioral trajectories\nwith MITRE-mapped drift direction',
     12, RGBColor(0xCC, 0xDD, 0xEE))


# ================================================================
# SLIDE 2: THE THREAT
# ================================================================
s = _content_slide('THREAT LANDSCAPE', 'The Threat: Two Apex Predators in U.S. Critical Infrastructure',
                   'Both groups are PRC state-sponsored. Both prepare for disruption. Both walked past commercial defenses undetected.',
                   CRED)

_box(s, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.5), fill=CARD, accent=ORG)
_txt(s, Inches(0.8), Inches(1.55), Inches(5.3), Inches(0.5),
     'VOLT TYPHOON  (CISA AA24-038A)', 16, NAVY, True, fn='Georgia')
_txt(s, Inches(0.8), Inches(2.1), Inches(5.3), Inches(0.3), 'MISSION', 10, ORG, True)
tf = _txt(s, Inches(0.8), Inches(2.4), Inches(5.3), Inches(0.7),
     'Pre-position in U.S. critical infrastructure (energy, water, transport, comms) for sabotage during geopolitical crisis.', 11, BODY)
_txt(s, Inches(0.8), Inches(3.2), Inches(5.3), Inches(0.3), 'SIGNATURE TTPs', 10, ORG, True)
tf = _txt(s, Inches(0.8), Inches(3.5), Inches(5.3), Inches(3.0),
     'Initial access via unpatched edge devices (Fortinet, Versa, Ivanti, PAN)', 10, BODY)
_p(tf, 'Living-off-the-Land: wmic, ntdsutil, netsh, rundll32, certutil', 10, BODY)
_p(tf, 'KV Botnet C2 (compromised SOHO routers)', 10, BODY)
_p(tf, 'Credential dumping: ntdsutil ifm, LSASS extraction', 10, BODY)
_p(tf, 'Lateral: SMB/RDP/WinRM with stolen admin credentials', 10, BODY)
_p(tf, 'Persistence: 5+ YEARS undetected dwell', 10, CRED, True)
_p(tf, '', 8, BODY)
_p(tf, 'No malware. No signatures. No IOCs.', 10, CRED, True)
_p(tf, 'Every tool used is a legitimate Windows binary.', 10, CRED, True)

_box(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.5), fill=CARD, accent=TEAL)
_txt(s, Inches(7.1), Inches(1.55), Inches(5.3), Inches(0.5),
     'SALT TYPHOON  (CISA AA25-239A)', 16, NAVY, True, fn='Georgia')
_txt(s, Inches(7.1), Inches(2.1), Inches(5.3), Inches(0.3), 'MISSION', 10, TEAL, True)
tf = _txt(s, Inches(7.1), Inches(2.4), Inches(5.3), Inches(0.7),
     'Espionage at carrier scale. Compromised 200+ telcos in 80+ countries. Accessed CALEA lawful intercept.', 11, BODY)
_txt(s, Inches(7.1), Inches(3.2), Inches(5.3), Inches(0.3), 'SIGNATURE TTPs', 10, TEAL, True)
tf = _txt(s, Inches(7.1), Inches(3.5), Inches(5.3), Inches(3.0),
     'Exploits known CVEs: Cisco IOS XE (CVE-2023-20198), Ivanti, PAN, Citrix', 10, BODY)
_p(tf, 'Creates local admin accounts on routers', 10, BODY)
_p(tf, 'Modifies ACLs to whitelist attacker IPs', 10, BODY)
_p(tf, 'TACACS+ credential capture (TCP/49)', 10, BODY)
_p(tf, 'GRE tunneling + JumbledPath malware for C2', 10, BODY)
_p(tf, 'CALEA wiretap access + CDR exfiltration', 10, BODY)
_p(tf, 'GhostSpider, Demodex, SnappyBee implants', 10, BODY)
_p(tf, '', 8, BODY)
_p(tf, '4+ YEARS undetected. 9 US telcos confirmed.', 10, CRED, True)

_txt(s, Inches(0.5), Inches(7.2), Inches(12.3), Inches(0.2),
     'Sources: CISA AA24-038A; CISA AA25-239A; FBI/NSA/Five-Eyes Joint Advisory; Microsoft Threat Intelligence; Trend Micro',
     9, SUB)


# ================================================================
# SLIDE 3: WHY CURRENT DEFENSES FAIL
# ================================================================
s = _content_slide('THE GAP', 'Why Current Defenses Fail Against These Adversaries',
                   'Two structural gaps -- one in preemption, one in detection -- that traditional tooling cannot close.',
                   CRED)

_box(s, Inches(0.5), Inches(1.5), Inches(5.8), Inches(2.3), fill=CARD, accent=CRED)
_txt(s, Inches(0.8), Inches(1.55), Inches(5.3), Inches(0.3), 'PREEMPTIVE LAYER FAILS', 12, CRED, True)
tf = _txt(s, Inches(0.8), Inches(1.9), Inches(5.3), Inches(1.8),
     '99% of firewall breaches are misconfiguration (Gartner)', 10, NAVY, True)
_p(tf, 'KEV patching covers only ~20% of attacks', 10, NAVY, True)
_p(tf, 'Pentest: point-in-time, incomplete, expensive', 10, BODY)
_p(tf, 'ASM: sees what exists, not what is REACHABLE', 10, BODY)
_p(tf, 'Config review: manual, stale within hours of completion', 10, BODY)
_p(tf, '2^8000 rule state space cannot be sampled adequately', 10, BODY)

_box(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(2.3), fill=CARD, accent=CRED)
_txt(s, Inches(7.1), Inches(1.55), Inches(5.3), Inches(0.3), 'DETECTION LAYER FAILS', 12, CRED, True)
tf = _txt(s, Inches(7.1), Inches(1.9), Inches(5.3), Inches(1.8),
     'Threshold-based UEBA: LoTL never crosses fixed thresholds', 10, NAVY, True)
_p(tf, 'Sigma rules: adversary downloads published rules, retool in 2 weeks', 10, NAVY, True)
_p(tf, 'Per-source silos: each system sees 1 identity, no cross-domain correlation', 10, BODY)
_p(tf, 'Anomaly score without direction: "87% anomalous" -- anomalous HOW?', 10, BODY)
_p(tf, '30-45 min analyst triage per alert (MTTR > dwell window)', 10, BODY)
_p(tf, 'No behavioral TRAJECTORY -- only point-in-time snapshots', 10, BODY)

_box(s, Inches(0.5), Inches(4.0), Inches(12.1), Inches(3.0), fill=CARD, accent=NAVY)
_txt(s, Inches(0.8), Inches(4.05), Inches(11.6), Inches(0.3),
     'WHAT THESE ADVERSARIES EXPLOIT', 14, NAVY, True)
gaps = [
    ('Configuration gaps exist that no one has PROVEN are closed', 'Both Typhoons enter via unpatched/misconfigured edge devices'),
    ('Valid credentials are indistinguishable from legitimate use', 'Stolen creds pass every config-based check by design'),
    ('LoTL tools are allowed by policy', 'wmic/ntdsutil/netsh are legitimate admin tools -- cannot be blanket-blocked'),
    ('Slow drift is invisible to fixed thresholds', '5+ year campaigns operate below every alert threshold'),
    ('Identity fragmentation hides cross-domain attack chains', 'One attacker appears as 10 unrelated identities'),
]
y = Inches(4.4)
for gap, typhoon in gaps:
    _txt(s, Inches(0.8), y, Inches(6.0), Inches(0.3), gap, 10, NAVY, True)
    _txt(s, Inches(6.9), y, Inches(5.5), Inches(0.3), typhoon, 10, SUB)
    y += Inches(0.42)

_bottom_bar(s, 'RESULT: $0 invested in tools that could have detected either campaign before disclosure.')


# ================================================================
# SLIDE 4: CYBERSECURITY LIFECYCLE — PREEMPTIVE vs REACTIVE
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_fill(s, WHITE)
_rect(s, Inches(0), Inches(0), W, Inches(0.3), fill=NAVY)
_txt(s, Inches(0.5), Inches(0.0), Inches(12.3), Inches(0.3),
     'THE CYBERSECURITY LIFECYCLE  --  PREEMPTIVE vs. REACTIVE', 11, WHITE, True)

_txt(s, Inches(0.5), Inches(0.45), Inches(12.3), Inches(0.5),
     'The Cybersecurity Lifecycle: Preemptive vs. Reactive', 24, NAVY, True, fn='Georgia')
_txt(s, Inches(0.5), Inches(0.9), Inches(12.3), Inches(0.3),
     'Where attacks are neutralized -- and where preemptive verification fills the gap no reactive tool can address.', 12, SUB)

# --- STAGE 1: Preempt & Prove (green border area) ---
GREEN = RGBColor(0x16, 0x6B, 0x3A)
LGREEN = RGBColor(0xE8, 0xF5, 0xE9)

_box(s, Inches(0.3), Inches(1.3), Inches(9.5), Inches(2.8), fill=LGREEN, border=GREEN)
_rect(s, Inches(0.5), Inches(1.35), Inches(3.0), Inches(0.3), fill=GREEN)
_txt(s, Inches(0.55), Inches(1.37), Inches(2.9), Inches(0.26), 'Stage 1: Preempt & Prove', 10, WHITE, True)
_txt(s, Inches(3.7), Inches(1.37), Inches(5.5), Inches(0.26),
     'The missing layer -- before any traffic flows', 9, GREEN)

# ~85% badge
_rect(s, Inches(0.6), Inches(1.8), Inches(1.8), Inches(1.8), fill=GREEN)
_txt(s, Inches(0.6), Inches(2.1), Inches(1.8), Inches(0.6), '~85%', 22, WHITE, True, PP_ALIGN.CENTER, 'Georgia')
_txt(s, Inches(0.6), Inches(2.7), Inches(1.8), Inches(0.5),
     'of attack surface\neliminated', 8, WHITE, False, PP_ALIGN.CENTER)

# Preemptive box details
_box(s, Inches(2.6), Inches(1.8), Inches(7.0), Inches(2.2), fill=CARD, border=GREEN)
_txt(s, Inches(2.8), Inches(1.82), Inches(6.6), Inches(0.3),
     'Preemptive posture verification -- formal mathematical defense', 11, GREEN, True)
tf = _txt(s, Inches(2.8), Inches(2.1), Inches(6.6), Inches(1.8),
     'Attack Intelligence: LLMs + NLP ingest threat intel (CVEs, IOCs, IOAs) to build MITRE ATT&CK attack graphs', 8.5, BODY)
_p(tf, 'Defense Intelligence: Formal math model of all NGFW configs (100K+ states, 2^8000 defense space) -- no sampling', 8.5, BODY)
_p(tf, 'Remediation Intelligence: Guardrailed agentic-AI prescribes vendor-specific, risk-prioritized fixes -- no new errors', 8.5, BODY)
_p(tf, 'Finds: shadows, conflicts, weak rules, redundancies, missing threat profiles, intent errors, multi-rule/device errors', 8.5, NAVY, True)
_p(tf, '', 4, SUB)
_p(tf, '99% of firewall breaches are caused by misconfiguration, not flaws in the firewall itself (Gartner)', 8, CRED)
_p(tf, 'No telemetry  |  No agents  |  Config-only pull  |  Zero false positives  |  Time to value < 4 weeks', 8, GREEN, True)

# Connector label
_txt(s, Inches(5.0), Inches(4.0), Inches(4.5), Inches(0.2),
     'Pushes verified, flawless rules down to...', 8, GREEN, True)

# Reactive label
_txt(s, Inches(0.3), Inches(4.2), Inches(9.5), Inches(0.15),
     'The live environment -- all tools below are reactive (detect & respond after breach)', 7.5, SUB, False, PP_ALIGN.CENTER)

# --- STAGE 2: Prevent (perimeter) ---
_rect(s, Inches(0.3), Inches(4.4), Inches(9.5), Inches(0.3), fill=ORG)
_txt(s, Inches(0.5), Inches(4.42), Inches(5.0), Inches(0.26),
     'Stage 2: Prevent (perimeter)', 10, WHITE, True)

s2_tools = [
    (Inches(0.5), 'NGFW', 'Palo Alto, Fortinet, Cisco'),
    (Inches(2.6), 'Vuln Scanners', 'Qualys, Tenable'),
    (Inches(4.7), 'IAM / PAM', 'CyberArk, Okta'),
    (Inches(6.8), 'ZTNA / CASB', 'Zscaler, Netskope'),
]
for x, name, vendors in s2_tools:
    _box(s, x, Inches(4.75), Inches(1.9), Inches(0.5), fill=CARD, border=ORG)
    _txt(s, x+Inches(0.1), Inches(4.77), Inches(1.7), Inches(0.22), name, 9, ORG, True)
    _txt(s, x+Inches(0.1), Inches(4.97), Inches(1.7), Inches(0.22), vendors, 7.5, SUB)

# ~10% badge
_rect(s, Inches(9.2), Inches(4.4), Inches(1.2), Inches(0.8), fill=ORG)
_txt(s, Inches(9.2), Inches(4.45), Inches(1.2), Inches(0.35), '~10%', 16, WHITE, True, PP_ALIGN.CENTER, 'Georgia')
_txt(s, Inches(9.2), Inches(4.85), Inches(1.2), Inches(0.3), 'blocked at\nperimeter', 7, WHITE, False, PP_ALIGN.CENTER)

_txt(s, Inches(3.5), Inches(5.3), Inches(6.0), Inches(0.15),
     'If attacker bypasses perimeter... (avg. breakout: 29 min)', 8, SUB)

# --- STAGE 3: Detect & Respond ---
_rect(s, Inches(0.3), Inches(5.5), Inches(9.5), Inches(0.3), fill=TEAL)
_txt(s, Inches(0.5), Inches(5.52), Inches(6.0), Inches(0.26),
     'Stage 3: Detect & respond (endpoints & network)', 10, WHITE, True)

s3_tools = [
    (Inches(0.5), 'EDR', 'CrowdStrike, Defender'),
    (Inches(2.6), 'NDR', 'Darktrace, ExtraHop'),
    (Inches(4.7), 'UEBA', 'Behavioral analytics'),
    (Inches(6.8), 'DLP', 'Symantec, Forcepoint'),
]
for x, name, vendors in s3_tools:
    _box(s, x, Inches(5.85), Inches(1.9), Inches(0.5), fill=CARD, border=TEAL)
    _txt(s, x+Inches(0.1), Inches(5.87), Inches(1.7), Inches(0.22), name, 9, TEAL, True)
    _txt(s, x+Inches(0.1), Inches(6.07), Inches(1.7), Inches(0.22), vendors, 7.5, SUB)

# ~4% badge
_rect(s, Inches(9.2), Inches(5.5), Inches(1.2), Inches(0.8), fill=TEAL)
_txt(s, Inches(9.2), Inches(5.55), Inches(1.2), Inches(0.35), '~4%', 16, WHITE, True, PP_ALIGN.CENTER, 'Georgia')
_txt(s, Inches(9.2), Inches(5.95), Inches(1.2), Inches(0.3), 'caught at\nendpoints', 7, WHITE, False, PP_ALIGN.CENTER)

# Connector label
_txt(s, Inches(3.5), Inches(6.4), Inches(5.0), Inches(0.15),
     'Sends alerts and logs...', 8, SUB)

# --- STAGE 4: SOC ---
_rect(s, Inches(0.3), Inches(6.55), Inches(9.5), Inches(0.3), fill=CRED)
_txt(s, Inches(0.5), Inches(6.57), Inches(6.0), Inches(0.26),
     'Stage 4: Visibility & triage (the SOC)', 10, WHITE, True)

s4_tools = [
    (Inches(0.5), 'SIEM', 'Splunk, Sentinel'),
    (Inches(2.6), 'SOAR', 'XSOAR, Splunk SOAR'),
]
for x, name, vendors in s4_tools:
    _box(s, x, Inches(6.88), Inches(1.9), Inches(0.4), fill=CARD, border=CRED)
    _txt(s, x+Inches(0.1), Inches(6.9), Inches(1.7), Inches(0.18), name, 9, CRED, True)
    _txt(s, x+Inches(0.1), Inches(7.06), Inches(1.7), Inches(0.18), vendors, 7.5, SUB)

# ~1% badge
_rect(s, Inches(9.2), Inches(6.55), Inches(1.2), Inches(0.7), fill=CRED)
_txt(s, Inches(9.2), Inches(6.6), Inches(1.2), Inches(0.35), '~1%', 16, WHITE, True, PP_ALIGN.CENTER, 'Georgia')
_txt(s, Inches(9.2), Inches(6.95), Inches(1.2), Inches(0.2), 'triaged by\nSOC analysts', 7, WHITE, False, PP_ALIGN.CENTER)

# --- Right side: The Takeaway ---
_box(s, Inches(10.6), Inches(1.3), Inches(2.5), Inches(5.9), fill=CARD, accent=GREEN)
_txt(s, Inches(10.8), Inches(1.35), Inches(2.1), Inches(0.3), 'THE TAKEAWAY', 11, GREEN, True)
tf = _txt(s, Inches(10.8), Inches(1.7), Inches(2.1), Inches(5.3),
     'CrowdStrike, Defender, Splunk are reactive -- alarms that go off when the burglar is already inside.', 9, BODY)
_p(tf, '', 4, SUB)
_p(tf, 'We ensure the locks on the doors (NGFWs) were actually installed correctly in the first place.', 9, NAVY, True)
_p(tf, '', 4, SUB)
_p(tf, 'Mathematically guaranteeing the burglar can never get in.', 9, GREEN, True)
_p(tf, '', 4, SUB)
_p(tf, 'We find what no other tool can:', 9, BODY)
_p(tf, '- Shadowed rules', 8.5, BODY)
_p(tf, '- Multi-rule conflicts', 8.5, BODY)
_p(tf, '- Missing threat profiles', 8.5, BODY)
_p(tf, '- Intent errors in 100K+ config states', 8.5, BODY)
_p(tf, '', 4, SUB)
_p(tf, 'Gartner predicts: by 2030, preemptive cyber will be 50% of IT security spending, up from <5% in 2024.', 8, GOLD, True)

_txt(s, Inches(0.5), Inches(7.3), Inches(9.0), Inches(0.15),
     'Gartner predicts: by 2030, preemptive cybersecurity will account for 50% of IT security spending, up from <5% in 2024.',
     7.5, SUB)


# ================================================================
# SLIDE 5: TWO-LAYER ANSWER
# ================================================================
s = _content_slide('THE ANSWER', 'The Two-Layer Answer: Preemptive + Behavioral',
                   '22CT Preemptive shrinks the attack surface mathematically. ACECARD detects what remains behaviorally.',
                   NAVY)

_box(s, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.7), fill=RIGOR_BG)
_rect(s, Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.5), fill=ORG)
_txt(s, Inches(0.7), Inches(1.52), Inches(5.4), Inches(0.45),
     'LAYER 1 -- 22CT PREEMPTIVE', 14, WHITE, True)
_txt(s, Inches(0.8), Inches(2.1), Inches(5.3), Inches(0.3), 'WHAT IT DOES', 10, ORG, True)
tf = _txt(s, Inches(0.8), Inches(2.4), Inches(5.3), Inches(2.5),
     'Builds a formal mathematical model of every NGFW, IPS, IdP, SASE, and WAF. Reasons exhaustively over 2^8000 state space to PROVE no attack path exists.', 11, BODY)
_p(tf, '', 8, BODY)
_p(tf, 'WHAT IT PROVES:', 10, ORG, True)
_p(tf, 'Every known CVE path is blocked (or compensating control identified)', 10, BODY)
_p(tf, 'Every threat-feed C2 destination is denied', 10, BODY)
_p(tf, 'Config drift is caught within 1 hour', 10, BODY)
_p(tf, 'Segmentation intent holds end-to-end', 10, BODY)
_p(tf, '', 8, BODY)
_p(tf, 'MATHEMATICAL GUARANTEE:', 10, NAVY, True)
_p(tf, 'No false positives. No false negatives.', 11, NAVY, True)
_p(tf, 'Complete reasoning. Not sampling.', 11, NAVY, True)

_box(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.7), fill=ACE_BG)
_rect(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(0.5), fill=TEAL)
_txt(s, Inches(7.0), Inches(1.52), Inches(5.4), Inches(0.45),
     'LAYER 2 -- ACECARD UEBA (BEHAVIORAL)', 14, WHITE, True)
_txt(s, Inches(7.1), Inches(2.1), Inches(5.3), Inches(0.3), 'WHAT IT DOES', 10, TEAL, True)
tf = _txt(s, Inches(7.1), Inches(2.4), Inches(5.3), Inches(2.5),
     'Embeds 5 behavioral signals into 1536-d vectors. Tracks behavioral TRAJECTORY. Detects drift toward MITRE-mapped threat concepts.', 11, BODY)
_p(tf, '', 8, BODY)
_p(tf, 'WHAT IT CATCHES:', 10, TEAL, True)
_p(tf, 'Compromised valid credentials (T1078)', 10, BODY)
_p(tf, 'Living-off-the-Land chains (T1059)', 10, BODY)
_p(tf, 'Slow APT drift (CUSUM catches sub-threshold accumulation)', 10, BODY)
_p(tf, 'Cross-domain lateral movement via entity fusion', 10, BODY)
_p(tf, '', 8, BODY)
_p(tf, 'STRUCTURAL INNOVATION:', 10, NAVY, True)
_p(tf, 'Not what is anomalous -- what they are BECOMING.', 11, NAVY, True)
_p(tf, 'Direction, not just magnitude.', 11, NAVY, True)

_bottom_bar(s, 'Together: either the attack surface never existed (Preemptive) OR the attack is detected within hours (ACECARD + ABAC)')


# ================================================================
# PART 1: 22CT PREEMPTIVE SLIDES
# ================================================================

# SLIDE 5: Preemptive Approach
s = _content_slide('PART 1  --  22CT PREEMPTIVE LAYER',
                   'Close the Door Before They Knock',
                   'Mathematically rigorous, preemptive cyber defense -- Complete, Continuous, Verifiable')

cols = [(Inches(0.5), 'Analyze All Known Threat Intelligence',
         'Terabytes of CVE, advisories, CTI, IOCs, IOAs, PSIRTs',
         'Attack Intelligence',
         'MITRE-enriched attack graphs built by ML/NLP/LLM pipeline. Full TTP coverage of Volt Typhoon (G1017) and Salt Typhoon. Auto-updates on new advisory publication.'),
        (Inches(4.8), 'Identify Every Defensive Control Gap',
         'Exponential combination of NGFW, IdP, IPS, SASE, WAF configs',
         'Defense Intelligence',
         'Symbolic Model of Computation. Formal mathematical model -- exhaustively reasons over the 2^8000 state space. No sampling. Complete reasoning. Every possible path analyzed.'),
        (Inches(9.1), 'Prescribe Correct Configurations',
         'Without creating new errors or business disruption',
         'Remediation Intelligence',
         'Grounded agentic-AI reasoners. Risk-prioritized, vendor-specific fixes. No false positives (mathematically proven). No false negatives (exhaustive analysis).')]

for x, prob, prob_desc, sol, sol_desc in cols:
    _box(s, x, Inches(1.5), Inches(4.0), Inches(2.0), fill=CARD, accent=CRED)
    _txt(s, x+Inches(0.3), Inches(1.55), Inches(3.5), Inches(0.25), 'PROBLEM', 10, CRED, True)
    _txt(s, x+Inches(0.3), Inches(1.8), Inches(3.5), Inches(0.5), prob, 13, NAVY, True)
    _txt(s, x+Inches(0.3), Inches(2.4), Inches(3.5), Inches(0.8), prob_desc, 10, BODY)

    _box(s, x, Inches(3.8), Inches(4.0), Inches(2.7), fill=CARD, accent=ORG)
    _txt(s, x+Inches(0.3), Inches(3.85), Inches(3.5), Inches(0.25), 'SOLUTION', 10, ORG, True)
    _txt(s, x+Inches(0.3), Inches(4.1), Inches(3.5), Inches(0.4), sol, 13, NAVY, True)
    _txt(s, x+Inches(0.3), Inches(4.6), Inches(3.5), Inches(1.7), sol_desc, 10.5, BODY)

_bottom_bar(s, "The world's first Mathematically Rigorous -- Complete, Preemptive, Continuous, Verifiable -- Cyber Defense Platform.  |  22nd Century Technologies Inc.")


# SLIDE A: Operational Flow
s = _content_slide('PART 1  --  22CT PREEMPTIVE  --  OPERATIONAL FLOW',
                   'How the Preemptive Platform Operates: Inputs -> Reasoning -> Continuous Remediation',
                   'SaaS, on-prem, and air-gapped deployments. Integrates with ServiceNow, JIRA, and major SOC ticketing.')

# Column headers
for x, label, color in [(Inches(0.5), 'INPUTS', NAVY),
                         (Inches(4.7), 'PREEMPTIVE ENGINE', ORG),
                         (Inches(9.2), 'CONTINUOUS PREEMPTIVE ASSURANCE', TEAL)]:
    _rect(s, x, Inches(1.35), Inches(4.0) if x < Inches(4.0) else Inches(4.2) if x < Inches(8.0) else Inches(4.0), Inches(0.32), fill=color)
    _txt(s, x+Inches(0.1), Inches(1.37), Inches(3.8) if x < Inches(4.0) else Inches(4.0) if x < Inches(8.0) else Inches(3.8), Inches(0.28), label, 10, WHITE, True)

# INPUTS column
inputs_a = [
    ('Threat Intelligence', 'CVE advisories, CTI feeds, IOCs, IOAs, PSIRTs'),
    ('Tech Controls Config', 'NGFW, IdP, IAM, IPS, EPP, SASE, WAF rules'),
    ('Grounded World Models', 'Vendor capabilities, MITRE ATT&CK / D3FEND'),
    ('Telemetry Schema', 'Asset inventory, topology, segmentation intent'),
]
y = Inches(1.75)
for name, desc in inputs_a:
    _box(s, Inches(0.5), y, Inches(4.0), Inches(0.9), fill=CARD, border=NAVY)
    _txt(s, Inches(0.7), y+Inches(0.05), Inches(3.6), Inches(0.28), name, 10, NAVY, True)
    _txt(s, Inches(0.7), y+Inches(0.38), Inches(3.6), Inches(0.45), desc, 8.5, BODY)
    y += Inches(0.97)

# ENGINE column
engine_steps = [
    ('1', 'Identify Risk-Prioritized Control Gaps',
     'Model attacks   |   Model defenses   |   Reason about gaps'),
    ('2', 'Generate Remediations',
     'Risk-prioritized fixes   |   Vendor-specific playbooks   |   Alerts to SOC'),
    ('3', 'Fix Control Gaps',
     'Apply remediations   |   Validate   |   Re-verify continuously'),
]
y = Inches(1.75)
for num, title, desc in engine_steps:
    _box(s, Inches(4.85), y, Inches(4.1), Inches(1.18), fill=RIGOR_BG, border=ORG)
    _rect(s, Inches(4.85), y, Inches(0.5), Inches(1.18), fill=ORG)
    _txt(s, Inches(4.87), y+Inches(0.38), Inches(0.46), Inches(0.4), num, 16, WHITE, True, PP_ALIGN.CENTER, 'Georgia')
    _txt(s, Inches(5.5), y+Inches(0.1), Inches(3.3), Inches(0.35), title, 10, NAVY, True)
    _txt(s, Inches(5.5), y+Inches(0.5), Inches(3.3), Inches(0.6), desc, 8.5, BODY)
    y += Inches(1.3)

# ASSURANCE column
assurance_sections = [
    ('Findings Review', [
        'Risk-prioritized gaps',
        'Vendor-specific remediation',
        'No false positives / No false negatives',
    ]),
    ('Continuous Remediation', [
        'SOC/control admin workflows',
        'ServiceNow / JIRA integration',
        'Auto re-verify on config change',
        'Auto re-verify on threat-feed update',
    ]),
]
y = Inches(1.75)
for title, items in assurance_sections:
    h = Inches(0.45 + len(items) * 0.35)
    _box(s, Inches(9.35), y, Inches(3.7), h, fill=ACE_BG, border=TEAL)
    _txt(s, Inches(9.5), y+Inches(0.05), Inches(3.4), Inches(0.3), title, 10, TEAL, True)
    yi = y + Inches(0.4)
    for item in items:
        _txt(s, Inches(9.6), yi, Inches(3.3), Inches(0.3), '- ' + item, 8.5, BODY)
        yi += Inches(0.3)
    y += h + Inches(0.2)

_bottom_bar(s, 'Currently supports leading N-S and E-W perimeter NGFWs (PAN, Fortinet, Cisco, Check Point). Identity, SASE, and additional vendors on roadmap.')


# SLIDE B: Volt Typhoon TTP-by-TTP
s = _content_slide('PART 1  --  22CT PREEMPTIVE  --  VOLT TYPHOON COVERAGE',
                   'How 22CT Preempts Volt Typhoon -- TTP by TTP',
                   'Each MITRE technique -> formal control gap query -> mathematically verified blocking configuration')

_rect(s, Inches(0.5), Inches(1.35), Inches(1.0), Inches(0.35), fill=NAVY)
_txt(s, Inches(0.55), Inches(1.37), Inches(0.9), Inches(0.3), 'MITRE ID', 9, WHITE, True)
_rect(s, Inches(1.6), Inches(1.35), Inches(2.5), Inches(0.35), fill=CRED)
_txt(s, Inches(1.65), Inches(1.37), Inches(2.4), Inches(0.3), 'TECHNIQUE', 9, WHITE, True)
_rect(s, Inches(4.2), Inches(1.35), Inches(8.6), Inches(0.35), fill=ORG)
_txt(s, Inches(4.25), Inches(1.37), Inches(8.5), Inches(0.3), '22CT PREEMPTIVE RESPONSE', 9, WHITE, True)

ttp_rows_v = [
    ('T1190', 'Exploit Public-Facing App', 'IPS signature + WAF rule verified for every edge CVE (Fortinet, Ivanti, Versa, PAN)'),
    ('T1133', 'External Remote Services', 'Zero-trust segmentation verified; MFA bypass paths surfaced and remediated'),
    ('T1059', 'Command & Scripting', 'EDR/script-control coverage proven on every endpoint in scope'),
    ('T1003', 'OS Credential Dumping', 'Egress filtering proven; ntdsutil blocked for non-admin via GPO verification'),
    ('T1021', 'Remote Services', 'SMB inspection verified; micro-segmentation proven; admin shares audited'),
    ('T1071', 'Application Layer Protocol', 'Threat-feed C2 IPs/domains denied across 50M+ rule combinations'),
    ('T1078', 'Valid Accounts', 'Credential-based lateral -> HANDED TO ACECARD behavioral layer'),
    ('T1562', 'Impair Defenses', 'Config Drift Mgmt catches logging/audit baseline deviation within 1 hour'),
    ('T1090', 'Proxy (KV Botnet)', 'SOHO router IPs from threat feeds denied all zones; botnet IOCs auto-updated'),
]
y = Inches(1.75)
for i, (mid, tech, response) in enumerate(ttp_rows_v):
    bg = CARD if i % 2 == 0 else WHITE
    _box(s, Inches(0.5), y, Inches(12.3), Inches(0.52), fill=bg)
    _txt(s, Inches(0.6), y+Pt(2), Inches(0.9), Inches(0.47), mid, 8.5, NAVY, True, fn='Cascadia Code')
    _txt(s, Inches(1.6), y+Pt(2), Inches(2.4), Inches(0.47), tech, 8.5, BODY)
    color = TEAL if mid == 'T1078' else RGBColor(0x1A, 0x4D, 0x8C)
    _txt(s, Inches(4.2), y+Pt(2), Inches(8.5), Inches(0.47), response, 8.5, color)
    y += Inches(0.54)

_txt(s, Inches(0.5), Inches(7.2), Inches(12.3), Inches(0.2),
     'Source: MITRE ATT&CK Group G1017; CISA AA24-038A', 9, SUB)


# SLIDE C: Salt Typhoon TTP-by-TTP
s = _content_slide('PART 1  --  22CT PREEMPTIVE  --  SALT TYPHOON COVERAGE',
                   'How 22CT Preempts Salt Typhoon -- TTP by TTP',
                   'Each MITRE technique -> formal control gap query -> mathematically verified blocking configuration')

_rect(s, Inches(0.5), Inches(1.35), Inches(1.0), Inches(0.35), fill=NAVY)
_txt(s, Inches(0.55), Inches(1.37), Inches(0.9), Inches(0.3), 'MITRE ID', 9, WHITE, True)
_rect(s, Inches(1.6), Inches(1.35), Inches(2.5), Inches(0.35), fill=CRED)
_txt(s, Inches(1.65), Inches(1.37), Inches(2.4), Inches(0.3), 'TECHNIQUE', 9, WHITE, True)
_rect(s, Inches(4.2), Inches(1.35), Inches(8.6), Inches(0.35), fill=TEAL)
_txt(s, Inches(4.25), Inches(1.37), Inches(8.5), Inches(0.3), '22CT PREEMPTIVE RESPONSE', 9, WHITE, True)

ttp_rows_s = [
    ('T1190', 'Exploit Public-Facing App', 'WAF/mgmt-plane ACL blocks /webui/ for CVE-2023-20198; compensating control verified'),
    ('T1136', 'Create Account', 'Config Drift detects unauthorized local account creation within 1 hour'),
    ('T1556', 'Modify Auth Process', 'AAA/TACACS+ config baseline comparison every hour; any delta = immediate gap finding'),
    ('T1040', 'Network Sniffing', 'On-box capture restricted to authorized roles; mgmt-plane access verified'),
    ('T1572', 'Protocol Tunneling', 'GRE peer authorization verified; unauthorized tunnel endpoints detected and flagged'),
    ('T1041', 'Exfil Over C2 Channel', 'Outbound to JumbledPath C2 IPs blocked across all zones; CALEA-zone segmentation proved'),
    ('T1562', 'Impair Defenses', 'Log clearing / ACL modification caught by drift management within detection window'),
    ('T1078', 'Valid Accounts', 'TACACS+ credential abuse -> HANDED TO ACECARD behavioral layer for drift detection'),
]
y = Inches(1.75)
for i, (mid, tech, response) in enumerate(ttp_rows_s):
    bg = CARD if i % 2 == 0 else WHITE
    _box(s, Inches(0.5), y, Inches(12.3), Inches(0.58), fill=bg)
    _txt(s, Inches(0.6), y+Pt(2), Inches(0.9), Inches(0.53), mid, 8.5, NAVY, True, fn='Cascadia Code')
    _txt(s, Inches(1.6), y+Pt(2), Inches(2.4), Inches(0.53), tech, 8.5, BODY)
    color = TEAL if mid == 'T1078' else RGBColor(0x1A, 0x4D, 0x8C)
    _txt(s, Inches(4.2), y+Pt(2), Inches(8.5), Inches(0.53), response, 8.5, color)
    y += Inches(0.6)

_txt(s, Inches(0.5), Inches(7.2), Inches(12.3), Inches(0.2),
     'Source: CISA AA25-239A; FBI/NSA/Five-Eyes Joint Advisory', 9, SUB)


# SLIDE D: Nine Use Cases
s = _content_slide('PART 1  --  22CT PREEMPTIVE  --  USE CASES',
                   'Nine Preemptive Use Cases -- Each Closes a Typhoon Attack Vector',
                   'Pre-built reasoning templates. Five directly preempt Volt + Salt Typhoon TTPs.')

uc_rows = [
    ('01', 'Preemptive Cyber Defense',
     'Continuous ML/NLP/LLM threat intel + config -> proof every TTP blocked',
     'STOPS: Volt T1190, T1133, T1071 / Salt T1190, T1556'),
    ('02', 'Configuration Drift Management',
     'Live vs approved baseline comparison every hour; any delta = gap finding',
     'STOPS: Salt ACL tampering, Guest Shell activation, local account creation'),
    ('03', 'Formal Policy Analysis',
     '2^8000 rule space: finds shadows, conflicts, weak inspection, intent errors',
     'STOPS: SMB lateral paths, shadowed inspection profiles, multi-rule bypasses'),
    ('04', 'Preemptive ASM (External)',
     'External attack surface mapped -> compensating controls verified -> proof',
     'STOPS: Edge CVE exploitation (Volt T1190, Salt T1190)'),
    ('05', 'Internal Attack Surface Mgmt',
     'Internal paths from compromised host to crown jewels identified and blocked',
     'STOPS: Cross-tenant pivots, lateral movement post-compromise'),
    ('06', 'NIST 800-53 / Zero-Trust',
     'SP800-207 zero-trust verification end-to-end across all control layers',
     'Compliance-as-code: automated evidence for FedRAMP, CMMC, DISA STIGs'),
    ('07', 'Vendor Migration Safety',
     'Before/after config analysis: proves no security regression during migration',
     'Use case: PAN->Fortinet, Cisco->Palo Alto -- zero blind spots'),
    ('08', 'M&A Due Diligence',
     'Combined posture analysis pre-merger identifies inherited misconfigurations',
     'Inherited misconfigs + shadow IT -> found before integration'),
    ('09', 'Third-Party / Supply Chain',
     'EASM across third-party ecosystem: supplier exposure mapped',
     'Supply-chain pivots identified before they become breaches'),
]

y = Inches(1.35)
for i, (num, name, desc, stops) in enumerate(uc_rows):
    bg = CARD if i % 2 == 0 else WHITE
    _box(s, Inches(0.5), y, Inches(12.3), Inches(0.56), fill=bg)
    _rect(s, Inches(0.5), y+Inches(0.05), Inches(0.55), Inches(0.46), fill=ORG)
    _txt(s, Inches(0.52), y+Inches(0.08), Inches(0.51), Inches(0.38), num, 11, WHITE, True, PP_ALIGN.CENTER, 'Arial Black')
    _txt(s, Inches(1.2), y+Pt(2), Inches(2.4), Inches(0.25), name, 9.5, NAVY, True)
    _txt(s, Inches(3.7), y+Pt(2), Inches(4.6), Inches(0.5), desc, 8.5, BODY)
    _txt(s, Inches(8.5), y+Pt(2), Inches(4.2), Inches(0.5), stops, 8, CRED if i < 5 else SUB)
    y += Inches(0.58)

_bottom_bar(s, 'Use cases 1-5 are direct Volt/Salt counters. Use cases 6-9 extend value across the broader cyber operating model.')


# SLIDE E: Proof Points
s = _content_slide('PART 1  --  22CT PREEMPTIVE  --  PROOF POINTS',
                   'In Production: Real Customer Findings',
                   'Selected findings from financial services, telecom, government, and critical infrastructure')

proof_cards = [
    ('FinTech | Asia', 'Closed SMB Lateral Movement Path',
     'SMB traffic profiles shadowed by 3 higher-priority rules -- inspection never applied. '
     'All lateral paths from user VLANs to DC management now formally verified closed.'),
    ('Communications | NA', 'Eliminated Known CVE Exposure',
     'IPS filters missing on rsync connections to backup servers. '
     'CVE-matched threat profiles enforced; compensating control deployed same day.'),
    ('Regional Bank | E. Asia', 'Closed RCE Attack Paths',
     'Both PAN and Fortinet allowed Windows Update traffic with no IPS inspection. '
     'Two-vendor blind spot closed; exploit path formally proved blocked.'),
    ('Critical Infra | SE Asia', 'Closed App Obfuscation Paths',
     'DMZ rules not blocking all targeted applications as intended -- policy intent error. '
     'Formal verification found 17 rule combinations allowing blocked app categories.'),
    ('Finance | E. Asia', 'NIST 800-53 SMTP Compliance',
     'SMTP from trusted to untrusted zones with no threat inspection profile. '
     'Control gap mapped to NIST 800-53 SC-7; remediation deployed within 4 hours.'),
    ('EU Sovereign Govt', 'Preemptive ASM for National Agency',
     'Proved compliance and coverage across 12 government agencies in a single unified analysis. '
     'First formal proof of cross-agency segmentation intent.'),
]

xs = [Inches(0.5), Inches(6.8)]
ys_p = [Inches(1.5), Inches(3.1), Inches(4.7)]
for i, (customer, title, desc) in enumerate(proof_cards):
    x = xs[i % 2]
    y = ys_p[i // 2]
    _box(s, x, y, Inches(5.8), Inches(1.4), fill=CARD, accent=ORG)
    _txt(s, x+Inches(0.3), y+Inches(0.05), Inches(5.3), Inches(0.25), customer, 9, ORG, True)
    _txt(s, x+Inches(0.3), y+Inches(0.3), Inches(5.3), Inches(0.35), title, 12, NAVY, True)
    _txt(s, x+Inches(0.3), y+Inches(0.7), Inches(5.3), Inches(0.6), desc, 9, BODY)

_bottom_bar(s, '20+ design partners across F500, federal, critical infrastructure, MSSPs, and nation states.')


# SLIDE 6: Formal Model - 2^8000
s = _content_slide('PART 1  --  22CT PREEMPTIVE  --  FORMAL MODEL',
                   'The Formal Model: Why 2^8000 Cannot Be Sampled',
                   'Traditional scanners sample configurations. Our platform reasons EXHAUSTIVELY over the entire state space.')

_box(s, Inches(0.5), Inches(1.5), Inches(5.8), Inches(3.0), fill=CARD, accent=CRED)
_txt(s, Inches(0.8), Inches(1.55), Inches(5.3), Inches(0.3), 'THE PROBLEM WITH SAMPLING', 12, CRED, True)
tf = _txt(s, Inches(0.8), Inches(1.9), Inches(5.3), Inches(2.4),
     'A typical enterprise NGFW has:', 11, BODY)
_p(tf, '  5,000 - 50,000 firewall rules', 11, BODY)
_p(tf, '  Each: src zone, dst zone, src IP, dst IP, port, app, user, action, profile', 10, SUB)
_p(tf, '  Interactions: 50,000^2 = 2.5 billion pairwise rule relationships', 11, BODY)
_p(tf, '  Multi-vendor: 3-10 firewalls in series/parallel paths', 11, BODY)
_p(tf, '', 8, SUB)
_p(tf, 'Total combinatorial state space: ~2^8000', 12, NAVY, True)
_p(tf, '', 8, SUB)
_p(tf, 'Sampling 1 billion paths/second for 1 year:', 10, SUB)
_p(tf, '  = 3.15 x 10^16 samples (vs 2^8000 = 10^2408 total)', 10, CRED, True)
_p(tf, '  Coverage: effectively 0.000...0%', 10, CRED, True)

_box(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(3.0), fill=CARD, accent=ORG)
_txt(s, Inches(7.1), Inches(1.55), Inches(5.3), Inches(0.3), '22CT: SYMBOLIC MODEL OF COMPUTATION', 12, ORG, True)
tf = _txt(s, Inches(7.1), Inches(1.9), Inches(5.3), Inches(2.4),
     'Instead of sampling, we build a mathematical MODEL:', 11, BODY)
_p(tf, '', 8, SUB)
_p(tf, '1. Parse all rules into symbolic representation', 11, BODY)
_p(tf, '2. Build constraint-satisfaction model', 11, BODY)
_p(tf, '3. For each threat (CVE, C2 dest, lateral path):', 11, BODY)
_p(tf, '   "Does ANY rule combination ALLOW this traffic?"', 11, NAVY, True)
_p(tf, '4. Solver PROVES answer is NO (safe) or YES (gap)', 11, BODY)
_p(tf, '', 8, SUB)
_p(tf, 'Result: COMPLETE coverage of 2^8000 in seconds', 12, ORG, True)
_p(tf, 'No sampling. No probabilistic reasoning.', 12, ORG, True)
_p(tf, 'Mathematical certainty.', 12, ORG, True)

_box(s, Inches(0.5), Inches(4.7), Inches(12.1), Inches(2.0), fill=CARD, accent=NAVY)
_txt(s, Inches(0.8), Inches(4.75), Inches(11.6), Inches(0.3), 'NOBODY ELSE HAS THIS', 14, NAVY, True)
approaches = [
    ('Traditional Pentest', '1-2 weeks, finds 5-15 paths', 'Point-in-time, incomplete, expensive', CRED),
    ('Vulnerability Scanner', 'Samples known CVEs against config', 'Misses rule interactions and shadows', CRED),
    ('Attack Surface Mgmt', 'Enumerates external-facing assets', 'Sees what exists, not what is REACHABLE', CRED),
    ('22CT FORMAL VERIFICATION', 'Models ALL paths through ALL rules', 'Complete, continuous, no gaps, vendor fixes', ORG),
]
y = Inches(5.1)
for tool, method, result, color in approaches:
    _rect(s, Inches(0.8), y, Inches(2.5), Inches(0.3), fill=color)
    _txt(s, Inches(0.85), y+Pt(1), Inches(2.4), Inches(0.28), tool, 10, WHITE, True, PP_ALIGN.CENTER)
    _txt(s, Inches(3.5), y+Pt(1), Inches(4.0), Inches(0.28), method, 10, BODY)
    _txt(s, Inches(7.6), y+Pt(1), Inches(4.8), Inches(0.28), result, 10, color, True)
    y += Inches(0.35)

_bottom_bar(s, 'Gartner 2024: "99% of firewall breaches are caused by misconfiguration, not firewall flaws." 22CT Preemptive eliminates that 99%.')


# SLIDE 7: Volt Typhoon phase-by-phase
s = _content_slide('PART 1  --  22CT PREEMPTIVE  --  STOPPING VOLT TYPHOON',
                   'How 22CT Stops Volt Typhoon -- Phase by Phase',
                   'Volt Typhoon needs configuration weakness to enter and pivot. 22CT proves those weaknesses do not exist.')

# Table header
_rect(s, Inches(0.5), Inches(1.4), Inches(1.5), Inches(0.4), fill=NAVY)
_txt(s, Inches(0.6), Inches(1.42), Inches(1.4), Inches(0.35), 'PHASE', 11, WHITE, True)
_rect(s, Inches(2.0), Inches(1.4), Inches(4.7), Inches(0.4), fill=CRED)
_txt(s, Inches(2.1), Inches(1.42), Inches(4.6), Inches(0.35), 'VOLT TYPHOON BEHAVIOR', 11, WHITE, True)
_rect(s, Inches(6.7), Inches(1.4), Inches(6.1), Inches(0.4), fill=ORG)
_txt(s, Inches(6.8), Inches(1.42), Inches(6.0), Inches(0.35), 'HOW 22CT STOPS IT', 11, WHITE, True)

rows = [
    ('PHASE 1\nInitial Access', 'Scans for unpatched Fortinet, Ivanti, Versa, PAN edge devices. Exploits CVE-2024-21887 etc.',
     'Ingests CISA/vendor PSIRTs. Proves every NGFW/VPN/SASE has IPS signature enforced. No patch = compensating control.'),
    ('PHASE 2\nPersistence', 'Web shells (.aspx/.jspx). C2 through KV botnet of compromised SOHO routers.',
     'Models 50M+ rule combinations. Proves all threat-feed IPs/domains blocked. WAF verified for upload inspection.'),
    ('PHASE 3\nCredential Access', 'ntdsutil dump NTDS.dit. LSASS extraction. Exfiltrates creds to C2.',
     'Verifies EDR/EPP coverage on DCs. Script-control blocks ntdsutil for non-admin. Egress filtering verified.'),
    ('PHASE 4\nLateral Movement', 'SMB/RDP/WinRM to domain controllers. Admin share access.',
     'Real finding: 3 rules shadowed SMB IPS. Micro-segmentation: user VLAN proven unable to reach DC mgmt.'),
    ('PHASE 5\nPre-positioning', 'Network mapping + credential harvest for future sabotage. 5+ year dwell.',
     'Continuous re-verification hourly + on every config change. Zero-trust segmentation re-proven.'),
]
y = Inches(1.9)
for i, (phase, attack, rigor) in enumerate(rows):
    bg = CARD if i % 2 == 0 else WHITE
    _box(s, Inches(0.5), y, Inches(12.3), Inches(1.0), fill=bg)
    _txt(s, Inches(0.6), y+Pt(2), Inches(1.4), Inches(0.95), phase, 9.5, NAVY, True)
    _txt(s, Inches(2.1), y+Pt(2), Inches(4.5), Inches(0.95), attack, 9, BODY)
    _txt(s, Inches(6.8), y+Pt(2), Inches(5.9), Inches(0.95), rigor, 9, RGBColor(0x1A, 0x4D, 0x8C))
    y += Inches(1.03)

_txt(s, Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.3),
     'Bottom line: 22CT Preemptive doesn\'t detect Volt Typhoon -- it makes the infrastructure impervious to their entry technique.',
     10, SUB)


# SLIDE 8: Salt Typhoon phase-by-phase
s = _content_slide('PART 1  --  22CT PREEMPTIVE  --  STOPPING SALT TYPHOON',
                   'How 22CT Stops Salt Typhoon -- Carrier-Scale Edge Defense',
                   'Salt Typhoon lives on misconfigured Cisco/Ivanti/PAN edge gear. 22CT verifies every config every hour.')

_rect(s, Inches(0.5), Inches(1.4), Inches(1.5), Inches(0.4), fill=NAVY)
_txt(s, Inches(0.6), Inches(1.42), Inches(1.4), Inches(0.35), 'PHASE', 11, WHITE, True)
_rect(s, Inches(2.0), Inches(1.4), Inches(4.7), Inches(0.4), fill=CRED)
_txt(s, Inches(2.1), Inches(1.42), Inches(4.6), Inches(0.35), 'SALT TYPHOON BEHAVIOR', 11, WHITE, True)
_rect(s, Inches(6.7), Inches(1.4), Inches(6.1), Inches(0.4), fill=ORG)
_txt(s, Inches(6.8), Inches(1.42), Inches(6.0), Inches(0.35), 'HOW 22CT STOPS IT', 11, WHITE, True)

rows_st = [
    ('PHASE 1\nEdge CVE Exploit', 'Cisco IOS XE web UI (CVE-2023-20198) for level-15 admin. 1000+ devices targeted.',
     'CVE-to-control mapping automatic. Proves WAF/mgmt-plane ACL blocks /webui/. Compensating control for unpatched.'),
    ('PHASE 2\nPersistence', 'Local admin accounts. Router ACL modifications. Guest Shell. 3+ years access.',
     'Config Drift: approved baseline compared hourly. Any unauthorized ACL/AAA change = immediate gap finding.'),
    ('PHASE 3\nTACACS+ Capture', 'Captures TACACS+ on TCP/49 via on-box pcap. Harvests creds. JumbledPath C2.',
     'Mgmt-plane segmentation verified. On-box capture restricted. JumbledPath C2 IPs blocked all zones.'),
    ('PHASE 4\nGRE Tunnel + C2', 'GRE tunnel between compromised routers. Encrypted C2. Multi-hop chains.',
     'GRE peer authorization verified. Unauthorized tunnel endpoints detected. Egress: only approved protocols.'),
    ('PHASE 5\nCALEA Exfil', 'CALEA lawful intercept access. CDR/metadata exfiltration over months.',
     'CALEA-zone segmentation proven end-to-end. DLP on CDR stores. Outbound restrictions verified.'),
]
y = Inches(1.9)
for i, (phase, attack, rigor) in enumerate(rows_st):
    bg = CARD if i % 2 == 0 else WHITE
    _box(s, Inches(0.5), y, Inches(12.3), Inches(1.0), fill=bg)
    _txt(s, Inches(0.6), y+Pt(2), Inches(1.4), Inches(0.95), phase, 9.5, NAVY, True)
    _txt(s, Inches(2.1), y+Pt(2), Inches(4.5), Inches(0.95), attack, 9, BODY)
    _txt(s, Inches(6.8), y+Pt(2), Inches(5.9), Inches(0.95), rigor, 9, RGBColor(0x1A, 0x4D, 0x8C))
    y += Inches(1.03)

_txt(s, Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.3),
     'Salt Typhoon: 200+ telcos, 80+ countries, 4+ years. Every phase requires a misconfigured edge device. 22CT proves correct.',
     10, SUB)


# SLIDE 9: Where Preemptive needs help -> ACECARD
s = _content_slide('PART 1  --  22CT PREEMPTIVE  --  HONEST ASSESSMENT',
                   'Where Preemptive Needs Help -- Why ACECARD Belongs Next to It',
                   'The preemptive layer closes the configuration attack surface. Three classes of attack remain.')

gaps_data = [
    ('1. Zero-Day / Unknown TTPs', CRED,
     'No signature, no IOC, no MITRE mapping yet = nothing for preemptive layer to prove against.',
     'ACECARD: Detects behavioral drift even with no label. Post-exploitation ALWAYS creates drift.'),
    ('2. Compromised Valid Accounts', ORG,
     'Both Typhoons abuse stolen valid credentials. Correctly-configured controls ALLOW this traffic by design.',
     'ACECARD: Cohort z-score + cosine drift catches abuse -- same credentials, different behavioral structure.'),
    ('3. Living-off-the-Land', RGBColor(0x6C, 0x2E, 0xD9),
     'wmic, ntdsutil, netsh, rundll32 -- legitimate admin tools. Cannot be blanket-blocked.',
     'ACECARD: 5-signal embedding captures the COMBINATION of LoTL tools + timing + targets. Tool-agnostic detection.'),
]

y = Inches(1.5)
for title, color, problem, acecard in gaps_data:
    _box(s, Inches(0.5), y, Inches(12.1), Inches(1.6), fill=CARD, accent=color)
    _txt(s, Inches(0.8), y+Inches(0.05), Inches(11.6), Inches(0.3), title, 14, color, True)
    _txt(s, Inches(0.8), y+Inches(0.4), Inches(11.6), Inches(0.5), problem, 11, BODY)
    _rect(s, Inches(0.8), y+Inches(0.95), Inches(0.08), Inches(0.5), fill=TEAL)
    _txt(s, Inches(1.0), y+Inches(0.95), Inches(11.4), Inches(0.55), acecard, 11, TEAL, True)
    y += Inches(1.7)

_bottom_bar(s, 'Configuration preemption is necessary but not sufficient. Behavioral detection closes the remaining surface.')


# ================================================================
# PART 2: ACECARD UEBA SLIDES
# ================================================================

# SLIDE 10: ACECARD section opener
s = _content_slide('PART 2  --  BEHAVIORAL LAYER  --  ACECARD UEBA',
                   'Why Current UEBA Misses Volt Typhoon',
                   'Both campaigns relied on Living-off-the-Land tradecraft. Both walked past every commercial UEBA.', TEAL)

_rect(s, Inches(0.5), Inches(1.6), Inches(3.5), Inches(1.2), fill=NAVY)
_txt(s, Inches(0.7), Inches(1.7), Inches(3.1), Inches(0.7), '>= 5 yrs', 36, ORG, True, fn='Georgia')
_txt(s, Inches(0.7), Inches(2.3), Inches(3.1), Inches(0.4),
     'Volt Typhoon undetected dwell\nin U.S. critical infrastructure', 10, WHITE)

_txt(s, Inches(4.2), Inches(1.6), Inches(8.3), Inches(0.3),
     'The Four Structural Gaps in Commercial UEBA', 14, NAVY, True)
_txt(s, Inches(4.2), Inches(1.9), Inches(8.3), Inches(0.5),
     'Each one is exactly the gap Volt Typhoon and Salt Typhoon exploit.', 11, BODY)

ueba_gaps = [
    ('01', 'Threshold-based scoring', 'Alerts fire only when a metric crosses a fixed line. LoTL never crosses any single threshold.',
     'Volt Typhoon process tree was within "normal" at every single data point.'),
    ('02', 'Per-source siloed analytics', 'Each tool sees one identity. AD, AAD, Okta, AWS, K8s -- separate systems, no correlation.',
     'Cross-domain attack chains never escalate above per-source noise floor.'),
    ('03', 'Static rules, decaying coverage', 'Sigma/CAR rules require maintenance. Adversary retool cycle: 2-4 weeks. Defender: 1-3 months.',
     'By the time a Sigma rule is published, the tradecraft has already evolved.'),
    ('04', 'Anomaly score without direction', '"User is 87% anomalous" -- anomalous HOW? No MITRE mapping. No kill-chain position.',
     'Mean time to investigate exceeds dwell window -- even when detection fires.'),
]
y = Inches(2.8)
for num, title, desc, consequence in ueba_gaps:
    _box(s, Inches(0.5) if num in ('01','03') else Inches(6.8),
         y if num in ('01','02') else y+Inches(2.0),
         Inches(5.8), Inches(1.8), fill=CARD, accent=CRED)
    xo = Inches(0.5) if num in ('01','03') else Inches(6.8)
    _rect(s, xo, y if num in ('01','02') else y+Inches(2.0), Inches(0.7), Inches(1.8), fill=CRED)
    _txt(s, xo+Inches(0.05), (y if num in ('01','02') else y+Inches(2.0))+Inches(0.2),
         Inches(0.6), Inches(0.4), num, 18, WHITE, True, PP_ALIGN.CENTER, 'Arial Black')
    _txt(s, xo+Inches(0.9), (y if num in ('01','02') else y+Inches(2.0))+Inches(0.1),
         Inches(4.7), Inches(0.3), title, 13, NAVY, True)
    _txt(s, xo+Inches(0.9), (y if num in ('01','02') else y+Inches(2.0))+Inches(0.4),
         Inches(4.7), Inches(0.8), desc, 10, BODY)
    _txt(s, xo+Inches(0.9), (y if num in ('01','02') else y+Inches(2.0))+Inches(1.25),
         Inches(4.7), Inches(0.45), 'CONSEQUENCE  ' + consequence, 9, CRED)


# SLIDE 11: Structural Innovations
s = _content_slide('PART 2  --  ACECARD UEBA  --  STRUCTURAL INNOVATIONS',
                   'Four Innovations That Close the Four Gaps',
                   'Each innovation directly addresses one of the four structural gaps that Volt/Salt Typhoon exploit.', TEAL)

innovations = [
    ('01', 'Behavioral TRAJECTORY, not threshold', TEAL,
     'Track HOW entities EVOLVE over time. CUSUM catches drift that never crosses a fixed threshold. 5+ year APT = detected in days.',
     'CLOSES GAP 1: Threshold-based scoring'),
    ('02', 'Entity FUSION across 10+ identity systems', TEAL,
     'One entity_uuid resolves AD, AAD, Okta, AWS, K8s, CrowdStrike, VPN, Splunk, PKI, TACACS+ into single behavioral entity.',
     'CLOSES GAP 2: Per-source silos'),
    ('03', 'Behavioral OBJECTIVES, not tool signatures', TEAL,
     'Detect the EFFECT of the attack (credential access, lateral movement) not the TOOL (ntdsutil vs secretsdump vs mimikatz).',
     'CLOSES GAP 3: Static rules / rule decay'),
    ('04', 'Drift DIRECTION with MITRE technique mapping', TEAL,
     'Not "87% anomalous" but "drifting toward lateral_movement (T1021) at 0.82 cosine + credential_dumping (T1003) at 0.54."',
     'CLOSES GAP 4: Anomaly without direction'),
]
y = Inches(1.5)
for num, title, color, desc, closes in innovations:
    _box(s, Inches(0.5), y, Inches(12.1), Inches(1.2), fill=CARD, accent=color)
    _rect(s, Inches(0.5), y, Inches(0.7), Inches(1.2), fill=TEAL)
    _txt(s, Inches(0.55), y+Inches(0.1), Inches(0.6), Inches(0.4), num, 18, WHITE, True, PP_ALIGN.CENTER, 'Arial Black')
    _txt(s, Inches(1.4), y+Inches(0.05), Inches(10.9), Inches(0.3), title, 13, NAVY, True)
    _txt(s, Inches(1.4), y+Inches(0.35), Inches(10.9), Inches(0.5), desc, 10.5, BODY)
    _txt(s, Inches(1.4), y+Inches(0.85), Inches(10.9), Inches(0.3), closes, 9, TEAL, True)
    y += Inches(1.3)

_bottom_bar(s, 'No commercial UEBA has ALL FOUR. Most have zero. This is the innovation gap we close.')


# SLIDE 12: Dual Embedding Architecture
s = _content_slide('PART 2  --  ACECARD UEBA  --  THE MATH',
                   'Dual Embedding Architecture: 6 Vectors = 9,216 Dimensions',
                   'Each signal individually embedded (1536-d) + one fused composite. Per-signal drift tracking.', TEAL)

_box(s, Inches(0.5), Inches(1.5), Inches(12.1), Inches(0.8), fill=ACE_BG)
_txt(s, Inches(0.7), Inches(1.55), Inches(11.7), Inches(0.25), 'THE KEY INSIGHT', 12, TEAL, True)
_txt(s, Inches(0.7), Inches(1.8), Inches(11.7), Inches(0.4),
     'Each of the 5 signals is INDIVIDUALLY embedded at 1536-d (per-signal drift tracking). '
     'All 5 are also concatenated into one COMPOSITE 1536-d vector (holistic entity state). '
     'Total: 6 vectors = 9,216 dimensions per entity per window.', 11, BODY)

sig_colors = [RGBColor(0x1E,0x4D,0x8C), RGBColor(0x6C,0x2E,0xD9), TEAL, ORG, CRED]
sig_names = ['AUTH', 'PROCESS', 'NETWORK', 'FILE', 'IDENTITY']
xp = [Inches(0.5), Inches(3.1), Inches(5.7), Inches(8.3), Inches(10.9)]
for i in range(5):
    _rect(s, xp[i], Inches(2.5), Inches(2.4), Inches(0.35), fill=sig_colors[i])
    _txt(s, xp[i], Inches(2.52), Inches(2.4), Inches(0.3), f'{sig_names[i]}', 11, WHITE, True, PP_ALIGN.CENTER)
    _box(s, xp[i], Inches(2.9), Inches(2.4), Inches(0.45), fill=CARD, border=sig_colors[i])
    _txt(s, xp[i]+Inches(0.1), Inches(2.95), Inches(2.2), Inches(0.35),
         'Individual 1536-d\nPer-signal drift', 9, BODY, al=PP_ALIGN.CENTER)

_rect(s, Inches(0.5), Inches(3.5), Inches(12.1), Inches(0.35), fill=NAVY)
_txt(s, Inches(0.7), Inches(3.52), Inches(11.7), Inches(0.3),
     'All 5 concatenated -> 1 COMPOSITE 1536-d vector (holistic entity behavioral state)', 11, GOLD, True, PP_ALIGN.CENTER)

# Three info boxes
_box(s, Inches(0.5), Inches(4.1), Inches(3.8), Inches(1.5), fill=CARD, accent=TEAL)
_txt(s, Inches(0.8), Inches(4.15), Inches(3.3), Inches(0.25), 'DIMENSIONS', 11, TEAL, True)
tf = _txt(s, Inches(0.8), Inches(4.4), Inches(3.3), Inches(1.1),
     '5 individual: 5 x 1,536 = 7,680', 10, BODY)
_p(tf, '1 composite: 1 x 1,536 = 1,536', 10, BODY)
_p(tf, 'TOTAL: 6 x 1,536 = 9,216 dims', 11, NAVY, True)

_box(s, Inches(4.6), Inches(4.1), Inches(4.0), Inches(1.5), fill=CARD, accent=ORG)
_txt(s, Inches(4.9), Inches(4.15), Inches(3.5), Inches(0.25), 'WHY INDIVIDUAL + COMPOSITE?', 11, ORG, True)
tf = _txt(s, Inches(4.9), Inches(4.4), Inches(3.5), Inches(1.1),
     'Composite: "entity behavior changed"', 10, BODY)
_p(tf, 'Individual: "AUTH drove 72% of change"', 10, BODY)
_p(tf, 'Direction: "PROCESS signal drifting', 10, BODY)
_p(tf, '  toward lotl_execution (T1059)"', 10, NAVY, True)

_box(s, Inches(8.8), Inches(4.1), Inches(3.8), Inches(1.5), fill=CARD, accent=NAVY)
_txt(s, Inches(9.1), Inches(4.15), Inches(3.3), Inches(0.25), 'DEPLOYMENT OPTIONS', 11, NAVY, True)
tf = _txt(s, Inches(9.1), Inches(4.4), Inches(3.3), Inches(1.1),
     'IL5: OpenAI text-embedding-3-small', 10, BODY)
_p(tf, 'IL6: Local Phi-4 (14B) or Mistral 7B', 10, BODY)
_p(tf, 'JWICS: Same SLM, air-gapped', 10, BODY)
_p(tf, 'Switch: env var only, no code change', 10, NAVY, True)

_box(s, Inches(0.5), Inches(5.8), Inches(12.1), Inches(0.5), fill=CARD, accent=GOLD)
_txt(s, Inches(0.8), Inches(5.85), Inches(11.6), Inches(0.4),
     'PRIVACY: No raw event content sent to embedding provider. Only summarized metrics. No usernames, IPs, paths. JWICS-safe.',
     10, GOLD, True)

_bottom_bar(s, 'NOBODY IN THE INDUSTRY embeds behavioral signals into a unified vector space with per-signal drift decomposition.')


# SLIDE 13: The Five Signals with metrics
s = _content_slide('PART 2  --  ACECARD UEBA  --  THE FIVE SIGNALS',
                   'Five Behavioral Signals: Sources, Metrics, Thresholds',
                   'Each signal summarizes one dimension of behavior. Together they capture the complete entity footprint.', TEAL)

signals = [
    ('AUTH', sig_colors[0], 'Win 4624/4625/4768, Okta, AAD, AWS CloudTrail',
     'logon_count, failure_rate, unique_hosts, off_hours_ratio, impossible_travel, mfa_skip_ratio',
     'Brute force, Pass-the-Hash, lateral movement, impossible travel, MFA bypass'),
    ('PROCESS', sig_colors[1], 'Sysmon EID 1/3/7/8, CrowdStrike, Defender, auditd',
     'unique_processes, cmdline_entropy, lolbin_count, parent_child_depth, unsigned_ratio, encoded_cmd',
     'LoTL chains (Volt Typhoon), code injection, encoded PowerShell, deep process trees'),
    ('NETWORK', sig_colors[2], 'NetFlow/IPFIX, Zeek, PAN-OS, FortiGate, VPC Flow',
     'unique_dest_ips, bytes_out_ratio, beacon_score, dns_query_rate, geo_anomaly, admin_share_access',
     'C2 beaconing, data exfiltration, DNS tunneling, scanning, GRE tunneling (Salt)'),
    ('FILE', sig_colors[3], 'Sysmon EID 11, Win 4663, DLP, cloud storage audit',
     'files_created, files_deleted, sensitive_access, archive_creates, extension_changes, shadow_copy_ops',
     'Ransomware, exfil staging, insider theft, credential file access, pre-encryption shadow delete'),
    ('IDENTITY', sig_colors[4], 'Win 4672/4728, CloudTrail IAM, K8s audit, AD change',
     'priv_escalations, group_adds, mfa_bypass, role_changes, service_acct_use, admin_actions',
     'Privilege escalation, Golden ticket, MFA bypass, service account abuse, insider admin actions'),
]

y = Inches(1.5)
for name, color, sources, metrics, attacks in signals:
    _rect(s, Inches(0.5), y, Inches(1.2), Inches(1.0), fill=color)
    _txt(s, Inches(0.55), y+Inches(0.2), Inches(1.1), Inches(0.5), name, 12, WHITE, True, PP_ALIGN.CENTER)
    _txt(s, Inches(1.8), y+Inches(0.02), Inches(4.5), Inches(0.3), sources, 9, SUB)
    _txt(s, Inches(1.8), y+Inches(0.3), Inches(4.5), Inches(0.35), metrics, 8.5, NAVY, True, fn='Cascadia Code')
    _txt(s, Inches(6.7), y+Inches(0.02), Inches(5.9), Inches(0.9), attacks, 9, BODY)
    y += Inches(1.05)

_bottom_bar(s, '30 metrics across 5 signals. ALL embedded into unified 1536-d space. Drift = attack. Direction = MITRE technique.')


# SLIDE F: The Math of Behavioral Drift
s = _content_slide('PART 2  --  ACECARD UEBA  --  THE MATH',
                   'The Math of Behavioral Drift',
                   'Each piece is well-known math. The innovation is the integration.', TEAL)

math_boxes = [
    ('A', 'Behavioral Embedding',
     'V_t = embed("auth: logon_count=47, failure_rate=0.032 | process: lolbin_count=0, cmdline_entropy=1.4 | '
     'network: unique_dest=3, beacon_score=0.05 | file: sensitive_access=0 | identity: priv_esc=0")',
     sig_colors[0]),
    ('B', 'Cosine Drift',
     'drift(t) = 1 - (V_t . V_baseline) / (||V_t|| * ||V_baseline||)\n'
     'Range: [0, 2]    Typical healthy: [0, 0.05]    Compromised: [0.06, 0.15+]',
     sig_colors[1]),
    ('C', 'CUSUM (Page 1954)',
     'S_t = max(0, S_{t-1} + (drift_t - mu - k))\n'
     'k = 0.5 * sigma  (slack)    h = 4 * sigma  (alarm threshold)\n'
     'ALARM when S_t > h. Catches sustained sub-threshold drift accumulated over windows.',
     sig_colors[2]),
    ('D', 'Drift Direction',
     'drift_vec = V_current - V_baseline\n'
     'proj(c) = cosine(drift_vec, embed(concept_c))    for c in {credential_dumping, lateral_movement, ...}\n'
     'Returns: argmax top-3 concepts with MITRE technique IDs',
     sig_colors[3]),
]

y = Inches(1.5)
for label, title, formula, color in math_boxes:
    _box(s, Inches(0.5), y, Inches(12.1), Inches(1.2), fill=CARD, accent=color)
    _rect(s, Inches(0.5), y, Inches(0.5), Inches(1.2), fill=color)
    _txt(s, Inches(0.52), y+Inches(0.35), Inches(0.46), Inches(0.5), label, 18, WHITE, True, PP_ALIGN.CENTER, 'Georgia')
    _txt(s, Inches(1.15), y+Inches(0.05), Inches(10.9), Inches(0.25), title, 12, NAVY, True)
    _txt(s, Inches(1.15), y+Inches(0.35), Inches(10.9), Inches(0.75), formula, 9, BODY, fn='Cascadia Code')
    y += Inches(1.3)

_bottom_bar(s, 'Cosine drift + CUSUM + drift direction: three formulas that together catch Volt + Salt Typhoon where thresholds fail.')


# SLIDE G: Pipeline & Trajectory
s = _content_slide('PART 2  --  ACECARD UEBA  --  PIPELINE & TRAJECTORY',
                   'Embedding Pipeline -> Trajectory Analysis Engine',
                   'Every entity, every window: ingest -> summarize -> embed -> analyze -> alert.', TEAL)

# Left: Pipeline
_rect(s, Inches(0.5), Inches(1.4), Inches(5.8), Inches(0.32), fill=TEAL)
_txt(s, Inches(0.6), Inches(1.42), Inches(5.6), Inches(0.28), 'EMBEDDING PIPELINE', 11, WHITE, True)

pipeline_steps = [
    ('1', 'INGEST', 'Raw events from SIEM, EDR, netflow, AD, IAM -- any source with a timestamp and entity ID'),
    ('2', 'WINDOW', 'Group by entity_uuid over configurable window (default 1 hour). 5 signal buckets.'),
    ('3', 'SUMMARIZE', 'Compute 30 metrics per window: logon_count, lolbin_count, beacon_score, etc.'),
    ('4', 'HASH', 'SHA-256 of metric snapshot. Tamper-detection baseline.'),
    ('5', 'EMBED', 'text-embedding-3-small (IL5) or local SLM (IL6/JWICS). 1536-d per signal + composite.'),
    ('6', 'STORE', 'entity_snapshots table: uuid, window, embedding, features, hmac_chain.'),
]
y = Inches(1.8)
for num, step, desc in pipeline_steps:
    _rect(s, Inches(0.5), y, Inches(0.5), Inches(0.55), fill=TEAL)
    _txt(s, Inches(0.52), y+Inches(0.08), Inches(0.46), Inches(0.35), num, 14, WHITE, True, PP_ALIGN.CENTER, 'Georgia')
    _box(s, Inches(1.1), y, Inches(5.1), Inches(0.55), fill=CARD, border=TEAL)
    _txt(s, Inches(1.25), y+Inches(0.02), Inches(1.1), Inches(0.25), step, 9, TEAL, True)
    _txt(s, Inches(2.45), y+Inches(0.02), Inches(3.65), Inches(0.5), desc, 8.5, BODY)
    y += Inches(0.62)

# Right: Trajectory metrics
_rect(s, Inches(6.9), Inches(1.4), Inches(5.9), Inches(0.32), fill=sig_colors[1])
_txt(s, Inches(7.0), Inches(1.42), Inches(5.7), Inches(0.28), 'TRAJECTORY ANALYSIS ENGINE', 11, WHITE, True)

metrics_t = [
    ('Cosine Drift', 'drift(t) = 1 - cos(V_t, V_baseline)  |  range [0,2]', sig_colors[0]),
    ('Velocity', 'v(t) = drift(t) - drift(t-1)  |  rate of change', sig_colors[1]),
    ('Acceleration', 'a(t) = v(t) - v(t-1)  |  second derivative', TEAL),
    ('CUSUM', 'S_t = max(0, S_{t-1} + drift_t - mu - k)  |  alarm: S > 4*sigma', ORG),
    ('Health Score', 'H = 100 * (1 - clamp(drift / h, 0, 1))  |  range [0, 100]', CRED),
]
y = Inches(1.8)
for name, formula, color in metrics_t:
    _box(s, Inches(6.9), y, Inches(5.9), Inches(0.55), fill=CARD, border=color)
    _rect(s, Inches(6.9), y, Inches(0.08), Inches(0.55), fill=color)
    _txt(s, Inches(7.1), y+Inches(0.04), Inches(2.2), Inches(0.25), name, 10, NAVY, True)
    _txt(s, Inches(7.1), y+Inches(0.28), Inches(5.6), Inches(0.25), formula, 8.5, BODY, fn='Cascadia Code')
    y += Inches(0.62)

_bottom_bar(s, 'Pipeline latency: event to analysis < 3 seconds. Throughput: 10K entities/hour per 4-vCPU worker.')


# SLIDE H: Detecting Volt Typhoon signal-by-signal
s = _content_slide('PART 2  --  ACECARD UEBA  --  DETECTING VOLT TYPHOON',
                   'How ACECARD Detects Volt Typhoon -- Signal by Signal',
                   '5 signals -> 1 fused trajectory. CUSUM alarms T+47min. Drift direction: lateral_movement (0.78) + credential_dumping (0.61). HIGH alert.',
                   TEAL)

_rect(s, Inches(0.5), Inches(1.35), Inches(1.1), Inches(0.38), fill=NAVY)
_txt(s, Inches(0.55), Inches(1.37), Inches(1.0), Inches(0.33), 'SIGNAL', 9, WHITE, True)
_rect(s, Inches(1.7), Inches(1.35), Inches(3.4), Inches(0.38), fill=CRED)
_txt(s, Inches(1.75), Inches(1.37), Inches(3.3), Inches(0.33), 'VOLT TYPHOON BEHAVIOR', 9, WHITE, True)
_rect(s, Inches(5.2), Inches(1.35), Inches(3.8), Inches(0.38), fill=TEAL)
_txt(s, Inches(5.25), Inches(1.37), Inches(3.7), Inches(0.33), 'ACECARD DRIFT SIGNAL', 9, WHITE, True)
_rect(s, Inches(9.1), Inches(1.35), Inches(3.7), Inches(0.38), fill=sig_colors[1])
_txt(s, Inches(9.15), Inches(1.37), Inches(3.6), Inches(0.33), 'KEY METRIC', 9, WHITE, True)

rows_vh = [
    ('1\nAUTH', sig_colors[0],
     'Stolen AD creds, 12 new hosts in 2h, off-hours logins, no MFA',
     'unique_hosts 0.3->12, off_hours_ratio 0->0.7',
     'drift=0.142, cohort z=6.17'),
    ('2\nPROCESS', sig_colors[1],
     'wmic, ntdsutil, netsh, rundll32, encoded PowerShell chains',
     'lolbin_count 0->7, cmdline_entropy 1.4->3.2',
     'CUSUM alarm @ T+47min'),
    ('3\nNETWORK', sig_colors[2],
     'C2 to KV botnet IPs, new destinations, regular beaconing',
     'unique_dest 3->18, beacon_score 0.05->0.42',
     'c2_beaconing (T1071) proj=0.71'),
    ('4\nFILE', sig_colors[3],
     'Web shells, NTDS.dit access, 7-Zip staging',
     'sensitive_access 0->24, archive_creates 0->3',
     'insider_hoarding (T1074) proj=0.58'),
    ('5\nIDENTITY', sig_colors[4],
     'Priv escalation, Domain Admins group adds, svc_acct spike',
     'priv_esc 0->3, svc_acct_use 5->47',
     'privilege_escalation (T1078) proj=0.66'),
]
y = Inches(1.8)
for sig, color, behavior, drift_sig, metric in rows_vh:
    bg = CARD if rows_vh.index((sig, color, behavior, drift_sig, metric)) % 2 == 0 else WHITE
    _box(s, Inches(0.5), y, Inches(12.3), Inches(0.9), fill=bg)
    _rect(s, Inches(0.5), y, Inches(1.1), Inches(0.9), fill=color)
    _txt(s, Inches(0.55), y+Inches(0.2), Inches(1.0), Inches(0.5), sig, 8.5, WHITE, True, PP_ALIGN.CENTER)
    _txt(s, Inches(1.7), y+Pt(2), Inches(3.35), Inches(0.85), behavior, 8.5, BODY)
    _txt(s, Inches(5.2), y+Pt(2), Inches(3.75), Inches(0.85), drift_sig, 8.5, TEAL, fn='Cascadia Code')
    _txt(s, Inches(9.1), y+Pt(2), Inches(3.65), Inches(0.85), metric, 8.5, NAVY, True)
    y += Inches(0.93)

_bottom_bar(s, '5 signals -> 1 fused trajectory. CUSUM alarms T+47min. Drift: lateral_movement (0.78) + credential_dumping (0.61). HIGH alert.')


# SLIDE I: Detecting Salt Typhoon signal-by-signal
s = _content_slide('PART 2  --  ACECARD UEBA  --  DETECTING SALT TYPHOON',
                   'How ACECARD Detects Salt Typhoon -- Signal by Signal',
                   'Salt Typhoon wins by going slow. CUSUM accumulates the small drifts. 30-day rolling baseline adapts.',
                   TEAL)

_rect(s, Inches(0.5), Inches(1.35), Inches(1.1), Inches(0.38), fill=NAVY)
_txt(s, Inches(0.55), Inches(1.37), Inches(1.0), Inches(0.33), 'SIGNAL', 9, WHITE, True)
_rect(s, Inches(1.7), Inches(1.35), Inches(3.4), Inches(0.38), fill=CRED)
_txt(s, Inches(1.75), Inches(1.37), Inches(3.3), Inches(0.33), 'SALT TYPHOON BEHAVIOR', 9, WHITE, True)
_rect(s, Inches(5.2), Inches(1.35), Inches(3.8), Inches(0.38), fill=TEAL)
_txt(s, Inches(5.25), Inches(1.37), Inches(3.7), Inches(0.33), 'ACECARD DRIFT SIGNAL', 9, WHITE, True)
_rect(s, Inches(9.1), Inches(1.35), Inches(3.7), Inches(0.38), fill=sig_colors[1])
_txt(s, Inches(9.15), Inches(1.37), Inches(3.6), Inches(0.33), 'KEY METRIC', 9, WHITE, True)

rows_sh = [
    ('1\nAUTH', sig_colors[0],
     'Stolen TACACS+ creds, slow host growth over 14 days, off-hours logins creep up',
     'unique_hosts 5->18 over 14d, off_hours_ratio 0.1->0.45',
     'Slow drift caught by CUSUM accumulation'),
    ('2\nPROCESS', sig_colors[1],
     'Guest Shell activation, JumbledPath ELF binary, tcpdump on router',
     'New cmdline patterns, unsigned_ratio spikes',
     'defense_evasion (T1562) proj=0.49'),
    ('3\nNETWORK', sig_colors[2],
     'GRE tunnels established, slow CDR exfil, encrypted C2 channels',
     'bytes_out grows linearly, beacon_score elevates gradually',
     'data_exfiltration (T1041) proj=0.44'),
    ('4\nFILE', sig_colors[3],
     'Router startup-config modifications outside change windows, CALEA zone writes',
     'config writes outside approved windows, archive_creates spike',
     'defense_evasion (T1562) pattern'),
    ('5\nIDENTITY', sig_colors[4],
     'Local privileged accounts created, AAA/TACACS+ config modifications',
     'identity-graph hash changes, priv_esc 0->3',
     'privilege_escalation (T1078) proj=0.71'),
]
y = Inches(1.8)
for sig, color, behavior, drift_sig, metric in rows_sh:
    bg = CARD if rows_sh.index((sig, color, behavior, drift_sig, metric)) % 2 == 0 else WHITE
    _box(s, Inches(0.5), y, Inches(12.3), Inches(0.9), fill=bg)
    _rect(s, Inches(0.5), y, Inches(1.1), Inches(0.9), fill=color)
    _txt(s, Inches(0.55), y+Inches(0.2), Inches(1.0), Inches(0.5), sig, 8.5, WHITE, True, PP_ALIGN.CENTER)
    _txt(s, Inches(1.7), y+Pt(2), Inches(3.35), Inches(0.85), behavior, 8.5, BODY)
    _txt(s, Inches(5.2), y+Pt(2), Inches(3.75), Inches(0.85), drift_sig, 8.5, TEAL, fn='Cascadia Code')
    _txt(s, Inches(9.1), y+Pt(2), Inches(3.65), Inches(0.85), metric, 8.5, NAVY, True)
    y += Inches(0.93)

_bottom_bar(s, 'Salt Typhoon wins by going slow. CUSUM accumulates sub-threshold drift. 30-day rolling baseline -- dwell of years -> alarms in days.')


# SLIDE J: Peer Cohort + Alert Tiers
s = _content_slide('PART 2  --  ACECARD UEBA  --  COHORT + ALERTING',
                   'Peer Cohort Comparison + Calibrated Alert Tiers',
                   'A sysadmin logging into 30 hosts is normal. A finance analyst doing the same is anomalous.', TEAL)

# Left: Peer Cohort
_rect(s, Inches(0.5), Inches(1.4), Inches(5.8), Inches(0.32), fill=TEAL)
_txt(s, Inches(0.6), Inches(1.42), Inches(5.6), Inches(0.28), 'PEER COHORT ANALYSIS', 11, WHITE, True)
_box(s, Inches(0.5), Inches(1.78), Inches(5.8), Inches(5.3), fill=CARD, border=TEAL)

_txt(s, Inches(0.7), Inches(1.85), Inches(5.4), Inches(0.28),
     'z = (entity_drift - cohort_mean) / cohort_std', 10, NAVY, True, fn='Cascadia Code')
_txt(s, Inches(0.7), Inches(2.2), Inches(5.4), Inches(0.25),
     'Example: jsmith (financial_analyst, 42 peers)', 10, SUB)

cohort_vals = [
    ('jsmith drift:', '0.142'),
    ('Cohort mean:', '0.031'),
    ('Cohort std:', '0.018'),
    ('z-score:', '6.17  >>> OUTLIER'),
]
y = Inches(2.55)
for label, val in cohort_vals:
    _txt(s, Inches(0.8), y, Inches(2.5), Inches(0.28), label, 10, BODY, True)
    _txt(s, Inches(3.2), y, Inches(3.0), Inches(0.28), val, 10, CRED if '>>>' in val else NAVY, True, fn='Cascadia Code')
    y += Inches(0.3)

_txt(s, Inches(0.7), Inches(3.6), Inches(5.3), Inches(0.25), 'Drift direction:', 10, BODY, True)
_txt(s, Inches(0.7), Inches(3.9), Inches(5.3), Inches(0.35),
     'lateral_movement (0.78, T1021)\ncredential_dumping (0.45, T1003)', 10, TEAL, True)
_txt(s, Inches(0.7), Inches(4.35), Inches(5.3), Inches(0.25), 'Action:', 10, BODY, True)
_txt(s, Inches(0.7), Inches(4.65), Inches(5.3), Inches(0.35), 'HIGH -> SOAR playbook triggered', 11, CRED, True)

# Right: Alert Tiers
_rect(s, Inches(6.9), Inches(1.4), Inches(5.9), Inches(0.32), fill=sig_colors[1])
_txt(s, Inches(7.0), Inches(1.42), Inches(5.7), Inches(0.28), 'CALIBRATED ALERT TIERS', 11, WHITE, True)

alert_tiers = [
    ('CRITICAL', CRED, 'Health < 20 OR CUSUM ALARM (4-sigma)', 'Immediate SOC investigation. All access revoked.'),
    ('HIGH', ORG, 'Health < 40 OR Velocity z > 3', 'Tier-2 triage within 15 min. Step-up MFA.'),
    ('MEDIUM', GOLD, 'Health < 70 + change-point detected', 'Tier-1 review within 1 hour. Watchlist.'),
    ('LOW', TEAL, 'Health 70-85 AND cohort z > 2', 'Log for trend analysis. Watchlist entry.'),
    ('INFO', SUB, 'Any CUSUM trigger (4-sigma threshold)', 'Log for trend analysis. Baseline audit.'),
]
y = Inches(1.78)
for tier, color, trigger, action in alert_tiers:
    _box(s, Inches(6.9), y, Inches(5.9), Inches(1.0), fill=CARD, border=color)
    _rect(s, Inches(6.9), y, Inches(1.3), Inches(1.0), fill=color)
    _txt(s, Inches(6.95), y+Inches(0.3), Inches(1.2), Inches(0.4), tier, 10, WHITE, True, PP_ALIGN.CENTER)
    _txt(s, Inches(8.35), y+Inches(0.05), Inches(4.3), Inches(0.4), trigger, 9, BODY)
    _txt(s, Inches(8.35), y+Inches(0.5), Inches(4.3), Inches(0.4), action, 9, color, True)
    y += Inches(1.08)

_bottom_bar(s, 'Cohort z-score: same behavior is normal for sysadmins, anomalous for finance. Context is everything.')


# SLIDE 14: CUSUM Algorithm
s = _content_slide('PART 2  --  ACECARD UEBA  --  CUSUM ALGORITHM',
                   'CUSUM: Catching What Thresholds Miss',
                   'Page\'s Cumulative Sum detects small, sustained behavioral shifts over time -- the Volt Typhoon signature.', TEAL)

_box(s, Inches(0.5), Inches(1.5), Inches(5.8), Inches(2.8), fill=CARD, accent=TEAL)
_txt(s, Inches(0.8), Inches(1.55), Inches(5.3), Inches(0.3), 'THE ALGORITHM', 12, TEAL, True)
tf = _txt(s, Inches(0.8), Inches(1.9), Inches(5.3), Inches(2.2),
     'S(t) = max(0, S(t-1) + drift(t) - (mu + k))', 12, NAVY, True, fn='Cascadia Code')
_p(tf, '', 6, SUB)
_p(tf, 'S(t) = CUSUM statistic at time t', 10, BODY)
_p(tf, 'drift(t) = cosine distance at time t', 10, BODY)
_p(tf, 'mu = expected mean drift (30-day rolling baseline)', 10, BODY)
_p(tf, 'k = slack = 0.5 * sigma', 10, BODY)
_p(tf, 'ALARM when S(t) > h = 4 * sigma', 11, CRED, True)
_p(tf, '', 6, SUB)
_p(tf, 'Key: S(t) resets to 0 when drift returns to normal.', 10, BODY)
_p(tf, 'No false accumulation from benign noise.', 10, BODY)

_box(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(2.8), fill=CARD, accent=ORG)
_txt(s, Inches(7.1), Inches(1.55), Inches(5.3), Inches(0.3), 'CONCRETE NUMBERS', 12, ORG, True)
tf = _txt(s, Inches(7.1), Inches(1.9), Inches(5.3), Inches(2.2),
     'Typical baseline sigma = 0.03', 11, BODY)
_p(tf, '', 6, SUB)
_p(tf, 'k = 0.5 * 0.03 = 0.015', 11, NAVY, True)
_p(tf, '  (tolerates drift up to baseline + 0.015)', 10, SUB)
_p(tf, 'h = 4 * 0.03 = 0.12', 11, NAVY, True)
_p(tf, '  (alarm at 0.12 accumulated excess)', 10, SUB)
_p(tf, '', 6, SUB)
_p(tf, 'Sustained drift of 0.06/window:', 11, BODY)
_p(tf, '  excess = 0.06 - 0.03 - 0.015 = 0.015/window', 10, BODY)
_p(tf, '  ALARM in 0.12/0.015 = 8 windows (8 hours)', 11, CRED, True)

_box(s, Inches(0.5), Inches(4.5), Inches(12.1), Inches(2.3), fill=CARD, accent=NAVY)
_txt(s, Inches(0.8), Inches(4.55), Inches(11.6), Inches(0.3),
     'WHY CUSUM CATCHES VOLT TYPHOON BUT THRESHOLDS DO NOT', 12, NAVY, True)

_rect(s, Inches(0.8), Inches(5.0), Inches(5.5), Inches(0.3), fill=CRED)
_txt(s, Inches(0.85), Inches(5.02), Inches(5.4), Inches(0.26), 'FIXED THRESHOLD = 0.15', 10, WHITE, True)
tf = _txt(s, Inches(0.8), Inches(5.35), Inches(5.5), Inches(1.2),
     'Volt Typhoon drift per window = 0.01-0.04', 10, BODY)
_p(tf, 'NEVER crosses 0.15. Ever.', 10, CRED, True)
_p(tf, 'Attacker operates for 5+ YEARS undetected.', 10, CRED, True)

_rect(s, Inches(6.8), Inches(5.0), Inches(5.5), Inches(0.3), fill=TEAL)
_txt(s, Inches(6.85), Inches(5.02), Inches(5.4), Inches(0.26), 'CUSUM (h = 4*sigma = 0.12)', 10, WHITE, True)
tf = _txt(s, Inches(6.8), Inches(5.35), Inches(5.5), Inches(1.2),
     'Same drift 0.03-0.04 ACCUMULATES over windows', 10, BODY)
_p(tf, 'CUSUM alarm in 4-8 days.', 10, TEAL, True)
_p(tf, 'Dwell time: YEARS -> DAYS.', 10, TEAL, True)

_bottom_bar(s, 'NO COMMERCIAL UEBA uses CUSUM on behavioral embeddings. This is original research applied to cybersecurity.')


# SLIDE 15: Drift Direction - 8 Concepts
s = _content_slide('PART 2  --  ACECARD UEBA  --  DRIFT DIRECTION',
                   'Drift Direction: 8 MITRE-Mapped Threat Concepts',
                   'Each concept is text-embedded into the same 1536-d space. Cosine projection tells WHAT the entity is BECOMING.', TEAL)

concepts = [
    ('credential_dumping', 'T1003', sig_colors[0], 'LSASS access, SAM reads, ntdsutil, DCSync, Kerberoasting'),
    ('lateral_movement', 'T1021', sig_colors[1], 'RDP/WinRM/SSH to new hosts, admin shares, PtH/PtT, unique_hosts spike'),
    ('data_exfiltration', 'T1041', sig_colors[2], 'Bytes-out inversion, archive staging, cloud uploads, DNS tunneling'),
    ('c2_beaconing', 'T1071', sig_colors[3], 'Regular-interval callbacks, DGA domains, encrypted channels, KV botnet'),
    ('lotl_execution', 'T1059', sig_colors[4], 'wmic/ntdsutil/netsh/rundll32 chains, encoded PowerShell, deep trees'),
    ('privilege_escalation', 'T1078', RGBColor(0x8B,0x5C,0xF6), 'Group adds (Domain Admins), sudo abuse, token manipulation, delegation'),
    ('defense_evasion', 'T1562', RGBColor(0xDB,0x27,0x77), 'AV exclusions, log clearing, agent tampering, timestomping'),
    ('insider_hoarding', 'T1074', RGBColor(0x06,0xB6,0xD4), 'Steady sensitive access increase, bulk downloads, data staging over months'),
]

y = Inches(1.5)
for name, tid, color, desc in concepts:
    _box(s, Inches(0.5), y, Inches(12.1), Inches(0.6), fill=CARD)
    _rect(s, Inches(0.5), y, Inches(2.4), Inches(0.28), fill=color)
    _txt(s, Inches(0.55), y+Pt(1), Inches(2.3), Inches(0.25), name, 10, WHITE, True, PP_ALIGN.CENTER)
    _txt(s, Inches(3.0), y+Pt(1), Inches(0.8), Inches(0.25), tid, 10, NAVY, True, fn='Cascadia Code')
    _txt(s, Inches(3.9), y+Inches(0.02), Inches(8.5), Inches(0.55), desc, 10, BODY)
    y += Inches(0.63)

_box(s, Inches(0.5), Inches(6.6), Inches(12.1), Inches(0.6), fill=CARD, accent=GOLD)
_txt(s, Inches(0.8), Inches(6.62), Inches(11.6), Inches(0.55),
     'projection = cosine(drift_vector, concept_embedding). Top 3 returned with MITRE IDs. '
     'Result: "Entity drifting toward lateral_movement (0.82, T1021) + lotl_execution (0.71, T1059)" -- actionable, mappable, auditable.',
     10, NAVY, True)

_bottom_bar(s, 'NO OTHER UEBA projects drift direction onto MITRE techniques. This is the core innovation.')


# SLIDE 16: Entity Fusion
s = _content_slide('PART 2  --  ACECARD UEBA  --  ENTITY FUSION',
                   'Entity Fusion: Solving Identity Fragmentation',
                   'One human = 10+ identifiers across 6+ systems. Without fusion, lateral movement is invisible.', TEAL)

_rect(s, Inches(0.5), Inches(1.5), Inches(12.1), Inches(0.35), fill=NAVY)
_txt(s, Inches(0.7), Inches(1.52), Inches(11.7), Inches(0.3),
     'ONE PERSON = 10+ IDENTIFIERS  (John Smith, Tier-2 SOC Analyst)', 12, WHITE, True, PP_ALIGN.CENTER)

ids = [
    ('Active Directory', 'CORP\\jsmith / S-1-5-21-...', sig_colors[0]),
    ('Azure AD (Entra)', 'john.smith@corp.mil / ObjectId', sig_colors[1]),
    ('Okta SSO', 'john.smith@corp.mil / 00u1a2b3c', sig_colors[2]),
    ('AWS IAM', 'arn:aws:iam::842:user/jsmith', sig_colors[3]),
    ('Kubernetes', 'system:sa:soc:jsmith-sa', sig_colors[4]),
    ('CrowdStrike', 'aid:abc123 / JSMITH-WS01', sig_colors[0]),
    ('VPN/RADIUS', 'jsmith@corp.mil (cert DN)', sig_colors[1]),
    ('Splunk', 'src_user=jsmith', sig_colors[2]),
    ('PKI/CAC', 'CN=SMITH.JOHN.Q.123456', sig_colors[3]),
    ('TACACS+', 'jsmith (net device auth)', sig_colors[4]),
]
y = Inches(2.0)
for system, ident, color in ids:
    _rect(s, Inches(0.5), y, Inches(1.6), Inches(0.25), fill=color)
    _txt(s, Inches(0.55), y+Pt(1), Inches(1.5), Inches(0.22), system, 8.5, WHITE, True, PP_ALIGN.CENTER)
    _txt(s, Inches(2.2), y+Pt(1), Inches(3.0), Inches(0.22), ident, 8.5, BODY, fn='Cascadia Code')
    y += Inches(0.27)

_txt(s, Inches(5.5), Inches(3.2), Inches(0.5), Inches(0.5), '=>', 28, ORG, True)

_box(s, Inches(6.0), Inches(2.0), Inches(6.5), Inches(3.0), fill=ACE_BG, accent=TEAL)
_txt(s, Inches(6.3), Inches(2.05), Inches(6.0), Inches(0.3), 'RESOLVED ENTITY', 12, TEAL, True)
_txt(s, Inches(6.3), Inches(2.35), Inches(6.0), Inches(0.25),
     'entity_uuid: e7a2b3c4-9f1d-4e8b-a2c7-3d5f', 10, NAVY, True, fn='Cascadia Code')
tf = _txt(s, Inches(6.3), Inches(2.7), Inches(6.0), Inches(2.0),
     'Resolution layers:', 10, NAVY, True)
_p(tf, '1. Deterministic: email match (AD + AAD + Okta + VPN)', 9.5, BODY)
_p(tf, '2. Deterministic: EDIPI from CAC matches HR record', 9.5, BODY)
_p(tf, '3. Probabilistic: CrowdStrike agent on hostname', 9.5, BODY)
_p(tf, '4. Graph: connected component = all 10 identifiers', 9.5, BODY)
_p(tf, '', 6, SUB)
_p(tf, 'All 5 signals from ALL sources feed ONE embedding.', 10, NAVY, True)
_p(tf, 'Lateral across 5 systems = 1 strong signal.', 10, NAVY, True)

# Before/After comparison
_box(s, Inches(0.5), Inches(5.2), Inches(5.8), Inches(1.7), fill=CARD, accent=CRED)
_txt(s, Inches(0.8), Inches(5.25), Inches(5.3), Inches(0.25), 'WITHOUT ENTITY FUSION', 11, CRED, True)
tf = _txt(s, Inches(0.8), Inches(5.5), Inches(5.3), Inches(1.2),
     'AD: 1 normal logon (score 0.3)', 10, BODY)
_p(tf, 'AAD: 1 normal sign-in (score 0.25)', 10, BODY)
_p(tf, 'AWS: 1 normal login (score 0.2)', 10, BODY)
_p(tf, 'K8s: 1 normal kubectl (score 0.15)', 10, BODY)
_p(tf, '4 alerts below threshold. No correlation.', 10, CRED, True)

_box(s, Inches(6.8), Inches(5.2), Inches(5.8), Inches(1.7), fill=CARD, accent=TEAL)
_txt(s, Inches(7.1), Inches(5.25), Inches(5.3), Inches(0.25), 'WITH ENTITY FUSION', 11, TEAL, True)
tf = _txt(s, Inches(7.1), Inches(5.5), Inches(5.3), Inches(1.2),
     'ONE entity, ALL events fused:', 10, BODY)
_p(tf, '4 cross-system logins in 1 hour', 10, BODY)
_p(tf, 'Composite drift: 0.23 (significant)', 10, BODY)
_p(tf, '1 HIGH alert: lateral_movement (0.82)', 10, TEAL, True)

_bottom_bar(s, 'NO OTHER UEBA resolves 10+ identity systems into a unified behavioral entity. This is infrastructure-level innovation.')


# SLIDE 17: 4-Component Signature
s = _content_slide('PART 2  --  ACECARD UEBA  --  ENTITY SIGNATURE',
                   '4-Component Entity Signature: Forensic-Grade Evidence',
                   'Per entity, per window: identity proof + behavioral state + explainability + tamper evidence.', TEAL)

comps = [
    ('1. IDENTITY-GRAPH HASH', sig_colors[0],
     'SHA-256 of entity resolution graph (all linked IDs + correlation edges + confidence)',
     'Proves WHO was resolved. Detects identity graph manipulation. Full audit trail.'),
    ('2. BEHAVIORAL EMBEDDING', sig_colors[1],
     '1536-d vector from embedding model (text-embedding-3-small or local SLM)',
     'IS the behavioral state. Enables drift, direction, peer comparison, kill-chain reconstruction.'),
    ('3. STRUCTURAL FEATURES', sig_colors[2],
     'JSONB document with all ~30 raw feature values: {"auth":{"logon_count":47,"failure_rate":0.032,...}}',
     'Explainability: which specific metric drove the drift? Auditable numbers at detection time.'),
    ('4. HMAC-CHAINED HISTORY', sig_colors[3],
     'HMAC-SHA256: HMAC(prev_hmac + embedding_hash + features_hash + timestamp)',
     'Tamper-evident chain. Modifying ANY past snapshot breaks ALL subsequent HMACs. Court-admissible.'),
]
y = Inches(1.5)
for label, color, what, purpose in comps:
    _box(s, Inches(0.5), y, Inches(12.1), Inches(1.15), fill=CARD, accent=color)
    _rect(s, Inches(0.5), y+Inches(0.05), Inches(3.0), Inches(0.3), fill=color)
    _txt(s, Inches(0.55), y+Inches(0.07), Inches(2.9), Inches(0.26), label, 10, WHITE, True, PP_ALIGN.CENTER)
    _txt(s, Inches(3.7), y+Inches(0.05), Inches(8.7), Inches(0.35), what, 10, BODY)
    _txt(s, Inches(3.7), y+Inches(0.5), Inches(8.7), Inches(0.55), purpose, 10, NAVY, True)
    y += Inches(1.25)

_box(s, Inches(0.5), Inches(6.5), Inches(12.1), Inches(0.6), fill=CARD, accent=GOLD)
_txt(s, Inches(0.8), Inches(6.52), Inches(11.6), Inches(0.55),
     'Together: WHO was resolved (hash) + WHAT they did (embedding) + WHY the embedding looks that way (features) + '
     'WHEN it happened with proof of non-tampering (HMAC chain). Forensic-grade. Court-admissible. JWICS-deployable.',
     10, NAVY, True)

_bottom_bar(s, 'NO OTHER UEBA produces tamper-evident, forensic-grade behavioral evidence chains. This is novel.')


# SLIDE 18: ABAC Trust Loop
s = _content_slide('PART 2  --  ACECARD UEBA  --  ABAC TRUST LOOP',
                   'ABAC Trust State Machine: Proportional Enforcement',
                   'Not binary lockout. Graduated, reversible, proportional response.', TEAL)

states = [
    ('TRUSTED', TEAL, 'health > 70, no CUSUM concern',
     'Normal operations. Standard access. Baseline monitoring.', 'Default state'),
    ('ELEVATED_WATCH', GOLD, 'health < 70 OR cusum > 2*sigma',
     'Step-up MFA. Increased log verbosity. Analyst notification.', 'Reversal: health > 75 for 24h'),
    ('RESTRICTED', ORG, 'health < 40 OR cusum > 3.5*sigma OR concept > 0.55',
     'Read-only access. Lateral blocked. Incident ticket created.', 'Reversal: analyst + health > 60 for 48h'),
    ('BLOCKED', CRED, 'health < 20 OR cusum ALARM (4*sigma) OR concept > 0.70',
     'All access revoked. Sessions terminated. IR team notified.', 'Reversal: IR team manual ONLY'),
]
y = Inches(1.5)
for label, color, trigger, action, reversal in states:
    _box(s, Inches(0.5), y, Inches(12.1), Inches(1.1), fill=CARD, accent=color)
    _rect(s, Inches(0.5), y+Inches(0.05), Inches(2.3), Inches(0.35), fill=color)
    _txt(s, Inches(0.55), y+Inches(0.07), Inches(2.2), Inches(0.3), label, 14, WHITE, True, PP_ALIGN.CENTER)
    _txt(s, Inches(3.0), y+Inches(0.05), Inches(1.5), Inches(0.2), 'TRIGGER:', 9, color, True)
    _txt(s, Inches(3.0), y+Inches(0.25), Inches(9.4), Inches(0.25), trigger, 10, BODY)
    _txt(s, Inches(3.0), y+Inches(0.5), Inches(1.5), Inches(0.2), 'ACTION:', 9, color, True)
    _txt(s, Inches(3.0), y+Inches(0.7), Inches(9.4), Inches(0.25), action, 10, BODY)
    _txt(s, Inches(9.0), y+Inches(0.85), Inches(3.4), Inches(0.2), reversal, 8.5, SUB)
    y += Inches(1.2)

_bottom_bar(s, 'Detection to enforcement: < 5 minutes. System improves with analyst feedback. False positive = auto-adjusts baseline.')


# SLIDE K: Combined Architecture
s = prs.slides.add_slide(prs.slide_layouts[6])
_fill(s, WHITE)
_rect(s, Inches(0), Inches(0), Inches(6.67), Inches(0.3), fill=ORG)
_rect(s, Inches(6.67), Inches(0), Inches(6.67), Inches(0.3), fill=TEAL)
_txt(s, Inches(0.5), Inches(0.0), Inches(12.3), Inches(0.3),
     'PART 3  --  INTEGRATION  --  COMBINED ARCHITECTURE', 11, WHITE, True)
_txt(s, Inches(0.5), Inches(0.45), Inches(12.3), Inches(0.55),
     'End-to-End Architecture: 22CT Preemptive + ACECARD UEBA', 24, NAVY, True, fn='Georgia')

# INPUTS column
_rect(s, Inches(0.5), Inches(1.15), Inches(2.8), Inches(0.3), fill=NAVY)
_txt(s, Inches(0.55), Inches(1.17), Inches(2.7), Inches(0.26), 'INPUTS', 11, WHITE, True)
inputs = ['Threat Intel (CVE, CTI, IOCs)', 'Defense Config (NGFW, IAM, SASE)',
          'Telemetry (logs, netflow, EDR)', 'Schema (MITRE ATT&CK/D3FEND)']
y = Inches(1.55)
for inp in inputs:
    _box(s, Inches(0.5), y, Inches(2.8), Inches(0.5), fill=CARD, border=NAVY)
    _txt(s, Inches(0.65), y+Inches(0.08), Inches(2.5), Inches(0.35), inp, 9.5, BODY)
    y += Inches(0.58)

# LAYER 1 column
_rect(s, Inches(3.5), Inches(1.15), Inches(4.2), Inches(0.3), fill=ORG)
_txt(s, Inches(3.55), Inches(1.17), Inches(4.1), Inches(0.26), 'LAYER 1 -- 22CT PREEMPTIVE', 11, WHITE, True)
l1_items = ['Attack Intelligence', 'Defense Intelligence', 'Remediation Intelligence']
y = Inches(1.55)
for item in l1_items:
    _box(s, Inches(3.5), y, Inches(4.2), Inches(0.42), fill=RIGOR_BG, border=ORG)
    _txt(s, Inches(3.65), y+Inches(0.06), Inches(3.9), Inches(0.3), item, 10, BODY)
    y += Inches(0.48)
_rect(s, Inches(3.5), y, Inches(4.2), Inches(0.3), fill=ORG)
_txt(s, Inches(3.55), y+Inches(0.03), Inches(4.1), Inches(0.24), 'OUTPUT', 10, WHITE, True)
y += Inches(0.35)
l1_out = ['Hardened posture (config verified)', 'Verified SOC playbooks', 'MITRE TTP proof']
for item in l1_out:
    _txt(s, Inches(3.65), y, Inches(4.0), Inches(0.3), item, 9, ORG, True)
    y += Inches(0.28)

# LAYER 2 column
_rect(s, Inches(7.9), Inches(1.15), Inches(4.9), Inches(0.3), fill=TEAL)
_txt(s, Inches(7.95), Inches(1.17), Inches(4.8), Inches(0.26), 'LAYER 2 -- ACECARD UEBA', 11, WHITE, True)
l2_items = ['5 Signal Summarizers', 'Trajectory Analysis', 'Drift Direction (8 MITRE concepts)']
y = Inches(1.55)
for item in l2_items:
    _box(s, Inches(7.9), y, Inches(4.9), Inches(0.42), fill=ACE_BG, border=TEAL)
    _txt(s, Inches(8.05), y+Inches(0.06), Inches(4.6), Inches(0.3), item, 10, BODY)
    y += Inches(0.48)
_rect(s, Inches(7.9), y, Inches(4.9), Inches(0.3), fill=TEAL)
_txt(s, Inches(7.95), y+Inches(0.03), Inches(4.8), Inches(0.24), 'OUTPUT', 10, WHITE, True)
y += Inches(0.35)
l2_out = ['5-tier alerts with MITRE mapping', 'ABAC trust enforcement', '7-chart dashboard + SOAR']
for item in l2_out:
    _txt(s, Inches(8.05), y, Inches(4.7), Inches(0.3), item, 9, TEAL, True)
    y += Inches(0.28)

# Shared vocabulary bar
_rect(s, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.3), fill=NAVY)
_txt(s, Inches(0.7), Inches(5.92), Inches(12.0), Inches(0.26),
     'Shared vocabulary: MITRE ATT&CK technique IDs  |  Entity UUIDs  |  Verified TTP coverage', 10, GOLD, True, PP_ALIGN.CENTER)

_bottom_bar(s, 'Two layers. One vocabulary. Zero gaps between preemption and detection.')


# SLIDE L: 4 Integration Points
s = _content_slide('PART 3  --  INTEGRATION  --  DATA FLOW',
                   '22CT Preemptive + ACECARD: The Data Flow',
                   'Four concrete integration points. Each makes both products stronger.')

int_pts = [
    ('1', '22CT MITRE Coverage Map -> ACECARD Concept Library',
     'ACECARD loads 22CT\'s verified TTP list as high-confidence concepts. '
     'Alert load drops 30-50% because unpreempted TTPs get priority weight.', ORG),
    ('2', '22CT Verified Findings -> ACECARD ABAC Trust Floor',
     'Entities that touch 22CT-flagged exposed assets are immediately elevated in ABAC. '
     'Step-up MFA triggered on assets with active preemptive finding.', TEAL),
    ('3', 'ACECARD Confirmed TPs -> 22CT Threat-Intel Feed',
     'Behavioral detections feed back as IOCs into 22CT\'s threat-intel pipeline. '
     'Loop: ACECARD detection -> 22CT blocks the C2 path for all other entities.', sig_colors[1]),
    ('4', 'Shared Vocabulary: MITRE Technique IDs & Entity UUIDs',
     'Single entity model across both layers. '
     'MITRE technique ID is the rosetta stone: preemption proof + behavioral alert share the same TTP label.', GOLD),
]

y = Inches(1.5)
for num, title, desc, color in int_pts:
    _box(s, Inches(0.5), y, Inches(12.1), Inches(1.2), fill=CARD, accent=color)
    _rect(s, Inches(0.5), y, Inches(0.6), Inches(1.2), fill=color)
    _txt(s, Inches(0.52), y+Inches(0.35), Inches(0.56), Inches(0.5), num, 22, WHITE, True, PP_ALIGN.CENTER, 'Georgia')
    _txt(s, Inches(1.25), y+Inches(0.05), Inches(10.9), Inches(0.35), title, 13, NAVY, True)
    _txt(s, Inches(1.25), y+Inches(0.45), Inches(10.9), Inches(0.65), desc, 10, BODY)
    y += Inches(1.3)

_bottom_bar(s, 'These four flows are why 22CT Preemptive + ACECARD is a system, not a bundle.')


# SLIDE M: Volt Typhoon Both-Layers Walkthrough
s = _content_slide('PART 3  --  INTEGRATION  --  VOLT TYPHOON WALKTHROUGH',
                   'Volt Typhoon Walkthrough: Both Layers in Action',
                   'Phase-by-phase: what 22CT preempts, what ACECARD detects, what gets blocked.')

_rect(s, Inches(0.5), Inches(1.4), Inches(1.6), Inches(0.38), fill=NAVY)
_txt(s, Inches(0.55), Inches(1.42), Inches(1.5), Inches(0.33), 'PHASE', 10, WHITE, True)
_rect(s, Inches(2.2), Inches(1.4), Inches(3.2), Inches(0.38), fill=CRED)
_txt(s, Inches(2.25), Inches(1.42), Inches(3.1), Inches(0.33), 'ATTACKER ACTION', 10, WHITE, True)
_rect(s, Inches(5.5), Inches(1.4), Inches(3.6), Inches(0.38), fill=ORG)
_txt(s, Inches(5.55), Inches(1.42), Inches(3.5), Inches(0.33), '22CT PREEMPTS', 10, WHITE, True)
_rect(s, Inches(9.2), Inches(1.4), Inches(3.6), Inches(0.38), fill=TEAL)
_txt(s, Inches(9.25), Inches(1.42), Inches(3.5), Inches(0.33), 'ACECARD DETECTS', 10, WHITE, True)

rows_vm = [
    ('Phase 1\n0-12h', 'SOHO router VPN cred reuse (KV botnet)', 'T1133 segmentation verified; zero-trust path proved', 'Auth: off_hours spikes, drift 0.08'),
    ('Phase 2\n12-36h', 'wmic, ntdsutil, netsh, rundll32 chain', 'T1059 EDR script-control coverage proved', 'Process: lolbin 0->7, CUSUM alarm T+47min'),
    ('Phase 3\n36-72h', 'RDP/WinRM to Domain Controllers', 'T1021 SMB shadow rule detected, micro-seg proved', 'Network: unique_dest 3->18, z=6.17'),
    ('Phase 4\n72-96h', 'Scheduled tasks, registry persistence', 'T1556 config drift catches baseline deviation', 'Identity: svc_acct spike, T1078 flag'),
    ('Phase 5\n96h+', 'Cred harvest, NTDS.dit, network mapping', 'T1003 egress filtering verified; ntdsutil blocked', 'File: sensitive_access 0->24, CRITICAL -> BLOCKED'),
]
y = Inches(1.85)
for i, (phase, attacker, preempt, acecard) in enumerate(rows_vm):
    bg = CARD if i % 2 == 0 else WHITE
    _box(s, Inches(0.5), y, Inches(12.3), Inches(0.95), fill=bg)
    _txt(s, Inches(0.6), y+Pt(2), Inches(1.5), Inches(0.9), phase, 8.5, NAVY, True)
    _txt(s, Inches(2.2), y+Pt(2), Inches(3.1), Inches(0.9), attacker, 8.5, BODY)
    _txt(s, Inches(5.5), y+Pt(2), Inches(3.5), Inches(0.9), preempt, 8.5, RGBColor(0x1A, 0x4D, 0x8C))
    _txt(s, Inches(9.2), y+Pt(2), Inches(3.5), Inches(0.9), acecard, 8.5, TEAL)
    y += Inches(0.98)

_bottom_bar(s, 'Layer 1 closes the door. Layer 2 watches the room. Together: Volt Typhoon has nowhere to hide.')


# SLIDE N: Salt Typhoon Both-Layers Walkthrough
s = _content_slide('PART 3  --  INTEGRATION  --  SALT TYPHOON WALKTHROUGH',
                   'Salt Typhoon Walkthrough: Both Layers in Action',
                   'Salt Typhoon wins by going slow. CUSUM accumulates the small drifts. 30-day rolling baseline adapts.')

_rect(s, Inches(0.5), Inches(1.4), Inches(1.6), Inches(0.38), fill=NAVY)
_txt(s, Inches(0.55), Inches(1.42), Inches(1.5), Inches(0.33), 'PHASE', 10, WHITE, True)
_rect(s, Inches(2.2), Inches(1.4), Inches(3.2), Inches(0.38), fill=CRED)
_txt(s, Inches(2.25), Inches(1.42), Inches(3.1), Inches(0.33), 'ATTACKER ACTION', 10, WHITE, True)
_rect(s, Inches(5.5), Inches(1.4), Inches(3.6), Inches(0.38), fill=ORG)
_txt(s, Inches(5.55), Inches(1.42), Inches(3.5), Inches(0.33), '22CT PREEMPTS', 10, WHITE, True)
_rect(s, Inches(9.2), Inches(1.4), Inches(3.6), Inches(0.38), fill=TEAL)
_txt(s, Inches(9.25), Inches(1.42), Inches(3.5), Inches(0.33), 'ACECARD DETECTS', 10, WHITE, True)

rows_sm = [
    ('Phase 1\n0-24h', 'CVE-2023-20198 IOS XE exploit (/webui/)', 'T1190 WAF/mgmt-plane ACL blocks /webui/', 'Network: anomalous mgmt-plane IP surfaced'),
    ('Phase 2\n1-7d', 'Local admin account + AAA config mod', 'T1556 config drift detects baseline diff <1hr', 'Identity: new SID linked, priv_esc flag'),
    ('Phase 3\n7-30d', 'TACACS+ capture on TCP/49 via pcap', 'Mgmt-plane segmentation, on-box capture restricted', 'Process: tcpdump outside window, CUSUM'),
    ('Phase 4\n30-90d', 'JumbledPath pivot chains, GRE tunnel', 'Outbound C2 IPs blocked; GRE peer verified', 'Composite: same entity across 4 hops'),
    ('Phase 5\n90d+', 'GRE exfil, CALEA CDR/metadata harvest', 'GRE peer verified; CALEA-zone segmentation proved', 'Network: data_exfil 0.44, CUSUM 4-sigma'),
]
y = Inches(1.85)
for i, (phase, attacker, preempt, acecard) in enumerate(rows_sm):
    bg = CARD if i % 2 == 0 else WHITE
    _box(s, Inches(0.5), y, Inches(12.3), Inches(0.95), fill=bg)
    _txt(s, Inches(0.6), y+Pt(2), Inches(1.5), Inches(0.9), phase, 8.5, NAVY, True)
    _txt(s, Inches(2.2), y+Pt(2), Inches(3.1), Inches(0.9), attacker, 8.5, BODY)
    _txt(s, Inches(5.5), y+Pt(2), Inches(3.5), Inches(0.9), preempt, 8.5, RGBColor(0x1A, 0x4D, 0x8C))
    _txt(s, Inches(9.2), y+Pt(2), Inches(3.5), Inches(0.9), acecard, 8.5, TEAL)
    y += Inches(0.98)

_bottom_bar(s, 'Salt Typhoon wins by going slow. CUSUM accumulates sub-threshold drift. 30-day rolling baseline -- dwell of years -> alarms in days.')


# ================================================================
# PART 3: INTEGRATION
# ================================================================

# SLIDE 19: Coverage Matrix
s = prs.slides.add_slide(prs.slide_layouts[6])
_fill(s, WHITE)
_rect(s, Inches(0), Inches(0), Inches(6.67), Inches(0.3), fill=ORG)
_rect(s, Inches(6.67), Inches(0), Inches(6.67), Inches(0.3), fill=TEAL)
_txt(s, Inches(0.5), Inches(0.0), Inches(12.3), Inches(0.3),
     'PART 3  --  INTEGRATION  --  COVERAGE MATRIX', 11, WHITE, True)
_txt(s, Inches(0.5), Inches(0.5), Inches(12.3), Inches(0.6),
     'Combined Coverage: Which Layer Owns What', 24, NAVY, True, fn='Georgia')
_txt(s, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.3),
     'Each adversary capability owned by the right layer. Preemption first, behavioral as backstop.', 13, SUB)

mx_rows = [
    ('Edge-device CVE exploitation', 'Y', 'Y', 'PREEMPT', 'Config is the gap -- formal proof patch enforced', ORG),
    ('Misconfigured / shadowed FW rules', 'Y', 'Y', 'PREEMPT', 'Math model finds shadows across multi-vendor', ORG),
    ('Router ACL tampering / drift', '', 'Y', 'PREEMPT', 'Config drift re-verifies hourly vs baseline', ORG),
    ('C2 to known malicious infra', 'Y', 'Y', 'BOTH', 'Preemptive blocks known; ACECARD catches unknown', RGBColor(0x6C,0x2E,0xD9)),
    ('Compromised valid accounts', 'Y', 'Y', 'ACECARD', 'Policy-allowed traffic; only drift reveals abuse', TEAL),
    ('Living-off-the-Land (T1059)', 'Y', 'Y', 'ACECARD', 'lotl_execution concept + process signal shift', TEAL),
    ('Cross-domain lateral movement', 'Y', 'Y', 'ACECARD', 'Entity fusion joins all identities -> one drift', TEAL),
    ('Slow APT (months/years)', 'Y', 'Y', 'ACECARD', 'CUSUM 4-sigma accumulates sub-threshold drift', TEAL),
    ('Zero-day / unknown TTPs', 'Y', 'Y', 'ACECARD', 'No MITRE label yet -> behavioral anomaly only', TEAL),
    ('Insider threat / data hoarding', '', '', 'ACECARD', 'insider_hoarding (T1074) + cohort z-score', TEAL),
]

_rect(s, Inches(0.5), Inches(1.4), Inches(12.1), Inches(0.35), fill=NAVY)
_txt(s, Inches(0.6), Inches(1.42), Inches(3.8), Inches(0.3), 'Adversary Capability', 9, WHITE, True)
_txt(s, Inches(4.5), Inches(1.42), Inches(0.5), Inches(0.3), 'VT', 9, WHITE, True, PP_ALIGN.CENTER)
_txt(s, Inches(5.1), Inches(1.42), Inches(0.5), Inches(0.3), 'ST', 9, WHITE, True, PP_ALIGN.CENTER)
_txt(s, Inches(5.7), Inches(1.42), Inches(1.2), Inches(0.3), 'Layer', 9, WHITE, True, PP_ALIGN.CENTER)
_txt(s, Inches(7.0), Inches(1.42), Inches(5.5), Inches(0.3), 'Why This Layer Owns It', 9, WHITE, True)

y = Inches(1.8)
for i, (cap, vt, st, layer, why, color) in enumerate(mx_rows):
    bg = CARD if i % 2 == 0 else WHITE
    _box(s, Inches(0.5), y, Inches(12.1), Inches(0.48), fill=bg)
    _txt(s, Inches(0.6), y+Pt(2), Inches(3.8), Inches(0.43), cap, 9.5, NAVY, True)
    _txt(s, Inches(4.5), y+Pt(2), Inches(0.5), Inches(0.43), vt, 9.5, CRED if vt else SUB, al=PP_ALIGN.CENTER)
    _txt(s, Inches(5.1), y+Pt(2), Inches(0.5), Inches(0.43), st, 9.5, CRED if st else SUB, al=PP_ALIGN.CENTER)
    _rect(s, Inches(5.75), y+Inches(0.08), Inches(1.1), Inches(0.24), fill=color)
    _txt(s, Inches(5.75), y+Inches(0.09), Inches(1.1), Inches(0.22), layer, 8.5, WHITE, True, PP_ALIGN.CENTER)
    _txt(s, Inches(7.0), y+Pt(2), Inches(5.5), Inches(0.43), why, 9, BODY)
    y += Inches(0.5)

_bottom_bar(s, 'Together: the attack surface never existed (Preemptive) OR the attack is detected and contained within hours (ACECARD + ABAC)')


# SLIDE O: Complementary Tools
s = _content_slide('PART 3  --  INTEGRATION  --  COMPLEMENTARY TOOLS',
                   'Complementary Tools: What Rounds Out Defense-in-Depth',
                   '22CT Preemptive + ACECARD is the core. These technologies close residual gaps.')

tools = [
    ('EDR/XDR', 'CrowdStrike, MS Defender, SentinelOne',
     'In-memory LOTL, process trees', 'Feeds ACECARD process signal'),
    ('Email Security', 'Proofpoint, MS Defender for O365, Abnormal',
     'Primary initial-access vector', 'Feeds identity signal'),
    ('NDR', 'Vectra, ExtraHop, Darktrace, Corelight',
     'East-west packet inspection', 'Feeds network signal'),
    ('SIEM', 'Splunk, Sentinel, Chronicle, Elastic',
     'Log aggregation, retention', 'ACECARD reads from SIEM'),
    ('SOAR', 'XSOAR, Splunk SOAR, Tines, Torq',
     'Automated playbooks', 'Receives alerts with MITRE + ABAC'),
    ('Identity/ZTNA', 'Okta, Entra, Ping, Zscaler',
     'Phishing-resistant MFA', 'Step-up MFA from ABAC'),
    ('OT/ICS', 'Dragos, Claroty, Nozomi, TXOne',
     'OT protocol awareness', 'Models IT/OT boundary'),
    ('Deception', 'Thinkst Canary, Acalvio',
     'High-fidelity tripwires', 'ACECARD correlates canary alerts with drift'),
]

tool_colors = [TEAL, ORG, sig_colors[1], NAVY, sig_colors[3], sig_colors[2], CRED, GOLD]
xs = [Inches(0.5), Inches(3.8), Inches(7.1), Inches(10.4)]
ys = [Inches(1.5), Inches(3.5)]
for i, (name, vendors, role, integration) in enumerate(tools):
    col = i % 4
    row = i // 4
    x = xs[col]
    y = ys[row]
    color = tool_colors[i]
    _box(s, x, y, Inches(3.1), Inches(1.7), fill=CARD, border=color)
    _rect(s, x, y, Inches(3.1), Inches(0.3), fill=color)
    _txt(s, x+Inches(0.1), y+Inches(0.03), Inches(2.9), Inches(0.24), name, 11, WHITE, True)
    _txt(s, x+Inches(0.1), y+Inches(0.35), Inches(2.9), Inches(0.3), vendors, 8, SUB)
    _txt(s, x+Inches(0.1), y+Inches(0.7), Inches(2.9), Inches(0.3), role, 9, BODY)
    _txt(s, x+Inches(0.1), y+Inches(1.05), Inches(2.9), Inches(0.55), 'Integration: ' + integration, 8.5, NAVY, True)

_bottom_bar(s, '22CT Preemptive + ACECARD is the core. EDR + Email + NDR close the remaining gaps for full defense-in-depth.')


# SLIDE P: Reference Architecture 5-Layer
s = _content_slide('PART 3  --  INTEGRATION  --  REFERENCE ARCHITECTURE',
                   'Reference Architecture: Comprehensive Defense Stack',
                   'Five layers. Every layer has a job.')

layer_data = [
    ('1. THREAT INTELLIGENCE', NAVY,
     'CISA/NSA/FBI Advisories  |  Vendor PSIRTs  |  Commercial CTI Feeds  |  MITRE ATT&CK / D3FEND'),
    ('2. PREEMPTIVE LAYER  (22CT Preemptive)', ORG,
     'Attack Intelligence  |  Defense Intelligence  |  Remediation Intelligence  |  Continuous Re-Verify (hourly + on config change)'),
    ('3. CONTROL PLANE', RGBColor(0x16, 0x6B, 0x3A),
     'NGFW / IPS / WAF  |  Identity (IAM/PAM/IdP)  |  EDR / XDR  |  Email / Web / DLP  |  ZTNA / SASE / CASB'),
    ('4. BEHAVIORAL LAYER  (ACECARD UEBA)', TEAL,
     '5 Signal Summarizers  |  1536-d Embedding  |  Trajectory Analysis  |  Drift Direction (8 MITRE concepts)  |  ABAC Trust Loop'),
    ('5. RESPONSE & AUDIT', sig_colors[1],
     'SIEM / Data Lake  |  SOAR / Playbooks  |  ServiceNow / JIRA Service Mgmt  |  Audit / Compliance / Forensics'),
]

y = Inches(1.5)
for label, color, content in layer_data:
    _rect(s, Inches(0.5), y, Inches(2.5), Inches(0.85), fill=color)
    _txt(s, Inches(0.6), y+Inches(0.1), Inches(2.3), Inches(0.65), label, 9, WHITE, True)
    _box(s, Inches(3.1), y, Inches(9.7), Inches(0.85), fill=CARD, border=color)
    _txt(s, Inches(3.3), y+Inches(0.2), Inches(9.3), Inches(0.55), content, 10, BODY)
    y += Inches(0.95)

_bottom_bar(s, 'Each layer has a defined owner and a defined job. No gaps. No overlaps. No ambiguity about what defeats what.')


# SLIDE Q: Verdict
s = _content_slide('PART 3  --  INTEGRATION  --  VERDICT',
                   'Verdict: What the Combined Stack Delivers')

_box(s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.5), fill=CARD, accent=TEAL)
_rect(s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(0.35), fill=TEAL)
_txt(s, Inches(0.7), Inches(1.52), Inches(5.6), Inches(0.3), 'WHAT IT COVERS', 13, WHITE, True)
covers = [
    'Every known MITRE TTP for Volt + Salt Typhoon -- preempted at config layer',
    'Living-off-the-Land tradecraft -- caught as behavioral drift',
    'Cross-domain identity fusion -- one trajectory across AD/AAD/AWS/Okta/K8s',
    'Zero-day & unmapped TTPs -- caught as deviation from baseline',
    'Insider threat & valid-account abuse -- cohort z-score + cosine drift',
    'Long-dwell stealth -- CUSUM 4-sigma on rolling 30-day baseline',
]
y = Inches(2.0)
for item in covers:
    _rect(s, Inches(0.65), y+Inches(0.06), Inches(0.12), Inches(0.12), fill=TEAL)
    _txt(s, Inches(0.9), y, Inches(5.4), Inches(0.35), item, 9.5, BODY)
    y += Inches(0.38)

_box(s, Inches(6.9), Inches(1.5), Inches(6.0), Inches(5.5), fill=CARD, accent=CRED)
_rect(s, Inches(6.9), Inches(1.5), Inches(6.0), Inches(0.35), fill=CRED)
_txt(s, Inches(7.1), Inches(1.52), Inches(5.6), Inches(0.3), 'REMAINING GAPS (BY DESIGN)', 13, WHITE, True)
gaps_q = [
    'Endpoint runtime enforcement -- pair with CrowdStrike / MS Defender',
    'Email & phishing -- pair with Proofpoint / Defender for O365',
    'OT/ICS protocols -- pair with Dragos / Claroty / Nozomi',
    'Physical & supply-chain -- pair with TPRM platforms',
    'Human factors -- awareness training + tabletops',
]
y = Inches(2.0)
for item in gaps_q:
    _rect(s, Inches(7.05), y+Inches(0.06), Inches(0.12), Inches(0.12), fill=CRED)
    _txt(s, Inches(7.3), y, Inches(5.4), Inches(0.35), item, 9.5, BODY)
    y += Inches(0.38)

_bottom_bar(s, '22CT Preemptive + ACECARD UEBA = comprehensive defense against Volt + Salt Typhoon. Add EDR + Email + NDR for full defense-in-depth.')


# SLIDE 20: Innovation Proof
s = prs.slides.add_slide(prs.slide_layouts[6])
_fill(s, WHITE)
_rect(s, Inches(0), Inches(0), Inches(6.67), Inches(0.3), fill=ORG)
_rect(s, Inches(6.67), Inches(0), Inches(6.67), Inches(0.3), fill=TEAL)
_txt(s, Inches(0.5), Inches(0.0), Inches(12.3), Inches(0.3),
     'PART 3  --  INNOVATION PROOF  --  NOBODY ELSE HAS THIS', 11, WHITE, True)
_txt(s, Inches(0.5), Inches(0.5), Inches(12.3), Inches(0.6),
     'What Exists vs. What We Built', 24, NAVY, True, fn='Georgia')
_txt(s, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.3),
     'Systematic comparison of our innovations against every commercial UEBA and preemptive platform.', 13, SUB)

innovations_proof = [
    ('Formal verification of 2^8000 config state', 'Pentests sample. ASM enumerates. Nobody PROVES.', '22CT PREEMPT', ORG),
    ('Zero false positives / zero false negatives', 'Mathematical guarantee, not probabilistic.', '22CT PREEMPT', ORG),
    ('Behavioral embedding in unified 1536-d space', 'Industry uses features/rules. We use vectors.', 'ACECARD ONLY', TEAL),
    ('Per-signal drift decomposition (5 individual)', 'Others embed once. We embed 6x for attribution.', 'ACECARD ONLY', TEAL),
    ('CUSUM on embedding trajectories', 'Nobody applies sequential change-point to vectors.', 'ACECARD ONLY', TEAL),
    ('Drift direction with MITRE technique mapping', '"Drifting toward T1021" -- no other UEBA does this.', 'ACECARD ONLY', TEAL),
    ('Entity fusion across 10+ identity systems', 'Commercial UEBA: 1-3 sources max. We fuse 10+.', 'ACECARD ONLY', TEAL),
    ('HMAC-chained forensic evidence', 'Tamper-evident behavioral chains. Novel.', 'ACECARD ONLY', TEAL),
    ('ABAC trust loop with behavioral triggers', 'Detection + enforcement in one system. <5 min.', 'ACECARD ONLY', TEAL),
    ('Preemptive + behavioral in one platform', 'Nobody combines formal config proof + behavioral AI.', 'COMBINED', NAVY),
]

y = Inches(1.4)
for innovation, why_unique, who, color in innovations_proof:
    _box(s, Inches(0.5), y, Inches(12.1), Inches(0.48), fill=CARD)
    _rect(s, Inches(0.5), y, Inches(0.08), Inches(0.48), fill=color)
    _txt(s, Inches(0.7), y+Pt(2), Inches(5.5), Inches(0.43), innovation, 10, NAVY, True)
    _txt(s, Inches(6.3), y+Pt(2), Inches(4.0), Inches(0.43), why_unique, 9.5, BODY)
    _rect(s, Inches(10.5), y+Inches(0.08), Inches(2.0), Inches(0.25), fill=color)
    _txt(s, Inches(10.5), y+Inches(0.09), Inches(2.0), Inches(0.23), who, 8.5, WHITE, True, PP_ALIGN.CENTER)
    y += Inches(0.51)

_bottom_bar(s, 'Combined: 10 innovations that do not exist anywhere else in the industry. This is not incremental improvement.')


# SLIDE 21: Key Numbers
s = prs.slides.add_slide(prs.slide_layouts[6])
_fill(s, WHITE)
_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.3), fill=NAVY)
_txt(s, Inches(0.5), Inches(0.0), Inches(12.3), Inches(0.3),
     'REFERENCE  --  KEY NUMBERS', 11, WHITE, True)
_txt(s, Inches(0.5), Inches(0.5), Inches(12.3), Inches(0.6),
     'Key Numbers to Memorize', 24, NAVY, True, fn='Georgia')

left = [
    ('5+ years', 'Volt Typhoon dwell time'),
    ('4+ years', 'Salt Typhoon dwell time'),
    ('200+', 'Telcos compromised by Salt Typhoon'),
    ('9', 'US telecoms confirmed compromised'),
    ('99%', 'FW breaches = misconfig (Gartner)'),
    ('~20%', 'Attacks stopped by KEV patching'),
    ('2^8000', 'Config state space (cannot be sampled)'),
    ('50M+', 'Rule combinations 22CT reasons over'),
    ('0', 'False positives (mathematical proof)'),
    ('0', 'False negatives (exhaustive reasoning)'),
    ('<1 hour', 'Config drift detection time'),
]
right = [
    ('1,536', 'Embedding dims per signal'),
    ('9,216', 'Total dims (6 x 1536) per entity'),
    ('5', 'Behavioral signals'),
    ('8', 'MITRE threat concepts'),
    ('4-sigma', 'CUSUM alarm threshold'),
    ('<3 sec', 'Event to analysis per entity'),
    ('10K/hr', 'Entities per 4-vCPU worker'),
    ('$90M', '22CT Army SOC/MDR contract'),
    ('800+', 'Cleared analysts on contract'),
    ('90 days', 'Contract to first detection'),
    ('IL5/IL6/JWICS', 'Three enclave support'),
]

for nums, xs in [(left, Inches(0.5)), (right, Inches(6.8))]:
    y = Inches(1.3)
    for num, meaning in nums:
        _box(s, xs, y, Inches(6.0), Inches(0.42), fill=CARD)
        _rect(s, xs+Inches(0.03), y+Inches(0.04), Inches(1.3), Inches(0.3), fill=NAVY)
        _txt(s, xs+Inches(0.08), y+Inches(0.06), Inches(1.2), Inches(0.26),
             num, 10.5, WHITE, True, PP_ALIGN.CENTER, fn='Cascadia Code')
        _txt(s, xs+Inches(1.5), y+Inches(0.06), Inches(4.3), Inches(0.26), meaning, 10, BODY)
        y += Inches(0.45)


# ================================================================
# SLIDE 22: CLOSING
# ================================================================
s = _title_slide(
    'BUILD YOUR DOME',
    'Preemptive closes the door. ACECARD watches the room.\nTogether: comprehensive defense achieved.',
    '22nd Century Technologies Inc.  |  Contact: ravindra.shukla@gmail.com  |  CFIC Augusta GA  |  May 6, 2026')
_txt(s, Inches(1), Inches(4.5), Inches(11.3), Inches(1.5),
     'PREEMPTIVE: Formal mathematical proof that every known TTP path is blocked.\n'
     'BEHAVIORAL: Continuous 1536-d trajectory intelligence that catches what preemption cannot.\n'
     'INTEGRATED: Shared MITRE vocabulary, trust-floor coupling, closed-loop threat intel.',
     12, RGBColor(0xA0, 0xB0, 0xC0), False, PP_ALIGN.CENTER)


# ================================================================
# APPENDIX A: CUSUM MATHEMATICAL DEEP DIVE
# ================================================================
s = _content_slide('APPENDIX A  --  MATHEMATICAL FOUNDATION',
                   'CUSUM Algorithm: Full Mathematical Derivation',
                   "Page's Cumulative Sum (1954) -- sequential change-point detection applied to behavioral embedding drift.", NAVY)

_box(s, Inches(0.5), Inches(1.5), Inches(5.8), Inches(2.5), fill=CARD, accent=TEAL)
_txt(s, Inches(0.8), Inches(1.55), Inches(5.3), Inches(0.3), 'STEP 1: BEHAVIORAL DRIFT MEASUREMENT', 11, TEAL, True)
tf = _txt(s, Inches(0.8), Inches(1.85), Inches(5.3), Inches(2.0),
     'For each entity e at time window t:', 10, BODY)
_p(tf, '', 4, SUB)
_p(tf, 'drift(e, t) = 1 - cosine_sim(V(e,t), V(e,t-1))', 11, NAVY, True, Pt(2))
_p(tf, '', 4, SUB)
_p(tf, 'V(e,t) = 1536-d composite behavioral embedding', 10, BODY)
_p(tf, 'V(e,t-1) = previous window embedding', 10, BODY)
_p(tf, 'Range: 0.0 (identical) to 2.0 (opposite)', 10, BODY)
_p(tf, '', 4, SUB)
_p(tf, 'Typical healthy entity: drift ~ 0.02-0.05 per window', 10, SUB)
_p(tf, 'Compromised entity: drift ~ 0.06-0.15 per window', 10, CRED)

_box(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(2.5), fill=CARD, accent=ORG)
_txt(s, Inches(7.1), Inches(1.55), Inches(5.3), Inches(0.3), 'STEP 2: BASELINE ESTIMATION', 11, ORG, True)
tf = _txt(s, Inches(7.1), Inches(1.85), Inches(5.3), Inches(2.0),
     'Rolling baseline over N=30 windows:', 10, BODY)
_p(tf, '', 4, SUB)
_p(tf, 'mu(e) = mean(drift(e, t-N .. t-1))', 11, NAVY, True, Pt(2))
_p(tf, 'sigma(e) = stdev(drift(e, t-N .. t-1))', 11, NAVY, True, Pt(2))
_p(tf, '', 4, SUB)
_p(tf, 'mu captures the entity\'s NORMAL drift rate', 10, BODY)
_p(tf, 'sigma captures how much that rate naturally varies', 10, BODY)
_p(tf, '', 4, SUB)
_p(tf, 'Example: mu=0.03, sigma=0.03', 10, SUB)
_p(tf, 'This entity normally drifts ~0.03/window +/- 0.03', 10, SUB)

_box(s, Inches(0.5), Inches(4.2), Inches(5.8), Inches(3.0), fill=CARD, accent=NAVY)
_txt(s, Inches(0.8), Inches(4.25), Inches(5.3), Inches(0.3), 'STEP 3: CUSUM ACCUMULATION', 11, NAVY, True)
tf = _txt(s, Inches(0.8), Inches(4.55), Inches(5.3), Inches(2.5),
     'Two parameters control sensitivity:', 10, BODY)
_p(tf, '', 4, SUB)
_p(tf, 'k (slack) = 0.5 * sigma  (allowable noise)', 11, NAVY, True, Pt(2))
_p(tf, 'h (threshold) = 4 * sigma  (alarm level)', 11, NAVY, True, Pt(2))
_p(tf, '', 4, SUB)
_p(tf, 'The CUSUM statistic:', 10, BODY)
_p(tf, 'S(t) = max(0,  S(t-1) + drift(t) - mu - k)', 12, CRED, True, Pt(2))
_p(tf, '', 4, SUB)
_p(tf, 'Key properties:', 10, BODY)
_p(tf, '  S(t) can only increase when drift > mu + k', 10, BODY)
_p(tf, '  S(t) resets to 0 when drift returns to normal', 10, BODY)
_p(tf, '  S(t) ACCUMULATES excess drift over time', 10, BODY)
_p(tf, '  ALARM when S(t) >= h', 10, CRED, True)

_box(s, Inches(6.8), Inches(4.2), Inches(5.8), Inches(3.0), fill=CARD, accent=CRED)
_txt(s, Inches(7.1), Inches(4.25), Inches(5.3), Inches(0.3), 'STEP 4: WORKED EXAMPLE (VOLT TYPHOON)', 11, CRED, True)
tf = _txt(s, Inches(7.1), Inches(4.55), Inches(5.3), Inches(2.5),
     'Entity baseline: mu=0.03, sigma=0.03', 10, BODY)
_p(tf, 'k = 0.5 * 0.03 = 0.015', 10, NAVY, True)
_p(tf, 'h = 4 * 0.03 = 0.12', 10, NAVY, True)
_p(tf, '', 4, SUB)
_p(tf, 'APT drift per window = 0.06 (below threshold!)', 10, BODY)
_p(tf, 'Fixed threshold = 0.15 --> NEVER fires', 10, CRED, True)
_p(tf, '', 4, SUB)
_p(tf, 'CUSUM accumulation:', 10, BODY)
_p(tf, '  Window 1: S = 0 + 0.06 - 0.03 - 0.015 = 0.015', 9, BODY)
_p(tf, '  Window 2: S = 0.015 + 0.015 = 0.030', 9, BODY)
_p(tf, '  Window 3: S = 0.030 + 0.015 = 0.045', 9, BODY)
_p(tf, '  Window 4: S = 0.045 + 0.015 = 0.060', 9, BODY)
_p(tf, '  ...', 9, SUB)
_p(tf, '  Window 8: S = 0.120 >= h = 0.12  -->  ALARM!', 10, CRED, True)
_p(tf, '', 4, SUB)
_p(tf, 'Dwell time: 5 YEARS (threshold) vs 8 HOURS (CUSUM)', 11, TEAL, True)


# ================================================================
# APPENDIX B: CUSUM vs THRESHOLD COMPARISON
# ================================================================
s = _content_slide('APPENDIX A (cont.)  --  CUSUM vs ALTERNATIVES',
                   'Why CUSUM Beats Every Alternative',
                   'Systematic comparison of sequential detection methods for behavioral embedding drift.', NAVY)

_box(s, Inches(0.5), Inches(1.5), Inches(12.1), Inches(0.6), fill=CARD, accent=TEAL)
_txt(s, Inches(0.8), Inches(1.52), Inches(11.6), Inches(0.55),
     'THE CORE INSIGHT: Slow APTs produce drift that NEVER exceeds a fixed threshold in any single window. '
     'CUSUM detects these by accumulating small excesses over time -- the mathematical equivalent of "death by a thousand cuts."',
     10, NAVY, True)

methods = [
    ('Fixed Threshold', 'Alert if drift(t) > T',
     'Simple, low compute', 'Misses slow drift entirely. T=0.15: Volt Typhoon undetected for 5+ years.', CRED),
    ('Z-Score', 'Alert if (drift(t)-mu)/sigma > z',
     'Adapts to entity baseline', 'Still single-window. z=3: needs 3-sigma spike. Misses gradual shift.', CRED),
    ('EWMA', 'Z(t) = lambda*drift(t) + (1-lambda)*Z(t-1)',
     'Smooths noise, catches trends', 'Lambda tuning fragile. Decays old signal. No formal ARL guarantee.', ORG),
    ('CUSUM (Ours)', 'S(t) = max(0, S(t-1)+drift(t)-mu-k)',
     'Accumulates excess. Resets on return.', 'Optimal: minimizes detection delay for given false alarm rate (Lorden 1971).', TEAL),
]

_rect(s, Inches(0.5), Inches(2.3), Inches(2.0), Inches(0.35), fill=NAVY)
_txt(s, Inches(0.55), Inches(2.32), Inches(1.9), Inches(0.3), 'METHOD', 10, WHITE, True)
_rect(s, Inches(2.5), Inches(2.3), Inches(3.0), Inches(0.35), fill=NAVY)
_txt(s, Inches(2.55), Inches(2.32), Inches(2.9), Inches(0.3), 'FORMULA', 10, WHITE, True)
_rect(s, Inches(5.5), Inches(2.3), Inches(3.0), Inches(0.35), fill=NAVY)
_txt(s, Inches(5.55), Inches(2.32), Inches(2.9), Inches(0.3), 'STRENGTH', 10, WHITE, True)
_rect(s, Inches(8.5), Inches(2.3), Inches(4.1), Inches(0.35), fill=NAVY)
_txt(s, Inches(8.55), Inches(2.32), Inches(4.0), Inches(0.3), 'WEAKNESS / VERDICT', 10, WHITE, True)

y = Inches(2.7)
for name, formula, strength, weakness, color in methods:
    bg = CARD if methods.index((name, formula, strength, weakness, color)) % 2 == 0 else WHITE
    _box(s, Inches(0.5), y, Inches(12.1), Inches(0.85), fill=bg)
    _rect(s, Inches(0.5), y, Inches(0.08), Inches(0.85), fill=color)
    _txt(s, Inches(0.7), y+Pt(2), Inches(1.7), Inches(0.4), name, 10, NAVY, True)
    _txt(s, Inches(2.5), y+Pt(2), Inches(3.0), Inches(0.8), formula, 9, BODY, fn='Cascadia Code')
    _txt(s, Inches(5.5), y+Pt(2), Inches(3.0), Inches(0.8), strength, 9, BODY)
    _txt(s, Inches(8.5), y+Pt(2), Inches(4.0), Inches(0.8), weakness, 9, color, True)
    y += Inches(0.88)

_box(s, Inches(0.5), Inches(6.3), Inches(12.1), Inches(0.8), fill=CARD, accent=GOLD)
_txt(s, Inches(0.8), Inches(6.32), Inches(11.6), Inches(0.25), 'THEORETICAL GUARANTEE (Lorden 1971, Moustakides 1986)', 10, GOLD, True)
_txt(s, Inches(0.8), Inches(6.6), Inches(11.6), Inches(0.4),
     'CUSUM is MINIMAX OPTIMAL: among all sequential tests with the same false-alarm rate (ARL0), CUSUM minimizes the '
     'worst-case detection delay. This is a proven mathematical result, not an empirical claim. No other method can detect faster '
     'without increasing false alarms.',
     10, NAVY, True)


# ================================================================
# APPENDIX C: MITRE ATT&CK TECHNIQUE REFERENCE
# ================================================================
s = _content_slide('APPENDIX B  --  MITRE ATT&CK REFERENCE',
                   'MITRE ATT&CK Technique IDs Used in This Deck',
                   'Industry-standard adversary behavior taxonomy maintained by MITRE Corporation (attack.mitre.org).', NAVY)

mitre_rows = [
    ('T1003', 'Credential Dumping', 'Credential Access', 'Stealing passwords, hashes, tickets from OS. Includes LSASS, SAM, NTDS.dit, DCSync, Kerberoasting.'),
    ('T1021', 'Remote Services', 'Lateral Movement', 'Moving between systems via RDP, SMB/Windows Admin Shares, SSH, WinRM. Key Volt Typhoon technique.'),
    ('T1041', 'Exfiltration Over C2', 'Exfiltration', 'Stealing data through the existing command-and-control channel rather than a separate connection.'),
    ('T1059', 'Command & Scripting', 'Execution', 'Running commands via PowerShell, cmd, wmic, VBScript. Core Living-off-the-Land technique.'),
    ('T1071', 'Application Layer Protocol', 'Command & Control', 'C2 communication over HTTP/S, DNS, or other standard protocols to blend with normal traffic.'),
    ('T1074', 'Data Staged', 'Collection', 'Collecting and staging data in a central location before exfiltration. Key insider threat indicator.'),
    ('T1078', 'Valid Accounts', 'Defense Evasion', 'Using stolen legitimate credentials. Passes every config-based check by design.'),
    ('T1562', 'Impair Defenses', 'Defense Evasion', 'Disabling or modifying security tools: AV exclusions, log clearing, agent tampering, firewall rules.'),
]

_rect(s, Inches(0.5), Inches(1.4), Inches(1.2), Inches(0.35), fill=NAVY)
_txt(s, Inches(0.55), Inches(1.42), Inches(1.1), Inches(0.3), 'ID', 10, WHITE, True, PP_ALIGN.CENTER)
_rect(s, Inches(1.7), Inches(1.4), Inches(2.3), Inches(0.35), fill=NAVY)
_txt(s, Inches(1.75), Inches(1.42), Inches(2.2), Inches(0.3), 'TECHNIQUE', 10, WHITE, True)
_rect(s, Inches(4.0), Inches(1.4), Inches(1.8), Inches(0.35), fill=NAVY)
_txt(s, Inches(4.05), Inches(1.42), Inches(1.7), Inches(0.3), 'TACTIC', 10, WHITE, True)
_rect(s, Inches(5.8), Inches(1.4), Inches(6.8), Inches(0.35), fill=NAVY)
_txt(s, Inches(5.85), Inches(1.42), Inches(6.7), Inches(0.3), 'DESCRIPTION & RELEVANCE', 10, WHITE, True)

y = Inches(1.8)
for tid, technique, tactic, desc in mitre_rows:
    bg = CARD if mitre_rows.index((tid, technique, tactic, desc)) % 2 == 0 else WHITE
    _box(s, Inches(0.5), y, Inches(12.1), Inches(0.6), fill=bg)
    _rect(s, Inches(0.5), y+Inches(0.1), Inches(1.1), Inches(0.3), fill=TEAL)
    _txt(s, Inches(0.55), y+Inches(0.12), Inches(1.0), Inches(0.26), tid, 10, WHITE, True, PP_ALIGN.CENTER, fn='Cascadia Code')
    _txt(s, Inches(1.75), y+Pt(2), Inches(2.2), Inches(0.55), technique, 10, NAVY, True)
    _txt(s, Inches(4.05), y+Pt(2), Inches(1.7), Inches(0.55), tactic, 9, SUB)
    _txt(s, Inches(5.85), y+Pt(2), Inches(6.7), Inches(0.55), desc, 9, BODY)
    y += Inches(0.63)

_box(s, Inches(0.5), Inches(6.9), Inches(12.1), Inches(0.4), fill=CARD, accent=GOLD)
_txt(s, Inches(0.8), Inches(6.92), Inches(11.6), Inches(0.35),
     'ACECARD drift direction projects behavioral drift onto these exact technique embeddings. '
     '"Drifting toward T1021 at 0.82" = actionable, auditable, MITRE-mapped intelligence.',
     10, NAVY, True)


# ================================================================
# APPENDIX D: ADDITIONAL MITRE TECHNIQUES (Volt/Salt Typhoon Specific)
# ================================================================
s = _content_slide('APPENDIX B (cont.)  --  VOLT & SALT TYPHOON ATT&CK MAPPING',
                   'Complete MITRE ATT&CK Mapping: Volt Typhoon (G1017) & Salt Typhoon',
                   'All techniques observed in CISA advisories AA24-038A and AA25-239A.', NAVY)

vt_techniques = [
    ('T1190', 'Exploit Public-Facing App', 'Fortinet, Ivanti, Versa, PAN-OS CVEs', 'VT'),
    ('T1133', 'External Remote Services', 'VPN/SASE exploitation for initial access', 'VT'),
    ('T1078', 'Valid Accounts', 'Stolen credentials from ntdsutil / LSASS', 'VT+ST'),
    ('T1059.001', 'PowerShell', 'Encoded commands, download cradles', 'VT'),
    ('T1047', 'WMI (wmic)', 'Remote execution via WMI for LoTL', 'VT'),
    ('T1003.003', 'NTDS (ntdsutil)', 'Domain controller NTDS.dit extraction', 'VT'),
    ('T1003.001', 'LSASS Memory', 'Credential harvesting from LSASS process', 'VT'),
    ('T1021.002', 'SMB/Admin Shares', 'Lateral via C$ and ADMIN$ shares', 'VT'),
    ('T1021.001', 'Remote Desktop', 'RDP lateral movement with stolen creds', 'VT+ST'),
    ('T1071.001', 'Web Protocols', 'HTTPS C2, KV Botnet communication', 'VT'),
    ('T1090', 'Proxy (KV Botnet)', 'Compromised SOHO routers as proxy infra', 'VT'),
    ('T1136.001', 'Create Local Account', 'Rogue admin accounts on Cisco IOS', 'ST'),
    ('T1556', 'Modify Auth Process', 'TACACS+ credential capture (TCP/49)', 'ST'),
    ('T1572', 'Protocol Tunneling', 'GRE tunnels between compromised routers', 'ST'),
    ('T1041', 'Exfil Over C2', 'CDR/metadata exfiltration, CALEA access', 'ST'),
    ('T1562.001', 'Disable Security Tools', 'Log clearing, ACL modification for evasion', 'VT+ST'),
]

_rect(s, Inches(0.5), Inches(1.4), Inches(1.4), Inches(0.3), fill=NAVY)
_txt(s, Inches(0.55), Inches(1.41), Inches(1.3), Inches(0.28), 'TECHNIQUE ID', 9, WHITE, True, PP_ALIGN.CENTER)
_rect(s, Inches(1.9), Inches(1.4), Inches(2.5), Inches(0.3), fill=NAVY)
_txt(s, Inches(1.95), Inches(1.41), Inches(2.4), Inches(0.28), 'NAME', 9, WHITE, True)
_rect(s, Inches(4.4), Inches(1.4), Inches(5.8), Inches(0.3), fill=NAVY)
_txt(s, Inches(4.45), Inches(1.41), Inches(5.7), Inches(0.28), 'CONTEXT', 9, WHITE, True)
_rect(s, Inches(10.2), Inches(1.4), Inches(2.4), Inches(0.3), fill=NAVY)
_txt(s, Inches(10.25), Inches(1.41), Inches(2.3), Inches(0.28), 'ACTOR', 9, WHITE, True, PP_ALIGN.CENTER)

y = Inches(1.73)
for tid, name, context, actor in vt_techniques:
    bg = CARD if vt_techniques.index((tid, name, context, actor)) % 2 == 0 else WHITE
    _box(s, Inches(0.5), y, Inches(12.1), Inches(0.32), fill=bg)
    actor_color = CRED if 'VT' in actor and 'ST' in actor else (ORG if 'VT' in actor else TEAL)
    _txt(s, Inches(0.6), y+Pt(1), Inches(1.2), Inches(0.28), tid, 8.5, NAVY, True, fn='Cascadia Code')
    _txt(s, Inches(1.95), y+Pt(1), Inches(2.4), Inches(0.28), name, 8.5, BODY)
    _txt(s, Inches(4.45), y+Pt(1), Inches(5.7), Inches(0.28), context, 8.5, SUB)
    _rect(s, Inches(10.3), y+Inches(0.03), Inches(1.0), Inches(0.22), fill=actor_color)
    _txt(s, Inches(10.3), y+Inches(0.04), Inches(1.0), Inches(0.2), actor, 8, WHITE, True, PP_ALIGN.CENTER)
    y += Inches(0.335)

_txt(s, Inches(0.5), Inches(7.15), Inches(12.3), Inches(0.2),
     'Sources: MITRE ATT&CK v15; CISA AA24-038A (Volt Typhoon); CISA AA25-239A (Salt Typhoon); Microsoft Threat Intelligence',
     8, SUB)


# ================================================================
# APPENDIX C: TECHNICAL GLOSSARY
# ================================================================
s = _content_slide('APPENDIX C  --  TECHNICAL GLOSSARY',
                   'Key Technical Terms & Metrics Explained',
                   'Plain-language definitions for every technical term used in this deck.', NAVY)

glossary_1 = [
    ('Cohort z-score', 'How far one entity\'s drift is from its PEER GROUP average. '
     'If all 50 sysadmins drift ~0.03 and one drifts 0.14, z = (0.14 - 0.03) / 0.018 = 6.17. '
     'z > 3 is extreme. This entity behaves radically differently from everyone with the same role.'),
    ('Drift (cosine distance)', 'How much an entity\'s behavioral embedding changed between two time windows. '
     'drift = 1 - cosine_similarity(V_now, V_previous). Range: 0.0 (identical) to 2.0 (opposite). '
     'Healthy: ~0.02-0.05. Compromised: ~0.06-0.15.'),
    ('CUSUM (Cumulative Sum)', 'Sequential change-point detection algorithm (Page, 1954). Accumulates small behavioral drifts over time. '
     'Catches slow APT campaigns that NEVER cross a fixed threshold in any single window. '
     'S(t) = max(0, S(t-1) + drift(t) - mu - k). Alarm when S(t) >= 4*sigma.'),
    ('Drift Direction (concept alignment)', 'Cosine projection of the drift VECTOR onto 8 MITRE threat concept embeddings. '
     'Answers "WHAT is this entity becoming?" not just "something changed." '
     'Example: "drifting toward lateral_movement (T1021) at cosine 0.82."'),
    ('1536-d Embedding', 'A 1,536-dimensional vector that captures the behavioral STATE of an entity. '
     'Generated by encoding behavioral metrics as natural language, then embedding via OpenAI text-embedding-3-small '
     '(or local SLM for IL6/JWICS). All entities live in the same vector space.'),
    ('Entity Fusion', 'Resolving 10+ identity system identifiers (AD, AAD, Okta, AWS, K8s, CrowdStrike, VPN, etc.) '
     'into a single entity_uuid. Without this, one attacker appears as 10 unrelated users.'),
]

y = Inches(1.5)
for term, definition in glossary_1:
    _box(s, Inches(0.5), y, Inches(12.1), Inches(0.88), fill=CARD)
    _rect(s, Inches(0.5), y, Inches(0.08), Inches(0.88), fill=TEAL)
    _txt(s, Inches(0.7), y+Pt(2), Inches(2.5), Inches(0.35), term, 11, NAVY, True)
    _txt(s, Inches(3.3), y+Pt(2), Inches(9.1), Inches(0.83), definition, 9.5, BODY)
    y += Inches(0.92)


# Glossary page 2 - Signal metrics
s = _content_slide('APPENDIX C (cont.)  --  SIGNAL METRICS GLOSSARY',
                   'Behavioral Signal Metrics: What Each Number Means',
                   'The raw metrics inside each of the 5 behavioral signals.', NAVY)

signal_glossary = [
    ('PROCESS SIGNAL METRICS', TEAL, [
        ('lolbin_count', 'Count of Living-off-the-Land Binaries executed (wmic, certutil, ntdsutil, rundll32, mshta, etc.). '
         'These are legitimate Windows tools weaponized by attackers. 0 is normal for most users.'),
        ('cmdline_entropy', 'Shannon entropy of command-line arguments. Low (1.0-2.0) = readable commands like "dir C:\\Users". '
         'High (3.0+) = Base64-encoded or obfuscated payloads. Volt Typhoon uses encoded PowerShell.'),
        ('parent_child_depth', 'Depth of the process tree. Normal: explorer > cmd > dir (depth 2). '
         'Attack: winword > cmd > powershell > wmic > rundll32 > cmd (depth 5+). Deep trees = chained execution.'),
        ('scripting_engine_calls', 'Count of PowerShell, VBScript, WScript, cscript invocations per window. '
         '0 is normal for non-admin users. 14 in one window = automated attack script running.'),
        ('unsigned_ratio', 'Fraction of executed binaries without valid code signatures. High ratio = custom malware or tampered tools.'),
    ]),
    ('AUTH SIGNAL METRICS', RGBColor(0x1E,0x4D,0x8C), [
        ('failure_rate', 'Ratio of failed to total logon attempts. Sudden spike = brute force or credential stuffing.'),
        ('off_hours_ratio', 'Fraction of logons outside normal working hours. Spike = compromised account used by attacker in different timezone.'),
        ('impossible_travel', 'Logons from locations that are geographically impossible given time between events. '
         'NYC at 2:00 PM, Beijing at 2:15 PM = compromised credential.'),
        ('mfa_skip_ratio', 'Fraction of sessions bypassing MFA. Token theft or legacy auth exploitation.'),
    ]),
    ('NETWORK SIGNAL METRICS', sig_colors[2], [
        ('beacon_score', 'Regularity of outbound connections. C2 beacons call home at fixed intervals. '
         'Score 0.0=random, 1.0=perfectly periodic. KV Botnet and JumbledPath have scores > 0.7.'),
        ('bytes_out_ratio', 'Outbound vs inbound data ratio. Normal users download more than upload. '
         'Ratio inversion (upload > download) = data exfiltration.'),
        ('dns_query_rate', 'DNS queries per hour. Spike = DNS tunneling for C2 or exfiltration.'),
    ]),
]

y = Inches(1.4)
for section, color, metrics in signal_glossary:
    _rect(s, Inches(0.5), y, Inches(12.1), Inches(0.3), fill=color)
    _txt(s, Inches(0.7), y+Pt(1), Inches(11.7), Inches(0.26), section, 10, WHITE, True)
    y += Inches(0.33)
    for metric, desc in metrics:
        _box(s, Inches(0.5), y, Inches(12.1), Inches(0.5), fill=CARD)
        _txt(s, Inches(0.7), y+Pt(1), Inches(2.2), Inches(0.45), metric, 9.5, NAVY, True, fn='Cascadia Code')
        _txt(s, Inches(3.0), y+Pt(1), Inches(9.4), Inches(0.45), desc, 9, BODY)
        y += Inches(0.5)
    y += Inches(0.08)


# ================================================================
import re as _re
def _safe_save(path, prs):
    try:
        prs.save(path)
        return path
    except PermissionError:
        base = _re.sub(r'(_v\d+)?\.pptx$', '', path)
        for n in range(2, 10):
            alt = f"{base}_v{n}.pptx"
            try:
                prs.save(alt)
                return alt
            except PermissionError:
                continue
        raise

saved = _safe_save(OUT, prs)
print(f"Saved: {saved}")
print(f"Slides: {len(prs.slides)}")
