# -*- coding: utf-8 -*-
"""UEBA deck - PART 1: problem-staging slides (why we need an advanced solution).
DLA-deck look and feel (Segoe UI, 22nd Century orange + navy + teal, thin top rule)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY=RGBColor(0x10,0x2a,0x4c); DARKNAVY=RGBColor(0x0c,0x20,0x3b); TOPRULE=RGBColor(0x06,0x24,0x4a)
BLUE=RGBColor(0x1f,0x5e,0x9c); TEAL=RGBColor(0x11,0x84,0x84); SLATE=RGBColor(0x5e,0x6a,0x75)
WHITE=RGBColor(0xff,0xff,0xff); LIGHT=RGBColor(0xec,0xf4,0xfc); GREEN=RGBColor(0x1a,0x72,0x4e)
RED=RGBColor(0xa6,0x1b,0x1b); ORANGE=RGBColor(0xe8,0x77,0x22); INK=RGBColor(0x1d,0x24,0x30)
LIGHTO=RGBColor(0xff,0xf5,0xeb); LIGHTT=RGBColor(0xea,0xf7,0xf6); GRAYL=RGBColor(0xd8,0xe0,0xec)
FT="Segoe UI"

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]; SW,SH=prs.slide_width,prs.slide_height

def slide(): return prs.slides.add_slide(BLANK)
def rect(s,x,y,w,h,fill,line=None):
    sp=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1)
    sp.shadow.inherit=False; return sp
def txt(s,x,y,w,h,text,size=18,color=NAVY,bold=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,italic=False,font=FT):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,ln in enumerate(text.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        r=p.add_run(); r.text=ln; r.font.size=Pt(size); r.font.bold=bold; r.font.italic=italic; r.font.color.rgb=color; r.font.name=font
    return tb
def bullets(s,x,y,w,h,items,size=15,color=INK,gap=8):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
    for i,(lvl,t,*c) in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.level=lvl; p.space_after=Pt(gap)
        r=p.add_run(); r.text=("•  " if lvl==0 else "–  ")+t
        r.font.size=Pt(size-(1 if lvl else 0)); r.font.color.rgb=(c[0] if c else color); r.font.name=FT
        if c and len(c)>1 and c[1]: r.font.bold=True
    return tb
def dtable(s,x,y,w,h,headers,rows,fs=11,widths=None):
    t=s.shapes.add_table(len(rows)+1,len(headers),x,y,w,h).table
    for j,hd in enumerate(headers):
        c=t.cell(0,j); c.text=""; r=c.text_frame.paragraphs[0].add_run(); r.text=hd; r.font.size=Pt(fs); r.font.bold=True; r.font.name=FT
    for i,row in enumerate(rows,1):
        for j,v in enumerate(row):
            c=t.cell(i,j); c.text=""; r=c.text_frame.paragraphs[0].add_run(); r.text=str(v); r.font.size=Pt(fs); r.font.name=FT
    if widths:
        for col,wd in enumerate(widths):
            for rr in t.rows: rr.cells[col].width=Inches(wd)
    return t
def header(s,title,kicker=None):
    rect(s,0,0,SW,Pt(7),TOPRULE)
    txt(s,Inches(0.5),Inches(0.28),Inches(12.3),Inches(0.7),title,size=23,color=NAVY,bold=True)
    if kicker: txt(s,Inches(0.5),Inches(0.98),Inches(12.3),Inches(0.35),kicker,size=12,color=SLATE)
    rect(s,Inches(0.5),Inches(1.4),Inches(2.6),Pt(3),ORANGE)
n=[0]
def fin(s):
    n[0]+=1
    rect(s,0,Inches(7.33),SW,Pt(1.2),GRAYL)
    txt(s,Inches(0.5),Inches(7.06),Inches(11),Inches(0.3),"CONFIDENTIAL  |  22nd Century Technologies  |  V-Intelligence Behavioral UEBA",size=9,color=SLATE)
    txt(s,Inches(12.4),Inches(7.06),Inches(0.7),Inches(0.3),str(n[0]),size=9,color=SLATE,align=PP_ALIGN.RIGHT)

# --- generate the real CUSUM/drift chart (LOTL operator USR-042) ---
import matplotlib; matplotlib.use("Agg")
import matplotlib.pyplot as plt
import pandas as pd
FIGP="docs/wp_figs"; os.makedirs(FIGP,exist_ok=True); CHART=f"{FIGP}/ueba_cusum.png"; HAVE_CHART=False
try:
    tr=pd.read_csv("data/tier3_results/weekly_zone_trajectories.csv")
    v=tr[tr.user_id=="USR-042"].sort_values("week_idx")
    wk=v.week_idx.values; ap=v.access_pattern_drift.values; cd=v.composite_drift.values
    cc=[]; c=0.0
    for m in cd: c=max(0.0,c+m-0.02); cc.append(c)
    fig,ax=plt.subplots(figsize=(9.4,3.5))
    ax.plot(wk,ap,color="#118484",lw=1.7,label="access_pattern drift")
    ax.plot(wk,cd,color="#1f5e9c",lw=1.7,label="composite drift")
    ax.plot(wk,cc,color="#e87722",lw=2.1,ls="--",label="CUSUM (accumulated)")
    ax.axhline(0.05,color="#5e6a75",lw=0.9,ls=":",label="CUSUM fire line (0.05)")
    ax.set_xlabel("Week",fontsize=9); ax.set_ylabel("Drift / CUSUM",fontsize=9)
    ax.set_title("LOTL operator (USR-042): slow, sub-threshold drift accumulates over time",fontsize=10,color="#102a4c",fontweight="bold")
    ax.legend(fontsize=8,loc="upper left"); ax.grid(alpha=0.22); ax.tick_params(labelsize=8)
    for sp in ["top","right"]: ax.spines[sp].set_visible(False)
    plt.tight_layout(); plt.savefig(CHART,dpi=150,bbox_inches="tight"); plt.close(); HAVE_CHART=True
    print("chart generated:",CHART)
except Exception as e:
    print("chart gen skipped:",e)

# ---- 1 TITLE ----
s=slide(); rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.22),SH,ORANGE); rect(s,Inches(8.5),0,SW-Inches(8.5),SH,DARKNAVY)
txt(s,Inches(0.7),Inches(1.5),Inches(7.4),Inches(1.8),"Behavioral UEBA for the Age of AI-Enabled Threats",size=33,color=NAVY,bold=True)
txt(s,Inches(0.7),Inches(3.35),Inches(7.4),Inches(0.6),"Why detect-and-respond no longer defends critical infrastructure",size=18,color=BLUE)
rect(s,Inches(0.7),Inches(4.2),Inches(2.4),Pt(3),ORANGE)
txt(s,Inches(0.7),Inches(5.7),Inches(7.4),Inches(0.5),"V-Intelligence  |  22nd Century Technologies",size=13,color=NAVY,italic=True)
txt(s,Inches(0.7),Inches(6.2),Inches(7.4),Inches(0.4),"Part 1: The Problem  -  why an advanced solution is needed",size=11,color=SLATE)
for i,(lbl,col) in enumerate([("Valid-account attacks",BLUE),("AI-speed adversaries",TEAL),("Alert-fatigued SOCs",ORANGE)]):
    rect(s,Inches(9.0),Inches(2.2)+Inches(1.05)*i,Inches(3.3),Inches(0.8),col)
    txt(s,Inches(9.2),Inches(2.2)+Inches(1.05)*i,Inches(2.9),Inches(0.8),lbl,size=14,color=WHITE,bold=True,anchor=MSO_ANCHOR.MIDDLE)
fin(s)

# ---- 2 HOW THE SOC WORKS TODAY ----
s=slide(); header(s,"The SOC Runs on SIEM and UEBA Today","Collect, correlate, alert, triage, respond")
steps=[("Collect","logs from many sources"),("Normalize","parse and enrich"),("Correlate","rules + statistical baselines"),("Alert","threshold or rule fires"),("Triage","analyst investigates"),("Respond","contain and remediate")]
x=Inches(0.5)
for i,(nm,desc) in enumerate(steps):
    rect(s,x,Inches(2.0),Inches(1.92),Inches(1.5),LIGHT,line=BLUE)
    txt(s,x+Inches(0.08),Inches(2.15),Inches(1.76),Inches(0.6),nm,size=13,color=NAVY,bold=True,align=PP_ALIGN.CENTER)
    txt(s,x+Inches(0.08),Inches(2.75),Inches(1.76),Inches(0.7),desc,size=10,color=INK,align=PP_ALIGN.CENTER)
    x+=Inches(2.05)
bullets(s,Inches(0.6),Inches(4.0),Inches(12),Inches(2.6),[
 (0,"The model assumes the attack produces a signal that crosses a rule or a threshold.",NAVY,True),
 (0,"It assumes a known signature, or a deviation large enough to stand out from a baseline.",),
 (0,"And it assumes a human analyst has the time to triage what fires.",),
],size=15,gap=10)
fin(s)

# ---- 3 WHERE IT BREAKS ----
s=slide(); header(s,"How Detection Works Today, and Where It Breaks")
rect(s,Inches(0.6),Inches(1.7),Inches(5.9),Inches(4.7),LIGHT); rect(s,Inches(0.6),Inches(1.7),Inches(5.9),Inches(0.7),BLUE)
txt(s,Inches(0.8),Inches(1.78),Inches(5.5),Inches(0.6),"How it works",size=16,color=WHITE,bold=True,anchor=MSO_ANCHOR.MIDDLE)
bullets(s,Inches(0.85),Inches(2.6),Inches(5.4),Inches(3.6),[
 (0,"Signature and correlation rules match known-bad patterns.",),
 (0,"UEBA adds statistical baselines (peer and self) and risk scores.",),
 (0,"Alerts are ranked by severity and pushed to the SOC queue.",),
],size=14,gap=10)
rect(s,Inches(6.8),Inches(1.7),Inches(5.9),Inches(4.7),LIGHTO); rect(s,Inches(6.8),Inches(1.7),Inches(5.9),Inches(0.7),ORANGE)
txt(s,Inches(7.0),Inches(1.78),Inches(5.5),Inches(0.6),"Where it breaks",size=16,color=WHITE,bold=True,anchor=MSO_ANCHOR.MIDDLE)
bullets(s,Inches(7.05),Inches(2.6),Inches(5.4),Inches(3.6),[
 (0,"Magnitude, not meaning: 47 files is 47 files, intent is lost.",RED,True),
 (0,"Thresholds miss slow, valid-account, living-off-the-land attacks that never spike.",RED,True),
 (0,"Per-source silos make one attacker look like ten users.",),
 (0,"Alert fatigue: analysts drown; real signal is buried in noise.",),
 (0,"Reactive: it detects after the fact, not before.",),
],size=14,gap=9)
fin(s)

# ---- 4 SIEM LANDSCAPE ----
s=slide(); header(s,"The SIEM / UEBA Landscape","Mature tools, the same structural limits")
dtable(s,Inches(0.5),Inches(1.7),Inches(12.3),Inches(3.6),
 ["Tool","Category","What it does"],
 [["Splunk Enterprise Security","SIEM","log search, correlation rules, dashboards"],
  ["IBM QRadar","SIEM","flow + log correlation, offense scoring"],
  ["Microsoft Sentinel","Cloud SIEM + SOAR","KQL analytics, automation, some UEBA"],
  ["Exabeam / Securonix","UEBA-led","behavioral baselining, risk scoring"],
  ["Elastic Security / Google SecOps","SIEM","search, detections, pipelines at scale"],
  ["LogRhythm / Devo / Sumo Logic / Rapid7","SIEM","log analytics and detections"]],fs=12,widths=[3.6,2.6,6.1])
txt(s,Inches(0.6),Inches(5.6),Inches(12),Inches(1.0),"All are excellent at collecting and searching logs and at matching known patterns. All share the same limit: they alert on magnitude and signatures, not on the meaning and direction of behavior, so the valid-account, slow, AI-speed adversary walks past.",size=14,color=NAVY,bold=True)
fin(s)

# ---- 5 THE INFLECTION: AI-ENABLED ATTACKERS ----
s=slide(); header(s,"The Inflection: AI-Enabled (Mythos-Level) Attackers","Adversaries now wield AI that finds and exploits vulnerabilities at machine speed")
bullets(s,Inches(0.6),Inches(1.7),Inches(12),Inches(5),[
 (0,"Offensive AI (Mythos-level) autonomously discovers vulnerabilities, chains low-severity weaknesses, and exploits them, faster than humans can respond.",NAVY,True),
 (0,"In 2025, an AI-orchestrated espionage campaign (Anthropic GTG-1002) ran largely autonomously against roughly 30 targets at machine speed.",),
 (0,"Threats now evolve faster than vulnerability disclosure and patch cycles; attackers exploit before a CVE or patch exists.",),
 (0,"By the time a detect-and-respond pipeline fires, the breach is already complete.",RED,True),
 (0,"Detect-and-respond, the model every SIEM is built on, cannot defend U.S. infrastructure against this.",RED,True),
],size=16,gap=13)
fin(s)

# ---- 6 OPERATIONAL DISRUPTION ----
s=slide(); header(s,"Mythos-Level Attacks Mean Operational Disruption","Not data loss alone: mission and infrastructure impact")
cards=[("Energy / Water","pre-positioned access enables disruption of power, water, and fuel on demand",ORANGE),
       ("Telecom","carrier-scale espionage and the ability to degrade communications",TEAL),
       ("Transport / Logistics","route, depot, and fulfillment disruption that stalls mobilization",BLUE),
       ("Defense / Mission","loss of decision time and readiness when it is needed most",NAVY)]
x=Inches(0.5)
for nm,desc,col in cards:
    rect(s,x,Inches(1.8),Inches(3.0),Inches(3.2),LIGHT); rect(s,x,Inches(1.8),Inches(3.0),Inches(0.8),col)
    txt(s,x+Inches(0.15),Inches(1.9),Inches(2.7),Inches(0.6),nm,size=14,color=WHITE,bold=True,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,x+Inches(0.15),Inches(2.8),Inches(2.7),Inches(2.0),desc,size=12.5,color=INK)
    x+=Inches(3.07)
txt(s,Inches(0.6),Inches(5.4),Inches(12),Inches(1.0),"Adversaries pre-position quietly, on valid credentials, to disrupt critical infrastructure at the moment of their choosing. The cost of a missed warning is operational, not just informational.",size=14,color=NAVY,bold=True)
fin(s)

# ---- 7 RECENT ATTACKS TOOLS COULD NOT STOP ----
s=slide(); header(s,"Recent Attacks Current Tools Could Not Stop")
dtable(s,Inches(0.5),Inches(1.6),Inches(12.3),Inches(4.4),
 ["Campaign","What it was","Why current tools missed it"],
 [["Volt Typhoon (2023-24)","living-off-the-land pre-positioning in U.S. critical infrastructure, 5+ year dwell","valid accounts, no malware, no signature, sub-threshold"],
  ["Salt Typhoon (2024-25)","telecom espionage at carrier scale","slow, valid-account, below every alert"],
  ["SolarWinds (2020)","trojanized signed update pushed to ~18,000 organizations","trusted, signed update, no anomaly to flag"],
  ["MOVEit (2023)","mass zero-day exploitation","no CVE or signature at time of exploit"],
  ["XZ Utils (2024)","multi-year supply-chain backdoor in a core library","trusted dependency, no known-bad version"],
  ["Colonial Pipeline (2021)","ransomware, fuel-supply disruption","credential access, detected too late to prevent impact"]],fs=11,widths=[2.5,4.9,4.9])
txt(s,Inches(0.6),Inches(6.2),Inches(12),Inches(0.6),"The common thread: valid credentials, no signature, no spike, or a trusted path. Magnitude-and-signature detection had nothing to fire on.",size=13,color=RED,bold=True)
fin(s)

# ---- 8 WHY WE NEED AN ADVANCED SOLUTION ----
s=slide(); rect(s,0,0,SW,SH,DARKNAVY); rect(s,0,Inches(2.5),SW,Pt(3),ORANGE)
txt(s,Inches(0.8),Inches(1.2),Inches(11.5),Inches(0.5),"WHY WE NEED AN ADVANCED SOLUTION",size=17,color=ORANGE,bold=True)
bullets(s,Inches(0.8),Inches(2.8),Inches(11.7),Inches(3.5),[
 (0,"Current SIEM and UEBA detect magnitude and signatures; they miss meaning and direction.",WHITE),
 (0,"They cannot catch valid-account, living-off-the-land attacks that cross no threshold.",WHITE),
 (0,"They are reactive, and AI-speed adversaries are gone before detect-and-respond fires.",WHITE),
 (0,"And they overwhelm the SOC with alerts instead of giving analysts a ranked, explainable verdict.",WHITE),
],size=17,gap=14)
txt(s,Inches(0.8),Inches(6.2),Inches(11.7),Inches(0.6),"We need behavior read as meaning, detection that is preemptive and precise, and a SOC that is helped, not buried.",size=15,color=ORANGE,bold=True)
fin(s)

# =====================================================================
# PART 2 -- THE SOLUTION (behavioral digital twin)
# =====================================================================

# ---- 9 SECTION DIVIDER ----
s=slide(); rect(s,0,0,SW,SH,DARKNAVY); rect(s,0,0,Inches(0.22),SH,ORANGE)
txt(s,Inches(0.8),Inches(2.5),Inches(11),Inches(0.5),"PART 2",size=15,color=ORANGE,bold=True)
txt(s,Inches(0.8),Inches(3.05),Inches(11.5),Inches(1.3),"The Solution: A Behavioral Digital Twin",size=34,color=WHITE,bold=True)
txt(s,Inches(0.8),Inches(4.5),Inches(11.5),Inches(0.7),"Read behavior as meaning, track where each entity is heading, hand the SOC a ranked, explainable verdict.",size=17,color=RGBColor(0xbf,0xd4,0xea))
fin(s)

# ---- 10 INPUT DATA / COLLECTION ----
s=slide(); header(s,"Input Data: the Logs the SOC Already Collects","No new sensors. The twin sits downstream of your existing SIEM.")
dtable(s,Inches(0.5),Inches(1.65),Inches(7.4),Inches(3.1),
 ["Log source","Signal it carries"],
 [["Identity / auth","who, when, from where, privilege use"],
  ["File / data","access, volume, sensitivity, movement"],
  ["Network / flow","destinations, beaconing, data egress"],
  ["DNS","resolution patterns, DGA, rare domains"],
  ["Endpoint / process","process lineage, living-off-the-land tools"]],fs=12,widths=[2.4,5.0])
rect(s,Inches(8.2),Inches(1.65),Inches(4.6),Inches(3.1),LIGHT); rect(s,Inches(8.2),Inches(1.65),Inches(4.6),Inches(0.6),TEAL)
txt(s,Inches(8.4),Inches(1.71),Inches(4.2),Inches(0.5),"Collection stages",size=14,color=WHITE,bold=True,anchor=MSO_ANCHOR.MIDDLE)
bullets(s,Inches(8.4),Inches(2.4),Inches(4.2),Inches(2.3),[
 (0,"Ingest from the SIEM / forwarders",),
 (0,"Normalize and parse to a common schema",),
 (0,"Group events into behavioral dimensions",),
 (0,"Serialize and embed; store as vectors",),
],size=12.5,gap=8)
txt(s,Inches(0.6),Inches(5.0),Inches(12.2),Inches(1.3),"Scale of the proof-of-concept: 250 entities, 485 days, about 14 million events across five log sources. Behavioral snapshots are stored as 1536-dimension vectors in a PostgreSQL pgvector database, so similarity and drift are queried directly in the database.",size=13.5,color=NAVY,bold=True)
fin(s)

# ---- 11 DIMENSIONS AND WHY ----
s=slide(); header(s,"The Dimensions We Model, and Why","Five behavioral zones, 23 features; identity as context")
zones=[("Identity","role, department, clearance, user type","context",SLATE),
       ("Access pattern","logins, sources, destinations, off-hours, methods (7 features)","action",BLUE),
       ("Data behavior","file volume, sensitivity, breadth, writes (6 features)","action",TEAL),
       ("Network footprint","egress bytes, destinations, DNS patterns (5 features)","action",GREEN),
       ("Risk posture","endpoint and process risk scores (5 features)","action",ORANGE)]
y=Inches(1.75)
for nm,desc,kind,col in zones:
    rect(s,Inches(0.5),y,Inches(0.18),Inches(0.78),col)
    rect(s,Inches(0.68),y,Inches(11.6),Inches(0.78),LIGHT)
    txt(s,Inches(0.85),y+Inches(0.06),Inches(3.0),Inches(0.66),nm,size=15,color=NAVY,bold=True,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,Inches(3.9),y+Inches(0.06),Inches(6.6),Inches(0.66),desc,size=13,color=INK,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,Inches(10.7),y+Inches(0.06),Inches(1.5),Inches(0.66),kind,size=11,color=col,bold=True,anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.RIGHT)
    y+=Inches(0.92)
txt(s,Inches(0.6),Inches(6.5),Inches(12.2),Inches(0.7),"Why these: an attacker on valid credentials betrays themselves not in one number but across zones at once. Identity sets the peer group; the four action zones carry the behavior we score against it.",size=13,color=NAVY,bold=True)
fin(s)

# ---- 12 AGGREGATION ----
s=slide(); header(s,"Aggregation: from Activity to a Comparable Window","Windowed so behavior can be compared over time and across peers")
rect(s,Inches(0.6),Inches(1.8),Inches(5.85),Inches(3.0),LIGHTT); rect(s,Inches(0.6),Inches(1.8),Inches(5.85),Inches(0.65),TEAL)
txt(s,Inches(0.8),Inches(1.86),Inches(5.5),Inches(0.55),"Today (proof-of-concept)",size=15,color=WHITE,bold=True,anchor=MSO_ANCHOR.MIDDLE)
bullets(s,Inches(0.85),Inches(2.65),Inches(5.4),Inches(2.0),[
 (0,"Weekly aggregation across 485 days",),
 (0,"Long horizon proves slow-drift detection",),
 (0,"One trajectory point per entity per week",),
],size=13.5,gap=9)
rect(s,Inches(6.85),Inches(1.8),Inches(5.95),Inches(3.0),LIGHTO); rect(s,Inches(6.85),Inches(1.8),Inches(5.95),Inches(0.65),ORANGE)
txt(s,Inches(7.05),Inches(1.86),Inches(5.5),Inches(0.55),"Design target",size=15,color=WHITE,bold=True,anchor=MSO_ANCHOR.MIDDLE)
bullets(s,Inches(7.1),Inches(2.65),Inches(5.5),Inches(2.0),[
 (0,"Hourly aggregation for near-real-time detection",ORANGE,True),
 (0,"Same pipeline, finer window: faster warning",),
 (0,"Shrinks dwell time from weeks toward hours",),
],size=13.5,gap=9)
txt(s,Inches(0.6),Inches(5.1),Inches(12.2),Inches(1.3),"The method is window-agnostic. The current build aggregates weekly because the study spans 485 days and targets slow living-off-the-land drift; the same serialize-embed-compose-score pipeline runs at an hourly cadence to move detection closer to real time.",size=13.5,color=NAVY,bold=True)
fin(s)

# ---- 13 THE VISUAL FLOW ----
s=slide(); header(s,"End-to-End Flow: Activity to Verdict")
flow=[("Activities","logins, files,\nflows, DNS",SLATE),("Input data","5 SIEM\nlog sources",BLUE),
      ("Dimensions","grouped into\n5 zones",TEAL),("Entity profile","23 features\nper window",GREEN),
      ("Digital twin","prose -> vector,\nattention-fused",NAVY),("Algorithms","CUSUM + two-\nlayer composite",ORANGE),
      ("Results","ranked,\nexplainable",RGBColor(0xc9,0x6a,0x1e))]
x=Inches(0.4); bw=Inches(1.62); gap=Inches(0.135)
for i,(nm,desc,col) in enumerate(flow):
    rect(s,x,Inches(2.6),bw,Inches(1.5),col)
    txt(s,x+Inches(0.05),Inches(2.72),bw-Inches(0.1),Inches(0.55),nm,size=12.5,color=WHITE,bold=True,align=PP_ALIGN.CENTER)
    txt(s,x+Inches(0.05),Inches(3.3),bw-Inches(0.1),Inches(0.7),desc,size=9.5,color=WHITE,align=PP_ALIGN.CENTER)
    if i<len(flow)-1:
        ar=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,x+bw,Inches(3.15),gap,Inches(0.4))
        ar.fill.solid(); ar.fill.fore_color.rgb=SLATE; ar.line.fill.background(); ar.shadow.inherit=False
    x=x+bw+gap
txt(s,Inches(0.6),Inches(4.6),Inches(12.2),Inches(0.6),"One unbroken path: the raw events the SOC already has become one ranked, explainable verdict, with the reason and the MITRE direction attached.",size=14,color=NAVY,bold=True)
bullets(s,Inches(0.6),Inches(5.4),Inches(12),Inches(1.4),[
 (0,"No new collection: it consumes existing logs.",),
 (0,"No labels required: the twin is unsupervised.",),
 (0,"Every score is traceable back to the behavior that caused it.",),
],size=13,gap=7)
fin(s)

# ---- 14 HOW THE TWIN IS BUILT ----
s=slide(); header(s,"How the Digital Twin Is Built","Raw metrics become meaning before they become math")
steps=[("Serialize","each zone's features\nbecome a sentence",TEAL),
       ("Embed","each sentence ->\n1536-d vector",BLUE),
       ("Compose","zones fused by\nattention weight",NAVY),
       ("Trajectory","one twin per\nwindow, over time",GREEN)]
x=Inches(0.5)
for i,(nm,desc,col) in enumerate(steps):
    rect(s,x,Inches(1.9),Inches(2.9),Inches(1.7),col)
    txt(s,x+Inches(0.12),Inches(2.05),Inches(2.66),Inches(0.5),nm,size=15,color=WHITE,bold=True,align=PP_ALIGN.CENTER)
    txt(s,x+Inches(0.12),Inches(2.6),Inches(2.66),Inches(0.9),desc,size=11.5,color=WHITE,align=PP_ALIGN.CENTER)
    if i<3:
        ar=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,x+Inches(2.9),Inches(2.55),Inches(0.16),Inches(0.4))
        ar.fill.solid(); ar.fill.fore_color.rgb=SLATE; ar.line.fill.background(); ar.shadow.inherit=False
    x+=Inches(3.06)
rect(s,Inches(0.6),Inches(4.1),Inches(12.2),Inches(2.3),LIGHT); rect(s,Inches(0.6),Inches(4.1),Inches(0.18),Inches(2.3),ORANGE)
bullets(s,Inches(0.95),Inches(4.3),Inches(11.6),Inches(2.0),[
 (0,"Meaning over magnitude: \"47 files\" becomes \"reaching far more sensitive data than peers, and widening.\"",NAVY,True),
 (0,"Composition, not concatenation: zones are fused by attention weight, so the signal that matters dominates.",),
 (0,"Drift is direction in meaning-space: the cosine distance between successive twins, tracked over time.",),
],size=13.5,gap=10)
fin(s)

# ---- 15 THE ALGORITHMS (TWO LAYERS) ----
s=slide(); header(s,"The Algorithms: Two Complementary Layers","Precision first, discovery alongside")
rect(s,Inches(0.6),Inches(1.7),Inches(5.9),Inches(4.6),LIGHTO); rect(s,Inches(0.6),Inches(1.7),Inches(5.9),Inches(0.75),ORANGE)
txt(s,Inches(0.8),Inches(1.78),Inches(5.5),Inches(0.6),"Layer A -- Precision (production)",size=15,color=WHITE,bold=True,anchor=MSO_ANCHOR.MIDDLE)
bullets(s,Inches(0.85),Inches(2.6),Inches(5.4),Inches(3.5),[
 (0,"Multi-front threat-profile detector.",NAVY,True),
 (0,"Matches entities to named known-bad fingerprints (mass collection, C2 beacon, DGA, recon fan-out, LOTL).",),
 (0,"Scored against role-group peers; fires only on a corroborated match.",),
 (0,"Result: 4 of 4 attackers at 0 false positives.",ORANGE,True),
],size=13,gap=9)
rect(s,Inches(6.8),Inches(1.7),Inches(5.95),Inches(4.6),LIGHTT); rect(s,Inches(6.8),Inches(1.7),Inches(5.95),Inches(0.75),TEAL)
txt(s,Inches(7.0),Inches(1.78),Inches(5.5),Inches(0.6),"Layer B -- Discovery (the twin)",size=15,color=WHITE,bold=True,anchor=MSO_ANCHOR.MIDDLE)
bullets(s,Inches(7.05),Inches(2.6),Inches(5.5),Inches(3.5),[
 (0,"CUSUM accumulates slow drift that no single week would trip.",NAVY,True),
 (0,"Direction projected onto MITRE ATT&CK concepts: names where the entity is heading.",),
 (0,"Composite score ranks every entity, no fingerprint needed.",),
 (0,"Surfaces 4 of 4 by rank where classical methods catch 0 of 4.",TEAL,True),
],size=13,gap=9)
fin(s)

# ---- 16 SLOW DRIFT MADE VISIBLE (real chart) ----
s=slide(); header(s,"Slow Drift, Made Visible","Real trajectory: what CUSUM accumulates that a weekly threshold never trips")
if HAVE_CHART:
    s.shapes.add_picture(CHART,Inches(0.55),Inches(1.75),width=Inches(8.7))
else:
    txt(s,Inches(0.55),Inches(3.0),Inches(8.7),Inches(0.6),"[trajectory chart]",size=14,color=SLATE)
rect(s,Inches(9.45),Inches(1.75),Inches(3.35),Inches(3.3),LIGHT); rect(s,Inches(9.45),Inches(1.75),Inches(0.16),Inches(3.3),ORANGE)
bullets(s,Inches(9.72),Inches(1.95),Inches(3.0),Inches(3.0),[
 (0,"Each week's drift is small, below any alert line.",NAVY,True),
 (0,"CUSUM sums the small, persistent pushes.",),
 (0,"It crosses the fire line and stays, so the slow operator is caught while still inside.",ORANGE,True),
],size=12,gap=10)
txt(s,Inches(0.55),Inches(6.3),Inches(12.2),Inches(0.7),"This is the signal classical point-anomaly methods miss: no single week is anomalous, but the direction is, sustained and accumulating.",size=13.5,color=NAVY,bold=True)
fin(s)

# ---- 17 RESULTS ----
s=slide(); header(s,"Results: the Blind Test","Four nation-state-style attackers hidden among 250 entities")
dtable(s,Inches(0.5),Inches(1.65),Inches(12.3),Inches(2.7),
 ["Method","Catches","Note"],
 [["Precision detector (Layer A)","4 of 4 at 0 false positives","named fingerprint per attacker"],
  ["Discovery twin (Layer B)","4 of 4 by rank (10.6% operating point)","no fingerprint; names MITRE direction"],
  ["Isolation Forest / One-Class SVM / LOF","0 of 4","point-anomaly, no behavior over time"],
  ["Z-score baseline","1 of 4","magnitude only"]],fs=12,widths=[4.4,3.6,4.3])
bullets(s,Inches(0.6),Inches(4.7),Inches(12.2),Inches(1.7),[
 (0,"Insider caught by mass collection + rare destinations; Slow APT by a 386-day C2 beacon + 160 DGA domains.",),
 (0,"Living-off-the-land operator caught by anomalous process use; telecom-pivot by reconnaissance fan-out.",),
],size=13,gap=9)
txt(s,Inches(0.6),Inches(6.4),Inches(12.2),Inches(0.5),"Controlled, synthetic-data proof-of-concept (250 entities, 485 days). Indicative, not a guarantee of field performance.",size=11.5,color=SLATE,italic=True)
fin(s)

# ---- 17 SOC PAYOFF ----
s=slide(); header(s,"What This Gives the SOC Analyst","Less noise, earlier warning, a verdict you can act on")
cards=[("Workload","One ranked queue, not a flood. The analyst sees the few entities that matter, each already explained.",ORANGE),
       ("Quick detection","Slow, valid-account drift is caught while it builds, not after exfiltration. CUSUM warns weeks early.",TEAL),
       ("Action","Every alert carries its reason and a MITRE direction, so triage starts at \"what and where,\" not \"is this real?\"",BLUE)]
x=Inches(0.5)
for nm,desc,col in cards:
    rect(s,x,Inches(1.8),Inches(4.05),Inches(3.4),LIGHT); rect(s,x,Inches(1.8),Inches(4.05),Inches(0.8),col)
    txt(s,x+Inches(0.2),Inches(1.9),Inches(3.7),Inches(0.6),nm,size=16,color=WHITE,bold=True,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,x+Inches(0.2),Inches(2.85),Inches(3.7),Inches(2.2),desc,size=13,color=INK)
    x+=Inches(4.12)
txt(s,Inches(0.6),Inches(5.55),Inches(12.2),Inches(1.2),"The shift: from chasing thousands of magnitude alerts to acting on a short, ranked, explainable list, so the analyst spends time on response, not triage.",size=14.5,color=NAVY,bold=True)
fin(s)

# ---- 18 CLOSING ----
s=slide(); rect(s,0,0,SW,SH,DARKNAVY); rect(s,0,Inches(2.4),SW,Pt(3),ORANGE)
txt(s,Inches(0.8),Inches(1.1),Inches(11.5),Inches(0.5),"THE BEHAVIORAL ANSWER",size=16,color=ORANGE,bold=True)
bullets(s,Inches(0.8),Inches(2.7),Inches(11.7),Inches(3.0),[
 (0,"Reads behavior as meaning, not magnitude, from the logs you already collect.",WHITE),
 (0,"Catches the valid-account, living-off-the-land attacker that crosses no threshold.",WHITE),
 (0,"Precision layer at 0 false positives; discovery twin that needs no prior fingerprint.",WHITE),
 (0,"Hands the analyst a ranked, explainable verdict with its MITRE direction attached.",WHITE),
],size=17,gap=14)
txt(s,Inches(0.8),Inches(6.1),Inches(11.7),Inches(0.7),"From detect-and-respond, too late, to read-the-meaning and warn early.",size=16,color=ORANGE,bold=True)
fin(s)

out="WP DLA/Presentation/V-Intelligence_UEBA_Deck.pptx"
prs.save(out)
print("saved:",out,"| slides:",len(prs.slides._sldIdLst))
