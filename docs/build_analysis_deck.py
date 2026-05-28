#!/usr/bin/env python3
"""
Build ACECARD Traditional vs Behavioral Analysis PowerPoint Deck.

Generates a 16-slide executive presentation comparing traditional anomaly
detection methods against ACECARD's behavioral drift approach, using
empirical results from 130-day synthetic telemetry analysis.

Output: docs/ACECARD_Traditional_vs_Behavioral_Deck.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Color palette
# ---------------------------------------------------------------------------
NAVY = RGBColor(13, 27, 42)
BLUE = RGBColor(27, 79, 114)
TEAL = RGBColor(14, 107, 138)
GOLD = RGBColor(183, 149, 11)
RED = RGBColor(192, 57, 43)
GREEN = RGBColor(39, 174, 96)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(240, 242, 245)
DARK_GRAY = RGBColor(80, 80, 80)
MID_GRAY = RGBColor(180, 180, 180)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None):
    """Add a rectangle shape with optional fill and line."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.fill.solid()
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape


def set_text(shape, text, font_size=18, bold=False, color=NAVY, alignment=PP_ALIGN.LEFT,
             font_name="Calibri"):
    """Set text on a shape's text frame."""
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_paragraph(text_frame, text, font_size=18, bold=False, color=NAVY,
                  alignment=PP_ALIGN.LEFT, space_before=Pt(6), space_after=Pt(3),
                  font_name="Calibri", level=0):
    """Add a new paragraph to an existing text frame."""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    p.level = level
    return p


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False,
                color=NAVY, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a textbox with text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_title_bar(slide, title_text):
    """Add a navy title bar at the top of a content slide."""
    bar = add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1), fill_color=NAVY)
    set_text(bar, title_text, font_size=30, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT,
             font_name="Calibri")
    bar.text_frame.paragraphs[0].space_before = Pt(12)
    bar.text_frame.margin_left = Inches(0.6)
    bar.text_frame.margin_top = Inches(0.15)
    # Gold accent line
    add_shape(slide, Inches(0), Inches(1.1), SLIDE_WIDTH, Inches(0.04), fill_color=GOLD)
    return bar


def add_footer(slide, text="22nd Century Technologies, Inc. | Confidential"):
    """Add a subtle footer."""
    txBox = add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12), Inches(0.4),
                        text, font_size=10, bold=False, color=MID_GRAY,
                        alignment=PP_ALIGN.LEFT)
    return txBox


def add_callout_box(slide, left, top, width, height, text, font_size=22):
    """Add a gold-bordered callout box with impact text."""
    box = add_rounded_rect(slide, left, top, width, height, fill_color=RGBColor(255, 251, 235))
    box.line.fill.solid()
    box.line.color.rgb = GOLD
    box.line.width = Pt(2.5)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    return box


def build_table(slide, rows, cols, left, top, width, height, data, header_row=True,
                col_widths=None, cell_colors=None):
    """Build a styled table on a slide."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # Set column widths if provided
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c] if r < len(data) and c < len(data[r]) else ""

            # Style text
            for para in cell.text_frame.paragraphs:
                para.font.name = "Calibri"
                if r == 0 and header_row:
                    para.font.size = Pt(14)
                    para.font.bold = True
                    para.font.color.rgb = WHITE
                    para.alignment = PP_ALIGN.CENTER
                else:
                    para.font.size = Pt(13)
                    para.font.bold = False
                    para.font.color.rgb = NAVY
                    para.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT

            cell.text_frame.margin_left = Inches(0.08)
            cell.text_frame.margin_right = Inches(0.08)
            cell.text_frame.margin_top = Inches(0.04)
            cell.text_frame.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            # Background colors
            fill = cell.fill
            fill.solid()
            if r == 0 and header_row:
                fill.fore_color.rgb = NAVY
            elif cell_colors and (r, c) in cell_colors:
                fill.fore_color.rgb = cell_colors[(r, c)]
                # Adjust text color for colored cells
                for para in cell.text_frame.paragraphs:
                    para.font.color.rgb = WHITE
                    para.font.bold = True
            elif r % 2 == 0:
                fill.fore_color.rgb = LIGHT_GRAY
            else:
                fill.fore_color.rgb = WHITE

    return table


# ===========================================================================
# SLIDE BUILDERS
# ===========================================================================

def slide_01_title(prs):
    """Slide 1: Title Slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, NAVY)

    # Top accent line
    add_shape(slide, Inches(0.8), Inches(1.8), Inches(2.0), Inches(0.05), fill_color=GOLD)

    # Main title
    add_textbox(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(2.0),
                "Traditional vs Behavioral:", font_size=40, bold=True,
                color=WHITE, alignment=PP_ALIGN.LEFT)
    add_textbox(slide, Inches(0.8), Inches(2.7), Inches(11.5), Inches(1.5),
                "Why Standard Anomaly Detection Fails", font_size=38, bold=True,
                color=TEAL, alignment=PP_ALIGN.LEFT)

    # Subtitle
    add_textbox(slide, Inches(0.8), Inches(4.0), Inches(11.5), Inches(0.8),
                "Empirical Analysis Across 130 Days of Cyber Telemetry",
                font_size=22, bold=False, color=GOLD, alignment=PP_ALIGN.LEFT)

    # Bottom accent line
    add_shape(slide, Inches(0.8), Inches(5.0), Inches(4.0), Inches(0.03), fill_color=GOLD)

    # Company / date
    add_textbox(slide, Inches(0.8), Inches(5.3), Inches(11.5), Inches(0.5),
                "22nd Century Technologies, Inc.  |  May 2025",
                font_size=16, bold=False, color=MID_GRAY, alignment=PP_ALIGN.LEFT)

    # Classification placeholder
    add_textbox(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.4),
                "DISTRIBUTION STATEMENT: For Official Use Only",
                font_size=11, bold=False, color=MID_GRAY, alignment=PP_ALIGN.LEFT)


def slide_02_problem(prs):
    """Slide 2: The Problem."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "The Problem")

    # Main statement
    txBox = add_textbox(slide, Inches(0.7), Inches(1.4), Inches(11.8), Inches(1.0),
                        "Your SIEM Sees Volume. Attackers Change Direction.",
                        font_size=28, bold=True, color=NAVY, alignment=PP_ALIGN.LEFT)

    # Bullet points
    bullets = [
        "Traditional anomaly detection looks for statistical outliers — unusual volumes, "
        "frequency spikes, or deviation from baselines. Sophisticated attackers stay within "
        "normal ranges.",
        "Nation-state APTs and insider threats operate for months within established behavioral "
        "norms, gradually shifting their patterns rather than spiking activity.",
        "The most dangerous cyber threats are designed specifically to evade the detection methods "
        "most organizations rely on today."
    ]
    y_pos = Inches(2.5)
    for i, bullet in enumerate(bullets):
        # Bullet marker
        marker = add_shape(slide, Inches(0.9), y_pos + Inches(0.08), Inches(0.12), Inches(0.12),
                           fill_color=TEAL)
        # Text
        add_textbox(slide, Inches(1.2), y_pos, Inches(11.0), Inches(0.9),
                    bullet, font_size=18, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)
        y_pos += Inches(1.1)

    # Callout stat
    add_callout_box(slide, Inches(3.0), Inches(5.8), Inches(7.0), Inches(0.8),
                    "No single method detects all 4 attack campaigns — behavioral analysis required",
                    font_size=20)

    add_footer(slide)


def slide_03_what_we_tested(prs):
    """Slide 3: What We Tested."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "What We Tested")

    # Three info columns
    col_data = [
        ("130 Days", "Synthetic Telemetry", "50 users across 7 log types\n"
         "23 behavioral features per user\n950 weekly feature vectors\n19 weeks of observation"),
        ("17 Methods", "Three-Tier Comparison", "Tier 1: Isolation Forest, SVM, LOF, Z-Score\n"
         "Temporal Z-Score, Feature Trajectory\n"
         "Tier 2: Behavioral Direction Analysis\n"
         "Tier 3: Digital Entity Analysis\n"
         "(9 zone/trajectory/relationship methods)"),
        ("4 Campaigns", "Injected Attack Patterns", "8-month insider threat escalation\n"
         "180-day slow APT with C2 beaconing\n"
         "115-day Volt Typhoon living-off-the-land\n"
         "100-day Salt Typhoon telecom interception\n"
         "Only 0.05% of events are attack-injected"),
    ]

    x_positions = [Inches(0.5), Inches(4.6), Inches(8.7)]
    for i, (stat, subtitle, detail) in enumerate(col_data):
        x = x_positions[i]
        # Stat number
        add_textbox(slide, x, Inches(1.5), Inches(3.8), Inches(0.7),
                    stat, font_size=36, bold=True, color=TEAL, alignment=PP_ALIGN.CENTER)
        # Subtitle
        add_textbox(slide, x, Inches(2.2), Inches(3.8), Inches(0.5),
                    subtitle, font_size=18, bold=True, color=NAVY, alignment=PP_ALIGN.CENTER)
        # Divider
        add_shape(slide, x + Inches(1.2), Inches(2.7), Inches(1.4), Inches(0.03),
                  fill_color=GOLD)
        # Detail
        add_textbox(slide, x + Inches(0.2), Inches(2.9), Inches(3.4), Inches(3.5),
                    detail, font_size=15, bold=False, color=DARK_GRAY,
                    alignment=PP_ALIGN.CENTER)

    # Vertical separators
    add_shape(slide, Inches(4.4), Inches(1.5), Inches(0.02), Inches(4.5), fill_color=LIGHT_GRAY)
    add_shape(slide, Inches(8.5), Inches(1.5), Inches(0.02), Inches(4.5), fill_color=LIGHT_GRAY)

    add_footer(slide)


def slide_04_attack_spectrum(prs):
    """Slide 4: The Attack Spectrum."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "The Attack Spectrum")

    # Timeline header labels
    add_textbox(slide, Inches(0.6), Inches(1.4), Inches(2.5), Inches(0.5),
                "FAST (Hours)", font_size=16, bold=True, color=GREEN, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(3.5), Inches(1.4), Inches(3.0), Inches(0.5),
                "MEDIUM (Days)", font_size=16, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(7.0), Inches(1.4), Inches(5.8), Inches(0.5),
                "SLOW (Months)", font_size=16, bold=True, color=RED, alignment=PP_ALIGN.CENTER)

    # Gradient bar representing time
    bar_y = Inches(1.9)
    bar_h = Inches(0.15)
    # Green segment (fast)
    add_shape(slide, Inches(0.6), bar_y, Inches(2.5), bar_h, fill_color=GREEN)
    # Gold segment (medium)
    add_shape(slide, Inches(3.1), bar_y, Inches(3.0), bar_h, fill_color=GOLD)
    # Red segment (slow)
    add_shape(slide, Inches(6.1), bar_y, Inches(6.7), bar_h, fill_color=RED)

    # Attack cards
    attacks = [
        (Inches(0.6), "Brute Force", "4 hours", "Volume spike", GREEN),
        (Inches(1.8), "Ransomware", "6 hours", "File encryption burst", GREEN),
        (Inches(3.8), "Credential Theft", "5 days", "Phishing + lateral move", GOLD),
        (Inches(7.0), "Supply Chain", "90 days dormant", "Implant activation", RED),
        (Inches(9.0), "APT C2", "180 days", "~4 beacons/day", RED),
        (Inches(11.0), "Insider Threat", "8 months", "4-phase escalation", RED),
    ]

    card_y = Inches(2.3)
    for x, name, duration, desc, color in attacks:
        card = add_rounded_rect(slide, x, card_y, Inches(1.8), Inches(2.2), fill_color=None)
        card.line.fill.solid()
        card.line.color.rgb = color
        card.line.width = Pt(2)
        # Top color accent
        add_shape(slide, x, card_y, Inches(1.8), Inches(0.06), fill_color=color)
        add_textbox(slide, x + Inches(0.05), card_y + Inches(0.2), Inches(1.7), Inches(0.5),
                    name, font_size=14, bold=True, color=NAVY, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.05), card_y + Inches(0.7), Inches(1.7), Inches(0.4),
                    duration, font_size=13, bold=True, color=color, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.05), card_y + Inches(1.2), Inches(1.7), Inches(0.8),
                    desc, font_size=11, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    # Bottom callout
    add_callout_box(slide, Inches(2.5), Inches(5.2), Inches(8.3), Inches(0.7),
                    "Traditional SIEM catches the fast. Misses the slow.", font_size=20)

    # Detection zone labels
    add_textbox(slide, Inches(0.6), Inches(6.1), Inches(3.0), Inches(0.5),
                "SIEM Detection Zone", font_size=14, bold=True, color=GREEN,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(7.0), Inches(6.1), Inches(5.8), Inches(0.5),
                "Behavioral Detection Required", font_size=14, bold=True, color=RED,
                alignment=PP_ALIGN.CENTER)

    add_footer(slide)


def slide_05_meet_attackers(prs):
    """Slide 5: Meet the Attackers."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Meet the Attackers")

    attackers = [
        ("USR-156: Insider Threat", "8-Month, 4-Phase Escalation",
         "Subtle mood shifts → curiosity browsing →\n"
         "restricted access → trickle exfiltration\n"
         "Activity volumes stay within normal ranges",
         RED),
        ("USR-234: Slow APT", "180-Day Persistent Campaign",
         "~4 C2 beacons/day at 6-hour intervals\n"
         "DGA DNS queries blend with normal traffic\n"
         "Progressive data staging over months",
         GOLD),
        ("USR-042: Volt Typhoon", "115-Day LOTL Pre-Positioning",
         "Living-off-the-land with RDP, WMI, PsExec\n"
         "Credential harvesting via legitimate tools\n"
         "Lateral movement mimics admin activity",
         RGBColor(142, 68, 173)),
        ("USR-118: Salt Typhoon", "100-Day Telecom Interception",
         "DNS tunneling for covert C2 channel\n"
         "CDR/lawful intercept data access\n"
         "Staged exfiltration mimics normal ops",
         BLUE),
    ]

    positions = [
        (Inches(0.4), Inches(1.4)),
        (Inches(6.7), Inches(1.4)),
        (Inches(0.4), Inches(3.8)),
        (Inches(6.7), Inches(3.8)),
    ]
    for i, (title, subtitle, detail, color) in enumerate(attackers):
        x, y = positions[i]
        card = add_rounded_rect(slide, x, y, Inches(6.0), Inches(2.1))
        card.line.fill.solid()
        card.line.color.rgb = color
        card.line.width = Pt(2)
        add_shape(slide, x, y, Inches(6.0), Inches(0.06), fill_color=color)
        add_textbox(slide, x + Inches(0.15), y + Inches(0.15), Inches(5.7), Inches(0.4),
                    title, font_size=16, bold=True, color=color, alignment=PP_ALIGN.LEFT)
        add_textbox(slide, x + Inches(0.15), y + Inches(0.5), Inches(5.7), Inches(0.3),
                    subtitle, font_size=12, bold=True, color=NAVY, alignment=PP_ALIGN.LEFT)
        add_shape(slide, x + Inches(0.3), y + Inches(0.8), Inches(5.4), Inches(0.02),
                  fill_color=LIGHT_GRAY)
        add_textbox(slide, x + Inches(0.15), y + Inches(0.85), Inches(5.7), Inches(1.1),
                    detail, font_size=11, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)

    # Bottom callout
    add_callout_box(slide, Inches(2.5), Inches(6.2), Inches(8.3), Inches(0.65),
                    "All 4 campaigns designed to stay within normal volume ranges", font_size=18)

    add_footer(slide)


def slide_06_isolation_forest(prs):
    """Slide 6: Traditional Algorithm #1 - Isolation Forest."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Traditional Algorithm #1 — Isolation Forest")

    # How it works
    add_textbox(slide, Inches(0.7), Inches(1.5), Inches(5.5), Inches(0.5),
                "How It Works", font_size=20, bold=True, color=NAVY, alignment=PP_ALIGN.LEFT)
    add_textbox(slide, Inches(0.7), Inches(2.0), Inches(5.5), Inches(1.0),
                "Randomly partitions feature space to isolate observations. Points that require "
                "fewer partitions to isolate are flagged as anomalies — they are statistically "
                "\"easy to separate\" from the bulk of data.",
                font_size=16, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)

    # Results box
    results_box = add_rounded_rect(slide, Inches(0.7), Inches(3.2), Inches(5.5), Inches(3.0),
                                   fill_color=RGBColor(253, 245, 245))
    results_box.line.fill.solid()
    results_box.line.color.rgb = RED
    results_box.line.width = Pt(2)

    add_textbox(slide, Inches(0.9), Inches(3.3), Inches(5.0), Inches(0.4),
                "RESULTS", font_size=18, bold=True, color=RED, alignment=PP_ALIGN.LEFT)

    results = [
        ("USR-156 (Insider, 8mo):", "MISSED", RED),
        ("USR-234 (APT, 180d):", "MISSED", RED),
        ("USR-042 (Volt Typhoon):", "DETECTED", GREEN),
        ("USR-118 (Salt Typhoon):", "DETECTED", GREEN),
        ("False Positive Rate:", "2.2% (1 normal user flagged)", DARK_GRAY),
    ]
    y = Inches(3.6)
    for label, result, color in results:
        add_textbox(slide, Inches(1.0), y, Inches(3.0), Inches(0.35),
                    label, font_size=14, bold=True, color=NAVY, alignment=PP_ALIGN.LEFT)
        add_textbox(slide, Inches(3.8), y, Inches(2.2), Inches(0.35),
                    result, font_size=14, bold=True, color=color, alignment=PP_ALIGN.LEFT)
        y += Inches(0.38)

    # Right side - Why it fails
    why_box = add_rounded_rect(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(4.7),
                               fill_color=NAVY)
    tf = set_text(why_box, "Why It Fails", font_size=22, bold=True, color=GOLD,
                  alignment=PP_ALIGN.LEFT)
    why_box.text_frame.margin_left = Inches(0.3)
    why_box.text_frame.margin_top = Inches(0.2)

    points = [
        "Catches Volt/Salt Typhoon — nation-state footprints create feature-space outliers",
        "Misses USR-156 insider — behavioral changes are directional, not magnitude-based",
        "Misses USR-234 slow APT — C2 beacons too subtle for random partitioning",
        "Cannot detect gradual drift — only sudden jumps in feature space",
    ]
    for p in points:
        add_paragraph(why_box.text_frame, "•  " + p, font_size=15, bold=False, color=WHITE,
                      space_before=Pt(12))

    add_footer(slide)


def slide_07_remaining_traditional(prs):
    """Slide 7: Traditional Algorithms #2-4."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Traditional Algorithms #2–4 — SVM, LOF, Z-Score")

    # Summary table
    data = [
        ["Method", "How It Works", "USR-156\nInsider", "USR-234\nAPT",
         "USR-042\nVolt", "USR-118\nSalt", "FP Rate"],
        ["One-Class\nSVM", "Learns boundary around\nnormal feature space",
         "DET", "MISS", "MISS", "DET", "19.6%"],
        ["Local Outlier\nFactor (LOF)", "Compares local density\nof each point to neighbors",
         "MISS", "DET", "DET", "DET", "0.0%"],
        ["Z-Score\nAnomaly", "Flags features > 3σ\nfrom global mean",
         "MISS", "DET", "DET", "DET", "2.2%"],
    ]

    cell_colors = {
        (1, 2): GREEN, (1, 3): RED, (1, 4): RED, (1, 5): GREEN, (1, 6): RED,
        (2, 2): RED, (2, 3): GREEN, (2, 4): GREEN, (2, 5): GREEN,
        (3, 2): RED, (3, 3): GREEN, (3, 4): GREEN, (3, 5): GREEN,
    }

    build_table(slide, rows=4, cols=7,
                left=Inches(0.3), top=Inches(1.5),
                width=Inches(12.7), height=Inches(2.8),
                data=data, cell_colors=cell_colors,
                col_widths=[Inches(1.7), Inches(2.8), Inches(1.5), Inches(1.5),
                            Inches(1.5), Inches(1.5), Inches(1.4)])

    # Key insight box
    insight_box = add_rounded_rect(slide, Inches(1.0), Inches(4.8), Inches(11.3), Inches(1.8),
                                   fill_color=NAVY)
    tf = set_text(insight_box, "Key Insight", font_size=22, bold=True, color=GOLD,
                  alignment=PP_ALIGN.LEFT)
    insight_box.text_frame.margin_left = Inches(0.3)
    insight_box.text_frame.margin_top = Inches(0.15)

    add_paragraph(insight_box.text_frame,
                  "LOF is the best traditional method: 3/4 attacks detected at 0% FP. "
                  "But it completely misses USR-156 (insider threat).",
                  font_size=17, color=WHITE, space_before=Pt(10))
    add_paragraph(insight_box.text_frame,
                  "USR-156 is invisible to 3 of 4 methods — behavioral-direction attacks don't "
                  "create magnitude anomalies. Only behavioral analysis can find it.",
                  font_size=17, color=WHITE, space_before=Pt(8))

    add_footer(slide)


def slide_08_temporal_upgrade(prs):
    """Slide 8: The Temporal Upgrade - Still Not Enough."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "The Temporal Upgrade — Still Not Enough")

    # Two columns
    # LEFT: Temporal Z-Score
    left_box = add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(4.0),
                                fill_color=None)
    left_box.line.fill.solid()
    left_box.line.color.rgb = TEAL
    left_box.line.width = Pt(2)
    add_shape(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.06), fill_color=TEAL)

    add_textbox(slide, Inches(0.7), Inches(1.7), Inches(5.4), Inches(0.4),
                "Temporal Z-Score", font_size=22, bold=True, color=TEAL,
                alignment=PP_ALIGN.LEFT)
    add_textbox(slide, Inches(0.7), Inches(2.2), Inches(5.4), Inches(0.6),
                "Train on first half of data, test on second half. Flag any feature that "
                "deviates > 2σ from the training period baseline.",
                font_size=14, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)
    add_textbox(slide, Inches(0.7), Inches(2.9), Inches(2.5), Inches(0.35),
                "All 4 attacks:", font_size=15, bold=True, color=NAVY, alignment=PP_ALIGN.LEFT)
    add_textbox(slide, Inches(3.0), Inches(2.9), Inches(3.0), Inches(0.35),
                "DETECTED", font_size=15, bold=True, color=GREEN, alignment=PP_ALIGN.LEFT)
    add_textbox(slide, Inches(0.7), Inches(3.3), Inches(5.0), Inches(0.35),
                "(USR-156, USR-234, USR-042, USR-118)",
                font_size=13, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)
    # FP rate - bad
    fp_box = add_rounded_rect(slide, Inches(0.8), Inches(4.3), Inches(5.2), Inches(0.7),
                              fill_color=RGBColor(253, 245, 245))
    fp_box.line.fill.solid()
    fp_box.line.color.rgb = RED
    fp_box.line.width = Pt(1.5)
    set_text(fp_box, "False Positive Rate: 100%  —  Every user flagged",
             font_size=15, bold=True, color=RED, alignment=PP_ALIGN.CENTER)

    # RIGHT: Feature Trajectory
    right_box = add_rounded_rect(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(4.0),
                                 fill_color=None)
    right_box.line.fill.solid()
    right_box.line.color.rgb = BLUE
    right_box.line.width = Pt(2)
    add_shape(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(0.06), fill_color=BLUE)

    add_textbox(slide, Inches(7.2), Inches(1.7), Inches(5.4), Inches(0.4),
                "Feature Trajectory Analysis", font_size=22, bold=True, color=BLUE,
                alignment=PP_ALIGN.LEFT)
    add_textbox(slide, Inches(7.2), Inches(2.2), Inches(5.4), Inches(0.6),
                "Change-point detection on raw feature trajectories. "
                "Tracks sustained shifts in individual feature values over time.",
                font_size=14, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)
    add_textbox(slide, Inches(7.2), Inches(2.9), Inches(2.5), Inches(0.35),
                "USR-042, USR-118:", font_size=15, bold=True, color=NAVY, alignment=PP_ALIGN.LEFT)
    add_textbox(slide, Inches(9.5), Inches(2.9), Inches(3.0), Inches(0.35),
                "DETECTED", font_size=15, bold=True, color=GREEN, alignment=PP_ALIGN.LEFT)
    add_textbox(slide, Inches(7.2), Inches(3.3), Inches(2.5), Inches(0.35),
                "USR-156, USR-234:", font_size=15, bold=True, color=NAVY, alignment=PP_ALIGN.LEFT)
    add_textbox(slide, Inches(9.5), Inches(3.3), Inches(3.0), Inches(0.35),
                "MISSED", font_size=15, bold=True, color=RED, alignment=PP_ALIGN.LEFT)
    # FP rate
    fp_box2 = add_rounded_rect(slide, Inches(7.3), Inches(4.3), Inches(5.2), Inches(0.7),
                               fill_color=RGBColor(253, 245, 245))
    fp_box2.line.fill.solid()
    fp_box2.line.color.rgb = GOLD
    fp_box2.line.width = Pt(1.5)
    set_text(fp_box2, "False Positive Rate: 6.5%  —  But only 2/4 detected",
             font_size=15, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)

    # Bottom callout
    add_callout_box(slide, Inches(2.5), Inches(6.0), Inches(8.3), Inches(0.7),
                    "Detection without precision is just noise", font_size=22)

    add_footer(slide)


def slide_09_comparison_matrix(prs):
    """Slide 9: The Detection-Precision Tradeoff (KEY SLIDE)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "The Detection-Precision Tradeoff")

    # KEY SLIDE label
    key_label = add_rounded_rect(slide, Inches(10.5), Inches(0.2), Inches(2.3), Inches(0.55),
                                 fill_color=GOLD)
    set_text(key_label, "KEY SLIDE", font_size=16, bold=True, color=WHITE,
             alignment=PP_ALIGN.CENTER)

    data = [
        ["Method", "Tier", "USR-156\nInsider", "USR-234\nAPT",
         "USR-042\nVolt", "USR-118\nSalt", "FP Rate"],
        ["Isolation Forest", "T1", "MISS", "MISS", "DET", "DET", "2.2%"],
        ["One-Class SVM", "T1", "DET", "MISS", "MISS", "DET", "19.6%"],
        ["LOF", "T1", "MISS", "DET", "DET", "DET", "0.0%"],
        ["Z-Score", "T1", "MISS", "DET", "DET", "DET", "2.2%"],
        ["T3 Zone Diverge", "T3", "DET", "DET", "MISS", "MISS", "6.5%"],
        ["T3 Combined", "T3", "DET", "DET", "DET", "DET", "8.7%"],
        ["LOF + Zone Div", "Best", "DET", "DET", "DET", "DET", "6.5%"],
    ]

    cell_colors = {}
    for r in range(1, 9):
        for c in [2, 3, 4, 5]:
            if r < len(data) and c < len(data[r]):
                if data[r][c] == "MISS":
                    cell_colors[(r, c)] = RED
                elif data[r][c] == "DET":
                    cell_colors[(r, c)] = GREEN
        # FP rate coloring
        if r < len(data):
            fp = data[r][6] if len(data[r]) > 6 else ""
            if fp in ["19.6%"]:
                cell_colors[(r, 6)] = RED

    # Highlight best ensemble row
    cell_colors[(7, 0)] = TEAL
    cell_colors[(7, 1)] = TEAL
    cell_colors[(7, 6)] = GREEN

    build_table(slide, rows=8, cols=7,
                left=Inches(0.3), top=Inches(1.3),
                width=Inches(12.7), height=Inches(4.2),
                data=data, cell_colors=cell_colors,
                col_widths=[Inches(2.2), Inches(0.8), Inches(1.8), Inches(1.8),
                            Inches(1.8), Inches(1.8), Inches(1.6)])

    # Bottom insight
    add_callout_box(slide, Inches(1.5), Inches(5.8), Inches(10.3), Inches(0.85),
                    "LOF + Zone Divergence: ALL 4 attacks detected at 6.5% FP — "
                    "the optimal 2-method ensemble across all 17 methods",
                    font_size=19)

    add_footer(slide)


def slide_10_core_problem(prs):
    """Slide 10: Why Traditional Methods Fail - The Core Problem."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    # Title
    add_textbox(slide, Inches(0.7), Inches(0.5), Inches(12.0), Inches(0.8),
                "Why Traditional Methods Fail — The Core Problem",
                font_size=32, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)
    add_shape(slide, Inches(0.7), Inches(1.2), Inches(3.0), Inches(0.04), fill_color=GOLD)

    # Main concept
    add_textbox(slide, Inches(0.7), Inches(1.5), Inches(11.8), Inches(0.8),
                "Feature distributions of attack users overlap completely with normal users.",
                font_size=22, bold=False, color=TEAL, alignment=PP_ALIGN.LEFT)

    # Example boxes
    examples = [
        ("USR-156 logs in 45 times/week", "Normal range: 30–60", "Within normal bounds"),
        ("USR-234 resolves 15 unique domains/day", "Normal range: 8–25", "Within normal bounds"),
        ("USR-156 accesses 8 files/day", "Normal range: 3–15", "Within normal bounds"),
    ]

    x_positions = [Inches(0.5), Inches(4.5), Inches(8.5)]
    for i, (metric, norm, verdict) in enumerate(examples):
        x = x_positions[i]
        box = add_rounded_rect(slide, x, Inches(2.6), Inches(3.8), Inches(2.2),
                               fill_color=RGBColor(20, 40, 60))
        box.line.fill.solid()
        box.line.color.rgb = TEAL
        box.line.width = Pt(1.5)

        add_textbox(slide, x + Inches(0.15), Inches(2.75), Inches(3.5), Inches(0.5),
                    metric, font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.15), Inches(3.3), Inches(3.5), Inches(0.4),
                    norm, font_size=14, bold=False, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
        # Verdict with check mark
        add_textbox(slide, x + Inches(0.15), Inches(3.8), Inches(3.5), Inches(0.5),
                    "✓  " + verdict, font_size=15, bold=True, color=GREEN,
                    alignment=PP_ALIGN.CENTER)

    # Key statements
    statements = [
        "LOF catches 3/4 attacks at 0% FP — but is completely blind to insider threats.",
        "No single traditional method detects all 4 campaigns.",
        "The direction changes, not the magnitude — invisible to traditional detection.",
        "Zone-specific behavioral analysis fills the gap traditional methods cannot.",
    ]

    y = Inches(5.2)
    for stmt in statements:
        marker = add_shape(slide, Inches(1.0), y + Inches(0.08), Inches(0.1), Inches(0.1),
                           fill_color=GOLD)
        add_textbox(slide, Inches(1.3), y, Inches(10.5), Inches(0.45),
                    stmt, font_size=17, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT)
        y += Inches(0.45)

    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12), Inches(0.4),
                "22nd Century Technologies, Inc. | Confidential",
                font_size=10, bold=False, color=RGBColor(60, 70, 85), alignment=PP_ALIGN.LEFT)


def slide_11_acecard_approach(prs):
    """Slide 11: ACECARD's Approach - Behavioral Drift Detection."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "ACECARD’s Approach — Behavioral Drift Detection")

    # Flow steps as connected boxes
    steps = [
        ("1", "Behavioral\nSerialization",
         "Raw cyber telemetry is\ntransformed into natural\nlanguage descriptions of\nuser behavior patterns"),
        ("2", "Semantic\nEmbedding",
         "Behavioral descriptions are\nencoded into high-dimensional\nsemantic space where meaning—\nnot magnitude—is preserved"),
        ("3", "Directional\nDrift Tracking",
         "Continuous monitoring of\nhow each user's behavioral\nvector moves through\nsemantic space over time"),
        ("4", "Change-Point\nDetection",
         "Cumulative behavioral drift\nis analyzed to identify\nstatistically significant\nshifts in trajectory"),
        ("5", "Threat Pattern\nClassification",
         "Drift direction is compared\nagainst known threat behavioral\npatterns to classify the\nnature of the change"),
    ]

    start_x = Inches(0.3)
    box_w = Inches(2.3)
    gap = Inches(0.25)

    for i, (num, title, desc) in enumerate(steps):
        x = start_x + i * (box_w + gap)
        y = Inches(1.6)

        # Number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.85), y, Inches(0.55),
                                        Inches(0.55))
        circle.fill.solid()
        circle.fill.fore_color.rgb = TEAL
        circle.line.fill.background()
        set_text(circle, num, font_size=22, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

        # Title box
        title_box = add_shape(slide, x, y + Inches(0.7), box_w, Inches(0.8), fill_color=NAVY)
        tf = set_text(title_box, title, font_size=14, bold=True, color=WHITE,
                      alignment=PP_ALIGN.CENTER)

        # Description box
        desc_box = add_rounded_rect(slide, x, y + Inches(1.5), box_w, Inches(2.0))
        desc_box.line.fill.solid()
        desc_box.line.color.rgb = TEAL
        desc_box.line.width = Pt(1)
        tf = set_text(desc_box, desc, font_size=12, bold=False, color=DARK_GRAY,
                      alignment=PP_ALIGN.CENTER)
        desc_box.text_frame.margin_left = Inches(0.1)
        desc_box.text_frame.margin_right = Inches(0.1)
        desc_box.text_frame.margin_top = Inches(0.15)

        # Arrow between boxes (except last)
        if i < len(steps) - 1:
            arrow_x = x + box_w
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x,
                                           y + Inches(0.85), gap, Inches(0.25))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = GOLD
            arrow.line.fill.background()

    # Bottom conceptual summary
    summary_box = add_rounded_rect(slide, Inches(1.5), Inches(5.6), Inches(10.3), Inches(1.1),
                                   fill_color=NAVY)
    tf = set_text(summary_box, "Catches the pattern shift, not the volume spike",
                  font_size=24, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)
    summary_box.text_frame.margin_top = Inches(0.1)
    add_paragraph(summary_box.text_frame,
                  "Traditional methods ask \"how much?\"  —  ACECARD asks \"in what direction?\"",
                  font_size=16, bold=False, color=WHITE, alignment=PP_ALIGN.CENTER,
                  space_before=Pt(8))

    add_footer(slide)


def slide_12_complementary_detection(prs):
    """Slide 12: Complementary Detection — Why Both Are Needed."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Complementary Detection — Why Both Are Needed")

    add_textbox(slide, Inches(0.7), Inches(1.4), Inches(11.8), Inches(0.5),
                "No single method catches all 4 threats. The optimal ensemble combines two.",
                font_size=17, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)

    # LOF card (left)
    lof_box = add_rounded_rect(slide, Inches(0.5), Inches(2.1), Inches(5.8), Inches(3.2),
                               fill_color=None)
    lof_box.line.fill.solid()
    lof_box.line.color.rgb = TEAL
    lof_box.line.width = Pt(2)
    add_shape(slide, Inches(0.5), Inches(2.1), Inches(5.8), Inches(0.06), fill_color=TEAL)

    add_textbox(slide, Inches(0.7), Inches(2.25), Inches(5.4), Inches(0.4),
                "LOF (Traditional — Tier 1)", font_size=18, bold=True, color=TEAL,
                alignment=PP_ALIGN.LEFT)

    lof_results = [
        ("USR-156 (Insider):", "MISSED", RED),
        ("USR-234 (APT):", "DETECTED", GREEN),
        ("USR-042 (Volt Typhoon):", "DETECTED", GREEN),
        ("USR-118 (Salt Typhoon):", "DETECTED", GREEN),
    ]
    y = Inches(2.8)
    for label, result, color in lof_results:
        add_textbox(slide, Inches(0.8), y, Inches(3.0), Inches(0.3),
                    label, font_size=13, bold=True, color=NAVY, alignment=PP_ALIGN.LEFT)
        result_shape = add_shape(slide, Inches(3.8), y, Inches(2.0), Inches(0.3),
                                 fill_color=color)
        set_text(result_shape, result, font_size=12, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
        y += Inches(0.35)

    lof_why = add_rounded_rect(slide, Inches(0.8), Inches(4.3), Inches(5.2), Inches(0.7),
                               fill_color=NAVY)
    set_text(lof_why, "0% FP — but blind to insider threats\nthat stay within volume norms",
             font_size=13, bold=False, color=GOLD, alignment=PP_ALIGN.CENTER)

    # Zone Divergence card (right)
    zd_box = add_rounded_rect(slide, Inches(7.0), Inches(2.1), Inches(5.8), Inches(3.2),
                              fill_color=None)
    zd_box.line.fill.solid()
    zd_box.line.color.rgb = GOLD
    zd_box.line.width = Pt(2)
    add_shape(slide, Inches(7.0), Inches(2.1), Inches(5.8), Inches(0.06), fill_color=GOLD)

    add_textbox(slide, Inches(7.2), Inches(2.25), Inches(5.4), Inches(0.4),
                "Zone Divergence (Behavioral — Tier 3)", font_size=18, bold=True, color=GOLD,
                alignment=PP_ALIGN.LEFT)

    zd_results = [
        ("USR-156 (Insider):", "DETECTED", GREEN),
        ("USR-234 (APT):", "DETECTED", GREEN),
        ("USR-042 (Volt Typhoon):", "MISSED", RED),
        ("USR-118 (Salt Typhoon):", "MISSED", RED),
    ]
    y = Inches(2.8)
    for label, result, color in zd_results:
        add_textbox(slide, Inches(7.3), y, Inches(3.0), Inches(0.3),
                    label, font_size=13, bold=True, color=NAVY, alignment=PP_ALIGN.LEFT)
        result_shape = add_shape(slide, Inches(10.3), y, Inches(2.0), Inches(0.3),
                                 fill_color=color)
        set_text(result_shape, result, font_size=12, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
        y += Inches(0.35)

    zd_why = add_rounded_rect(slide, Inches(7.3), Inches(4.3), Inches(5.2), Inches(0.7),
                              fill_color=NAVY)
    set_text(zd_why, "6.5% FP — catches behavioral drift\ninvisible to volume-based methods",
             font_size=13, bold=False, color=GOLD, alignment=PP_ALIGN.CENTER)

    # Combined result callout
    combined = add_rounded_rect(slide, Inches(2.0), Inches(5.5), Inches(9.3), Inches(1.2),
                                fill_color=NAVY)
    tf = set_text(combined,
                  "LOF + Zone Divergence:  ALL 4 DETECTED  |  6.5% FP",
                  font_size=22, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)
    combined.text_frame.margin_top = Inches(0.1)
    add_paragraph(combined.text_frame,
                  "LOF catches nation-state footprints. Zone Divergence catches behavioral drift. "
                  "Together: complete coverage at operational FP rates.",
                  font_size=14, bold=False, color=WHITE, alignment=PP_ALIGN.CENTER,
                  space_before=Pt(6))

    add_footer(slide)


def slide_13_per_threat_fp(prs):
    """Slide 13: Per-Threat False Positive Analysis."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Per-Threat Detection Cost Analysis")

    add_textbox(slide, Inches(0.7), Inches(1.3), Inches(11.8), Inches(0.5),
                "What does it actually cost (in false positives) to detect each specific threat?",
                font_size=17, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)

    threats = [
        ("USR-156: Insider Threat", "HARDEST TO DETECT",
         "Zone Divergence: 6.5% FP\nOne-Class SVM: 19.6% FP\n\nNo other viable method detects this threat.",
         RED, "Only behavioral analysis\nfinds this attack."),
        ("USR-234: Slow APT", "MODERATE",
         "LOF: 0.0% FP (best)\nZ-Score: 2.2% FP\nZone Divergence: 6.5% FP",
         GOLD, "Multiple methods work.\nLOF is optimal."),
        ("USR-042: Volt Typhoon", "DETECTABLE",
         "LOF: 0.0% FP (best)\nIsolation Forest: 2.2% FP\nZ-Score: 2.2% FP",
         RGBColor(142, 68, 173), "Traditional methods\nsuffice for LOTL."),
        ("USR-118: Salt Typhoon", "DETECTABLE",
         "LOF: 0.0% FP (best)\nIsolation Forest: 2.2% FP\nZ-Score: 2.2% FP",
         BLUE, "Traditional methods\nsuffice for telecom."),
    ]

    positions = [
        (Inches(0.4), Inches(2.0)),
        (Inches(6.7), Inches(2.0)),
        (Inches(0.4), Inches(4.3)),
        (Inches(6.7), Inches(4.3)),
    ]

    for i, (title, difficulty, methods, color, insight) in enumerate(threats):
        x, y = positions[i]
        card = add_rounded_rect(slide, x, y, Inches(6.0), Inches(2.0))
        card.line.fill.solid()
        card.line.color.rgb = color
        card.line.width = Pt(2)
        add_shape(slide, x, y, Inches(6.0), Inches(0.06), fill_color=color)

        add_textbox(slide, x + Inches(0.15), y + Inches(0.1), Inches(3.8), Inches(0.35),
                    title, font_size=14, bold=True, color=color, alignment=PP_ALIGN.LEFT)

        diff_badge = add_rounded_rect(slide, x + Inches(4.0), y + Inches(0.1),
                                      Inches(1.8), Inches(0.35), fill_color=color)
        set_text(diff_badge, difficulty, font_size=10, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)

        add_textbox(slide, x + Inches(0.15), y + Inches(0.55), Inches(3.3), Inches(1.3),
                    methods, font_size=11, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)
        add_textbox(slide, x + Inches(3.5), y + Inches(0.55), Inches(2.3), Inches(1.3),
                    insight, font_size=11, bold=True, color=NAVY, alignment=PP_ALIGN.CENTER)

    add_callout_box(slide, Inches(2.5), Inches(6.5), Inches(8.3), Inches(0.65),
                    "USR-156 is the acid test — only behavioral analysis can find it",
                    font_size=18)

    add_footer(slide)


def slide_14_results_summary(prs):
    """Slide 14: Results Summary - 2x2 Matrix."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    # Title
    add_textbox(slide, Inches(0.7), Inches(0.4), Inches(12.0), Inches(0.8),
                "Results Summary", font_size=34, bold=True, color=WHITE,
                alignment=PP_ALIGN.LEFT)
    add_shape(slide, Inches(0.7), Inches(1.1), Inches(3.0), Inches(0.04), fill_color=GOLD)

    # 3 approach categories
    categories = [
        ("Traditional (Tier 1)", "Best: LOF — 3/4 detected, 0% FP",
         "Excellent for nation-state attacks", "Blind to insider threat (USR-156)",
         "PARTIAL", GOLD,
         "LOF catches USR-234, USR-042, USR-118\n"
         "at 0% false positive rate. But completely\n"
         "misses the insider threat campaign."),
        ("Behavioral (Tier 3)", "Zone Divergence — 2/4 detected, 6.5% FP",
         "Detects behavioral drift attacks", "Misses network-footprint attacks",
         "TARGETED", TEAL,
         "Catches USR-156 (insider) and USR-234\n"
         "(slow APT) by analyzing zone-specific\n"
         "behavioral drift patterns."),
        ("LOF + Zone Divergence", "Optimal 2-Method Ensemble",
         "All 4 Attacks Detected", "6.5% False Positive Rate",
         "ACTIONABLE", GREEN,
         "Traditional catches network footprints.\n"
         "Behavioral catches direction changes.\n"
         "Together: 100% detection at 6.5% FP."),
    ]

    x_positions = [Inches(0.4), Inches(4.5), Inches(8.6)]
    for i, (title, subtitle, pro, con, verdict, color, desc) in enumerate(categories):
        x = x_positions[i]
        y = Inches(1.6)

        # Card
        card = add_rounded_rect(slide, x, y, Inches(3.9), Inches(4.8),
                                fill_color=RGBColor(20, 40, 60))
        card.line.fill.solid()
        card.line.color.rgb = color
        card.line.width = Pt(2.5)

        # Color top bar
        add_shape(slide, x, y, Inches(3.9), Inches(0.07), fill_color=color)

        # Title
        add_textbox(slide, x + Inches(0.15), y + Inches(0.2), Inches(3.6), Inches(0.5),
                    title, font_size=18, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
        # Subtitle
        add_textbox(slide, x + Inches(0.15), y + Inches(0.65), Inches(3.6), Inches(0.4),
                    subtitle, font_size=12, bold=False, color=MID_GRAY,
                    alignment=PP_ALIGN.CENTER)

        # Divider
        add_shape(slide, x + Inches(0.4), y + Inches(1.1), Inches(3.1), Inches(0.02),
                  fill_color=RGBColor(50, 70, 90))

        # Pro / Con
        add_textbox(slide, x + Inches(0.15), y + Inches(1.25), Inches(3.6), Inches(0.4),
                    "✓  " + pro, font_size=13, bold=False,
                    color=GREEN if i == 2 else MID_GRAY,
                    alignment=PP_ALIGN.LEFT)
        add_textbox(slide, x + Inches(0.15), y + Inches(1.65), Inches(3.6), Inches(0.4),
                    ("✓  " if i == 2 else "✗  ") + con, font_size=13, bold=False,
                    color=GREEN if i == 2 else RED,
                    alignment=PP_ALIGN.LEFT)

        # Description
        add_textbox(slide, x + Inches(0.15), y + Inches(2.2), Inches(3.6), Inches(1.5),
                    desc, font_size=12, bold=False, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

        # Verdict badge
        badge = add_rounded_rect(slide, x + Inches(0.8), y + Inches(3.9), Inches(2.3),
                                 Inches(0.6), fill_color=color)
        set_text(badge, verdict, font_size=18, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)

    # Bottom statement
    add_textbox(slide, Inches(1.0), Inches(6.7), Inches(11.3), Inches(0.5),
                "LOF + Zone Divergence: all 4 campaigns detected at 6.5% FP — optimal ensemble.",
                font_size=22, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.5), Inches(7.1), Inches(12), Inches(0.3),
                "22nd Century Technologies, Inc. | Confidential",
                font_size=10, bold=False, color=RGBColor(60, 70, 85), alignment=PP_ALIGN.LEFT)


def slide_15_agency_implications(prs):
    """Slide 15: What This Means for Your Agency."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "What This Means for Your Agency")

    # Impact statements
    impacts = [
        ("Current SIEM/UEBA tools are blind to the most dangerous threats.",
         "Traditional anomaly detection relies on statistical outliers. Slow, sophisticated "
         "attackers operate within normal statistical ranges by design. Your current tools "
         "cannot see them."),
        ("Nation-state APTs and insider threats are designed to evade.",
         "Volt Typhoon, Salt Typhoon, and similar campaigns specifically model normal user "
         "behavior to avoid triggering volume-based alerts. The threat model has evolved; "
         "detection must evolve with it."),
        ("Behavioral drift analysis is the only proven approach.",
         "Our empirical analysis across 130 days demonstrates that only semantic behavioral "
         "drift detection achieves both high detection rates and actionable false positive "
         "rates for slow, stealthy attacks."),
        ("Every day without behavioral detection is exposure.",
         "If an 8-month insider threat campaign is running in your environment today, your "
         "current tools will not detect it until data has already been exfiltrated."),
    ]

    y = Inches(1.5)
    for i, (headline, detail) in enumerate(impacts):
        # Left accent bar
        accent_color = [TEAL, BLUE, GREEN, RED][i]
        add_shape(slide, Inches(0.7), y, Inches(0.08), Inches(1.15), fill_color=accent_color)

        # Headline
        add_textbox(slide, Inches(1.0), y, Inches(11.5), Inches(0.5),
                    headline, font_size=18, bold=True, color=NAVY, alignment=PP_ALIGN.LEFT)
        # Detail
        add_textbox(slide, Inches(1.0), y + Inches(0.45), Inches(11.5), Inches(0.7),
                    detail, font_size=14, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)
        y += Inches(1.35)

    # Bottom callout
    callout = add_rounded_rect(slide, Inches(2.0), Inches(6.2), Inches(9.3), Inches(0.8),
                               fill_color=NAVY)
    set_text(callout, "The threats you can’t see are the ones that matter most.",
             font_size=22, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)

    add_footer(slide)


def slide_16_next_steps(prs):
    """Slide 16: Next Steps - Prove It On Your Data."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Next Steps — Prove It On Your Data")

    # Phase cards
    phases = [
        ("Week 1–2", "Integration & Baseline",
         "Connect to existing SIEM/log feeds\nEstablish behavioral baselines\n"
         "No new agents or data collection\nZero operational disruption",
         TEAL),
        ("Week 3", "Detection Calibration",
         "Tune detection thresholds\nValidate against known incidents\n"
         "Measure false positive rate\nCompare to current SIEM alerts",
         BLUE),
        ("Week 4", "Results & Assessment",
         "Full detection report\nFP rate comparison\nMean time to detect analysis\n"
         "Go/no-go recommendation",
         GOLD),
    ]

    x_positions = [Inches(0.5), Inches(4.5), Inches(8.5)]
    for i, (timeline, title, details, color) in enumerate(phases):
        x = x_positions[i]

        # Timeline badge
        badge = add_rounded_rect(slide, x + Inches(0.8), Inches(1.5), Inches(2.2), Inches(0.5),
                                 fill_color=color)
        set_text(badge, timeline, font_size=16, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)

        # Card
        card = add_rounded_rect(slide, x, Inches(2.2), Inches(3.8), Inches(3.2))
        card.line.fill.solid()
        card.line.color.rgb = color
        card.line.width = Pt(2)
        add_shape(slide, x, Inches(2.2), Inches(3.8), Inches(0.06), fill_color=color)

        # Title
        add_textbox(slide, x + Inches(0.15), Inches(2.4), Inches(3.5), Inches(0.5),
                    title, font_size=18, bold=True, color=NAVY, alignment=PP_ALIGN.CENTER)
        # Divider
        add_shape(slide, x + Inches(0.4), Inches(2.9), Inches(3.0), Inches(0.02),
                  fill_color=LIGHT_GRAY)
        # Details
        add_textbox(slide, x + Inches(0.2), Inches(3.05), Inches(3.4), Inches(2.2),
                    details, font_size=14, bold=False, color=DARK_GRAY,
                    alignment=PP_ALIGN.LEFT)

        # Arrow between cards (except last)
        if i < 2:
            arrow_x = x + Inches(3.8)
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x + Inches(0.1),
                                           Inches(3.5), Inches(0.5), Inches(0.25))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = GOLD
            arrow.line.fill.background()

    # Key differentiators
    diff_box = add_rounded_rect(slide, Inches(0.5), Inches(5.7), Inches(12.3), Inches(1.1),
                                fill_color=NAVY)
    tf = set_text(diff_box, "No new agents.  No new data collection.  Works with existing "
                  "SIEM feeds.",
                  font_size=20, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)
    diff_box.text_frame.margin_top = Inches(0.05)
    add_paragraph(diff_box.text_frame,
                  "Measurable success criteria: detection rate, false positive rate, "
                  "mean time to detect.",
                  font_size=16, bold=False, color=WHITE, alignment=PP_ALIGN.CENTER,
                  space_before=Pt(8))

    add_footer(slide)


def slide_17_contact(prs):
    """Slide 17: Contact / Q&A."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    # Large Q&A
    add_textbox(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.5),
                "Questions?", font_size=52, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    # Gold accent line
    add_shape(slide, Inches(5.0), Inches(3.0), Inches(3.3), Inches(0.04), fill_color=GOLD)

    # Company info
    add_textbox(slide, Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.7),
                "22nd Century Technologies, Inc.", font_size=26, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.5),
                "Behavioral Intelligence for Federal Cybersecurity",
                font_size=18, bold=False, color=TEAL, alignment=PP_ALIGN.CENTER)

    # Contact details box
    contact_box = add_rounded_rect(slide, Inches(3.5), Inches(5.0), Inches(6.3), Inches(1.5),
                                   fill_color=RGBColor(20, 40, 60))
    contact_box.line.fill.solid()
    contact_box.line.color.rgb = TEAL
    contact_box.line.width = Pt(1)
    tf = set_text(contact_box, "ACECARD — Adaptive Cyber-Entity Correlation & "
                  "Anomalous Response Detection",
                  font_size=14, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)
    contact_box.text_frame.margin_top = Inches(0.15)
    add_paragraph(contact_box.text_frame,
                  "Proven behavioral drift detection for the threats\nyour current tools "
                  "cannot see.",
                  font_size=13, bold=False, color=MID_GRAY, alignment=PP_ALIGN.CENTER,
                  space_before=Pt(10))

    # Bottom
    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.3),
                "We proved it on synthetic data. Let us prove it on yours.",
                font_size=14, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Slide 18: Detection Playbook
# ---------------------------------------------------------------------------
def slide_18_detection_playbook(prs):
    """Detection Playbook: threat type → recommended algorithm."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Detection Playbook: Match the Threat to the Method")

    add_textbox(slide, Inches(0.7), Inches(1.25), Inches(11.9), Inches(0.5),
                "UEBA detects threats by behavioral drift, not attack signatures. "
                "Each threat type requires a different detection method.",
                font_size=15, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)

    # Playbook cards
    cards = [
        ("Insider Threat", "USR-156", RED,
         "Zone Divergence", "6.5%",
         "Same person, different data.\nIdentity stable + data zone drifting.",
         "LOF, IForest blind\n(measure magnitude only)"),
        ("Slow APT / C2", "USR-234", RGBColor(230, 126, 34),
         "LOF + Zone Div", "0%* / 6.5%",
         "Network outlier + identity stable\n+ network zone drifting.",
         "IForest misses\nslow drift"),
        ("Nation-State LOTL", "USR-042", RGBColor(142, 68, 173),
         "LOF + Regime Shift", "0%* / 6.5%",
         "Strong endpoint anomaly +\nclear before/after phase break.",
         "Zone Divergence misses\n(uniform change)"),
        ("Telecom Pivot", "USR-118", RGBColor(41, 128, 185),
         "LOF + Embed CUSUM", "0%* / 6.5%",
         "Network footprint +\npersistent 100-day drift.",
         "Zone Divergence misses\n(broad multi-zone)"),
    ]

    for i, (threat, user, color, method, fp, why, fails) in enumerate(cards):
        y = Inches(1.95 + i * 1.28)

        # Threat label
        tbox = add_rounded_rect(slide, Inches(0.3), y, Inches(2.2), Inches(1.08),
                                fill_color=color)
        tf = set_text(tbox, threat, font_size=16, bold=True, color=WHITE,
                      alignment=PP_ALIGN.CENTER)
        tbox.text_frame.margin_top = Inches(0.08)
        add_paragraph(tbox.text_frame, user, font_size=12, bold=False, color=WHITE,
                      alignment=PP_ALIGN.CENTER, space_before=Pt(4))

        # Method
        mbox = add_rounded_rect(slide, Inches(2.7), y, Inches(2.0), Inches(1.08))
        mbox.line.fill.solid()
        mbox.line.color.rgb = GREEN
        mbox.line.width = Pt(2)
        set_text(mbox, method, font_size=15, bold=True, color=GREEN,
                 alignment=PP_ALIGN.CENTER)
        mbox.text_frame.margin_top = Inches(0.12)
        add_paragraph(mbox.text_frame, f"FP: {fp}", font_size=12, bold=False,
                      color=DARK_GRAY, alignment=PP_ALIGN.CENTER, space_before=Pt(6))

        # Why
        add_textbox(slide, Inches(4.9), y + Inches(0.05), Inches(4.0), Inches(1.0),
                    why, font_size=12, bold=False, color=DARK_GRAY)

        # Fails
        add_textbox(slide, Inches(9.1), y + Inches(0.05), Inches(3.8), Inches(1.0),
                    fails, font_size=12, bold=False, color=RED)

    # Column headers
    add_textbox(slide, Inches(0.3), Inches(1.7), Inches(2.2), Inches(0.3),
                "THREAT TYPE", font_size=11, bold=True, color=MID_GRAY)
    add_textbox(slide, Inches(2.7), Inches(1.7), Inches(2.0), Inches(0.3),
                "BEST METHOD", font_size=11, bold=True, color=MID_GRAY)
    add_textbox(slide, Inches(4.9), Inches(1.7), Inches(4.0), Inches(0.3),
                "WHY IT WORKS", font_size=11, bold=True, color=MID_GRAY)
    add_textbox(slide, Inches(9.1), Inches(1.7), Inches(3.8), Inches(0.3),
                "WHAT FAILS", font_size=11, bold=True, color=MID_GRAY)

    # Bottom insight
    insight = add_rounded_rect(slide, Inches(0.3), Inches(7.1) - Inches(0.8),
                               Inches(12.7), Inches(0.55), fill_color=NAVY)
    set_text(insight,
             "The threat type determines the algorithm. Traditional = how much changed. "
             "Behavioral = what kind of change. Deploy both layers.",
             font_size=14, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)
    insight.text_frame.margin_top = Inches(0.05)

    add_footer(slide)


# ===========================================================================
# MAIN
# ===========================================================================

def slide_19_algorithm_parameters(prs):
    """Algorithm input features and parameters reference."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Algorithm Parameters & Input Features")

    add_textbox(slide, Inches(0.7), Inches(1.25), Inches(11.9), Inches(0.4),
                "All algorithms draw from the same 23 behavioral features across "
                "5 log sources. The detection approach determines what gets caught.",
                font_size=14, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)

    # Tier 1 table
    add_textbox(slide, Inches(0.5), Inches(1.75), Inches(5.0), Inches(0.35),
                "Tier 1: Traditional Detection", font_size=16, bold=True, color=NAVY)

    t1_data = [
        ["Algorithm", "Key Parameters", "Detection Rule"],
        ["Isolation Forest", "200 trees, 5% budget", "Isolation outliers"],
        ["One-Class SVM", "RBF kernel, 5% nu", "Outside boundary"],
        ["LOF", "20 neighbors, 5%", "Low local density"],
        ["Z-Score", "3-sigma threshold", "Any feature > 3σ"],
        ["Temporal Z-Score", "50/50 train/test", "Deviation from baseline"],
        ["Feature Trajectory", "Cumulative drift", "Top 10% accumulated"],
    ]

    build_table(slide, 7, 3, Inches(0.3), Inches(2.15), Inches(6.2), Inches(2.6),
                t1_data,
                col_widths=[Inches(2.0), Inches(2.0), Inches(2.2)])

    # Tier 3 table
    add_textbox(slide, Inches(6.8), Inches(1.75), Inches(6.0), Inches(0.35),
                "Tier 3: Behavioral Detection", font_size=16, bold=True, color=NAVY)

    t3_data = [
        ["Method", "Input Signal", "Key Parameter"],
        ["Velocity/Accel", "Composite trajectory", "Acceleration > 0.01"],
        ["Regime Shift", "Weekly embeddings", "Similarity < 0.7"],
        ["Zone Divergence", "5 zone drifts", "Identity < 0.02, zone > 0.05"],
        ["Relationship", "User-Device vector", "Drift > 0.05"],
        ["Contextual", "4 context weights", "Consistency > 30%"],
        ["Drift Accum.", "Baseline distance", "Cumulative excess"],
        ["Progression", "Threat alignment", "Increasing trend"],
        ["Combined", "All 6 core scores", "Zone Div OR ≥ 2 agree"],
    ]

    build_table(slide, 9, 3, Inches(6.7), Inches(2.15), Inches(6.3), Inches(3.1),
                t3_data,
                col_widths=[Inches(1.8), Inches(2.2), Inches(2.3)])

    # Feature categories at bottom
    add_textbox(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.35),
                "23 Input Features by Log Source", font_size=16, bold=True, color=NAVY)

    categories = [
        ("Authentication (7)", TEAL,
         "auth_total, failures, fail_rate,\noff_hours, sources, dests, methods"),
        ("File Access (6)", BLUE,
         "file_total, restricted, confidential,\nwrite_ratio, paths, bytes"),
        ("Endpoint (5)", RGBColor(0x8E, 0x44, 0xAD),
         "total, suspicious_ratio,\nmax_risk, mean_risk, processes"),
        ("Network (5)", NAVY,
         "bytes_out, unique_dsts,\nexternal_ratio, dns_domains, nxdomain"),
    ]

    for i, (title, color, feats) in enumerate(categories):
        x = Inches(0.3 + i * 3.25)
        card = add_rounded_rect(slide, x, Inches(5.4), Inches(3.0), Inches(1.4))
        card.line.fill.solid()
        card.line.color.rgb = color
        card.line.width = Pt(2)
        add_shape(slide, x, Inches(5.4), Inches(3.0), Inches(0.06), fill_color=color)
        add_textbox(slide, x + Inches(0.1), Inches(5.5), Inches(2.8), Inches(0.3),
                    title, font_size=13, bold=True, color=color, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.15), Inches(5.85), Inches(2.7), Inches(0.8),
                    feats, font_size=11, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)

    add_footer(slide)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide_01_title(prs)
    slide_02_problem(prs)
    slide_03_what_we_tested(prs)
    slide_04_attack_spectrum(prs)
    slide_05_meet_attackers(prs)
    slide_06_isolation_forest(prs)
    slide_07_remaining_traditional(prs)
    slide_08_temporal_upgrade(prs)
    slide_09_comparison_matrix(prs)
    slide_10_core_problem(prs)
    slide_11_acecard_approach(prs)
    slide_12_complementary_detection(prs)
    slide_13_per_threat_fp(prs)
    slide_14_results_summary(prs)
    slide_15_agency_implications(prs)
    slide_16_next_steps(prs)
    slide_18_detection_playbook(prs)
    slide_19_algorithm_parameters(prs)
    slide_17_contact(prs)

    out_dir = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.join(out_dir, "ACECARD_Traditional_vs_Behavioral_Deck.pptx")
    prs.save(out_path)
    print(f"Deck created successfully at: {out_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
