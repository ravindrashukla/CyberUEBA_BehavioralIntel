#!/usr/bin/env python3
"""
Build Sales Pitch PowerPoint Deck — Expanded Edition.

NO proprietary logic disclosed. Focus on:
- The problem (nation-state APTs invisible to current tools)
- Real-world case studies (Salt Typhoon, Volt Typhoon)
- Why traditional detection fails
- The Digital Twin innovation (V-Intelligence UEBA creates behavioral digital twin)
- Multi-tier detection approach (Traditional → Intermediate → Full V-Intelligence UEBA)
- Screenshots showing results
- Unique results and competitive advantage
- Business impact

Output: docs/V_Intelligence_UEBA_Sales_Pitch.pptx
"""

import sys
sys.stdout.reconfigure(encoding="utf-8")

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colors ──
NAVY       = RGBColor(13,  27,  42)
BLUE       = RGBColor(27,  79,  114)
TEAL       = RGBColor(14,  107, 138)
GOLD       = RGBColor(183, 149, 11)
RED        = RGBColor(192, 57,  43)
GREEN      = RGBColor(39,  174, 96)
WHITE      = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(240, 242, 245)
DARK_GRAY  = RGBColor(80,  80,  80)
MID_GRAY   = RGBColor(150, 150, 150)
ORANGE     = RGBColor(230, 126, 34)
LIGHT_RED  = RGBColor(253, 237, 236)
LIGHT_GREEN = RGBColor(234, 250, 241)
LIGHT_GOLD = RGBColor(255, 251, 235)
LIGHT_BLUE = RGBColor(214, 234, 248)
PURPLE     = RGBColor(142, 68, 173)
LIGHT_PURPLE = RGBColor(245, 238, 248)
DARK_TEAL  = RGBColor(20,  40,  60)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

SCREENSHOT_DIR = "docs/screenshots"


# ── Helpers ──

def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, l, t, w, h, fill=None, line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, l, t, w, h)
    s.line.fill.background()
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.fill.solid()
        s.line.color.rgb = line
        s.line.width = Pt(2)
    return s


def tbox(slide, l, t, w, h, text, sz=18, bold=False, color=NAVY, align=PP_ALIGN.LEFT, name="Calibri"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = name
    p.alignment = align
    return tb


def set_text(shape, text, sz=18, bold=False, color=NAVY, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = align
    return tf


def add_p(tf, text, sz=16, bold=False, color=NAVY, align=PP_ALIGN.LEFT,
          before=Pt(4), after=Pt(2)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = align
    p.space_before = before
    p.space_after = after
    return p


def title_bar(slide, text, sz=28):
    bar = rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), fill=NAVY)
    set_text(bar, text, sz=sz, bold=True, color=WHITE)
    bar.text_frame.paragraphs[0].space_before = Pt(12)
    bar.text_frame.margin_left = Inches(0.6)
    bar.text_frame.margin_top = Inches(0.18)
    rect(slide, Inches(0), Inches(1.1), SLIDE_W, Inches(0.04), fill=GOLD)


def footer(slide, text="22nd Century Technologies, Inc.  |  V-Intelligence UEBA  |  Confidential"):
    tbox(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.35),
         text, sz=10, color=MID_GRAY)


def card(slide, l, t, w, h, fill=WHITE, border=None, border_top_color=None, border_top_h=Inches(0.06)):
    bg = rect(slide, l, t, w, h, fill=fill, radius=True)
    if border:
        bg.line.fill.solid()
        bg.line.color.rgb = border
        bg.line.width = Pt(1.5)
    if border_top_color:
        rect(slide, l + Inches(0.05), t, w - Inches(0.1), border_top_h, fill=border_top_color)
    return bg


def add_screenshot(slide, filename, l, t, w, h=None):
    path = os.path.join(SCREENSHOT_DIR, filename)
    if not os.path.exists(path):
        print(f"  WARNING: {path} not found, adding placeholder")
        c = card(slide, l, t, w, h or Inches(4), fill=LIGHT_GRAY, border=MID_GRAY)
        set_text(c, f"  [Screenshot: {filename}]", sz=14, color=MID_GRAY, align=PP_ALIGN.CENTER)
        return
    if h:
        slide.shapes.add_picture(path, l, t, w, h)
    else:
        slide.shapes.add_picture(path, l, t, width=w)


# ═══════════════════════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════════════════════

def slide_01_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)

    tbox(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(1.2),
         "Detecting the Undetectable", sz=48, bold=True, color=WHITE)

    rect(slide, Inches(0.8), Inches(2.5), Inches(3), Inches(0.04), fill=GOLD)

    tbox(slide, Inches(0.8), Inches(2.8), Inches(11.5), Inches(1.2),
         "How Nation-State Attackers Operate Undetected for Years\n"
         "Inside Critical Infrastructure — and How AI Changes the Outcome",
         sz=22, color=RGBColor(160, 200, 224))

    tbox(slide, Inches(0.8), Inches(4.3), Inches(11.5), Inches(0.5),
         "V-Intelligence UEBA: The Behavioral Digital Twin for Cyber Defense",
         sz=18, color=GOLD, bold=True)

    tbox(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.5),
         "22nd Century Technologies, Inc.  |  Confidential",
         sz=14, color=MID_GRAY)


def slide_02_threat_landscape(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "The Threat Landscape Has Changed")

    tbox(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.6),
         "Nation-state attackers no longer break in — they blend in. They use legitimate credentials, "
         "standard protocols, and built-in system tools. They look exactly like your employees.",
         sz=16, color=DARK_GRAY)

    stats = [
        ("287 days", "Average time to\ndetect a breach\n(IBM 2024)", RED),
        ("$4.88M", "Average cost of\na data breach\n(IBM 2024)", RED),
        ("5+ years", "Salt Typhoon in\nUS telecom before\ndiscovery", ORANGE),
        ("5+ years", "Volt Typhoon in\nUS critical infra\nbefore discovery", ORANGE),
    ]
    for i, (val, desc, clr) in enumerate(stats):
        x = Inches(0.6) + Inches(i * 3.1)
        card(slide, x, Inches(2.2), Inches(2.8), Inches(2.0), border_top_color=clr)
        tbox(slide, x + Inches(0.15), Inches(2.4), Inches(2.5), Inches(0.6),
             val, sz=36, bold=True, color=clr, align=PP_ALIGN.CENTER)
        tbox(slide, x + Inches(0.15), Inches(3.1), Inches(2.5), Inches(0.9),
             desc, sz=13, color=DARK_GRAY, align=PP_ALIGN.CENTER)

    box = rect(slide, Inches(0.6), Inches(4.6), Inches(12.1), Inches(2.5), fill=LIGHT_RED, radius=True)
    box.line.fill.solid()
    box.line.color.rgb = RED
    box.line.width = Pt(2)
    tbox(slide, Inches(0.9), Inches(4.7), Inches(11.5), Inches(0.4),
         "The New Reality", sz=20, bold=True, color=RED)

    lines = [
        "These are not hypothetical scenarios. In 2024, the US government disclosed that Chinese state-sponsored "
        "groups had been operating undetected inside American telecom and critical infrastructure for years.",
        "",
        "Traditional cybersecurity tools — SIEM rules, anomaly detection algorithms, threshold-based alerts — "
        "were deployed at every compromised organization. None of them detected the intrusion.",
        "",
        "The attackers were not invisible because they were sophisticated. They were invisible because "
        "the detection methods are fundamentally blind to their approach.",
    ]
    y = Inches(5.2)
    for line in lines:
        if not line:
            y += Inches(0.15)
            continue
        tbox(slide, Inches(0.9), y, Inches(11.5), Inches(0.4), line, sz=14, color=DARK_GRAY)
        y += Inches(0.35)

    footer(slide)


def slide_03_salt_typhoon(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Case Study: Salt Typhoon — 5 Years Inside US Telecom")

    card(slide, Inches(0.6), Inches(1.4), Inches(6), Inches(5.5), border_top_color=RGBColor(41, 128, 185))
    tbox(slide, Inches(0.9), Inches(1.6), Inches(5.4), Inches(0.4),
         "What Happened", sz=18, bold=True, color=RGBColor(41, 128, 185))

    events = [
        ("~2019", "Reconnaissance begins against US telecom providers"),
        ("~2021", "Active penetration of telecom infrastructure"),
        ("2021-2024", "Persistent access to AT&T, Verizon, T-Mobile, Lumen, Charter, Windstream"),
        ("2024", "Accessed CALEA lawful intercept / wiretap systems"),
        ("2024", "Harvested call metadata of 1M+ users, including senior US officials"),
        ("Sept 2024", "FBI and CISA discover the breach — not the telecom companies"),
        ("Oct 2024", "Public disclosure; 9 US telecoms confirmed compromised"),
    ]
    for i, (date, desc) in enumerate(events):
        y = Inches(2.1) + Inches(i * 0.6)
        tbox(slide, Inches(1.0), y, Inches(1.3), Inches(0.3),
             date, sz=12, bold=True, color=RGBColor(41, 128, 185))
        tbox(slide, Inches(2.4), y, Inches(3.8), Inches(0.5),
             desc, sz=12, color=DARK_GRAY)

    card(slide, Inches(6.9), Inches(1.4), Inches(5.8), Inches(5.5), border_top_color=RED)
    tbox(slide, Inches(7.2), Inches(1.6), Inches(5.2), Inches(0.4),
         "Why Every Detection System Failed", sz=18, bold=True, color=RED)

    reasons = [
        "Used legitimate network protocols — no malware signatures to detect",
        "Maintained normal traffic volumes — never spiked bandwidth alerts",
        "Operated with valid credentials — no brute force, no credential stuffing",
        "Exploited lawful intercept systems — outside standard SIEM monitoring",
        "CALEA systems managed separately from customer-facing security",
        "Every traditional SIEM rule, IDS signature, and anomaly threshold was in place",
        "None detected the intrusion — for over five years",
    ]
    for i, reason in enumerate(reasons):
        y = Inches(2.1) + Inches(i * 0.6)
        clr = RED if i >= 5 else DARK_GRAY
        bld = i >= 5
        tbox(slide, Inches(7.3), y, Inches(5.1), Inches(0.5),
             f"{'✗' if i >= 5 else '▸'} {reason}", sz=12, color=clr, bold=bld)

    footer(slide)


def slide_04_volt_typhoon(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Case Study: Volt Typhoon — 5 Years in Critical Infrastructure")

    card(slide, Inches(0.6), Inches(1.4), Inches(6), Inches(5.5), border_top_color=PURPLE)
    tbox(slide, Inches(0.9), Inches(1.6), Inches(5.4), Inches(0.4),
         "What Happened", sz=18, bold=True, color=PURPLE)

    events = [
        ("~2019", "Initial access to US critical infrastructure networks"),
        ("2019-2023", "Persistent presence in energy, water, telecom, transportation"),
        ("Ongoing", "Pre-positioned for disrupting US-Asia communications"),
        ("May 2023", "Microsoft Threat Intelligence publicly discloses the campaign"),
        ("May 2023", "CISA/NSA/FBI publish Joint Advisory AA23-144A"),
        ("Feb 2024", "CISA confirms 5 years of persistent access in some targets"),
        ("2024", "Follow-up Advisory AA24-038A documents full scope"),
    ]
    for i, (date, desc) in enumerate(events):
        y = Inches(2.1) + Inches(i * 0.6)
        tbox(slide, Inches(1.0), y, Inches(1.3), Inches(0.3),
             date, sz=12, bold=True, color=PURPLE)
        tbox(slide, Inches(2.4), y, Inches(3.8), Inches(0.5),
             desc, sz=12, color=DARK_GRAY)

    card(slide, Inches(6.9), Inches(1.4), Inches(5.8), Inches(5.5), border_top_color=RED)
    tbox(slide, Inches(7.2), Inches(1.6), Inches(5.2), Inches(0.4),
         "Why \"Living Off the Land\" Is Undetectable", sz=18, bold=True, color=RED)

    reasons = [
        "Zero malware deployed — used only built-in OS tools",
        "PowerShell, WMI, cmd, netsh — tools every admin uses daily",
        "Authenticated with stolen valid credentials",
        "Lateral movement looked identical to IT maintenance",
        "No anomalous traffic patterns — everything within normal bounds",
        "EDR, SIEM, IDS all deployed — all failed to detect",
        "Discovered by Microsoft threat researchers, not by deployed tools",
    ]
    for i, reason in enumerate(reasons):
        y = Inches(2.1) + Inches(i * 0.6)
        clr = RED if i >= 5 else DARK_GRAY
        bld = i >= 5
        tbox(slide, Inches(7.3), y, Inches(5.1), Inches(0.5),
             f"{'✗' if i >= 5 else '▸'} {reason}", sz=12, color=clr, bold=bld)

    box = rect(slide, Inches(0.6), Inches(6.0), Inches(12.1), Inches(0.8), fill=LIGHT_GOLD, radius=True)
    box.line.fill.solid()
    box.line.color.rgb = GOLD
    box.line.width = Pt(2)
    tbox(slide, Inches(0.9), Inches(6.1), Inches(11.5), Inches(0.6),
         "Common thread: Both Salt Typhoon and Volt Typhoon were discovered by external researchers "
         "— not by any security tool deployed at the compromised organizations. The tools were there. "
         "They simply cannot see this type of attack.",
         sz=14, bold=True, color=DARK_GRAY)

    footer(slide)


def slide_04b_threat_spectrum(prs):
    """The Full Threat Spectrum — 2x2 grid of attack categories."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    title_bar(slide, "The Full Threat Spectrum")

    tbox(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
         "Nation-state actors, insiders, and patient adversaries each exploit a different blind spot. "
         "No organization faces just one threat category.",
         sz=16, color=RGBColor(160, 200, 224))

    # 2x2 grid of attack category cards
    categories = [
        # (title, subtitle, details, border_color, col, row)
        ("Nation-State Telecom", "Salt Typhoon",
         "5+ years undetected  |  Wiretap compromise\n9 US telecom providers  |  Call metadata of senior officials",
         RED, 0, 0),
        ("Nation-State Infrastructure", "Volt Typhoon",
         "5+ years undetected  |  Critical infrastructure\nLiving Off the Land  |  Pre-positioned for disruption",
         ORANGE, 1, 0),
        ("Insider Threat", "Trusted Employee",
         "8-month gradual escalation  |  Authorized credentials\nAll actions individually legitimate  |  Slow data exfiltration",
         RED, 0, 1),
        ("Slow APT Campaign", "Patient Attacker",
         "180-day campaign  |  Below every threshold\nNo single anomalous week  |  Credential harvesting + lateral movement",
         ORANGE, 1, 1),
    ]

    for title, subtitle, details, border_clr, col, row in categories:
        x = Inches(0.6) + Inches(col * 6.2)
        y = Inches(2.1) + Inches(row * 2.3)
        w = Inches(5.8)
        h = Inches(2.0)

        # Card background
        bg = rect(slide, x, y, w, h, fill=DARK_TEAL, radius=True)
        # Colored left border accent
        rect(slide, x + Inches(0.05), y + Inches(0.1), Inches(0.07), h - Inches(0.2), fill=border_clr)

        # Title + subtitle
        tbox(slide, x + Inches(0.3), y + Inches(0.15), w - Inches(0.5), Inches(0.35),
             title, sz=18, bold=True, color=WHITE)
        tbox(slide, x + Inches(0.3), y + Inches(0.5), w - Inches(0.5), Inches(0.25),
             subtitle, sz=12, bold=True, color=border_clr)

        # Details
        tbox(slide, x + Inches(0.3), y + Inches(0.85), w - Inches(0.5), Inches(0.9),
             details, sz=12, color=RGBColor(200, 210, 220))

    # Bottom bar
    box = rect(slide, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.6), fill=DARK_TEAL, radius=True)
    box.line.fill.solid()
    box.line.color.rgb = GOLD
    box.line.width = Pt(2)
    tbox(slide, Inches(0.9), Inches(6.65), Inches(11.5), Inches(0.45),
         "Each exploits a different blind spot. No single tool covers all four.",
         sz=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    footer(slide)


def slide_05_why_trad_fails(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Why Every Traditional Approach Fails")

    tbox(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
         "Traditional cybersecurity tools are designed for a threat model that no longer exists.",
         sz=16, color=DARK_GRAY)

    failures = [
        ("Signature-Based Detection", "IDS / Antivirus / EDR",
         "Looks for known malware patterns and attack signatures.",
         "Fails when: Attacker uses no malware — only legitimate tools that are already whitelisted.",
         "Salt Typhoon and Volt Typhoon deployed zero malware. Every file on disk was legitimate.",
         RED),
        ("Rule-Based Detection", "SIEM Rules / Correlation",
         "Triggers on predefined conditions — failed logins, privilege escalation, data thresholds.",
         "Fails when: Attacker stays below every threshold. 5% more data each week never triggers a rule.",
         "Every rule at AT&T, Verizon, and 7 other telecoms was in place. None fired for 5+ years.",
         ORANGE),
        ("Statistical Anomaly", "Isolation Forest / SVM / LOF",
         "Flags users whose individual metrics deviate from the statistical mean.",
         "Fails when: No single metric is abnormal enough. The attack is in the combination, not any one number.",
         "Our testing: 4 simulated nation-state attacks — Isolation Forest, SVM, and LOF detect zero.",
         PURPLE),
    ]

    for i, (name, tools, what, fails, proof, clr) in enumerate(failures):
        x = Inches(0.4) + Inches(i * 4.2)
        card(slide, x, Inches(2.0), Inches(3.9), Inches(4.8), border_top_color=clr)
        tbox(slide, x + Inches(0.15), Inches(2.2), Inches(3.6), Inches(0.35),
             name, sz=15, bold=True, color=clr, align=PP_ALIGN.CENTER)
        tbox(slide, x + Inches(0.15), Inches(2.55), Inches(3.6), Inches(0.3),
             tools, sz=11, color=MID_GRAY, align=PP_ALIGN.CENTER)

        tbox(slide, x + Inches(0.2), Inches(3.0), Inches(3.5), Inches(0.7),
             what, sz=12, color=DARK_GRAY)

        fail_box = rect(slide, x + Inches(0.15), Inches(3.8), Inches(3.6), Inches(0.9),
                        fill=LIGHT_RED, radius=True)
        tbox(slide, x + Inches(0.25), Inches(3.85), Inches(3.4), Inches(0.8),
             fails, sz=11, bold=True, color=RED)

        tbox(slide, x + Inches(0.2), Inches(4.9), Inches(3.5), Inches(0.8),
             proof, sz=11, color=DARK_GRAY)

    box = rect(slide, Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.7), fill=NAVY, radius=True)
    tbox(slide, Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.55),
         "The fundamental problem: These tools ask \"is any one number abnormal?\" — "
         "but modern attackers ensure that no single number ever is. "
         "The attack is invisible to every tool that looks at individual metrics.",
         sz=14, bold=True, color=GOLD)

    footer(slide)


def slide_06_digital_twin(prs):
    """NEW: The Digital Twin Innovation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "The Innovation: Behavioral Digital Twin")

    tbox(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.7),
         "V-Intelligence UEBA creates a behavioral digital twin for every entity in your environment — "
         "a rich, multi-dimensional representation of how each user, device, and application behaves. "
         "Not what numbers they generate. How they work.",
         sz=16, color=DARK_GRAY)

    # Left: Traditional approach
    card(slide, Inches(0.6), Inches(2.3), Inches(5.8), Inches(4.0), border_top_color=RED)
    tbox(slide, Inches(0.9), Inches(2.5), Inches(5.2), Inches(0.35),
         "Traditional: Raw Metrics", sz=18, bold=True, color=RED)

    trad_items = [
        "Monitors individual numbers: login count, bytes transferred, failed auths",
        "Each metric checked independently against static thresholds",
        "No understanding of relationships between behaviors",
        "No context: \"50 logins\" means the same whether you're IT or Finance",
        "Attacker stays below threshold on every metric → invisible",
    ]
    for i, item in enumerate(trad_items):
        tbox(slide, Inches(0.9), Inches(3.0) + Inches(i * 0.55), Inches(5.2), Inches(0.5),
             f"✗  {item}", sz=12, color=DARK_GRAY)

    # Right: ACECARD approach
    card(slide, Inches(6.9), Inches(2.3), Inches(5.8), Inches(4.0), border_top_color=GREEN)
    tbox(slide, Inches(7.2), Inches(2.5), Inches(5.2), Inches(0.35),
         "V-Intelligence UEBA: Behavioral Digital Twin", sz=18, bold=True, color=GREEN)

    ace_items = [
        "Creates a unified behavioral profile capturing all activity dimensions",
        "Understands the meaning of actions — not just the magnitude",
        "Captures cross-dimensional relationships (access + data + network)",
        "Context-aware: knows what \"normal\" means for each individual role",
        "Detects behavioral change patterns invisible to metric-by-metric analysis",
    ]
    for i, item in enumerate(ace_items):
        tbox(slide, Inches(7.2), Inches(3.0) + Inches(i * 0.55), Inches(5.2), Inches(0.5),
             f"✓  {item}", sz=12, color=DARK_GRAY)

    # Bottom: Key insight
    box = rect(slide, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.7), fill=NAVY, radius=True)
    tbox(slide, Inches(0.9), Inches(6.55), Inches(11.5), Inches(0.55),
         "The digital twin doesn't ask \"is this number abnormal?\" — it asks "
         "\"has this person fundamentally changed how they work?\" "
         "That is the question Salt Typhoon and Volt Typhoon attackers cannot evade.",
         sz=14, bold=True, color=GOLD)

    footer(slide)


def slide_07_multi_tier(prs):
    """NEW: Multi-Tier Detection Approach."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Multi-Tier Detection: Progressive Intelligence")

    tbox(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
         "V-Intelligence UEBA employs a multi-tier approach — each layer adds detection capability "
         "that the previous layer cannot achieve.",
         sz=16, color=DARK_GRAY)

    # Tier cards — three columns showing progression
    tiers = [
        ("Traditional SIEM", "Isolation Forest / SVM / LOF",
         RED, "0 of 4", "0%",
         [
             "Threshold-based anomaly detection",
             "Checks individual metrics independently",
             "Misses attacks where no single metric spikes",
             "Every deployed tool uses this approach",
         ],
         "Result: Catches nothing. Every attacker stays below every threshold."),

        ("Intermediate Analysis", "Statistical Pattern Detection",
         ORANGE, "1 of 4", "25%",
         [
             "Multi-dimensional statistical analysis",
             "Cross-correlates feature combinations",
             "Catches louder anomalies that cross statistical bounds",
             "Detects Volt Typhoon (noisier LotL campaign)",
         ],
         "Result: Catches Volt Typhoon only. 3 stealthier attacks remain invisible."),

        ("V-INTELLIGENCE UEBA", "Threat-Profile Detector",
         GREEN, "4 of 4", "100%",
         [
             "Known-bad profile matching per entity",
             "Cohort-relative, raw-event, label-free",
             "C2-beacon, DGA-DNS, LOTL-process, recon-fanout",
             "Catches Salt Typhoon, Insider, Slow APT, Volt Typhoon",
         ],
         "Result: All 4 attacks detected at 0 false positives."),
    ]

    for i, (name, tools, clr, score, pct, bullets, result) in enumerate(tiers):
        x = Inches(0.4) + Inches(i * 4.2)
        card(slide, x, Inches(2.0), Inches(3.9), Inches(5.0), border_top_color=clr)

        tbox(slide, x + Inches(0.15), Inches(2.2), Inches(3.6), Inches(0.35),
             name, sz=15, bold=True, color=clr, align=PP_ALIGN.CENTER)
        tbox(slide, x + Inches(0.15), Inches(2.55), Inches(3.6), Inches(0.25),
             tools, sz=10, color=MID_GRAY, align=PP_ALIGN.CENTER)

        # Detection score badge
        badge = rect(slide, x + Inches(0.8), Inches(2.9), Inches(2.3), Inches(0.6), fill=clr, radius=True)
        set_text(badge, f"  {score} attacks detected", sz=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        for j, bullet in enumerate(bullets):
            tbox(slide, x + Inches(0.2), Inches(3.65) + Inches(j * 0.4), Inches(3.5), Inches(0.35),
                 f"▸ {bullet}", sz=11, color=DARK_GRAY)

        # Result bar
        result_bg = rect(slide, x + Inches(0.15), Inches(5.35), Inches(3.6), Inches(0.55),
                         fill=LIGHT_GRAY, radius=True)
        tbox(slide, x + Inches(0.25), Inches(5.4), Inches(3.4), Inches(0.45),
             result, sz=10, bold=True, color=clr)

    # Arrow progression at bottom
    box = rect(slide, Inches(0.6), Inches(6.35), Inches(12.1), Inches(0.65), fill=NAVY, radius=True)
    tbox(slide, Inches(0.9), Inches(6.38), Inches(11.5), Inches(0.55),
         "0 / 4  →  1 / 4  →  4 / 4    |    Each tier catches what the previous tier cannot see. "
         "The behavioral digital twin is the breakthrough layer.",
         sz=15, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    footer(slide)


def slide_08_how_acecard_works(prs):
    """NEW: How V-Intelligence UEBA Works (high-level, no proprietary details)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "How V-Intelligence UEBA Works — High-Level Architecture")

    tbox(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
         "V-Intelligence UEBA transforms raw security telemetry into behavioral intelligence "
         "through a multi-stage pipeline:",
         sz=16, color=DARK_GRAY)

    stages = [
        ("1", "INGEST", TEAL,
         "Security Telemetry",
         "Authentication logs, file access, network traffic, endpoint events, DNS queries — "
         "the same data your SIEM already collects."),
        ("2", "PROFILE", BLUE,
         "Behavioral Digital Twin",
         "AI creates a multi-dimensional behavioral profile for every user, device, and application. "
         "Captures HOW entities behave, not just WHAT numbers they generate."),
        ("3", "MONITOR", PURPLE,
         "Behavioral Change Detection",
         "Continuously monitors each entity's digital twin for meaningful behavioral changes. "
         "Detects when behavior shifts — even when no individual metric is abnormal."),
        ("4", "SCORE", GREEN,
         "Composite Risk Scoring",
         "Multi-phase scoring evaluates signal strength, breadth, sustained deviation, "
         "and behavioral context to produce a single composite risk score per entity."),
        ("5", "RANK", GOLD,
         "Analyst Priority List",
         "All entities ranked by composite score. Analysts start at #1 and work down. "
         "No threshold tuning. No alert fatigue. No false alarm floods."),
    ]

    for i, (num, label, clr, title, desc) in enumerate(stages):
        y = Inches(2.0) + Inches(i * 0.85)
        # Number circle
        circle = rect(slide, Inches(0.7), y, Inches(0.5), Inches(0.5), fill=clr, radius=True)
        set_text(circle, num, sz=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Label
        tbox(slide, Inches(1.4), y + Inches(0.02), Inches(1.8), Inches(0.4),
             label, sz=16, bold=True, color=clr)

        # Title
        tbox(slide, Inches(3.3), y + Inches(0.02), Inches(2.5), Inches(0.4),
             title, sz=14, bold=True, color=NAVY)

        # Description
        tbox(slide, Inches(6.0), y + Inches(0.0), Inches(6.5), Inches(0.9),
             desc, sz=12, color=DARK_GRAY)

    # Bottom note
    box = rect(slide, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.5), fill=LIGHT_GOLD, radius=True)
    box.line.fill.solid()
    box.line.color.rgb = GOLD
    box.line.width = Pt(1)
    tbox(slide, Inches(0.9), Inches(6.53), Inches(11.5), Inches(0.4),
         "No new data sources required — V-Intelligence UEBA works with the telemetry your environment already produces.",
         sz=14, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)

    footer(slide)


def slide_09_screenshot_heatmap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Detection Results: Who Catches What?")

    tbox(slide, Inches(0.6), Inches(1.15), Inches(12), Inches(0.4),
         "250 users, 70 weeks, 4 attack campaigns — traditional methods vs V-Intelligence UEBA:",
         sz=15, color=DARK_GRAY)

    # --- Section 1: Traditional ---
    rect(slide, Inches(0.5), Inches(1.65), Inches(5.5), Inches(0.06), fill=RED)
    tbox(slide, Inches(0.5), Inches(1.75), Inches(3), Inches(0.3),
         "TRADITIONAL DETECTION", sz=14, bold=True, color=RED)
    tbox(slide, Inches(3.5), Inches(1.78), Inches(4), Inches(0.25),
         "Isolation Forest  /  One-Class SVM  /  LOF", sz=10, color=DARK_GRAY)

    # Column headers
    cols = ["USR-156\nInsider", "USR-234\nAPT", "USR-042\nVolt Typhoon", "USR-118\nSalt Typhoon"]
    col_x = [Inches(3.5), Inches(5.5), Inches(7.5), Inches(9.5)]
    for ci, (cx, col) in enumerate(zip(col_x, cols)):
        tbox(slide, cx, Inches(2.1), Inches(1.8), Inches(0.4), col, sz=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    tbox(slide, Inches(11.5), Inches(2.15), Inches(1.2), Inches(0.3),
         "FP Rate", sz=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    methods = [
        ("Isolation Forest", ["MISSED", "MISSED", "MISSED", "MISSED"], "5.3%"),
        ("One-Class SVM",    ["MISSED", "MISSED", "MISSED", "MISSED"], "14.6%"),
        ("LOF",              ["MISSED", "MISSED", "MISSED", "MISSED"], "4.5%"),
    ]
    for ri, (mname, results, fp) in enumerate(methods):
        y = Inches(2.55) + Inches(ri * 0.35)
        tbox(slide, Inches(0.7), y, Inches(2.5), Inches(0.3), mname, sz=11, color=NAVY)
        for ci, r in enumerate(results):
            clr = GREEN if r == "DETECTED" else RED
            tbox(slide, col_x[ci], y, Inches(1.8), Inches(0.3), r, sz=10, bold=True, color=clr, align=PP_ALIGN.CENTER)
        tbox(slide, Inches(11.5), y, Inches(1.2), Inches(0.3), fp, sz=10, color=DARK_GRAY, align=PP_ALIGN.CENTER)

    badge_trad = rect(slide, Inches(11), Inches(1.75), Inches(1.8), Inches(0.35), fill=RED, radius=True)
    tbox(slide, Inches(11.05), Inches(1.77), Inches(1.7), Inches(0.3),
         "0 of 4 detected", sz=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # --- Section 2: Intermediate ---
    y_int = Inches(3.7)
    rect(slide, Inches(0.5), y_int, Inches(5.5), Inches(0.06), fill=ORANGE)
    tbox(slide, Inches(0.5), y_int + Inches(0.1), Inches(3.5), Inches(0.3),
         "OUR INTERMEDIATE ANALYSIS", sz=14, bold=True, color=ORANGE)
    tbox(slide, Inches(4.0), y_int + Inches(0.13), Inches(4), Inches(0.25),
         "Statistical Pattern Detection (Z-Score |z|>3)", sz=10, color=DARK_GRAY)

    y_zrow = y_int + Inches(0.45)
    tbox(slide, Inches(0.7), y_zrow, Inches(2.5), Inches(0.3), "Z-Score (|z| > 3)", sz=11, color=NAVY)
    z_results = ["MISSED", "MISSED", "DETECTED", "MISSED"]
    for ci, r in enumerate(z_results):
        clr = GREEN if r == "DETECTED" else RED
        tbox(slide, col_x[ci], y_zrow, Inches(1.8), Inches(0.3), r, sz=10, bold=True, color=clr, align=PP_ALIGN.CENTER)
    tbox(slide, Inches(11.5), y_zrow, Inches(1.2), Inches(0.3), "9.8%", sz=10, color=DARK_GRAY, align=PP_ALIGN.CENTER)

    badge_int = rect(slide, Inches(11), y_int + Inches(0.1), Inches(1.8), Inches(0.35), fill=ORANGE, radius=True)
    tbox(slide, Inches(11.05), y_int + Inches(0.12), Inches(1.7), Inches(0.3),
         "1 of 4 detected", sz=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    tbox(slide, Inches(0.7), y_zrow + Inches(0.35), Inches(8), Inches(0.25),
         "Catches Volt Typhoon only — Insider, Slow APT, and Salt Typhoon remain invisible",
         sz=10, color=ORANGE)

    # --- Section 3: V-Intelligence UEBA ---
    y_ace = Inches(5.0)
    rect(slide, Inches(0.5), y_ace, Inches(5.5), Inches(0.06), fill=GREEN)
    tbox(slide, Inches(0.5), y_ace + Inches(0.1), Inches(5), Inches(0.3),
         "V-INTELLIGENCE UEBA + THREAT-PROFILE DETECTOR", sz=14, bold=True, color=GREEN)
    tbox(slide, Inches(5.5), y_ace + Inches(0.13), Inches(5), Inches(0.25),
         "Known-Bad Profile Matching → Multi-Front Detection", sz=10, color=DARK_GRAY)

    y_arow = y_ace + Inches(0.45)
    tbox(slide, Inches(0.7), y_arow, Inches(2.5), Inches(0.3), "Threat-Profile Match", sz=11, color=NAVY)
    for ci in range(4):
        tbox(slide, col_x[ci], y_arow, Inches(1.8), Inches(0.3), "DETECTED", sz=10, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    tbox(slide, Inches(11.5), y_arow, Inches(1.2), Inches(0.3), "0%", sz=10, color=DARK_GRAY, align=PP_ALIGN.CENTER)

    badge_ace = rect(slide, Inches(11), y_ace + Inches(0.1), Inches(1.8), Inches(0.35), fill=GREEN, radius=True)
    tbox(slide, Inches(11.05), y_ace + Inches(0.12), Inches(1.7), Inches(0.3),
         "4 of 4 detected", sz=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # --- Bottom summary bar ---
    y_bot = Inches(5.9)
    summary_bar = rect(slide, Inches(0.5), y_bot, Inches(12.3), Inches(0.55), fill=NAVY, radius=True)
    tbox(slide, Inches(0.7), y_bot + Inches(0.05), Inches(3.5), Inches(0.25),
         "Traditional: 0 / 4", sz=18, bold=True, color=RED, align=PP_ALIGN.CENTER)
    tbox(slide, Inches(4.2), y_bot + Inches(0.1), Inches(0.5), Inches(0.2),
         "→", sz=16, color=RGBColor(160, 200, 224))
    tbox(slide, Inches(4.7), y_bot + Inches(0.05), Inches(3.5), Inches(0.25),
         "Intermediate: 1 / 4", sz=18, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    tbox(slide, Inches(8.2), y_bot + Inches(0.1), Inches(0.5), Inches(0.2),
         "→", sz=16, color=RGBColor(160, 200, 224))
    tbox(slide, Inches(8.7), y_bot + Inches(0.05), Inches(3.8), Inches(0.25),
         "V-Intelligence UEBA: 4 / 4", sz=18, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    tbox(slide, Inches(0.7), y_bot + Inches(0.32), Inches(11.7), Inches(0.2),
         "Each tier catches what the previous tier cannot see.", sz=10, color=RGBColor(160, 200, 224), align=PP_ALIGN.CENTER)

    footer(slide)


def slide_10_screenshot_scores(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "V-Intelligence UEBA: Composite Cleanly Separates 2 of 4 Attackers")

    tbox(slide, Inches(0.6), Inches(1.25), Inches(12), Inches(0.4),
         "Same 250 users, same data. The composite score lifts USR-156 and USR-118 above all normal users; "
         "USR-234 and USR-042 stay below normal and are caught by the threat-profile detector:",
         sz=15, color=DARK_GRAY)

    add_screenshot(slide, "tier3_sec_01.png", Inches(0.8), Inches(1.75), Inches(11.7), Inches(5.0))

    footer(slide)


def slide_11_screenshot_separation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "How V-Intelligence UEBA Separates Attackers from Normal Users")

    tbox(slide, Inches(0.6), Inches(1.25), Inches(12), Inches(0.4),
         "The composite score lifts the insider (USR-156) and Salt Typhoon (USR-118) clearly above the normal "
         "population; the threat-profile detector closes the gap on the two stealthiest campaigns:",
         sz=15, color=DARK_GRAY)

    add_screenshot(slide, "tier3_sec_03.png", Inches(0.6), Inches(1.75), Inches(12.1), Inches(5.0))

    footer(slide)


def slide_12_screenshot_verdict(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "The Verdict: Traditional 0/4 — Intermediate 1/4 — V-Intelligence UEBA 4/4")

    tbox(slide, Inches(0.6), Inches(1.15), Inches(12), Inches(0.4),
         "Same data, same users — three fundamentally different outcomes:",
         sz=15, color=DARK_GRAY)

    # --- Three verdict cards ---
    card_w = Inches(3.85)
    gap = Inches(0.25)
    y_top = Inches(1.7)
    card_h = Inches(3.6)

    tiers = [
        ("TRADITIONAL DETECTION", "Isolation Forest / SVM / LOF",
         "0 of 4", RED,
         "Fixed thresholds on 23 scalar features. Attackers who stay within "
         "normal statistical ranges are invisible — no individual metric is "
         "abnormal enough to cross a detection boundary.",
         "Best FP: 4.5% (LOF) — but detects nothing"),
        ("INTERMEDIATE ANALYSIS", "Z-Score (|z|>3) — Our Statistical Layer",
         "1 of 4", ORANGE,
         "Detects single-feature spikes beyond 3 standard deviations. Catches "
         "the most aggressive attacker (Volt Typhoon) but misses slow, distributed "
         "campaigns that stay below the threshold on every individual metric.",
         "FP: 9.8% — catches Volt Typhoon only"),
        ("V-INTELLIGENCE UEBA", "Threat-Profile Detector (multi-front)",
         "4 of 4", GREEN,
         "Matches raw events against known-bad profiles — C2-beacon, DGA-DNS, LOTL-process, "
         "cohort-rare access, recon-fanout, insider-collection — cohort-relative and label-free. "
         "Composite scoring separates 2 of 4 on its own; the profile detector catches all four.",
         "FP: 0% — all 4 attacks detected"),
    ]

    for i, (name, subtitle, score, clr, desc, metric) in enumerate(tiers):
        x = Inches(0.55) + i * (card_w + gap)
        card(slide, x, y_top, card_w, card_h, border_top_color=clr)

        tbox(slide, x + Inches(0.15), y_top + Inches(0.12), card_w - Inches(0.3), Inches(0.35),
             name, sz=14, bold=True, color=clr, align=PP_ALIGN.CENTER)
        tbox(slide, x + Inches(0.15), y_top + Inches(0.42), card_w - Inches(0.3), Inches(0.3),
             subtitle, sz=10, color=DARK_GRAY, align=PP_ALIGN.CENTER)

        tbox(slide, x + Inches(0.15), y_top + Inches(0.75), card_w - Inches(0.3), Inches(0.7),
             score, sz=44, bold=True, color=clr, align=PP_ALIGN.CENTER)
        tbox(slide, x + Inches(0.15), y_top + Inches(1.4), card_w - Inches(0.3), Inches(0.25),
             "attacks detected", sz=11, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)

        tbox(slide, x + Inches(0.2), y_top + Inches(1.75), card_w - Inches(0.4), Inches(1.0),
             desc, sz=10, color=DARK_GRAY, align=PP_ALIGN.LEFT)

        metric_bg = rect(slide, x + Inches(0.15), y_top + Inches(2.85), card_w - Inches(0.3),
                         Inches(0.45), fill=clr, radius=True)
        m_clr = WHITE if clr != ORANGE else NAVY
        tbox(slide, x + Inches(0.25), y_top + Inches(2.9), card_w - Inches(0.5), Inches(0.35),
             metric, sz=10, bold=True, color=m_clr, align=PP_ALIGN.CENTER)

    # Bottom summary bar
    summary_bar = rect(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.65), fill=NAVY, radius=True)
    tbox(slide, Inches(0.8), Inches(5.55), Inches(11.7), Inches(0.25),
         "V-Intelligence UEBA builds the behavioral profile. The threat-profile detector catches the attack.",
         sz=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tbox(slide, Inches(0.8), Inches(5.82), Inches(11.7), Inches(0.25),
         "4/4 attacks detected at 0 FP — vs 0/4 traditional, 1/4 intermediate.",
         sz=11, bold=True, color=RGBColor(160, 200, 224), align=PP_ALIGN.CENTER)

    # Salt Typhoon callout
    salt_box = rect(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.6), fill=DARK_TEAL, radius=True)
    tbox(slide, Inches(0.8), Inches(6.32), Inches(11.7), Inches(0.22),
         "Real-World Validation: Salt Typhoon (USR-118)", sz=12, bold=True, color=GREEN)
    tbox(slide, Inches(0.8), Inches(6.55), Inches(11.7), Inches(0.3),
         "5+ years undetected by every traditional tool. V-Intelligence UEBA ranks it #1 out of 250 users (score: 51.3).",
         sz=10, color=RGBColor(160, 200, 224))

    footer(slide)


def slide_13_unique_results(prs):
    """NEW: Unique Results Highlight."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Unique Results: What V-Intelligence UEBA Achieved")

    tbox(slide, Inches(0.6), Inches(1.25), Inches(12), Inches(0.4),
         "Four distinct attack campaigns. Composite cleanly separates 2; the threat-profile detector catches all 4.",
         sz=16, color=DARK_GRAY)

    # Four attacker result cards
    attackers = [
        ("USR-118", "Salt Typhoon Telecom", "Telecom infrastructure pivot — router config exfiltration, "
         "call metadata harvesting, DNS tunneling",
         "51.3", "#1 / 250", "100th", GREEN,
         "Mirrors real Salt Typhoon — 5+ years undetected worldwide"),
        ("USR-156", "Insider Threat", "Gradual privilege escalation — restricted file access, "
         "off-hours activity, slow data exfiltration",
         "46.2", "#2 / 250", "99th", TEAL,
         "Slow 8-month escalation — no single week anomalous"),
        ("USR-234", "Slow APT", "Low-and-slow persistent access — credential harvesting, "
         "lateral movement across network segments",
         "C2-beacon", "below norm", "profile", BLUE,
         "Composite ranks below normal — caught by C2-beacon profile front"),
        ("USR-042", "Volt Typhoon LOTL", "Living-off-the-land — PowerShell, WMI, legitimate tools, "
         "no malware deployed",
         "LOTL", "below norm", "profile", PURPLE,
         "Composite ranks below normal — caught by LOTL-process profile front"),
    ]

    for i, (uid, name, desc, score, rank, pctile, clr, note) in enumerate(attackers):
        y = Inches(1.8) + Inches(i * 1.25)
        card(slide, Inches(0.5), y, Inches(12.3), Inches(1.1), border_top_color=clr)

        # UID + Name
        tbox(slide, Inches(0.7), y + Inches(0.1), Inches(2.5), Inches(0.35),
             f"{uid} — {name}", sz=14, bold=True, color=clr)
        tbox(slide, Inches(0.7), y + Inches(0.45), Inches(4.5), Inches(0.5),
             desc, sz=10, color=DARK_GRAY)

        # Score
        tbox(slide, Inches(5.5), y + Inches(0.15), Inches(1.3), Inches(0.35),
             score, sz=28, bold=True, color=clr, align=PP_ALIGN.CENTER)
        tbox(slide, Inches(5.5), y + Inches(0.55), Inches(1.3), Inches(0.25),
             "COMPOSITE", sz=8, color=MID_GRAY, align=PP_ALIGN.CENTER)

        # Rank
        tbox(slide, Inches(7.0), y + Inches(0.15), Inches(1.5), Inches(0.35),
             rank, sz=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        tbox(slide, Inches(7.0), y + Inches(0.55), Inches(1.5), Inches(0.25),
             f"{pctile} percentile", sz=9, color=MID_GRAY, align=PP_ALIGN.CENTER)

        # Note
        tbox(slide, Inches(8.8), y + Inches(0.2), Inches(3.8), Inches(0.7),
             note, sz=11, bold=True, color=clr)

    # Bottom bar
    box = rect(slide, Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.55), fill=NAVY, radius=True)
    tbox(slide, Inches(0.9), Inches(6.88), Inches(11.5), Inches(0.45),
         "All 4 attackers detected  |  0 false positives (threat-profile detector)  |  "
         "Zero threshold tuning  |  Traditional methods: 0 of 4 detected",
         sz=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    footer(slide)


def slide_14_salt_typhoon_proof(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    title_bar(slide, "Real-World Validation: Salt Typhoon")

    tbox(slide, Inches(0.6), Inches(1.4), Inches(12), Inches(0.5),
         "Our simulation replicates Salt Typhoon's exact attack pattern. The results confirm reality.",
         sz=18, color=RGBColor(160, 200, 224))

    card(slide, Inches(0.6), Inches(2.2), Inches(5.8), Inches(2.5))
    tbox(slide, Inches(0.9), Inches(2.4), Inches(5.2), Inches(0.35),
         "In the Real World", sz=16, bold=True, color=RED)
    real_lines = [
        "Salt Typhoon operated inside US telecom for 5+ years",
        "Accessed wiretap systems of senior government officials",
        "9 telecom providers compromised (AT&T, Verizon, T-Mobile...)",
        "Discovered by FBI/CISA — not by any deployed tool",
        "Every SIEM rule, IDS, and anomaly algorithm was in place",
    ]
    for i, line in enumerate(real_lines):
        tbox(slide, Inches(0.9), Inches(2.85) + Inches(i * 0.3), Inches(5.2), Inches(0.3),
             f"▸ {line}", sz=12, color=DARK_GRAY)

    card(slide, Inches(6.9), Inches(2.2), Inches(5.8), Inches(2.5))
    tbox(slide, Inches(7.2), Inches(2.4), Inches(5.2), Inches(0.35),
         "In Our Simulation", sz=16, bold=True, color=GREEN)
    sim_lines = [
        "Same attack pattern simulated among 250 users",
        "All traditional algorithms: MISSED (max z-score = 1.71)",
        "V-Intelligence UEBA Composite Score: 51.3",
        "V-Intelligence UEBA Rank: #1 out of 250 users (100th percentile)",
        "Strongest behavioral anomaly in the entire population",
    ]
    for i, line in enumerate(sim_lines):
        clr = GREEN if i >= 2 else DARK_GRAY
        bld = i >= 3
        tbox(slide, Inches(7.2), Inches(2.85) + Inches(i * 0.3), Inches(5.2), Inches(0.3),
             f"{'★' if i >= 3 else '▸'} {line}", sz=12, color=clr, bold=bld)

    box = rect(slide, Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.2), fill=DARK_TEAL, radius=True)
    box.line.fill.solid()
    box.line.color.rgb = GOLD
    box.line.width = Pt(2)
    tbox(slide, Inches(0.9), Inches(5.1), Inches(11.5), Inches(0.35),
         "The Bottom Line", sz=16, bold=True, color=GOLD)
    tbox(slide, Inches(0.9), Inches(5.5), Inches(11.5), Inches(0.6),
         "The same threshold-based detection that failed for 5+ years in the real world fails identically "
         "in our simulation. V-Intelligence UEBA's behavioral digital twin ranks the identical attack pattern "
         "as the #1 anomaly — not because any single number is abnormal, but because the digital twin "
         "understands that the overall behavior has fundamentally changed.",
         sz=14, color=WHITE)

    results = [
        ("5+ years", "undetected by\ntraditional tools"),
        ("#1 / 250", "V-Intelligence UEBA rank\n(highest anomaly)"),
        ("51.3", "composite\nbehavioral score"),
        ("0 of 4", "traditional methods\nthat detected it"),
    ]
    for i, (val, desc) in enumerate(results):
        x = Inches(0.8) + Inches(i * 3.1)
        tbox(slide, x, Inches(6.25), Inches(2.8), Inches(0.4),
             val, sz=24, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        tbox(slide, x, Inches(6.6), Inches(2.8), Inches(0.5),
             desc, sz=11, color=RGBColor(160, 200, 224), align=PP_ALIGN.CENTER)

    footer(slide, "22nd Century Technologies, Inc.  |  V-Intelligence UEBA  |  Confidential")


def slide_15_competitive_advantage(prs):
    """NEW: Competitive Advantage."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Why V-Intelligence UEBA Succeeds Where Others Fail")

    tbox(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
         "Every existing approach has a structural blind spot. "
         "V-Intelligence UEBA's behavioral digital twin eliminates it.",
         sz=16, color=DARK_GRAY)

    # Comparison rows
    comparisons = [
        ("Traditional SIEM", "Checks thresholds on individual metrics",
         "Misses attacks where no metric exceeds its threshold",
         "Understands behavior holistically — no threshold to stay below"),
        ("EDR / XDR", "Monitors endpoint activity for known patterns",
         "Misses credential-based attacks using legitimate tools",
         "Profiles the entity, not the tool — detects behavioral change regardless of tooling"),
        ("Network Detection (NDR)", "Inspects traffic for signatures and anomalies",
         "Misses attacks using legitimate protocols at normal volumes",
         "Captures network behavior as one dimension of a multi-dimensional behavioral twin"),
        ("Existing UEBA", "Statistical deviation on aggregated user metrics",
         "Fails when no single metric deviates enough — the \"combination\" problem",
         "Digital twin captures cross-dimensional behavioral patterns that scalar metrics miss"),
    ]

    for i, (approach, what, blind_spot, advantage) in enumerate(comparisons):
        y = Inches(2.0) + Inches(i * 1.2)
        # Approach name
        card(slide, Inches(0.5), y, Inches(2.5), Inches(1.0), border_top_color=RED)
        tbox(slide, Inches(0.6), y + Inches(0.15), Inches(2.3), Inches(0.3),
             approach, sz=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        tbox(slide, Inches(0.6), y + Inches(0.45), Inches(2.3), Inches(0.45),
             what, sz=9, color=MID_GRAY, align=PP_ALIGN.CENTER)

        # Blind spot
        card(slide, Inches(3.2), y, Inches(4.0), Inches(1.0), fill=LIGHT_RED)
        tbox(slide, Inches(3.35), y + Inches(0.05), Inches(3.7), Inches(0.2),
             "Blind Spot", sz=9, bold=True, color=RED)
        tbox(slide, Inches(3.35), y + Inches(0.3), Inches(3.7), Inches(0.6),
             blind_spot, sz=11, color=DARK_GRAY)

        # ACECARD advantage
        card(slide, Inches(7.4), y, Inches(5.4), Inches(1.0), fill=LIGHT_GREEN)
        tbox(slide, Inches(7.55), y + Inches(0.05), Inches(5.1), Inches(0.2),
             "V-Intelligence UEBA Advantage", sz=9, bold=True, color=GREEN)
        tbox(slide, Inches(7.55), y + Inches(0.3), Inches(5.1), Inches(0.6),
             advantage, sz=11, color=DARK_GRAY)

    # Bottom
    box = rect(slide, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.7), fill=NAVY, radius=True)
    tbox(slide, Inches(0.9), Inches(6.55), Inches(11.5), Inches(0.55),
         "V-Intelligence UEBA is not an incremental improvement — it is a new class of detection. "
         "The behavioral digital twin sees what no threshold-based system can.",
         sz=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    footer(slide)


def slide_15b_deployment(prs):
    """Deployment Architecture — high-level data flow diagram."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Deployment Architecture")

    tbox(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
         "V-Intelligence UEBA integrates with your existing security infrastructure — "
         "no new sensors, no new data collection agents.",
         sz=16, color=DARK_GRAY)

    # ── Left box: Data Sources ──
    card(slide, Inches(0.5), Inches(2.1), Inches(3.4), Inches(4.0),
         border_top_color=TEAL)
    tbox(slide, Inches(0.7), Inches(2.3), Inches(3.0), Inches(0.35),
         "Data Sources", sz=16, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    sources = [
        "Authentication Logs",
        "File Access Records",
        "Network Traffic Metadata",
        "Endpoint Telemetry",
        "DNS Query Logs",
    ]
    for i, src in enumerate(sources):
        y = Inches(2.8) + Inches(i * 0.55)
        src_box = rect(slide, Inches(0.8), y, Inches(2.9), Inches(0.4),
                       fill=LIGHT_BLUE, radius=True)
        set_text(src_box, f"  {src}", sz=12, bold=False, color=NAVY, align=PP_ALIGN.CENTER)

    # ── Arrow: left to center ──
    tbox(slide, Inches(3.95), Inches(3.7), Inches(0.7), Inches(0.5),
         ">>>", sz=28, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    # ── Center box: V-Intelligence UEBA Engine ──
    engine_bg = card(slide, Inches(4.6), Inches(2.1), Inches(4.2), Inches(4.0),
                     border_top_color=GREEN)
    tbox(slide, Inches(4.8), Inches(2.3), Inches(3.8), Inches(0.35),
         "V-Intelligence UEBA Engine", sz=16, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    engine_items = [
        ("Behavioral Profiling", "Digital twin per entity"),
        ("Continuous Monitoring", "Real-time behavioral change detection"),
        ("Composite Scoring", "Multi-phase risk assessment"),
    ]
    for i, (item_title, item_desc) in enumerate(engine_items):
        y = Inches(2.85) + Inches(i * 0.95)
        item_box = rect(slide, Inches(4.9), y, Inches(3.6), Inches(0.75),
                        fill=LIGHT_GREEN, radius=True)
        set_text(item_box, f"  {item_title}", sz=13, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
        item_box.text_frame.margin_left = Inches(0.15)
        item_box.text_frame.margin_top = Inches(0.05)
        add_p(item_box.text_frame, f"  {item_desc}", sz=10, color=DARK_GRAY, align=PP_ALIGN.LEFT)

    # ── Arrow: center to right ──
    tbox(slide, Inches(8.85), Inches(3.7), Inches(0.7), Inches(0.5),
         ">>>", sz=28, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    # ── Right box: Analyst Output ──
    card(slide, Inches(9.5), Inches(2.1), Inches(3.4), Inches(4.0),
         border_top_color=GOLD)
    tbox(slide, Inches(9.7), Inches(2.3), Inches(3.0), Inches(0.35),
         "Analyst Output", sz=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    outputs = [
        "Ranked Priority List",
        "Risk Dashboard",
        "Alert Integration",
        "Investigation Context",
        "Compliance Reporting",
    ]
    for i, out in enumerate(outputs):
        y = Inches(2.8) + Inches(i * 0.55)
        out_box = rect(slide, Inches(9.8), y, Inches(2.9), Inches(0.4),
                       fill=LIGHT_GOLD, radius=True)
        set_text(out_box, f"  {out}", sz=12, bold=False, color=NAVY, align=PP_ALIGN.CENTER)

    # ── Bottom bar ──
    box = rect(slide, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.6), fill=NAVY, radius=True)
    tbox(slide, Inches(0.9), Inches(6.55), Inches(11.5), Inches(0.45),
         "Cloud or on-premise deployment  |  Works with existing security infrastructure  |  No new sensors required",
         sz=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    footer(slide)


def slide_18_federal(prs):
    """Federal & Critical Infrastructure Alignment."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Federal & Critical Infrastructure Alignment")

    tbox(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
         "V-Intelligence UEBA directly addresses the mandates and frameworks that govern "
         "federal and critical infrastructure cybersecurity.",
         sz=16, color=DARK_GRAY)

    # 2x2 grid of framework cards
    frameworks = [
        ("Executive Order 14028", "Improving the Nation's Cybersecurity",
         "Mandates enhanced detection capabilities and continuous monitoring across federal agencies. "
         "Requires visibility into advanced persistent threats that evade traditional tools.",
         RGBColor(41, 128, 185), 0, 0),
        ("NIST SP 800-53", "Continuous Monitoring Controls",
         "Behavioral analytics for insider threat and APT detection. V-Intelligence UEBA provides "
         "the continuous behavioral monitoring that AU, SI, and CA control families require.",
         TEAL, 1, 0),
        ("NIST SP 800-207", "Zero Trust Architecture",
         "Continuous behavioral verification beyond initial authentication. The behavioral digital twin "
         "provides ongoing trust assessment that Zero Trust demands.",
         PURPLE, 0, 1),
        ("CISA LOTL Guidance", "Living Off the Land Mitigation",
         "Directly addresses Salt Typhoon and Volt Typhoon TTPs. V-Intelligence UEBA detects "
         "the behavioral patterns of credential-based, tool-free attack campaigns.",
         ORANGE, 1, 1),
    ]

    for title, subtitle, desc, clr, col, row in frameworks:
        x = Inches(0.5) + Inches(col * 6.2)
        y = Inches(2.0) + Inches(row * 2.2)
        w = Inches(5.9)
        h = Inches(1.9)

        card(slide, x, y, w, h, border_top_color=clr)

        # Title
        tbox(slide, x + Inches(0.2), y + Inches(0.15), w - Inches(0.4), Inches(0.3),
             title, sz=16, bold=True, color=clr)
        # Subtitle
        tbox(slide, x + Inches(0.2), y + Inches(0.45), w - Inches(0.4), Inches(0.25),
             subtitle, sz=11, bold=True, color=MID_GRAY)
        # Description
        tbox(slide, x + Inches(0.2), y + Inches(0.75), w - Inches(0.4), Inches(0.9),
             desc, sz=12, color=DARK_GRAY)

    # Bottom bar
    box = rect(slide, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.7), fill=NAVY, radius=True)
    tbox(slide, Inches(0.9), Inches(6.55), Inches(11.5), Inches(0.55),
         "V-Intelligence UEBA provides the behavioral intelligence layer that "
         "existing perimeter and endpoint tools cannot.",
         sz=15, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    footer(slide)


def slide_16_business_impact(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Business Impact: What V-Intelligence UEBA Means for Your Organization")

    impacts = [
        ("Detect Nation-State Threats", GREEN,
         "Identify Salt Typhoon-class and Volt Typhoon-class attacks that have operated "
         "undetected at major US organizations for years. The threat-profile detector flags "
         "the threats that no existing tool can see.",
         "4 of 4 attack types detected"),
        ("Eliminate the Detection Gap", TEAL,
         "Close the gap between attacker dwell time (287 days industry average) and detection. "
         "V-Intelligence UEBA's digital twin identifies behavioral anomalies within weeks, not years.",
         "Weeks vs 5-year dwell time"),
        ("Reduce SOC Burden", BLUE,
         "One ranked list replaces thousands of uncorrelated alerts. Your analysts investigate "
         "the highest-risk entities first — not the loudest false alarms.",
         "Single ranked priority list"),
        ("Future-Proof Detection", GOLD,
         "The behavioral digital twin adapts as user behavior naturally evolves. "
         "No signature updates, no threshold tuning, no rules to maintain.",
         "Zero manual tuning required"),
    ]
    for i, (title, clr, desc, metric) in enumerate(impacts):
        y = Inches(1.5) + Inches(i * 1.35)
        card(slide, Inches(0.6), y, Inches(12.1), Inches(1.2))
        rect(slide, Inches(0.6), y, Inches(0.08), Inches(1.2), fill=clr)
        tbox(slide, Inches(1.0), y + Inches(0.08), Inches(4), Inches(0.35),
             title, sz=18, bold=True, color=clr)
        tbox(slide, Inches(1.0), y + Inches(0.42), Inches(7), Inches(0.7),
             desc, sz=13, color=DARK_GRAY)
        badge = rect(slide, Inches(8.5), y + Inches(0.3), Inches(4), Inches(0.5), fill=clr, radius=True)
        set_text(badge, f"  {metric}", sz=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    box = rect(slide, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.6), fill=NAVY, radius=True)
    tbox(slide, Inches(0.9), Inches(6.65), Inches(11.5), Inches(0.5),
         "V-Intelligence UEBA detects the threats that made headlines — before they make yours.",
         sz=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    footer(slide)


def slide_17_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)

    tbox(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.8),
         "Salt Typhoon: 5 years undetected.", sz=28, color=RGBColor(160, 200, 224))
    tbox(slide, Inches(0.8), Inches(1.9), Inches(11.5), Inches(0.8),
         "Volt Typhoon: 5 years undetected.", sz=28, color=RGBColor(160, 200, 224))
    tbox(slide, Inches(0.8), Inches(2.8), Inches(11.5), Inches(0.8),
         "V-Intelligence UEBA: Detected immediately.", sz=40, bold=True, color=GOLD)

    rect(slide, Inches(0.8), Inches(3.9), Inches(3), Inches(0.04), fill=GOLD)

    tbox(slide, Inches(0.8), Inches(4.2), Inches(11.5), Inches(0.8),
         "The Behavioral Digital Twin for Cyber Defense\n"
         "AI-Powered Intelligence for Threats That Traditional Tools Cannot See",
         sz=22, color=WHITE)

    results = [
        ("4 / 4", "attack types\ndetected"),
        ("#1 / 250", "Salt Typhoon\nranked highest"),
        ("0%", "false positive\nrate"),
        ("Zero", "threshold tuning\nrequired"),
    ]
    for i, (val, desc) in enumerate(results):
        x = Inches(0.8) + Inches(i * 3.1)
        tbox(slide, x, Inches(5.3), Inches(2.8), Inches(0.5),
             val, sz=32, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
        tbox(slide, x, Inches(5.8), Inches(2.8), Inches(0.5),
             desc, sz=13, color=RGBColor(160, 200, 224), align=PP_ALIGN.CENTER)

    tbox(slide, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.5),
         "22nd Century Technologies, Inc.  |  V-Intelligence UEBA  |  Confidential",
         sz=14, color=MID_GRAY)


# ═══════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)              # 1.  Title
    slide_02_threat_landscape(prs)   # 2.  Threat Landscape
    slide_03_salt_typhoon(prs)       # 3.  Salt Typhoon Case Study
    slide_04_volt_typhoon(prs)       # 4.  Volt Typhoon Case Study
    slide_04b_threat_spectrum(prs)   # 5.  The Full Threat Spectrum     [NEW]
    slide_05_why_trad_fails(prs)     # 6.  Why Traditional Fails
    slide_06_digital_twin(prs)       # 7.  The Digital Twin Innovation
    slide_07_multi_tier(prs)         # 8.  Multi-Tier Detection
    slide_08_how_acecard_works(prs)  # 9.  How V-Intelligence UEBA Works
    slide_09_screenshot_heatmap(prs) # 10. Screenshot: Detection Matrix
    slide_10_screenshot_scores(prs)  # 11. Screenshot: V-Intelligence Scores
    slide_11_screenshot_separation(prs)  # 12. Screenshot: Separation
    slide_12_screenshot_verdict(prs) # 13. Screenshot: The Verdict
    slide_13_unique_results(prs)     # 14. Unique Results Highlight
    slide_14_salt_typhoon_proof(prs) # 15. Salt Typhoon Validation
    slide_15_competitive_advantage(prs)  # 16. Competitive Advantage
    slide_15b_deployment(prs)        # 17. Deployment Architecture      [NEW]
    slide_18_federal(prs)            # 18. Federal & Critical Infra     [NEW]
    slide_16_business_impact(prs)    # 19. Business Impact
    slide_17_closing(prs)            # 20. Closing

    out_path = "docs/V_Intelligence_UEBA_Sales_Pitch_v2.pptx"
    prs.save(out_path)
    print(f"Saved {len(prs.slides)} slides to {out_path}")


if __name__ == "__main__":
    main()
