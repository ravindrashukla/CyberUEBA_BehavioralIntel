"""5-slide capability/partnership brief for the Brandon Pugh (R Street Institute)
meeting. Theme: Preemptive Cyber Defense + V-Intelligence UEBA behavioral
detection as a two-layer answer for critical-infrastructure defense.
Branding: V-Intelligence UEBA (22nd Century Technologies Inc. parent in footer).
Reuses the exact palette / card style of the 22CT Typhoon reference deck."""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

BASE = os.path.join(os.path.dirname(__file__),
                    "VIntelligence_Preemptive_Behavioral_Brief.pptx")

# === palette (from reference deck) ===
DKNAVY = RGBColor(0x07, 0x14, 0x2A)
NAVY   = RGBColor(0x0B, 0x1F, 0x3A)
BODY   = RGBColor(0x2D, 0x37, 0x48)
SUB    = RGBColor(0x5A, 0x68, 0x78)
ORG    = RGBColor(0xF2, 0x6B, 0x1F)
TEAL   = RGBColor(0x08, 0x91, 0xB2)
CRED   = RGBColor(0xC0, 0x39, 0x2B)
GREEN  = RGBColor(0x1E, 0x7A, 0x44)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
CARD   = RGBColor(0xF7, 0xF8, 0xFA)
PEACH  = RGBColor(0xFD, 0xE5, 0xD4)
CYAN   = RGBColor(0xE0, 0xF2, 0xF7)
GOLD   = RGBColor(0xB8, 0x86, 0x2A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = prs.slide_width, prs.slide_height


def _fill(s, c):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = c

def _rect(s, l, t, w, h, fill=None):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.background()
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background(); sh.shadow.inherit = False
    return sh

def _box(s, l, t, w, h, fill=None, border=None, accent=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.background()
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border:
        sh.line.color.rgb = border; sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    if accent:
        _rect(s, l, t, Inches(0.1), h, fill=accent)
    return sh

def _txt(s, l, t, w, h, text, sz=11, c=BODY, b=False, al=PP_ALIGN.LEFT, fn='Calibri'):
    tf = s.shapes.add_textbox(l, t, w, h).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b
    p.font.name = fn; p.alignment = al
    return tf

def _p(tf, text, sz=11, c=BODY, b=False, sp=Pt(4), al=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b
    p.font.name = 'Calibri'; p.space_before = sp; p.alignment = al
    return p

def _content_slide(part_text, title, subtitle=None, bar_color=ORG):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _fill(s, WHITE)
    _rect(s, Inches(0), Inches(0), W, Inches(0.3), fill=bar_color)
    _txt(s, Inches(0.5), Inches(0.0), Inches(12.3), Inches(0.3), part_text, 11, WHITE, True)
    _txt(s, Inches(0.5), Inches(0.45), Inches(12.3), Inches(0.6), title, 24, NAVY, True, fn='Georgia')
    if subtitle:
        _txt(s, Inches(0.5), Inches(1.02), Inches(12.3), Inches(0.35), subtitle, 13, SUB)
    return s

def _bottom_bar(s, text, y=6.85):
    _rect(s, Inches(0.5), Inches(y), Inches(12.3), Inches(0.45), fill=NAVY)
    _txt(s, Inches(0.7), Inches(y - 0.02), Inches(12.0), Inches(0.5), text, 11.5, GOLD, True, PP_ALIGN.CENTER)


# ================================================================
# SLIDE 1 — TITLE
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_fill(s, DKNAVY)
_rect(s, Inches(0), Inches(0), W, Inches(0.18), fill=ORG)
_rect(s, Inches(0), H - Inches(0.18), W, Inches(0.18), fill=TEAL)
_txt(s, Inches(1), Inches(1.7), Inches(11.3), Inches(0.5),
     'V-INTELLIGENCE UEBA', 30, RGBColor(0x39, 0xC0, 0xC8), True, PP_ALIGN.CENTER, 'Georgia')
_txt(s, Inches(1), Inches(2.55), Inches(11.3), Inches(1.0),
     'Preemptive Cyber Defense + Behavioral Trajectory Intelligence',
     22, WHITE, True, PP_ALIGN.CENTER, 'Georgia')
_txt(s, Inches(1), Inches(3.5), Inches(11.3), Inches(0.5),
     'A Two-Layer Answer for Critical-Infrastructure Defense', 15, RGBColor(0xA0, 0xB0, 0xC0),
     False, PP_ALIGN.CENTER)

_txt(s, Inches(1.2), Inches(4.6), Inches(5.4), Inches(1.4),
     'LAYER 1 — PREEMPTIVE', 13, ORG, True)
_txt(s, Inches(1.2), Inches(4.95), Inches(5.4), Inches(1.4),
     'Mathematically prove every known attack path\nis closed — before any traffic flows.',
     12, RGBColor(0xCC, 0xDD, 0xEE))
_txt(s, Inches(7.0), Inches(4.6), Inches(5.4), Inches(1.4),
     'LAYER 2 — BEHAVIORAL', 13, RGBColor(0x39, 0xC0, 0xC8), True)
_txt(s, Inches(7.0), Inches(4.95), Inches(5.4), Inches(1.4),
     'Continuous 1536-d behavioral trajectories\nwith MITRE-mapped drift direction.',
     12, RGBColor(0xCC, 0xDD, 0xEE))

_txt(s, Inches(1), Inches(6.55), Inches(11.3), Inches(0.4),
     'Prepared for discussion with Brandon Pugh  ·  R Street Institute', 12,
     RGBColor(0x90, 0xA0, 0xB0), False, PP_ALIGN.CENTER)
_txt(s, Inches(1), Inches(6.9), Inches(11.3), Inches(0.3),
     '22nd Century Technologies Inc. (TSCTI)  ·  UNCLASSIFIED', 10,
     RGBColor(0x60, 0x70, 0x80), False, PP_ALIGN.CENTER)


# ================================================================
# SLIDE 2 — THE STAKES
# ================================================================
s = _content_slide('THE STAKES', 'Apex Adversaries Are Already Inside Critical Infrastructure',
                   'PRC state actors walked past every commercial defense — and stayed for years.', CRED)

_box(s, Inches(0.5), Inches(1.5), Inches(5.85), Inches(3.05), fill=CARD, accent=ORG)
_txt(s, Inches(0.8), Inches(1.58), Inches(5.3), Inches(0.4), 'VOLT TYPHOON  ·  CISA AA24-038A', 14, NAVY, True, fn='Georgia')
tf = _txt(s, Inches(0.8), Inches(2.08), Inches(5.3), Inches(2.4),
     'Mission: pre-position in U.S. energy, water, transport and comms for sabotage during crisis.', 11, BODY)
_p(tf, 'Living-off-the-Land — wmic, ntdsutil, netsh, rundll32. No malware, no signatures, no IOCs.', 11, BODY)
_p(tf, 'Every tool used is a legitimate Windows binary.', 11, BODY)
_p(tf, '5+ YEARS undetected dwell in U.S. critical infrastructure.', 11.5, CRED, True, sp=Pt(8))

_box(s, Inches(6.95), Inches(1.5), Inches(5.85), Inches(3.05), fill=CARD, accent=TEAL)
_txt(s, Inches(7.25), Inches(1.58), Inches(5.3), Inches(0.4), 'SALT TYPHOON  ·  CISA AA25-239A', 14, NAVY, True, fn='Georgia')
tf = _txt(s, Inches(7.25), Inches(2.08), Inches(5.3), Inches(2.4),
     'Mission: espionage at carrier scale across U.S. and allied telecom.', 11, BODY)
_p(tf, '200+ telcos in 80+ countries; 9 U.S. carriers confirmed (AT&T, Verizon, T-Mobile).', 11, BODY)
_p(tf, 'Accessed CALEA lawful-intercept systems used for senior U.S. officials.', 11, BODY)
_p(tf, '4+ YEARS undetected. Slow, sub-threshold, patient.', 11.5, CRED, True, sp=Pt(8))

# policy band
_box(s, Inches(0.5), Inches(4.75), Inches(12.3), Inches(1.95), fill=PEACH, accent=ORG)
_txt(s, Inches(0.8), Inches(4.83), Inches(11.7), Inches(0.35), 'THE POLICY INFLECTION', 13, ORG, True)
tf = _txt(s, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.4),
     'Defense is shifting from reactive to preemptive. Gartner projects preemptive cybersecurity will grow '
     'from <5% of security spending in 2024 to 50% by 2030. Yet 99% of firewall breaches are still '
     'misconfiguration, and threshold-based detection never fires on slow, credential-borne intrusions.', 11.5, BODY)
_p(tf, 'The structural lesson of Volt + Salt Typhoon: we must prove the door is closed — and watch the room for what '
       'slips through. Neither layer alone is sufficient.', 11.5, NAVY, True, sp=Pt(8))

_bottom_bar(s, 'Years of undetected access vs. immediate identification — the gap is operational, not theoretical.')


# ================================================================
# SLIDE 3 — WHY TODAY'S TOOLS MISS IT
# ================================================================
s = _content_slide('THE GAP', 'Why Current Tooling Cannot Close It',
                   'Two structural gaps — one in prevention, one in detection.', CRED)

_box(s, Inches(0.5), Inches(1.5), Inches(5.85), Inches(2.55), fill=CARD, accent=CRED)
_txt(s, Inches(0.8), Inches(1.56), Inches(5.3), Inches(0.3), 'PREVENTION FAILS', 12, CRED, True)
tf = _txt(s, Inches(0.8), Inches(1.92), Inches(5.3), Inches(2.0),
     'Pentests sample; they never prove. Point-in-time and incomplete.', 10.5, BODY)
_p(tf, 'ASM sees what exists, not what is reachable.', 10.5, BODY)
_p(tf, 'Config review is manual — stale within hours.', 10.5, BODY)
_p(tf, 'The rule state space is too large to sample adequately.', 10.5, BODY)
_p(tf, 'KEV patching covers only ~20% of attacks.', 10.5, NAVY, True)

_box(s, Inches(6.95), Inches(1.5), Inches(5.85), Inches(2.55), fill=CARD, accent=CRED)
_txt(s, Inches(7.25), Inches(1.56), Inches(5.3), Inches(0.3), 'DETECTION FAILS', 12, CRED, True)
tf = _txt(s, Inches(7.25), Inches(1.92), Inches(5.3), Inches(2.0),
     'Threshold UEBA: LoTL never crosses a fixed line.', 10.5, BODY)
_p(tf, 'Sigma rules: adversaries retool faster than rules ship.', 10.5, BODY)
_p(tf, 'Per-source silos: one attacker looks like 10 users.', 10.5, BODY)
_p(tf, '"87% anomalous" — anomalous HOW? No direction.', 10.5, BODY)
_p(tf, 'Only point-in-time snapshots — no trajectory.', 10.5, NAVY, True)

_box(s, Inches(0.5), Inches(4.25), Inches(12.3), Inches(2.4), fill=CARD, accent=NAVY)
_txt(s, Inches(0.8), Inches(4.31), Inches(11.7), Inches(0.3), 'WHAT THE ADVERSARY EXPLOITS', 13, NAVY, True)
rows = [
    ('Configuration gaps no one has PROVEN are closed', 'Both Typhoons enter via unpatched / misconfigured edge devices'),
    ('Valid credentials indistinguishable from legitimate use', 'Stolen creds pass every config-based check by design'),
    ('Living-off-the-Land tools allowed by policy', 'wmic / ntdsutil / netsh cannot be blanket-blocked'),
    ('Slow drift invisible to fixed thresholds', '5+ year campaigns operate below every alert threshold'),
    ('Identity fragmentation hides cross-domain chains', 'One attacker appears as 10 unrelated identities'),
]
y = Inches(4.72)
for a, b in rows:
    _txt(s, Inches(0.8), y, Inches(6.0), Inches(0.3), a, 10.5, NAVY, True)
    _txt(s, Inches(6.9), y, Inches(5.7), Inches(0.3), b, 10.5, SUB)
    y += Inches(0.37)

_bottom_bar(s, 'Result: $0 invested in tooling that could have detected either campaign before public disclosure.')


# ================================================================
# SLIDE 4 — THE TWO-LAYER ANSWER
# ================================================================
s = _content_slide('THE ANSWER', 'Two Layers, One Shared Vocabulary',
                   'Preemptive shrinks the attack surface mathematically. V-Intelligence detects what remains.', ORG)

# Layer 1
_box(s, Inches(0.5), Inches(1.5), Inches(5.85), Inches(3.55), fill=PEACH, accent=ORG)
_txt(s, Inches(0.8), Inches(1.58), Inches(5.3), Inches(0.4), 'LAYER 1 — PREEMPTIVE DEFENSE', 14, ORG, True, fn='Georgia')
tf = _txt(s, Inches(0.8), Inches(2.08), Inches(5.3), Inches(2.9),
     'Builds a formal mathematical model of every firewall, IdP, IPS, SASE and WAF and reasons '
     'exhaustively over the full configuration state space — no sampling.', 11, BODY)
_p(tf, 'Proves every known CVE path is blocked or compensated.', 11, NAVY, True, sp=Pt(8))
_p(tf, 'Catches config drift within the hour.', 11, NAVY, True)
_p(tf, 'Complete reasoning — not probabilistic. Zero telemetry, zero agents.', 11, NAVY, True)
_p(tf, '"Close the door before they knock."', 11, ORG, True, sp=Pt(8))

# Layer 2
_box(s, Inches(6.95), Inches(1.5), Inches(5.85), Inches(3.55), fill=CYAN, accent=TEAL)
_txt(s, Inches(7.25), Inches(1.58), Inches(5.3), Inches(0.4), 'LAYER 2 — V-INTELLIGENCE UEBA', 14, TEAL, True, fn='Georgia')
tf = _txt(s, Inches(7.25), Inches(2.08), Inches(5.3), Inches(2.9),
     'Embeds 5 behavioral signals (auth, process, network, file, identity) into a unified 1536-d '
     'vector per entity, per window — and tracks the trajectory.', 11, BODY)
_p(tf, 'CUSUM change-point catches slow, sub-threshold drift.', 11, NAVY, True, sp=Pt(8))
_p(tf, 'Drift direction projects onto 8 MITRE ATT&CK concepts — not "anomalous," but "drifting toward lateral_movement (T1021)."', 11, NAVY, True)
_p(tf, 'Entity fusion unifies 10+ identity systems into one behavioral entity.', 11, NAVY, True)
_p(tf, '"Watch the room for what slips through."', 11, TEAL, True, sp=Pt(8))

# proof band
_box(s, Inches(0.5), Inches(5.25), Inches(12.3), Inches(1.4), fill=CARD, accent=GREEN)
_txt(s, Inches(0.8), Inches(5.31), Inches(11.7), Inches(0.3), 'VALIDATION — BLIND DETECTION TEST, 250 ENTITIES', 13, GREEN, True)
tf = _txt(s, Inches(0.8), Inches(5.68), Inches(11.7), Inches(0.9),
     'Against four stealth campaigns (insider, slow APT, Volt Typhoon LoTL, Salt Typhoon telecom): traditional SIEM '
     'caught 0 of 4; z-score caught 1 of 4; V-Intelligence + composite scoring caught 4 of 4 at 8.1% false-positive '
     '(cleanly separating 2 of 4); the multi-front threat-profile detector caught 4 of 4 at 0% false-positive. '
     'Salt Typhoon (max z-score 1.71 — invisible to thresholds) ranked #1 of 250.', 11, BODY)

_bottom_bar(s, 'Shared MITRE vocabulary + single entity model = a system, not a bundle. Zero gaps between prevention and detection.')


# ================================================================
# SLIDE 5 — PREEMPTIVE LAYER IN DEPTH
# ================================================================
s = _content_slide('PREEMPTIVE — LAYER 1', 'How Preemptive Defense Closes the Door',
                   'Prove every known attack path is blocked — before any traffic flows. Three intelligence pillars, one formal model.', ORG)

pillars = [
    ('ATTACK INTELLIGENCE', ORG,
     'Ingests terabytes of CVE, CTI, IOC/IOA feeds and advisories. ML + NLP + LLMs build deep, '
     'MITRE ATT&CK-enriched attack graphs of every campaign that matters — auto-updating on each new advisory.'),
    ('DEFENSE INTELLIGENCE', TEAL,
     'A symbolic model of computation of every firewall, IdP, IPS, SASE and WAF. Reasons exhaustively over the '
     'full configuration state space — no sampling. Pinpoints shadows, conflicts, weak rules and intent errors across vendors.'),
    ('REMEDIATION INTELLIGENCE', GREEN,
     'A guardrailed agentic-AI reasoner, grounded in each vendor\'s behavior, prescribes precise, risk-prioritized '
     'fixes — no new errors, no business disruption. Findings reviewed, then applied and continuously re-verified.'),
]
xs = [Inches(0.5), Inches(4.65), Inches(8.8)]
ws = [Inches(3.95), Inches(3.95), Inches(4.0)]
for x, w, (head, col, body) in zip(xs, ws, pillars):
    _box(s, x, Inches(1.5), w, Inches(2.85), fill=CARD, accent=col)
    _txt(s, x + Inches(0.25), Inches(1.58), w - Inches(0.45), Inches(0.4), head, 12.5, col, True)
    _txt(s, x + Inches(0.25), Inches(2.05), w - Inches(0.45), Inches(2.2), body, 10.5, BODY)

_box(s, Inches(0.5), Inches(4.55), Inches(12.3), Inches(1.6), fill=PEACH, accent=ORG)
_txt(s, Inches(0.8), Inches(4.62), Inches(11.7), Inches(0.3), 'WHY A MODEL — NOT A SCAN', 12, ORG, True)
tf = _txt(s, Inches(0.8), Inches(4.98), Inches(11.7), Inches(0.7),
     'A multi-vendor firewall path spans an astronomically large configuration state space (≫2^8000). Sampling a '
     'billion paths per second for a year covers effectively 0%. A formal model instead PROVES, for every threat: '
     'does any rule combination allow this traffic? Complete coverage — no false positives, no false negatives.', 11, BODY)
_txt(s, Inches(0.8), Inches(5.78), Inches(11.7), Inches(0.3),
     'No telemetry  ·  No agents  ·  Config-only pull  ·  Time-to-value < 2 weeks  ·  SaaS / on-prem / air-gapped',
     10.5, NAVY, True)

_bottom_bar(s, 'Gartner: 99% of firewall breaches are misconfiguration. Preemptive verification eliminates that surface mathematically.')


# ================================================================
# SLIDE 6 — PREEMPTIVE IN PRACTICE
# ================================================================
s = _content_slide('PREEMPTIVE — IN PRACTICE', 'Coverage, Proof, and Federal Fit',
                   'One formal model answers many questions — and finds what sampling tools never will.', ORG)

# Left — real findings
_box(s, Inches(0.5), Inches(1.5), Inches(5.85), Inches(3.7), fill=CARD, accent=CRED)
_txt(s, Inches(0.8), Inches(1.57), Inches(5.3), Inches(0.35), 'WHAT FORMAL VERIFICATION FINDS', 11.5, CRED, True)
tf = _txt(s, Inches(0.8), Inches(1.98), Inches(5.3), Inches(3.1),
     'SMB lateral-movement path shadowed by three interspersed rules — IPS / AntiVirus silently bypassed.', 10.5, BODY)
_p(tf, 'Two-vendor blind spot: PAN and Fortinet both allowed update traffic with no IPS inspection — an open RCE path.', 10.5, BODY, sp=Pt(7))
_p(tf, 'Unsanctioned-app rule partially shadowed and built on an incomplete object — credential and IP exposure.', 10.5, BODY, sp=Pt(7))
_p(tf, '~25% of internet-facing rules carried weak or no threat inspection — proven reachable into trusted zones.', 10.5, BODY, sp=Pt(7))
_p(tf, 'SMTP trusted→untrusted with no inspection profile — a NIST 800-53 SC-7 compliance gap, remediated same day.', 10.5, BODY, sp=Pt(7))

# Right — one model, nine use cases
_box(s, Inches(6.95), Inches(1.5), Inches(5.85), Inches(3.7), fill=CYAN, accent=TEAL)
_txt(s, Inches(7.25), Inches(1.57), Inches(5.3), Inches(0.35), 'ONE FORMAL MODEL → NINE USE CASES', 11.5, TEAL, True)
uses = [
    'Security posture — segmentation & inspection correctness',
    'External attack surface — ASM compliance, formally extended',
    'Change management — does the change meet its intent?',
    'Configuration drift — vs. an approved golden baseline',
    'Vendor migration — pre/post security posture equivalence',
    'Business / regulatory / infosec compliance',
    'Complete preemptive defense — full TTP coverage proof',
    'Third-party / supply-chain risk',
    'Precision cyber insurance — exact risk, not heuristics',
]
y = Inches(1.98)
for u in uses:
    _txt(s, Inches(7.25), y, Inches(5.35), Inches(0.3), '•  ' + u, 10, BODY)
    y += Inches(0.34)

_box(s, Inches(0.5), Inches(5.35), Inches(12.3), Inches(1.3), fill=DKNAVY)
_txt(s, Inches(0.8), Inches(5.43), Inches(11.7), Inches(0.3), 'FEDERAL FIT', 11.5, RGBColor(0x39, 0xC0, 0xC8), True)
_txt(s, Inches(0.8), Inches(5.78), Inches(11.7), Inches(0.8),
     'Maps to NIST CSF pre-attack functions (Govern · Identify · Protect), supports CISA CDM continuous assurance, '
     'and enables Policy-as-Code. Gartner projects preemptive cyber rising from <5% of security spend in 2024 to 50% by 2030.',
     11, WHITE)

_bottom_bar(s, 'Preemptive shrinks the attack surface to what is provably reachable. V-Intelligence then watches what remains.')


# ================================================================
# SLIDE 7 — WHY SEMANTIC BEATS SCALAR
# ================================================================
s = _content_slide('THE CORE INSIGHT', 'Why Behavioral Meaning Beats Scalar Counts',
                   'The same number can mean two completely different things. Embeddings see the difference; counters never will.', TEAL)

# Left — the blind spot
_box(s, Inches(0.5), Inches(1.5), Inches(5.85), Inches(4.55), fill=CARD, accent=CRED)
_txt(s, Inches(0.8), Inches(1.57), Inches(5.3), Inches(0.35), 'TRADITIONAL TOOLS COUNT — THEY DO NOT COMPREHEND', 11.5, CRED, True)
tf = _txt(s, Inches(0.8), Inches(1.98), Inches(5.3), Inches(1.3),
     'Two users each touch 47 files this week. One reads their own team\'s documents; the other scans 12 '
     'departments at 2 a.m. A scalar feature sees the same number — 47. The attacker is invisible.', 11, BODY)
_txt(s, Inches(0.8), Inches(3.35), Inches(5.3), Inches(0.3), 'EVERY TRADITIONAL METHOD HAS A FATAL FLAW', 10.5, NAVY, True)
weak = [
    ('Isolation Forest', 'misses slow drift within normal ranges'),
    ('One-Class SVM', 'high false positives in high dimensions'),
    ('Local Outlier Factor', 'sensitive to feature scaling'),
    ('Z-Score (|z|>3)', 'only catches single-feature spikes'),
    ('Feature CUSUM', 'flags everyone; no semantic context'),
]
y = Inches(3.72)
for m, w in weak:
    _txt(s, Inches(0.85), y, Inches(2.0), Inches(0.3), m, 10, CRED, True)
    _txt(s, Inches(2.95), y, Inches(3.2), Inches(0.3), w, 10, BODY)
    y += Inches(0.42)

# Right — what V-Intelligence does
_box(s, Inches(6.95), Inches(1.5), Inches(5.85), Inches(4.55), fill=CYAN, accent=TEAL)
_txt(s, Inches(7.25), Inches(1.57), Inches(5.3), Inches(0.35), 'V-INTELLIGENCE READS BEHAVIOR AS MEANING', 11.5, TEAL, True)
steps = [
    ('1.  Serialize', 'Raw metrics become prose: "USR-042 accessed 47 restricted files across 12 departments, 68% outside normal hours."'),
    ('2.  Embed', 'Text → a 1,536-dimensional vector that captures behavioral semantics, not just counts.'),
    ('3.  Five behavioral zones', 'Identity · Access Pattern · Data Behavior · Network Footprint · Risk Posture.'),
    ('4.  Drift', 'Cosine distance between consecutive weekly embeddings = the direction of behavioral change.'),
]
y = Inches(2.02)
for h, d in steps:
    _txt(s, Inches(7.25), y, Inches(5.3), Inches(0.3), h, 11, NAVY, True)
    _txt(s, Inches(7.25), y + Inches(0.32), Inches(5.3), Inches(0.62), d, 10.5, BODY)
    y += Inches(1.02)

_bottom_bar(s, 'Scalar features see the same number. Embeddings see different meaning — that is the whole difference between 0/4 and 4/4.')


# ================================================================
# SLIDE 6 — THE UEBA INNOVATION GRID
# ================================================================
s = _content_slide('THE INNOVATION', 'What No Other UEBA Does',
                   'Each technique is well-known. The innovation is the integration — and seven pieces exist nowhere else.', TEAL)

innov = [
    ('Behavioral embedding in a unified 1536-d space',
     'Industry scores features and rules. We embed behavior into vectors — meaning, not just magnitude.'),
    ('Per-signal drift decomposition (6 vectors / 9,216-d)',
     'Five signals embedded individually + one composite. Tells you "AUTH drove 72% of the change."'),
    ('CUSUM change-point on embedding trajectories',
     'Sequential detection on vectors — minimax-optimal (Lorden 1971). Dwell of years collapses to days.'),
    ('Drift direction projected onto MITRE ATT&CK',
     'Not "87% anomalous" but "drifting toward lateral_movement (T1021) at 0.82." Actionable, auditable.'),
    ('Entity fusion across 10+ identity systems',
     'AD, Entra, Okta, AWS, K8s, EDR, VPN, PKI, TACACS+ → one entity. An attacker stops looking like 10 users.'),
    ('HMAC-chained forensic evidence',
     'Tamper-evident behavioral history — altering any past snapshot breaks the chain. Court-admissible.'),
    ('ABAC trust loop with behavioral triggers',
     'Detection and proportional enforcement in one system — step-up MFA to lockout in under 5 minutes.'),
]
cols = [Inches(0.5), Inches(6.95)]
rowH = Inches(1.24)
step = Inches(1.30)
top0 = Inches(1.45)
for idx, (head, sub) in enumerate(innov):
    col = cols[idx % 2]
    top = top0 + (idx // 2) * step
    _box(s, col, top, Inches(5.85), rowH, fill=CYAN if idx % 2 else CARD, accent=TEAL if idx % 2 else ORG)
    _txt(s, col + Inches(0.3), top + Inches(0.06), Inches(5.35), Inches(0.5), f'{idx+1}.  {head}', 11.5, NAVY, True)
    _txt(s, col + Inches(0.3), top + Inches(0.55), Inches(5.35), Inches(0.66), sub, 10, BODY)

# 8th cell (right column, row 4) — summary callout
top = top0 + 3 * step
_box(s, Inches(6.95), top, Inches(5.85), rowH, fill=PEACH, accent=ORG)
_txt(s, Inches(7.25), top + Inches(0.06), Inches(5.35), Inches(0.45), 'WHY IT COMPOUNDS', 11.5, ORG, True)
_txt(s, Inches(7.25), top + Inches(0.5), Inches(5.35), Inches(0.7),
     'A normal user can score high on one dimension by chance. An attacker scores high across many — '
     'embedding + trajectory + direction + fusion reinforce into one defensible signal.', 10, BODY)

_bottom_bar(s, 'No commercial UEBA combines all seven. Most have zero. This is the differentiation, not an increment.')


# ================================================================
# SLIDE 6 — INSIDE THE ENGINE
# ================================================================
s = _content_slide('HOW IT WORKS', 'Inside the Engine: From Telemetry to Direction',
                   'Dual embedding builds the entity; composite scoring decides who is anomalous.', TEAL)

# Left: the math chain
_box(s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(3.35), fill=CARD, accent=TEAL)
_txt(s, Inches(0.8), Inches(1.56), Inches(5.5), Inches(0.35), 'THE BEHAVIORAL MATH', 12, TEAL, True)
tf = _txt(s, Inches(0.8), Inches(1.98), Inches(5.5), Inches(2.8),
     'A.  Embed — 5 signal summaries → text-embedding-3-small → 1536-d behavioral state per window.', 10.5, BODY)
_p(tf, 'B.  Drift — drift(t) = 1 − cos(V_t, V_baseline).  Healthy ~0.02–0.05; compromised 0.06–0.15.', 10.5, BODY, sp=Pt(7))
_p(tf, 'C.  CUSUM — S(t) = max(0, S(t−1) + drift − μ − k).  Alarm at 4σ. Accumulates sub-threshold drift.', 10.5, BODY, sp=Pt(7))
_p(tf, 'D.  Direction — proj(c) = cos(V_t − V_baseline, concept_c) over 8 MITRE concepts → top-3 + technique IDs.', 10.5, BODY, sp=Pt(7))
_p(tf, 'Dual embedding: 5 individual + 1 composite = 6 × 1,536 = 9,216 dims per entity, per window.', 10.5, NAVY, True, sp=Pt(9))

# Right: 5-phase composite scoring
_box(s, Inches(6.85), Inches(1.5), Inches(5.95), Inches(3.35), fill=CYAN, accent=ORG)
_txt(s, Inches(7.15), Inches(1.56), Inches(5.4), Inches(0.35), '5-PHASE COMPOSITE SCORING', 12, ORG, True)
phases = [
    ('Signal strength', 'How far the drift exceeds the entity\'s own baseline.'),
    ('Breadth', 'How many of the behavioral features moved at once.'),
    ('Sustained deviation', 'Whether the shift persists window over window.'),
    ('Context divergence', 'Whether it appears under multiple analytical lenses.'),
    ('Novelty persistence', 'Whether genuinely new behaviors keep recurring.'),
]
tf = _txt(s, Inches(7.15), Inches(2.0), Inches(5.4), Inches(2.5), '', 10.5, BODY)
first = True
for nm, ds in phases:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.font.name = 'Calibri'; p.space_before = Pt(6)
    r1 = p.add_run(); r1.text = f'{nm} — '; r1.font.bold = True; r1.font.size = Pt(10.5); r1.font.color.rgb = NAVY
    r2 = p.add_run(); r2.text = ds; r2.font.size = Pt(10.5); r2.font.color.rgb = BODY
_txt(s, Inches(7.15), Inches(4.42), Inches(5.4), Inches(0.35),
     'Fused into one ranked list. No single phase catches everyone.', 10.5, ORG, True)

# Bottom: the payoff
_box(s, Inches(0.5), Inches(5.05), Inches(12.3), Inches(1.6), fill=CARD, accent=GREEN)
_txt(s, Inches(0.8), Inches(5.11), Inches(11.7), Inches(0.3), 'THE PAYOFF — SEMANTIC SPACE SEPARATES WHAT FEATURE SPACE HIDES', 12.5, GREEN, True)
_txt(s, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.05),
     'In feature space the four stealth attackers sit inside the normal band — undetectable. In V-Intelligence\'s '
     'semantic embedding space, all four separate from the population. Composite scoring then ranks them: '
     'insider #2, slow APT #7, Volt Typhoon LoTL #24, Salt Typhoon #1 — all at 8.1% false-positive, where traditional '
     'methods caught zero.', 11, BODY)

_bottom_bar(s, 'V-Intelligence builds the digital entity. Composite scoring detects the anomaly. Direction makes it actionable.')


# ================================================================
# SLIDE 8 — THE PROOF: FOUR ATTACKERS, FOUR TELLS
# ================================================================
s = _content_slide('THE PROOF', 'Four Stealth Campaigns, Four Different Tells',
                   'No single method catches them all. Composite scoring catches each on the phase that exposes it — 250 entities, blind.', GREEN)

attackers = [
    ('USR-156', 'Insider Threat · 8-month', 'Rank #2 / 250', CRED,
     'Tell: Signal Strength 30 + Breadth 19', 'Strong, broad, sustained escalation across many features.'),
    ('USR-234', 'Slow APT · 180-day', 'Rank #7 / 250', ORG,
     'Tell: Novelty Persistence 13', 'Low magnitude, but novel C2 destinations recur week after week.'),
    ('USR-042', 'Volt Typhoon LoTL · 115-day', 'Rank #24 / 250', GOLD,
     'Tell: Breadth 11', 'Living-off-the-Land creates breadth, not magnitude.'),
    ('USR-118', 'Salt Typhoon Telecom · 100-day', 'Rank #1 / 250', TEAL,
     'Tell: Sustained 9.6 + Context 7.9', 'Max z-score 1.71 — invisible to thresholds — ranked #1 of 250.'),
]
xs = [Inches(0.5), Inches(3.62), Inches(6.74), Inches(9.86)]
cw = Inches(2.97)
for x, (nm, ty, rk, col, tell, why) in zip(xs, attackers):
    _box(s, x, Inches(1.5), cw, Inches(2.75), fill=CARD, accent=col)
    _txt(s, x + Inches(0.18), Inches(1.56), cw - Inches(0.3), Inches(0.3), nm, 13, NAVY, True, fn='Georgia')
    _txt(s, x + Inches(0.18), Inches(1.92), cw - Inches(0.3), Inches(0.5), ty, 9.5, SUB)
    _txt(s, x + Inches(0.18), Inches(2.38), cw - Inches(0.3), Inches(0.3), rk, 11, col, True)
    _txt(s, x + Inches(0.18), Inches(2.78), cw - Inches(0.3), Inches(0.4), tell, 9.5, NAVY, True)
    _txt(s, x + Inches(0.18), Inches(3.25), cw - Inches(0.3), Inches(0.9), why, 9.5, BODY)

# discrimination margin band
_box(s, Inches(0.5), Inches(4.45), Inches(6.0), Inches(1.55), fill=CYAN, accent=GREEN)
_txt(s, Inches(0.78), Inches(4.51), Inches(5.5), Inches(0.3), 'THE DISCRIMINATION MARGIN', 11.5, GREEN, True)
tf = _txt(s, Inches(0.78), Inches(4.88), Inches(5.5), Inches(1.05),
     'Normal users — average composite 6.6, max 21.5.', 10.5, BODY)
_p(tf, 'Attackers — lowest 13.7, average 32.7.', 10.5, NAVY, True, sp=Pt(5))
_p(tf, 'The gap between the strongest normal user and the weakest attacker is clean separation no single method achieves.', 10.5, BODY, sp=Pt(5))

# detection timeline band
_box(s, Inches(6.7), Inches(4.45), Inches(6.1), Inches(1.55), fill=CARD, accent=ORG)
_txt(s, Inches(6.98), Inches(4.51), Inches(5.6), Inches(0.3), 'TIME-TO-DETECT — INSIDER (USR-156)', 11.5, ORG, True)
tf = _txt(s, Inches(6.98), Inches(4.88), Inches(5.6), Inches(1.05),
     'Isolation Forest / SVM / LOF / Z-Score:  NEVER ALERT.', 10.5, CRED, True)
_p(tf, 'V-Intelligence + composite:  flags in Week 8 — before data-access escalation.', 10.5, GREEN, True, sp=Pt(5))
_p(tf, 'Direction tracked: baseline → off-hours access → cross-dept scope creep → data staging + exfiltration.', 10.5, BODY, sp=Pt(5))

_bottom_bar(s, '4 / 4 detected at 8.1% false-positive (threat-profile detector: 0% FP) — versus 0 / 4 traditional and 1 / 4 z-score. Earlier detection = smaller blast radius.')


# ================================================================
# SLIDE 9 — WHY IT MATTERS / PARTNERSHIP
# ================================================================
s = _content_slide('THE PATH FORWARD', 'Proven Foundation, National Relevance, Open to Partnership',
                   'A capability built on a fielded model — and a conversation worth having.', TEAL)

_box(s, Inches(0.5), Inches(1.5), Inches(3.95), Inches(3.7), fill=CARD, accent=ORG)
_txt(s, Inches(0.75), Inches(1.58), Inches(3.5), Inches(0.4), 'PROVEN', 13, ORG, True)
tf = _txt(s, Inches(0.75), Inches(2.0), Inches(3.5), Inches(3.1),
     'Behavioral trajectory engine fielded at DLA on the Entity Digital Model (500+ entities, CUSUM, 5 signals).', 10.5, BODY)
_p(tf, 'Same architecture, retargeted to cyber UEBA.', 10.5, BODY)
_p(tf, '4 / 4 stealth campaigns detected at 8.1% FP (threat-profile detector: 0% FP) in blind test.', 10.5, NAVY, True)
_p(tf, 'Explainable: WHAT is happening, WHERE it lives, WHO else is affected.', 10.5, BODY, sp=Pt(8))

_box(s, Inches(4.65), Inches(1.5), Inches(3.95), Inches(3.7), fill=CARD, accent=TEAL)
_txt(s, Inches(4.9), Inches(1.58), Inches(3.5), Inches(0.4), 'NATIONAL RELEVANCE', 13, TEAL, True)
tf = _txt(s, Inches(4.9), Inches(2.0), Inches(3.5), Inches(3.1),
     'Directly counters the Volt + Salt Typhoon playbook the nation is now contending with.', 10.5, BODY)
_p(tf, 'Aligns with the preemptive-defense shift (Gartner: 50% by 2030).', 10.5, BODY)
_p(tf, 'Deployable across IL5 / IL6 / JWICS; local model option for air-gapped enclaves.', 10.5, BODY)
_p(tf, 'No PII in embeddings — privacy-preserving by design.', 10.5, NAVY, True, sp=Pt(8))

_box(s, Inches(8.8), Inches(1.5), Inches(4.0), Inches(3.7), fill=CYAN, accent=GREEN)
_txt(s, Inches(9.05), Inches(1.58), Inches(3.55), Inches(0.4), 'WHERE R STREET FITS', 13, GREEN, True)
tf = _txt(s, Inches(9.05), Inches(2.0), Inches(3.55), Inches(3.1),
     'Independent perspective on the preemptive + behavioral framing for critical-infrastructure policy.', 10.5, BODY)
_p(tf, 'Thought-leadership on closing the LoTL / valid-account detection gap.', 10.5, BODY)
_p(tf, 'A grounded reference case for the reactive-to-preemptive transition.', 10.5, BODY)
_p(tf, 'Explore: briefing, joint commentary, or validation dialogue.', 10.5, NAVY, True, sp=Pt(8))

_box(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.25), fill=DKNAVY)
_txt(s, Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.5),
     'Preemptive closes the door. V-Intelligence watches the room. Together: detect behavioral shifts hours before IOCs fire — '
     'explain what is happening, act proportionally.', 13, WHITE, True, PP_ALIGN.CENTER)
_txt(s, Inches(0.8), Inches(6.25), Inches(11.7), Inches(0.3),
     '22nd Century Technologies Inc.  ·  Contact: ravindra.shukla@gmail.com', 11, RGBColor(0x90, 0xA8, 0xC0), False, PP_ALIGN.CENTER)


# ---- save (auto-version if locked/open) ----
out = BASE
n = 2
while True:
    try:
        prs.save(out)
        break
    except PermissionError:
        root, ext = os.path.splitext(BASE)
        out = f"{root}_v{n}{ext}"; n += 1

print("Saved:", out, "|", len(prs.slides._sldIdLst), "slides")
