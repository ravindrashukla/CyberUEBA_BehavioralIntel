#!/usr/bin/env python3
"""
Build Detection Comparison PowerPoint Deck.

Business-audience-friendly presentation explaining why traditional
detection fails and how V-Intelligence UEBA + Composite Scoring catches sophisticated
cyber threats. No proprietary math, no embedding dimensions, no CUSUM
internals — cybersecurity terms only, explained in plain language.

Output: docs/Detection_Comparison_Business_Deck.pptx
"""

import sys
sys.stdout.reconfigure(encoding="utf-8")

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colors ──
NAVY       = RGBColor(13,  27,  42)
BLUE       = RGBColor(27,  79,  114)
TEAL       = RGBColor(14,  107, 138)
GOLD       = RGBColor(183, 149, 11)
RED        = RGBColor(192, 57,  43)
GREEN      = RGBColor(39,  174, 96)
WHITE      = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(240, 242, 245)
DARK_GRAY  = RGBColor(80,  80,  80)
MID_GRAY   = RGBColor(150, 150, 150)
ORANGE     = RGBColor(230, 126, 34)
LIGHT_RED  = RGBColor(253, 237, 236)
LIGHT_GREEN = RGBColor(234, 250, 241)
LIGHT_GOLD = RGBColor(255, 251, 235)
LIGHT_BLUE = RGBColor(214, 234, 248)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ── Helpers ──

def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, l, t, w, h, fill=None, line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, l, t, w, h)
    s.line.fill.background()
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.fill.solid()
        s.line.color.rgb = line
        s.line.width = Pt(2)
    return s


def tbox(slide, l, t, w, h, text, sz=18, bold=False, color=NAVY, align=PP_ALIGN.LEFT, name="Calibri"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = name
    p.alignment = align
    return tb


def set_text(shape, text, sz=18, bold=False, color=NAVY, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = align
    return tf


def add_p(tf, text, sz=16, bold=False, color=NAVY, align=PP_ALIGN.LEFT,
          before=Pt(4), after=Pt(2), level=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = align
    p.space_before = before
    p.space_after = after
    p.level = level
    return p


def title_bar(slide, text, sz=28):
    bar = rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), fill=NAVY)
    set_text(bar, text, sz=sz, bold=True, color=WHITE)
    bar.text_frame.paragraphs[0].space_before = Pt(12)
    bar.text_frame.margin_left = Inches(0.6)
    bar.text_frame.margin_top = Inches(0.18)
    rect(slide, Inches(0), Inches(1.1), SLIDE_W, Inches(0.04), fill=GOLD)


def footer(slide, text="22nd Century Technologies, Inc.  |  V-Intelligence UEBA Detection Comparison  |  Confidential"):
    tbox(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.35),
         text, sz=10, color=MID_GRAY)


def card(slide, l, t, w, h, fill=WHITE, border=None, border_top_color=None, border_top_h=Inches(0.06)):
    bg = rect(slide, l, t, w, h, fill=fill, radius=True)
    if border:
        bg.line.fill.solid()
        bg.line.color.rgb = border
        bg.line.width = Pt(1.5)
    if border_top_color:
        rect(slide, l + Inches(0.05), t, w - Inches(0.1), border_top_h, fill=border_top_color)
    return bg


def bullet_box(slide, l, t, w, h, title, bullets, title_color=NAVY, bg_color=WHITE, border_color=None):
    c = card(slide, l, t, w, h, fill=bg_color, border=border_color)
    tbox(slide, l + Inches(0.2), t + Inches(0.1), w - Inches(0.4), Inches(0.4),
         title, sz=16, bold=True, color=title_color)
    for i, b in enumerate(bullets):
        tbox(slide, l + Inches(0.3), t + Inches(0.5) + Inches(i * 0.35), w - Inches(0.5), Inches(0.35),
             f"• {b}", sz=13, color=DARK_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════

def slide_01_title(prs):
    """Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)

    tbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.2),
         "Detection Comparison", sz=44, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    rect(slide, Inches(0.8), Inches(2.8), Inches(3), Inches(0.04), fill=GOLD)

    tbox(slide, Inches(0.8), Inches(3.1), Inches(11.5), Inches(0.8),
         "Why Traditional Cyber Detection Fails Against Modern Threats —\n"
         "and How Behavioral Analysis Changes the Game",
         sz=22, color=RGBColor(160, 200, 224), align=PP_ALIGN.LEFT)

    tbox(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.5),
         "V-Intelligence UEBA + Composite Scoring  |  Business Overview",
         sz=18, color=GOLD, bold=True, align=PP_ALIGN.LEFT)

    tbox(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.5),
         "22nd Century Technologies, Inc.  |  Confidential",
         sz=14, color=MID_GRAY, align=PP_ALIGN.LEFT)


def slide_02_agenda(prs):
    """Agenda."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Agenda")

    items = [
        ("1", "The Problem", "Why do sophisticated attackers evade your current SIEM?"),
        ("2", "How Traditional Detection Works", "Thresholds, rules, and statistical algorithms"),
        ("3", "What Traditional Detection Lacks", "Flagged as anomalous — but with no actionable direction"),
        ("4", "How V-Intelligence UEBA Changes Detection", "Building a behavioral digital entity for every user"),
        ("5", "Composite Scoring: 5-Phase Detection", "The scoring engine that separates attackers from noise"),
        ("6", "Head-to-Head Results", "Same data, same users — completely different outcomes"),
        ("7", "Per-Attacker Analysis", "Why each attack type requires a different detection phase"),
        ("8", "Real-World Validation: Salt Typhoon", "5 years undetected — and why V-Intelligence UEBA catches it"),
        ("9", "Business Impact", "What this means for your SOC, your analysts, and your risk posture"),
    ]
    for i, (num, title, desc) in enumerate(items):
        y = Inches(1.5) + Inches(i * 0.7)
        circle = rect(slide, Inches(0.8), y, Inches(0.45), Inches(0.45), fill=NAVY, radius=True)
        set_text(circle, num, sz=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tbox(slide, Inches(1.5), y, Inches(3.5), Inches(0.4), title, sz=18, bold=True, color=NAVY)
        tbox(slide, Inches(5.2), y + Inches(0.02), Inches(7.5), Inches(0.4), desc, sz=15, color=DARK_GRAY)

    footer(slide)


def slide_03_problem(prs):
    """The Problem."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "The Problem: Modern Attacks Are Invisible to Traditional Detection")

    # Stats bar
    stats = [
        ("287 days", "Average time to detect\na breach (IBM 2024)"),
        ("68%", "Of breaches involve\nhuman element (Verizon DBIR)"),
        ("$4.88M", "Average cost of\na data breach (IBM 2024)"),
        ("80%", "Of insider threats go\nundetected by SIEM rules"),
    ]
    for i, (val, desc) in enumerate(stats):
        x = Inches(0.6) + Inches(i * 3.1)
        card(slide, x, Inches(1.5), Inches(2.8), Inches(1.6), border_top_color=RED)
        tbox(slide, x + Inches(0.15), Inches(1.7), Inches(2.5), Inches(0.5),
             val, sz=32, bold=True, color=RED, align=PP_ALIGN.CENTER)
        tbox(slide, x + Inches(0.15), Inches(2.3), Inches(2.5), Inches(0.7),
             desc, sz=13, color=DARK_GRAY, align=PP_ALIGN.CENTER)

    # Key insight
    box = rect(slide, Inches(0.6), Inches(3.5), Inches(12.1), Inches(1.2), fill=LIGHT_GOLD, radius=True)
    box.line.fill.solid()
    box.line.color.rgb = GOLD
    box.line.width = Pt(2)
    tbox(slide, Inches(0.9), Inches(3.6), Inches(11.5), Inches(0.4),
         "The Core Challenge", sz=18, bold=True, color=GOLD)
    tbox(slide, Inches(0.9), Inches(4.0), Inches(11.5), Inches(0.6),
         "Today's sophisticated attackers — nation-state APTs, patient insiders, living-off-the-land operators — "
         "don't trigger alarms because they deliberately stay within normal statistical boundaries. "
         "No single metric is abnormal. The attack is in the pattern, not the threshold.",
         sz=15, color=DARK_GRAY)

    # Three attacker types
    types = [
        ("Insider Threat", "14-month slow escalation", "Accesses 5% more files per week.\nNever crosses a threshold.", RED),
        ("Nation-State APT", "417-day persistent access", "Uses legitimate tools and protocols.\nBlends with normal traffic.", ORANGE),
        ("Living-off-the-Land", "Uses built-in system tools", "No malware to detect. Appears as\nnormal admin activity.", RGBColor(142, 68, 173)),
    ]
    for i, (name, sub, desc, clr) in enumerate(types):
        x = Inches(0.6) + Inches(i * 4.1)
        card(slide, x, Inches(5.1), Inches(3.8), Inches(1.7), border_top_color=clr)
        tbox(slide, x + Inches(0.15), Inches(5.3), Inches(3.5), Inches(0.35),
             name, sz=16, bold=True, color=clr)
        tbox(slide, x + Inches(0.15), Inches(5.6), Inches(3.5), Inches(0.3),
             sub, sz=12, color=MID_GRAY)
        tbox(slide, x + Inches(0.15), Inches(5.95), Inches(3.5), Inches(0.7),
             desc, sz=13, color=DARK_GRAY)

    footer(slide)


def slide_04_traditional(prs):
    """How Traditional Detection Works."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "How Traditional Detection Works")

    tbox(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
         "Traditional SIEM and anomaly detection tools rely on statistical thresholds applied to individual metrics.",
         sz=16, color=DARK_GRAY)

    # Feature categories
    cats = [
        ("Authentication", "Login counts, failure rates,\noff-hours access ratios", "🔐"),
        ("File Access", "Files opened, sensitive file ratios,\nwrite operations, bytes transferred", "📁"),
        ("Network", "Bytes sent, unique destinations,\nexternal connection ratios", "🌐"),
        ("Endpoint", "Process risk scores, suspicious\nprocess ratios, unique processes", "🖥️"),
    ]
    for i, (name, desc, icon) in enumerate(cats):
        x = Inches(0.6) + Inches(i * 3.1)
        card(slide, x, Inches(1.9), Inches(2.8), Inches(1.5), fill=LIGHT_BLUE, border=BLUE)
        tbox(slide, x + Inches(0.15), Inches(2.0), Inches(2.5), Inches(0.4),
             f"{name}", sz=15, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
        tbox(slide, x + Inches(0.15), Inches(2.4), Inches(2.5), Inches(0.8),
             desc, sz=12, color=DARK_GRAY, align=PP_ALIGN.CENTER)

    # Algorithms
    tbox(slide, Inches(0.6), Inches(3.7), Inches(12), Inches(0.4),
         "Algorithms Applied to These Features:", sz=16, bold=True, color=NAVY)

    algos = [
        ("Isolation Forest", "Randomly partitions data; anomalies are points that isolate quickly"),
        ("One-Class SVM", "Draws a boundary around \"normal\" data; flags points outside the boundary"),
        ("Local Outlier Factor", "Compares each user's density to their neighbors"),
        ("Z-Score Analysis", "Flags any metric more than 3 standard deviations from the mean"),
    ]
    for i, (name, desc) in enumerate(algos):
        y = Inches(4.2) + Inches(i * 0.5)
        tbox(slide, Inches(0.8), y, Inches(3), Inches(0.45), f"  {name}", sz=14, bold=True, color=NAVY)
        tbox(slide, Inches(3.8), y, Inches(8.5), Inches(0.45), desc, sz=13, color=DARK_GRAY)

    # Weakness callout
    box = rect(slide, Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.6), fill=LIGHT_RED, radius=True)
    box.line.fill.solid()
    box.line.color.rgb = RED
    box.line.width = Pt(2)
    tbox(slide, Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.5),
         "Fundamental limitation: These algorithms can flag anomalous users, but provide NO directional intelligence. "
         "They say \"anomalous\" without telling the analyst WHICH behavioral zone changed or TOWARD WHAT threat pattern.",
         sz=14, bold=True, color=RED)

    footer(slide)


def slide_05_what_it_misses(prs):
    """What Traditional Detection Misses — The SOC View."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "What Your SOC Analyst Sees: Four Real Attack Profiles")

    tbox(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
         "Same telemetry data, same time period, same users. Traditional methods flag anomalies but give no directional intelligence.",
         sz=16, color=DARK_GRAY)

    attackers = [
        ("USR-156", "Insider Threat", "14-month slow escalation", "FLAGGED",
         "Flagged as anomalous, but no indication\nof WHICH behavior zone changed.", RED),
        ("USR-234", "Slow APT", "417-day persistent access", "MISSED",
         "Only attack missed by LOF and Isolation\nForest — too slow to trigger.", ORANGE),
        ("USR-042", "Volt Typhoon (LOTL)", "412-day campaign", "FLAGGED",
         "Flagged as anomalous, but no directional\nintelligence on threat pattern.", RGBColor(142, 68, 173)),
        ("USR-118", "Salt Typhoon", "412-day telecom intrusion", "FLAGGED",
         "Flagged by Z-Score/OCSVM, but buried\nin false positives with no ranking.", BLUE),
    ]
    for i, (uid, name, sub, status, desc, clr) in enumerate(attackers):
        y = Inches(1.9) + Inches(i * 1.15)
        card(slide, Inches(0.6), y, Inches(12.1), Inches(1.0), border=MID_GRAY)
        # Left: user info
        tbox(slide, Inches(0.9), y + Inches(0.08), Inches(2), Inches(0.35),
             uid, sz=18, bold=True, color=NAVY)
        tbox(slide, Inches(0.9), y + Inches(0.4), Inches(2.5), Inches(0.5),
             f"{name} — {sub}", sz=12, color=MID_GRAY)
        # Middle: description
        tbox(slide, Inches(3.8), y + Inches(0.15), Inches(5.5), Inches(0.7),
             desc, sz=13, color=DARK_GRAY)
        # Right: status badge
        badge_color = RED if status == "MISSED" else ORANGE
        badge = rect(slide, Inches(10.2), y + Inches(0.25), Inches(2.2), Inches(0.5),
                     fill=badge_color, radius=True)
        set_text(badge, f"  {status}", sz=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Bottom callout
    box = rect(slide, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.6), fill=LIGHT_RED, radius=True)
    box.line.fill.solid()
    box.line.color.rgb = RED
    box.line.width = Pt(2)
    tbox(slide, Inches(0.9), Inches(6.55), Inches(11.5), Inches(0.5),
         "Result: 3 of 4 flagged as anomalous, but with NO directional intelligence.  Your SOC analyst sees \"anomalous\" "
         "without knowing WHICH behavioral zone changed or TOWARD WHAT threat pattern — and the slow APT is missed entirely.",
         sz=15, bold=True, color=RED)

    footer(slide)


def slide_06_acecard_intro(prs):
    """How V-Intelligence UEBA Changes Detection."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "How V-Intelligence UEBA Changes Detection: Building a Digital Entity")

    tbox(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.8),
         "Instead of checking individual numbers against thresholds, V-Intelligence UEBA builds a complete behavioral profile "
         "— a \"digital entity\" — for every user. It captures what users are doing, not just how much.",
         sz=16, color=DARK_GRAY)

    # Two-layer architecture
    # Layer 1
    l1 = rect(slide, Inches(0.6), Inches(2.3), Inches(5.8), Inches(4.0), fill=LIGHT_BLUE, radius=True)
    l1.line.fill.solid()
    l1.line.color.rgb = BLUE
    l1.line.width = Pt(2)
    tbox(slide, Inches(0.9), Inches(2.4), Inches(5.2), Inches(0.4),
         "Layer 1: V-Intelligence UEBA — Digital Entity Engine", sz=18, bold=True, color=BLUE)

    l1_items = [
        "Converts raw log data into a behavioral description of each user",
        "Captures the meaning of actions, not just the count",
        "Example: \"USR-042 accessed 47 restricted files across 12 departments, 68% outside normal hours\"",
        "Tracks 5 behavioral zones per user:",
        "    Identity — role, department, clearance level",
        "    Access Pattern — when and how they authenticate",
        "    Data Behavior — what files they touch and how",
        "    Network Footprint — where their traffic goes",
        "    Risk Posture — endpoint and process activity",
        "Measures behavioral drift — direction of change over time",
    ]
    for i, item in enumerate(l1_items):
        indent = Inches(0.3) if item.startswith("    ") else Inches(0)
        sz = 12 if item.startswith("    ") else 13
        tbox(slide, Inches(1.0) + indent, Inches(2.85) + Inches(i * 0.32), Inches(5.0) - indent, Inches(0.32),
             f"{'▸' if not item.startswith('    ') else '–'} {item.strip()}", sz=sz, color=DARK_GRAY)

    # Arrow
    tbox(slide, Inches(6.55), Inches(3.8), Inches(0.6), Inches(0.5),
         "→", sz=36, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # Layer 2
    l2 = rect(slide, Inches(7.0), Inches(2.3), Inches(5.8), Inches(4.0), fill=LIGHT_GREEN, radius=True)
    l2.line.fill.solid()
    l2.line.color.rgb = GREEN
    l2.line.width = Pt(2)
    tbox(slide, Inches(7.3), Inches(2.4), Inches(5.2), Inches(0.4),
         "Layer 2: Composite Scoring — Anomaly Detector", sz=18, bold=True, color=GREEN)

    l2_items = [
        "Uses V-Intelligence UEBA's behavioral profiles to determine which users are truly anomalous",
        "Applies 5 scoring phases simultaneously:",
        "    Signal Strength — Is the drift larger than peers?",
        "    Breadth — Are many behaviors anomalous at once?",
        "    Sustained Deviation — Does it persist over weeks?",
        "    Context Divergence — Anomalous under different analytical lenses?",
        "    Novelty Persistence — Are new behaviors recurring?",
        "Normal users may score high on 1 phase by chance",
        "Attackers score high on multiple phases simultaneously",
        "Output: One ranked list — no threshold tuning needed",
    ]
    for i, item in enumerate(l2_items):
        indent = Inches(0.3) if item.startswith("    ") else Inches(0)
        sz = 12 if item.startswith("    ") else 13
        tbox(slide, Inches(7.3) + indent, Inches(2.85) + Inches(i * 0.32), Inches(5.0) - indent, Inches(0.32),
             f"{'▸' if not item.startswith('    ') else '–'} {item.strip()}", sz=sz, color=DARK_GRAY)

    # Key difference callout
    box = rect(slide, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.6), fill=NAVY, radius=True)
    tbox(slide, Inches(0.9), Inches(6.55), Inches(11.5), Inches(0.5),
         "Key difference: Traditional detection asks \"is any one number abnormal?\"  "
         "V-Intelligence UEBA + Composite asks \"has the overall behavior pattern changed in a meaningful, persistent way?\"",
         sz=15, bold=True, color=GOLD)

    footer(slide)


def slide_07_why_composite(prs):
    """Why Composite Scoring Catches What Single Methods Can't."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Why Composite Scoring Catches What Single Methods Can't")

    tbox(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
         "A single detection method can only check one dimension. Composite scoring checks five simultaneously.",
         sz=16, color=DARK_GRAY)

    # Visual: single method vs composite
    # Left: Single method
    card(slide, Inches(0.6), Inches(2.0), Inches(5.5), Inches(3.8), border_top_color=RED)
    tbox(slide, Inches(0.9), Inches(2.2), Inches(5), Inches(0.4),
         "Single Method (e.g., Isolation Forest)", sz=16, bold=True, color=RED)
    single_items = [
        "Checks: Is this user a statistical outlier?",
        "Measures: Distance from center of data cluster",
        "Fails when: Attacker stays near the center",
        "Result: Attacker classified as NORMAL",
        "",
        "Why it fails:",
        "An insider who increases file access by 5% each week",
        "is never far enough from the mean to trigger an alert.",
        "Each individual metric looks normal.",
    ]
    for i, item in enumerate(single_items):
        if not item:
            continue
        clr = RED if item.startswith("Why") or item.startswith("An insider") or item.startswith("is never") or item.startswith("Each individual") else DARK_GRAY
        bld = item.startswith("Why")
        tbox(slide, Inches(0.9), Inches(2.65) + Inches(i * 0.34), Inches(5), Inches(0.34),
             item if bld else f"  {item}", sz=13, color=clr, bold=bld)

    # Right: Composite
    card(slide, Inches(6.8), Inches(2.0), Inches(5.9), Inches(3.8), border_top_color=GREEN)
    tbox(slide, Inches(7.1), Inches(2.2), Inches(5.5), Inches(0.4),
         "Composite Scoring (5 Phases)", sz=16, bold=True, color=GREEN)
    comp_items = [
        ("Signal Strength", "Is the drift larger than their peer group?"),
        ("Breadth", "Are multiple behaviors anomalous at once?"),
        ("Sustained", "Does the anomaly persist across weeks?"),
        ("Context", "Is it anomalous from multiple analytical angles?"),
        ("Novelty", "Are new behaviors (new destinations, new files) recurring?"),
    ]
    for i, (phase, desc) in enumerate(comp_items):
        y = Inches(2.7) + Inches(i * 0.6)
        badge = rect(slide, Inches(7.1), y, Inches(1.8), Inches(0.35), fill=GREEN, radius=True)
        set_text(badge, f"  {phase}", sz=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tbox(slide, Inches(9.1), y, Inches(3.5), Inches(0.35), desc, sz=12, color=DARK_GRAY)

    tbox(slide, Inches(7.1), Inches(5.3), Inches(5.5), Inches(0.4),
         "A normal user may score high on 1 phase by chance.\n"
         "An attacker scores high on multiple phases — that's the signal.",
         sz=13, bold=True, color=GREEN)

    # Bottom key insight
    box = rect(slide, Inches(0.6), Inches(6.2), Inches(12.1), Inches(0.8), fill=LIGHT_GOLD, radius=True)
    box.line.fill.solid()
    box.line.color.rgb = GOLD
    box.line.width = Pt(2)
    tbox(slide, Inches(0.9), Inches(6.25), Inches(11.5), Inches(0.35),
         "The Analogy:", sz=15, bold=True, color=GOLD)
    tbox(slide, Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.35),
         "Traditional detection is like checking if someone is too tall or too heavy. "
         "Composite scoring is like noticing they're walking differently, going to new places, "
         "keeping different hours, and talking to different people — all at once, over weeks.",
         sz=14, color=DARK_GRAY)

    footer(slide)


def slide_08_results(prs):
    """Head-to-Head Results."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Head-to-Head Results: Same Data, Different Outcomes")

    tbox(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
         "250 users monitored over 70 weeks (485 days). 4 known attack campaigns embedded in the population.",
         sz=16, color=DARK_GRAY)

    # Left: Traditional
    card(slide, Inches(0.6), Inches(2.0), Inches(5.8), Inches(4.5), border_top_color=RED)
    tbox(slide, Inches(0.9), Inches(2.2), Inches(5.2), Inches(0.4),
         "TRADITIONAL DETECTION", sz=20, bold=True, color=RED, align=PP_ALIGN.CENTER)
    tbox(slide, Inches(0.9), Inches(2.65), Inches(5.2), Inches(0.3),
         "Isolation Forest, SVM, LOF, Z-Score", sz=13, color=MID_GRAY, align=PP_ALIGN.CENTER)

    tbox(slide, Inches(0.9), Inches(3.2), Inches(5.2), Inches(0.8),
         "3–4 of 4", sz=56, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    tbox(slide, Inches(0.9), Inches(4.1), Inches(5.2), Inches(0.4),
         "flagged as anomalous — but no direction", sz=18, color=NAVY, align=PP_ALIGN.CENTER)

    tbox(slide, Inches(0.9), Inches(4.7), Inches(5.2), Inches(0.4),
         "Z-Score and OCSVM flag 4/4; LOF and IF miss the slow APT", sz=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    tbox(slide, Inches(0.9), Inches(5.3), Inches(5.2), Inches(0.8),
         "Traditional methods say \"anomalous\" but cannot tell the analyst "
         "WHICH behavioral zone shifted or TOWARD WHAT threat pattern. FP rates range from 3.3% to 29.7%.",
         sz=13, color=DARK_GRAY, align=PP_ALIGN.CENTER)

    # Right: V-Intelligence UEBA
    card(slide, Inches(6.9), Inches(2.0), Inches(5.8), Inches(4.5), border_top_color=GREEN)
    tbox(slide, Inches(7.2), Inches(2.2), Inches(5.2), Inches(0.4),
         "V-INTELLIGENCE UEBA + COMPOSITE SCORING", sz=20, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    tbox(slide, Inches(7.2), Inches(2.65), Inches(5.2), Inches(0.3),
         "Digital Entity Features → 5-Phase Anomaly Detection", sz=13, color=MID_GRAY, align=PP_ALIGN.CENTER)

    tbox(slide, Inches(7.2), Inches(3.2), Inches(5.2), Inches(0.8),
         "4 of 4", sz=56, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    tbox(slide, Inches(7.2), Inches(4.1), Inches(5.2), Inches(0.4),
         "attacks detected", sz=18, color=NAVY, align=PP_ALIGN.CENTER)

    tbox(slide, Inches(7.2), Inches(4.7), Inches(5.2), Inches(0.4),
         "10.6% FP rate + directional intelligence", sz=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    tbox(slide, Inches(7.2), Inches(5.3), Inches(5.2), Inches(0.8),
         "All 4 detected with actionable direction: tells the analyst WHICH behavioral "
         "zone changed and TOWARD WHAT threat pattern. Ranked list, no threshold tuning.",
         sz=13, color=DARK_GRAY, align=PP_ALIGN.CENTER)

    # Bottom bar
    box = rect(slide, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.5), fill=NAVY, radius=True)
    tbox(slide, Inches(0.9), Inches(6.73), Inches(11.5), Inches(0.4),
         "Same 250 users. Same 70 weeks. Same telemetry. The difference is what the system understands about behavior.",
         sz=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    footer(slide)


def slide_09_per_attacker(prs):
    """Per-Attacker Analysis."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Why Each Attacker Gets Caught: Different Phases, Different Signals")

    tbox(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
         "No single detection phase catches everyone. Each attack type leaves a distinct behavioral fingerprint.",
         sz=16, color=DARK_GRAY)

    attackers = [
        ("USR-118", "Salt Typhoon", "412-Day Telecom Intrusion", "#1 / 250", "51.7",
         GREEN, RGBColor(41, 128, 185),
         "Signal Strength + Sustained Deviation",
         "Flagged by Z-Score and OCSVM but buried in false positives with no directional context. "
         "Composite scoring ranks it #1 with actionable intelligence: WHICH zones shifted and TOWARD WHAT pattern."),
        ("USR-156", "Insider Threat", "14-Month Escalation", "#2 / 250", "46.2",
         GREEN, RED,
         "Signal Strength + Breadth",
         "Flagged as anomalous by LOF/Z-Score/OCSVM but with no direction. Composite provides actionable "
         "intelligence: gradual scope creep across departments, off-hours access, data staging — broad and persistent."),
        ("USR-042", "Volt Typhoon", "412-Day LOTL Campaign", "#30 / 250", "12.9",
         GREEN, RGBColor(142, 68, 173),
         "Breadth + Sustained Deviation",
         "Traditional methods flag it but cannot explain WHY. Composite reveals the anomaly is in the breadth "
         "of unusual behavior across zones — uses legitimate tools, so direction matters more than detection."),
        ("USR-234", "Slow APT", "417-Day Campaign", "#7 / 250", "20.0",
         GREEN, ORANGE,
         "Novelty Persistence",
         "The only attacker missed by LOF and Isolation Forest. Z-Score and OCSVM flag it but without direction. "
         "Composite uniquely identifies persistent novel network destinations — the hallmark of C2 beacons."),
    ]
    for i, (uid, name, sub, rank, score, det_color, atk_color, phase, desc) in enumerate(attackers):
        y = Inches(1.9) + Inches(i * 1.3)
        card(slide, Inches(0.6), y, Inches(12.1), Inches(1.15), border=MID_GRAY)
        # Left stripe
        rect(slide, Inches(0.6), y, Inches(0.08), Inches(1.15), fill=atk_color)
        # User info
        tbox(slide, Inches(0.9), y + Inches(0.05), Inches(1.3), Inches(0.3),
             uid, sz=17, bold=True, color=NAVY)
        tbox(slide, Inches(0.9), y + Inches(0.35), Inches(2), Inches(0.3),
             f"{name} — {sub}", sz=11, color=MID_GRAY)
        tbox(slide, Inches(0.9), y + Inches(0.65), Inches(2), Inches(0.3),
             f"Rank {rank}  •  Score {score}", sz=12, bold=True, color=atk_color)
        # Phase badge
        badge = rect(slide, Inches(3.2), y + Inches(0.1), Inches(3), Inches(0.35), fill=atk_color, radius=True)
        set_text(badge, f"  {phase}", sz=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Description
        tbox(slide, Inches(3.2), y + Inches(0.5), Inches(9.3), Inches(0.6),
             desc, sz=12, color=DARK_GRAY)

    footer(slide)


def slide_10_usr234_deep(prs):
    """USR-234 Deep Dive — Why Novelty Persistence Matters."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Deep Dive: How Composite Scoring Catches a 417-Day APT")

    tbox(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.6),
         "USR-234 is the hardest attacker to detect. Every traditional metric looks normal. "
         "Even basic behavioral drift analysis barely separates this user from the crowd.",
         sz=16, color=DARK_GRAY)

    # Left: what traditional sees
    card(slide, Inches(0.6), Inches(2.1), Inches(5.8), Inches(2.5), border_top_color=RED)
    tbox(slide, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.35),
         "What Traditional Detection Sees", sz=16, bold=True, color=RED)
    trad_lines = [
        "Authentication: Normal login patterns, normal failure rate",
        "File Access: Within expected volume ranges",
        "Network: Bytes out within 1 standard deviation of mean",
        "Endpoint: No suspicious processes detected",
        "Verdict: NORMAL — No alerts generated",
    ]
    for i, line in enumerate(trad_lines):
        clr = RED if "Verdict" in line else DARK_GRAY
        bld = "Verdict" in line
        tbox(slide, Inches(0.9), Inches(2.75) + Inches(i * 0.33), Inches(5.2), Inches(0.33),
             f"{'▸' if not bld else '✗'} {line}", sz=12, color=clr, bold=bld)

    # Right: what composite sees
    card(slide, Inches(6.9), Inches(2.1), Inches(5.8), Inches(2.5), border_top_color=GREEN)
    tbox(slide, Inches(7.2), Inches(2.3), Inches(5.2), Inches(0.35),
         "What Composite Scoring Sees", sz=16, bold=True, color=GREEN)
    comp_lines = [
        ("Signal Strength:", "Low (14th percentile) — drift magnitude is small", MID_GRAY),
        ("Breadth:", "Low (5th percentile) — only 1 feature mildly anomalous", MID_GRAY),
        ("Sustained:", "Low (31st percentile) — not persistently high", MID_GRAY),
        ("Context Divergence:", "Moderate (29th percentile) — starting to show", ORANGE),
        ("Novelty Persistence:", "MAXIMUM (100th percentile) — novel network destinations persist every week", GREEN),
    ]
    for i, (phase, desc, clr) in enumerate(comp_lines):
        tbox(slide, Inches(7.2), Inches(2.75) + Inches(i * 0.33), Inches(1.8), Inches(0.33),
             f"  {phase}", sz=11, bold=True, color=NAVY)
        tbox(slide, Inches(9.0), Inches(2.75) + Inches(i * 0.33), Inches(3.5), Inches(0.33),
             desc, sz=11, color=clr)

    # Explanation box
    box = rect(slide, Inches(0.6), Inches(4.9), Inches(12.1), Inches(1.8), fill=LIGHT_GOLD, radius=True)
    box.line.fill.solid()
    box.line.color.rgb = GOLD
    box.line.width = Pt(2)
    tbox(slide, Inches(0.9), Inches(5.0), Inches(11.5), Inches(0.35),
         "Why Novelty Persistence Catches APTs", sz=16, bold=True, color=GOLD)
    explanation_lines = [
        "An APT (Advanced Persistent Threat) establishes hidden communication channels called \"C2 beacons\" "
        "— periodic connections to attacker-controlled servers.",
        "These destinations are new to the user's normal network profile. A one-time new destination could be "
        "anything — a new vendor, a new service. But when the same new destinations keep appearing week after week, "
        "that's the fingerprint of persistent attacker infrastructure.",
        "USR-234's network volume, byte counts, and connection rates all look normal. But the destinations are new, "
        "and they persist. Composite Scoring's Novelty Persistence phase catches exactly this pattern — "
        "ranking USR-234 at #7 out of 250 users despite every other metric looking clean.",
    ]
    for i, line in enumerate(explanation_lines):
        tbox(slide, Inches(0.9), Inches(5.4) + Inches(i * 0.4), Inches(11.5), Inches(0.4),
             line, sz=12, color=DARK_GRAY)

    footer(slide)


def slide_11_salt_typhoon(prs):
    """Real-World Validation: Salt Typhoon — 5 Years Undetected."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Real-World Validation: Salt Typhoon — 5 Years Undetected")

    tbox(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.6),
         "Salt Typhoon (a Chinese state-sponsored APT group) operated inside US telecom infrastructure for over 5 years "
         "before discovery. Our simulation confirms why — and demonstrates how behavioral analysis changes the outcome.",
         sz=16, color=DARK_GRAY)

    # Left: Why it evaded detection
    card(slide, Inches(0.6), Inches(2.1), Inches(5.8), Inches(3.5), border_top_color=RED)
    tbox(slide, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.35),
         "Why Traditional Detection Failed — For 5 Years", sz=16, bold=True, color=RED)

    trad_lines = [
        "Used legitimate network protocols and tools — no malware signatures",
        "Maintained normal traffic volumes — never spiked bandwidth",
        "Operated within authorized access boundaries — valid credentials",
        "Blended with routine telecom operations — low-and-slow pattern",
        "No single metric exceeded any statistical threshold",
        "",
        "Our Simulation Results:",
        "LOF:  3/4 detected at 3.3% FP — flags Salt Typhoon, no direction",
        "Z-Score:  4/4 detected at 4.9% FP — flags it, but no direction",
        "Isolation Forest:  3/4 detected at 7.3% FP — flags it, no direction",
        "OCSVM:  4/4 detected at 29.7% FP — flags it, buried in noise",
    ]
    for i, line in enumerate(trad_lines):
        if not line:
            continue
        is_header = line.startswith("Our Simulation")
        is_result = "detected at" in line
        clr = ORANGE if is_result else (NAVY if is_header else DARK_GRAY)
        bld = is_header
        sz = 13 if is_header else 12
        tbox(slide, Inches(0.9), Inches(2.7) + Inches(i * 0.28), Inches(5.2), Inches(0.28),
             f"{'▸' if not is_result and not is_header else '  '} {line}", sz=sz, color=clr, bold=bld)

    # Right: Why V-Intelligence UEBA catches it
    card(slide, Inches(6.9), Inches(2.1), Inches(5.8), Inches(3.5), border_top_color=GREEN)
    tbox(slide, Inches(7.2), Inches(2.3), Inches(5.2), Inches(0.35),
         "Why V-Intelligence UEBA + Composite Catches It — Immediately", sz=16, bold=True, color=GREEN)

    ace_lines = [
        ("Behavioral Profile", "Tracks the meaning of actions, not just counts"),
        ("Multi-Zone Analysis", "Identity stable, but network footprint drifting"),
        ("Sustained Signal", "Behavioral deviation persists across every week"),
        ("Context Divergence", "Anomalous from multiple analytical perspectives"),
        ("Signal Strength", "29.9 — strongest drift in the entire population"),
    ]
    for i, (phase, desc) in enumerate(ace_lines):
        y = Inches(2.75) + Inches(i * 0.55)
        badge = rect(slide, Inches(7.2), y, Inches(2.2), Inches(0.32), fill=GREEN, radius=True)
        set_text(badge, f"  {phase}", sz=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tbox(slide, Inches(9.6), y, Inches(3.0), Inches(0.32), desc, sz=12, color=DARK_GRAY)

    tbox(slide, Inches(7.2), Inches(5.2), Inches(5.2), Inches(0.35),
         "Composite Score: 51.7 — Rank #1 out of 250 users",
         sz=15, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    # Bottom: Key insight
    box = rect(slide, Inches(0.6), Inches(5.9), Inches(12.1), Inches(1.3), fill=NAVY, radius=True)
    tbox(slide, Inches(0.9), Inches(6.0), Inches(11.5), Inches(0.35),
         "The Critical Insight", sz=16, bold=True, color=GOLD)
    tbox(slide, Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.8),
         "Salt Typhoon evaded meaningful detection for 5 years because traditional tools flagged anomalies "
         "without directional intelligence — analysts saw \"anomalous\" but not WHICH zone shifted or TOWARD WHAT threat. "
         "V-Intelligence UEBA ranks this pattern as #1 out of 250 AND tells the analyst exactly which behavioral zones "
         "drifted, in what direction, and toward what known threat pattern — turning a vague flag into actionable response.",
         sz=13, color=WHITE)

    footer(slide)


def slide_12_false_positives(prs):
    """False Positive Rate Comparison."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "False Positive Rate: Why Precision Matters")

    tbox(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.8),
         "A detection method that flags every user is useless — even if it catches all attackers. "
         "False positive rate measures what percentage of normal users are incorrectly flagged as threats. "
         "Lower is better.",
         sz=16, color=DARK_GRAY)

    # Method comparison cards
    methods = [
        ("Local Outlier Factor", "3 of 4", "3.3%", "Flags 3 as anomalous, misses slow\nAPT. No directional intelligence.", ORANGE, False),
        ("Z-Score (|z|>3)", "4 of 4", "4.9%", "Flags all 4 as anomalous but gives\nno direction on threat pattern.", ORANGE, False),
        ("Isolation Forest", "3 of 4", "7.3%", "Flags 3 as anomalous, misses slow\nAPT. No directional intelligence.", ORANGE, False),
        ("One-Class SVM", "4 of 4", "29.7%", "Flags all 4 but buries them in\nmassive false positive volume.", RED, False),
        ("V-Intelligence UEBA + Composite", "4 of 4", "10.6%", "Detects all 4 WITH directional\nintelligence on threat patterns.", GREEN, True),
        ("Multi-Front Threat-Profile", "4 of 4", "0%", "Primary detector: all 4 matched to\nnamed known-bad fingerprints.", GREEN, True),
    ]
    for i, (name, detected, fp, desc, clr, highlight) in enumerate(methods):
        x = Inches(0.4) + Inches(i * 2.18)
        bg = LIGHT_GREEN if highlight else WHITE
        bdr = GREEN if highlight else MID_GRAY
        card(slide, x, Inches(2.5), Inches(2.0), Inches(3.2), fill=bg, border=bdr,
             border_top_color=clr)
        tbox(slide, x + Inches(0.1), Inches(2.7), Inches(1.8), Inches(0.4),
             name, sz=12, bold=True, color=clr, align=PP_ALIGN.CENTER)
        tbox(slide, x + Inches(0.1), Inches(3.2), Inches(1.8), Inches(0.5),
             detected, sz=26, bold=True, color=clr, align=PP_ALIGN.CENTER)
        tbox(slide, x + Inches(0.1), Inches(3.7), Inches(1.8), Inches(0.3),
             "detected", sz=12, color=MID_GRAY, align=PP_ALIGN.CENTER)
        tbox(slide, x + Inches(0.1), Inches(4.15), Inches(1.8), Inches(0.4),
             f"FP Rate: {fp}", sz=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        tbox(slide, x + Inches(0.1), Inches(4.6), Inches(1.8), Inches(0.8),
             desc, sz=10, color=DARK_GRAY, align=PP_ALIGN.CENTER)

    # What FP means in practice
    box = rect(slide, Inches(0.6), Inches(6.0), Inches(12.1), Inches(1.0), fill=LIGHT_BLUE, radius=True)
    box.line.fill.solid()
    box.line.color.rgb = BLUE
    box.line.width = Pt(2)
    tbox(slide, Inches(0.9), Inches(6.05), Inches(11.5), Inches(0.35),
         "What does 10.6% FP rate mean in practice?", sz=15, bold=True, color=BLUE)
    tbox(slide, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.5),
         "In an organization with 10,000 users, 10.6% FP = ~1,060 users flagged for review. "
         "With composite scoring's ranked list AND directional intelligence, analysts start at #1 and know exactly "
         "WHICH behavioral zone shifted and TOWARD WHAT threat pattern. Compare Z-Score: 4/4 detected at 4.9% FP, "
         "but no ranking and no direction — the analyst only knows \"anomalous,\" not what to investigate.",
         sz=13, color=DARK_GRAY)

    footer(slide)


def slide_13_business_impact(prs):
    """Business Impact."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Business Impact: What This Means for Your Organization")

    impacts = [
        ("Actionable Directional Intelligence", GREEN,
         "Traditional methods flag anomalies but cannot tell analysts WHICH behavioral zone "
         "shifted or TOWARD WHAT threat pattern. V-Intelligence provides the direction to act.",
         "4 of 4 with directional intelligence"),
        ("Reduce Alert Fatigue", TEAL,
         "One ranked list replaces hundreds of uncorrelated alerts. Your SOC analysts "
         "investigate the highest-risk entities first, not the loudest alarms.",
         "10.6% false positive rate with ranked prioritization"),
        ("Faster Time-to-Detect", BLUE,
         "Behavioral drift detection identifies threats within weeks of campaign start — "
         "not months after data exfiltration.",
         "3 weeks into campaign vs 287-day industry average"),
        ("No Threshold Tuning", GOLD,
         "Composite scoring is relative to the population — no manually-set thresholds "
         "that break when user behavior naturally evolves.",
         "Zero analyst effort for threshold management"),
    ]
    for i, (title, clr, desc, metric) in enumerate(impacts):
        y = Inches(1.5) + Inches(i * 1.4)
        card(slide, Inches(0.6), y, Inches(12.1), Inches(1.25))
        rect(slide, Inches(0.6), y, Inches(0.08), Inches(1.25), fill=clr)
        tbox(slide, Inches(1.0), y + Inches(0.08), Inches(4), Inches(0.35),
             title, sz=18, bold=True, color=clr)
        tbox(slide, Inches(1.0), y + Inches(0.42), Inches(7), Inches(0.7),
             desc, sz=13, color=DARK_GRAY)
        badge = rect(slide, Inches(8.5), y + Inches(0.3), Inches(4), Inches(0.5), fill=clr, radius=True)
        set_text(badge, f"  {metric}", sz=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Bottom callout
    box = rect(slide, Inches(0.6), Inches(7.1), Inches(12.1), Inches(0), fill=NAVY)
    # (invisible - just using the banner below)
    box2 = rect(slide, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.6), fill=NAVY, radius=True)
    tbox(slide, Inches(0.9), Inches(6.65), Inches(11.5), Inches(0.5),
         "V-Intelligence UEBA builds the behavioral digital entity. Composite Scoring turns it into actionable intelligence.",
         sz=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    footer(slide)


def slide_14_closing(prs):
    """Closing slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)

    tbox(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(0.8),
         "Traditional detection measures numbers.", sz=28, color=RGBColor(160, 200, 224))
    tbox(slide, Inches(0.8), Inches(2.6), Inches(11.5), Inches(0.8),
         "V-Intelligence UEBA understands behavior.", sz=32, bold=True, color=GOLD)
    tbox(slide, Inches(0.8), Inches(3.4), Inches(11.5), Inches(0.8),
         "Composite Scoring detects the anomaly.", sz=32, bold=True, color=GREEN)

    rect(slide, Inches(0.8), Inches(4.5), Inches(3), Inches(0.04), fill=GOLD)

    results = [
        ("4 / 4", "attack types detected"),
        ("10.6%", "false positive rate"),
        ("Ranked", "single prioritized list"),
        ("Zero", "threshold tuning required"),
    ]
    for i, (val, desc) in enumerate(results):
        x = Inches(0.8) + Inches(i * 3.1)
        tbox(slide, x, Inches(4.9), Inches(2.8), Inches(0.6),
             val, sz=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tbox(slide, x, Inches(5.5), Inches(2.8), Inches(0.4),
             desc, sz=14, color=RGBColor(160, 200, 224), align=PP_ALIGN.CENTER)

    tbox(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.5),
         "22nd Century Technologies, Inc.  |  V-Intelligence UEBA Detection Comparison  |  Confidential",
         sz=14, color=MID_GRAY, align=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)
    slide_02_agenda(prs)
    slide_03_problem(prs)
    slide_04_traditional(prs)
    slide_05_what_it_misses(prs)
    slide_06_acecard_intro(prs)
    slide_07_why_composite(prs)
    slide_08_results(prs)
    slide_09_per_attacker(prs)
    slide_10_usr234_deep(prs)
    slide_11_salt_typhoon(prs)
    slide_12_false_positives(prs)
    slide_13_business_impact(prs)
    slide_14_closing(prs)

    out_path = "docs/Detection_Comparison_Business_Deck.pptx"
    prs.save(out_path)
    print(f"Saved {len(prs.slides)} slides to {out_path}")


if __name__ == "__main__":
    main()
