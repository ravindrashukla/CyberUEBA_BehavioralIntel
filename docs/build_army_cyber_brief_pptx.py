"""Army Cyber briefing — PowerPoint (content slides only, no title slide).
Run:  python docs/build_army_cyber_brief_pptx.py
Out:  docs/VIntelligence_Army_Cyber_Brief.pptx

6 content slides: problem, Salt/Volt difficulty, why current tools fall short,
how we solved it, results, preemptive cyber + advanced use cases.
Results are SYNTHETIC and labeled as such.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
NAVY = RGBColor(0x0D, 0x1B, 0x2A)
BLUE = RGBColor(0x1B, 0x4F, 0x72)
RED = RGBColor(0xB9, 0x1C, 0x1C)
GREEN = RGBColor(0x1E, 0x8A, 0x49)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x55, 0x55, 0x55)
LBLUE = RGBColor(0xEA, 0xF2, 0xFB)
LRED = RGBColor(0xFD, 0xEC, 0xEA)
LGREEN = RGBColor(0xE6, 0xF4, 0xEA)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def _title_band(s, num, title):
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.05))
    band.fill.solid(); band.fill.fore_color.rgb = NAVY; band.line.fill.background()
    band.shadow.inherit = False
    tf = band.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.45); tf.margin_right = Inches(0.4)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = f"{num}    {title}"
    r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = WHITE


def _bullets(s, items, top=1.35, height=4.7, intro=None, size=17):
    box = s.shapes.add_textbox(Inches(0.55), Inches(top), Inches(12.2), Inches(height))
    tf = box.text_frame; tf.word_wrap = True
    first = True
    if intro:
        p = tf.paragraphs[0]; first = False
        r = p.add_run(); r.text = intro; r.font.size = Pt(15); r.font.italic = True; r.font.color.rgb = GRAY
        p.space_after = Pt(10)
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.space_after = Pt(9)
        rb = p.add_run(); rb.text = "▪  "; rb.font.size = Pt(size); rb.font.color.rgb = BLUE; rb.font.bold = True
        if isinstance(it, tuple):
            r1 = p.add_run(); r1.text = it[0]; r1.font.size = Pt(size); r1.font.bold = True; r1.font.color.rgb = NAVY
            r2 = p.add_run(); r2.text = it[1]; r2.font.size = Pt(size); r2.font.color.rgb = RGBColor(0x2D, 0x37, 0x48)
        else:
            r = p.add_run(); r.text = it; r.font.size = Pt(size); r.font.color.rgb = RGBColor(0x2D, 0x37, 0x48)
    return box


def _callout(s, text, fill=LBLUE, txt=NAVY, top=6.35):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(top), Inches(12.2), Inches(0.85))
    box.fill.solid(); box.fill.fore_color.rgb = fill; box.line.color.rgb = txt; box.line.width = Pt(1)
    box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.3); tf.margin_right = Inches(0.3)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = txt


def slide(num, title, bullets, intro=None, callout=None, callout_fill=LBLUE, callout_txt=NAVY):
    s = prs.slides.add_slide(BLANK)
    _title_band(s, num, title)
    _bullets(s, bullets, intro=intro)
    if callout:
        _callout(s, callout, fill=callout_fill, txt=callout_txt)
    return s


# ---------- 1. The Problem ----------
slide("1", "The most dangerous threats stay inside “normal”",
      intro="The intrusions and insiders that do the most damage keep every individual signal inside "
            "its normal range. The damage is in the DIRECTION and accumulation of behavioral change "
            "— what an entity is becoming — not in any threshold being crossed.",
      bullets=[
          ("Insider: ", "logs in to different things, not more — shifts public→restricted content at the same volume."),
          ("Nation-state: ", "beacons a few times a day for months with valid tools and credentials — no malware, nothing loud."),
          ("Per-metric tools ", "score each signal on its own axis, so a change of KIND at constant volume is invisible."),
          ("The result: ", "the threat is missed, or it's an undifferentiated “anomalous” alert with no direction — alert fatigue, manual triage."),
      ],
      callout="Defenders ask “is any single metric abnormal?”  The right question is “what is this entity becoming?”")

# ---------- 2. Salt & Volt Typhoon ----------
slide("2", "Why Salt & Volt Typhoon are so hard to detect",
      intro="These PRC state-sponsored campaigns are the defining example — and exactly what V-Intelligence is built to surface.",
      bullets=[
          ("Volt Typhoon: ", "living-off-the-land in U.S. critical infrastructure — UNDETECTED for at least 5 years (CISA AA24-038A)."),
          ("Salt Typhoon: ", "deep, persistent compromise of U.S. telecommunications infrastructure (U.S. Government reporting, 2024–25)."),
          ("No signature: ", "living-off-the-land — built-in tools, valid accounts, constant volume; nothing breaches a threshold."),
          ("The only signal ", "is a slow change in behavioral DIRECTION over months — which per-metric tools never compute."),
      ],
      callout="A five-year dwell time is the predictable result of watching magnitude instead of direction.",
      callout_fill=LRED, callout_txt=RED)

# ---------- 3. Why current tools fall short ----------
slide("3", "Why current tools are not good enough",
      bullets=[
          ("Commercial UEBA ", "(Exabeam, Securonix, Microsoft Sentinel): statistical / p-value baselining, rarity- & magnitude-based scores, count/TF-IDF peer grouping."),
          ("Their GenAI ", "features are SOC copilots and conversational layers — assistants, not the detection engine."),
          ("Detection gap: ", "on a controlled benchmark of 4 advanced campaigns, magnitude-based detectors caught at most 1 of 4."),
          ("Intelligence gap: ", "even when one fires, the alert says only “anomalous” — no zone, no direction, no MITRE mapping."),
      ],
      callout="Two gaps, one root cause: they miss the stealthy threats AND can't explain the rare hit.")

# ---------- 4. How we solved it ----------
slide("4", "How we solved it — a digital twin of every entity",
      bullets=[
          ("Digital twin: ", "every entity (user, device, application, segment, session) is a behavioral profile tracked as a trajectory over time."),
          ("Behavior-as-language: ", "serialize behavior to text and embed it into one unified semantic space — LLM representational power, coupled with rigorous math."),
          ("Magnitude AND direction: ", "detect the DIRECTION of drift and project it onto named threat concepts mapped to MITRE ATT&CK — explainable."),
          ("Fused & evasion-resistant: ", "peer-cohort comparison + multi-phase fusion → one ranked, explained verdict."),
      ],
      callout="An LLM uses that space to understand language; V-Intelligence uses it to understand behavior.",
      callout_fill=LGREEN, callout_txt=GREEN)

# ---------- 5. Results (table) ----------
s = prs.slides.add_slide(BLANK)
_title_band(s, "5", "Results (controlled synthetic evaluation)")
ib = s.shapes.add_textbox(Inches(0.55), Inches(1.3), Inches(12.2), Inches(0.7))
itf = ib.text_frame; itf.word_wrap = True
ip = itf.paragraphs[0]; ir = ip.add_run()
ir.text = "Evaluation: 250 entities over ~485 days, with 4 embedded advanced campaigns — a patient insider, a slow C2 intrusion, and Volt- and Salt-Typhoon-class nation-state campaigns."
ir.font.size = Pt(14); ir.font.italic = True; ir.font.color.rgb = GRAY
# table
rows, cols = 4, 4
tbl_sh = s.shapes.add_table(rows, cols, Inches(0.55), Inches(2.25), Inches(12.2), Inches(2.2))
tbl = tbl_sh.table
hdr = ["Approach", "Campaigns caught (of 4)", "False-positive rate", "Directional explanation"]
data = [
    ["Traditional magnitude-based\n(LOF, IForest, OCSVM, Z-Score)", "0 – 1 of 4", "4.5% – 14.6%", "None"],
    ["V-Intelligence composite", "4 of 4", "10.6%", "Zone + direction + MITRE"],
    ["V-Intelligence threat-profile detector", "4 of 4", "0%", "Named technique"],
]
for j, h in enumerate(hdr):
    c = tbl.cell(0, j); c.text = h
    c.fill.solid(); c.fill.fore_color.rgb = NAVY
    pr = c.text_frame.paragraphs[0]; pr.runs[0].font.size = Pt(13); pr.runs[0].font.bold = True; pr.runs[0].font.color.rgb = WHITE
for i, row in enumerate(data, start=1):
    for j, v in enumerate(row):
        c = tbl.cell(i, j); c.text = v
        pr = c.text_frame.paragraphs[0]; pr.runs[0].font.size = Pt(13)
        if i >= 2:
            pr.runs[0].font.bold = True; pr.runs[0].font.color.rgb = GREEN if j > 0 else NAVY
        else:
            pr.runs[0].font.color.rgb = RGBColor(0x2D, 0x37, 0x48)
bb = s.shapes.add_textbox(Inches(0.55), Inches(4.7), Inches(12.2), Inches(1.5))
btf = bb.text_frame; btf.word_wrap = True
for k, it in enumerate([
    ("Primary detection: ", "the multi-front threat-profile detector caught all four by named technique (C2-beacon, DGA, LOTL-process, cohort-rare access, recon-fanout, insider-collection) at ZERO false positives."),
    ("Composite: ", "embedding composite also caught all four — including BOTH Volt and Salt Typhoon — at 10.6% FP, though it cleanly separates only two of the four (USR-118, USR-156); AUC ≈ 0.98."),
    ("The contrast: ", "traditional detectors caught at most one and ranked the stealthy campaigns among ordinary users."),
]):
    p = btf.paragraphs[0] if k == 0 else btf.add_paragraph(); p.space_after = Pt(8)
    rb = p.add_run(); rb.text = "▪  "; rb.font.size = Pt(16); rb.font.color.rgb = BLUE; rb.font.bold = True
    r1 = p.add_run(); r1.text = it[0]; r1.font.size = Pt(16); r1.font.bold = True; r1.font.color.rgb = NAVY
    r2 = p.add_run(); r2.text = it[1]; r2.font.size = Pt(16); r2.font.color.rgb = RGBColor(0x2D, 0x37, 0x48)
cap = s.shapes.add_textbox(Inches(0.55), Inches(6.75), Inches(12.2), Inches(0.5))
cp = cap.text_frame.paragraphs[0]; cr = cp.add_run()
cr.text = "All figures are from synthetic data and must be re-validated on operational telemetry — a scoped pilot is the recommended next step."
cr.font.size = Pt(11); cr.font.italic = True; cr.font.color.rgb = GRAY

# ---------- 6. Preemptive + use cases ----------
slide("6", "Preemptive cybersecurity & advanced use cases",
      bullets=[
          ("Preemptive (get left of impact): ", "the signal is a directional drift that builds over weeks/months — an EARLY-warning capability that closes the Volt-Typhoon dwell window."),
          ("Cyber early warning: ", "insider / account-takeover before exfiltration; configuration-drift; continuous monitoring for RMF / ATO."),
          ("One architecture, many missions: ", "supplier & supply-chain risk, DLA demand forecasting, prescriptive maintenance, fraud / improper-payments, counter-intelligence."),
          ("Deployable: ", "model- and cloud-agnostic, air-gap capable, US-developed and patent-pending — complements existing tools, no vendor lock-in."),
      ],
      callout="Same engine, many missions — detect what an entity is becoming, and act before impact.",
      callout_fill=LGREEN, callout_txt=GREEN)

out = os.path.join(HERE, "VIntelligence_Army_Cyber_Brief.pptx")
prs.save(out)
print("Saved:", out, "| slides:", len(prs.slides._sldIdLst))
