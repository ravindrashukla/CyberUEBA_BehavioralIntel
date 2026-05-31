#!/usr/bin/env python3
"""
Build V-Intelligence UEBA Technical Deep Dive PowerPoint Deck.

Generates a 20-slide technical presentation covering synthetic data design,
algorithm comparison, and empirical results from the V-Intelligence UEBA behavioral
analysis pipeline.  Safe for federal government audience — no proprietary
math, no implementation internals, no vendor-branding constraints.

Output: docs/V_Intelligence_UEBA_Technical_Deep_Dive.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Color palette  (identical to build_analysis_deck.py)
# ---------------------------------------------------------------------------
NAVY       = RGBColor(13,  27,  42)
BLUE       = RGBColor(27,  79,  114)
TEAL       = RGBColor(14,  107, 138)
GOLD       = RGBColor(183, 149, 11)
RED        = RGBColor(192, 57,  43)
GREEN      = RGBColor(39,  174, 96)
WHITE      = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(240, 242, 245)
DARK_GRAY  = RGBColor(80,  80,  80)
MID_GRAY   = RGBColor(180, 180, 180)
ORANGE     = RGBColor(230, 126, 34)

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


# ---------------------------------------------------------------------------
# Primitive helpers
# ---------------------------------------------------------------------------

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.fill.solid()
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape


def set_text(shape, text, font_size=18, bold=False, color=NAVY,
             alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_paragraph(text_frame, text, font_size=18, bold=False, color=NAVY,
                  alignment=PP_ALIGN.LEFT, space_before=Pt(6), space_after=Pt(3),
                  font_name="Calibri", level=0):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    p.level = level
    return p


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False,
                color=NAVY, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_title_bar(slide, title_text, font_size=28):
    bar = add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1), fill_color=NAVY)
    set_text(bar, title_text, font_size=font_size, bold=True, color=WHITE,
             alignment=PP_ALIGN.LEFT)
    bar.text_frame.paragraphs[0].space_before = Pt(12)
    bar.text_frame.margin_left = Inches(0.6)
    bar.text_frame.margin_top = Inches(0.18)
    add_shape(slide, Inches(0), Inches(1.1), SLIDE_WIDTH, Inches(0.04), fill_color=GOLD)
    return bar


def add_footer(slide, text="22nd Century Technologies, Inc.  |  V-Intelligence UEBA Technical Analysis  |  Confidential"):
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.35),
                text, font_size=10, bold=False, color=MID_GRAY, alignment=PP_ALIGN.LEFT)


def add_callout_box(slide, left, top, width, height, text, font_size=20,
                    bg=RGBColor(255, 251, 235), border=None):
    if border is None:
        border = GOLD
    box = add_rounded_rect(slide, left, top, width, height, fill_color=bg)
    box.line.fill.solid()
    box.line.color.rgb = border
    box.line.width = Pt(2.5)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = border
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    return box


def build_table(slide, rows, cols, left, top, width, height, data,
                header_row=True, col_widths=None, cell_colors=None,
                body_font_size=13, header_font_size=14):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c] if r < len(data) and c < len(data[r]) else ""

            for para in cell.text_frame.paragraphs:
                para.font.name = "Calibri"
                if r == 0 and header_row:
                    para.font.size = Pt(header_font_size)
                    para.font.bold = True
                    para.font.color.rgb = WHITE
                    para.alignment = PP_ALIGN.CENTER
                else:
                    para.font.size = Pt(body_font_size)
                    para.font.bold = False
                    para.font.color.rgb = NAVY
                    para.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT

            cell.text_frame.margin_left  = Inches(0.07)
            cell.text_frame.margin_right = Inches(0.07)
            cell.text_frame.margin_top   = Inches(0.04)
            cell.text_frame.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            fill = cell.fill
            fill.solid()
            if r == 0 and header_row:
                fill.fore_color.rgb = NAVY
            elif cell_colors and (r, c) in cell_colors:
                fill.fore_color.rgb = cell_colors[(r, c)]
                for para in cell.text_frame.paragraphs:
                    para.font.color.rgb = WHITE
                    para.font.bold = True
            elif r % 2 == 0:
                fill.fore_color.rgb = LIGHT_GRAY
            else:
                fill.fore_color.rgb = WHITE

    return table


# ===========================================================================
# SLIDE BUILDERS
# ===========================================================================

# ---------------------------------------------------------------------------
# Slide 1: Title
# ---------------------------------------------------------------------------
def slide_01_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    # Gold accent rule above title
    add_shape(slide, Inches(0.8), Inches(1.6), Inches(2.5), Inches(0.05), fill_color=GOLD)

    add_textbox(slide, Inches(0.8), Inches(1.8), Inches(11.8), Inches(0.9),
                "V-Intelligence UEBA Technical Deep Dive",
                font_size=42, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)

    add_textbox(slide, Inches(0.8), Inches(2.65), Inches(11.8), Inches(0.7),
                "Synthetic Data, Detection Algorithms & Empirical Results",
                font_size=26, bold=False, color=TEAL, alignment=PP_ALIGN.LEFT)

    add_shape(slide, Inches(0.8), Inches(3.5), Inches(5.5), Inches(0.04), fill_color=GOLD)

    add_textbox(slide, Inches(0.8), Inches(3.75), Inches(11.8), Inches(0.55),
                "130-Day Controlled Analysis  |  50 Users  |  17 Detection Methods  |  4 Attack Campaigns",
                font_size=18, bold=False, color=GOLD, alignment=PP_ALIGN.LEFT)

    add_textbox(slide, Inches(0.8), Inches(4.6), Inches(11.8), Inches(0.5),
                "22nd Century Technologies, Inc.  |  V-Intelligence UEBA Program  |  May 2025",
                font_size=15, bold=False, color=MID_GRAY, alignment=PP_ALIGN.LEFT)

    add_textbox(slide, Inches(0.8), Inches(6.6), Inches(11.8), Inches(0.4),
                "DISTRIBUTION STATEMENT: For Official Use Only",
                font_size=11, bold=False, color=MID_GRAY, alignment=PP_ALIGN.LEFT)


# ---------------------------------------------------------------------------
# Slide 2: Executive Summary
# ---------------------------------------------------------------------------
def slide_02_exec_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Executive Summary")

    add_textbox(slide, Inches(0.7), Inches(1.3), Inches(11.9), Inches(0.5),
                "Punchline: No single method detects all 4 attack campaigns. The optimal 2-method "
                "ensemble (LOF + Zone Divergence) achieves 100% detection at 6.5% FP.",
                font_size=16, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)

    findings = [
        (TEAL,  "Finding 1 — Traditional Methods Miss Insider Threats",
                "LOF detects 3/4 attacks at 0% FP — best single traditional method. But no "
                "traditional algorithm detects USR-156 (8-month insider threat) at viable FP rates."),
        (ORANGE, "Finding 2 — Temporal Methods Over-Alert",
                "Temporal Z-Score detects all 4 attacks but at 100% FP — every user flagged. "
                "Feature Trajectory achieves 6.5% FP but only detects 2 of 4 attacks."),
        (GREEN,  "Finding 3 — Zone-Based Behavioral Analysis Fills the Gap",
                "Tier 3 Zone Divergence detects the insider (USR-156) and slow APT (USR-234) "
                "by analyzing which behavioral dimensions drift. 6.5% FP."),
        (NAVY,   "Finding 4 — LOF + Zone Divergence = Optimal Ensemble",
                "2-method ensemble detects ALL 4 campaigns at 6.5% FP — LOF catches network "
                "footprint attacks, Zone Divergence catches behavioral-direction attacks."),
    ]

    y = Inches(2.05)
    for color, headline, detail in findings:
        add_shape(slide, Inches(0.7), y, Inches(0.1), Inches(0.85), fill_color=color)
        add_textbox(slide, Inches(1.0), y, Inches(11.6), Inches(0.4),
                    headline, font_size=16, bold=True, color=NAVY, alignment=PP_ALIGN.LEFT)
        add_textbox(slide, Inches(1.0), y + Inches(0.38), Inches(11.6), Inches(0.5),
                    detail, font_size=13, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)
        y += Inches(1.05)

    add_callout_box(slide, Inches(1.8), Inches(6.3), Inches(9.7), Inches(0.65),
                    "LOF + Zone Divergence: all 4 campaigns detected at 6.5% FP — optimal ensemble",
                    font_size=18)
    add_footer(slide)


# ---------------------------------------------------------------------------
# Slide 3: Synthetic Data Overview
# ---------------------------------------------------------------------------
def slide_03_synthetic_data_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Synthetic Data Overview")

    stats = [
        ("130", "Days of Telemetry", "Jan 1 – May 10, 2025"),
        ("50",  "Simulated Users",   "Employees, contractors, privileged"),
        ("7",   "Log Source Types",  "Auth, File, Network, Email, Web, Privilege, Endpoint"),
        ("23",  "Behavioral Features","Extracted per user per weekly window"),
        ("950", "Feature Vectors",   "50 users × 19 weekly observation windows"),
        ("0.05%", "Attack Event Rate", "1,877 injected events in ~3.5M total"),
    ]

    positions = [
        (Inches(0.5),  Inches(1.5)),
        (Inches(4.6),  Inches(1.5)),
        (Inches(8.7),  Inches(1.5)),
        (Inches(0.5),  Inches(4.05)),
        (Inches(4.6),  Inches(4.05)),
        (Inches(8.7),  Inches(4.05)),
    ]

    for (stat, label, sublabel), (x, y) in zip(stats, positions):
        box = add_rounded_rect(slide, x, y, Inches(3.9), Inches(2.2))
        box.line.fill.solid()
        box.line.color.rgb = TEAL
        box.line.width = Pt(1.5)
        add_shape(slide, x, y, Inches(3.9), Inches(0.06), fill_color=TEAL)

        add_textbox(slide, x + Inches(0.1), y + Inches(0.15), Inches(3.7), Inches(0.75),
                    stat, font_size=36, bold=True, color=TEAL, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.1), y + Inches(0.9), Inches(3.7), Inches(0.45),
                    label, font_size=16, bold=True, color=NAVY, alignment=PP_ALIGN.CENTER)
        add_shape(slide, x + Inches(1.1), y + Inches(1.35), Inches(1.7), Inches(0.03),
                  fill_color=GOLD)
        add_textbox(slide, x + Inches(0.1), y + Inches(1.5), Inches(3.7), Inches(0.55),
                    sublabel, font_size=12, bold=False, color=DARK_GRAY,
                    alignment=PP_ALIGN.CENTER)

    # Vertical separators
    add_shape(slide, Inches(4.45), Inches(1.5), Inches(0.02), Inches(4.9), fill_color=LIGHT_GRAY)
    add_shape(slide, Inches(8.55), Inches(1.5), Inches(0.02), Inches(4.9), fill_color=LIGHT_GRAY)

    add_footer(slide)


# ---------------------------------------------------------------------------
# Slide 4: Log Types & Feature Engineering
# ---------------------------------------------------------------------------
def slide_04_log_types_feature_engineering(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Log Types & Feature Engineering")

    add_textbox(slide, Inches(0.7), Inches(1.25), Inches(11.9), Inches(0.45),
                "Seven distinct telemetry sources are ingested and normalized. "
                "23 behavioral features are computed per user per weekly window.",
                font_size=16, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)

    log_types = [
        ("Authentication",       "5 features",  TEAL,  "Login patterns, failure rates,\noff-hours activity"),
        ("File Access",          "4 features",  BLUE,  "File counts, sensitivity\nscore, restricted access"),
        ("Network",              "5 features",  NAVY,  "Data volumes, DNS diversity,\nexternal connections"),
        ("Email",                "4 features",  GOLD,  "Sent volume, external ratio,\nattachment behavior"),
        ("Web Proxy",            "3 features",  TEAL,  "Request volume, block rate,\ndomain diversity"),
        ("Privilege Operations", "1 feature",   BLUE,  "Escalation events,\nadmin action count"),
        ("Endpoint",             "1 feature",   NAVY,  "Admin actions on\nendpoint systems"),
    ]

    positions = [
        (Inches(0.35),  Inches(2.0)),
        (Inches(2.2),   Inches(2.0)),
        (Inches(4.05),  Inches(2.0)),
        (Inches(5.9),   Inches(2.0)),
        (Inches(7.75),  Inches(2.0)),
        (Inches(9.6),   Inches(2.0)),
        (Inches(11.45), Inches(2.0)),
    ]

    for (name, feat_count, color, desc), (x, y) in zip(log_types, positions):
        box_w = Inches(1.75)
        box = add_rounded_rect(slide, x, y, box_w, Inches(4.6))
        box.line.fill.solid()
        box.line.color.rgb = color
        box.line.width = Pt(1.5)
        add_shape(slide, x, y, box_w, Inches(0.06), fill_color=color)

        add_textbox(slide, x + Inches(0.05), y + Inches(0.15), box_w - Inches(0.1),
                    Inches(0.7), name, font_size=13, bold=True, color=NAVY,
                    alignment=PP_ALIGN.CENTER)
        badge = add_rounded_rect(slide, x + Inches(0.3), y + Inches(0.9), Inches(1.15),
                                 Inches(0.42), fill_color=color)
        set_text(badge, feat_count, font_size=12, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)

        add_shape(slide, x + Inches(0.2), y + Inches(1.45), box_w - Inches(0.4),
                  Inches(0.02), fill_color=LIGHT_GRAY)
        add_textbox(slide, x + Inches(0.05), y + Inches(1.6), box_w - Inches(0.1),
                    Inches(2.8), desc, font_size=11, bold=False, color=DARK_GRAY,
                    alignment=PP_ALIGN.CENTER)

    # Arrow showing flow
    add_shape(slide, Inches(0.4), Inches(6.85), Inches(12.5), Inches(0.04), fill_color=GOLD)
    add_textbox(slide, Inches(0.4), Inches(6.9), Inches(6.0), Inches(0.35),
                "7 log types  →  23 behavioral features  →  950 weekly vectors",
                font_size=14, bold=True, color=NAVY, alignment=PP_ALIGN.LEFT)

    add_footer(slide)


# ---------------------------------------------------------------------------
# Slide 5: Feature Detail Table
# ---------------------------------------------------------------------------
def slide_05_feature_detail_table(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "23 Behavioral Features — Complete Reference")

    data = [
        ["Feature Name",            "Source",              "What It Captures"],
        ["auth_total",              "Authentication",      "Total authentication events per week"],
        ["auth_failures",           "Authentication",      "Count of failed authentication attempts"],
        ["auth_failure_rate",       "Authentication",      "Ratio of failures to total auth events"],
        ["auth_off_hours",          "Authentication",      "Auth events outside business hours"],
        ["unique_auth_sources",     "Authentication",      "Distinct source IPs/devices used for auth"],
        ["file_total",              "File Access",         "Total file access events"],
        ["files_restricted",        "File Access",         "Accesses to restricted/sensitive files"],
        ["file_sensitivity_score",  "File Access",         "Weighted score of file classification levels accessed"],
        ["unique_file_paths",       "File Access",         "Number of distinct directory paths visited"],
        ["bytes_out",               "Network",             "Total outbound data volume (bytes)"],
        ["bytes_in",                "Network",             "Total inbound data volume (bytes)"],
        ["dns_unique_domains",      "Network",             "Distinct DNS domains resolved"],
        ["connections_external",    "Network",             "External (non-internal) network connections"],
        ["avg_session_duration",    "Network",             "Mean session length in seconds"],
        ["emails_sent",             "Email",               "Total outbound emails sent"],
        ["emails_external",         "Email",               "Emails sent to external recipients"],
        ["attachment_count",        "Email",               "Emails containing attachments"],
        ["emails_with_large_attachments", "Email",         "Emails with attachments above size threshold"],
        ["web_requests",            "Web Proxy",           "Total web proxy requests"],
        ["web_blocked",             "Web Proxy",           "Requests blocked by proxy policy"],
        ["unique_domains_visited",  "Web Proxy",           "Distinct domains accessed via proxy"],
        ["privilege_escalations",   "Privilege Operations","Count of privilege escalation events"],
        ["admin_actions",           "Endpoint/Privilege",  "Administrative actions performed"],
    ]

    source_colors = {
        "Authentication": TEAL,
        "File Access":    BLUE,
        "Network":        NAVY,
        "Email":          GOLD,
        "Web Proxy":      TEAL,
        "Privilege Operations": BLUE,
        "Endpoint/Privilege": NAVY,
    }
    # Alternating source blocks look fine with standard alternating rows; no extra color needed

    build_table(slide, rows=24, cols=3,
                left=Inches(0.4), top=Inches(1.25),
                width=Inches(12.5), height=Inches(5.55),
                data=data,
                col_widths=[Inches(2.9), Inches(2.7), Inches(6.9)],
                body_font_size=11, header_font_size=13)

    add_footer(slide)


# ---------------------------------------------------------------------------
# Slide 6: Attack Scenarios
# ---------------------------------------------------------------------------
def slide_06_attack_scenarios(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Injected Attack Scenarios")

    add_textbox(slide, Inches(0.7), Inches(1.25), Inches(11.9), Inches(0.4),
                "Four attack campaigns injected into 50-user population. "
                "Attack events = 0.05% of total telemetry volume.",
                font_size=15, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)

    PURPLE = RGBColor(142, 68, 173)
    attacks = [
        ("USR-156 — Insider Threat", "8-Month Escalation  |  Mar 2025",
         [("Phase 1", "Subtle mood shift, off-hours increase"),
          ("Phase 2", "Curiosity browsing outside scope"),
          ("Phase 3", "Restricted file access"),
          ("Phase 4", "Trickle exfiltration")],
         "Volumes stay in normal ranges", RED),
        ("USR-234 — Slow APT", "180-Day Campaign  |  Apr 2025",
         [("C2 Pattern", "~4 beacons/day at 6hr intervals"),
          ("DNS", "DGA queries blend with normal"),
          ("Staging", "Progressive data accumulation"),
          ("Exfil", "Encrypted channel exfiltration")],
         "Subtle network footprint", GOLD),
        ("USR-042 — Volt Typhoon", "115-Day LOTL  |  Jan 2025",
         [("Tools", "RDP, WMI, PsExec (legitimate)"),
          ("Creds", "Credential harvesting via admin"),
          ("Lateral", "Movement mimics admin activity"),
          ("Infra", "Pre-positioning for future ops")],
         "Blends with normal admin behavior", PURPLE),
        ("USR-118 — Salt Typhoon", "100-Day Telecom  |  Jan 2025",
         [("C2", "DNS tunneling for covert channel"),
          ("Access", "CDR/lawful intercept data"),
          ("Exfil", "Staged exfiltration over weeks"),
          ("Cover", "Mimics telecom operations")],
         "Mimics legitimate telecom ops", BLUE),
    ]

    positions = [
        (Inches(0.3), Inches(1.75)),
        (Inches(6.7), Inches(1.75)),
        (Inches(0.3), Inches(4.55)),
        (Inches(6.7), Inches(4.55)),
    ]

    for (title, subtitle, phases, insight, color), (x, y) in zip(attacks, positions):
        card = add_rounded_rect(slide, x, y, Inches(6.2), Inches(2.55))
        card.line.fill.solid()
        card.line.color.rgb = color
        card.line.width = Pt(2)
        add_shape(slide, x, y, Inches(6.2), Inches(0.06), fill_color=color)

        add_textbox(slide, x + Inches(0.15), y + Inches(0.08), Inches(5.9), Inches(0.38),
                    title, font_size=15, bold=True, color=color, alignment=PP_ALIGN.LEFT)
        add_textbox(slide, x + Inches(0.15), y + Inches(0.42), Inches(5.9), Inches(0.3),
                    subtitle, font_size=11, bold=True, color=NAVY, alignment=PP_ALIGN.LEFT)
        add_shape(slide, x + Inches(0.3), y + Inches(0.72), Inches(5.6), Inches(0.02),
                  fill_color=LIGHT_GRAY)

        py = y + Inches(0.78)
        for label, desc in phases:
            add_textbox(slide, x + Inches(0.2), py, Inches(1.2), Inches(0.28),
                        label, font_size=10, bold=True, color=color, alignment=PP_ALIGN.LEFT)
            add_textbox(slide, x + Inches(1.4), py, Inches(4.6), Inches(0.28),
                        desc, font_size=10, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)
            py += Inches(0.28)

        ins_box = add_rounded_rect(slide, x + Inches(0.2), py + Inches(0.05), Inches(5.8),
                                   Inches(0.35), fill_color=RGBColor(245, 245, 250))
        ins_box.line.fill.solid()
        ins_box.line.color.rgb = color
        ins_box.line.width = Pt(1)
        set_text(ins_box, insight, font_size=10, bold=True, color=color,
                 alignment=PP_ALIGN.CENTER)
    add_footer(slide)


# ---------------------------------------------------------------------------
# Slide 7: Attack Design Philosophy
# ---------------------------------------------------------------------------
def slide_07_attack_design_philosophy(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_textbox(slide, Inches(0.7), Inches(0.35), Inches(12.0), Inches(0.75),
                "Attack Design Philosophy — Why These Threats Are Hard to Detect",
                font_size=30, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)
    add_shape(slide, Inches(0.7), Inches(1.05), Inches(3.5), Inches(0.04), fill_color=GOLD)

    add_textbox(slide, Inches(0.7), Inches(1.25), Inches(11.9), Inches(0.5),
                "All four attacks are specifically engineered to stay within the statistical "
                "profile of normal users.",
                font_size=18, bold=False, color=TEAL, alignment=PP_ALIGN.LEFT)

    contrasts = [
        ("Volume-Based Signal",   "Activity counts, bytes, event frequencies",
         "Magnitude Anomaly",     "USR-042/USR-118 create network footprints; USR-234 C2 beacons raise DNS counts"),
        ("Direction-Based Signal","Which resources accessed, what changed over time",
         "Behavioral Anomaly",    "USR-156 access patterns shift without volume spike; USR-042 LOTL blends with admin"),
    ]

    y = Inches(2.0)
    for sig_type, sig_desc, anom_type, anom_desc in contrasts:
        left_box = add_rounded_rect(slide, Inches(0.5), y, Inches(6.0), Inches(1.55),
                                    fill_color=RGBColor(20, 40, 60))
        left_box.line.fill.solid()
        left_box.line.color.rgb = TEAL
        left_box.line.width = Pt(1.5)
        add_textbox(slide, Inches(0.7), y + Inches(0.1), Inches(5.6), Inches(0.4),
                    sig_type, font_size=16, bold=True, color=TEAL, alignment=PP_ALIGN.LEFT)
        add_textbox(slide, Inches(0.7), y + Inches(0.5), Inches(5.6), Inches(0.9),
                    sig_desc, font_size=14, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT)

        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                       Inches(6.6), y + Inches(0.6),
                                       Inches(0.55), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GOLD
        arrow.line.fill.background()

        right_box = add_rounded_rect(slide, Inches(7.3), y, Inches(5.6), Inches(1.55),
                                     fill_color=RGBColor(20, 40, 60))
        right_box.line.fill.solid()
        right_box.line.color.rgb = GOLD
        right_box.line.width = Pt(1.5)
        add_textbox(slide, Inches(7.5), y + Inches(0.1), Inches(5.2), Inches(0.4),
                    anom_type, font_size=16, bold=True, color=GOLD, alignment=PP_ALIGN.LEFT)
        add_textbox(slide, Inches(7.5), y + Inches(0.5), Inches(5.2), Inches(0.9),
                    anom_desc, font_size=14, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT)
        y += Inches(1.85)

    # Why this matters box
    why = add_rounded_rect(slide, Inches(0.5), Inches(5.95), Inches(12.3), Inches(1.1),
                           fill_color=RGBColor(10, 25, 38))
    why.line.fill.solid()
    why.line.color.rgb = GOLD
    why.line.width = Pt(2)
    set_text(why, "Four attack types, two signal categories — no single method covers both.",
             font_size=22, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)
    why.text_frame.margin_top = Inches(0.05)
    add_paragraph(why.text_frame,
                  "LOF + Zone Divergence: each covers what the other misses.",
                  font_size=16, bold=False, color=WHITE, alignment=PP_ALIGN.CENTER,
                  space_before=Pt(8))

    add_textbox(slide, Inches(0.5), Inches(7.1), Inches(12), Inches(0.3),
                "22nd Century Technologies, Inc.  |  V-Intelligence UEBA Technical Analysis  |  Confidential",
                font_size=10, bold=False, color=RGBColor(60, 70, 85), alignment=PP_ALIGN.LEFT)


# ---------------------------------------------------------------------------
# Slide 8: Traditional Methods (IForest, SVM, LOF, Z-Score)
# ---------------------------------------------------------------------------
def slide_08_traditional_methods(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Algorithms 1–4: Traditional Static Detection Methods")

    data = [
        ["Algorithm",          "Core Approach",
         "USR-156\nInsider", "USR-234\nAPT", "USR-042\nVolt", "USR-118\nSalt",
         "FP Rate"],
        ["Isolation Forest",   "Randomly partition;\nisolate outliers",
         "MISS", "MISS", "DET", "DET", "2.2%"],
        ["One-Class SVM",      "Learn boundary around\nnormal data",
         "DET", "MISS", "MISS", "DET", "19.6%"],
        ["Local Outlier Factor","Compare local density\nto neighbors",
         "MISS", "DET", "DET", "DET", "0.0%"],
        ["Z-Score (|z|>3)",    "Flag features >3σ\nfrom mean",
         "MISS", "DET", "DET", "DET", "2.2%"],
    ]

    cell_colors = {}
    for r in [1, 2, 3, 4]:
        for c in [2, 3, 4, 5]:
            if data[r][c] == "MISS":
                cell_colors[(r, c)] = RED
            else:
                cell_colors[(r, c)] = GREEN
    cell_colors[(2, 6)] = RED  # SVM FP rate is high

    build_table(slide, rows=5, cols=7,
                left=Inches(0.3), top=Inches(1.3),
                width=Inches(12.7), height=Inches(3.0),
                data=data, cell_colors=cell_colors,
                col_widths=[Inches(2.0), Inches(2.5), Inches(1.4), Inches(1.4),
                            Inches(1.4), Inches(1.4), Inches(1.4)],
                body_font_size=11, header_font_size=12)

    # Insight panel
    insight = add_rounded_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.2),
                               fill_color=NAVY)
    insight.line.fill.background()
    set_text(insight, "Key Insight: LOF is the Best Single Traditional Method",
             font_size=20, bold=True, color=GOLD, alignment=PP_ALIGN.LEFT)
    insight.text_frame.margin_left = Inches(0.35)
    insight.text_frame.margin_top = Inches(0.15)

    pts = [
        "LOF detects 3/4 attacks at 0% FP by comparing local density — USR-234 C2 beacons, "
        "USR-042 lateral movement, and USR-118 DNS tunneling all create detectable density changes.",
        "But USR-156's feature values remain within normal ranges throughout the entire 8-month "
        "insider campaign — no traditional method reliably catches this.",
        "Only behavioral analysis (Tier 3 Zone Divergence) can detect USR-156 by identifying "
        "which behavioral dimensions are drifting, not just overall magnitude.",
    ]
    for pt in pts:
        add_paragraph(insight.text_frame, "•  " + pt, font_size=13, bold=False,
                      color=WHITE, space_before=Pt(8))

    add_footer(slide)


# ---------------------------------------------------------------------------
# Slide 9: Temporal Methods
# ---------------------------------------------------------------------------
def slide_09_temporal_methods(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Algorithms 5–6: Temporal Change Detection Methods")

    add_textbox(slide, Inches(0.7), Inches(1.25), Inches(11.9), Inches(0.42),
                "Temporal methods add a time dimension: they compare each user's recent behavior "
                "against their own historical baseline rather than against the population.",
                font_size=15, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)

    # Temporal Z-Score box
    lx = Inches(0.5)
    lbox = add_rounded_rect(slide, lx, Inches(1.9), Inches(6.0), Inches(4.0))
    lbox.line.fill.solid()
    lbox.line.color.rgb = TEAL
    lbox.line.width = Pt(2)
    add_shape(slide, lx, Inches(1.9), Inches(6.0), Inches(0.07), fill_color=TEAL)

    add_textbox(slide, lx + Inches(0.2), Inches(2.05), Inches(5.6), Inches(0.45),
                "Algorithm 5 — Temporal Z-Score", font_size=19, bold=True, color=TEAL,
                alignment=PP_ALIGN.LEFT)
    add_textbox(slide, lx + Inches(0.2), Inches(2.6), Inches(5.6), Inches(0.6),
                "Train on first half of the time window; test on the second half. "
                "Flag features that deviate more than 2σ from the user's own training-period baseline.",
                font_size=13, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)

    add_textbox(slide, lx + Inches(0.2), Inches(3.3), Inches(2.8), Inches(0.38),
                "All 4 attacks:", font_size=15, bold=True, color=NAVY, alignment=PP_ALIGN.LEFT)
    r_box = add_shape(slide, lx + Inches(3.1), Inches(3.3),
                      Inches(2.6), Inches(0.38), fill_color=GREEN)
    set_text(r_box, "DETECTED", font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    fp_box = add_rounded_rect(slide, lx + Inches(0.3), Inches(4.3), Inches(5.3), Inches(0.8),
                              fill_color=RGBColor(253, 240, 240))
    fp_box.line.fill.solid()
    fp_box.line.color.rgb = RED
    fp_box.line.width = Pt(1.5)
    set_text(fp_box, "False Positive Rate: 100.0%\nAll 46 normal users flagged",
             font_size=13, bold=True, color=RED, alignment=PP_ALIGN.CENTER)
    fp_box.text_frame.margin_top = Inches(0.08)

    # Feature Trajectory box
    rx = Inches(7.1)
    rbox = add_rounded_rect(slide, rx, Inches(1.9), Inches(6.0), Inches(4.0))
    rbox.line.fill.solid()
    rbox.line.color.rgb = BLUE
    rbox.line.width = Pt(2)
    add_shape(slide, rx, Inches(1.9), Inches(6.0), Inches(0.07), fill_color=BLUE)

    add_textbox(slide, rx + Inches(0.2), Inches(2.05), Inches(5.6), Inches(0.45),
                "Algorithm 6 — Feature-Level Change Detection", font_size=19, bold=True,
                color=BLUE, alignment=PP_ALIGN.LEFT)
    add_textbox(slide, rx + Inches(0.2), Inches(2.6), Inches(5.6), Inches(0.6),
                "Tracks each feature's trajectory over time and flags sustained "
                "directional shifts. Detects persistent changes in individual feature "
                "values across weeks.",
                font_size=13, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)

    add_textbox(slide, rx + Inches(0.2), Inches(3.3), Inches(2.8), Inches(0.38),
                "USR-042, USR-118:", font_size=15, bold=True, color=NAVY, alignment=PP_ALIGN.LEFT)
    r_box2 = add_shape(slide, rx + Inches(3.1), Inches(3.3),
                       Inches(2.6), Inches(0.38), fill_color=GREEN)
    set_text(r_box2, "DETECTED", font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, rx + Inches(0.2), Inches(3.7), Inches(2.8), Inches(0.38),
                "USR-156, USR-234:", font_size=15, bold=True, color=NAVY, alignment=PP_ALIGN.LEFT)
    r_box3 = add_shape(slide, rx + Inches(3.1), Inches(3.7),
                       Inches(2.6), Inches(0.38), fill_color=RED)
    set_text(r_box3, "MISSED", font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    fp_box2 = add_rounded_rect(slide, rx + Inches(0.3), Inches(4.3), Inches(5.3), Inches(0.8),
                               fill_color=RGBColor(255, 252, 240))
    fp_box2.line.fill.solid()
    fp_box2.line.color.rgb = GOLD
    fp_box2.line.width = Pt(1.5)
    set_text(fp_box2, "False Positive Rate: 6.5%\n3 of 46 normal users flagged (but 2/4 missed)",
             font_size=13, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)
    fp_box2.text_frame.margin_top = Inches(0.08)

    add_callout_box(slide, Inches(1.8), Inches(6.3), Inches(9.7), Inches(0.65),
                    "Detection without precision is operational noise, not actionable intelligence",
                    font_size=17)
    add_footer(slide)


# ---------------------------------------------------------------------------
# Slide 10: V-Intelligence UEBA Behavioral Analysis
# ---------------------------------------------------------------------------
def slide_10_acecard_behavioral(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "V-Intelligence UEBA Behavioral Drift Direction Analysis")

    add_textbox(slide, Inches(0.7), Inches(1.25), Inches(11.9), Inches(0.5),
                "V-Intelligence UEBA analyzes the trajectory of behavioral change — "
                "not whether a user is an outlier today, but whether their behavior "
                "is consistently moving in a threat-consistent direction over time.",
                font_size=15, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)

    # Five-step pipeline (high level, no implementation details)
    steps = [
        ("1", "Behavioral\nSerialization",
         "Telemetry is transformed\ninto structured behavioral\ndescriptions for each\nuser per time window"),
        ("2", "Semantic\nRepresentation",
         "Behavioral descriptions\nare encoded into a\nsemantic space preserving\nmeaning, not just magnitude"),
        ("3", "Trajectory\nTracking",
         "Each user's behavioral\nvector is tracked week\nover week to reveal\ndirectional movement"),
        ("4", "Change-Point\nIdentification",
         "Sustained cumulative\nbehavioral drift is\nanalyzed to identify\nstatistically significant shifts"),
        ("5", "Threat Pattern\nAlignment",
         "Drift direction is\ncompared against known\nthreat behavioral profiles\nto classify the change"),
    ]

    start_x = Inches(0.3)
    box_w = Inches(2.35)
    gap = Inches(0.25)

    for i, (num, title, desc) in enumerate(steps):
        x = start_x + i * (box_w + gap)
        y_base = Inches(2.05)

        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                        x + Inches(0.87), y_base,
                                        Inches(0.55), Inches(0.55))
        circle.fill.solid()
        circle.fill.fore_color.rgb = TEAL
        circle.line.fill.background()
        set_text(circle, num, font_size=20, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

        title_box = add_shape(slide, x, y_base + Inches(0.65), box_w, Inches(0.75),
                              fill_color=NAVY)
        set_text(title_box, title, font_size=13, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
        title_box.text_frame.margin_top = Inches(0.08)

        desc_box = add_rounded_rect(slide, x, y_base + Inches(1.45), box_w, Inches(2.05))
        desc_box.line.fill.solid()
        desc_box.line.color.rgb = TEAL
        desc_box.line.width = Pt(1)
        set_text(desc_box, desc, font_size=12, bold=False, color=DARK_GRAY,
                 alignment=PP_ALIGN.CENTER)
        desc_box.text_frame.margin_left = Inches(0.1)
        desc_box.text_frame.margin_right = Inches(0.1)
        desc_box.text_frame.margin_top = Inches(0.15)

        if i < len(steps) - 1:
            arrow_x = x + box_w
            arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                         arrow_x, y_base + Inches(0.85),
                                         gap, Inches(0.25))
            arr.fill.solid()
            arr.fill.fore_color.rgb = GOLD
            arr.line.fill.background()

    summary = add_rounded_rect(slide, Inches(1.0), Inches(5.75), Inches(11.3), Inches(1.0),
                               fill_color=NAVY)
    set_text(summary, "Catches the pattern shift, not the volume spike",
             font_size=22, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)
    summary.text_frame.margin_top = Inches(0.08)
    add_paragraph(summary.text_frame,
                  "Traditional methods ask \"how much?\"  —  V-Intelligence UEBA asks \"in what direction?\"",
                  font_size=15, bold=False, color=WHITE, alignment=PP_ALIGN.CENTER,
                  space_before=Pt(6))

    add_footer(slide)


# ---------------------------------------------------------------------------
# Slide 11: Full Results Comparison Table (KEY SLIDE)
# ---------------------------------------------------------------------------
def slide_11_full_results_table(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Full Results: Three-Tier Detection Comparison", font_size=26)

    # KEY SLIDE badge
    key_badge = add_rounded_rect(slide, Inches(10.8), Inches(0.2), Inches(2.0), Inches(0.5),
                                 fill_color=GOLD)
    set_text(key_badge, "KEY SLIDE", font_size=14, bold=True, color=WHITE,
             alignment=PP_ALIGN.CENTER)

    data = [
        ["Method",            "Tier",
         "USR-156\nInsider", "USR-234\nAPT", "USR-042\nVolt", "USR-118\nSalt",
         "TP", "FP Rate"],
        ["Isolation Forest",   "T1", "MISS", "MISS", "DET",  "DET",  "2", "2.2%"],
        ["One-Class SVM",      "T1", "DET",  "MISS", "MISS", "DET",  "2", "19.6%"],
        ["LOF",                "T1", "MISS", "DET",  "DET",  "DET",  "3", "0.0%"],
        ["Z-Score",            "T1", "MISS", "DET",  "DET",  "DET",  "3", "2.2%"],
        ["Temporal Z-Score",   "T1", "DET",  "DET",  "DET",  "DET",  "4", "100%"],
        ["Feature Trajectory", "T1", "MISS", "MISS", "DET",  "DET",  "2", "6.5%"],
        ["V-Intelligence UEBA Direction",  "T2", "MISS", "MISS", "MISS", "MISS", "0", "4.3%"],
        ["T3 Regime Shift",    "T3", "MISS", "MISS", "DET",  "DET",  "2", "6.5%"],
        ["T3 Zone Divergence", "T3", "DET",  "DET",  "MISS", "MISS", "2", "6.5%"],
        ["T3 Combined",        "T3", "DET",  "DET",  "DET",  "DET",  "4", "8.7%"],
        ["LOF + Zone Div",     "ENS","DET",  "DET",  "DET",  "DET",  "4", "6.5%"],
    ]

    cell_colors = {}
    for r in range(1, len(data)):
        for c in [2, 3, 4, 5]:
            if data[r][c] == "MISS":
                cell_colors[(r, c)] = RED
            elif data[r][c] == "DET":
                cell_colors[(r, c)] = GREEN
        fp = data[r][7]
        if fp in ("100%", "19.6%"):
            cell_colors[(r, 7)] = RED
        elif fp in ("0.0%", "2.2%", "4.3%", "6.5%", "8.7%"):
            cell_colors[(r, 7)] = GREEN

    # Highlight best ensemble row
    cell_colors[(11, 0)] = TEAL
    cell_colors[(11, 1)] = TEAL

    build_table(slide, rows=12, cols=8,
                left=Inches(0.2), top=Inches(1.2),
                width=Inches(12.9), height=Inches(5.0),
                data=data, cell_colors=cell_colors,
                col_widths=[Inches(2.2), Inches(0.7), Inches(1.4), Inches(1.4),
                            Inches(1.4), Inches(1.4), Inches(0.7), Inches(1.2)],
                body_font_size=10, header_font_size=11)

    add_callout_box(slide, Inches(1.0), Inches(6.4), Inches(11.3), Inches(0.6),
                    "LOF + Zone Divergence: ALL 4 attacks detected at 6.5% FP — optimal 2-method ensemble",
                    font_size=17)
    add_footer(slide)


# ---------------------------------------------------------------------------
# Slide 12: Traditional Methods Deep Dive
# ---------------------------------------------------------------------------
def slide_12_traditional_deep_dive(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Traditional Methods — Why They Detect APT but Miss the Insider")

    add_textbox(slide, Inches(0.7), Inches(1.25), Inches(11.9), Inches(0.45),
                "All four traditional algorithms share the same structural limitation: "
                "they measure feature-space distance from normal, not trajectory direction over time.",
                font_size=15, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)

    # Left: Why APT is detected
    lbox = add_rounded_rect(slide, Inches(0.5), Inches(1.95), Inches(6.0), Inches(3.5))
    lbox.line.fill.solid()
    lbox.line.color.rgb = GREEN
    lbox.line.width = Pt(2)
    add_shape(slide, Inches(0.5), Inches(1.95), Inches(6.0), Inches(0.07), fill_color=GREEN)

    add_textbox(slide, Inches(0.7), Inches(2.1), Inches(5.6), Inches(0.45),
                "USR-234 (APT) — Why DETECTED", font_size=18, bold=True, color=GREEN,
                alignment=PP_ALIGN.LEFT)
    apt_points = [
        "C2 beaconing generates ~4 network connections every 6 hours",
        "This raises dns_unique_domains and connections_external above the "
        "normal population range",
        "These elevated counts are clear statistical outliers in the feature space",
        "All four algorithms independently identify this as an anomaly",
        "This is a volume anomaly — exactly what traditional methods are designed for",
    ]
    y = Inches(2.65)
    for pt in apt_points:
        add_shape(slide, Inches(0.75), y + Inches(0.1), Inches(0.1), Inches(0.1),
                  fill_color=GREEN)
        add_textbox(slide, Inches(1.0), y, Inches(5.2), Inches(0.38),
                    pt, font_size=12, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)
        y += Inches(0.42)

    # Right: Why Insider is missed
    rbox = add_rounded_rect(slide, Inches(7.0), Inches(1.95), Inches(6.0), Inches(3.5))
    rbox.line.fill.solid()
    rbox.line.color.rgb = RED
    rbox.line.width = Pt(2)
    add_shape(slide, Inches(7.0), Inches(1.95), Inches(6.0), Inches(0.07), fill_color=RED)

    add_textbox(slide, Inches(7.2), Inches(2.1), Inches(5.6), Inches(0.45),
                "USR-156 (Insider) — Why MISSED", font_size=18, bold=True, color=RED,
                alignment=PP_ALIGN.LEFT)
    insider_points = [
        "Login count: 45/week — normal range is 30–60; not flagged",
        "File accesses: 8/day — normal range 3–15; not flagged",
        "Off-hours access: slightly elevated but within population norms",
        "No single feature value is a statistical outlier at any point",
        "The attack is a direction change, not a magnitude spike",
    ]
    y = Inches(2.65)
    for pt in insider_points:
        add_shape(slide, Inches(7.25), y + Inches(0.1), Inches(0.1), Inches(0.1),
                  fill_color=RED)
        add_textbox(slide, Inches(7.5), y, Inches(5.2), Inches(0.38),
                    pt, font_size=12, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)
        y += Inches(0.42)

    insight = add_rounded_rect(slide, Inches(0.5), Inches(5.65), Inches(12.3), Inches(1.05),
                               fill_color=NAVY)
    set_text(insight, "Core Limitation: Traditional methods evaluate a snapshot in feature space.",
             font_size=18, bold=True, color=GOLD, alignment=PP_ALIGN.LEFT)
    insight.text_frame.margin_left = Inches(0.3)
    insight.text_frame.margin_top = Inches(0.1)
    add_paragraph(insight.text_frame,
                  "They cannot detect that a user is gradually moving toward a more dangerous "
                  "behavioral profile over weeks and months.",
                  font_size=14, bold=False, color=WHITE, space_before=Pt(6))

    add_footer(slide)


# ---------------------------------------------------------------------------
# Slide 13: Temporal Methods Deep Dive
# ---------------------------------------------------------------------------
def slide_13_temporal_deep_dive(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Temporal Methods — Detection Without Precision")

    add_textbox(slide, Inches(0.7), Inches(1.25), Inches(11.9), Inches(0.45),
                "Temporal methods detect all 4 attack campaigns — but flag the majority "
                "of the normal user population as well, rendering the alerts non-actionable.",
                font_size=15, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)

    # Good news / bad news
    good_box = add_rounded_rect(slide, Inches(0.5), Inches(1.95), Inches(5.8), Inches(2.5))
    good_box.line.fill.solid()
    good_box.line.color.rgb = GREEN
    good_box.line.width = Pt(2)
    add_shape(slide, Inches(0.5), Inches(1.95), Inches(5.8), Inches(0.07), fill_color=GREEN)
    add_textbox(slide, Inches(0.7), Inches(2.1), Inches(5.4), Inches(0.45),
                "What Works", font_size=18, bold=True, color=GREEN, alignment=PP_ALIGN.LEFT)
    good_pts = ["Temporal comparison catches all 4 attack campaigns",
                "Detects USR-156's insider shift and USR-234's C2 changes",
                "Also catches USR-042 (Volt Typhoon) and USR-118 (Salt Typhoon)"]
    for idx, pt in enumerate(good_pts):
        add_shape(slide, Inches(0.75), Inches(2.65 + idx * 0.42), Inches(0.1), Inches(0.1), fill_color=GREEN)
        add_textbox(slide, Inches(1.0), Inches(2.65 + idx * 0.42),
                    Inches(5.0), Inches(0.38), pt, font_size=12, bold=False, color=DARK_GRAY,
                    alignment=PP_ALIGN.LEFT)

    bad_box = add_rounded_rect(slide, Inches(7.0), Inches(1.95), Inches(5.8), Inches(2.5))
    bad_box.line.fill.solid()
    bad_box.line.color.rgb = RED
    bad_box.line.width = Pt(2)
    add_shape(slide, Inches(7.0), Inches(1.95), Inches(5.8), Inches(0.07), fill_color=RED)
    add_textbox(slide, Inches(7.2), Inches(2.1), Inches(5.4), Inches(0.45),
                "Why It Fails Operationally", font_size=18, bold=True, color=RED,
                alignment=PP_ALIGN.LEFT)
    bad_pts = ["Normal employees also shift behavior over months",
               "Org-wide changes (policy, project ramp-ups) affect everyone",
               "Cannot distinguish threat drift from benign evolution",
               "46 of 46 normal users flagged — every alert is probably a false alarm"]
    y = Inches(2.65)
    for pt in bad_pts:
        add_shape(slide, Inches(7.25), y + Inches(0.1), Inches(0.1), Inches(0.1), fill_color=RED)
        add_textbox(slide, Inches(7.5), y, Inches(5.0), Inches(0.38),
                    pt, font_size=12, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)
        y += Inches(0.42)

    # FP comparison bar-style visual
    add_textbox(slide, Inches(0.7), Inches(4.65), Inches(11.9), Inches(0.4),
                "False Positive Breakdown:", font_size=16, bold=True, color=NAVY,
                alignment=PP_ALIGN.LEFT)

    bars = [
        ("Temporal Z-Score",              46, 100.0, RED),
        ("Feature Trajectory Top-10%",     3,   6.5, GREEN),
    ]
    y_bar = Inches(5.15)
    max_bar_w = Inches(6.0)
    for method, count, rate, color in bars:
        add_textbox(slide, Inches(0.7), y_bar, Inches(3.5), Inches(0.38),
                    method, font_size=12, bold=False, color=NAVY, alignment=PP_ALIGN.LEFT)
        bar_w = Inches(6.0 * rate / 100.0)
        add_shape(slide, Inches(4.3), y_bar, bar_w, Inches(0.3), fill_color=color)
        add_textbox(slide, Inches(4.3) + bar_w + Inches(0.1), y_bar, Inches(1.5), Inches(0.32),
                    f"{count} users ({rate}%)", font_size=12, bold=True, color=color,
                    alignment=PP_ALIGN.LEFT)
        y_bar += Inches(0.45)

    add_callout_box(slide, Inches(1.5), Inches(6.5), Inches(10.3), Inches(0.55),
                    "Feature Trajectory Top-10% brings FP to 6.5% but detects only 2 of 4 — "
                    "misses the insider and slow APT entirely",
                    font_size=14)
    add_footer(slide)


# ---------------------------------------------------------------------------
# Slide 14: V-Intelligence UEBA Direction Results
# ---------------------------------------------------------------------------
def slide_14_acecard_direction_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Zone Divergence — The Insider Detection Breakthrough")

    add_textbox(slide, Inches(0.7), Inches(1.25), Inches(11.9), Inches(0.45),
                "Tier 3 Zone Divergence detects attacks by analyzing WHICH behavioral zone "
                "drifts while identity remains stable — catching threats invisible to all "
                "traditional methods and single-composite embeddings.",
                font_size=15, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)

    data = [
        ["Method",           "USR-156\n(Insider)", "USR-234\n(Slow APT)",
         "USR-042\n(Volt Typh)", "USR-118\n(Salt Typh)", "FP Rate"],
        ["LOF (Tier 1)",     "MISSED",  "DETECTED", "DETECTED", "DETECTED", "0.0%"],
        ["Zone Divergence",  "DETECTED", "DETECTED", "MISSED",  "MISSED",   "6.5%"],
        ["LOF + Zone Div",   "DETECTED", "DETECTED", "DETECTED", "DETECTED", "6.5%"],
    ]
    cell_colors = {}
    for r in [1, 2, 3]:
        cell_colors[(r, 0)] = TEAL
        for c in [1, 2, 3, 4]:
            val = data[r][c]
            cell_colors[(r, c)] = GREEN if val == "DETECTED" else RED

    build_table(slide, rows=4, cols=6,
                left=Inches(0.5), top=Inches(2.0),
                width=Inches(12.3), height=Inches(1.5),
                data=data, cell_colors=cell_colors,
                col_widths=[Inches(2.2), Inches(2.0), Inches(2.0), Inches(2.0),
                            Inches(2.0), Inches(1.5)],
                body_font_size=13, header_font_size=12)

    # Why Zone Divergence works / Why single composite fails
    lbox = add_rounded_rect(slide, Inches(0.5), Inches(3.8), Inches(5.9), Inches(2.5))
    lbox.line.fill.solid()
    lbox.line.color.rgb = GREEN
    lbox.line.width = Pt(2)
    add_shape(slide, Inches(0.5), Inches(3.8), Inches(5.9), Inches(0.07), fill_color=GREEN)
    add_textbox(slide, Inches(0.7), Inches(3.95), Inches(5.5), Inches(0.45),
                "Why Zone Divergence Catches the Insider", font_size=18, bold=True, color=GREEN,
                alignment=PP_ALIGN.LEFT)
    pts_l = [
        "Decomposes behavior into 5 zones: identity, access, data, network, risk",
        "USR-156: identity drift = 0.00, data_behavior drift = 0.33 — "
        "accessing different data while looking like the same person",
        "USR-234: identity drift = 0.00, network drift = 0.28 — "
        "C2 beaconing creates a zone-specific signature",
    ]
    y = Inches(4.5)
    for pt in pts_l:
        add_shape(slide, Inches(0.75), y + Inches(0.08), Inches(0.1), Inches(0.1),
                  fill_color=GREEN)
        add_textbox(slide, Inches(1.0), y, Inches(5.1), Inches(0.55),
                    pt, font_size=12, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)
        y += Inches(0.55)

    rbox = add_rounded_rect(slide, Inches(7.0), Inches(3.8), Inches(5.9), Inches(2.5))
    rbox.line.fill.solid()
    rbox.line.color.rgb = RGBColor(180, 180, 180)
    rbox.line.width = Pt(2)
    add_shape(slide, Inches(7.0), Inches(3.8), Inches(5.9), Inches(0.07),
              fill_color=RGBColor(180, 180, 180))
    add_textbox(slide, Inches(7.2), Inches(3.95), Inches(5.5), Inches(0.45),
                "Why Single-Composite Embedding Fails", font_size=18, bold=True, color=DARK_GRAY,
                alignment=PP_ALIGN.LEFT)
    pts_r = [
        "Tier 2 averages all behavioral zones into one composite vector",
        "Zone-specific drift signals are diluted — data_behavior drift "
        "is averaged with stable identity, access, and risk zones",
        "Result: V-Intelligence UEBA Direction detects 0 of 4 attacks at threshold",
    ]
    y = Inches(4.5)
    for pt in pts_r:
        add_shape(slide, Inches(7.25), y + Inches(0.08), Inches(0.1), Inches(0.1),
                  fill_color=DARK_GRAY)
        add_textbox(slide, Inches(7.5), y, Inches(5.1), Inches(0.55),
                    pt, font_size=12, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)
        y += Inches(0.55)

    add_callout_box(slide, Inches(1.5), Inches(6.55), Inches(10.3), Inches(0.55),
                    "LOF + Zone Divergence: 4 of 4 attacks detected at 6.5% FP — "
                    "the optimal 2-method ensemble",
                    font_size=18)
    add_footer(slide)


# ---------------------------------------------------------------------------
# Slide 15: Combined Approach
# ---------------------------------------------------------------------------
def slide_15_combined_approach(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Optimal Ensemble — LOF + Zone Divergence")

    add_textbox(slide, Inches(0.7), Inches(1.25), Inches(11.9), Inches(0.45),
                "The four attack types require different detection strategies. "
                "The optimal 2-method ensemble covers both signal categories.",
                font_size=15, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)

    # Coverage boxes
    components = [
        ("Local Outlier Factor (Tier 1)",   BLUE,
         "USR-234, USR-042, USR-118",
         "Density-based detection catches network-footprint anomalies from "
         "C2 beaconing, lateral movement, and DNS tunneling. 0% FP."),
        ("Zone Divergence (Tier 3)",  TEAL,
         "USR-156, USR-234",
         "Behavioral zone analysis detects which dimensions are drifting — "
         "catches insider data_behavior drift and APT network_footprint drift. 6.5% FP."),
    ]

    x_pos = [Inches(0.5), Inches(6.9)]
    for i, (method, color, target, desc) in enumerate(components):
        x = x_pos[i]
        box = add_rounded_rect(slide, x, Inches(1.95), Inches(6.0), Inches(3.0))
        box.line.fill.solid()
        box.line.color.rgb = color
        box.line.width = Pt(2)
        add_shape(slide, x, Inches(1.95), Inches(6.0), Inches(0.07), fill_color=color)

        add_textbox(slide, x + Inches(0.2), Inches(2.1), Inches(5.6), Inches(0.45),
                    method, font_size=18, bold=True, color=color, alignment=PP_ALIGN.LEFT)
        add_textbox(slide, x + Inches(0.2), Inches(2.55), Inches(5.6), Inches(0.4),
                    f"Catches: {target}", font_size=15, bold=True, color=NAVY,
                    alignment=PP_ALIGN.LEFT)
        add_textbox(slide, x + Inches(0.2), Inches(3.05), Inches(5.6), Inches(0.9),
                    desc, font_size=13, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)

    # Plus sign between boxes
    add_textbox(slide, Inches(6.3), Inches(3.0), Inches(0.5), Inches(0.6),
                "+", font_size=40, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)

    # Combined result
    combined = add_rounded_rect(slide, Inches(0.5), Inches(5.15), Inches(12.3), Inches(1.45),
                                fill_color=NAVY)
    combined.line.fill.background()
    set_text(combined, "Combined Result: LOF + Zone Divergence",
             font_size=20, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)
    combined.text_frame.margin_top = Inches(0.1)
    add_paragraph(combined.text_frame,
                  "USR-156: DET  |  USR-234: DET  |  USR-042: DET  |  USR-118: DET",
                  font_size=17, bold=False, color=WHITE, alignment=PP_ALIGN.CENTER,
                  space_before=Pt(8))
    add_paragraph(combined.text_frame,
                  "False Positive Rate: 6.5%  (3 of 46 normal users)  "
                  "— operationally investigable",
                  font_size=15, bold=False, color=TEAL, alignment=PP_ALIGN.CENTER,
                  space_before=Pt(6))

    add_footer(slide)


# ---------------------------------------------------------------------------
# Slide 16: Complementary Detection Visualization
# ---------------------------------------------------------------------------
def slide_16_complementary_visualization(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Complementary Detection — Each Method Covers a Different Attack Surface")

    # Three-column coverage grid
    cols_data = [
        ("LOF\n(Traditional — Tier 1)",
         BLUE, [
             ("USR-156 (Insider)", "MISSED", RED),
             ("USR-234 (APT)", "DETECTED", GREEN),
             ("USR-042 (Volt)", "DETECTED", GREEN),
             ("USR-118 (Salt)", "DETECTED", GREEN),
             ("FP Rate", "0.0%", GREEN),
         ]),
        ("Zone Divergence\n(Behavioral — Tier 3)",
         TEAL, [
             ("USR-156 (Insider)", "DETECTED", GREEN),
             ("USR-234 (APT)", "DETECTED", GREEN),
             ("USR-042 (Volt)", "MISSED", RED),
             ("USR-118 (Salt)", "MISSED", RED),
             ("FP Rate", "6.5%", GREEN),
         ]),
        ("LOF + Zone Divergence\nOptimal Ensemble",
         GOLD, [
             ("USR-156 (Insider)", "DETECTED", GREEN),
             ("USR-234 (APT)", "DETECTED", GREEN),
             ("USR-042 (Volt)", "DETECTED", GREEN),
             ("USR-118 (Salt)", "DETECTED", GREEN),
             ("FP Rate", "6.5%", GREEN),
         ]),
    ]

    x_positions = [Inches(0.35), Inches(4.65), Inches(8.95)]
    for i, (title, color, rows) in enumerate(cols_data):
        x = x_positions[i]
        box = add_rounded_rect(slide, x, Inches(1.3), Inches(4.1), Inches(4.8))
        box.line.fill.solid()
        box.line.color.rgb = color
        box.line.width = Pt(2)
        add_shape(slide, x, Inches(1.3), Inches(4.1), Inches(0.07), fill_color=color)

        add_textbox(slide, x + Inches(0.15), Inches(1.45), Inches(3.8), Inches(0.7),
                    title, font_size=14, bold=True, color=NAVY, alignment=PP_ALIGN.CENTER)
        add_shape(slide, x + Inches(0.3), Inches(2.15), Inches(3.5), Inches(0.02),
                  fill_color=LIGHT_GRAY)

        y = Inches(2.3)
        for label, value, vcolor in rows:
            add_textbox(slide, x + Inches(0.2), y, Inches(2.2), Inches(0.38),
                        label, font_size=12, bold=False, color=DARK_GRAY,
                        alignment=PP_ALIGN.LEFT)
            val_box = add_shape(slide, x + Inches(2.4), y, Inches(1.5), Inches(0.32),
                                fill_color=vcolor)
            set_text(val_box, value, font_size=11, bold=True, color=WHITE,
                     alignment=PP_ALIGN.CENTER)
            y += Inches(0.52)

    # Separator lines
    add_shape(slide, Inches(4.55), Inches(1.3), Inches(0.02), Inches(4.8), fill_color=LIGHT_GRAY)
    add_shape(slide, Inches(8.85), Inches(1.3), Inches(0.02), Inches(4.8), fill_color=LIGHT_GRAY)

    add_callout_box(slide, Inches(1.5), Inches(6.4), Inches(10.3), Inches(0.65),
                    "LOF catches network footprints. Zone Divergence catches behavioral drift. Together: 100%.",
                    font_size=18)
    add_footer(slide)


# ---------------------------------------------------------------------------
# Slide 17: False Positive Rate Comparison
# ---------------------------------------------------------------------------
def slide_17_fp_comparison(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "False Positive Rate Comparison — All Methods")

    add_textbox(slide, Inches(0.7), Inches(1.25), Inches(11.9), Inches(0.4),
                "Lower FP rate = more actionable alerts. "
                "46 normal users in the population (4 attack users).",
                font_size=15, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)

    methods = [
        ("Isolation Forest",          2.2,   1,  GREEN),
        ("One-Class SVM",             19.6,  9,  RED),
        ("LOF",                       0.0,   0,  GREEN),
        ("Z-Score (|z|>3)",           2.2,   1,  GREEN),
        ("Temporal Z-Score",          100.0, 46, RED),
        ("Feature Trajectory Top-10%", 6.5,   3,  GREEN),
        ("V-Intelligence UEBA Direction",         4.3,   2,  GREEN),
        ("T3 Zone Divergence",        6.5,   3,  GREEN),
        ("T3 Combined",               8.7,   4,  GREEN),
        ("LOF + Zone Divergence",     6.5,   3,  TEAL),
    ]

    y = Inches(1.9)
    bar_max = Inches(7.5)
    label_w = Inches(3.4)
    bar_start = Inches(3.8)
    count_start = bar_start + bar_max + Inches(0.15)

    for method, rate, count, color in methods:
        add_textbox(slide, Inches(0.6), y, label_w, Inches(0.36),
                    method, font_size=12, bold=False, color=NAVY, alignment=PP_ALIGN.LEFT)
        bar_w = Inches(7.5 * rate / 100.0)
        add_shape(slide, bar_start, y + Inches(0.04), bar_w, Inches(0.28), fill_color=color)
        add_textbox(slide, bar_start + bar_w + Inches(0.1), y, Inches(1.8), Inches(0.36),
                    f"{rate}%  ({count} users)", font_size=12, bold=True, color=color,
                    alignment=PP_ALIGN.LEFT)
        y += Inches(0.45)

    # Legend
    add_shape(slide, Inches(0.6), Inches(6.7), Inches(0.2), Inches(0.2), fill_color=GREEN)
    add_textbox(slide, Inches(0.85), Inches(6.68), Inches(2.0), Inches(0.28),
                "Operationally viable", font_size=11, bold=False, color=DARK_GRAY,
                alignment=PP_ALIGN.LEFT)
    add_shape(slide, Inches(3.0), Inches(6.7), Inches(0.2), Inches(0.2), fill_color=TEAL)
    add_textbox(slide, Inches(3.25), Inches(6.68), Inches(2.5), Inches(0.28),
                "Combined (best overall)", font_size=11, bold=False, color=DARK_GRAY,
                alignment=PP_ALIGN.LEFT)
    add_shape(slide, Inches(5.8), Inches(6.7), Inches(0.2), Inches(0.2), fill_color=RED)
    add_textbox(slide, Inches(6.05), Inches(6.68), Inches(2.5), Inches(0.28),
                "Non-actionable alert volume", font_size=11, bold=False, color=DARK_GRAY,
                alignment=PP_ALIGN.LEFT)

    add_footer(slide)


# ---------------------------------------------------------------------------
# Slide 18: Key Findings Summary
# ---------------------------------------------------------------------------
def slide_18_key_findings(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_textbox(slide, Inches(0.7), Inches(0.35), Inches(12.0), Inches(0.7),
                "Key Findings", font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)
    add_shape(slide, Inches(0.7), Inches(1.0), Inches(3.0), Inches(0.04), fill_color=GOLD)

    findings = [
        (1, GOLD,  "LOF is the Best Traditional Method — But Cannot Detect Insiders",
         "LOF detected 3 of 4 attacks (APT, Volt Typhoon, Salt Typhoon) at 0% FP. But USR-156 "
         "(insider threat) is completely invisible to all traditional methods — behavioral "
         "direction attacks stay within normal volume ranges by design."),
        (2, ORANGE, "No Single Method Detects All Four Campaign Types",
         "Across 17 methods in 3 tiers, no single algorithm detects all 4 attacks at viable FP. "
         "Temporal Z-Score catches all 4 but at 100% FP — unusable. Each method has blind spots "
         "in different signal categories."),
        (3, TEAL,  "Zone-Specific Behavioral Analysis Fills the Detection Gap",
         "Tier 3 Zone Divergence detects USR-156 and USR-234 by analyzing which behavioral "
         "dimensions are drifting (data_behavior, network_footprint) while identity remains "
         "stable. This is invisible to aggregate-level traditional analysis."),
        (4, GREEN, "LOF + Zone Divergence = Optimal 2-Method Ensemble",
         "This combination detects ALL 4 attack campaigns at only 6.5% FP. LOF catches "
         "network-footprint attacks, Zone Divergence catches behavioral-direction attacks. "
         "The two methods are perfectly complementary."),
    ]

    y = Inches(1.25)
    for num, color, headline, detail in findings:
        num_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                            Inches(0.5), y + Inches(0.1),
                                            Inches(0.55), Inches(0.55))
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = color
        num_circle.line.fill.background()
        set_text(num_circle, str(num), font_size=18, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)

        add_textbox(slide, Inches(1.2), y, Inches(11.4), Inches(0.38),
                    headline, font_size=16, bold=True, color=color, alignment=PP_ALIGN.LEFT)
        add_textbox(slide, Inches(1.2), y + Inches(0.38), Inches(11.4), Inches(0.65),
                    detail, font_size=13, bold=False, color=MID_GRAY, alignment=PP_ALIGN.LEFT)
        y += Inches(1.3)

    add_textbox(slide, Inches(0.5), Inches(7.1), Inches(12), Inches(0.3),
                "22nd Century Technologies, Inc.  |  V-Intelligence UEBA Technical Analysis  |  Confidential",
                font_size=10, bold=False, color=RGBColor(60, 70, 85), alignment=PP_ALIGN.LEFT)


# ---------------------------------------------------------------------------
# Slide 19: Implications for Federal Cybersecurity
# ---------------------------------------------------------------------------
def slide_19_federal_implications(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Implications for Federal Cybersecurity Programs")

    implications = [
        (TEAL,   "Current UEBA and SIEM Tooling Has a Structural Detection Gap",
                 "If your deployed anomaly detection relies on statistical outlier detection "
                 "or threshold-based alerting, behavioral-direction attacks — including slow "
                 "insider threats and patient APT campaigns — will not generate alerts. "
                 "This is not a configuration problem; it is a method limitation."),
        (RED,    "Nation-State Campaigns Are Designed to Exploit This Gap",
                 "Adversaries operating at the level of Volt Typhoon and Salt Typhoon explicitly "
                 "model normal user behavior. Their operations are designed to remain within "
                 "statistical norms. The detection method must evolve to match the threat model."),
        (GREEN,  "Behavioral Analysis Provides a Proven, Complementary Layer",
                 "This analysis demonstrates that behavioral drift analysis can detect "
                 "insider threats and slow adversarial campaigns at operationally viable "
                 "false positive rates — without requiring new data collection or "
                 "additional agents on endpoints."),
        (GOLD,   "A Combined Detection Architecture Is the Operationally Sound Path",
                 "The empirical finding is clear: no single algorithm achieves both full "
                 "coverage and actionable precision. A layered detection architecture "
                 "that pairs traditional anomaly methods with behavioral drift analysis "
                 "is the only approach validated across both attack surface types."),
    ]

    y = Inches(1.45)
    for color, headline, detail in implications:
        add_shape(slide, Inches(0.65), y, Inches(0.1), Inches(1.05), fill_color=color)
        add_textbox(slide, Inches(0.95), y, Inches(11.7), Inches(0.4),
                    headline, font_size=15, bold=True, color=NAVY, alignment=PP_ALIGN.LEFT)
        add_textbox(slide, Inches(0.95), y + Inches(0.4), Inches(11.7), Inches(0.65),
                    detail, font_size=13, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)
        y += Inches(1.3)

    add_footer(slide)


# ---------------------------------------------------------------------------
# Slide 20: Next Steps
# ---------------------------------------------------------------------------
def slide_20_next_steps(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Next Steps — 4-Week Demonstration on Your Environment")

    add_textbox(slide, Inches(0.7), Inches(1.25), Inches(11.9), Inches(0.42),
                "We proved these results on synthetic telemetry. "
                "We can prove them on your data in 4 weeks.",
                font_size=16, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)

    phases = [
        ("Week 1–2", "Integration & Baseline",
         "Connect to existing SIEM or log feeds\n"
         "Establish behavioral baselines per entity\n"
         "No new agents or data collection required\n"
         "Zero operational disruption to production",
         TEAL),
        ("Week 3", "Detection Calibration",
         "Tune detection thresholds to your environment\n"
         "Validate against any known historical incidents\n"
         "Measure false positive rate on live population\n"
         "Compare directly to current SIEM alert volume",
         BLUE),
        ("Week 4", "Results & Decision",
         "Full detection report with empirical metrics\n"
         "FP rate versus current alerting baseline\n"
         "Mean time-to-detect analysis\n"
         "Go/no-go recommendation with supporting data",
         GOLD),
    ]

    x_positions = [Inches(0.5), Inches(4.55), Inches(8.6)]
    for i, (timeline, title, details, color) in enumerate(phases):
        x = x_positions[i]

        badge = add_rounded_rect(slide, x + Inches(0.7), Inches(1.9), Inches(2.4), Inches(0.48),
                                 fill_color=color)
        set_text(badge, timeline, font_size=15, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)

        card = add_rounded_rect(slide, x, Inches(2.55), Inches(3.9), Inches(3.3))
        card.line.fill.solid()
        card.line.color.rgb = color
        card.line.width = Pt(2)
        add_shape(slide, x, Inches(2.55), Inches(3.9), Inches(0.06), fill_color=color)

        add_textbox(slide, x + Inches(0.15), Inches(2.72), Inches(3.6), Inches(0.45),
                    title, font_size=17, bold=True, color=NAVY, alignment=PP_ALIGN.CENTER)
        add_shape(slide, x + Inches(0.4), Inches(3.17), Inches(3.1), Inches(0.02),
                  fill_color=LIGHT_GRAY)
        add_textbox(slide, x + Inches(0.2), Inches(3.3), Inches(3.5), Inches(2.2),
                    details, font_size=13, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)

        if i < 2:
            arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                         x + Inches(3.9) + Inches(0.05),
                                         Inches(3.9), Inches(0.5), Inches(0.28))
            arr.fill.solid()
            arr.fill.fore_color.rgb = GOLD
            arr.line.fill.background()

    diff_box = add_rounded_rect(slide, Inches(0.5), Inches(6.05), Inches(12.3), Inches(1.0),
                                fill_color=NAVY)
    set_text(diff_box,
             "No new agents.  No new data collection.  Works with existing log infrastructure.",
             font_size=18, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)
    diff_box.text_frame.margin_top = Inches(0.08)
    add_paragraph(diff_box.text_frame,
                  "Success criteria: detection rate, false positive rate, mean time to detect — "
                  "all measurable and reported.",
                  font_size=14, bold=False, color=WHITE, alignment=PP_ALIGN.CENTER,
                  space_before=Pt(6))

    add_footer(slide)


# ---------------------------------------------------------------------------
# Slide 21: UEBA Framework
# ---------------------------------------------------------------------------
def slide_21_ueba_framework(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "UEBA: Why Behavioral Analytics Detects What Rules Cannot")

    add_textbox(slide, Inches(0.7), Inches(1.25), Inches(11.9), Inches(0.45),
                "User and Entity Behavior Analytics (UEBA) — detect threats by behavioral drift, "
                "not attack signatures.",
                font_size=15, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)

    # Left: What is UEBA
    lbox = add_rounded_rect(slide, Inches(0.4), Inches(1.85), Inches(5.8), Inches(2.4))
    lbox.line.fill.solid()
    lbox.line.color.rgb = BLUE
    lbox.line.width = Pt(2)
    add_shape(slide, Inches(0.4), Inches(1.85), Inches(5.8), Inches(0.06), fill_color=BLUE)

    add_textbox(slide, Inches(0.6), Inches(2.0), Inches(5.4), Inches(0.35),
                "What is UEBA?", font_size=19, bold=True, color=BLUE)
    add_textbox(slide, Inches(0.6), Inches(2.4), Inches(5.4), Inches(1.6),
                "Learn what normal looks like for every user and device. "
                "Flag when behavior drifts from that baseline.\n\n"
                "SIEM rules ask: 'Did a known bad thing happen?'\n"
                "UEBA asks: 'Is this entity behaving differently than before?'\n\n"
                "Modern attackers use valid credentials, legitimate tools, "
                "and authorized access. They don't trigger signatures — "
                "the only signal is behavioral change over time.",
                font_size=13, bold=False, color=DARK_GRAY)

    # Right: Why SIEM fails
    rbox = add_rounded_rect(slide, Inches(6.5), Inches(1.85), Inches(6.4), Inches(2.4))
    rbox.line.fill.solid()
    rbox.line.color.rgb = RED
    rbox.line.width = Pt(2)
    add_shape(slide, Inches(6.5), Inches(1.85), Inches(6.4), Inches(0.06), fill_color=RED)

    add_textbox(slide, Inches(6.7), Inches(2.0), Inches(6.0), Inches(0.35),
                "Why SIEM Rules Fail", font_size=19, bold=True, color=RED)
    add_textbox(slide, Inches(6.7), Inches(2.4), Inches(6.0), Inches(1.6),
                "Insider: All access uses valid credentials — no rule violation\n"
                "Slow APT: C2 is small, periodic, encrypted — looks like HTTPS\n"
                "LOTL: No malware — uses PowerShell, WMI, RDP from day one\n"
                "Telecom: Authorized maintenance credentials to network devices\n\n"
                "Common thread: attacker's identity stays the same, "
                "but their behavior changes.",
                font_size=13, bold=False, color=DARK_GRAY)

    # Bottom: 4 threat types with UEBA signals
    threat_data = [
        ["Threat", "Attack Pattern", "UEBA Signal", "Duration"],
        ["Insider Threat", "Gradual data access\nescalation", "Identity stable +\ndata zone drifting", "8 months"],
        ["Slow APT / C2", "Persistent C2\nbeaconing", "Identity stable +\nnetwork zone drifting", "180 days"],
        ["Nation-State LOTL", "Admin tools\nrepurposed", "All zones shift\n(phase break)", "115 days"],
        ["Telecom Pivot", "Infrastructure\ntargeting", "Broad multi-zone\ncumulative drift", "100 days"],
    ]

    threat_colors = {}
    threat_row_colors = [RED, RGBColor(0xE6, 0x7E, 0x22), RGBColor(0x8E, 0x44, 0xAD), RGBColor(0x29, 0x80, 0xB9)]
    for r in range(1, 5):
        threat_colors[(r, 0)] = threat_row_colors[r - 1]

    build_table(slide, 5, 4, Inches(0.4), Inches(4.5), Inches(12.5), Inches(2.3),
                threat_data, col_widths=[Inches(2.4), Inches(3.5), Inches(4.2), Inches(2.4)],
                cell_colors=threat_colors, body_font_size=12)

    add_footer(slide)


# ---------------------------------------------------------------------------
# Slide 22: Detection Playbook
# ---------------------------------------------------------------------------
def slide_22_detection_playbook(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Detection Playbook: Match the Threat to the Algorithm")

    add_textbox(slide, Inches(0.7), Inches(1.25), Inches(11.9), Inches(0.45),
                "No single algorithm catches all 4 threat types. "
                "The threat type determines which method to deploy.",
                font_size=15, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)

    playbook_data = [
        ["Threat Type", "Best Method", "FP Rate", "Why It Works", "What Fails"],
        ["Insider\n(USR-156)", "Zone\nDivergence", "6.5%",
         "Only method that sees\n'same person, different data'",
         "LOF, IForest, Z-Score\n(measure magnitude only)"],
        ["Slow APT\n(USR-234)", "LOF + Zone\nDivergence", "0%* /\n6.5%",
         "LOF: network outlier\nZone Div: network zone drifting",
         "IForest\n(not sensitive to slow drift)"],
        ["LOTL\n(USR-042)", "LOF +\nRegime Shift", "0%* /\n6.5%",
         "LOF: endpoint anomaly\nRegime Shift: phase break",
         "Zone Divergence\n(uniform, not zone-specific)"],
        ["Telecom\n(USR-118)", "LOF +\nEmbed CUSUM", "0%* /\n6.5%",
         "LOF: network footprint\nCUSUM: persistent 100-day drift",
         "Zone Divergence\n(broad multi-zone change)"],
    ]

    pb_colors = {}
    row_colors = [RED, RGBColor(0xE6, 0x7E, 0x22), RGBColor(0x8E, 0x44, 0xAD), RGBColor(0x29, 0x80, 0xB9)]
    for r in range(1, 5):
        pb_colors[(r, 0)] = row_colors[r - 1]
    for r in range(1, 5):
        pb_colors[(r, 2)] = GREEN

    build_table(slide, 5, 5, Inches(0.3), Inches(1.85), Inches(12.7), Inches(2.8),
                playbook_data,
                col_widths=[Inches(1.8), Inches(2.0), Inches(1.2), Inches(4.0), Inches(3.7)],
                cell_colors=pb_colors, body_font_size=12, header_font_size=13)

    # Layered strategy
    add_textbox(slide, Inches(0.5), Inches(4.9), Inches(12.3), Inches(0.35),
                "Layered Deployment Strategy", font_size=18, bold=True, color=NAVY)

    layers = [
        ("Layer 1: Magnitude\n(LOF)", BLUE,
         "How much changed?\nCatches LOTL, Telecom, APT\nAlways-on baseline"),
        ("Layer 2: Direction\n(Zone Divergence)", RGBColor(0x8E, 0x44, 0xAD),
         "What kind of change?\nCatches Insider, APT\nZone-decomposed drift"),
        ("Layer 3: Accumulation\n(CUSUM + Regime)", TEAL,
         "How long has it persisted?\nCatches long campaigns\nPhase change detection"),
    ]

    for i, (title, color, desc) in enumerate(layers):
        x = Inches(0.4 + i * 4.3)
        card = add_rounded_rect(slide, x, Inches(5.35), Inches(3.9), Inches(1.5))
        card.line.fill.solid()
        card.line.color.rgb = color
        card.line.width = Pt(2)
        add_shape(slide, x, Inches(5.35), Inches(3.9), Inches(0.06), fill_color=color)
        add_textbox(slide, x + Inches(0.15), Inches(5.5), Inches(3.6), Inches(0.5),
                    title, font_size=15, bold=True, color=color, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.15), Inches(6.0), Inches(3.6), Inches(0.7),
                    desc, font_size=12, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    footnote = add_rounded_rect(slide, Inches(0.4), Inches(6.95), Inches(12.5), Inches(0.35),
                                fill_color=NAVY)
    set_text(footnote,
             "Minimum viable deployment: LOF + Zone Divergence → all 4 threats detected at 6.5% FP",
             font_size=14, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)

    add_footer(slide)


# ---------------------------------------------------------------------------
# Slide 23: Tier 1 Algorithm Parameter Reference
# ---------------------------------------------------------------------------
def slide_23_tier1_parameters(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Tier 1 Algorithm Parameters & Input Features")

    add_textbox(slide, Inches(0.7), Inches(1.25), Inches(11.9), Inches(0.42),
                "All 6 Tier 1 algorithms operate on the same 23 behavioral features "
                "extracted from 5 log sources. Each algorithm applies a different mathematical "
                "approach to find anomalies.",
                font_size=14, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)

    data = [
        ["Algorithm", "Type", "Key Parameters", "Input", "Detection Rule"],
        ["Isolation\nForest", "Static",
         "200 trees\n5% anomaly budget",
         "23 features\n(normalized)",
         "Partition outliers\nthat isolate quickly"],
        ["One-Class\nSVM", "Static",
         "RBF kernel\n5% outlier fraction",
         "23 features\n(normalized)",
         "Outside learned\nnormality boundary"],
        ["LOF", "Static",
         "20 neighbors\n5% anomaly budget",
         "23 features\n(normalized)",
         "Low local density\nvs neighbors"],
        ["Z-Score", "Static",
         "3-sigma threshold",
         "23 features\n(normalized)",
         "Any feature\n> 3 std dev"],
        ["Temporal\nZ-Score", "Temporal",
         "50/50 train/test split\n3-sigma threshold",
         "23 features\nper week",
         "Test-period deviation\nfrom own baseline"],
        ["Feature\nTrajectory", "Temporal",
         "Cumulative drift\n3-week minimum run",
         "Weekly composite\nfeature drift",
         "Top 10% by\naccumulated change"],
    ]

    cell_colors = {}
    for r in range(1, 5):
        cell_colors[(r, 1)] = TEAL
    for r in range(5, 7):
        cell_colors[(r, 1)] = BLUE

    build_table(slide, 7, 5, Inches(0.3), Inches(1.85), Inches(12.7), Inches(3.2),
                data, cell_colors=cell_colors,
                col_widths=[Inches(1.6), Inches(1.2), Inches(2.8), Inches(2.8), Inches(4.3)],
                body_font_size=11, header_font_size=12)

    # Feature categories
    add_textbox(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.35),
                "23 Input Features by Log Source", font_size=17, bold=True, color=NAVY)

    categories = [
        ("Authentication (7)", TEAL,
         "auth_total, auth_failed,\nfail_rate, off_hours,\nunique_sources,\nunique_dests, methods"),
        ("File Access (6)", BLUE,
         "file_total, restricted_ratio,\nconfidential_ratio,\nwrite_ratio, unique_paths,\ntotal_bytes"),
        ("Endpoint (5)", RGBColor(0x8E, 0x44, 0xAD),
         "endpoint_total,\nsuspicious_ratio,\nmax_risk, mean_risk,\nunique_processes"),
        ("Network (5)", NAVY,
         "net_bytes_out,\nunique_dsts,\nexternal_ratio,\ndns_domains, nxdomain"),
    ]

    for i, (title, color, feats) in enumerate(categories):
        x = Inches(0.3 + i * 3.25)
        card = add_rounded_rect(slide, x, Inches(5.6), Inches(3.0), Inches(1.6))
        card.line.fill.solid()
        card.line.color.rgb = color
        card.line.width = Pt(2)
        add_shape(slide, x, Inches(5.6), Inches(3.0), Inches(0.06), fill_color=color)
        add_textbox(slide, x + Inches(0.1), Inches(5.72), Inches(2.8), Inches(0.35),
                    title, font_size=14, bold=True, color=color, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.15), Inches(6.1), Inches(2.7), Inches(1.0),
                    feats, font_size=11, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)

    add_footer(slide)


# ---------------------------------------------------------------------------
# Slide 24: Tier 3 Method-to-Zone Reference
# ---------------------------------------------------------------------------
def slide_24_tier3_parameters(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Tier 3 Detection Methods — Inputs & Parameters")

    add_textbox(slide, Inches(0.7), Inches(1.25), Inches(11.9), Inches(0.42),
                "Each Tier 3 method receives specific zone/trajectory/relationship inputs "
                "and applies independent detection logic. All use top-10% rank thresholding.",
                font_size=14, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)

    data = [
        ["Method", "Input Signal", "What It Measures", "Key Parameter", "Detection"],
        ["Velocity /\nAcceleration", "Composite trajectory\n(weekly snapshots)",
         "How fast behavior\nis changing direction",
         "Acceleration > 0.01\nwith consistent trend",
         "Top 10% by\nweighted score"],
        ["Regime\nShift", "Consecutive weekly\nembeddings",
         "Fundamental phase\nchange in behavior",
         "Similarity < 0.7\n= phase break",
         "Top 10% by\ninstability score"],
        ["Zone\nDivergence", "5 per-zone\ntrajectory drifts",
         "Which zone drifts\nwhile identity is stable",
         "Identity < 0.02 (stable)\nZone > 0.05 (drifting)",
         "Top 10% OR\nstandalone flag"],
        ["Relationship\nDrift", "User-Device\ninteraction vector",
         "Changing interaction\npattern between entities",
         "Relationship\ndrift > 0.05",
         "Top 10% by\ndrift magnitude"],
        ["Contextual\n(Multi-Context)", "4 context-weighted\ncompositions",
         "Attack-specific signals\namplified by context",
         "Threat consistency\n> 30% of weeks",
         "Top 10% by\nbest context"],
        ["Embedding\nDrift Accum.", "Baseline vs weekly\nsnapshot distance",
         "Persistent slow drift\nthat accumulates",
         "Cumulative excess\ndrift from baseline",
         "Top 10% by\naccumulated drift"],
        ["Behavioral\nProgression", "Per-zone weekly\nthreat alignment",
         "Monotonically increasing\nthreat-like behavior",
         "Kendall trend\ncorrelation > 0",
         "Top 10% by\ntrend strength"],
        ["Combined\n(Ensemble)", "All 6 core method\nscores",
         "Corroborating evidence\nfrom multiple signals",
         "Zone Div standalone\nOR ≥ 2 methods agree",
         "Zone Div alone\nOR consensus"],
    ]

    cell_colors = {}
    method_colors = [TEAL, RGBColor(0x8E, 0x44, 0xAD), RED, BLUE,
                     GOLD, RGBColor(0x8E, 0x44, 0xAD), TEAL, NAVY]
    for r in range(1, 9):
        cell_colors[(r, 0)] = method_colors[r - 1]

    build_table(slide, 9, 5, Inches(0.3), Inches(1.85), Inches(12.7), Inches(4.0),
                data, cell_colors=cell_colors,
                col_widths=[Inches(1.8), Inches(2.2), Inches(2.5), Inches(3.0), Inches(3.2)],
                body_font_size=10, header_font_size=11)

    # Zone input reference
    zone_bar = add_rounded_rect(slide, Inches(0.3), Inches(6.05), Inches(12.7), Inches(1.2),
                                fill_color=NAVY)
    set_text(zone_bar,
             "5 Behavioral Zones (input to Zone Divergence, Contextual, Progression)",
             font_size=15, bold=True, color=GOLD, alignment=PP_ALIGN.LEFT)
    zone_bar.text_frame.margin_left = Inches(0.3)
    zone_bar.text_frame.margin_top = Inches(0.1)

    zone_text = (
        "Identity (role, dept, clearance)  |  Access (auth patterns)  |  "
        "Data (file access behavior)  |  Network (traffic, DNS)  |  Risk (endpoint health)"
    )
    add_paragraph(zone_bar.text_frame, zone_text,
                  font_size=13, bold=False, color=WHITE, space_before=Pt(6))
    add_paragraph(zone_bar.text_frame,
                  "Context-adaptive weights: Normal (equal) → Insider (data=40%) → "
                  "APT Hunt (network=40%) → Privilege Audit (risk=40%)",
                  font_size=12, bold=False, color=RGBColor(0xBB, 0xBB, 0xBB),
                  space_before=Pt(4))

    add_footer(slide)


# ===========================================================================
# MAIN
# ===========================================================================

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide_01_title(prs)
    slide_02_exec_summary(prs)
    slide_03_synthetic_data_overview(prs)
    slide_04_log_types_feature_engineering(prs)
    slide_05_feature_detail_table(prs)
    slide_06_attack_scenarios(prs)
    slide_07_attack_design_philosophy(prs)
    slide_08_traditional_methods(prs)
    slide_09_temporal_methods(prs)
    slide_10_acecard_behavioral(prs)
    slide_11_full_results_table(prs)
    slide_12_traditional_deep_dive(prs)
    slide_13_temporal_deep_dive(prs)
    slide_14_acecard_direction_results(prs)
    slide_15_combined_approach(prs)
    slide_16_complementary_visualization(prs)
    slide_17_fp_comparison(prs)
    slide_18_key_findings(prs)
    slide_19_federal_implications(prs)
    slide_20_next_steps(prs)
    slide_21_ueba_framework(prs)
    slide_22_detection_playbook(prs)
    slide_23_tier1_parameters(prs)
    slide_24_tier3_parameters(prs)

    out_dir  = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.join(out_dir, "V_Intelligence_UEBA_Technical_Deep_Dive.pptx")
    prs.save(out_path)
    print(f"Deck created: {out_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
