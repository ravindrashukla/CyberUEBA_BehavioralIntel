# -*- coding: utf-8 -*-
"""Extensive Pentagon working-session deck (Army Cyber: Esquibel / DiGrezio / Herrmann).
Four-layer, one-company framing; digital twin as the anchor. No partner names."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# DLA-deck-matched palette (Segoe UI, 22nd Century orange + navy + teal + green)
NAVY=RGBColor(0x10,0x2a,0x4c); DARKNAVY=RGBColor(0x0c,0x20,0x3b); TOPRULE=RGBColor(0x06,0x24,0x4a)
BLUE=RGBColor(0x1f,0x5e,0x9c); TEAL=RGBColor(0x11,0x84,0x84); SLATE=RGBColor(0x5e,0x6a,0x75)
WHITE=RGBColor(0xff,0xff,0xff); LIGHT=RGBColor(0xec,0xf4,0xfc); GREEN=RGBColor(0x1a,0x72,0x4e)
RED=RGBColor(0xa6,0x1b,0x1b); ORANGE=RGBColor(0xe8,0x77,0x22); AMBER=RGBColor(0xe8,0x77,0x22); GOLD=RGBColor(0xe8,0x77,0x22)
L1=RGBColor(0x1f,0x5e,0x9c); L2=RGBColor(0x11,0x84,0x84); L3=RGBColor(0x1a,0x72,0x4e); L4=RGBColor(0xc9,0x6a,0x1e)
FT="Segoe UI"

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]
SW,SH=prs.slide_width,prs.slide_height

def slide():
    return prs.slides.add_slide(BLANK)
def rect(s,x,y,w,h,fill,line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sp=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,x,y,w,h)
    sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1)
    sp.shadow.inherit=False
    return sp
def txt(s,x,y,w,h,text,size=18,color=NAVY,bold=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,italic=False,font="Segoe UI"):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
    tf.vertical_anchor=anchor
    lines=text.split("\n")
    for i,ln in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; r=p.add_run(); r.text=ln
        r.font.size=Pt(size); r.font.bold=bold; r.font.italic=italic; r.font.color.rgb=color; r.font.name=font
    return tb
def bullets(s,x,y,w,h,items,size=16,color=RGBColor(0x22,0x22,0x22),gap=6):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
    for i,(lvl,t,*c) in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.level=lvl; p.space_after=Pt(gap)
        r=p.add_run(); r.text=(chr(0x2022)+" " if lvl==0 else "   - ")+t
        r.font.size=Pt(size-(2 if lvl else 0)); r.font.color.rgb=(c[0] if c else color); r.font.name="Segoe UI"
        if c and len(c)>1 and c[1]: r.font.bold=True
    return tb
def header(s,title,kicker=None):
    rect(s,0,0,SW,Pt(7),TOPRULE)
    txt(s,Inches(0.5),Inches(0.28),Inches(12.3),Inches(0.7),title,size=23,color=NAVY,bold=True)
    if kicker:
        txt(s,Inches(0.5),Inches(0.98),Inches(12.3),Inches(0.35),kicker,size=12,color=SLATE)
    rect(s,Inches(0.5),Inches(1.4),Inches(2.6),Pt(3),ORANGE)
def footer(s,n):
    rect(s,0,Inches(7.33),SW,Pt(1.2),RGBColor(0xd8,0xe0,0xec))
    txt(s,Inches(0.5),Inches(7.06),Inches(11),Inches(0.3),"CONFIDENTIAL  |  22nd Century Technologies  |  V-Intelligence Army Cyber Working Session",size=9,color=SLATE)
    txt(s,Inches(12.4),Inches(7.06),Inches(0.7),Inches(0.3),str(n),size=9,color=SLATE,align=PP_ALIGN.RIGHT)

def dtable(s,x,y,w,h,headers,rows,fs=11):
    t=s.shapes.add_table(len(rows)+1,len(headers),x,y,w,h).table
    for j,hd in enumerate(headers):
        c=t.cell(0,j); c.text=hd
        rn=c.text_frame.paragraphs[0].runs[0]; rn.font.size=Pt(fs); rn.font.bold=True; rn.font.name="Segoe UI"
    for i,row in enumerate(rows,1):
        for j,v in enumerate(row):
            c=t.cell(i,j); c.text=str(v)
            r0=c.text_frame.paragraphs[0].runs[0]; r0.font.size=Pt(fs); r0.font.name="Segoe UI"
    return t

n=0
def fin(s):
    global n; n+=1; footer(s,n)

# ---- 1 TITLE ----
s=slide(); rect(s,0,0,SW,SH,WHITE)
rect(s,0,0,Inches(0.22),SH,ORANGE)
rect(s,Inches(8.5),0,SW-Inches(8.5),SH,DARKNAVY)
txt(s,Inches(0.7),Inches(1.55),Inches(7.4),Inches(1.7),"Defense in Depth for the Age of AI-Enabled Attacks",size=33,color=NAVY,bold=True)
txt(s,Inches(0.7),Inches(3.35),Inches(7.4),Inches(0.6),"A Four-Layer Architecture under One MITRE Vocabulary",size=18,color=BLUE)
txt(s,Inches(0.7),Inches(4.05),Inches(7.4),Inches(0.5),"Army Cyber Working Session  |  In-Person Technical Deep-Dive",size=14,color=SLATE)
rect(s,Inches(0.7),Inches(4.75),Inches(2.4),Pt(3),ORANGE)
txt(s,Inches(0.7),Inches(5.6),Inches(7.4),Inches(0.5),"Prepared for U.S. Army and Pentagon Cyber Mission Defenders",size=13,color=NAVY,italic=True)
txt(s,Inches(0.7),Inches(6.1),Inches(7.4),Inches(0.4),"V-Intelligence  |  22nd Century Technologies",size=11,color=SLATE)
_chips=[("One platform",BLUE),("Four layers",TEAL),("Behavioral core: 0 FP",ORANGE)]
yy=Inches(2.2)
for lbl,col in _chips:
    rect(s,Inches(9.0),yy,Inches(3.3),Inches(0.8),col)
    txt(s,Inches(9.2),yy,Inches(2.9),Inches(0.8),lbl,size=15,color=WHITE,bold=True,anchor=MSO_ANCHOR.MIDDLE)
    yy+=Inches(1.05)
fin(s)

# ---- 2 WHO'S IN THE ROOM ----
s=slide(); header(s,"Who Is in the Room","Three lenses we must satisfy in one session")
cards=[("Dr. Judy Esquibel","Senior Cyber Advisor, ARCYBER","Deeply technical operator. Detection fidelity, living-off-the-land, false-positive economics. Our toughest and most valuable seat.",L2),
       ("LTC Micah DiGrezio","Army acquisition","Can it be bought, integrated, fielded, sustained? Data, footprint, contract path.",L1),
       ("Ms. Kellsie Herrmann","Army cyber / technology policy","Why it matters, authorities, readiness. The strategic significance she can carry upward.",L3)]
x=Inches(0.5)
for nm,role,desc,col in cards:
    rect(s,x,Inches(1.5),Inches(4.0),Inches(4.6),LIGHT)
    rect(s,x,Inches(1.5),Inches(4.0),Inches(0.9),col)
    txt(s,x+Inches(0.2),Inches(1.6),Inches(3.6),Inches(0.7),nm,size=17,color=WHITE,bold=True,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,x+Inches(0.2),Inches(2.5),Inches(3.6),Inches(0.5),role,size=13,color=NAVY,bold=True)
    txt(s,x+Inches(0.2),Inches(3.1),Inches(3.6),Inches(2.8),desc,size=13,color=RGBColor(0x33,0x33,0x33))
    x+=Inches(4.27)
txt(s,Inches(0.5),Inches(6.0),Inches(12.3),Inches(0.9),"Audience: the Army cyber enterprise, the office of the Army's Principal Cyber Advisor. A strong showing is a direct line to a bounded pilot on real Army data.",size=13,color=SLATE,italic=True)
fin(s)

# ---- 3 THESIS ----
s=slide(); rect(s,0,0,SW,SH,LIGHT); rect(s,0,0,Inches(0.25),SH,GOLD)
txt(s,Inches(1.0),Inches(1.4),Inches(11.3),Inches(0.5),"THE THESIS WE ARE PROVING",size=16,color=BLUE,bold=True)
txt(s,Inches(1.0),Inches(2.1),Inches(11.3),Inches(2.6),
 "No single tool stops a Volt or a Salt, because the adversary moves across the whole "
 "lifecycle on valid credentials.",size=32,color=NAVY,bold=True)
txt(s,Inches(1.0),Inches(4.5),Inches(11.3),Inches(1.5),
 "We present one platform, four layers, one MITRE vocabulary, and the behavioral core runs live today.",
 size=22,color=GREEN,bold=True)
txt(s,Inches(1.0),Inches(6.1),Inches(11.3),Inches(0.5),"Golden rule: national problem first, then architecture, then proof. Show what runs. Label roadmap, out loud.",size=13,color=SLATE,italic=True)
fin(s)

# ---- 4 AGENDA ----
s=slide(); header(s,"Agenda","In-person, half-day (~4.5 hours), one company, one platform")
rows=[("0","Open: problem, arc, honesty boundary","20 min","Session lead"),
      ("1","Layer 2 - Behavioral Digital Twin (theory, live demo, results)","60 min","Twin lead"),
      ("2","Layer 1 - Preemptive Network Assurance","60 min","L1 lead"),
      ("3","Layer 3 - API, Data & Agentic-AI Runtime","60 min","L3 lead"),
      ("4","Layer 4 - Code & Supply-Chain (roadmap)","15 min","Session lead"),
      ("5","Integration + federal fit + the ask","30 min","Session lead"),
      ("-","Buffer","20 min","All")]
y=Inches(1.5)
for i,(num,blk,tm,lead) in enumerate(rows):
    bg=LIGHT if i%2==0 else WHITE
    rect(s,Inches(0.5),y,Inches(12.3),Inches(0.62),bg)
    txt(s,Inches(0.6),y,Inches(0.6),Inches(0.62),num,size=15,color=NAVY,bold=True,anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.CENTER)
    txt(s,Inches(1.3),y,Inches(8.0),Inches(0.62),blk,size=14,color=RGBColor(0x22,0x22,0x22),anchor=MSO_ANCHOR.MIDDLE)
    txt(s,Inches(9.4),y,Inches(1.3),Inches(0.62),tm,size=13,color=SLATE,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,Inches(10.8),y,Inches(2.0),Inches(0.62),lead,size=13,color=BLUE,anchor=MSO_ANCHOR.MIDDLE)
    y+=Inches(0.66)
fin(s)

# ---- 5 PROBLEM: INSIDE ----
s=slide(); header(s,"The Adversary Is Already Inside","Volt Typhoon, Salt Typhoon: valid accounts, no malware")
bullets(s,Inches(0.6),Inches(1.5),Inches(7.4),Inches(5),[
 (0,"Apex nation-state operators are resident in U.S. critical infrastructure now.",NAVY,True),
 (0,"They enter through unproven configuration gaps and unpatched edge devices.",),
 (0,"They operate with valid credentials and living-off-the-land techniques that cross no threshold.",),
 (0,"They fragment one identity across ten or more systems, and stay below every alert for years.",),
 (0,"Common technique: MITRE T1078, Valid Accounts. They log in; they do not break in.",GREEN,True),
],size=16,gap=12)
rect(s,Inches(8.4),Inches(1.6),Inches(4.4),Inches(4.6),LIGHT)
txt(s,Inches(8.6),Inches(1.8),Inches(4.0),Inches(0.5),"Army stake",size=15,color=BLUE,bold=True)
txt(s,Inches(8.6),Inches(2.4),Inches(4.0),Inches(3.6),
 "These campaigns pre-position in the commercial critical infrastructure that mobilization and power "
 "projection depend on, and they target the enterprise and tactical networks, and the IL5, IL6, and "
 "JWICS enclaves, that ARCYBER defends.",size=14,color=RGBColor(0x33,0x33,0x33))
fin(s)

# ---- 6 PROBLEM: AI TIMELINE ----
s=slide(); header(s,"Offensive AI Has Collapsed the Timeline","Detect-and-respond is too late against machine-speed attacks")
bullets(s,Inches(0.6),Inches(1.6),Inches(12),Inches(5),[
 (0,"Autonomous agents reconnoiter, find vulnerabilities, and exploit them at machine speed.",NAVY,True),
 (0,"A single operator can now act with the reach of a team.",),
 (0,"Threats evolve faster than vulnerability disclosure and patch cycles can respond.",),
 (0,"By the time an attack is detected, it is already too late.",RED,True),
 (0,"The only durable answer is preemptive defense plus behavioral detection, before and after the breach.",GREEN,True),
],size=18,gap=14)
fin(s)

# ---- 7 TWO GAPS ----
s=slide(); header(s,"Two Structural Gaps the Prior Model Leaves Open")
rect(s,Inches(0.6),Inches(1.6),Inches(5.9),Inches(4.5),LIGHT); rect(s,Inches(0.6),Inches(1.6),Inches(5.9),Inches(0.8),L1)
txt(s,Inches(0.8),Inches(1.7),Inches(5.5),Inches(0.6),"1. Preventive-proof gap",size=18,color=WHITE,bold=True,anchor=MSO_ANCHOR.MIDDLE)
txt(s,Inches(0.8),Inches(2.6),Inches(5.5),Inches(3.3),
 "Prevention tools sample the configuration and attack surface; they never prove it is closed. The "
 "adversary needs only the one path no one proved was shut.",size=15,color=RGBColor(0x33,0x33,0x33))
rect(s,Inches(6.8),Inches(1.6),Inches(5.9),Inches(4.5),LIGHT); rect(s,Inches(6.8),Inches(1.6),Inches(5.9),Inches(0.8),L2)
txt(s,Inches(7.0),Inches(1.7),Inches(5.5),Inches(0.6),"2. Behavioral-detection gap",size=18,color=WHITE,bold=True,anchor=MSO_ANCHOR.MIDDLE)
txt(s,Inches(7.0),Inches(2.6),Inches(5.5),Inches(3.3),
 "Detection tools threshold on magnitude and watch for spikes. A valid-account, living-off-the-land "
 "operator produces no spike and stays below every alert for years.",size=15,color=RGBColor(0x33,0x33,0x33))
txt(s,Inches(0.6),Inches(6.3),Inches(12),Inches(0.5),"The architecture closes both: preventive proof at the edge (Layer 1), behavioral watch inside (Layer 2).",size=14,color=NAVY,bold=True)
fin(s)

# ---- 8 ARCHITECTURE ----
s=slide(); header(s,"The Architecture: Four Layers, One Vocabulary")
data=[("Layer 1","Preemptive Network Assurance","Prove the door is closed before traffic flows","Fielded",L1),
      ("Layer 2","Behavioral Entity Intelligence (digital twin)","Catch valid-account, living-off-the-land behavior inside","LIVE demo",L2),
      ("Layer 3","Application, API & Data Runtime + Agentic AI","Guard APIs and data; govern the AI agents","Integrated",L3),
      ("Layer 4","Code & Supply-Chain Assurance","Find novel zero-days before deploy","In development",L4)]
y=Inches(1.55)
for ln,nm,desc,mat,col in data:
    rect(s,Inches(0.6),y,Inches(12.1),Inches(1.15),LIGHT)
    rect(s,Inches(0.6),y,Inches(2.0),Inches(1.15),col)
    txt(s,Inches(0.7),y,Inches(1.8),Inches(1.15),ln,size=17,color=WHITE,bold=True,anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.CENTER)
    txt(s,Inches(2.8),y+Inches(0.12),Inches(7.2),Inches(0.5),nm,size=15,color=NAVY,bold=True)
    txt(s,Inches(2.8),y+Inches(0.62),Inches(7.2),Inches(0.45),desc,size=12.5,color=RGBColor(0x44,0x44,0x44))
    mc=GREEN if mat in ("Live","LIVE demo","Fielded") else (AMBER if mat=="Integrated" else RED)
    txt(s,Inches(10.2),y,Inches(2.4),Inches(1.15),mat,size=13,color=mc,bold=True,anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.CENTER)
    y+=Inches(1.22)
txt(s,Inches(0.6),Inches(6.5),Inches(12),Inches(0.4),"One shared vocabulary across all four layers: MITRE ATT&CK + ATLAS. One platform, human-in-the-loop.",size=13,color=BLUE,bold=True)
fin(s)

# ---- 9 LAYER 2 WHAT ----
s=slide(); header(s,"Layer 2 - Behavioral Detection: Two Complementary Detectors","A precision threat-profile engine (production) and a discovery twin")
bullets(s,Inches(0.6),Inches(1.5),Inches(12),Inches(5),[
 (0,"Precision layer (production): a multi-front threat-profile detector matches named known-bad fingerprints, scored against role-group peers and fired only on corroboration. 4 of 4 attackers at 0 false positives.",GREEN,True),
 (0,"Discovery layer: a semantic behavioral digital twin reads behavior as meaning, tracks trajectory, and names direction in MITRE terms, an unsupervised net needing no prior fingerprint.",NAVY,True),
 (0,"The threat-profile detector gives the precise what-fired; the twin (next slides) gives the explainable how and where the entity is heading.",),
 (0,"Behavioral snapshots are stored as vectors in a PostgreSQL pgvector database; detection reads from the database.",),
],size=15,gap=14)
fin(s)

# ---- 9B PRODUCTION RESULT ----
s=slide(); header(s,"Layer 2 - Production Result: 4 of 4 at 0 False Positives","The threat-profile detector, by named fingerprint")
dtable(s,Inches(0.5),Inches(1.6),Inches(12.3),Inches(2.6),
 ["Attacker","Named fingerprint(s)","Fronts"],
 [["Insider (USR-156)","mass_collection (z 5.9) + cohort_rare_dst (76 IPs)","2"],
  ["Slow APT (USR-234)","c2_beacon (386-day persistent dst) + dga_dns (160 domains)","2"],
  ["LOTL / Volt (USR-042)","lotl_process (z 4.5)","1"],
  ["Salt (USR-118)","recon_fanout (z 8.2)","1"]],fs=12)
rect(s,Inches(0.5),Inches(4.5),Inches(12.3),Inches(1.5),NAVY)
txt(s,Inches(0.7),Inches(4.7),Inches(11.9),Inches(1.2),"4 of 4 caught, 0 false positives, precision 100 percent. Each profile is a measurable signature of a known technique, scored against the entity's role-group peers and fired only on a corroborated match.",size=15,color=WHITE,bold=True)
fin(s)

# ---- DT-1 PIPELINE ----
s=slide(); header(s,"Digital Twin in Detail: The Pipeline","Raw logs become a living behavioral representation, rebuilt weekly")
steps=[("Raw logs","auth, file, network, DNS, endpoint"),("5 zones","23 features as prose"),("Embed","each zone to 1536-d"),("Compose","attention-weighted twin vector"),("Trajectory","one composite per week"),("Drift + direction","CUSUM + MITRE projection"),("Verdict","one ranked, explainable score")]
x=Inches(0.45)
for nm,desc in steps:
    rect(s,x,Inches(2.5),Inches(1.66),Inches(2.4),LIGHT,line=BLUE)
    txt(s,x+Inches(0.07),Inches(2.62),Inches(1.52),Inches(0.85),nm,size=12.5,color=NAVY,bold=True,align=PP_ALIGN.CENTER)
    txt(s,x+Inches(0.07),Inches(3.5),Inches(1.52),Inches(1.3),desc,size=10,color=RGBColor(0x33,0x33,0x33),align=PP_ALIGN.CENTER)
    x+=Inches(1.78)
txt(s,Inches(0.5),Inches(5.4),Inches(12.3),Inches(1.0),"The same transformation runs on real customer logs; only the source of the raw events changes. Every step in this pipeline runs today.",size=14,color=GREEN,bold=True)
fin(s)

# ---- DT-2 ZONES ----
s=slide(); header(s,"Digital Twin in Detail: The Five Behavioral Zones","23 numeric features across four action zones, plus identity as context")
zones=[("identity","role, department, clearance, tenure","context: what normal looks like here",SLATE),
       ("access_pattern","auth volume, fail rate, off-hours, unique destinations","credential abuse, lateral movement",L1),
       ("data_behavior","file count, restricted and confidential ratios, bytes","collection and exfiltration staging",L2),
       ("network_footprint","bytes out, unique destinations, external ratio, DNS","C2 beaconing, tunneling, exfiltration",L3),
       ("risk_posture","endpoint events, suspicious ratio, unique processes","living-off-the-land tooling, malware",L4)]
y=Inches(1.5)
for nm,feats,catch,col in zones:
    rect(s,Inches(0.6),y,Inches(12.1),Inches(0.95),LIGHT); rect(s,Inches(0.6),y,Inches(2.6),Inches(0.95),col)
    txt(s,Inches(0.7),y,Inches(2.4),Inches(0.95),nm,size=14,color=WHITE,bold=True,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,Inches(3.4),y+Inches(0.06),Inches(5.2),Inches(0.85),feats,size=11.5,color=RGBColor(0x33,0x33,0x33),anchor=MSO_ANCHOR.MIDDLE)
    txt(s,Inches(8.8),y+Inches(0.06),Inches(3.8),Inches(0.85),catch,size=11.5,color=NAVY,anchor=MSO_ANCHOR.MIDDLE)
    y+=Inches(1.02)
fin(s)

# ---- DT-3 MEANING ----
s=slide(); header(s,"Digital Twin in Detail: Meaning Over Magnitude","Why prose and embedding, not raw counts")
bullets(s,Inches(0.6),Inches(1.5),Inches(7.2),Inches(5),[
 (0,"Each zone's attributes fuse into one sentence, then embed; the meaning of every feature is carried into one vector.",NAVY,True),
 (0,"Meaning lives in the relationships among attributes, not any single number.",),
 (0,"Same 250 MB of file activity: 180 public files (benign) versus 6 restricted files with high write (exfil staging) embed in completely different regions.",),
 (0,"Individually benign access signals (off-hours, broad, error-free) only read as lateral movement once fused.",),
],size=15,gap=12)
rect(s,Inches(8.2),Inches(1.6),Inches(4.6),Inches(4.5),LIGHT)
txt(s,Inches(8.4),Inches(1.8),Inches(4.2),Inches(0.5),"Direction, not just magnitude",size=15,color=BLUE,bold=True)
txt(s,Inches(8.4),Inches(2.5),Inches(4.2),Inches(3.4),"We do not say '87 percent anomalous'. We say the entity is drifting toward a named MITRE threat concept, for example living-off-the-land [T1218, T1053] or c2_beacon [T1071, T1573].",size=14,color=RGBColor(0x33,0x33,0x33))
fin(s)

# ---- DT-4 DETECTION ----
s=slide(); header(s,"Digital Twin in Detail: Detection Over Time","Drift, CUSUM, direction, composite, novelty")
bullets(s,Inches(0.6),Inches(1.5),Inches(12),Inches(5),[
 (0,"Drift: how far the twin moved week to week (cosine distance between weekly composites).",),
 (0,"CUSUM: accumulates small, sub-threshold drift above a noise floor; a slow APT that never trips a per-week line still fires.",NAVY,True),
 (0,"Direction: the drift vector is projected onto 14 threat and 2 benign concepts; the best match names the MITRE technique.",),
 (0,"Composite: a five-phase score (signal strength, breadth, sustained, context divergence, novelty) yields one ranked verdict.",),
 (0,"Novelty: a novel command-and-control IP recurring week after week is the beacon signature, and it is what catches the stealthiest attacker.",GREEN,True),
],size=15,gap=12)
fin(s)

# ---- DT-5 SERIALIZATION ----
s=slide(); header(s,"Digital Twin in Detail: From Raw Logs to a Vector","Each zone fuses its attributes into one sentence, then embeds into 1536-d")
txt(s,Inches(0.6),Inches(1.35),Inches(12),Inches(0.4),"Real zone serializations (verbatim, entity USR-156):",size=14,color=NAVY,bold=True)
ser=("identity:  User USR-156: role=Analyst, dept=Marketing, clearance=internal\n"
     "access:    auth_events=16, fail_rate=0.0625, off_hours=0.1250, unique_dests=5\n"
     "data:      file_accesses=10, restricted=0.0, confidential=0.2, bytes=124,933,027\n"
     "network:   bytes_out=125,407,529, unique_dsts=22, external_ratio=0.1042\n"
     "risk:      endpoint_events=8, suspicious=0.0, max_risk=30.0, processes=7")
rect(s,Inches(0.6),Inches(1.95),Inches(12.1),Inches(2.7),LIGHT)
txt(s,Inches(0.8),Inches(2.15),Inches(11.7),Inches(2.4),ser,size=12,color=RGBColor(0x22,0x22,0x22),font="Consolas")
txt(s,Inches(0.6),Inches(4.9),Inches(12),Inches(0.9),"Each line becomes one 1536-d vector. Five zones, five vectors, composed into one twin. The meaning of every feature in a zone is carried into that zone's embedding.",size=14,color=GREEN,bold=True)
fin(s)

# ---- DT-6 COMPOSITE WORKED ----
s=slide(); header(s,"Digital Twin in Detail: The Composite Score, Worked","Group-relative z-scores, fused by one formula")
txt(s,Inches(0.6),Inches(1.35),Inches(12),Inches(0.5),"Each user's features become z-scores versus the normals in their OWN role group (a higher bar than a global threshold).",size=13,color=NAVY,bold=True)
rect(s,Inches(0.6),Inches(2.05),Inches(12.1),Inches(1.0),LIGHT)
txt(s,Inches(0.8),Inches(2.2),Inches(11.7),Inches(0.75),"composite = signal_strength + 0.5 x breadth + 0.3 x sustained + 0.5 x ctx_spread + 0.3 x ctx_max + novelty",size=13.5,color=NAVY,font="Consolas",bold=True)
rect(s,Inches(0.6),Inches(3.35),Inches(12.1),Inches(1.5),RGBColor(0xf2,0xf6,0xfc))
txt(s,Inches(0.8),Inches(3.55),Inches(11.7),Inches(1.1),"Volt:   6.40 + 5.50 + 1.14 + 0.23 + 0.43 + 0.00  =  13.70   (rank #24, stealthiest)\nSalt:  29.89 + 9.00 + 2.87 + 4.32 + 2.39 + 2.80  =  51.27   (rank #1)",size=13,color=RGBColor(0x22,0x22,0x22),font="Consolas")
txt(s,Inches(0.6),Inches(5.1),Inches(12),Inches(0.8),"Every composite re-derives live from the formula, in front of the room, from the results file.",size=14,color=GREEN,bold=True)
fin(s)

# ---- DT-7 FOUR ATTACKERS DATA ----
s=slide(); header(s,"Digital Twin in Detail: The Four Attackers, by the Numbers","Why each scored where it did")
dtable(s,Inches(0.5),Inches(1.5),Inches(12.3),Inches(2.5),
 ["Attacker","signal","breadth","sustained","ctx spr/max","novelty","composite (rank)"],
 [["Salt","29.89","18","9.58","8.65 / 7.94","2.8","51.27 (#1)"],
  ["Insider","30.02","19","8.37","4.34 / 6.80","0.0","46.24 (#2)"],
  ["Slow APT","4.47","1","1.57","1.48 / 0.85","13.0","19.44 (#7)"],
  ["Volt","6.40","11","3.80","0.47 / 1.43","0.0","13.70 (#24)"]],fs=12)
txt(s,Inches(0.6),Inches(4.6),Inches(12.1),Inches(1.7),
 "Read it: Slow APT is invisible on every magnitude column (signal 4.47, breadth 1) and is caught by novelty ALONE "
 "(13.0, a novel C2 IP every week for 60 weeks). Volt has breadth (11) of small signals and no novelty: the "
 "living-off-the-land signature. No single number flags either; the twin separates them by the full picture.",
 size=15,color=NAVY)
fin(s)

# ---- 10 LAYER 2 PROOF ----
s=slide(); header(s,"Layer 2 - The Discovery Net (the twin)","Unsupervised, by rank: 4 of 4 surfaced at an 8.1% operating point")
rect(s,Inches(0.6),Inches(1.5),Inches(7.3),Inches(4.7),LIGHT)
hdr=["Attacker","Rank","Composite"]
rowsP=[("Salt Typhoon","#1","51.27",L2),("Insider","#2","46.24",L2),("Slow APT","#7","19.44",AMBER),("Volt Typhoon","#24","13.70",RED)]
txt(s,Inches(0.8),Inches(1.65),Inches(7,),Inches(0.4),"Four-attacker blind test (risk-ranked of 250)",size=14,color=NAVY,bold=True)
yy=Inches(2.2)
for a,rk,cp,col in rowsP:
    txt(s,Inches(0.9),yy,Inches(3.5),Inches(0.5),a,size=15,color=col,bold=True)
    txt(s,Inches(4.6),yy,Inches(1.2),Inches(0.5),rk,size=15,color=NAVY,bold=True)
    txt(s,Inches(6.0),yy,Inches(1.6),Inches(0.5),cp,size=15,color=SLATE)
    yy+=Inches(0.62)
rect(s,Inches(8.2),Inches(1.5),Inches(4.6),Inches(4.7),NAVY)
txt(s,Inches(8.4),Inches(1.9),Inches(4.2),Inches(1.0),"4 of 4",size=44,color=WHITE,bold=True,align=PP_ALIGN.CENTER)
txt(s,Inches(8.4),Inches(3.0),Inches(4.2),Inches(0.8),"surfaced by rank at an 8.1% operating point (value is coverage + MITRE direction)",size=14,color=RGBColor(0xcf,0xdd,0xf2),align=PP_ALIGN.CENTER)
txt(s,Inches(8.4),Inches(3.9),Inches(4.2),Inches(1.5),"Classical methods (Isolation Forest, One-Class SVM, LOF): 0 of 4.  z-score: 1 of 4.",size=15,color=RGBColor(0xff,0xe0,0xe0),align=PP_ALIGN.CENTER,bold=True)
txt(s,Inches(0.6),Inches(6.4),Inches(12),Inches(0.5),"Controlled, synthetic-data proof of concept, indicative rather than a guarantee of field performance.",size=12,color=SLATE,italic=True)
fin(s)

# ---- 11 LAYER 2 DATA POINTS ----
s=slide(); header(s,"Layer 2 - The Two Data Points That Carry the Room")
rect(s,Inches(0.6),Inches(1.6),Inches(5.9),Inches(4.4),LIGHT); rect(s,Inches(0.6),Inches(1.6),Inches(5.9),Inches(0.8),AMBER)
txt(s,Inches(0.8),Inches(1.7),Inches(5.5),Inches(0.6),"Slow APT: invisible to thresholds",size=17,color=WHITE,bold=True,anchor=MSO_ANCHOR.MIDDLE)
txt(s,Inches(0.8),Inches(2.6),Inches(5.5),Inches(3.2),
 "Signal strength 4.47, breadth 1: invisible to every magnitude method. Caught by novelty alone, a "
 "novel command-and-control IP recurring every week for 60 weeks. That one signal lifts it to rank #7.",size=15,color=RGBColor(0x33,0x33,0x33))
rect(s,Inches(6.8),Inches(1.6),Inches(5.9),Inches(4.4),LIGHT); rect(s,Inches(6.8),Inches(1.6),Inches(5.9),Inches(0.8),RED)
txt(s,Inches(7.0),Inches(1.7),Inches(5.5),Inches(0.6),"Volt: living off the land",size=17,color=WHITE,bold=True,anchor=MSO_ANCHOR.MIDDLE)
txt(s,Inches(7.0),Inches(2.6),Inches(5.5),Inches(3.2),
 "Many small elevations, no single dramatic one, no novelty. Caught by breadth across fronts, the "
 "signature of legitimate admin tooling used maliciously. The stealthiest scorer, and still surfaced.",size=15,color=RGBColor(0x33,0x33,0x33))
txt(s,Inches(0.6),Inches(6.3),Inches(12),Inches(0.5),"We present a risk-ranked list, not a threshold. The claim is recall at the top of the ranking.",size=14,color=NAVY,bold=True)
fin(s)

# ---- 12 LAYER 1 ----
s=slide(); header(s,"Layer 1 - Preemptive Network Assurance","The prevention half: prove the edge is closed before traffic flows")
bullets(s,Inches(0.6),Inches(1.5),Inches(12),Inches(5),[
 (0,"Incumbents sample the configuration state space; they cannot prove segmentation is correct.",),
 (0,"99% of firewall breaches are misconfiguration; the state space is too large to sample.",RED,True),
 (0,"Our wedge: an exhaustive formal model of the control estate, queried against MITRE-enriched attack models.",NAVY,True),
 (0,"Uses residual capacity in the estate to block attacks from live intelligence; drift caught within the hour.",),
 (0,"Agentless, configuration-only; first actionable findings in under 15 days.",GREEN,True),
],size=16,gap=12)
fin(s)

# ---- 13 LAYER 1 PROOF ----
s=slide(); header(s,"Layer 1 - Verified Preemptive Findings","Real findings human review and existing tools had missed")
bullets(s,Inches(0.6),Inches(1.5),Inches(12),Inches(5),[
 (0,"Credential and IP-theft path closed: a shadowed, incomplete rule set allowed an unsanctioned app from internal subnets, against CISO intent.",),
 (0,"Lateral-movement path closed: three higher-priority rules formed a group shadow letting legacy SMBv1.0 pass without inspection.",),
 (0,"External attack surface formally assured for a national cyber-defense agency, beyond what EASM can simulate.",),
 (0,"Regulatory compliance verified against a live national malicious-IP and domain threat feed.",),
 (0,"Infosec compliance gaps diagnosed by formal verification of firewall configurations.",),
],size=15,gap=11)
fin(s)

# ---- 14 LAYER 3 ----
s=slide(); header(s,"Layer 3 - Application, API & Data Runtime + Agentic AI","The exfiltration altitude, and the agentic-AI blind spot")
bullets(s,Inches(0.6),Inches(1.5),Inches(12),Inches(5),[
 (0,"Where Salt actually collects and exfiltrates: APIs, microservices, data, and AI-agent pathways.",NAVY,True),
 (0,"WAFs miss API nuance and business-logic abuse; 95% of API attacks come from authenticated sources.",),
 (0,"Agentic AI is the new blind spot: AI agents talk over APIs; ~99% of AI-related vulnerabilities are API-related.",RED,True),
 (0,"Our wedge: unified discover-comply-protect on a behavioral/identity core, with runtime-native mitigation.",GREEN,True),
 (0,"Agentic-AI governance: secure agent access, block AI-driven bots, real-time PII and credential redaction.",),
 (0,"Maturity: integrated capability, presented as architecture; not a live demo in this session.",AMBER,True),
],size=16,gap=12)
fin(s)

# ---- 15 LAYER 4 ----
s=slide(); header(s,"Layer 4 - Code & Supply-Chain Assurance","Close the lifecycle at the source (presented as in development)")
bullets(s,Inches(0.6),Inches(1.5),Inches(12),Inches(5),[
 (0,"The exploitable code or dependency that grants initial access in the first place.",NAVY,True),
 (0,"Incumbents match known CVEs; they cannot find novel zero-days and are reactive to disclosure.",),
 (0,"Supply-chain injection (SolarWinds, XZ Utils) slips past dependency scanners checking known-bad versions.",),
 (0,"Our wedge: AI-driven novel zero-day discovery pre-deploy, gated through multi-stage CI/CD assurance.",GREEN,True),
 (0,"Honest maturity: this layer is in development. We present the approach, not a running product.",AMBER,True),
],size=16,gap=12)
fin(s)

# ---- 16 INTEGRATION ----
s=slide(); header(s,"Integration: One Typhoon Across All Four Layers","The hand-off is the moat no point tool can claim")
steps=[("Layer 4 - Code","Removes the exploitable code that grants initial access",L4),
       ("Layer 1 - Edge","Proves the misconfigured edge path is closed",L1),
       ("Layer 2 - Behavior","Catches the valid-account behavior once inside",L2),
       ("Layer 3 - API/Data","Catches the API and data exfiltration, and the agents",L3)]
x=Inches(0.6)
for nm,desc,col in steps:
    rect(s,x,Inches(2.0),Inches(2.85),Inches(2.8),col)
    txt(s,x+Inches(0.15),Inches(2.15),Inches(2.55),Inches(0.7),nm,size=15,color=WHITE,bold=True)
    txt(s,x+Inches(0.15),Inches(2.9),Inches(2.55),Inches(1.8),desc,size=13,color=WHITE)
    x+=Inches(3.05)
txt(s,Inches(0.6),Inches(5.3),Inches(12),Inches(1.0),
 "Each layer catches what the others structurally cannot. Covering all four is completeness, not redundancy. "
 "Anything less leaves a seam the adversary already knows how to use.",size=16,color=NAVY,bold=True)
fin(s)

# ---- 17 FEDERAL FIT ----
s=slide(); header(s,"Federal Fit","Fieldable on what ARCYBER already collects")
bullets(s,Inches(0.6),Inches(1.5),Inches(12),Inches(5),[
 (0,"Runs on the logs you already collect: sign-in, network, endpoint, file, identity.",NAVY,True),
 (0,"Containerized; deployable across IL5, IL6, and JWICS enclaves.",),
 (0,"No PII held in the behavioral vector; offline embedding today, local semantic model as a Phase-1 build.",),
 (0,"Layer 1 is agentless and configuration-only; first findings in under 15 days.",GREEN,True),
 (0,"One platform, one vendor, one accountable integrator: not four separate procurements.",BLUE,True),
],size=16,gap=12)
fin(s)

# ---- 18 HONESTY / VALIDATION ----
s=slide(); header(s,"Validation Boundaries (stated up front)","Candor is our advantage; this room punishes overclaim")
bullets(s,Inches(0.6),Inches(1.5),Inches(12),Inches(5),[
 (0,"Layer 2: the 4-of-4 result is a synthetic-data proof of concept on 250 entities, not yet field-validated.",),
 (0,"Layer 1: customer findings are representative engagements; they demonstrate the method, not a guarantee.",),
 (0,"Layer 4: in development; the approach, not a running product.",),
 (0,"Identity fusion, tamper-evident forensics, automated response: roadmap, not running today.",),
 (0,"Everything maps to MITRE ATT&CK and ATLAS, but is not yet third-party audited.",),
],size=16,gap=12)
fin(s)

# ---- 19 THE ASK ----
s=slide(); rect(s,0,0,SW,SH,NAVY); rect(s,0,Inches(2.3),SW,Pt(3),GOLD)
txt(s,Inches(0.9),Inches(1.3),Inches(11.5),Inches(0.5),"THE ASK",size=18,color=GOLD,bold=True)
txt(s,Inches(0.9),Inches(2.0),Inches(11.5),Inches(2.0),
 "A bounded Phase-1 pilot on real Army telemetry: advisory, human-in-the-loop, measured side by side "
 "against current tooling.",size=26,color=WHITE,bold=True)
bullets(s,Inches(0.9),Inches(4.2),Inches(11.5),Inches(2.5),[
 (0,"Layer 1: agentless config-only assessment of one enclave's edge estate; first findings in under 15 days.",WHITE),
 (0,"Layer 2: blind behavioral evaluation on a bounded Army enclave dataset at a fixed false-positive point.",WHITE),
 (0,"Environment: an IL5 or IL6 enclave, extensible to JWICS. Path: OTA prototype maturing toward a program of record.",WHITE),
],size=15,gap=10)
fin(s)

# ---- 20 WHY IT MATTERS ----
s=slide(); header(s,"Why This Matters","It satisfies all three seats at once")
cards=[("Esquibel","Measured detection fidelity: 4 of 4 at 8.1% where classical methods catch 0 of 4, on the logs she already has.",L2),
       ("DiGrezio","One fieldable platform from one vendor, not four procurements. Agentless start, clear acquisition path.",L1),
       ("Herrmann","A complete answer to the national valid-account / critical-infrastructure gap she can carry upward.",L3)]
x=Inches(0.5)
for nm,desc,col in cards:
    rect(s,x,Inches(1.6),Inches(4.0),Inches(4.4),LIGHT); rect(s,x,Inches(1.6),Inches(4.0),Inches(0.8),col)
    txt(s,x+Inches(0.2),Inches(1.7),Inches(3.6),Inches(0.6),nm,size=18,color=WHITE,bold=True,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,x+Inches(0.2),Inches(2.7),Inches(3.6),Inches(3.0),desc,size=14,color=RGBColor(0x33,0x33,0x33))
    x+=Inches(4.27)
txt(s,Inches(0.5),Inches(6.3),Inches(12),Inches(0.5),"Mapped 1:1 to how the adversary operates. The whole lifecycle, one vocabulary, under human control.",size=14,color=NAVY,bold=True)
fin(s)

# ---- 21 NUMBERS CARD ----
s=slide(); header(s,"Numbers to Keep Straight","The reference card")
bullets(s,Inches(0.6),Inches(1.5),Inches(12),Inches(5),[
 (0,"Dataset: 250 entities, 485 days (70 weekly snapshots), ~14M events, 5 log sources.",),
 (0,"Production (threat-profile detector): 4 of 4 at 0 false positives, by named fingerprint. Discovery (twin): 4 of 4 by rank at 8.1% FP; classical baselines 0 of 4, z-score 1 of 4.",NAVY,True),
 (0,"Ranks: Salt #1 (51.27), Insider #2 (46.24), Slow APT #7 (19.44), Volt #24 (13.70).",),
 (0,"Separation is by rank, not a clean gap: top normal user outscores Slow APT and Volt; present a ranked list.",AMBER,True),
 (0,"Layer 1: 99% of firewall breaches are misconfiguration; first findings in under 15 days.",),
 (0,"Layer 3: ~99% of AI-related vulnerabilities are API-related; 95% of API attacks are authenticated.",),
],size=15,gap=11)
fin(s)

# ---- 22 CLOSING ----
s=slide(); rect(s,0,0,SW,SH,NAVY); rect(s,0,Inches(2.4),SW,Pt(3),GOLD)
txt(s,Inches(0.9),Inches(1.6),Inches(11.5),Inches(2.2),
 "One platform, four layers, one MITRE vocabulary, covering the full attack lifecycle. The behavioral "
 "core is ours and runs live, on the logs you already collect, catching four nation-state campaigns your "
 "current tools score as normal.",size=24,color=WHITE,bold=True)
txt(s,Inches(0.9),Inches(4.4),Inches(11.5),Inches(1.0),
 "The other layers close the rest of the lifecycle. Layer 4 is the honest roadmap. The right next step is "
 "a bounded pilot on real Army telemetry.",size=18,color=RGBColor(0xcf,0xdd,0xf2))
txt(s,Inches(0.9),Inches(6.2),Inches(11.5),Inches(0.5),"Thank you.  |  V-Intelligence",size=16,color=GOLD,bold=True)
fin(s)

out="WP DLA/Presentation/Pentagon_Working_Session_Deck.pptx"
prs.save(out)
print("saved:",out,"| slides:",len(prs.slides._sldIdLst))
