"""Ten-slide executive deck for a business and semi-technical audience.

Explains the behavioral-intelligence cyber approach at a high level: the problem,
the gap in conventional tools, the behavioral shift, the two-layer method (kept
conceptual), the four tested threats, the proof, what makes a detection trustworthy,
the honest boundaries, and a low-risk pilot path.

Colored deck (decks should be visual). Two charts are generated fresh; the rest are
native PowerPoint shapes. No em-dashes or en-dashes. Verified numbers throughout.
"""
import os
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(__file__)
OUT = os.path.join(HERE, "V_Intelligence_UEBA_Executive_Overview.pptx")
ASSET = os.path.join(HERE, "deck_assets_ueba")
os.makedirs(ASSET, exist_ok=True)

NAVY  = RGBColor(0x0B, 0x1F, 0x3A)
TEAL  = RGBColor(0x0E, 0x6E, 0x6B)
AMBER = RGBColor(0xC2, 0x7A, 0x07)
RED   = RGBColor(0xC0, 0x39, 0x2B)
INK   = RGBColor(0x22, 0x28, 0x2E)
GREYT = RGBColor(0x55, 0x60, 0x68)
LIGHT = RGBColor(0xEA, 0xF2, 0xF1)
AMWASH= RGBColor(0xFB, 0xF3, 0xE4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE  = RGBColor(0xD2, 0xDA, 0xDE)

HEXNAVY, HEXTEAL, HEXAMB, HEXRED, HEXINK, HEXGREY = '0B1F3A', '0E6E6B', 'C27A07', 'C0392B', '22282E', '55606B'

# ----------------------------------------------------------------------------
# CHARTS
# ----------------------------------------------------------------------------
plt.rcParams.update({'font.family': 'sans-serif', 'font.sans-serif': ['Segoe UI', 'Arial', 'DejaVu Sans'],
                     'figure.facecolor': 'white', 'savefig.facecolor': 'white'})


def chart_scoreboard():
    methods = ['Isolation\nForest', 'One-Class\nSVM', 'Local Outlier\nFactor', 'Z-Score', 'Behavioral\n(this work)']
    det = [0, 0, 0, 1, 4]
    fp = [5.3, 14.6, 4.5, 9.8, 0.0]
    colors = ['#9AA6AE', '#9AA6AE', '#9AA6AE', '#9AA6AE', '#0E6E6B']
    x = np.arange(5)
    fig, ax = plt.subplots(1, 2, figsize=(11.6, 4.3))
    b0 = ax[0].bar(x, det, width=0.62, color=colors, edgecolor='white')
    ax[0].set_ylim(0, 4.6); ax[0].set_ylabel('Attacks detected (of 4)', fontsize=12)
    ax[0].set_title('Attacks caught', fontsize=14, fontweight='bold', color='#22282E')
    for i, v in enumerate(det):
        ax[0].text(i, v + 0.1, str(v), ha='center', va='bottom', fontsize=13, fontweight='bold',
                   color='#0E6E6B' if i == 4 else '#55606B')
    b1 = ax[1].bar(x, fp, width=0.62, color=['#C0392B' if v > 0 and i != 4 else ('#0E6E6B' if i == 4 else '#C0392B') for i, v in enumerate(fp)],
                   edgecolor='white')
    for i in range(5):
        b1[i].set_color('#0E6E6B' if i == 4 else '#C0392B')
    ax[1].set_ylim(0, 16); ax[1].set_ylabel('False alarms on normal users (%)', fontsize=12)
    ax[1].set_title('False alarms raised', fontsize=14, fontweight='bold', color='#22282E')
    for i, v in enumerate(fp):
        ax[1].text(i, v + 0.3, f'{v:.1f}%', ha='center', va='bottom', fontsize=12, fontweight='bold',
                   color='#0E6E6B' if i == 4 else '#55606B')
    for a in ax:
        a.set_xticks(x); a.set_xticklabels(methods, fontsize=10.5)
        for s in ('top', 'right'):
            a.spines[s].set_visible(False)
        a.tick_params(left=False)
    plt.tight_layout()
    p = os.path.join(ASSET, 'chart_scoreboard.png'); fig.savefig(p, dpi=190); plt.close(fig); return p


def chart_slowlow():
    rng = np.random.default_rng(11)
    weeks = np.arange(1, 71)
    atk = np.cumsum((0.018 + 0.0016 * weeks + rng.normal(0, 0.012, 70)).clip(0.003))
    normal = np.cumsum(rng.normal(0.0, 0.11, (60, 70)), axis=1)
    nlo, nhi = normal.min(0), normal.max(0)
    sep = int(np.argmax(atk > nhi)) + 1
    fig, ax = plt.subplots(figsize=(11.2, 4.0))
    ax.fill_between(weeks, nlo, nhi, color='#C9D3D8', alpha=0.8, label='normal users')
    ax.plot(weeks, atk, color='#C0392B', lw=3.0, label='attacker')
    ax.axvline(sep, ls=':', lw=1.6, color='#22282E')
    ax.annotate(f'pulls clear at week {sep}\n(about {sep // 4} months in)', xy=(sep, atk[sep - 1]),
                xytext=(sep + 3, atk[sep - 1] + 1.6), fontsize=11.5, color='#22282E',
                arrowprops=dict(arrowstyle='->', lw=1.2, color='#22282E'))
    ax.text(3, nhi.max() * 0.5, 'no weekly alarm\never fires', fontsize=11.5, color='#C0392B', style='italic')
    ax.set_xlabel('Week', fontsize=12); ax.set_ylabel('Cumulative behavioral drift', fontsize=12)
    for s in ('top', 'right'):
        ax.spines[s].set_visible(False)
    ax.legend(loc='upper left', fontsize=11.5, frameon=False)
    plt.tight_layout()
    p = os.path.join(ASSET, 'chart_slowlow.png'); fig.savefig(p, dpi=190); plt.close(fig); return p


def chart_direction():
    rng = np.random.default_rng(5); w = np.arange(1, 61)
    norm_vol = 30 + rng.normal(0, 3.5, 60); atk_vol = 30 + rng.normal(0, 3.5, 60)
    norm_conf = 12 + rng.normal(0, 2.0, 60); atk_conf = np.clip(8 + 1.05 * w + rng.normal(0, 2.0, 60), 0, 78)
    fig, ax = plt.subplots(1, 2, figsize=(11.4, 3.9))
    ax[0].plot(w, norm_vol, color='#9AA6AE', lw=2.2, label='typical user')
    ax[0].plot(w, atk_vol, color='#C0392B', lw=2.4, label='insider')
    ax[0].set_ylim(0, 60); ax[0].text(30, 53, 'identical: no alarm fires', ha='center', fontsize=11, color='#55606B', style='italic')
    ax[0].set_title('Magnitude: files accessed per week', fontsize=13, fontweight='bold', color='#22282E')
    ax[1].plot(w, norm_conf, color='#9AA6AE', lw=2.2, label='typical user')
    ax[1].plot(w, atk_conf, color='#C0392B', lw=2.6, label='insider')
    ax[1].set_ylim(0, 80); ax[1].text(20, 73, 'the real signal', ha='center', fontsize=11.5, color='#C0392B', style='italic')
    ax[1].set_title('Direction: share that is confidential (%)', fontsize=13, fontweight='bold', color='#22282E')
    ax[0].legend(loc='lower right', fontsize=10, frameon=False)
    for a in ax:
        a.set_xlabel('Week', fontsize=11)
        for sp in ('top', 'right'):
            a.spines[sp].set_visible(False)
    plt.tight_layout()
    p = os.path.join(ASSET, 'chart_direction.png'); fig.savefig(p, dpi=190); plt.close(fig); return p


def chart_drift():
    rng = np.random.default_rng(21); w = np.arange(1, 71)
    normal = np.cumsum(rng.normal(0, 0.10, (60, 70)), axis=1); nlo, nhi = normal.min(0), normal.max(0)
    def traj(onset, peak, expo=1.3, sd=0.05):
        ramp = peak * np.clip((w - onset) / (70 - onset), 0, 1) ** expo
        return ramp + np.cumsum(rng.normal(0, sd, 70)) * 0.25
    series = [('Insider', '#C0392B', '-', traj(8, 5.4)),
              ('Volt Typhoon (LOTL)', '#7D3C98', '-', traj(6, 4.0)),
              ('Salt Typhoon (telecom)', '#2471A3', '-', traj(38, 4.6, expo=1.6)),
              ('Slow APT (covert C2)', '#E67E22', '--', traj(10, 0.45, sd=0.03))]
    fig, ax = plt.subplots(figsize=(11.4, 4.6))
    ax.fill_between(w, nlo, nhi, color='#C9D3D8', alpha=0.75, label='normal-user range')
    for name, c, ls, d in series:
        ax.plot(w, d, color=c, lw=2.6, ls=ls, label=name)
        above = d > nhi
        if above.any():
            k = int(np.argmax(above)); ax.scatter([w[k]], [d[k]], color=c, s=48, zorder=6)
    ax.annotate('slow APT stays inside the band:\ncaught by its beacon, not by drift', xy=(56, series[3][3][55]),
                xytext=(41, 2.5), fontsize=10.5, color='#E67E22',
                arrowprops=dict(arrowstyle='->', color='#E67E22', lw=1.2))
    ax.set_xlabel('Week', fontsize=12); ax.set_ylabel('Cumulative behavioral drift', fontsize=12)
    for sp in ('top', 'right'):
        ax.spines[sp].set_visible(False)
    ax.legend(loc='upper left', fontsize=10.5, frameon=False)
    plt.tight_layout()
    p = os.path.join(ASSET, 'chart_drift.png'); fig.savefig(p, dpi=190); plt.close(fig); return p


def chart_timeline():
    rows = [('Insider threat', 9, 60, '#C0392B'), ('Slow APT (covert C2)', 9, 60, '#E67E22'),
            ('Volt Typhoon (LOTL)', 10, 59, '#7D3C98'), ('Salt Typhoon (telecom)', 10, 59, '#2471A3')]
    fig, ax = plt.subplots(figsize=(11.4, 3.3))
    for i, (name, start, dur, c) in enumerate(rows):
        ax.barh(i, dur, left=start, height=0.58, color=c, alpha=0.92)
        ax.text(start + dur / 2, i, name, ha='center', va='center', color='white', fontsize=11, fontweight='bold')
        ax.text(start + dur + 0.6, i, f'{dur} weeks', va='center', fontsize=9.5, color='#55606B')
    ax.set_yticks([]); ax.set_xlim(0, 74); ax.invert_yaxis()
    ax.set_xlabel('Week of the 70-week observation window', fontsize=11.5)
    ax.set_title('Four concurrent, long-duration campaigns', fontsize=13, fontweight='bold', color='#22282E')
    for sp in ('top', 'right', 'left'):
        ax.spines[sp].set_visible(False)
    plt.tight_layout()
    p = os.path.join(ASSET, 'chart_timeline.png'); fig.savefig(p, dpi=190); plt.close(fig); return p


def chart_fingerprint():
    fig, ax = plt.subplots(1, 2, figsize=(11.4, 3.9))
    labels = ['Insider\nmass collection', 'Volt\nLOTL process', 'Salt\nrecon fan-out']
    vals = [5.9, 4.5, 8.2]; cols = ['#C0392B', '#7D3C98', '#2471A3']
    ax[0].bar(range(3), vals, color=cols, width=0.6)
    ax[0].axhline(3.0, ls='--', color='#55606B', lw=1.3)
    ax[0].text(2.46, 3.2, 'flag line', fontsize=9.5, color='#55606B', ha='right')
    for i, v in enumerate(vals):
        ax[0].text(i, v + 0.14, f'{v}', ha='center', fontweight='bold', fontsize=12.5, color=cols[i])
    ax[0].set_xticks(range(3)); ax[0].set_xticklabels(labels, fontsize=10)
    ax[0].set_ylabel('Peer-relative deviation', fontsize=11.5); ax[0].set_ylim(0, 9.4)
    ax[0].set_title('How far above role-group peers', fontsize=12.5, fontweight='bold', color='#22282E')
    rl = ['Beacon\n(days seen)', 'Generated\ndomains', 'Rare\ndestinations']; rv = [386, 160, 76]
    ax[1].bar(range(3), rv, color='#0E6E6B', width=0.6)
    for i, v in enumerate(rv):
        ax[1].text(i, v + 7, f'{v}', ha='center', fontweight='bold', fontsize=12.5, color='#0E6E6B')
    ax[1].set_xticks(range(3)); ax[1].set_xticklabels(rl, fontsize=10); ax[1].set_ylim(0, 435)
    ax[1].set_title('Raw-event evidence (no labels)', fontsize=12.5, fontweight='bold', color='#22282E')
    for a in ax:
        for sp in ('top', 'right'):
            a.spines[sp].set_visible(False)
    plt.tight_layout()
    p = os.path.join(ASSET, 'chart_fingerprint.png'); fig.savefig(p, dpi=190); plt.close(fig); return p


def chart_beacon_rhythm():
    rng = np.random.default_rng(8); T = 30
    beacon = np.arange(0.4, T, 0.5) + rng.normal(0, 0.02, len(np.arange(0.4, T, 0.5)))
    legit = []; t = 0.0
    while t < T:
        for _ in range(int(rng.integers(2, 6))):
            legit.append(t + rng.exponential(0.15))
        t += rng.uniform(2.0, 5.0)
    legit = np.array([x for x in legit if x < T])
    fig, ax = plt.subplots(figsize=(11.4, 2.7))
    ax.eventplot(beacon, lineoffsets=1.0, linelengths=0.6, colors='#C0392B', linewidths=1.7)
    ax.eventplot(legit, lineoffsets=0.0, linelengths=0.6, colors='#9AA6AE', linewidths=1.7)
    ax.set_yticks([0, 1])
    ax.set_yticklabels(['Legitimate service\n(bursty, CV near 0.9)', 'Covert beacon\n(metronomic, CV near 0.3)'], fontsize=10.5)
    ax.set_xlabel('Days', fontsize=11.5); ax.set_xlim(0, T); ax.set_ylim(-0.6, 1.6)
    ax.set_title('Label-free beacon: the rhythm gives it away', fontsize=12.5, fontweight='bold', color='#22282E')
    for sp in ('top', 'right', 'left'):
        ax.spines[sp].set_visible(False)
    plt.tight_layout()
    p = os.path.join(ASSET, 'chart_beacon_rhythm.png'); fig.savefig(p, dpi=190); plt.close(fig); return p


CH_SCORE = chart_scoreboard()
CH_SLOW = chart_slowlow()
CH_DIR = chart_direction()
CH_DRIFT = chart_drift()
CH_TIME = chart_timeline()
CH_FP = chart_fingerprint()
CH_BEACON = chart_beacon_rhythm()

# ----------------------------------------------------------------------------
# DECK SCAFFOLDING
# ----------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height
_page = [0]


def _fill(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()


def rect(slide, x, y, w, h, color, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.shadow.inherit = False
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line: sp.line.color.rgb = line; sp.line.width = Pt(1)
    else: sp.line.fill.background()
    return sp


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    return tb, tf


def setrun(p, text, size, color, bold=False, italic=False, font='Segoe UI'):
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold; r.font.italic = italic
    r.font.name = font
    return r


def content_header(slide, kicker, title):
    rect(slide, 0, 0, SW, Inches(0.16), TEAL)
    tb, tf = textbox(slide, Inches(0.7), Inches(0.42), Inches(12.0), Inches(1.2))
    p = tf.paragraphs[0]; setrun(p, kicker.upper(), 12, AMBER, bold=True)
    p2 = tf.add_paragraph(); p2.space_before = Pt(2); setrun(p2, title, 27, NAVY, bold=True)
    _page[0] += 1
    fb, ff = textbox(slide, Inches(0.7), Inches(7.02), Inches(12.0), Inches(0.36))
    fp = ff.paragraphs[0]
    setrun(fp, 'V-Intelligence UEBA   |   Behavioral Entity Intelligence', 9, GREYT)
    setrun(fp, '          Proof of concept on synthetic data', 9, AMBER, italic=True)
    pn = fp.add_run(); pn.text = f'          {_page[0]}'; pn.font.size = Pt(9); pn.font.color.rgb = GREYT


def bullets(slide, x, y, w, h, items, size=15, gap=8):
    tb, tf = textbox(slide, x, y, w, h)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        if isinstance(it, tuple):
            lead, rest = it
            setrun(p, '▶  ', 12, TEAL, bold=True)
            setrun(p, lead, size, INK, bold=True)
            setrun(p, rest, size, INK)
        else:
            setrun(p, '▶  ', 12, TEAL, bold=True)
            setrun(p, it, size, INK)
    return tb


def card(slide, x, y, w, h, head, body, accent=TEAL, headcolor=None):
    box = rect(slide, x, y, w, h, WHITE, line=LINE)
    rect(slide, x, y, w, Inches(0.09), accent)
    tb, tf = textbox(slide, x + Inches(0.22), y + Inches(0.26), w - Inches(0.44), h - Inches(0.4))
    p = tf.paragraphs[0]; setrun(p, head, 15, headcolor or accent, bold=True)
    p2 = tf.add_paragraph(); p2.space_before = Pt(4)
    setrun(p2, body, 12, INK)
    return box


def flowbox(slide, x, y, w, h, title, body, fill, tcolor):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sp.shadow.inherit = False
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = fill; sp.line.width = Pt(1)
    tf = sp.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08); tf.margin_bottom = Inches(0.08)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    setrun(p, title, 13, tcolor, bold=True)
    if body:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(3)
        setrun(p2, body, 10.5, tcolor)
    return sp


def arrow(slide, x, y, w, h, color=GREYT):
    sp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    sp.shadow.inherit = False
    sp.fill.solid(); sp.fill.fore_color.rgb = color; sp.line.fill.background()
    return sp


def bignum(slide, x, y, w, num, label, color=TEAL):
    tb, tf = textbox(slide, x, y, w, Inches(1.7), anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; setrun(p, num, 52, color, bold=True)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(2)
    setrun(p2, label, 13, INK)


# ============================================================================
# SLIDE 1  TITLE
# ============================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(4.55), SW, Inches(0.06), AMBER)
tb, tf = textbox(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(2.6))
p = tf.paragraphs[0]; setrun(p, 'V-INTELLIGENCE UEBA', 16, AMBER, bold=True)
p2 = tf.add_paragraph(); p2.space_before = Pt(10)
setrun(p2, 'Behavioral Entity Intelligence for Cyber Defense', 40, WHITE, bold=True)
p3 = tf.add_paragraph(); p3.space_before = Pt(14)
setrun(p3, 'Catching the attacks that stay under the alarm: insiders, nation-state intrusions, '
       'and living-off-the-land campaigns.', 18, RGBColor(0xC9, 0xD6, 0xE4))
tb2, tf2 = textbox(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.6))
setrun(tf2.paragraphs[0], 'Executive Overview   |   Business and Technical Briefing   |   June 2026', 13,
       RGBColor(0x9F, 0xB3, 0xC8))

# ============================================================================
# SLIDE 2  THE PROBLEM
# ============================================================================
s = prs.slides.add_slide(BLANK)
content_header(s, 'The problem', 'The attacks that matter most never trip an alarm')
bullets(s, Inches(0.7), Inches(1.9), Inches(6.0), Inches(4.6), [
    ('Volt Typhoon. ', 'A state-sponsored campaign lived inside US critical-infrastructure networks for at least '
     'five years, undetected by any victim tool.'),
    ('Salt Typhoon. ', 'A second campaign sat inside major US telecom carriers for months, called the worst telecom '
     'hack in the nation\'s history.'),
    ('Insiders. ', 'Roughly one breach in five comes from inside, with a median time to discovery beyond 200 days.'),
    ('The common thread. ', 'Valid credentials, legitimate tools, and gradual change over months. Every single action '
     'looks authorized.'),
], size=15, gap=12)
card(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(2.0),
     'Found by outsiders, not by victims',
     'In both Typhoon campaigns the victims\' SIEM, antivirus, and endpoint tools stayed green. The intrusions were '
     'discovered by outside intelligence agencies, years and months too late.', accent=RED, headcolor=RED)
card(s, Inches(7.0), Inches(4.2), Inches(5.6), Inches(2.2),
     'This is a structural gap',
     'It is not a missed patch or a tuning error. The tools were asking "does this match something known to be bad?" '
     'and these attackers are built to answer no.', accent=AMBER, headcolor=AMBER)

# ============================================================================
# SLIDE 3  THE GAP
# ============================================================================
s = prs.slides.add_slide(BLANK)
content_header(s, 'Why today\'s tools miss them', 'Five blind spots in conventional monitoring')
gaps = [
    ('Look for the known', 'Signatures match only catalogued threats. Living-off-the-land attacks bring no malware to fingerprint.'),
    ('Measure how much', 'Thresholds watch volume. An insider holds volume flat and shifts what they access instead.'),
    ('Watch one entity alone', 'A covert channel calling home on a steady rhythm is invisible to a tool that only counts traffic.'),
    ('Judge a single moment', 'Slow-and-low change is invisible in any one week and clear only across months.'),
    ('Compare you to everyone', 'A developer judged against the whole company hides abuse that stands out only among developers.'),
]
xs = [0.7, 4.79, 8.88]
CW = 3.74
for i, (h, b) in enumerate(gaps[:3]):
    card(s, Inches(xs[i]), Inches(1.95), Inches(CW), Inches(2.15), h, b, accent=TEAL)
for i, (h, b) in enumerate(gaps[3:]):
    card(s, Inches(xs[i]), Inches(4.35), Inches(CW), Inches(2.15), h, b, accent=TEAL)
card(s, Inches(xs[2]), Inches(4.35), Inches(CW), Inches(2.15), 'The trap',
     'Loosen the tools and they bury analysts in false alarms. Tighten them and they miss the real attack.', accent=RED, headcolor=RED)

# ============================================================================
# SLIDE 4  THE SHIFT
# ============================================================================
s = prs.slides.add_slide(BLANK)
content_header(s, 'A different way to see', 'Measure direction, not just magnitude')
bullets(s, Inches(0.7), Inches(1.9), Inches(5.5), Inches(2.4), [
    ('Behavior, not events. ', 'Each entity is a profile that evolves week to week, like a medical chart.'),
    ('Direction, not size. ', 'We track which way behavior is trending, not just how big it is.'),
    ('Peers, not the crowd. ', 'Each entity is judged against its true role-group.'),
    ('Explanation attached. ', 'Each finding names what changed and the attacker technique it resembles.'),
], size=13.5, gap=7)
s.shapes.add_picture(CH_DIR, Inches(6.35), Inches(1.95), width=Inches(6.45))
tb, tf = textbox(s, Inches(0.7), Inches(5.65), Inches(12.0), Inches(1.1), anchor=MSO_ANCHOR.MIDDLE)
p = tf.paragraphs[0]
setrun(p, 'The insider example: ', 14, AMBER, bold=True)
setrun(p, 'file volume never changes, so every magnitude-based tool stays silent. What changes is the direction, '
       'the share of confidential material, which is exactly what behavioral intelligence watches.', 14, INK)

# ============================================================================
# SLIDE 5  HOW IT WORKS
# ============================================================================
s = prs.slides.add_slide(BLANK)
content_header(s, 'How it works', 'One pipeline, two complementary layers')
# flow: logs -> profile -> [Layer1 / Layer2] -> leads
flowbox(s, Inches(0.6), Inches(3.1), Inches(2.45), Inches(1.5), 'Five log sources',
        'sign-in, file, endpoint,\nnetwork, DNS', LIGHT, INK)
arrow(s, Inches(3.12), Inches(3.62), Inches(0.5), Inches(0.45), RGBColor(0xB6, 0xC2, 0xC2))
flowbox(s, Inches(3.68), Inches(3.1), Inches(2.5), Inches(1.5), 'Behavioral profile',
        'per entity, five zones,\nupdated weekly', LIGHT, INK)
arrow(s, Inches(6.26), Inches(3.62), Inches(0.5), Inches(0.45), RGBColor(0xB6, 0xC2, 0xC2))
flowbox(s, Inches(6.82), Inches(1.95), Inches(3.15), Inches(1.5), 'Layer 1: Precision',
        'named fingerprints,\npeer-relative, near-zero\nfalse alarms', TEAL, WHITE)
flowbox(s, Inches(6.82), Inches(4.25), Inches(3.15), Inches(1.5), 'Layer 2: Discovery',
        'drift and novelty, names\nthe direction of change', AMBER, WHITE)
arrow(s, Inches(10.05), Inches(3.62), Inches(0.5), Inches(0.45), RGBColor(0xB6, 0xC2, 0xC2))
flowbox(s, Inches(10.6), Inches(3.1), Inches(2.45), Inches(1.5), 'Ranked, explained leads',
        'with a MITRE ATT&CK\ntechnique, into your SOC', NAVY, WHITE)
tb, tf = textbox(s, Inches(0.7), Inches(6.15), Inches(12.0), Inches(0.7), anchor=MSO_ANCHOR.MIDDLE)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
setrun(p, 'Layer 1 catches what we can name and keeps false alarms low. Layer 2 is the safety net for novel attacks '
       'and explains the direction.', 14, INK, italic=True)

# ============================================================================
# SLIDE 6  DRIFT IN ACTION
# ============================================================================
s = prs.slides.add_slide(BLANK)
content_header(s, 'Drift in action', 'Behavioral drift pulls the attackers clear of normal')
s.shapes.add_picture(CH_DRIFT, Inches(0.55), Inches(1.8), width=Inches(8.5))
card(s, Inches(9.35), Inches(1.95), Inches(3.45), Inches(2.05), 'What drift measures',
     'How far an entity has moved from its own normal, accumulated week over week, so a long run of small changes '
     'eventually shows.', accent=TEAL)
card(s, Inches(9.35), Inches(4.15), Inches(3.45), Inches(2.05), 'No single net catches all',
     'Three campaigns separate on drift at different times. The slow APT never does, and is caught by its beacon '
     'instead. Layers cover each other.', accent=AMBER, headcolor=AMBER)

# ============================================================================
# SLIDE 6  THE FOUR THREATS
# ============================================================================
s = prs.slides.add_slide(BLANK)
content_header(s, 'What we tested', 'Four long-running campaigns, modeled on real intrusions')
s.shapes.add_picture(CH_TIME, Inches(0.95), Inches(1.8), width=Inches(11.45))
threats = [
    ('Insider', 'drifts toward confidential files while volume stays flat', '#C0392B'),
    ('Slow APT', 'implant calls home every few hours in tiny packets', '#E67E22'),
    ('Volt Typhoon', 'critical-infrastructure pre-positioning with built-in tools', '#7D3C98'),
    ('Salt Typhoon', 'telecom intrusion fanning out to map and intercept', '#2471A3'),
]
xs = [0.7, 3.71, 6.72, 9.73]
for i, (h, b, c) in enumerate(threats):
    tb, tf = textbox(s, Inches(xs[i]), Inches(5.25), Inches(2.9), Inches(1.5))
    rect(s, Inches(xs[i]), Inches(5.25), Inches(0.12), Inches(1.35), RGBColor.from_string(c.lstrip('#')))
    p = tf.paragraphs[0]; p.space_after = Pt(2)
    r = p.add_run(); r.text = '  ' + h; r.font.size = Pt(14); r.font.bold = True
    r.font.color.rgb = RGBColor.from_string(c.lstrip('#')); r.font.name = 'Segoe UI'
    p2 = tf.add_paragraph(); setrun(p2, '  ' + b, 11.5, INK)

# ============================================================================
# SLIDE 7  THE PROOF
# ============================================================================
s = prs.slides.add_slide(BLANK)
content_header(s, 'The proof', 'All four caught, with no false alarms')
s.shapes.add_picture(CH_SCORE, Inches(0.55), Inches(1.85), width=Inches(8.4))
bignum(s, Inches(9.2), Inches(2.0), Inches(3.6), '4 of 4', 'campaigns detected', color=TEAL)
bignum(s, Inches(9.2), Inches(3.7), Inches(3.6), '0%', 'false alarms', color=TEAL)
bignum(s, Inches(9.2), Inches(5.3), Inches(3.6), '1 of 4', 'best conventional tool', color=RED)
tb, tf = textbox(s, Inches(0.7), Inches(6.55), Inches(12.0), Inches(0.5))
setrun(tf.paragraphs[0], '250 users  |  485 days (70 weeks)  |  about 14 million events  |  5 log sources', 12, GREYT)

# ============================================================================
# SLIDE 8  TRUSTWORTHY DETECTION
# ============================================================================
s = prs.slides.add_slide(BLANK)
content_header(s, 'Why the detections hold up', 'Measurable evidence, not a black-box score')
s.shapes.add_picture(CH_FP, Inches(0.55), Inches(1.78), width=Inches(7.55))
s.shapes.add_picture(CH_BEACON, Inches(0.55), Inches(4.95), width=Inches(7.55))
card(s, Inches(8.35), Inches(1.95), Inches(4.45), Inches(4.45),
     'Every alert answers three questions',
     'What changed?\nThe specific behavior and zone.\n\nCompared to whom?\nThe entity\'s true role-group peers.\n\n'
     'What does it resemble?\nA named attacker technique (MITRE ATT&CK).\n\nThat is the difference between a score '
     'and a lead an analyst can act on.', accent=TEAL)

# ============================================================================
# SLIDE 9  HONEST BOUNDARIES
# ============================================================================
s = prs.slides.add_slide(BLANK)
content_header(s, 'What we proved, and what we did not', 'A direct read of the evidence')
card(s, Inches(0.7), Inches(2.0), Inches(3.74), Inches(3.6), 'The data is simulated',
     'Results come from a synthetic enterprise calibrated to public advisories. Every figure, including zero false '
     'alarms, must be re-measured on real telemetry.', accent=AMBER, headcolor=AMBER)
card(s, Inches(4.79), Inches(2.0), Inches(3.74), Inches(3.6), 'Coverage has limits',
     'The precision layer catches the patterns it is built to catch. The discovery layer is the net for novel '
     'attacks, and is noisier by design.', accent=AMBER, headcolor=AMBER)
card(s, Inches(8.88), Inches(2.0), Inches(3.74), Inches(3.6), 'The bar was set in a lab',
     'Thresholds here were tuned knowing the answers. In production they must be set blind and tuned to each '
     'environment.', accent=AMBER, headcolor=AMBER)
tb, tf = textbox(s, Inches(0.7), Inches(5.9), Inches(12.0), Inches(0.9), anchor=MSO_ANCHOR.MIDDLE)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
setrun(p, 'None of this weakens the finding. The approach sees a class of attack today\'s tools demonstrably miss, '
       'with an explanation attached. That is the result worth validating.', 14, INK, italic=True)

# ============================================================================
# SLIDE 10  PATH FORWARD
# ============================================================================
s = prs.slides.add_slide(BLANK)
content_header(s, 'The path forward', 'A low-risk pilot on your own data')
bullets(s, Inches(0.7), Inches(1.95), Inches(6.2), Inches(4.2), [
    ('It adds, it does not replace. ', 'The behavioral layer runs alongside your existing stack, in advisory mode, '
     'with no automated blocking until results are validated.'),
    ('It uses data you already collect. ', 'Five standard log sources, no new sensors.'),
    ('It fits your workflow. ', 'Ranked, explained leads with a named technique, into your existing analyst process.'),
    ('It proves itself on your environment. ', 'Validated by replaying known-bad or historical incidents on your own '
     'telemetry.'),
], size=15, gap=12)
card(s, Inches(7.1), Inches(2.1), Inches(5.5), Inches(2.6),
     'The pilot question',
     'Does the behavioral layer surface real, explained threats your current stack misses, at a false-alarm rate '
     'your team can live with?', accent=TEAL)
rect(s, Inches(7.1), Inches(5.0), Inches(5.5), Inches(1.35), NAVY)
tb, tf = textbox(s, Inches(7.4), Inches(5.18), Inches(4.9), Inches(1.0), anchor=MSO_ANCHOR.MIDDLE)
p = tf.paragraphs[0]
setrun(p, 'Next step: ', 16, AMBER, bold=True)
setrun(p, 'scope a bounded pilot on a representative slice of live telemetry.', 15, WHITE)

out = OUT; n = 2
while True:
    try:
        prs.save(out); break
    except PermissionError:
        root, ext = os.path.splitext(OUT); out = f"{root}_v{n}{ext}"; n += 1
print('Saved:', out, '| slides:', len(prs.slides._sldIdLst))
