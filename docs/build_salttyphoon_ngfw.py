"""Salt Typhoon TTPs vs NGFW Coverage — technical coverage deck.
Color-coded YES / PARTIAL / NO matrices mapping G1045 techniques to
verifiable NGFW controls. Reuses the house palette/card style."""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

BASE = os.path.join(os.path.dirname(__file__), "SaltTyphoon_NGFW_Coverage.pptx")

DKNAVY = RGBColor(0x07, 0x14, 0x2A)
NAVY   = RGBColor(0x0B, 0x1F, 0x3A)
BODY   = RGBColor(0x2D, 0x37, 0x48)
SUB    = RGBColor(0x5A, 0x68, 0x78)
ORG    = RGBColor(0xF2, 0x6B, 0x1F)
TEAL   = RGBColor(0x08, 0x91, 0xB2)
CRED   = RGBColor(0xC0, 0x39, 0x2B)
GREEN  = RGBColor(0x1E, 0x7A, 0x44)
AMBER  = RGBColor(0xB5, 0x7A, 0x12)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
CARD   = RGBColor(0xF7, 0xF8, 0xFA)
PEACH  = RGBColor(0xFD, 0xE5, 0xD4)
CYAN   = RGBColor(0xE0, 0xF2, 0xF7)
GOLD   = RGBColor(0xB8, 0x86, 0x2A)
GREEN_BG = RGBColor(0xD9, 0xEE, 0xDF)
AMBER_BG = RGBColor(0xFB, 0xEA, 0xCC)
RED_BG   = RGBColor(0xF6, 0xDA, 0xD6)
LINE     = RGBColor(0xD8, 0xDD, 0xE3)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = prs.slide_width, prs.slide_height

VC = {'YES': (GREEN, GREEN_BG), 'PARTIAL': (AMBER, AMBER_BG), 'NO': (CRED, RED_BG)}


def _fill(s, c):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = c

def _rect(s, l, t, w, h, fill=None):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.background()
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background(); sh.shadow.inherit = False
    return sh

def _box(s, l, t, w, h, fill=None, accent=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.background()
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background(); sh.shadow.inherit = False
    if accent:
        _rect(s, l, t, Inches(0.1), h, fill=accent)
    return sh

def _txt(s, l, t, w, h, text, sz=11, c=BODY, b=False, al=PP_ALIGN.LEFT, fn='Calibri'):
    tf = s.shapes.add_textbox(l, t, w, h).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b
    p.font.name = fn; p.alignment = al
    return tf

def _p(tf, text, sz=11, c=BODY, b=False, sp=Pt(4)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b
    p.font.name = 'Calibri'; p.space_before = sp
    return p

def _content(part, title, sub=None, bar=TEAL):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _fill(s, WHITE)
    _rect(s, Inches(0), Inches(0), W, Inches(0.3), fill=bar)
    _txt(s, Inches(0.5), Inches(0.0), Inches(12.3), Inches(0.3), part, 11, WHITE, True)
    _txt(s, Inches(0.5), Inches(0.42), Inches(12.3), Inches(0.55), title, 23, NAVY, True, fn='Georgia')
    if sub:
        _txt(s, Inches(0.5), Inches(0.97), Inches(12.3), Inches(0.32), sub, 12.5, SUB)
    return s

def _bottom(s, text, y=6.9):
    _rect(s, Inches(0.5), Inches(y), Inches(12.3), Inches(0.42), fill=NAVY)
    _txt(s, Inches(0.7), Inches(y - 0.02), Inches(12.0), Inches(0.45), text, 11, GOLD, True, PP_ALIGN.CENTER)

def _cell(cell, text, sz=8.5, c=BODY, b=False, fill=None, al=PP_ALIGN.LEFT):
    if fill:
        cell.fill.solid(); cell.fill.fore_color.rgb = fill
    else:
        cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
    cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.06)
    cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = al
    p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b; p.font.name = 'Calibri'

def _rowest(notes):
    n = len(notes)
    return 0.30 if n < 60 else (0.44 if n < 130 else (0.60 if n < 210 else 0.76))

def add_table(s, left, top, width, label, label_color, rows):
    _txt(s, left, top, width, Inches(0.26), label, 12, label_color, True)
    ttop = top + Inches(0.3)
    nrows = len(rows) + 1
    gtbl = s.shapes.add_table(nrows, 3, left, ttop, width, Inches(0.3 * nrows)).table
    gtbl.first_row = False; gtbl.horz_banding = False
    gtbl.columns[0].width = Inches(2.55)
    gtbl.columns[1].width = Inches(1.15)
    gtbl.columns[2].width = width - Inches(3.70)
    # header
    _cell(gtbl.cell(0, 0), 'Technique (MITRE ATT&CK)', 8.5, WHITE, True, NAVY)
    _cell(gtbl.cell(0, 1), 'NGFW', 8.5, WHITE, True, NAVY, PP_ALIGN.CENTER)
    _cell(gtbl.cell(0, 2), 'Notes', 8.5, WHITE, True, NAVY)
    gtbl.rows[0].height = Inches(0.28)
    for i, (ttp, verdict, notes) in enumerate(rows, start=1):
        vc, vbg = VC[verdict]
        _cell(gtbl.cell(i, 0), ttp, 8.5, NAVY, True, CARD)
        _cell(gtbl.cell(i, 1), verdict, 8.5, vc, True, vbg, PP_ALIGN.CENTER)
        _cell(gtbl.cell(i, 2), notes, 8.0, BODY, False, WHITE)
        gtbl.rows[i].height = Inches(_rowest(notes))
    est = 0.3 + 0.28 + sum(_rowest(n) for _, _, n in rows)
    return top + Inches(0.3 + est) + Inches(0.12)


# ============ SLIDE 1 — TITLE ============
s = prs.slides.add_slide(prs.slide_layouts[6])
_fill(s, DKNAVY)
_rect(s, Inches(0), Inches(0), W, Inches(0.18), fill=CRED)
_rect(s, Inches(0), H - Inches(0.18), W, Inches(0.18), fill=TEAL)
_txt(s, Inches(1), Inches(2.0), Inches(11.3), Inches(0.9),
     'Salt Typhoon TTPs vs NGFW Coverage', 34, WHITE, True, PP_ALIGN.CENTER, 'Georgia')
_txt(s, Inches(1), Inches(3.1), Inches(11.3), Inches(0.5),
     'Mapping G1045 techniques to verifiable next-gen firewall controls', 16,
     RGBColor(0xA8, 0xB6, 0xC4), False, PP_ALIGN.CENTER)
_txt(s, Inches(1), Inches(4.1), Inches(11.3), Inches(0.9),
     'PRC state-sponsored telecom espionage  ·  CISA AA25-239A  ·  Enterprise blast-radius containment view',
     13, RGBColor(0x90, 0xA0, 0xB0), False, PP_ALIGN.CENTER)
_txt(s, Inches(1), Inches(6.55), Inches(11.3), Inches(0.4),
     'Current as of May 29, 2026  ·  Sources: CISA AA25-239A, Cisco Talos, MITRE ATT&CK G1045',
     10, RGBColor(0x60, 0x70, 0x80), False, PP_ALIGN.CENTER)


# ============ SLIDE 2 — BACKGROUND ============
s = _content('THREAT BACKGROUND', 'Who Salt Typhoon Is — and What They Did',
             'PRC state-sponsored espionage against telecom carriers, active since at least 2021.', CRED)
facts = [
    ('600+', 'organizations across 80+ countries notified by the FBI as targeted'),
    ('9', 'U.S. carriers confirmed compromised (7 named: Verizon, AT&T, T-Mobile, Charter, Lumen, Consolidated, Windstream)'),
    ('CALEA', 'lawful-intercept infrastructure reached — the court-authorized surveillance system'),
    ('3+ yrs', 'dwell time in some compromised carrier networks'),
]
x = Inches(0.5)
for big, desc in facts:
    _box(s, x, Inches(1.5), Inches(2.97), Inches(1.7), fill=CARD, accent=CRED)
    _txt(s, x + Inches(0.15), Inches(1.62), Inches(2.7), Inches(0.55), big, 26, CRED, True, PP_ALIGN.CENTER, 'Georgia')
    _txt(s, x + Inches(0.18), Inches(2.25), Inches(2.62), Inches(0.9), desc, 9.5, BODY)
    x += Inches(3.12)

_box(s, Inches(0.5), Inches(3.45), Inches(12.3), Inches(3.2), fill=CARD, accent=NAVY)
_txt(s, Inches(0.8), Inches(3.52), Inches(11.7), Inches(0.3), 'THE TRADECRAFT', 13, NAVY, True)
tf = _txt(s, Inches(0.8), Inches(3.9), Inches(11.7), Inches(2.6),
     'Establish on carrier-grade routers using stolen credentials, then modify ACLs and AAA configurations to maintain persistent access.', 11.5, BODY)
_p(tf, 'Abuse built-in services — Guest Shell and sshd_operns — for on-device persistence.', 11.5, BODY, sp=Pt(7))
_p(tf, 'Run the custom JumbledPath utility from jump hosts to capture traffic and disable logging on compromised devices.', 11.5, BODY, sp=Pt(7))
_p(tf, 'Exfiltrate call records, metadata, voice and SMS over GRE tunnels built between compromised routers.', 11.5, BODY, sp=Pt(7))
_p(tf, 'Tracked as MITRE ATT&CK Group G1045; detailed in CISA AA25-239A (Aug 2025) and Cisco Talos "Weathering the Storm" (Feb 2025).', 11, NAVY, True, sp=Pt(9))

_bottom(s, 'No malware required to enter — the entry path is stolen credentials and configuration abuse on the carrier core.')


# ============ SLIDE 3 — SCOPE & FRAMING ============
s = _content('SCOPE & FRAMING', 'The Enterprise NGFW Story Is Blast-Radius Containment',
             'Salt Typhoon\'s primary victim is the carrier. Most of its activity is not in path for any enterprise firewall.', TEAL)

_box(s, Inches(0.5), Inches(1.5), Inches(5.95), Inches(3.6), fill=RED_BG, accent=CRED)
_txt(s, Inches(0.8), Inches(1.57), Inches(5.4), Inches(0.35), 'NOT IN ENTERPRISE NGFW PATH', 12, CRED, True)
tf = _txt(s, Inches(0.8), Inches(1.98), Inches(5.4), Inches(3.0),
     'Activity inside compromised carrier infrastructure is invisible to any enterprise firewall:', 11, BODY)
_p(tf, '• Router-to-router GRE tunneling in the carrier backbone', 10.5, BODY, sp=Pt(6))
_p(tf, '• On-box AAA / ACL edits and config changes', 10.5, BODY, sp=Pt(5))
_p(tf, '• Guest Shell / sshd_operns enablement on the device', 10.5, BODY, sp=Pt(5))
_p(tf, '• On-box packet capture (CDR / CALEA) and log clearing', 10.5, BODY, sp=Pt(5))
_p(tf, 'The CALEA compromise is a policy / regulatory matter — no NGFW configuration prevents it.', 10.5, CRED, True, sp=Pt(9))

_box(s, Inches(6.85), Inches(1.5), Inches(5.95), Inches(3.6), fill=GREEN_BG, accent=GREEN)
_txt(s, Inches(7.15), Inches(1.57), Inches(5.4), Inches(0.35), 'WHERE THE NGFW ADDS VALUE', 12, GREEN, True)
tf = _txt(s, Inches(7.15), Inches(1.98), Inches(5.4), Inches(3.0),
     'For an enterprise consuming carrier services, the firewall job is downstream — contain a carrier compromise inside your own network:', 11, BODY)
_p(tf, '• Default-deny zero-trust segmentation at the carrier-to-enterprise handoff', 10.5, BODY, sp=Pt(6))
_p(tf, '• Egress allowlists from sensitive and management zones', 10.5, BODY, sp=Pt(5))
_p(tf, '• Port / protocol denies for observed Salt Typhoon channels', 10.5, BODY, sp=Pt(5))
_p(tf, '• Continuous config-drift detection on boundary and handoff rules', 10.5, BODY, sp=Pt(5))
_p(tf, 'Carrier-side coverage is possible too — but only if verification is deployed at the carrier.', 10.5, GREEN, True, sp=Pt(9))

_bottom(s, 'Honest scoping: the firewall contains the blast radius downstream; it does not defend the carrier core.')


# ============ SLIDE 4 — WHAT CONFIG HYGIENE PREVENTS ============
s = _content('PREEMPTIVE VALUE', 'What Formal Configuration Hygiene Would Have Prevented',
             'Several Salt Typhoon paths trace to configuration gaps a formal NGFW analysis is designed to find before exploitation.', ORG)
items = [
    ('GRE / IPsec peer authorization', 'Absent explicit peer authorization is a primary persistence channel for the tunneling. Flagged as a coverage gap on any onboarded management plane.'),
    ('Zero-trust segmentation verification', 'At the carrier↔enterprise boundary, closes the lateral pivot path that enables downstream compromise.'),
    ('Egress allowlist verification', 'On management interfaces, blocks the configuration exfiltration performed over TFTP and FTP.'),
    ('Continuous configuration-drift detection', 'Flags unauthorized changes to segmentation, egress policy, or carrier-handoff rules the moment they appear in live config.'),
    ('Reachability-graph (symbolic) verification', 'Unlike rule-syntax or signature audits, verifies the reachability implied by the live policy — surfacing shadowed allow-rules and segmentation gaps no single rule reveals.'),
]
y = Inches(1.5)
for i, (h, d) in enumerate(items):
    _box(s, Inches(0.5), y, Inches(12.3), Inches(0.98), fill=CARD, accent=ORG if i % 2 == 0 else TEAL)
    _txt(s, Inches(0.8), y + Inches(0.08), Inches(4.0), Inches(0.85), h, 12, NAVY, True)
    _txt(s, Inches(5.0), y + Inches(0.08), Inches(7.6), Inches(0.85), d, 10.5, BODY)
    y += Inches(1.04)

_bottom(s, 'The symbolic model verifies what the policy actually allows — not just what each rule says in isolation.', y=6.85)


# ============ SLIDE 5 — COVERAGE SCORECARD ============
s = _content('COVERAGE SCORECARD', 'NGFW Coverage at a Glance',
             '26 technique-mappings across 8 ATT&CK tactics (some techniques mapped under two tactics).', TEAL)
score = [('YES — fully blockable', 6, GREEN, GREEN_BG),
         ('PARTIAL — partially blockable', 11, AMBER, AMBER_BG),
         ('NO — out of NGFW scope', 9, CRED, RED_BG)]
x = Inches(0.5)
for label, n, col, bg in score:
    _box(s, x, Inches(1.5), Inches(3.97), Inches(2.0), fill=bg, accent=col)
    _txt(s, x + Inches(0.2), Inches(1.7), Inches(3.6), Inches(0.9), str(n), 48, col, True, PP_ALIGN.CENTER, 'Georgia')
    _txt(s, x + Inches(0.2), Inches(2.85), Inches(3.6), Inches(0.5), label, 12, NAVY, True, PP_ALIGN.CENTER)
    x += Inches(4.13)

_box(s, Inches(0.5), Inches(3.75), Inches(12.3), Inches(2.9), fill=CARD, accent=NAVY)
_txt(s, Inches(0.8), Inches(3.83), Inches(11.7), Inches(0.3), 'HOW TO READ THIS', 13, NAVY, True)
tf = _txt(s, Inches(0.8), Inches(4.22), Inches(11.7), Inches(2.3),
     'The 9 NO verdicts and many of the 11 PARTIAL verdicts sit inside the compromised carrier\'s gear — AAA edits, '
     'Guest Shell, on-box capture, log clearing — invisible to any enterprise firewall.', 11.5, BODY)
_p(tf, 'The 6 YES verdicts are the enterprise blast-radius wins: zero-trust segmentation at the carrier handoff (T1199), '
       'config-dump egress denies (T1602.002), non-standard-port denies (T1571), and C2 / unencrypted-exfil blocks (T1041, T1048.003).', 11.5, NAVY, True, sp=Pt(8))
_p(tf, 'Low-and-slow exfiltration (T1029) is weak on NGFWs by design — a behavioral layer (V-Intelligence UEBA) or NDR is the right control.', 11.5, BODY, sp=Pt(8))

_bottom(s, 'A firewall cannot defend a router it never sees — but it can stop the compromise from spreading into the enterprise.')


# ============ SLIDES 6-9 — TTP TABLES ============
IA = [
    ('T1078 Valid Accounts', 'PARTIAL', 'Enforce source-IP / geo restrictions and require MFA on VPN/SSL-VPN logins. Cannot detect that the credential itself is stolen.'),
    ('T1190 Exploit Public-Facing App', 'PARTIAL', 'IPS signatures for the relevant edge CVEs (Cisco IOS XE, Smart Install, PAN) plus mgmt-plane ACLs blocking /webui/ and TCP/4786. Enterprise NGFWs are rarely in path for carrier-core devices; coverage applies mainly to enterprise-edge devices in the same families.'),
]
CA = [
    ('T1555 Credentials from Password Stores', 'NO', 'Offline decryption of exfiltrated device configs. No NGFW visibility into off-box activity.'),
    ('T1040 Network Sniffing (TACACS+ TCP/49)', 'PARTIAL', 'Segment the management VRF and restrict TACACS+ flows to authorized device-AAA pairs. Cannot stop sniffing that runs on a compromised router itself.'),
    ('T1552 Unsecured Credentials', 'NO', 'Credentials stored in device configs. On-device — no NGFW visibility.'),
    ('T1556.004 Modify Auth Process: Network Device Auth', 'PARTIAL', 'TACACS+ redirect captures admin creds when AAA traffic is rerouted. Egress allowlist on AAA/TACACS+/RADIUS to the authorized server set blocks the redirect. (Also maps to Persistence.)'),
    ('T1110.002 Brute Force: Password Cracking', 'NO', 'Offline cracking of captured TACACS+/RADIUS hashes. No NGFW visibility into off-box activity.'),
]
PE = [
    ('T1098.004 Account Manipulation: SSH Authorized Keys', 'NO', 'On-device modification of SSH authorized_keys.'),
    ('T1136 Create Account (local privileged)', 'NO', 'On-device account creation.'),
    ('T1505 Server Software Component (Guest Shell, sshd_operns)', 'PARTIAL', 'Deny TCP/57722 inbound and outbound across all zones. Cannot prevent Guest Shell enablement on the device itself.'),
    ('T1556.004 Modify Auth Process: Network Device Auth', 'PARTIAL', 'Egress allowlist on AAA/TACACS+/RADIUS catches redirect to an attacker-controlled AAA server. Cannot prevent the on-device config change.'),
]
DE = [
    ('T1027 Obfuscated Files (encrypted PCAPs)', 'NO', 'Once encrypted, the NGFW cannot inspect content.'),
    ('T1070 Indicator Removal (log clearing)', 'NO', 'On-device action against local logs.'),
    ('T1562 Impair Defenses (logging mods)', 'NO', 'On-device logging configuration changes.'),
    ('T1021.004 Remote Services: SSH (loopback rewrite to bypass ACLs)', 'PARTIAL', 'An independent NGFW with allow-rules bound to authorized loopback ranges defeats the bypass at the next hop. Cannot prevent the router rewriting its own loopback.'),
]
LM = [
    ('T1021 Remote Services (SSH on non-standard ports)', 'YES', 'Deny SSH on non-standard ports across zones. App-ID detection of SSH on unexpected ports. Force bastion-only paths.'),
    ('JumbledPath multi-hop chains (behavior; not in ATT&CK)', 'PARTIAL', 'Egress filtering from network-device zones to internet; threat-feed denies on known JumbledPath C2. Cannot stop intra-zone chaining among already-compromised devices.'),
    ('T1199 Trusted Relationship (carrier→enterprise pivot)', 'YES', 'Default-deny zero-trust segmentation at trust boundaries with explicit per-service allow-rules. Highest-value NGFW control for enterprise customers downstream of a compromised carrier; not applicable to the carrier itself.'),
]
COL = [
    ('T1040 Network Sniffing (on-box CDR/CALEA capture)', 'NO', 'Capture runs on the compromised carrier device itself. No NGFW visibility.'),
    ('T1602.002 Data from Config Repository: Device Config Dump', 'YES', 'Deny outbound TFTP/FTP/HTTP from network-device management interfaces to anything outside the authorized config-mgmt server set. (Value where the NGFW is in path of the device\'s mgmt egress.)'),
]
C2 = [
    ('JumbledPath multi-hop chains (behavior; not in ATT&CK)', 'PARTIAL', 'Same controls as Lateral Movement: egress filtering plus threat-feed denies on C2 destinations.'),
    ('T1571 Non-Standard Port (sshd_operns 57722)', 'YES', 'Port-based denies for 57722 and other observed non-standard ports.'),
    ('T1572 Protocol Tunneling (GRE and IPsec)', 'PARTIAL', 'Explicit GRE/IPsec peer authorization with default-deny tunnels except between named endpoints. Effective where the NGFW sits between device and tunnel peer; router-to-router tunneling inside a carrier backbone bypasses any enterprise NGFW.'),
]
EX = [
    ('T1029 Scheduled Transfer (low-and-slow)', 'PARTIAL', 'Threat-feed destination blocks help. Volumetric / timing anomaly detection is weak on NGFWs — V-Intelligence UEBA or NDR is the right layer.'),
    ('T1041 Exfiltration Over C2 Channel', 'YES', 'Threat-feed blocks on known C2 IPs, domains and ASNs. Egress allowlists from sensitive zones (CALEA, mgmt VRF, lawful-intercept gear).'),
    ('T1048.003 Exfil Over Unencrypted Non-C2 Protocol (TFTP/FTP)', 'YES', 'Deny outbound TFTP/FTP/SFTP from network devices to the internet.'),
]

s = _content('TTP COVERAGE  ·  1 of 4', 'Initial Access  &  Credential Access', bar=TEAL)
y = add_table(s, Inches(0.5), Inches(1.35), Inches(12.3), 'INITIAL ACCESS', ORG, IA)
add_table(s, Inches(0.5), y, Inches(12.3), 'CREDENTIAL ACCESS', ORG, CA)

s = _content('TTP COVERAGE  ·  2 of 4', 'Persistence  &  Defense Evasion', bar=TEAL)
y = add_table(s, Inches(0.5), Inches(1.35), Inches(12.3), 'PERSISTENCE', ORG, PE)
add_table(s, Inches(0.5), y, Inches(12.3), 'DEFENSE EVASION', ORG, DE)

s = _content('TTP COVERAGE  ·  3 of 4', 'Lateral Movement  &  Collection', bar=TEAL)
y = add_table(s, Inches(0.5), Inches(1.35), Inches(12.3), 'LATERAL MOVEMENT', ORG, LM)
add_table(s, Inches(0.5), y, Inches(12.3), 'COLLECTION', ORG, COL)

s = _content('TTP COVERAGE  ·  4 of 4', 'Command & Control  &  Exfiltration', bar=TEAL)
y = add_table(s, Inches(0.5), Inches(1.35), Inches(12.3), 'COMMAND AND CONTROL', ORG, C2)
add_table(s, Inches(0.5), y, Inches(12.3), 'EXFILTRATION', ORG, EX)


# ============ SLIDE 10 — OUT OF SCOPE ============
s = _content('SCOPE BOUNDARY', 'What Sits Outside Any Enterprise NGFW',
             'The NO verdicts — and many PARTIAL ones — live inside the carrier\'s gear.', CRED)
_box(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.6), fill=RED_BG, accent=CRED)
tf = _txt(s, Inches(0.8), Inches(1.65), Inches(11.7), Inches(2.3),
     'AAA modification, ACL tampering, Guest Shell and sshd_operns enablement, on-box packet capture, and log clearing all occur '
     'on the router itself — invisible to any firewall outside the carrier.', 12.5, BODY)
_p(tf, 'Closing those gaps requires carrier-side configuration verification (which can be performed if deployed at the carrier).', 12.5, NAVY, True, sp=Pt(10))
_p(tf, 'The CALEA lawful-intercept compromise is a policy and regulatory matter, not a control-plane one — no NGFW configuration prevents it.', 12.5, CRED, True, sp=Pt(10))

_box(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.0), fill=CYAN, accent=TEAL)
_txt(s, Inches(0.8), Inches(4.42), Inches(11.7), Inches(0.3), 'THE TWO-LAYER TAKEAWAY', 13, TEAL, True)
tf = _txt(s, Inches(0.8), Inches(4.8), Inches(11.7), Inches(1.4),
     'NGFW configuration hygiene contains the enterprise blast radius and closes the verifiable boundary controls. '
     'Behavioral detection (V-Intelligence UEBA) and NDR cover the residual — low-and-slow exfiltration, credential abuse, '
     'and cross-domain movement that no single firewall rule reveals.', 12, BODY)

_bottom(s, 'Configuration verification proves the door is closed. Behavioral detection watches what the carrier compromise drags in.')


# ============ SLIDE 11 — SOURCES ============
s = _content('SOURCES', 'References', bar=NAVY)
srcs = [
    ('CISA AA25-239A', 'cisa.gov/news-events/cybersecurity-advisories/aa25-239a  (August 27, 2025)'),
    ('Cisco Talos — "Weathering the Storm"', 'blog.talosintelligence.com/salt-typhoon-analysis/  (February 2025)'),
    ('MITRE ATT&CK Group G1045 (Salt Typhoon)', 'attack.mitre.org/groups/G1045/'),
    ('Wall Street Journal', 'U.S. Wiretap Systems Targeted in China-Linked Hack  (October 2024)'),
    ('BleepingComputer', 'Charter and Windstream among nine US telecoms hacked by China  (December 2024)'),
]
y = Inches(1.6)
for t, u in srcs:
    _box(s, Inches(0.5), y, Inches(12.3), Inches(0.82), fill=CARD, accent=TEAL)
    _txt(s, Inches(0.8), y + Inches(0.12), Inches(4.6), Inches(0.6), t, 12, NAVY, True)
    _txt(s, Inches(5.5), y + Inches(0.14), Inches(7.1), Inches(0.6), u, 10.5, SUB)
    y += Inches(0.92)
_txt(s, Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.3), 'Current as of May 29, 2026.', 10, SUB)


# ---- save (auto-version if locked) ----
out = BASE; n = 2
while True:
    try:
        prs.save(out); break
    except PermissionError:
        root, ext = os.path.splitext(BASE); out = f"{root}_v{n}{ext}"; n += 1
print("Saved:", out, "|", len(prs.slides._sldIdLst), "slides")
