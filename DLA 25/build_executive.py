"""Top-executive deck (6 slides) for V-Intelligence UEBA, on the demand-deck template.
Storyline mirrors the app's Detection Comparison flow, at executive altitude:
  1. Three tools, three verdicts on the same 4 attacks (0/4 · 1/4 · 4/4)
  2. Same logs, read as living behavioral entities
  3. Signal separation — on any given day the attacker looks normal (dual-lens CUSUM)
  4. Multi-phase fingerprint — no normal user is extreme on every front (radar)
  5. Every attack caught, at the false-positive cost shown (composite + 0-FP known-bad)
  6. What this enables for the mission
Reuses figures generated this session.
"""
import copy
import os
import shutil
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

SRC = 'Legacy_vs_EDM_Demand_Forecast_Example-1-ver 2.pptx'
DST = 'UEBA_VectorIntelligence_Executive_Brief.pptx'
shutil.copyfile(SRC, DST)
FOOTER = 'CONFIDENTIAL  |  22nd Century Technologies  |  V-Intelligence UEBA — Executive Brief'
NAVY, BODY, BLUE, ORANGE, GREY = '102A4C', '1D2430', '1F5E9C', 'F36C21', '5E6A75'
D3_KEEP = {'Text 0', 'Text 1', 'Text 2', 'Text 20', 'Text 21', 'Text 22'}

# Slide 1 (demand slide 1 skeleton): three tools, three verdicts
S1 = {
    'Text 0': 'Same four stealth attacks. Three tools, three verdicts.',
    'Text 2': 'Inside authorized access, the most dangerous threats never trip a threshold — so most tools miss them, or bury the signal in false alarms.',
    'Text 4': 'TRADITIONAL SIEM', 'Text 5': 'point-anomaly tools',
    'Text 8': 'Isolation Forest / SVM / LOF', 'Text 11': 'Scores each user vs a threshold', 'Text 14': '0 of 4 caught',
    'Text 17': 'INTERMEDIATE', 'Text 18': 'statistical z-score',
    'Text 21': 'Flags the biggest outliers', 'Text 24': 'Catches 1 of 4',
    'Text 27': '…and alarms on nearly everyone', 'Text 30': 'no threat context',
    'Text 33': 'V-INTELLIGENCE', 'Text 34': 'behavioral entities',
    'Text 35': 'Cohort + drift + known-bad\nfused into one ranked score',
    'Text 36': '4 of 4 caught',
    'Text 38': 'Why it matters',
    'Text 39': 'The tools most agencies run today see all four attackers as normal employees. V-Intelligence catches every one — and names why.',
    'Text 40': FOOTER, 'Text 41': '1',
}

# Slide 2 (demand slide 2 skeleton): same logs -> behavioral entities
S2 = {
    'Text 0': 'Same logs, read as living behavioral entities',
    'Text 2': 'No new data. The identical security logs become connected, time-aware entities that reveal where a user is heading — not just what crossed a line today.',
    'Text 4': 'SAME LOGS', 'Text 5': 'no new data needed',
    'Text 8': 'Sign-in & access logs', 'Text 11': 'Network & endpoint records', 'Text 14': 'User & device context',
    'Text 17': 'BEHAVIORAL TWIN', 'Text 18': 'same logs, in context',
    'Text 20': 'Tracked over time (drift)', 'Text 22': 'Compared to its cohort (peers)',
    'Text 24': 'Matched to known-bad profiles', 'Text 26': 'Explained in plain language',
    'Text 29': 'DETECTION', 'Text 30': 'one ranked verdict',
    'Text 33': 'Behavioral risk score', 'Text 36': '+ named technique evidence',
    'Text 39': 'Calibrated, auditable', 'Text 42': 'Catches low-and-slow attacks',
    'Text 44': 'Key point',
    'Text 45': 'Same data, a richer model — it captures what a user HAS done plus where they’re HEADING.',
    'Text 46': FOOTER, 'Text 47': '2',
}


def set_text(shape, text):
    tf = shape.text_frame
    runs = [r for p in tf.paragraphs for r in p.runs]
    align = tf.paragraphs[0].alignment
    size = bold = ital = name = color = None
    if runs:
        f = runs[0].font
        size, bold, ital, name = f.size, f.bold, f.italic, f.name
        try:
            if f.color and f.color.type == 1:
                color = str(f.color.rgb)
        except Exception:
            pass
    tf.clear()
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        if size: r.font.size = size
        if bold is not None: r.font.bold = bold
        if ital is not None: r.font.italic = ital
        if name: r.font.name = name
        if color: r.font.color.rgb = RGBColor.from_string(color)


def apply(slide, mapping):
    by_name = {sh.name: sh for sh in slide.shapes}
    for nm, txt in mapping.items():
        if nm in by_name:
            set_text(by_name[nm], txt)


def recolor(slide, name, hexstr):
    for sh in slide.shapes:
        if sh.name == name:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    r.font.color.rgb = RGBColor.from_string(hexstr)


def duplicate_slide_from(prs, src):
    new = prs.slides.add_slide(src.slide_layout)
    for shp in list(new.shapes):
        shp._element.getparent().remove(shp._element)
    for shp in src.shapes:
        new.shapes._spTree.append(copy.deepcopy(shp._element))
    return new


def clean(slide, keep):
    for sh in list(slide.shapes):
        if sh.name not in keep:
            sh._element.getparent().remove(sh._element)


def header(slide, title, subtitle, page):
    apply(slide, {'Text 0': title, 'Text 2': subtitle, 'Text 21': FOOTER, 'Text 22': page})
    for sh in slide.shapes:
        if sh.name == 'Text 2':
            sh.height = Inches(0.8)


def textbox(slide, l, t, w, h, lines, align=None):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (txt, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        if align is not None:
            p.alignment = align
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size); r.font.color.rgb = RGBColor.from_string(color); r.font.bold = bold
    return tb


prs = Presentation(DST)
d1, d2, d3 = prs.slides[0], prs.slides[1], prs.slides[2]

# Slide 1 & 2: content via the two 3-panel skeletons
apply(d1, S1); recolor(d1, 'Text 38', BLUE)
apply(d2, S2)

# Strip d3 chart/cards, then clone the chart-free slide for slides 4-6
clean(d3, D3_KEEP)
s_radar = duplicate_slide_from(prs, d3)
s_result = duplicate_slide_from(prs, d3)
s_enable = duplicate_slide_from(prs, d3)
s_signal = d3

# Slide 3: signal separation (dual-lens CUSUM)
header(s_signal, 'On any given day, the attacker looks normal',
       'Two lenses read the same logs: raw numbers catch the noisy attack first; the AI “meaning” lens catches the subtle ones ~30 weeks sooner. The stealth APT hides from both on drift alone — which is why we fuse multiple fronts.', '3')
s_signal.shapes.add_picture('feature_cusum.png', Inches(0.45), Inches(1.72), width=Inches(5.95))
s_signal.shapes.add_picture('embedding_cusum.png', Inches(6.92), Inches(1.72), width=Inches(5.95))
textbox(s_signal, 0.45, 5.5, 12.4, 0.35,
        [('First clear detection (★):   USR-118 wk 36   ·   USR-156 wk 4   ·   USR-042 wk 15   ·   USR-234 — caught by known-bad profile, not drift', 11, NAVY, True)], align=2)
textbox(s_signal, 0.45, 5.92, 12.4, 0.5,
        [('Each lens has a blind spot. Running them together — plus the cohort and known-bad fronts — is what catches all four.', 12, BODY, False)], align=2)

# Slide 4: radar
header(s_radar, 'No normal user is extreme on every front at once',
       'Each attacker pushes past the normal cluster on a different behavioral phase. That multi-phase extremity is the fingerprint a single threshold can’t see.', '4')
s_radar.shapes.add_picture('radar.png', Inches(0.4), Inches(1.6), width=Inches(5.2))
textbox(s_radar, 6.1, 1.75, 6.7, 4.6, [
    ('Grey = normal users (median & 75th percentile): a tight cluster near the center.', 12.5, BODY, False),
    ('The five behavioral phases:', 13, NAVY, True),
    ('•  Signal Strength — how far behavior departs from peers', 12, BODY, False),
    ('•  Breadth — how many behaviors shifted at once', 12, BODY, False),
    ('•  Sustained Deviation — how long the shift persisted', 12, BODY, False),
    ('•  Context Divergence — how unlike its role-group it became', 12, BODY, False),
    ('•  Novelty Persistence — new connections (e.g., C2) that keep recurring', 12, BODY, False),
    ('USR-156 dominates Signal Strength & Breadth; USR-234 spikes on Novelty Persistence (its C2 beacon). No normal user reaches these extremes on multiple phases at once.', 12, NAVY, False),
])

# Slide 5: result
header(s_result, 'Every attack caught — at the false-positive cost shown',
       'Fused into one score, all four campaigns rank above the line that catches all four. A complementary known-bad detector catches the same four at zero false positives.', '5')
s_result.shapes.add_picture('composite_rank.png', Inches(0.45), Inches(1.85), width=Inches(6.7))
textbox(s_result, 7.45, 1.9, 5.4, 1.4, [
    ('0 → 4 of 4', 40, ORANGE, True),
    ('stealth attacks caught that legacy scored as normal', 13, BODY, False),
])
textbox(s_result, 7.45, 3.7, 5.4, 2.4, [
    ('•  Known-bad profiles: 4 of 4 at 0 false positives', 13, NAVY, True),
    ('•  Composite ranking: 4 of 4 at ~8% false positives', 13, NAVY, True),
    ('•  Subtle attacks flagged up to ~30 weeks earlier', 13, NAVY, True),
    ('•  Every alert names the behavior that fired — auditable', 13, NAVY, True),
])

# Slide 6: what it enables
header(s_enable, 'What this enables for the mission',
       'From a daily flood of alerts to one ranked, explainable verdict — on the logs you already collect.', '6')
textbox(s_enable, 0.7, 1.85, 12.0, 3.2, [
    ('Triage by risk — one risk-ranked list of entities, not thousands of threshold alerts.', 15, NAVY, True),
    ('Explainable verdicts — every flag names the behavior, mapped to MITRE ATT&CK.', 15, NAVY, True),
    ('Slow-and-low coverage — catches insiders and APTs that stay inside every threshold for months.', 15, NAVY, True),
    ('Deployable where you operate — containerized; runs on-premises or in an enclave.', 15, NAVY, True),
])
textbox(s_enable, 0.7, 5.3, 12.0, 1.0, [
    ('Bottom line: catch what rules miss, before it becomes a breach — with evidence a SOC analyst and an auditor can both trust.', 15, ORANGE, True),
])

prs.save(DST)
print('wrote', DST, '-', len(prs.slides), 'slides')
