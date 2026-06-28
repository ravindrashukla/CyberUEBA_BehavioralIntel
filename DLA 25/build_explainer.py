"""Standalone UEBA explainer deck (4 slides), built on the demand-deck template:
  1. Contrast        — legacy single-user rules miss 4/4; behavioral twin catches them
  2. Signal Separation — dual-lens CUSUM (raw vs AI), first-detection stars
  3. Multi-phase radar — attackers extreme on different phases; normals cluster center
  4. Fused composite  — all 5 phases -> one score; 4 of 4 at 8% FP; business outcome
Requires the four PNGs from gen_explainer_figs.py in this folder.
"""
import copy
import os
import shutil
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

SRC = 'Legacy_vs_EDM_Demand_Forecast_Example-1-ver 2.pptx'
DST = 'UEBA_VectorIntelligence_Explainer.pptx'
shutil.copyfile(SRC, DST)
FOOTER = 'CONFIDENTIAL  |  22nd Century Technologies  |  EDM Vector Intelligence — Cyber UEBA'
NAVY, BODY, BLUE, ORANGE, GREY = '102A4C', '1D2430', '1F5E9C', 'F36C21', '5E6A75'
# Each demand slide numbers its footer/page shapes differently.
D2_KEEP = {'Text 0', 'Text 1', 'Text 2', 'Text 46', 'Text 47'}   # demand slide 2
D3_KEEP = {'Text 0', 'Text 1', 'Text 2', 'Text 20', 'Text 21', 'Text 22'}  # demand slide 3

# Slide 1 — contrast (cloned from demand slide 1: INPUT/METHOD/OUTPUT skeleton)
S1 = {
    'Text 0': 'Cyber UEBA: catching what single-user rules miss',
    'Text 2': 'Legacy scores one user against a fixed line and misses all four stealth attacks. V-Intelligence reads the same logs as behavioral entities — and catches every one.',
    'Text 4': 'LEGACY', 'Text 5': 'one user vs a threshold',
    'Text 8': 'Single-user history only', 'Text 11': 'Fixed alert thresholds', 'Text 14': '0 of 4 attacks flagged',
    'Text 17': 'BEHAVIORAL TWIN', 'Text 18': 'same logs, in context',
    'Text 21': 'Cohort-relative deviation', 'Text 24': 'Self-drift over time',
    'Text 27': 'Known-bad profile match', 'Text 30': 'Evidence-backed, explainable',
    'Text 33': 'RESULT', 'Text 34': 'caught, with evidence',
    'Text 35': 'What a user HAS done\nplus where they’re HEADING',
    'Text 36': '4 of 4 caught · 0 false positives',
    'Text 38': 'Why it matters',
    'Text 39': 'Point-anomaly tools saw normal employees; the behavioral twin named each attack by the behavior that gave it away.',
    'Text 40': FOOTER, 'Text 41': '1',
}


def set_text(shape, text):
    tf = shape.text_frame
    runs = [r for p in tf.paragraphs for r in p.runs]
    align = tf.paragraphs[0].alignment
    size = bold = ital = name = color = None
    if runs:
        f = runs[0].font
        size, bold, ital, name = f.size, f.bold, f.italic, f.name
        try:
            if f.color and f.color.type == 1:
                color = str(f.color.rgb)
        except Exception:
            pass
    tf.clear()
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        if size: r.font.size = size
        if bold is not None: r.font.bold = bold
        if ital is not None: r.font.italic = ital
        if name: r.font.name = name
        if color: r.font.color.rgb = RGBColor.from_string(color)


def apply(slide, mapping):
    by_name = {sh.name: sh for sh in slide.shapes}
    for nm, txt in mapping.items():
        if nm in by_name:
            set_text(by_name[nm], txt)


def duplicate_slide_from(prs, src):
    new = prs.slides.add_slide(src.slide_layout)
    for shp in list(new.shapes):
        shp._element.getparent().remove(shp._element)
    for shp in src.shapes:
        new.shapes._spTree.append(copy.deepcopy(shp._element))
    return new


def clean(slide, keep):
    for sh in list(slide.shapes):
        if sh.name not in keep:
            sh._element.getparent().remove(sh._element)


def header(slide, title, subtitle, page, footer_name='Text 21', page_name='Text 22'):
    apply(slide, {'Text 0': title, 'Text 2': subtitle, footer_name: FOOTER, page_name: page})
    for sh in slide.shapes:
        if sh.name == 'Text 2':
            sh.height = Inches(0.8)


def textbox(slide, l, t, w, h, lines, align=None, anchor=None):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    if anchor:
        tf.vertical_anchor = anchor
    for i, (txt, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(3)
        if align is not None:
            p.alignment = align
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size); r.font.color.rgb = RGBColor.from_string(color); r.font.bold = bold
    return tb


prs = Presentation(DST)
# Repurpose all three demand slides (no deletion -> no orphaned-part name clash),
# duplicate only the chart-free d3 for the 4th slide.
d1, d2, d3 = prs.slides[0], prs.slides[1], prs.slides[2]

# ---- Slide 1: contrast ----
apply(d1, S1)
for sh in d1.shapes:
    if sh.name == 'Text 38':
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = RGBColor.from_string(BLUE)

clean(d2, D2_KEEP)                 # signal slide canvas (demand slide 2)
clean(d3, D3_KEEP)                 # radar slide canvas (demand slide 3)
s_cmp = duplicate_slide_from(prs, d3)   # slide 4 (copy of chart-free d3)
s_sig, s_rad = d2, d3

# ---- Slide 2: signal separation ----
header(s_sig, 'Signal separation: two lenses, different attacks caught first',
       'Over the same 70-week campaign, the raw-numbers lens flags the noisy volume attack earliest; the AI “meaning” lens flags the subtle insider & stealth-hacker ~30 weeks sooner. Neither catches the slow APT alone.',
       '2', footer_name='Text 46', page_name='Text 47')
s_sig.shapes.add_picture('feature_cusum.png', Inches(0.45), Inches(1.72), width=Inches(5.95))
s_sig.shapes.add_picture('embedding_cusum.png', Inches(6.92), Inches(1.72), width=Inches(5.95))
textbox(s_sig, 0.45, 5.5, 12.4, 0.35,
        [('First sustained detection (★):    156 → feature wk 39 / AI wk 4      ·      042 → wk 47 / wk 15      ·      118 → wk 36 / wk 60      ·      234 → not flagged by drift alone', 11, NAVY, True)],
        align=2)
textbox(s_sig, 0.45, 5.92, 12.4, 0.5,
        [('Each lens has a blind spot — so V-Intelligence runs both and fuses them with the cohort and known-bad fronts.', 12, BODY, False)],
        align=2)

# ---- Slide 3: radar ----
header(s_rad, 'No normal user is extreme on every front at once',
       'Each attacker pushes past the normal cluster on a different phase — that multi-phase extremity is the fingerprint a single threshold can’t see.', '3')
s_rad.shapes.add_picture('radar.png', Inches(0.4), Inches(1.6), width=Inches(5.2))
textbox(s_rad, 6.1, 1.75, 6.7, 4.6, [
    ('Grey = normal users (median & 75th percentile): a tight cluster near the center.', 12, BODY, False),
    ('The five phases, in plain terms:', 12.5, NAVY, True),
    ('•  Signal Strength — how far behavior departs from peers', 11.5, BODY, False),
    ('•  Breadth — how many behaviors shifted at once', 11.5, BODY, False),
    ('•  Sustained Deviation — how long the shift persisted', 11.5, BODY, False),
    ('•  Context Divergence — how unlike its role-group it became', 11.5, BODY, False),
    ('•  Novelty Persistence — new connections (e.g., C2) that keep recurring', 11.5, BODY, False),
    ('USR-156 dominates Signal Strength & Breadth; USR-234 spikes on Novelty Persistence (its C2 beacon); USR-118 stretches Context Divergence. No normal user reaches these extremes on multiple phases at once.', 11.5, NAVY, False),
])

# ---- Slide 4: fused composite + outcome ----
header(s_cmp, 'Fused into one score: 4 of 4, at 8% false positives',
       'The five phases combine into a single composite — every attacker rises above the line that catches all four, surfacing the two stealth movers legacy scored as normal.', '4')
s_cmp.shapes.add_picture('composite_rank.png', Inches(0.45), Inches(1.85), width=Inches(6.7))
textbox(s_cmp, 7.45, 1.9, 5.4, 1.4, [
    ('0 → 4 of 4', 40, ORANGE, True),
    ('stealth attacks caught that legacy scored as normal', 13, BODY, False),
])
textbox(s_cmp, 7.45, 3.7, 5.4, 2.4, [
    ('•  39 legacy false alarms → 0 on the threat-profile match', 13, NAVY, True),
    ('•  Subtle attacks flagged up to ~30 weeks earlier', 13, NAVY, True),
    ('•  Every alert names the behavior that fired — auditable', 13, NAVY, True),
])

prs.save(DST)
print('wrote', DST, '-', len(prs.slides), 'slides')
