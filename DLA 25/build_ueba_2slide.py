"""Compressed 2-slide UEBA use case.
Slide 1: 'Same logs, two outcomes' — Legacy -> Behavioral Twin -> Result (clones
         the demand deck's slide-1 INPUT/METHOD/OUTPUT skeleton).
Slide 2: the proof — chart + hero number + benefit cards (the demand deck's
         slide-3 skeleton, with its existing chart reused).
The middle 'digital twin' slide is dropped. Numbers are the project's verified
figures (legacy 0/4 caught + 39/246 false-flagged; behavioral 4/4 at 0 FP).
"""
import os
import shutil
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt, Emu

SRC = 'Legacy_vs_EDM_Demand_Forecast_Example-1-ver 2.pptx'
DST = 'Legacy_vs_EDM_UEBA_Detection_2slide.pptx'
shutil.copyfile(SRC, DST)
FOOTER = 'CONFIDENTIAL  |  22nd Century Technologies  |  EDM Vector Intelligence — Cyber UEBA'

# Slide 1 (clone of demand slide 1): LEGACY -> BEHAVIORAL TWIN -> RESULT
S1 = {
    'Text 0': 'Same logs, two outcomes',
    'Text 2': 'Legacy scores one user against a fixed line. Behavioral entities read the same logs in context — and catch all four stealth attacks.',
    'Text 4': 'LEGACY',
    'Text 5': 'one user vs a threshold',
    'Text 8': 'Single-user history only',
    'Text 11': 'Fixed alert thresholds',
    'Text 14': '0 of 4 attacks flagged',
    'Text 17': 'BEHAVIORAL TWIN',
    'Text 18': 'same logs, in context',
    'Text 21': 'Cohort-relative deviation',
    'Text 24': 'Self-drift over time',
    'Text 27': 'Known-bad profile match',
    'Text 30': 'Evidence-backed, explainable',
    'Text 33': 'RESULT',
    'Text 34': 'caught, with evidence',
    'Text 35': 'What a user HAS done\nplus where they’re HEADING',
    'Text 36': '4 of 4 caught · 0 false positives',
    'Text 38': 'Why it matters',
    'Text 39': 'Point-anomaly tools saw normal employees; the behavioral twin named each attack by the behavior that gave it away.',
    'Text 40': FOOTER,
    'Text 41': '1',
}

# Slide 2 (the former slide 3 skeleton): two analytics charts as images.
# The native chart + right-side cards are deleted; we keep title/subtitle/footer.
S2 = {
    'Text 0': 'Vector Intelligence: per-week → cumulative drift → composite score',
    'Text 2': 'Per-week, an attacker resembles a common user; accumulating the signal reveals persistent drift; the composite fuses every front — 4 of 4 caught.',
    'Text 20': 'Synthetic / representative demo data — to be revalidated on customer SOC data before production claims.',
    'Text 21': FOOTER,
    'Text 22': '2',
}
# Shapes from the source slide 3 to remove (native chart + the 3 benefit cards)
S2_DELETE = ['Shape 3', 'Text 4', 'Chart 0',
             'Shape 5', 'Text 6', 'Text 7',
             'Shape 8', 'Shape 9', 'Text 10', 'Text 11',
             'Shape 12', 'Shape 13', 'Text 14', 'Text 15',
             'Shape 16', 'Shape 17', 'Text 18', 'Text 19']


def set_text(shape, text):
    tf = shape.text_frame
    runs = [r for p in tf.paragraphs for r in p.runs]
    align = tf.paragraphs[0].alignment
    size = bold = ital = name = color = None
    if runs:
        f = runs[0].font
        size, bold, ital, name = f.size, f.bold, f.italic, f.name
        try:
            if f.color and f.color.type == 1:
                color = str(f.color.rgb)
        except Exception:
            pass
    tf.clear()
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        if size: r.font.size = size
        if bold is not None: r.font.bold = bold
        if ital is not None: r.font.italic = ital
        if name: r.font.name = name
        if color: r.font.color.rgb = RGBColor.from_string(color)


def apply(slide, mapping):
    by_name = {sh.name: sh for sh in slide.shapes}
    for nm, txt in mapping.items():
        if nm not in by_name:
            raise SystemExit('MISSING SHAPE %r' % nm)
        set_text(by_name[nm], txt)


prs = Presentation(DST)
# Drop the middle 'digital twin' slide -> leaves [slide1 skeleton, slide3 skeleton]
sldIdLst = prs.slides._sldIdLst
ids = list(sldIdLst)
sldIdLst.remove(ids[1])

s_legacy, s_proof = prs.slides[0], prs.slides[1]
apply(s_legacy, S1)
# recolor bottom-bar header red -> brand blue (positive note here)
for sh in s_legacy.shapes:
    if sh.name == 'Text 38':
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = RGBColor.from_string('1F5E9C')

apply(s_proof, S2)
# Remove the native chart + benefit cards, then drop in the two analytics images.
for sh in list(s_proof.shapes):
    if sh.name in S2_DELETE:
        sh._element.getparent().remove(sh._element)

IMGS = ['per_week_drift.png', 'cumulative_drift.png', 'composite_score.png']
for fn in IMGS:
    if not os.path.exists(fn):
        raise SystemExit('Missing %s — run generate_figs.py first (with DB env).' % fn)

IMG_TOP, IMG_W = Inches(1.72), Inches(4.05)
for i, fn in enumerate(IMGS):
    s_proof.shapes.add_picture(fn, Inches(0.35 + i * 4.29), IMG_TOP, width=IMG_W)

# Bottom takeaway line
tb = s_proof.shapes.add_textbox(Inches(0.35), Inches(5.15), Inches(12.6), Inches(0.8))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run()
r.text = ('On any given week the stealth APT (USR-234) sits inside the normal range 97% of the time. '
          'Accumulating the signal — then fusing every front into one composite score — separates all '
          'four attackers from 250 users, at 8% false positives.')
r.font.size = Pt(12); r.font.color.rgb = RGBColor.from_string('1D2430')

prs.save(DST)
print('wrote', DST, '-', len(prs.slides), 'slides')
