"""Build the UEBA 'Legacy vs EDM' use-case deck by cloning the demand-forecast
deck and swapping text per shape name (preserves all colors/gradients/layout).
Numbers are the project's verified figures:
  legacy point-anomaly (IF/SVM/LOF): 0/4 caught, 39/246 normals false-flagged
  EDM behavioral (threat-profile):   4/4 caught, 0 false flags
"""
import copy
import shutil
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

SRC = 'Legacy_vs_EDM_Demand_Forecast_Example-1-ver 2.pptx'
DST = 'Legacy_vs_EDM_UEBA_Detection_Example-ver2.pptx'
shutil.copyfile(SRC, DST)

FOOTER = 'CONFIDENTIAL  |  22nd Century Technologies  |  EDM Vector Intelligence — Cyber UEBA'

# Per-slide {shape_name: new_text}. Shapes not listed are left untouched.
S1 = {
    'Text 0': 'Legacy detection: a transparent single-user baseline',
    'Text 2': 'Useful and explainable — yet it only sees this user’s own activity against a fixed threshold.',
    'Text 5': 'one user, fixed rules',
    'Text 8': 'This user’s sign-in history',
    'Text 11': '90-day window',
    'Text 14': 'Static alert thresholds',
    'Text 18': 'count, threshold, flag',
    'Text 21': 'Failed-login counters',
    'Text 24': 'Signature / rule match',
    'Text 27': 'Fixed anomaly thresholds',
    'Text 30': 'Point-anomaly scoring',
    'Text 34': 'an alert, or silence',
    'Text 35': 'Per-user rule hits\nBinary alert / no-alert',
    'Text 36': '4 STEALTH ATTACKS · 0 of 4 flagged',
    'Text 39': 'No peer/cohort context   ·   slow attacks stay under thresholds   ·   thresholds are fixed',
    'Text 40': FOOTER,
}
S2 = {
    'Text 0': 'Same logs, re-imagined as living behavioral entities',
    'Text 2': 'The identical records become connected, time-aware entities — ready for GenAI-grade detection.',
    'Text 8': 'Same sign-in & access logs',
    'Text 11': 'Network & endpoint records',
    'Text 14': 'User & device context',
    'Text 17': 'BEHAVIORAL TWIN',
    'Text 22': 'Learns from peers (its cohort)',
    'Text 26': 'Matches known-bad profiles',
    'Text 29': 'DETECTION',
    'Text 33': 'Behavioral risk model',
    'Text 36': '+ cohort drift signal',
    'Text 39': 'Calibrated, evidence-backed',
    'Text 42': 'Catches low-and-slow attacks',
    'Text 45': 'Same data, a richer model — it captures what a user HAS done plus where they’re HEADING.',
    'Text 46': FOOTER,
}
S3 = {
    'Text 2': 'On the same logs, Vector Intelligence catches all 4 stealth attacks — earlier warning, fewer false alarms.',
    'Text 4': 'Detection on the same logs — Legacy vs EDM',
    'Text 6': '0 → 4 of 4',
    'Text 7': 'stealth attacks caught that legacy scored as normal',
    'Text 11': 'Catches a slow breach before any single rule trips.',
    'Text 14': 'Right-sized alerts',
    'Text 15': '39 legacy false alarms → 0 on the profile match.',
    'Text 19': 'Shows which behavior fired — cohort, drift, known-bad profile.',
    'Text 20': 'Synthetic / representative demo data — to be revalidated on customer SOC data before production claims.',
    'Text 21': FOOTER,
}


# Slide 4 — how it works (cloned from slide 1's INPUT/METHOD/OUTPUT skeleton)
S4 = {
    'Text 0': 'How the Behavioral Twin catches what rules miss',
    'Text 2': 'Three complementary fronts turn the same logs into named, evidence-backed detections.',
    'Text 4': 'THREE FRONTS',
    'Text 5': 'complementary views',
    'Text 8': 'Cohort-relative deviation',
    'Text 11': 'Self-drift over time',
    'Text 14': 'Raw-event fingerprints',
    'Text 17': 'KNOWN-BAD PROFILES',
    'Text 18': 'measurable fingerprints',
    'Text 21': 'C2-beacon rhythm',
    'Text 24': 'DGA-DNS domains',
    'Text 27': 'Cohort-rare destinations',
    'Text 30': 'Mass data collection',
    'Text 33': 'NAMED ALERT',
    'Text 34': 'technique + evidence',
    'Text 35': 'Which behavior fired\nWhy it’s abnormal vs peers',
    'Text 36': '4 of 4 caught · 0 false positives',
    'Text 38': 'Why it holds up',
    'Text 39': 'Each flag is a named technique with raw evidence — not an unexplained anomaly score.',
    'Text 40': FOOTER,
    'Text 41': '4',
}

# Slide 5 — the four attacks, each caught by name (cloned from slide 1)
S5 = {
    'Text 0': 'The four stealth attacks — each caught by name',
    'Text 2': 'Same logs legacy scored as normal; every campaign matched a measurable known-bad profile.',
    'Text 4': 'LEGACY SAW',
    'Text 5': 'point-anomaly tools',
    'Text 8': 'USR-156, USR-234:  normal',
    'Text 11': 'USR-042, USR-118:  normal',
    'Text 14': '0 of 4 flagged',
    'Text 17': 'PROFILE THAT FIRED',
    'Text 18': 'named technique match',
    'Text 21': 'USR-156 · mass collection + rare dst',
    'Text 24': 'USR-234 · C2 beacon + DGA domains',
    'Text 27': 'USR-042 · living-off-the-land',
    'Text 30': 'USR-118 · recon fan-out',
    'Text 33': 'RESULT',
    'Text 34': 'named & evidence-backed',
    'Text 35': 'Every campaign matched a\nmeasurable known-bad profile',
    'Text 36': '4 of 4 caught · 0 false positives',
    'Text 38': 'Bottom line',
    'Text 39': 'Each attack is named by the behavior that gave it away — auditable, not a black-box score.',
    'Text 40': FOOTER,
    'Text 41': '5',
}


def duplicate_slide_from(prs, src_slide):
    new = prs.slides.add_slide(src_slide.slide_layout)
    for shp in list(new.shapes):
        shp._element.getparent().remove(shp._element)
    for shp in src_slide.shapes:
        new.shapes._spTree.append(copy.deepcopy(shp._element))
    return new


def recolor(shape, hexstr):
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            r.font.color.rgb = RGBColor.from_string(hexstr)


def set_text(shape, text):
    tf = shape.text_frame
    runs = [r for p in tf.paragraphs for r in p.runs]
    align = tf.paragraphs[0].alignment
    size = bold = ital = name = color = None
    if runs:
        f = runs[0].font
        size, bold, ital, name = f.size, f.bold, f.italic, f.name
        try:
            if f.color and f.color.type == 1:
                color = str(f.color.rgb)
        except Exception:
            pass
    tf.clear()
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        if size: r.font.size = size
        if bold is not None: r.font.bold = bold
        if ital is not None: r.font.italic = ital
        if name: r.font.name = name
        if color: r.font.color.rgb = RGBColor.from_string(color)


prs = Presentation(DST)
for slide, mapping in zip(prs.slides, [S1, S2, S3]):
    by_name = {sh.name: sh for sh in slide.shapes}
    for nm, txt in mapping.items():
        if nm not in by_name:
            raise SystemExit('MISSING SHAPE %r' % nm)
        set_text(by_name[nm], txt)
    # chart only on slide 3
    for sh in slide.shapes:
        if sh.has_chart:
            cd = CategoryChartData()
            cd.categories = ['Attacks caught (of 4)', 'Stealth attacks missed', 'Normal users false-flagged']
            cd.add_series('Legacy', (0, 4, 39))
            cd.add_series('EDM Vector Intelligence', (4, 0, 0))
            sh.chart.replace_data(cd)

# Slides 4 & 5: clone slide 1's skeleton, swap text, recolor the bottom-bar header
# from alarm-red to brand-blue (it's a positive note on these slides).
base = prs.slides[0]
for mapping in [S4, S5]:
    new = duplicate_slide_from(prs, base)
    by_name = {sh.name: sh for sh in new.shapes}
    for nm, txt in mapping.items():
        if nm not in by_name:
            raise SystemExit('MISSING SHAPE %r on cloned slide' % nm)
        set_text(by_name[nm], txt)
    recolor(by_name['Text 38'], '1F5E9C')

prs.save(DST)
print('wrote', DST, '-', len(prs.slides), 'slides')
