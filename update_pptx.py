"""Update ACECARD PowerPoint decks with v3.0 composite scoring results."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

NAVY = RGBColor(0x0D, 0x1B, 0x2A)
BLUE = RGBColor(0x1B, 0x4F, 0x72)
GOLD = RGBColor(0xB7, 0x95, 0x0B)
RED = RGBColor(0xC0, 0x39, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x6C, 0x75, 0x7D)
GREEN = RGBColor(0x27, 0xAE, 0x60)


def add_slide_from_layout(prs, layout_idx=1):
    """Add a new slide using the specified layout."""
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)


def add_text_box(slide, left, top, width, height, text, font_size=12,
                 bold=False, color=NAVY, alignment=PP_ALIGN.LEFT):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tf


def add_paragraph(text_frame, text, font_size=11, bold=False, color=NAVY,
                  alignment=PP_ALIGN.LEFT, space_before=Pt(4)):
    """Add a paragraph to an existing text frame."""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    p.space_before = space_before
    return p


def update_deep_dive():
    """Update the Technical Deep Dive deck."""
    path = r'c:\Users\shuklar\ClaudeCode\CyberUEBA_BehavioralIntel\docs\ACECARD_Technical_Deep_Dive.pptx'
    prs = Presentation(path)

    # Add new slide: Composite Scoring Results (after slide 14, before Key Findings)
    slide = add_slide_from_layout(prs, 1)

    # Title
    tf = add_text_box(slide, 0.5, 0.3, 9, 0.6,
                      "v3.0: Multi-Phase Composite Scoring",
                      font_size=24, bold=True, color=NAVY)

    # Subtitle
    add_text_box(slide, 0.5, 0.85, 9, 0.4,
                 "5-phase detection system replaces single-method ensembles. All 4 attacks detected at 8.5% FP.",
                 font_size=12, color=GRAY)

    # Left column — Phase descriptions
    tf = add_text_box(slide, 0.5, 1.5, 4.5, 4.5, "", font_size=11)
    tf.paragraphs[0].text = "Five Detection Phases"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = BLUE

    phases = [
        ("Phase 1: Signal Strength", "Top-3 group-relative z-scores"),
        ("Phase 2: Signal Breadth", "Count of features with z > 1.5"),
        ("Phase 3: Sustained Deviation", "Late-period zone drift z-scores"),
        ("Phase 4: Context Divergence", "Cross-context composite spread"),
        ("Phase 5: Novelty Persistence", "C2 beacon via persistent novel IPs"),
    ]
    for phase_name, desc in phases:
        add_paragraph(tf, phase_name, font_size=11, bold=True, color=NAVY, space_before=Pt(8))
        add_paragraph(tf, desc, font_size=10, color=GRAY, space_before=Pt(2))

    add_paragraph(tf, "", font_size=8, space_before=Pt(12))
    add_paragraph(tf, "Group-relative scoring compares each user to their role group",
                  font_size=10, bold=True, color=BLUE, space_before=Pt(4))
    add_paragraph(tf, "(admin, security, developer, business, executive)",
                  font_size=9, color=GRAY, space_before=Pt(2))

    # Right column — Results
    tf2 = add_text_box(slide, 5.3, 1.5, 4.2, 4.5, "", font_size=11)
    tf2.paragraphs[0].text = "Detection Results (250 Users)"
    tf2.paragraphs[0].font.size = Pt(14)
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.color.rgb = BLUE

    results = [
        ("USR-118 (Salt Typhoon)", "Rank #1  |  Score: 51.3"),
        ("USR-156 (Insider)", "Rank #2  |  Score: 46.2"),
        ("USR-234 (Slow APT)", "Rank #7  |  Score: 19.4  (Novelty: 13.0)"),
        ("USR-042 (Volt Typhoon)", "Rank #24  |  Score: 13.7"),
    ]
    for name, detail in results:
        add_paragraph(tf2, name, font_size=11, bold=True, color=NAVY, space_before=Pt(8))
        add_paragraph(tf2, detail, font_size=10, color=GRAY, space_before=Pt(2))

    add_paragraph(tf2, "", font_size=8, space_before=Pt(12))
    add_paragraph(tf2, "4/4 Attacks Detected", font_size=16, bold=True, color=GREEN, space_before=Pt(4))
    add_paragraph(tf2, "8.5% False Positive Rate (21 FP / 246 normal)",
                  font_size=11, bold=True, color=NAVY, space_before=Pt(4))
    add_paragraph(tf2, "100% Recall at 90th percentile threshold",
                  font_size=10, color=GRAY, space_before=Pt(2))

    # Footer
    add_text_box(slide, 0.5, 6.8, 9, 0.3,
                 "22nd Century Technologies, Inc.  |  ACECARD v3.0  |  Confidential",
                 font_size=8, color=GRAY, alignment=PP_ALIGN.CENTER)

    # Add second new slide: Novelty Persistence
    slide2 = add_slide_from_layout(prs, 1)

    tf = add_text_box(slide2, 0.5, 0.3, 9, 0.6,
                      "Novelty Persistence: Detecting What Embeddings Cannot",
                      font_size=22, bold=True, color=NAVY)

    add_text_box(slide2, 0.5, 0.85, 9, 0.4,
                 "LLM embeddings treat IP addresses as generic tokens. C2 beacons require direct feature detection.",
                 font_size=12, color=GRAY)

    # The problem
    tf = add_text_box(slide2, 0.5, 1.5, 4.5, 2.5, "", font_size=11)
    tf.paragraphs[0].text = "The Embedding Limitation"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RED

    add_paragraph(tf, "USR-234 (Slow APT) has a C2 beacon IP present in qualitative features for 60 consecutive weeks.",
                  font_size=10, color=NAVY, space_before=Pt(8))
    add_paragraph(tf, "But embedding models don't understand IP semantics — 198.51.100.47 produces the same type of vector as any other IP.",
                  font_size=10, color=NAVY, space_before=Pt(4))
    add_paragraph(tf, "Result: USR-234 was rank #99/250 with embedding-only detection.",
                  font_size=10, bold=True, color=RED, space_before=Pt(8))

    # The solution
    tf2 = add_text_box(slide2, 5.3, 1.5, 4.2, 2.5, "", font_size=11)
    tf2.paragraphs[0].text = "The Fix: Direct Feature Extraction"
    tf2.paragraphs[0].font.size = Pt(14)
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.color.rgb = GREEN

    add_paragraph(tf2, "Track novel external IPs per user over time.",
                  font_size=10, color=NAVY, space_before=Pt(8))
    add_paragraph(tf2, "A novel IP persisting for 60/60 post-baseline weeks is unambiguously suspicious.",
                  font_size=10, color=NAVY, space_before=Pt(4))
    add_paragraph(tf2, "Normal users: <5 weeks persistence average.",
                  font_size=10, color=NAVY, space_before=Pt(4))
    add_paragraph(tf2, "Result: USR-234 jumps from rank #99 to rank #7.",
                  font_size=10, bold=True, color=GREEN, space_before=Pt(8))

    # Key insight
    tf3 = add_text_box(slide2, 0.5, 4.3, 9, 2, "", font_size=11)
    tf3.paragraphs[0].text = "Architectural Principle"
    tf3.paragraphs[0].font.size = Pt(14)
    tf3.paragraphs[0].font.bold = True
    tf3.paragraphs[0].font.color.rgb = BLUE

    add_paragraph(tf3, "Hybrid detection combines embedding-derived signals (behavioral drift, context divergence) with direct numeric features (novelty persistence).",
                  font_size=11, color=NAVY, space_before=Pt(8))
    add_paragraph(tf3, "Neither approach alone achieves 4/4 detection. The composite scorer fuses both into a single ranked score.",
                  font_size=11, color=NAVY, space_before=Pt(4))
    add_paragraph(tf3, "This principle applies broadly: when the embedding model has known blind spots, augment with direct feature extraction.",
                  font_size=10, color=GRAY, space_before=Pt(4))

    add_text_box(slide2, 0.5, 6.8, 9, 0.3,
                 "22nd Century Technologies, Inc.  |  ACECARD v3.0  |  Confidential",
                 font_size=8, color=GRAY, alignment=PP_ALIGN.CENTER)

    # Update slide 18 (Key Findings) — now at index 19 after 2 insertions
    # Find and update the Key Findings slide
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if "Key Findings" in para.text and i >= 15:
                        # Update Finding 4
                        for shape2 in slide.shapes:
                            if shape2.has_text_frame:
                                for p in shape2.text_frame.paragraphs:
                                    if "LOF + Zone Divergence = Optimal" in p.text:
                                        p.text = "v3.0 Composite Scoring: 4/4 Attacks at 8.5% FP"
                                    elif "This combination detects ALL 4" in p.text:
                                        p.text = "Multi-phase composite scoring combines signal strength, breadth, sustained deviation, context divergence, and novelty persistence into a single ranked score. Detects all 4 campaigns including the Slow APT (previously rank #99) via C2 beacon novelty persistence."
                        break

    prs.save(path)
    print(f"Updated: {path}")
    print(f"  Added 2 new slides (Composite Scoring, Novelty Persistence)")
    print(f"  Total slides: {len(prs.slides)}")


def update_traditional_vs_behavioral():
    """Update the Traditional vs Behavioral deck."""
    path = r'c:\Users\shuklar\ClaudeCode\CyberUEBA_BehavioralIntel\docs\ACECARD_Traditional_vs_Behavioral_Deck.pptx'
    prs = Presentation(path)

    # Add new slide before "Results Summary" (slide 14)
    slide = add_slide_from_layout(prs, 1)

    tf = add_text_box(slide, 0.5, 0.3, 9, 0.6,
                      "v3.0 Composite Detection: The Full Picture",
                      font_size=24, bold=True, color=NAVY)

    add_text_box(slide, 0.5, 0.85, 9, 0.4,
                 "Multi-phase scoring unifies traditional and behavioral detection into a single ranked output.",
                 font_size=12, color=GRAY)

    # Before vs After
    tf = add_text_box(slide, 0.5, 1.5, 4.3, 3.5, "", font_size=11)
    tf.paragraphs[0].text = "Before (v1.0)"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RED

    before_items = [
        "17 individual detection methods",
        "No single method catches all 4 attacks",
        "Best ensemble (LOF + Zone Divergence): 4/4 at 6.5% FP",
        "But requires manual method selection per threat type",
        "Slow APT (USR-234): rank #99 with z-score methods",
    ]
    for item in before_items:
        add_paragraph(tf, f"• {item}", font_size=10, color=NAVY, space_before=Pt(6))

    tf2 = add_text_box(slide, 5.3, 1.5, 4.2, 3.5, "", font_size=11)
    tf2.paragraphs[0].text = "After (v3.0)"
    tf2.paragraphs[0].font.size = Pt(14)
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.color.rgb = GREEN

    after_items = [
        "Single composite score per user",
        "5 detection phases auto-fused",
        "4/4 attacks detected at 8.5% FP",
        "No manual method selection needed",
        "Slow APT (USR-234): rank #7 via novelty persistence",
    ]
    for item in after_items:
        add_paragraph(tf2, f"• {item}", font_size=10, color=NAVY, space_before=Pt(6))

    # Bottom banner
    tf3 = add_text_box(slide, 0.5, 5.3, 9, 1.2, "", font_size=11)
    tf3.paragraphs[0].text = "From 17 Methods to 1 Score"
    tf3.paragraphs[0].font.size = Pt(16)
    tf3.paragraphs[0].font.bold = True
    tf3.paragraphs[0].font.color.rgb = GOLD
    tf3.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_paragraph(tf3, "The composite scorer replaces the need to manually select algorithms per threat type. "
                  "Analysts receive a single ranked list of users by behavioral anomaly, with phase decomposition "
                  "explaining WHY each user scored high.",
                  font_size=10, color=NAVY, alignment=PP_ALIGN.CENTER, space_before=Pt(8))

    add_text_box(slide, 0.5, 6.8, 9, 0.3,
                 "22nd Century Technologies, Inc.  |  ACECARD v3.0  |  Confidential",
                 font_size=8, color=GRAY, alignment=PP_ALIGN.CENTER)

    prs.save(path)
    print(f"Updated: {path}")
    print(f"  Added 1 new slide (Composite Detection)")
    print(f"  Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    update_deep_dive()
    update_traditional_vs_behavioral()
    print("\nDone. Both decks updated.")
