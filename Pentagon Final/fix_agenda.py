# -*- coding: utf-8 -*-
"""Restyle the Agenda slide (slide 2) to match the deck's content slides (6-10):
Segoe UI, navy #102A4C headings, #222222 body, gray #5B6470 breaks/footer, orange rule #E87722.
Content unchanged - look and feel only. No em/en dashes."""
import os, shutil
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.oxml.ns import qn

DECK=r"Pentagon - 22nd Century Technologies - V1.0.pptx"
BAK=r"C:/Users/shuklar/AppData/Local/Temp/claude/c--Users-shuklar-ClaudeCode-CyberUEBA-BehavioralIntel-Enhanced/a6818674-bde5-4462-be21-ffe7ac102bf2/scratchpad/Pentagon_V1.0_prestyle.pptx"
shutil.copyfile(DECK,BAK)

NAVY=RGBColor(0x10,0x2A,0x4C); INK=RGBColor(0x22,0x22,0x22); GRAY=RGBColor(0x5B,0x64,0x70)
ORANGE=RGBColor(0xE8,0x77,0x22); TOPBAR=RGBColor(0x00,0x21,0x40); FOOTRULE=RGBColor(0xC9,0xD3,0xDF)
FONT="Segoe UI"; EMU=914400
FOOTER="CONFIDENTIAL  |  22nd Century Technologies  |  Defense in Depth for the Age of AI-Enabled Attacks  ·  UNCLASSIFIED // FOUO"

def solid(sh,rgb):
    try:
        sh.fill.solid(); sh.fill.fore_color.rgb=rgb
    except Exception: pass
    try:
        sh.line.color.rgb=rgb
    except Exception: pass
def set_indent(p, marL_in=0.0, indent_in=0.0, nobul=True):
    pPr=p._p.get_or_add_pPr()
    pPr.set('marL',str(int(marL_in*EMU))); pPr.set('indent',str(int(indent_in*EMU)))
    if nobul:
        for tag in ('a:buChar','a:buAutoNum','a:buNone'):
            for e in pPr.findall(qn(tag)): pPr.remove(e)
        pPr.append(pPr.makeelement(qn('a:buNone'),{}))
def run(p,text,size=11,bold=False,italic=False,color=INK):
    r=p.add_run(); r.text=text; r.font.size=Pt(size); r.font.bold=bold; r.font.italic=italic
    r.font.name=FONT; r.font.color.rgb=color; return r

prs=Presentation(DECK)
s=prs.slides[1]
sh=s.shapes

# --- title ---
tp=sh[0].text_frame.paragraphs[0]
for r in tp.runs: r.font.name=FONT; r.font.bold=True; r.font.italic=False; r.font.color.rgb=NAVY
# --- orange rule / footer rule / top bar ---
solid(sh[1],ORANGE)      # title underline rule
solid(sh[2],FOOTRULE)    # footer rule
solid(sh[4],TOPBAR)      # top bar
# --- footer text ---
ftf=sh[3].text_frame; ftf.clear()
fp=ftf.paragraphs[0]; run(fp,FOOTER,size=8,color=GRAY)
# --- page number ---
for r in sh[5].text_frame.paragraphs[0].runs: r.font.name=FONT; r.font.color.rgb=GRAY

# --- body rebuild (content identical, deck styling) ---
tf=sh[6].text_frame
try: tf.auto_size=MSO_AUTO_SIZE.NONE
except Exception: pass
tf.word_wrap=True; tf.clear()
_first=[True]
def newp(before=0,after=2,align=PP_ALIGN.LEFT,marL=0.0,indent=0.0):
    p=tf.paragraphs[0] if _first[0] else tf.add_paragraph(); _first[0]=False
    p.alignment=align; p.space_before=Pt(before); p.space_after=Pt(after); set_indent(p,marL,indent); return p
def head(text,before=9):
    p=newp(before=before,after=3); run(p,text,size=12.5,bold=True,color=NAVY)
def item(num,text,before=0):
    p=newp(before=before,after=2,marL=0.32,indent=-0.32)
    run(p,num+"   ",size=11,color=INK); run(p,text,size=11,color=INK)
def layer(label,desc,before=1):
    p=newp(before=before,after=3,marL=0.32,indent=-0.32)
    run(p,label+"  ",size=11,bold=True,color=NAVY); run(p,desc,size=11,color=INK)
def brk(text):
    p=newp(before=7,after=7,align=PP_ALIGN.CENTER); run(p,text,size=11,bold=True,color=GRAY)

head("Part 0  -  Introduction",before=0)
item("1.","22nd Century - offerings and capabilities (15-20 min)")
head("Part 1  -  Problem statement")
item("1.","Current challenges: AI-enabled threats, insider threats, and the shortcomings of detect-and-respond (2-3 min)")
item("2.","Today's tools leave a gap at every layer - network, behavioral / UEBA, app / API / data, code, supply chain (2-3 min)")
item("3.","No single tool covers this; a multi-layered platform is the answer (2-3 min)")
brk("Break")
head("Part 2  -  Technical deep dive (innovation & demo)")
layer("Layer 1  -  Preemptive network defense.","Respond to AI-enabled threats at the network layer before they enter (50 min)")
layer("Layer 2  -  Behavioral digital twin.","User and entity anomaly detection - meaning over magnitude, two-layer algorithms (30-40 min)")
brk("Lunch")
layer("Layer 3  -  App / API / data protection.","Defend against bots and agents; protect sensitive data and applications at runtime (45 min)")
layer("Layer 4  -  Code vulnerability.","Novel zero-day discovery pre-deploy (1 slide, demo, 5-7 min)")
layer("Layer 5  -  Supply-chain assurance and containment.","1 slide, demo (5-7 min)")
head("Part 3  -  Next steps & implementation roadmap")
item("1.","MVP pilot on real Army telemetry, ROM, and contract vehicles (BD - Reddy input)")
item("2.","Implementation steps (open for discussion based on scope)")

try:
    prs.save(DECK); print("saved",os.path.abspath(DECK))
except PermissionError:
    alt=r"C:/Users/shuklar/AppData/Local/Temp/claude/c--Users-shuklar-ClaudeCode-CyberUEBA-BehavioralIntel-Enhanced/a6818674-bde5-4462-be21-ffe7ac102bf2/scratchpad/Pentagon_V1.0_agenda_styled.pptx"
    prs.save(alt); print("LOCKED - saved fallback",alt)
