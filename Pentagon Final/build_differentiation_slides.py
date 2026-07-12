# -*- coding: utf-8 -*-
"""Technical gap-analysis slides: overview + one 'gap / how we cover / boundary' slide per layer.
Enriched to the depth of the Word write-up: detailed gap, mechanism steps, and a boundary/evidence strip.
Hardcore-technical, not sales. Navy Mythos. No em/en dashes."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY=RGBColor(0x0E,0x28,0x41); NAVY2=RGBColor(0x1F,0x5E,0x9C); BLUE=RGBColor(0x15,0x60,0x82)
GBLUE=RGBColor(0x2F,0x5F,0xA0); MAGENTA=RGBColor(0xA0,0x2B,0x93); GREEN=RGBColor(0x1B,0x6E,0x2E)
ORANGE=RGBColor(0xE8,0x87,0x1E); AMBER=RGBColor(0xB9,0x6A,0x12); SLATE=RGBColor(0x5E,0x6C,0x7A)
STEEL=RGBColor(0x1F,0x5E,0x7A); RISK=RGBColor(0xA5,0x34,0x2F)
INK=RGBColor(0x24,0x2C,0x38); WHITE=RGBColor(0xff,0xff,0xff); BORDER=RGBColor(0xC4,0xCF,0xDB)
LGRAY=RGBColor(0xF3,0xF5,0xF8); LNAVY=RGBColor(0xE9,0xF1,0xFA); STRIP=RGBColor(0xED,0xF0,0xF4)
FOOT=RGBColor(0x5E,0x6C,0x7A)
FT="Aptos"; FTH="Aptos Display"; FOOTER="CONFIDENTIAL   |   22nd Century Technologies   |   Technical deep-dive   ·   UNCLASSIFIED // FOUO"

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]; SW,SH=prs.slide_width,prs.slide_height
def rect(s,x,y,w,h,fill,rounded=True,line=None):
    sp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,x,y,w,h)
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1)
    sp.shadow.inherit=False; return sp
def grad(sp,c1,c2):
    try:
        sp.fill.gradient(); g=sp.fill.gradient_stops; g[0].color.rgb=c1; g[0].position=0; g[1].color.rgb=c2; g[1].position=1
    except Exception: sp.fill.solid(); sp.fill.fore_color.rgb=c1
def txt(s,x,y,w,h,runs,size=9,color=INK,bold=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,italic=False,font=FT,sp_after=3,leading=1.02):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,ln in enumerate(runs if isinstance(runs,list) else [runs]):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(sp_after); p.line_spacing=leading
        if isinstance(ln,tuple):
            r=p.add_run(); r.text=ln[0]; r.font.size=Pt(size); r.font.bold=True; r.font.color.rgb=ln[2] if len(ln)>2 else color; r.font.name=font
            r2=p.add_run(); r2.text=ln[1]; r2.font.size=Pt(size); r2.font.bold=False; r2.font.color.rgb=color; r2.font.name=font; r2.font.italic=italic
        else:
            r=p.add_run(); r.text=ln; r.font.size=Pt(size); r.font.bold=bold; r.font.italic=italic; r.font.color.rgb=color; r.font.name=font
    return tb
def notes(s,text): s.notes_slide.notes_text_frame.text=text
def header(s,title,sub):
    rect(s,0,0,SW,SH,WHITE,rounded=False)
    txt(s,Inches(0.5),Inches(0.24),Inches(12.4),Inches(0.66),title,size=23,color=NAVY,bold=True,font=FTH)
    bar=rect(s,Inches(0.52),Inches(0.94),Inches(2.6),Inches(0.07),None); grad(bar,GBLUE,MAGENTA)
    txt(s,Inches(0.5),Inches(1.05),Inches(12.4),Inches(0.44),sub,size=12,color=BLUE,italic=True)
    txt(s,Inches(0.5),Inches(7.28),Inches(12.33),Inches(0.2),FOOTER,size=8,color=FOOT,align=PP_ALIGN.CENTER)

def gap_slide(title,sub,gaps,cover,prop,bnd,note,bnd_label="BOUNDARY",bnd_color=RISK,prop_label="NET PROPERTY:"):
    s=prs.slides.add_slide(BLANK); header(s,title,sub)
    CY=1.62; CH=3.9
    # LEFT: the gap (emphasis) - wider card, amber accent
    lx=Inches(0.5); lw=Inches(6.75)
    rect(s,lx,Inches(CY),lw,Inches(CH),LGRAY,line=BORDER)
    rect(s,lx,Inches(CY),Inches(0.09),Inches(CH),AMBER,rounded=False)
    rect(s,lx,Inches(CY),lw,Inches(0.46),SLATE)
    txt(s,lx,Inches(CY),lw,Inches(0.46),"THE GAP IN CURRENT TOOLS",size=11.5,color=WHITE,bold=True,font=FTH,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,lx+Inches(0.26),Inches(CY+0.58),lw-Inches(0.46),Inches(CH-0.66),
        [("- "+a+"  ",b,AMBER) for a,b in gaps],size=9.5,color=INK,sp_after=7,leading=1.04)
    # RIGHT: how we cover
    rx=Inches(7.5); rw=Inches(5.33)
    rect(s,rx,Inches(CY),rw,Inches(CH),LNAVY,line=BORDER)
    hd2=rect(s,rx,Inches(CY),rw,Inches(0.46),None); grad(hd2,NAVY,NAVY2)
    txt(s,rx,Inches(CY),rw,Inches(0.46),"HOW WE COVER IT",size=11.5,color=WHITE,bold=True,font=FTH,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,rx+Inches(0.24),Inches(CY+0.58),rw-Inches(0.46),Inches(CH-0.66),
        [(f"{i+1}  "+a+"  ",b,NAVY) for i,(a,b) in enumerate(cover)],size=9.5,color=INK,sp_after=8,leading=1.05)
    # BOUNDARY / EVIDENCE strip
    sy=5.62; ssh=0.6
    rect(s,Inches(0.5),Inches(sy),Inches(12.33),Inches(ssh),STRIP,line=BORDER)
    rect(s,Inches(0.5),Inches(sy),Inches(0.09),Inches(ssh),bnd_color,rounded=False)
    txt(s,Inches(0.72),Inches(sy),Inches(11.9),Inches(ssh),[(bnd_label+"    ",bnd,bnd_color)],
        size=9,color=SLATE,italic=True,anchor=MSO_ANCHOR.MIDDLE,leading=1.03)
    # NET PROPERTY banner
    bl=rect(s,Inches(0.5),Inches(6.36),Inches(12.33),Inches(0.62),None); grad(bl,NAVY,NAVY2)
    tb=s.shapes.add_textbox(Inches(0.75),Inches(6.36),Inches(11.9),Inches(0.62)); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=prop_label+"  "; r.font.size=Pt(10.5); r.font.bold=True; r.font.color.rgb=ORANGE; r.font.name=FTH
    r2=p.add_run(); r2.text=prop; r2.font.size=Pt(11.5); r2.font.bold=True; r2.font.color.rgb=WHITE; r2.font.name=FT
    notes(s,note); return s

# ============ OVERVIEW ============
s=prs.slides.add_slide(BLANK); header(s,"The gap at each layer, and how we close it",
    "Where current tools stop short (technically) and the mechanism that covers it. Formal-methods layers prove a property; semantic layers classify meaning.")
rows=[
 ("L1","Preemptive Network Defense","Verifiers prove packet reachability; BAS only samples a scenario library - none certify control efficacy or a true negative","Symbolic proof the config defeats every modeled technique (soundness over config + corpus)",NAVY),
 ("L2","Behavioral Twin","SIEM / UEBA threshold on magnitude; slow low-magnitude drift and drift DIRECTION go unseen","Brings the math that powers LLMs to the security layer; reads what behavior means and where it is heading, not raw numbers  (detail next slide)",GBLUE),
 ("L3","App & API Protection","WAAP alerts on deviation from a learned normal and enforces at a vendor edge - intent + air-gap unmet","Classify request intent; enforce inline in the request path with no traffic egress",NAVY2),
 ("L4","Agentic Governance","Scopes are hand-authored RBAC; learned baselines need a warm-up; supply-chain caught in scans","Persona compiled to a least-privilege MCP scope; spec is the baseline; registry enforced inline",NAVY),
 ("L5","CVE + Containment","Reachability is heuristic (no soundness); virtual patch is a signature; host and agent live on separate planes","Sound un-reachability proof; proof-carrying patch; one host + agent containment policy",NAVY2),
]
hy=1.72
txt(s,Inches(1.52),Inches(hy),Inches(5.5),Inches(0.3),"THE GAP IN CURRENT TOOLS",size=9,color=AMBER,bold=True,font=FTH)
txt(s,Inches(7.72),Inches(hy),Inches(5.0),Inches(0.3),"HOW WE COVER IT",size=9,color=NAVY,bold=True,font=FTH)
TY=2.06; RH=0.86
for i,(code,name,gap,cover,col) in enumerate(rows):
    y=TY+i*RH
    if i%2==0: rect(s,Inches(0.5),Inches(y),Inches(12.33),Inches(RH-0.06),LGRAY,line=None)
    rect(s,Inches(0.62),Inches(y+0.16),Inches(0.78),Inches(0.5),col)
    txt(s,Inches(0.62),Inches(y+0.14),Inches(0.78),Inches(0.5),code,size=12.5,color=WHITE,bold=True,font=FTH,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,Inches(1.52),Inches(y+0.05),Inches(6.0),Inches(0.25),name,size=8,color=SLATE,bold=True)
    txt(s,Inches(1.52),Inches(y+0.26),Inches(6.05),Inches(RH-0.32),gap,size=9,color=INK,leading=1.01,anchor=MSO_ANCHOR.TOP)
    txt(s,Inches(7.72),Inches(y+0.08),Inches(5.05),Inches(RH-0.16),cover,size=9.5,color=NAVY,bold=True,anchor=MSO_ANCHOR.MIDDLE,leading=1.03)
# discipline strip
dy=6.28
rect(s,Inches(0.5),Inches(dy),Inches(12.33),Inches(0.44),STRIP,line=BORDER)
rect(s,Inches(0.5),Inches(dy),Inches(0.09),Inches(0.44),RISK,rounded=False)
txt(s,Inches(0.72),Inches(dy),Inches(11.9),Inches(0.44),
    [("DISCIPLINE    ","Every guarantee is stated with its boundary (config + corpus, TCB, cohort + concept basis); probabilistic AI is kept outside any formal claim; every number comes from a primary source.",RISK)],
    size=8.5,color=SLATE,italic=True,anchor=MSO_ANCHOR.MIDDLE)
band=rect(s,Inches(0.5),Inches(6.8),Inches(12.33),Inches(0.4),None); grad(band,NAVY,NAVY2)
txt(s,Inches(0.75),Inches(6.8),Inches(11.9),Inches(0.4),
    "Formal-methods layers (L1, L5) discharge a proof over a bounded model; semantic layers (L2, L3, L4) classify intent and behavior, not surface counts.",
    size=10,color=WHITE,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
notes(s,"Gap-first framing. Every 'how we cover' is bounded to a model or analysis boundary. Verify all vendor claims against current sources.")

# ============ L1 ============
gap_slide(
 "Layer 1: control-efficacy verification",
 "Preemptive network defense - the Layer 1 innovation. Current tools test reachability or sample attacks; none prove the control config defeats the technique.",
 [("Reachability verifiers (Forward, Batfish):","SMT-based, but model the packet-forwarding graph (ACLs, routes, NAT) - they answer 'can A reach B:port', not whether the IPS + EPP + WAF policy stops a technique."),
  ("Breach-and-attack simulation (SafeBreach, AttackIQ):","executes a finite scenario library against live infra - coverage equals the library, and a blocked result is one observation, not a proof over all paths."),
  ("Scanners and CTEM (Tenable, XM Cyber, RedSeal):","enumerate CVEs and build best-effort attack graphs, then simulate - heuristic, no soundness, no true-negative."),
  ("Structural gap:","all are positive-evidence methods that find a hit; none can certify a true negative (prove no evasion path exists), which is what 'our controls defeat this technique' requires.")],
 [("Symbolic model of the controls.","NGFW / IPS / EPP / WAF / IdP rule-sets are lifted into one policy model of allow / deny / detect semantics; each ATT&CK technique becomes a state transition with pre- and post-conditions."),
  ("Solve for an evasion path.","An SMT solver searches for any technique chain the composed controls fail to block; a satisfying assignment is a concrete evasion path."),
  ("Return a proof from config alone.","UNSAT = a soundness certificate (no evasion path in the model); SAT = the counter-path and the failing control. No agents, no live traffic.")],
 "a soundness certificate that no modeled technique chain evades the controls - bounded to the config and the encoded ATT&CK corpus.",
 "Soundness holds over the modeled configuration + encoded ATT&CK corpus - you cannot prove absence of a technique you never modeled. The LLM is confined to config ingestion + human-readable remediation, outside the proof core.",
 "BOUNDARY: reachability verifiers are the nearest formal analog but model forwarding only. Any probabilistic step inside the proof voids 'provable' - keep the LLM strictly outside it.")

# ============ L2 ============
gap_slide(
 "Layer 2: the math behind LLMs, applied to behavior",
 "Behavioral Intelligence - the Layer 2 innovation. We bring the mathematics that gives LLMs their power to read meaning to the security layer - reading what behavior MEANS and where it is heading, not how big any single number is.",
 [("SIEM / UEBA (Splunk, Sentinel, Exabeam):","represent behavior as numeric aggregates and threshold on magnitude - a metric can be normal-sized while the KIND of activity has shifted; they see 'how much', never 'in what direction'."),
  ("Network anomaly (Darktrace):","scores deviation on flows and packets - magnitude of deviation, not movement in a semantic model of the entity's role behavior."),
  ("Global-baseline ML (Vectra, Gurucul):","compare each entity to a population threshold - a shift abnormal for one role but ordinary population-wide is missed."),
  ("Structural gap:","none embed behavior itself as the representation (embeddings elsewhere serve search / RAG), so slow low-magnitude drift stays under every single-window threshold.")],
 [("Unified semantic space (meaning + context).","5 behavioral zones serialized to prose, embedded to 1536-d, composed by context-adaptive attention into ONE entity vector - the embedding holds meaning and context, not counts."),
  ("A math model drives the LLM space.","Drift velocity + cosine projection onto 16 named concepts (direction: 'what is it becoming') + feature / embedding-CUSUM (persistence) -> MITRE top-3. The concept alignment IS the detector."),
  ("The whole exceeds the sum of parts.","Each signal stays below its own threshold for a patient adversary; only the fused multi-phase composite crosses the line - risk is visible only in combination."),
  ("Two detectors, cohort-relative.","Cumulative composite (discovery) + known-bad profile classifier (precision): C2-beacon, DGA, LOTL, recon fan-out, insider collection - each raw-event corroborated.")],
 "no single action looks suspicious on its own. We read a person's whole behavior together and catch the overall direction it is drifting - so a patient attacker who keeps every number in the normal range is still caught.",
 "Live DB, never hardcoded: baselines 0/4  ·  z-score 1/4 (noisy)  ·  composite 4/4 but separates only 2/4 at 10.6% FP  ·  threat-profile 4/4 at 0 FP. Acid test: USR-234 (slow APT) stays inside both drift bands.",
 "EVIDENCE: all metrics computed live from PostgreSQL :5438. USR-234 is the acid test - any 'drift catches all 4' claim is wrong. Prose serialization of z-scores + meaning before embedding is intentional design, not circular leakage.",
 bnd_label="EVIDENCE",bnd_color=STEEL,prop_label="IN PLAIN TERMS:")

# ============ L3 ============
gap_slide(
 "Layer 3: intent classification, enforced inline",
 "Application & API protection - the Layer 3 innovation. Current WAAP learns 'normal' and enforces at a vendor edge; intent and air-gapped enforcement go unmet.",
 [("Anomaly-baseline WAAP (Salt, Imperva, Akamai):","must converge a learned 'normal' per API, then alert on deviation - a warm-up requirement, and an output that is a deviation score, not the request's intent."),
  ("Edge / CDN-coupled enforcement (Akamai, Cloudflare, F5):","places the enforcement point in the vendor cloud - traffic must transit it, a non-starter for air-gapped and ATO-bounded enclaves."),
  ("Detect-only designs (Salt):","split detection from blocking - they need a downstream WAF or gateway to enforce."),
  ("Stateless per-request rules:","miss business-logic abuse that only appears across an API call sequence.")],
 [("Classify intent, not deviation.","A cross-tenant global model + a per-endpoint local model classify the request's intent - no 'normal' to converge first."),
  ("Enforce in the request path.","Block / rate-limit / deceive executed inline with no traffic leaving the boundary - deployable in air-gapped enclaves."),
  ("Score the sequence; attest the agent.","Requests judged within their session sequence (business-logic abuse); hardware Secure-Enclave attestation a cloud agent cannot virtualize.")],
 "intent is classified and enforced in the data path - no learned baseline to converge, no traffic leaving the enclave.",
 "Layer 3 innovation: keep claims at the mechanism level (dual-model, inline, sequence scoring, enclave attestation). Verify scale figures against the current datasheet; do not cite uncorroborated numbers.",
 "BOUNDARY: do not cite the deck's 117M bot-IP / '12+ telcos' figures; state architecture, not unverified scale.")

# ============ L4 ============
gap_slide(
 "Layer 4: the persona as the enforcement spec",
 "Agentic AI governance - the Layer 4 innovation. Current tools hand-author agent scopes and learn a baseline; the spec itself is not the control.",
 [("Hand-authored scopes (Entra Agent ID, Zscaler registry, Noma):","write least-privilege manually per agent - the scope does not follow from the agent's declared job, so it drifts from intent and scales poorly."),
  ("Learned-baseline monitors:","need a warm-up window to converge 'normal' before they can flag an off-spec action - blind during onboarding."),
  ("SASE / transport controls:","inspect the network, not the agent's tool-call semantics - they cannot see which MCP tool was invoked with which parameters."),
  ("MCP supply-chain (rug-pull, tool-poisoning):","is caught in a periodic scan, not at the moment of the call.")],
 [("Compile the persona to a scope.","A plain-English job description compiles into a least-privilege MCP tool-scope (tools + parameters) - the allow-list is generated from intent, not hand-written."),
  ("Make the spec the baseline.","A per-persona virtual MCP endpoint proxies the agent; the first off-spec tool call is a deterministic violation - no warm-up, no probabilistic score."),
  ("Enforce registry integrity inline.","A crypto-signed, version-pinned MCP registry blocks rug-pull / tool-poisoning at call time, not in a scan report.")],
 "the declared persona is the enforcement spec - a violation is the first tool call outside scope, decided deterministically with no learning window.",
 "Durable technical slivers: NL-to-scope compilation and spec-as-baseline (deterministic, zero warm-up). The AI-gateway category is not unique - the depth inside it is the point.",
 "BOUNDARY: do not claim the category (Palo Alto Prisma AIRS, Zscaler AI Broker exist). Have a per-call enforcement latency figure ready - inline checking cost is the obvious objection.")

# ============ L5 ============
gap_slide(
 "Layer 5: provable un-reachability + unified containment",
 "Code / CVE and containment. Current tools estimate reachability and split host from agent; neither yields a proof nor one policy.",
 [("SCA reachability (Endor, Snyk, Contrast):","uses heuristic / observed call graphs and markets 'fewer false positives' - it cannot claim zero false negatives, so it cannot justify deferring a patch."),
  ("Virtual patching (Trend Micro, Fortinet):","applies a signature block at the perimeter - not a proof the vulnerable code path is closed."),
  ("Microsegmentation (Illumio, Guardicore):","contains hosts at the network but cannot revoke a token or tool permission - an agent inside a trusted host never crosses a segment."),
  ("Agent-identity tools:","revoke tokens but enforce nothing at the network - two separate control planes for one incident.")],
 [("Prove un-reachability (soundness).","Sound call-graph analysis: UNSAT reachability = no false negatives over the boundary -> a CVE can be deferred defensibly on a mission system."),
  ("Proof-carrying virtual patch.","A reachable sink is closed by a patch that carries a proof it preserves behavior - not a WAF signature at the perimeter."),
  ("One containment policy.","Network microsegmentation + identity / token / tool-permission revocation under one verified policy, sequenced by attack path - host or agent.")],
 "provable un-reachability (soundness over the analysis boundary) and one policy that revokes a host segment or an agent's token.",
 "'Formally verified' is bounded to a named TCB (seL4-style) - the audience funded HACMS / seL4 and will probe it. Soundness (zero false negatives) is the wedge over heuristic reachability.",
 "BOUNDARY: state the verification boundary explicitly (backup slide). Never say 'reachability' or 'virtual patching' unqualified - always 'sound' / 'provable' / 'formally verified'.")

out=os.path.join(os.path.dirname(os.path.abspath(__file__)),"Layer_Differentiation_Slides.pptx")
try:
    prs.save(out); print("saved",os.path.abspath(out),"| slides:",len(prs.slides._sldIdLst))
except PermissionError:
    alt=r"C:/Users/shuklar/AppData/Local/Temp/claude/c--Users-shuklar-ClaudeCode-CyberUEBA-BehavioralIntel-Enhanced/a6818674-bde5-4462-be21-ffe7ac102bf2/scratchpad/Layer_Differentiation_Slides.pptx"
    prs.save(alt); print("LOCKED - saved fallback",alt,"| slides:",len(prs.slides._sldIdLst))
